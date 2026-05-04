# Copyright (C) 2024-2026 Koto AI. All rights reserved.
# SPDX-License-Identifier: LicenseRef-Koto-Proprietary
"""
Koto 全格式文件解析模块 — Phase 1 BFF 管线核心
支持格式: DOCX / XLSX / PPTX / PDF
每个解析函数接收文件路径或字节流,输出标准化 JSON 供前端多态渲染器消费。
"""

from __future__ import annotations

import base64
import io
import logging
import math
import mimetypes
import os
import re
import uuid
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# DOCX → Semantic HTML
# ─────────────────────────────────────────────────────────────────────────────

# ── Shared image compression settings ────────────────────────────────────────
_MAX_IMG_DIMENSION = 1200  # px — max width or height
_MAX_IMG_BYTES = 300 * 1024   # 300 KB threshold for triggering compression
_MAX_BLOB_BYTES = 15 * 1024 * 1024  # 15 MB hard limit — blobs larger than this are skipped entirely
_MAX_PPTX_BYTES = 100 * 1024 * 1024  # 100 MB PPTX size cap (likely contains embedded video)
_DOCX_PREVIEW_TARGET_PAGES = 3
_DOCX_PREVIEW_UNITS_PER_PAGE = 34
_DOCX_PREVIEW_MAX_TABLE_ROWS = 18


def _compress_image_bytes(
    img_bytes: bytes, content_type: str = "image/png"
) -> tuple[bytes, str]:
    """
    Compress image bytes if they exceed _MAX_IMG_BYTES.

    Returns ``(possibly_compressed_bytes, mime_type)`` — unchanged if Pillow is
    not installed or the image is already small enough.
    """
    if len(img_bytes) <= _MAX_IMG_BYTES:
        return img_bytes, content_type
    try:
        from PIL import Image as PILImage

        pil_img = PILImage.open(io.BytesIO(img_bytes))
        w, h = pil_img.size

        # Downscale if either dimension exceeds limit
        if w > _MAX_IMG_DIMENSION or h > _MAX_IMG_DIMENSION:
            ratio = min(_MAX_IMG_DIMENSION / w, _MAX_IMG_DIMENSION / h)
            pil_img = pil_img.resize(
                (int(w * ratio), int(h * ratio)), PILImage.LANCZOS
            )

        # Convert to RGB (JPEG doesn't support alpha)
        if pil_img.mode in ("RGBA", "P", "LA"):
            bg = PILImage.new("RGB", pil_img.size, (255, 255, 255))
            if pil_img.mode == "P":
                pil_img = pil_img.convert("RGBA")
            bg.paste(pil_img, mask=pil_img.split()[-1])
            pil_img = bg

        buf = io.BytesIO()
        pil_img.save(buf, format="JPEG", quality=82, optimize=True)
        compressed = buf.getvalue()
        logger.info(
            "[_compress_image_bytes] %dx%d → %dx%d, %.0f KB → %.0f KB",
            w, h, pil_img.size[0], pil_img.size[1],
            len(img_bytes) / 1024, len(compressed) / 1024,
        )
        return compressed, "image/jpeg"
    except ImportError:
        return img_bytes, content_type
    except Exception as exc:
        logger.debug("[_compress_image_bytes] failed (using original): %s", exc)
        return img_bytes, content_type


def _extract_table_styles(docx_path: str) -> list[dict]:
    """
    Extract table cell formatting metadata from a DOCX file using python-docx.

    Returns a list — one entry per table, positionally indexed:
        [
          {  # table 0
            (row, col): {
              "bg":      "RRGGBB" | None,   # cell background fill
              "colspan": int,               # w:gridSpan value
              "rowspan": int,               # logical rowspan (from w:vMerge analysis)
              "bold":    bool,              # header-row bold hint
            },
            ...
          },
          ...  # table 1, table 2, …
        ]

    Falls back to [] silently if python-docx or lxml is unavailable.
    """
    try:
        from docx import Document
        from lxml import etree
    except ImportError:
        return []

    WNS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    def _w(tag: str) -> str:
        return f"{{{WNS}}}{tag}"

    try:
        doc = Document(docx_path)
    except Exception as exc:
        logger.warning("[_extract_table_styles] 无法打开文件: %s", exc)
        return []

    result = []
    for tbl in doc.tables:
        tbl_data: dict = {}

        # ── Column widths from w:tblGrid/w:gridCol (twips → percentages) ──
        # tbl._tbl is the lxml element for the <w:tbl> node.
        tbl_xml = tbl._tbl
        grid_cols = tbl_xml.findall(f".//{_w('gridCol')}")
        if grid_cols:
            raw_widths = []
            for gc in grid_cols:
                w_val = gc.get(_w("w"))
                try:
                    raw_widths.append(max(1, int(w_val or 0)))
                except (ValueError, TypeError):
                    raw_widths.append(1)
            total_w = sum(raw_widths) or 1
            # Store as percentages rounded to 1 decimal; sum ≈ 100 %
            tbl_data["__col_pcts"] = [round(w / total_w * 100, 1) for w in raw_widths]

        for ri, row in enumerate(tbl.rows):
            for ci, cell in enumerate(row.cells):
                tc = cell._tc

                # ── Background colour (w:shd w:fill) ────────────────────
                shd = tc.find(f".//{_w('shd')}")
                bg = None
                if shd is not None:
                    fill = shd.get(_w("fill"))
                    if fill and fill.upper() not in ("AUTO", "FFFFFF", ""):
                        bg = fill.upper()

                # ── Colspan (w:gridSpan) ─────────────────────────────────
                grid_span_el = tc.find(f".//{_w('gridSpan')}")
                colspan = int(grid_span_el.get(_w("val"), 1)) if grid_span_el is not None else 1

                # ── Rowspan: count how many rows below start with w:vMerge (no val attr) ──
                # The first cell of a vertical merge has w:vMerge val="restart";
                # continuation cells have w:vMerge with no val.
                v_merge_el = tc.find(f".//{_w('vMerge')}")
                is_merge_start = (
                    v_merge_el is not None
                    and v_merge_el.get(_w("val"), "") == "restart"
                )
                rowspan = 1
                if is_merge_start:
                    # Count how many rows below continue this merge at the same col
                    for future_ri in range(ri + 1, len(tbl.rows)):
                        if ci >= len(tbl.rows[future_ri].cells):
                            break
                        future_tc = tbl.rows[future_ri].cells[ci]._tc
                        future_vm = future_tc.find(f".//{_w('vMerge')}")
                        if future_vm is not None and future_vm.get(_w("val"), "") == "":
                            rowspan += 1
                        else:
                            break

                # ── Skip continuation vMerge cells (they'll be expressed via rowspan) ──
                is_merge_continuation = (
                    v_merge_el is not None
                    and v_merge_el.get(_w("val"), "") == ""
                )
                if is_merge_continuation:
                    continue  # nothing to write; the HTML rowspan covers it

                tbl_data[(ri, ci)] = {
                    "bg": bg,
                    "colspan": colspan,
                    "rowspan": rowspan,
                }

        result.append(tbl_data)

    return result


def _inject_table_styles(html: str, table_styles: list[dict]) -> str:
    """
    Post-process mammoth's HTML to inject cell background colours, colspan,
    and rowspan attributes that mammoth discards.

    Only modifies cells for which _extract_table_styles() found metadata;
    leaves everything else untouched.
    """
    if not table_styles:
        return html

    try:
        from bs4 import BeautifulSoup
    except ImportError:
        return html

    soup = BeautifulSoup(html, "html.parser")
    tables = soup.find_all("table")

    for ti, tbl_tag in enumerate(tables):
        if ti >= len(table_styles):
            break
        style_map = table_styles[ti]
        if not style_map:
            continue

        # ── Inject <colgroup> for proportional column widths ──────────────────
        col_pcts: list[float] = style_map.get("__col_pcts", [])  # type: ignore[assignment]
        if col_pcts and not tbl_tag.find("colgroup"):
            colgroup = soup.new_tag("colgroup")
            for pct in col_pcts:
                col = soup.new_tag("col")
                col["style"] = f"width:{pct:.1f}%"
                colgroup.append(col)
            tbl_tag.insert(0, colgroup)

        rows = tbl_tag.find_all("tr")

        # ── Also inject width% on first-row cells (colgroup may be stripped by
        #    the editor's deserializer; inline td widths survive reliably) ──
        if col_pcts and rows:
            first_row_cells = rows[0].find_all(["td", "th"])
            for ci_col, cell_tag in enumerate(first_row_cells):
                if ci_col < len(col_pcts):
                    existing = cell_tag.get("style", "")
                    w_css = f"width:{col_pcts[ci_col]:.1f}%;"
                    # Don't double-inject
                    if "width:" not in existing:
                        cell_tag["style"] = (existing.rstrip(";") + ";" + w_css).lstrip(";")

        for ri, row_tag in enumerate(rows):
            cells = row_tag.find_all(["td", "th"])
            for ci, cell_tag in enumerate(cells):
                meta = style_map.get((ri, ci))
                if not meta:
                    continue

                # colspan / rowspan
                if meta.get("colspan", 1) > 1:
                    cell_tag["colspan"] = str(meta["colspan"])
                if meta.get("rowspan", 1) > 1:
                    cell_tag["rowspan"] = str(meta["rowspan"])

                # background colour
                bg = meta.get("bg")
                if bg:
                    existing = cell_tag.get("style", "")
                    bg_css = f"background-color:#{bg.lower()};"
                    cell_tag["style"] = (existing.rstrip(";") + ";" + bg_css).lstrip(";")

    # ── Remove entirely-empty trailing columns ─────────────────────────────
    # The editor sometimes pads rows with phantom empty cells; stripping
    # them here (Python side) avoids reconciler issues on the frontend.
    for tbl_tag in soup.find_all("table"):
        rows = tbl_tag.find_all("tr")
        if not rows:
            continue
        max_cols = max(
            (len(r.find_all(["td", "th"], recursive=False)) for r in rows),
            default=0,
        )
        if max_cols < 2:
            continue
        for ci in range(max_cols - 1, 0, -1):
            all_empty = all(
                len(cells) <= ci
                or (not cells[ci].get_text(strip=True) and not cells[ci].find("img"))
                for cells in (r.find_all(["td", "th"], recursive=False) for r in rows)
            )
            if all_empty:
                for row in rows:
                    cells = row.find_all(["td", "th"], recursive=False)
                    if ci < len(cells):
                        cells[ci].decompose()
            else:
                break

    return str(soup)


def _get_floating_image_srcs(docx_path: str) -> list[str]:
    """
    Return base64 data URIs for floating (anchor-positioned) images in a DOCX.

    mammoth only converts <wp:inline> images; <wp:anchor> images (floated to a
    fixed page position, e.g. a resume photo) are silently skipped.
    Uses namespace-aware ElementTree parsing so attribute order never matters.
    """
    import zipfile as _zipfile
    import xml.etree.ElementTree as _ET

    # Full OOXML namespace URIs — must match exactly (not prefix aliases)
    IMG_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    WP_NS    = "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"
    A_NS     = "http://schemas.openxmlformats.org/drawingml/2006/main"
    R_NS     = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    result: list[str] = []
    try:
        with _zipfile.ZipFile(docx_path, "r") as z:
            file_list = z.namelist()

            # Locate document.xml and its rels (case-insensitive for robustness)
            doc_name = next(
                (n for n in file_list if n.lower() == "word/document.xml"), None
            )
            rels_name = next(
                (n for n in file_list if n.lower() == "word/_rels/document.xml.rels"), None
            )
            if not doc_name or not rels_name:
                return result

            # ── rel_id → media path via ElementTree (handles ANY attribute order) ──
            rels_root = _ET.fromstring(z.read(rels_name))
            rel_map: dict[str, str] = {}
            for rel in rels_root:
                rid   = rel.get("Id", "")
                rtype = rel.get("Type", "")
                tgt   = rel.get("Target", "")
                if rid and IMG_TYPE in rtype and tgt:
                    rel_map[rid] = tgt if tgt.startswith("word/") else f"word/{tgt}"

            # ── Walk document.xml, separate inline vs anchor image embed IDs ──
            doc_root = _ET.fromstring(z.read(doc_name))

            inline_ids: set[str] = set()
            floating_ids: list[str] = []

            for inline in doc_root.iter(f"{{{WP_NS}}}inline"):
                for blip in inline.iter(f"{{{A_NS}}}blip"):
                    eid = blip.get(f"{{{R_NS}}}embed")
                    if eid:
                        inline_ids.add(eid)

            for anchor in doc_root.iter(f"{{{WP_NS}}}anchor"):
                for blip in anchor.iter(f"{{{A_NS}}}blip"):
                    eid = blip.get(f"{{{R_NS}}}embed")
                    if eid:
                        floating_ids.append(eid)

            seen: set[str] = set()
            for rel_id in floating_ids:
                if rel_id in inline_ids or rel_id in seen:
                    continue
                seen.add(rel_id)
                media_path = rel_map.get(rel_id)
                if not media_path:
                    logger.debug(
                        "[floating_imgs] no media path for rel_id=%s; rel_map keys=%s",
                        rel_id, list(rel_map)[:8],
                    )
                    continue
                try:
                    img_bytes = z.read(media_path)
                    mime = mimetypes.guess_type(media_path)[0] or "image/png"
                    img_bytes, mime = _compress_image_bytes(img_bytes, mime)
                    b64 = base64.b64encode(img_bytes).decode("ascii")
                    result.append(f"data:{mime};base64,{b64}")
                except Exception as e:
                    logger.debug("[floating_imgs] failed to read %s: %s", media_path, e)
    except Exception as exc:
        logger.warning("[_get_floating_image_srcs] %s", exc)
    return result


def _unwrap_layout_tables(html: str) -> str:
    """
    Unwrap single-column layout tables into flat content sections.

    Many Word documents (especially resumes) wrap *all* content in a 1-column
    table purely for layout control.  When rendered in the editor these tables
    show ugly borders and break natural page flow.

    Criteria for "layout table":
      - Every row has exactly 1 cell (colspan is OK — still single visual column)
      - No rowspan > 1 (vertical merges indicate a multi-row data table)
      - Table has >= 3 rows (prevents stripping intentional single-row tables)
      - No header row (<th>)

    Each cell's inner HTML is extracted and wrapped in
    ``<div class="koto-layout-section">...</div>`` to preserve grouping.
    Empty cells (spacer rows) become ``<div class="koto-layout-spacer"></div>``.
    """
    try:
        from bs4 import BeautifulSoup, Tag
    except ImportError:
        return html

    soup = BeautifulSoup(html, "html.parser")
    tables = soup.find_all("table")
    changed = False

    for tbl in tables:
        rows = tbl.find_all("tr", recursive=False)
        if len(rows) < 3:
            continue  # too small — probably a real data table

        # Check: is every row single-cell with no <th>?
        is_layout = True
        for row in rows:
            cells = row.find_all(["td", "th"], recursive=False)
            if len(cells) != 1:
                is_layout = False
                break
            cell = cells[0]
            if cell.name == "th":
                is_layout = False
                break
            # colspan is fine (cell just spans full width) — but rowspan > 1
            # means a complex multi-row structure, not pure layout.
            rs = int(cell.get("rowspan", 1) or 1)
            if rs > 1:
                is_layout = False
                break

        if not is_layout:
            continue

        # Skip tables that have visible borders (inline style from rich renderer)
        tbl_style = tbl.get("style", "")
        if "border:" in tbl_style and "border:none" not in tbl_style:
            continue

        fragments = []
        for row in rows:
            cell = row.find(["td", "th"], recursive=False)
            if not cell:
                continue
            inner = cell.decode_contents().strip()
            if not inner or inner in ("<br/>", "<br>", "<br >"):
                div = soup.new_tag("div")
                div["class"] = "koto-layout-spacer"
                fragments.append(div)
            else:
                div = soup.new_tag("div")
                div["class"] = "koto-layout-section"
                # Preserve cell background colour as section background
                cell_style = cell.get("style", "")
                if "background" in cell_style:
                    div["style"] = cell_style
                div.append(BeautifulSoup(inner, "html.parser"))
                fragments.append(div)

        # Replace the <table> in-place
        for frag in reversed(fragments):
            tbl.insert_after(frag)
        tbl.decompose()
        changed = True

    return str(soup) if changed else html


def _deduplicate_images(html: str) -> str:
    """
    Remove duplicate inline base64 images from HTML.

    Compares the first 200 chars of each data URI.  If two <img> tags share
    the same base64 prefix (i.e. are the same image), only the first is kept.
    """
    # Quick check — fewer than 2 images means nothing to deduplicate
    if html.count("<img") < 2:
        return html

    seen_prefixes: set[str] = set()
    result_parts: list[str] = []
    pos = 0

    for m in re.finditer(r"<img\s[^>]*>", html, re.IGNORECASE):
        img_tag = m.group()
        src_m = re.search(r'src="(data:[^"]{0,200})', img_tag)
        if src_m:
            prefix = src_m.group(1)
            if prefix in seen_prefixes:
                # Duplicate — skip this <img> tag
                result_parts.append(html[pos:m.start()])
                pos = m.end()
                continue
            seen_prefixes.add(prefix)

        result_parts.append(html[pos:m.end()])
        pos = m.end()

    result_parts.append(html[pos:])
    return "".join(result_parts)


def _docx_to_rich_html(
    file_path: str,
    *,
    progressive_preview: bool = False,
) -> tuple[str, list[dict], dict[str, Any]]:
    """
    将 DOCX 转换为保留完整 Word 格式的内联样式 HTML。

    使用 python-docx 直接读取 OOXML，提取:
      - 段落对齐/缩进/行距/段前段后
      - 标题级别 (Heading 1-6 / 标题 1-6) → <h1>-<h6>
      - 行内字体名称、大小、颜色、粗体、斜体、下划线、删除线、上下标
      - 超链接 → <a href="...">
      - 有序/无序列表 → <ol>/<ul> 按 numId 分组
      - 表格: 列宽、合并单元格、边框、单元格背景色、垂直对齐
      - 内联图片 (wp:inline) 和浮动图片 (wp:anchor) → base64 <img>
      - 页眉/页脚 → <div class="koto-header">/<div class="koto-footer">
    """
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_ROW_HEIGHT_RULE, WD_CELL_VERTICAL_ALIGNMENT
        from docx.shared import Pt, RGBColor
    except ImportError:
        raise RuntimeError("python-docx 未安装，请执行: pip install python-docx")

    # ── Heading style name → tag map ────────────────────────────────────────
    _HEADING_MAP: dict[str, str] = {
        "heading 1": "h1", "heading 2": "h2", "heading 3": "h3",
        "heading 4": "h4", "heading 5": "h5", "heading 6": "h6",
        # Some documents use style IDs/names without spaces.
        "heading1": "h1", "heading2": "h2", "heading3": "h3",
        "heading4": "h4", "heading5": "h5", "heading6": "h6",
        "标题 1": "h1", "标题 2": "h2", "标题 3": "h3",
        "标题 4": "h4", "标题 5": "h5", "标题 6": "h6",
        # WPS / Word Chinese no-space variants (e.g. "标题1" not "标题 1")
        "标题1": "h1", "标题2": "h2", "标题3": "h3",
        "标题4": "h4", "标题5": "h5", "标题6": "h6",
        # Common alternative Chinese heading names
        "一级标题": "h1", "二级标题": "h2", "三级标题": "h3",
        "标题": "h1",
        # subheading
        "subheading 1": "h2", "subheading 2": "h3",
    }

    _HEADING_TYPOGRAPHY_FALLBACKS: dict[str, dict[str, str]] = {
        "h1": {"font-size": "16.0pt", "font-weight": "bold"},
        "h2": {"font-size": "14.0pt", "font-weight": "bold"},
        "h3": {"font-size": "12.0pt", "font-weight": "bold"},
        "h4": {"font-size": "11.0pt", "font-weight": "bold"},
        "h5": {"font-size": "10.5pt", "font-weight": "bold"},
        "h6": {"font-size": "10.5pt", "font-weight": "bold"},
    }

    _ALIGN_MAP = {
        WD_ALIGN_PARAGRAPH.LEFT: "left",
        WD_ALIGN_PARAGRAPH.CENTER: "center",
        WD_ALIGN_PARAGRAPH.RIGHT: "right",
        WD_ALIGN_PARAGRAPH.JUSTIFY: "justify",
        WD_ALIGN_PARAGRAPH.DISTRIBUTE: "justify",
    }

    _CN_FONT_MAP_P: dict[str, str] = {
        "黑体": "SimHei", "宋体": "SimSun", "楷体": "KaiTi",
        "仿宋": "FangSong", "微软雅黑": "Microsoft YaHei",
        "华文中宋": "STZhongsong", "华文宋体": "STSong",
        "华文黑体": "STHeiti", "华文楷体": "STKaiti",
        "华文仿宋": "STFangsong",
    }

    _EMPTY_STYLE_DEFAULTS: dict[str, Any] = {
        "style_name": "",
        "style_id": "",
        "text_align": None,
        "space_before": None,
        "space_after": None,
        "line_height": None,
        "font_size": None,
        "font_family": None,
        "font_weight": None,
        "font_style": None,
        "outline_tag": None,
    }
    _para_style_ref_cache: dict[int, Any] = {}
    _style_defaults_cache: dict[int, dict[str, Any]] = {}

    # ── Helpers ─────────────────────────────────────────────────────────────
    def _twips_to_pt(twips: int | None) -> float | None:
        """Convert EMU twips (1/20 pt) to pt."""
        if twips is None:
            return None
        return round(twips / 20, 2)

    def _emu_to_px(emu: int) -> int:
        """Convert EMU to pixels (96 dpi)."""
        return max(1, round(emu / 914400 * 96))

    def _rgb_to_hex(rgb: RGBColor | None) -> str | None:
        if rgb is None:
            return None
        try:
            return "#{:02X}{:02X}{:02X}".format(int(rgb.red), int(rgb.green), int(rgb.blue))
        except Exception:
            return None

    def _norm_style_key(val: str | None) -> str:
        """Normalise style key text for robust cross-language matching."""
        if not val:
            return ""
        return str(val).strip().lower().replace(" ", "").replace("_", "")

    def _is_heading_style_key(val: str | None) -> str | None:
        """Return h1..h6 if the style key represents a heading style."""
        if not val:
            return None
        _raw = str(val).strip().lower()
        if _raw in _HEADING_MAP:
            return _HEADING_MAP[_raw]
        _k = _norm_style_key(val)
        # heading1 / heading2 / ...
        if _k.startswith("heading") and len(_k) > 7 and _k[7:].isdigit():
            _lv = int(_k[7:])
            if 1 <= _lv <= 6:
                return f"h{_lv}"
        # 标题1 / 标题2 / ...
        if _k.startswith("标题"):
            _digits = "".join(ch for ch in _k if ch.isdigit())
            if _digits:
                _lv = int(_digits)
                if 1 <= _lv <= 6:
                    return f"h{_lv}"
        return None

    def _extract_toc_level_from_style(val: str | None) -> str:
        """Extract TOC level from style-like text, defaulting to "1"."""
        if not val:
            return "1"
        _s = str(val)
        for _ch in _s:
            if _ch.isdigit():
                return _ch
        return "1"

    def _apply_heading_typography_fallback(tag: str, css: dict[str, str], inner_html: str) -> None:
        """Keep semantic headings readable when DOCX outline metadata has no font props."""
        _fallback = _HEADING_TYPOGRAPHY_FALLBACKS.get((tag or "").lower())
        if not _fallback:
            return

        _inner_lower = inner_html.lower()
        if "font-size" not in css and "font-size:" not in _inner_lower:
            css["font-size"] = _fallback["font-size"]
        if "font-weight" not in css and "font-weight:" not in _inner_lower:
            css["font-weight"] = _fallback["font-weight"]

    def _p_elem_has_toc_anchor(p_el) -> bool:
        """True when paragraph has internal TOC hyperlink targets like #_Toc*."""
        if p_el is None:
            return False
        try:
            for _hl in p_el.findall(qn("w:hyperlink")):
                _a = (_hl.get(qn("w:anchor")) or "").lower()
                if _a.startswith("_toc"):
                    return True
        except Exception:
            return False
        return False

    def _p_elem_has_toc_field(p_el) -> bool:
        """True when paragraph carries a TOC/PAGEREF field code."""
        if p_el is None:
            return False
        try:
            _instr = " ".join(
                (_node.text or "")
                for _node in p_el.iter(qn("w:instrText"))
            )
            _instr_norm = re.sub(r"\s+", " ", _instr).strip().upper()
            if not _instr_norm:
                return False
            return (
                " TOC " in f" {_instr_norm} "
                or (" PAGEREF " in f" {_instr_norm} " and "_TOC" in _instr_norm)
            )
        except Exception:
            return False

    def _p_elem_text_content(p_el) -> str:
        """Extract visible paragraph text from runs for structural heuristics."""
        if p_el is None:
            return ""
        try:
            _text = "".join((_node.text or "") for _node in p_el.iter(qn("w:t")))
            return re.sub(r"\s+", " ", _text).strip()
        except Exception:
            return ""

    def _p_elem_looks_like_toc_line(p_el) -> bool:
        """Heuristic for field-updated TOC lines that only keep visible text."""
        _text = _p_elem_text_content(p_el)
        if not _text or len(_text) > 160:
            return False
        return re.match(r"^.+?\d{1,4}$", _text) is not None

    def _line_spacing_to_css(ls, ls_rule) -> str | None:
        from docx.enum.text import WD_LINE_SPACING

        if ls is None:
            return None

        _FIXED_MULT = {
            WD_LINE_SPACING.SINGLE: 1.0,
            WD_LINE_SPACING.ONE_POINT_FIVE: 1.5,
            WD_LINE_SPACING.DOUBLE: 2.0,
        }
        if ls_rule in _FIXED_MULT:
            return f"{_FIXED_MULT[ls_rule]}"
        if ls_rule in (WD_LINE_SPACING.MULTIPLE, None):
            try:
                return f"{round(float(ls), 4)}"
            except (TypeError, ValueError):
                pass
            try:
                return f"{round(ls.pt, 2)}pt"
            except Exception:
                return None
        try:
            return f"{round(ls.pt, 2)}pt"
        except Exception:
            return None

    def _resolve_para_style_ref(para):
        _cache_key = id(getattr(para, "_element", para))
        if _cache_key in _para_style_ref_cache:
            return _para_style_ref_cache[_cache_key]
        try:
            _style = para.style if para.style else None
        except Exception:
            _style = None
        _para_style_ref_cache[_cache_key] = _style
        return _style

    def _read_on_off_prop(el) -> bool | None:
        """Return OOXML on/off state for elements like <w:b/> or <w:i w:val="0"/>."""
        if el is None:
            return None
        try:
            _val = el.get(qn("w:val"))
        except Exception:
            _val = None
        if _val is None or str(_val).strip() == "":
            return True
        _norm = str(_val).strip().lower()
        if _norm in ("0", "false", "off", "no"):
            return False
        return True

    def _read_rpr_font_props(rpr_el) -> dict[str, Any]:
        """Extract default font props from a run-properties XML element."""
        _props: dict[str, Any] = {
            "font_size": None,
            "font_family": None,
            "font_weight": None,
            "font_weight_set": False,
            "font_style": None,
            "font_style_set": False,
        }
        if rpr_el is None:
            return _props

        try:
            _sz = rpr_el.find(qn("w:sz"))
            if _sz is None:
                _sz = rpr_el.find(qn("w:szCs"))
            if _sz is not None:
                _raw = _sz.get(qn("w:val")) or _sz.get("val") or ""
                if _raw:
                    _props["font_size"] = round(int(_raw) / 2, 1)
        except Exception:
            pass

        try:
            _rFonts = rpr_el.find(qn("w:rFonts"))
            if _rFonts is not None:
                for _key in ("w:eastAsia", "w:ascii", "w:hAnsi", "w:cs"):
                    _font_name = _rFonts.get(qn(_key), "") or ""
                    if _font_name:
                        _props["font_family"] = _CN_FONT_MAP_P.get(_font_name, _font_name)
                        break
        except Exception:
            pass

        try:
            _bold_state = _read_on_off_prop(rpr_el.find(qn("w:b")))
            if _bold_state is None:
                _bold_state = _read_on_off_prop(rpr_el.find(qn("w:bCs")))
            if _bold_state is not None:
                _props["font_weight_set"] = True
                _props["font_weight"] = "bold" if _bold_state else None
        except Exception:
            pass

        try:
            _italic_state = _read_on_off_prop(rpr_el.find(qn("w:i")))
            if _italic_state is None:
                _italic_state = _read_on_off_prop(rpr_el.find(qn("w:iCs")))
            if _italic_state is not None:
                _props["font_style_set"] = True
                _props["font_style"] = "italic" if _italic_state else None
        except Exception:
            pass

        return _props

    def _resolve_style_defaults(style_ref) -> dict[str, Any]:
        if style_ref is None or not hasattr(style_ref, "_element"):
            return _EMPTY_STYLE_DEFAULTS

        _cache_key = id(style_ref._element)
        _cached = _style_defaults_cache.get(_cache_key)
        if _cached is not None:
            return _cached

        _resolved = dict(_EMPTY_STYLE_DEFAULTS)
        try:
            _resolved["style_name"] = style_ref.name or ""
        except Exception:
            pass
        try:
            _resolved["style_id"] = style_ref.style_id or ""
        except Exception:
            pass

        if _resolved["outline_tag"] is None:
            try:
                _st_pPr = style_ref._element.find(qn("w:pPr"))
                if _st_pPr is not None:
                    _olvl = _st_pPr.find(qn("w:outlineLvl"))
                    if _olvl is not None:
                        _val = _olvl.get(qn("w:val"))
                        if _val is not None:
                            _lvl = int(_val)
                            if 0 <= _lvl <= 5:
                                _resolved["outline_tag"] = f"h{_lvl + 1}"
                            elif _lvl == 9:
                                _resolved["outline_tag"] = None
            except Exception:
                pass

        _style = style_ref
        _font_weight_locked = False
        _font_style_locked = False
        while _style:
            try:
                _spf = _style.paragraph_format
            except Exception:
                _spf = None

            try:
                _rPr = _style._element.find(qn("w:rPr"))
            except Exception:
                _rPr = None
            _rpr_props = _read_rpr_font_props(_rPr)

            if _spf is not None:
                if _resolved["text_align"] is None:
                    try:
                        _align = _ALIGN_MAP.get(_spf.alignment)
                        if _align:
                            _resolved["text_align"] = _align
                    except Exception:
                        pass
                if _resolved["space_before"] is None:
                    try:
                        _val = _spf.space_before
                        if _val is not None:
                            _resolved["space_before"] = _twips_to_pt(_val.twips)
                    except Exception:
                        pass
                if _resolved["space_after"] is None:
                    try:
                        _val = _spf.space_after
                        if _val is not None:
                            _resolved["space_after"] = _twips_to_pt(_val.twips)
                    except Exception:
                        pass
                if _resolved["line_height"] is None:
                    try:
                        _line_height = _line_spacing_to_css(
                            _spf.line_spacing,
                            _spf.line_spacing_rule,
                        )
                        if _line_height is not None:
                            _resolved["line_height"] = _line_height
                    except Exception:
                        pass

            if _resolved["font_size"] is None:
                try:
                    _fsize = _style.font.size
                    if _fsize is not None:
                        _resolved["font_size"] = round(_fsize.pt, 1)
                except Exception:
                    pass
                if _resolved["font_size"] is None and _rpr_props["font_size"] is not None:
                    _resolved["font_size"] = _rpr_props["font_size"]

            if not _font_weight_locked:
                _bold_state = None
                try:
                    if _style.font.bold is True:
                        _bold_state = True
                    elif _style.font.bold is False:
                        _bold_state = False
                except Exception:
                    pass
                if _bold_state is None and _rpr_props["font_weight_set"]:
                    _bold_state = _rpr_props["font_weight"] == "bold"
                if _bold_state is not None:
                    _font_weight_locked = True
                    _resolved["font_weight"] = "bold" if _bold_state else None

            if not _font_style_locked:
                _italic_state = None
                try:
                    if _style.font.italic is True:
                        _italic_state = True
                    elif _style.font.italic is False:
                        _italic_state = False
                except Exception:
                    pass
                if _italic_state is None and _rpr_props["font_style_set"]:
                    _italic_state = _rpr_props["font_style"] == "italic"
                if _italic_state is not None:
                    _font_style_locked = True
                    _resolved["font_style"] = "italic" if _italic_state else None

            if _resolved["font_family"] is None:
                try:
                    if _rpr_props["font_family"]:
                        _resolved["font_family"] = _rpr_props["font_family"]
                    if _resolved["font_family"] is None:
                        _fn = _style.font.name
                        if _fn:
                            _resolved["font_family"] = _CN_FONT_MAP_P.get(_fn, _fn)
                except Exception:
                    pass

            if _resolved["text_align"] is None:
                try:
                    _pPr = _style._element.find(qn("w:pPr"))
                    if _pPr is not None:
                        _jc = _pPr.find(qn("w:jc"))
                        if _jc is not None:
                            _align = _jc.get(qn("w:val"), "")
                            if _align in ("left", "center", "right", "justify"):
                                _resolved["text_align"] = _align
                            elif _align == "distribute":
                                _resolved["text_align"] = "justify"
                except Exception:
                    pass

            if (
                _resolved["space_before"] is not None
                and _resolved["space_after"] is not None
                and _resolved["line_height"] is not None
                and _resolved["font_size"] is not None
                and _resolved["font_family"] is not None
                and _resolved["outline_tag"] is not None
            ):
                break

            try:
                _style = _style.base_style
            except Exception:
                _style = None

        _style_defaults_cache[_cache_key] = _resolved
        return _resolved

    def _detect_toc_info(para=None, p_el=None, style_ref=None) -> tuple[bool, str]:
        """Detect whether a paragraph is a TOC entry and return (is_toc, level)."""
        level = "1"
        style_candidates: list[str] = []

        if para is not None:
            _style_defaults = _resolve_style_defaults(
                style_ref if style_ref is not None else _resolve_para_style_ref(para)
            )
            try:
                style_candidates.append(_style_defaults.get("style_name") or "")
            except Exception:
                pass
            try:
                style_candidates.append(_style_defaults.get("style_id") or "")
            except Exception:
                pass
            try:
                if p_el is None:
                    p_el = para._element
            except Exception:
                pass

        if p_el is not None:
            try:
                _pPr = p_el.find(qn("w:pPr"))
                if _pPr is not None:
                    _pSt = _pPr.find(qn("w:pStyle"))
                    if _pSt is not None:
                        style_candidates.append(_pSt.get(qn("w:val")) or "")
            except Exception:
                pass

        _has_toc_anchor = _p_elem_has_toc_anchor(p_el)
        _has_toc_field = _p_elem_has_toc_field(p_el)
        _has_tab = False
        if p_el is not None:
            try:
                _has_tab = p_el.find(".//" + qn("w:tab")) is not None
            except Exception:
                _has_tab = False
        _looks_like_toc_line = _p_elem_looks_like_toc_line(p_el)
        _has_toc_signal = _has_toc_anchor or _has_toc_field or _has_tab or _looks_like_toc_line

        for _sv in style_candidates:
            _sv_norm = _norm_style_key(_sv)
            if not _sv_norm:
                continue
            if "toc" in _sv_norm or "目录" in _sv_norm or "tableofcontents" in _sv_norm:
                if not _has_toc_signal:
                    continue
                level = _extract_toc_level_from_style(_sv)
                return True, level

        # Fallback for custom style names: internal TOC anchor + tab run pattern.
        if p_el is not None:
            try:
                if _has_tab and (_has_toc_anchor or _has_toc_field or _looks_like_toc_line):
                    return True, "1"
            except Exception:
                pass

        return False, level

    def _color_from_xml(elem) -> str | None:
        """Extract theme/RGB color from a w:color XML element."""
        if elem is None:
            return None
        val = elem.get(qn("w:val"))
        if val and val.lower() not in ("auto", "none", ""):
            if len(val) == 6:
                return "#" + val.upper()
        return None

    def _encode_image_part(part) -> str:
        """Encode an image relationship part to base64 data URI."""
        try:
            img_bytes = part.blob
            content_type = part.content_type or "image/png"
            img_bytes, content_type = _compress_image_bytes(img_bytes, content_type)
            b64 = base64.b64encode(img_bytes).decode("ascii")
            return f"data:{content_type};base64,{b64}"
        except Exception as exc:
            logger.debug("[RichHtml] image encode failed: %s", exc)
            return ""

    def _inline_img_html(drawing_elem, doc) -> str:
        """Render wp:inline drawing as <img> HTML."""
        try:
            blipFill = drawing_elem.find(".//" + qn("a:blip"))
            if blipFill is None:
                return ""
            rId = blipFill.get(qn("r:embed"))
            if not rId:
                return ""
            # Navigate to the part that owns this drawing
            # drawing_elem is inside a <w:r> inside document body
            part = doc.part.related_parts.get(rId)
            if part is None:
                return ""
            src = _encode_image_part(part)
            if not src:
                return ""
            # Get display size — width only; height omitted so the browser
            # preserves the image's natural aspect ratio.
            ext = drawing_elem.find(".//" + qn("wp:extent"))
            style = ""
            if ext is not None:
                cx = int(ext.get("cx", 0))
                if cx:
                    w_px = _emu_to_px(cx)
                    style = f' style="width:{w_px}px;max-width:100%"'
            return f'<img src="{src}" alt=""{style} />'
        except Exception as exc:
            logger.debug("[RichHtml] inline img failed: %s", exc)
            return ""

    def _anchor_img_html(drawing_elem, doc) -> str:
        """Render wp:anchor (floating) drawing as <img> HTML.

        Reads OOXML positioning attributes to reproduce Word's layout:
          - wp:positionH / wp:positionV  → CSS float or margin
          - wp:wrapSquare/Tight          → float (left / right)
          - wp:wrapTopAndBottom          → display:block, centred
          - wp:wrapNone                  → inline-block (behind/in-front)
        Falls back to float:right when no positioning info is present.
        """
        try:
            blipFill = drawing_elem.find(".//" + qn("a:blip"))
            if blipFill is None:
                return ""
            rId = blipFill.get(qn("r:embed"))
            if not rId:
                return ""
            part = doc.part.related_parts.get(rId)
            if part is None:
                return ""
            src = _encode_image_part(part)
            if not src:
                return ""

            # ── Size (width only; height omitted to preserve natural ratio) ──
            ext = drawing_elem.find(".//" + qn("wp:extent"))
            size_style = ""
            if ext is not None:
                cx = int(ext.get("cx", 0))
                if cx:
                    w_px = _emu_to_px(cx)
                    size_style = f"width:{w_px}px;"

            # ── Horizontal alignment ───────────────────────────────────────
            # Read <wp:positionH> / <wp:positionV> to determine layout intent.
            pos_h = drawing_elem.find(qn("wp:positionH"))
            h_align = None
            if pos_h is not None:
                align_el = pos_h.find(qn("wp:align"))
                if align_el is not None and align_el.text:
                    h_align = align_el.text.strip().lower()  # left | center | right | inside | outside

            # ── Wrap mode ─────────────────────────────────────────────────
            # Determine CSS positioning from wrap type.
            has_wrap_square  = drawing_elem.find(qn("wp:wrapSquare"))  is not None
            has_wrap_tight   = drawing_elem.find(qn("wp:wrapTight"))   is not None
            has_wrap_none    = drawing_elem.find(qn("wp:wrapNone"))    is not None
            has_wrap_tb      = drawing_elem.find(qn("wp:wrapTopAndBottom")) is not None

            # ── Build CSS ─────────────────────────────────────────────────
            if has_wrap_none:
                # Behind/in-front of text: render as inline-block centered
                pos_style = "display:inline-block;vertical-align:middle;margin:4px 8px;"
            elif has_wrap_tb:
                # Top-and-bottom wrap: block element, centred on the page
                pos_style = "display:block;margin:10px auto;"
            elif has_wrap_square or has_wrap_tight:
                # Float based on declared horizontal alignment
                if h_align in ("left", "inside"):
                    pos_style = "float:left;margin:0 14px 10px 0;"
                else:
                    # right / outside / unspecified → float right (Word default)
                    pos_style = "float:right;margin:0 0 10px 14px;"
            elif h_align == "center":
                pos_style = "display:block;margin:10px auto;"
            elif h_align == "left":
                pos_style = "float:left;margin:0 14px 10px 0;"
            else:
                # No meaningful positioning data — use float:right as default
                pos_style = "float:right;margin:0 0 10px 14px;"

            full_style = pos_style + size_style + "max-width:100%;"
            return f'<img src="{src}" alt="" style="{full_style}" />'
        except Exception as exc:
            logger.debug("[RichHtml] anchor img failed: %s", exc)
            return ""

    def _run_html(run, doc) -> str:
        """Convert a single w:r run to HTML span(s), handling images and hyperlinks."""
        # Check for drawing (inline image)
        drawing = run._element.find(qn("w:drawing"))
        if drawing is not None:
            inline = drawing.find(qn("wp:inline"))
            anchor = drawing.find(qn("wp:anchor"))
            if inline is not None:
                return _inline_img_html(inline, doc)
            if anchor is not None:
                return _anchor_img_html(anchor, doc)

        text = run.text or ""

        # Resolve enclosing paragraph element once; needed for robust TOC/tab detection.
        _p_el = run._element.getparent()
        while _p_el is not None:
            _local = _p_el.tag.split("}")[-1] if "}" in _p_el.tag else _p_el.tag
            if _local == "p":
                break
            _p_el = _p_el.getparent()

        _is_toc_run, _ = _detect_toc_info(p_el=_p_el)

        _has_tab_elem = run._element.find(qn("w:tab")) is not None
        _has_tab_char = "\t" in text
        _tab_html = ""
        if _has_tab_elem or _has_tab_char:
            _tab_html = '<span class="koto-toc-tab"></span>' if _is_toc_run else ('\u00a0' * 6)

        if not text:
            # Tab character → render appropriately based on context.
            # In TOC entries, render the tab as a CSS flex-spacer span so
            # the page number sits at the right margin (Word dot-leader effect).
            if _has_tab_elem:
                return _tab_html
            return ""

        def _esc(_t: str) -> str:
            return (_t.replace("&", "&amp;")
                      .replace("<", "&lt;")
                      .replace(">", "&gt;")
                      .replace('"', "&quot;"))

        tab_segments: list[str] | None = None
        if _has_tab_char:
            # Keep tab placeholders as direct siblings between text fragments,
            # otherwise flex layout (dot leaders + right page number) breaks.
            tab_segments = [_esc(seg) for seg in text.split("\t")]
            text = ""
        else:
            text = _esc(text)

        # Collect inline styles
        styles: list[str] = []
        f = run.font

        # Font family
        # python-docx's run.font.name only returns the Latin/ASCII font.
        # East-Asian font names are stored in w:rFonts w:eastAsia and must
        # be read directly from the XML.  We also map common Chinese font
        # names to their CSS-equivalent ASCII names where needed.
        _CN_FONT_MAP: dict[str, str] = {
            "黑体": "SimHei",
            "宋体": "SimSun",
            "楷体": "KaiTi",
            "仿宋": "FangSong",
            "微软雅黑": "Microsoft YaHei",
            "华文中宋": "STZhongsong",
            "华文宋体": "STSong",
            "华文黑体": "STHeiti",
            "华文楷体": "STKaiti",
            "华文仿宋": "STFangsong",
            "方正书宋": "FZShuSong-Z01",
            "方正黑体": "FZHei-B01",
        }
        fn = f.name
        ea_font: str = ""
        try:
            _rPr = run._element.find(qn("w:rPr"))
            if _rPr is not None:
                _rFonts = _rPr.find(qn("w:rFonts"))
                if _rFonts is not None:
                    ea_font = _rFonts.get(qn("w:eastAsia"), "") or \
                               _rFonts.get(qn("w:cs"), "") or ""
        except Exception:
            pass
        # Normalise Chinese names to CSS names
        fn = _CN_FONT_MAP.get(fn, fn) if fn else fn
        ea_font = _CN_FONT_MAP.get(ea_font, ea_font) if ea_font else ea_font
        # Build font-family stack: prefer Latin+EastAsian together
        if fn and ea_font and fn != ea_font:
            styles.append(f"font-family:'{fn}','{ea_font}'")
        elif ea_font:
            styles.append(f"font-family:'{ea_font}'")
        elif fn:
            styles.append(f"font-family:'{fn}'")

        # Font size (Pt object → float pt value)
        # Skip for TOC runs — CSS normalizes sizes per-level to avoid inconsistency
        if not _is_toc_run:
            try:
                fsize = f.size
                if fsize is not None:
                    pt_val = round(fsize.pt, 1)
                    styles.append(f"font-size:{pt_val}pt")
            except Exception:
                pass

        # Color
        try:
            color_xml = run._element.find(qn("w:rPr") + "/" + qn("w:color"))
            if color_xml is None:
                # try direct rPr child
                rpr = run._element.find(qn("w:rPr"))
                if rpr is not None:
                    color_xml = rpr.find(qn("w:color"))
            hex_color = _color_from_xml(color_xml)
            if hex_color:
                styles.append(f"color:{hex_color}")
        except Exception:
            try:
                rgb = f.color.rgb
                if rgb:
                    styles.append(f"color:#{rgb}")
            except Exception:
                pass

        # Bold — implement OOXML toggle property semantics to prevent phantom bold.
        # When a paragraph's pPr/rPr has <w:b/> set AND a run also has <w:b/>,
        # the run's bold TOGGLES the paragraph-level bold OFF (net result: not bold).
        # python-docx's run.font.bold misses this toggle and always returns True.
        # Skip bold entirely for TOC runs — CSS normalizes per-level.
        decorations: list[str] = []
        _run_rpr = run._element.find(qn("w:rPr"))
        _b_el = _run_rpr.find(qn("w:b")) if _run_rpr is not None else None
        if _b_el is not None and not _is_toc_run:
            _b_val = (_b_el.get(qn("w:val")) or "1").lower()
            _run_bold_on = _b_val not in ("0", "false", "off")
            # Check paragraph pPr/rPr/w:b as the toggle base
            _para_ppr_bold = False
            if _p_el is not None:
                _p_pPr = _p_el.find(qn("w:pPr"))
                if _p_pPr is not None:
                    _p_rpr2 = _p_pPr.find(qn("w:rPr"))
                    if _p_rpr2 is not None:
                        _p_b = _p_rpr2.find(qn("w:b"))
                        if _p_b is not None:
                            _p_b_val = (_p_b.get(qn("w:val")) or "1").lower()
                            _para_ppr_bold = _p_b_val not in ("0", "false", "off")
            # Toggle: run bold XOR paragraph pPr bold
            if _run_bold_on and not _para_ppr_bold:
                styles.append("font-weight:bold")
        if f.italic:
            styles.append("font-style:italic")
        if f.underline and not _is_toc_run:
            decorations.append("underline")
        if f.strike:
            decorations.append("line-through")
        if decorations:
            styles.append(f"text-decoration:{' '.join(decorations)}")

        # Superscript / Subscript
        if f.superscript:
            styles.append("vertical-align:super;font-size:smaller")
        elif f.subscript:
            styles.append("vertical-align:sub;font-size:smaller")

        # Highlight color (background)
        try:
            highlight_xml = None
            rpr2 = run._element.find(qn("w:rPr"))
            if rpr2 is not None:
                highlight_xml = rpr2.find(qn("w:highlight"))
            if highlight_xml is not None:
                hval = highlight_xml.get(qn("w:val"), "")
                _HL_MAP = {
                    "yellow": "#FFFF00", "green": "#00FF00", "cyan": "#00FFFF",
                    "magenta": "#FF00FF", "blue": "#0000FF", "red": "#FF0000",
                    "darkBlue": "#00008B", "darkCyan": "#008B8B",
                    "darkGreen": "#006400", "darkMagenta": "#8B008B",
                    "darkRed": "#8B0000", "darkYellow": "#808000",
                    "darkGray": "#A9A9A9", "lightGray": "#D3D3D3",
                }
                if hval in _HL_MAP:
                    styles.append(f"background-color:{_HL_MAP[hval]}")
        except Exception:
            pass

        if tab_segments is not None:
            if styles:
                style_str = ";".join(styles)
                _styled = [f'<span style="{style_str}">{seg}</span>' if seg else '' for seg in tab_segments]
                return _tab_html.join(_styled)
            return _tab_html.join(tab_segments)

        if styles:
            style_str = ";".join(styles)
            return f'<span style="{style_str}">{text}</span>'
        return text

    def _para_style(para, style_ref=None) -> dict[str, str]:
        """Extract paragraph CSS properties as a dict."""
        css: dict[str, str] = {}
        pf = para.paragraph_format
        _p_pPr = None
        _para_rpr_props = _read_rpr_font_props(None)
        try:
            _p_pPr = para._element.find(qn("w:pPr"))
            if _p_pPr is not None:
                _para_rpr_props = _read_rpr_font_props(_p_pPr.find(qn("w:rPr")))
        except Exception:
            _p_pPr = None
            _para_rpr_props = _read_rpr_font_props(None)
        _style_defaults = _resolve_style_defaults(
            style_ref if style_ref is not None else _resolve_para_style_ref(para)
        )

        # Alignment — read from paragraph XML directly, then immediate style XML only.
        # Avoid walking the full style inheritance chain: a distant ancestor
        # style (e.g. a title/heading base style) with "center" would otherwise
        # propagate to all body text paragraphs that happen to inherit from it.
        _align_val = None
        try:
            if _p_pPr is not None:
                _p_jc = _p_pPr.find(qn("w:jc"))
                if _p_jc is not None:
                    _align_val = _p_jc.get(qn("w:val"), "")
        except Exception:
            pass
        if not _align_val:
            # Check immediate style only (no inheritance walk)
            _sr = style_ref if style_ref is not None else _resolve_para_style_ref(para)
            if _sr is not None and hasattr(_sr, "_element"):
                try:
                    _st_pPr = _sr._element.find(qn("w:pPr"))
                    if _st_pPr is not None:
                        _st_jc = _st_pPr.find(qn("w:jc"))
                        if _st_jc is not None:
                            _align_val = _st_jc.get(qn("w:val"), "")
                except Exception:
                    pass
        if _align_val in ("center", "right", "justify", "distribute"):
            css["text-align"] = "justify" if _align_val == "distribute" else _align_val
        # "left" is the CSS/browser default; no need to explicitly set it.

        # Space before/after (twips → pt)
        # First try paragraph-level XML; fall back to named style
        def _resolve_spacing(attr_name):
            """Return pt value or None, walking paragraph → style → base_style chain."""
            try:
                val = getattr(pf, attr_name)
                if val is not None:
                    return _twips_to_pt(val.twips)
            except Exception:
                pass
            return _style_defaults.get(attr_name)

        sb = _resolve_spacing("space_before")
        sa = _resolve_spacing("space_after")
        if sb is not None:
            css["margin-top"] = f"{sb}pt"
        if sa is not None:
            css["margin-bottom"] = f"{sa}pt"

        # Line spacing — paragraph XML first, then named style
        _line_height = None
        try:
            _line_height = _line_spacing_to_css(pf.line_spacing, pf.line_spacing_rule)
        except Exception:
            pass
        if _line_height is None:
            _line_height = _style_defaults.get("line_height")
        if _line_height is not None:
            css["line-height"] = _line_height

        # ── Font-size from paragraph style chain ─────────────────────
        # When runs have no explicit font-size, they inherit from the
        # paragraph's block-level style.  Setting font-size on the <p>/<h>
        # tag ensures correct inheritance without bloating every <span>.
        # Also prevents browser-default heading sizes (h1=2em etc.) from
        # overriding the DOCX style's intended size.
        _fs = _para_rpr_props.get("font_size")
        if _fs is None:
            _fs = _style_defaults.get("font_size")
        if _fs is not None:
            css["font-size"] = f"{_fs}pt"

        # ── Font-family from paragraph style chain ────────────────────
        _ff = _para_rpr_props.get("font_family") or _style_defaults.get("font_family")
        if _ff:
            css["font-family"] = f"'{_ff}'"

        if _para_rpr_props.get("font_weight_set"):
            _fw = _para_rpr_props.get("font_weight")
        else:
            # Only apply bold from the paragraph's *direct* named style, not from
            # ancestor styles.  Walking the full inheritance chain causes all body
            # paragraphs to appear bold when a base style (Normal / 正文) has
            # w:b set — a common mishap in WPS / Word document templates.
            _fw = None
            _direct_sr = style_ref if style_ref is not None else _resolve_para_style_ref(para)
            if _direct_sr is not None:
                try:
                    if _direct_sr.font.bold is True:
                        _fw = "bold"
                except Exception:
                    pass
        if _fw:
            css["font-weight"] = _fw

        if _para_rpr_props.get("font_style_set"):
            _fi = _para_rpr_props.get("font_style")
        else:
            _fi = _style_defaults.get("font_style")
        if _fi:
            css["font-style"] = _fi

        # First-line indent
        try:
            fli = pf.first_line_indent
            if fli is not None and fli != 0:
                pt_val = _twips_to_pt(fli.twips)
                if pt_val is not None:
                    css["text-indent"] = f"{pt_val}pt"
        except Exception:
            pass

        # Left indent
        try:
            li = pf.left_indent
            if li is not None and li != 0:
                pt_val = _twips_to_pt(li.twips)
                if pt_val is not None:
                    css["padding-left"] = f"{pt_val}pt"
        except Exception:
            pass

        return css

    def _para_html(para, doc, tag: str = "p", style_ref=None) -> str:
        """Render a paragraph to an HTML block element."""
        style_ref = style_ref if style_ref is not None else _resolve_para_style_ref(para)
        css = _para_style(para, style_ref=style_ref)

        # ── Scan for bookmark IDs (used as TOC link targets) ─────────
        bm_id = None
        for _bm in para._element.findall(qn("w:bookmarkStart")):
            _name = _bm.get(qn("w:name"), "")
            if _name and not _name.startswith("_GoBack"):
                bm_id = _name
                break

        # ── Detect TOC style → CSS class for front-end styling ───────
        toc_class = ""
        _is_toc, _toc_level = _detect_toc_info(
            para=para,
            p_el=para._element,
            style_ref=style_ref,
        )
        if _is_toc:
            toc_class = f"koto-toc-{_toc_level}"

        # Collect run HTML
        inner_parts: list[str] = []

        # Iterate XML children to handle hyperlinks inline
        for child in para._element:
            tag_name = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag_name == "hyperlink":
                # Extract URL
                rId = child.get(qn("r:id"))
                url = ""
                if rId:
                    try:
                        url = para.part.relationships[rId].target_ref
                    except Exception:
                        pass
                if not url:
                    # w:anchor hyperlink
                    anchor_val = child.get(qn("w:anchor"), "")
                    if anchor_val:
                        url = "#" + anchor_val
                link_inner = ""
                from docx.text.run import Run
                for r_elem in child.findall(qn("w:r")):
                    run_obj = Run(r_elem, para)
                    link_inner += _run_html(run_obj, doc)
                if link_inner:
                    esc_url = url.replace('"', "&quot;")
                    # Internal links (#anchor) stay in the editor; external get _blank
                    target_attr = '' if url.startswith('#') else ' target="_blank"'
                    if toc_class:
                        # TOC links need their own flex layout so koto-toc-tab spacer
                        # stretches between the entry text and the right-aligned page number.
                        link_style = 'display:flex;align-items:baseline;flex:1;min-width:0;color:#1155CC;'
                    else:
                        link_style = 'color:#1155CC;text-decoration:underline;'
                    inner_parts.append(
                        f'<a href="{esc_url}"{target_attr} '
                        f'style="{link_style}">'
                        f'{link_inner}</a>'
                    )
            elif tag_name == "r":
                from docx.text.run import Run
                run_obj = Run(child, para)
                inner_parts.append(_run_html(run_obj, doc))

        inner = "".join(inner_parts)
        if not inner.strip():
            inner = "<br/>"

        _apply_heading_typography_fallback(tag, css, inner)

        style_str = ";".join(f"{k}:{v}" for k, v in css.items())
        # For TOC entries: strip per-paragraph font-size, font-weight, and
        # line-height — the global CSS normalizes these per level so inconsistent
        # Word style inheritance doesn't make entries look different from each other.
        # Keep only font-family and spacing. Then inject display:flex.
        if toc_class:
            # Rebuild style without the properties CSS will normalize
            _skip = {"font-size", "font-weight", "line-height", "margin-top",
                     "margin-bottom", "display", "align-items"}
            toc_parts = [p for p in css.items() if p[0] not in _skip]
            style_str = "display:flex;align-items:baseline;" + ";".join(f"{k}:{v}" for k, v in toc_parts)
        style_attr = f' style="{style_str}"' if style_str else ""
        id_attr = f' id="{bm_id}"' if bm_id else ""
        class_attr = f' class="{toc_class}"' if toc_class else ""
        return f"<{tag}{id_attr}{class_attr}{style_attr}>{inner}</{tag}>"

    def _table_has_visible_borders(tbl_elem) -> bool:
        """Check if a table has any non-nil visible border in w:tblBorders or its table style."""
        tblPr = tbl_elem.find(qn("w:tblPr"))
        if tblPr is None:
            return False
        # Check explicit tblBorders first
        tblBorders = tblPr.find(qn("w:tblBorders"))
        if tblBorders is not None:
            border_tags = ["w:top", "w:left", "w:bottom", "w:right",
                           "w:insideH", "w:insideV"]
            for bt in border_tags:
                b_elem = tblBorders.find(qn(bt))
                if b_elem is not None:
                    val = b_elem.get(qn("w:val"), "none")
                    if val not in ("none", "nil", ""):
                        return True
        # Also check the referenced table style for border definitions
        tblStyle_el = tblPr.find(qn("w:tblStyle"))
        if tblStyle_el is not None:
            # We can't access doc.styles here, but we'll trust _get_tblStyle_border_defaults
            # to resolve this in _table_html. Return True conservatively so CSS is not suppressed.
            return True
        return False

    def _border_elem_to_css(b_el) -> str:
        """Convert a single w:top/w:left/etc. border element to a CSS border value.
        Returns 'none' if explicitly disabled, or a string like '0.5pt solid #000'.
        Returns '' (empty) if the element is missing (meaning: no explicit override).
        """
        if b_el is None:
            return ""
        val = b_el.get(qn("w:val"), "none")
        if val in ("none", "nil", ""):
            return "none"
        sz = b_el.get(qn("w:sz"), "4")
        clr = b_el.get(qn("w:color"), "auto")
        try:
            sz_pt = round(int(sz) / 8, 2)
        except (TypeError, ValueError):
            sz_pt = 0.5
        clr_css = f"#{clr}" if (clr and clr.lower() != "auto") else "#000"
        return f"{sz_pt}pt solid {clr_css}"

    def _get_table_border_defaults(tbl_elem) -> dict:
        """Extract table-level border defaults from w:tblBorders.

        Returns a dict with keys: top, bottom, left, right, insideH, insideV.
        Each value is a CSS border string like '0.5pt solid #000', 'none',
        or None when the element is absent.
        """
        defaults = {"top": None, "bottom": None, "left": None, "right": None,
                    "insideH": None, "insideV": None}
        tblPr = tbl_elem.find(qn("w:tblPr"))
        if tblPr is None:
            return defaults
        tblBorders = tblPr.find(qn("w:tblBorders"))
        if tblBorders is None:
            return defaults
        for key, wtag in (("top", "w:top"), ("bottom", "w:bottom"),
                          ("left", "w:left"), ("right", "w:right"),
                          ("insideH", "w:insideH"), ("insideV", "w:insideV")):
            b_el = tblBorders.find(qn(wtag))
            css = _border_elem_to_css(b_el)
            if css != "":   # only store when element actually exists
                defaults[key] = css
        return defaults

    def _get_tblStyle_border_defaults(tbl_elem, doc) -> dict:
        """Look up the table's referenced style (w:tblStyle) and return its
        tblBorders as a border-defaults dict (same format as _get_table_border_defaults).
        Returns all-None dict if no style, no borders, or any error.
        """
        defaults = {"top": None, "bottom": None, "left": None, "right": None,
                    "insideH": None, "insideV": None}
        try:
            tblPr = tbl_elem.find(qn("w:tblPr"))
            if tblPr is None:
                return defaults
            tblStyle_el = tblPr.find(qn("w:tblStyle"))
            if tblStyle_el is None:
                return defaults
            style_id = tblStyle_el.get(qn("w:val"))
            if not style_id or doc is None:
                return defaults
            # Find the matching style object in the document
            target_style = None
            for st in doc.styles:
                if st.style_id == style_id:
                    target_style = st
                    break
            if target_style is None:
                return defaults
            # Read w:tblBorders from the style's XML element
            style_tblBorders = target_style._element.find(".//" + qn("w:tblBorders"))
            if style_tblBorders is None:
                return defaults
            for key, wtag in (("top", "w:top"), ("bottom", "w:bottom"),
                              ("left", "w:left"), ("right", "w:right"),
                              ("insideH", "w:insideH"), ("insideV", "w:insideV")):
                b_el = style_tblBorders.find(qn(wtag))
                css = _border_elem_to_css(b_el)
                if css != "":
                    defaults[key] = css
        except Exception:
            pass
        return defaults

    def _table_html(tbl, doc, *, max_rows: int | None = None) -> str:
        """Render a docx Table to an HTML <table>."""
        tbl_elem = tbl._tbl

        # Detect border style (check style-level too)
        has_borders = _table_has_visible_borders(tbl_elem)

        # Collect column widths from <w:tblGrid>
        tblGrid = tbl_elem.find(qn("w:tblGrid"))
        col_widths_twips: list[int] = []
        if tblGrid is not None:
            for gc in tblGrid.findall(qn("w:gridCol")):
                w_val = gc.get(qn("w:w"))
                try:
                    col_widths_twips.append(int(w_val))
                except (TypeError, ValueError):
                    col_widths_twips.append(0)

        total_w = sum(col_widths_twips) or 1

        # Extract table-level border defaults for per-cell inheritance
        tbl_border_defs = _get_table_border_defaults(tbl_elem)
        # If the table has no explicit tblBorders, inherit from the referenced table style
        if all(v is None for v in tbl_border_defs.values()):
            style_borders = _get_tblStyle_border_defaults(tbl_elem, doc)
            if any(v is not None for v in style_borders.values()):
                tbl_border_defs = style_borders
        all_row_elems = tbl_elem.findall(qn("w:tr"))
        table_is_truncated = bool(max_rows and len(all_row_elems) > max_rows)
        render_row_elems = all_row_elems[:max_rows] if table_is_truncated else all_row_elems
        num_rows = len(render_row_elems)
        num_cols_grid = len(col_widths_twips) or 1

        # Pre-compute pixel widths per grid column (for ProseMirror colwidth).
        # ProseMirror expects absolute pixel values.  1 twip = 1/20 pt; 1 pt = 96/72 px → px = twips/15.
        col_widths_px: list[int] = [max(1, round(t / 15)) for t in col_widths_twips] if col_widths_twips else []

        # Table-level default cell margins (OOXML defaults: 0pt vert, 5.4pt horiz)
        _WORD_DEFAULT_CELL_MAR = {"top": 0.0, "bottom": 0.0, "left": 5.4, "right": 5.4}
        _tbl_cell_mar: dict[str, float] = dict(_WORD_DEFAULT_CELL_MAR)
        _tblPr = tbl_elem.find(qn("w:tblPr"))
        if _tblPr is not None:
            _tblCellMar = _tblPr.find(qn("w:tblCellMar"))
            if _tblCellMar is not None:
                for _side in ("top", "bottom", "left", "right"):
                    _el = _tblCellMar.find(qn(f"w:{_side}"))
                    if _el is None and _side == "left":
                        _el = _tblCellMar.find(qn("w:start"))
                    if _el is None and _side == "right":
                        _el = _tblCellMar.find(qn("w:end"))
                    if _el is not None:
                        try:
                            _tbl_cell_mar[_side] = round(int(_el.get(qn("w:w"), 0)) / 20, 2)
                        except (TypeError, ValueError):
                            pass

        parts: list[str] = ["<table"]
        # border-collapse:collapse is always needed so cell borders merge cleanly.
        # Don't set a table-level border — all borders are expressed per-cell.
        parts.append(' class="koto-docx-table" style="border-collapse:collapse;width:100%;">')

        # Emit <colgroup>/<col> for visual percentage widths (non-TipTap renderers).
        # TipTap strips these, so we also emit data-colwidth on each <td> below.
        if col_widths_twips:
            parts.append("<colgroup>")
            for w_twips in col_widths_twips:
                pct = round(w_twips / total_w * 100, 2)
                parts.append(f'<col style="width:{pct}%">')
            parts.append("</colgroup>")

        # Track merged cells: set of (row_idx, col_idx) that are "occupied"
        # We use raw XML iteration to avoid python-docx's grid-filling behaviour:
        # row.cells repeats the same _Cell object for each grid column it spans,
        # causing duplicate <td> entries for horizontally merged cells.
        # row_elem.findall(qn("w:tc")) returns exactly one element per cell definition.
        for ri, row_elem in enumerate(render_row_elems):
            row_h_rule = None
            row_h = None
            try:
                trPr = row_elem.find(qn("w:trPr"))
                if trPr is not None:
                    trH = trPr.find(qn("w:trHeight"))
                    if trH is not None:
                        row_h = trH.get(qn("w:val"))
                        row_h_rule = trH.get(qn("w:hRule"), "auto")
            except Exception:
                pass

            row_style = ""
            if row_h and row_h_rule == "exact":
                try:
                    h_pt = _twips_to_pt(int(row_h))
                    row_style = f' style="height:{h_pt}pt"'
                except Exception:
                    pass

            parts.append(f"<tr{row_style}>")
            tc_elems = row_elem.findall(qn("w:tc"))
            for ci, cell_elem in enumerate(tc_elems):

                # Skip vertically merged continuation cells
                tcPr = cell_elem.find(qn("w:tcPr"))
                vmerge = tcPr.find(qn("w:vMerge")) if tcPr is not None else None
                if vmerge is not None and vmerge.get(qn("w:val")) != "restart":
                    continue

                # colspan / rowspan
                colspan = 1
                rowspan = 1
                cell_style_parts: list[str] = []

                if tcPr is not None:
                    # colspan via gridSpan
                    gs = tcPr.find(qn("w:gridSpan"))
                    if gs is not None:
                        try:
                            colspan = int(gs.get(qn("w:val"), 1))
                        except (TypeError, ValueError):
                            colspan = 1

                    # rowspan: count how many rows below have vMerge (continue)
                    try:
                        col_idx_in_grid = 0
                        # Find this cell's column index in the grid
                        for prev_tc in tc_elems:
                            if prev_tc is cell_elem:
                                break
                            prev_gs_elem = prev_tc.find(qn("w:tcPr") + "/" + qn("w:gridSpan"))
                            if prev_gs_elem is not None:
                                try:
                                    col_idx_in_grid += int(prev_gs_elem.get(qn("w:val"), 1))
                                except (TypeError, ValueError):
                                    col_idx_in_grid += 1
                            else:
                                col_idx_in_grid += 1

                        # Walk subsequent rows looking for vMerge continuation
                        for next_row_elem in render_row_elems[ri + 1:]:
                            next_cells = next_row_elem.findall(qn("w:tc"))
                            target_ci = 0
                            found_continue = False
                            for nc in next_cells:
                                if target_ci == col_idx_in_grid:
                                    nc_vmerge = nc.find(qn("w:tcPr") + "/" + qn("w:vMerge"))
                                    if nc_vmerge is not None and nc_vmerge.get(qn("w:val")) != "restart":
                                        rowspan += 1
                                        found_continue = True
                                    break
                                nc_gs = nc.find(qn("w:tcPr") + "/" + qn("w:gridSpan"))
                                try:
                                    target_ci += int(nc_gs.get(qn("w:val"), 1)) if nc_gs is not None else 1
                                except (TypeError, ValueError):
                                    target_ci += 1
                            if not found_continue:
                                break
                    except Exception:
                        pass

                    # Background color
                    shd = tcPr.find(qn("w:shd"))
                    if shd is not None:
                        fill = shd.get(qn("w:fill"), "")
                        if fill and fill.upper() not in ("AUTO", "FFFFFF", ""):
                            cell_style_parts.append(f"background-color:#{fill}")

                    # Cell width (convert dxa twips → percentage of table width)
                    tcW = tcPr.find(qn("w:tcW"))
                    if tcW is not None:
                        w_val = tcW.get(qn("w:w"))
                        w_type = tcW.get(qn("w:type"), "dxa")
                        try:
                            if w_type == "pct" and w_val:
                                pct = round(int(w_val) / 5000 * 100, 2)
                                cell_style_parts.append(f"width:{pct}%")
                            elif w_type == "dxa" and w_val and total_w > 1:
                                pct = round(int(w_val) / total_w * 100, 2)
                                cell_style_parts.append(f"width:{pct}%")
                        except (TypeError, ValueError):
                            pass

                    # Vertical alignment
                    vAlign = tcPr.find(qn("w:vAlign"))
                    if vAlign is not None:
                        v_val = vAlign.get(qn("w:val"), "")
                        _VA_MAP = {"top": "top", "center": "middle", "bottom": "bottom"}
                        va_css = _VA_MAP.get(v_val)
                        if va_css:
                            cell_style_parts.append(f"vertical-align:{va_css}")

                # Cell borders — compute effective border per side by combining
                # cell-level overrides with table-level defaults.
                # Mirrors Word's border resolution order:
                #   1. Explicit tcBorders override → use as-is
                #   2. Outer edges (first/last row/col) → table top/bottom/left/right default
                #   3. Inner edges → table insideH / insideV default
                #   4. No definition anywhere → emit 'none' to suppress CSS fallback
                tcBorders = tcPr.find(qn("w:tcBorders")) if tcPr is not None else None

                # Determine this cell's grid column index for outer-edge detection
                try:
                    _cell_grid_col = 0
                    for _prev in tc_elems:
                        if _prev is cell_elem:
                            break
                        _prev_gs = _prev.find(qn("w:tcPr") + "/" + qn("w:gridSpan"))
                        _cell_grid_col += int(_prev_gs.get(qn("w:val"), 1)) if _prev_gs is not None else 1
                    _cell_grid_end_col = _cell_grid_col + colspan - 1
                except Exception:
                    _cell_grid_col = ci
                    _cell_grid_end_col = ci

                _is_first_row = (ri == 0)
                _is_last_row  = (ri == num_rows - 1)
                _is_first_col = (_cell_grid_col == 0)
                _is_last_col  = (_cell_grid_end_col >= num_cols_grid - 1)

                _side_map = {
                    "top":    ("w:top",    "insideH" if not _is_first_row else "top"),
                    "bottom": ("w:bottom", "insideH" if not _is_last_row  else "bottom"),
                    "left":   ("w:left",   "insideV" if not _is_first_col else "left"),
                    "right":  ("w:right",  "insideV" if not _is_last_col  else "right"),
                }
                _border_css_values: list[str] = []
                for side, (wtag, tbl_def_key) in _side_map.items():
                    # 1. Cell-level explicit override
                    cell_b = None
                    if tcBorders is not None:
                        cell_b = _border_elem_to_css(tcBorders.find(qn(wtag)))
                    # 2. Table-level default (fallback when no cell border or sub-element absent)
                    if cell_b is None or cell_b == "":
                        cell_b = tbl_border_defs.get(tbl_def_key)
                    # 3. Emit (including 'none' to suppress CSS fallback border)
                    if cell_b:
                        _border_css_values.append(cell_b)
                        cell_style_parts.append(f"border-{side}:{cell_b}")
                    else:
                        _border_css_values.append("none")
                        cell_style_parts.append(f"border-{side}:none")

                # Cell padding — per-cell <w:tcMar> overrides table-level <w:tblCellMar>,
                # which itself overrides the OOXML default (0pt top/bottom, 5.4pt left/right).
                _cell_pad = dict(_tbl_cell_mar)
                _tcMar = tcPr.find(qn("w:tcMar")) if tcPr is not None else None
                if _tcMar is not None:
                    for _side in ("top", "bottom", "left", "right"):
                        _el = _tcMar.find(qn(f"w:{_side}"))
                        if _el is None and _side == "left":
                            _el = _tcMar.find(qn("w:start"))
                        if _el is None and _side == "right":
                            _el = _tcMar.find(qn("w:end"))
                        if _el is not None:
                            try:
                                _cell_pad[_side] = round(int(_el.get(qn("w:w"), 0)) / 20, 2)
                            except (TypeError, ValueError):
                                pass
                _pt = _cell_pad["top"]; _pb = _cell_pad["bottom"]
                _pl = _cell_pad["left"]; _pr = _cell_pad["right"]
                if _pt == _pb == 0 and _pl == _pr:
                    cell_style_parts.append(f"padding:0 {_pr}pt")
                elif _pt == _pb and _pl == _pr:
                    cell_style_parts.append(f"padding:{_pt}pt {_pr}pt")
                else:
                    cell_style_parts.append(f"padding:{_pt}pt {_pr}pt {_pb}pt {_pl}pt")

                cell_style = ";".join(cell_style_parts)
                td_attrs = f' style="{cell_style}"'
                if _border_css_values and all(_val.strip().lower() == "none" for _val in _border_css_values):
                    td_attrs += ' data-koto-borderless-cell="true"'
                if colspan > 1:
                    td_attrs += f' colspan="{colspan}"'
                if rowspan > 1:
                    td_attrs += f' rowspan="{rowspan}"'

                # Emit data-colwidth for ProseMirror table column resize.
                # colwidth is an array of pixel values, one per grid column the cell spans.
                if col_widths_px:
                    try:
                        _cw_start = 0
                        for _prev in tc_elems:
                            if _prev is cell_elem:
                                break
                            _prev_gs = _prev.find(qn("w:tcPr") + "/" + qn("w:gridSpan"))
                            _cw_start += int(_prev_gs.get(qn("w:val"), 1)) if _prev_gs is not None else 1
                        _cw_vals = col_widths_px[_cw_start:_cw_start + colspan]
                        if _cw_vals:
                            td_attrs += f' data-colwidth="{",".join(str(v) for v in _cw_vals)}"'
                    except Exception:
                        pass

                # Inner cell content: paragraphs (iterate raw XML to avoid
                # table-in-cell issues with python-docx's Paragraph() objects)
                cell_inner: list[str] = []
                for p_elem in cell_elem.findall(qn("w:p")):
                    from docx.text.paragraph import Paragraph as _Paragraph
                    cell_para = _Paragraph(p_elem, doc)
                    cell_inner.append(_para_html(cell_para, doc, "p"))
                # Ensure every cell has at least one child so ProseMirror
                # normalisation does not create phantom empty cells.
                if not cell_inner:
                    cell_inner = ["<p></p>"]
                parts.append(f"<td{td_attrs}>{''.join(cell_inner)}</td>")

            parts.append("</tr>")
        if table_is_truncated:
            remaining_rows = max(0, len(all_row_elems) - len(render_row_elems))
            parts.append(
                "<tr data-koto-preview-more=\"true\">"
                f"<td colspan=\"{num_cols_grid}\" "
                'style="padding:6pt 5.4pt;border-top:1px dashed #aeb7c6;'
                'border-bottom:none;border-left:none;border-right:none;'
                'background:#f6f8fb;color:#5f6b7a;font-style:italic;">'
                f"表格其余 {remaining_rows} 行正在后台加载…"
                "</td></tr>"
            )
        parts.append("</table>")
        return "".join(parts)

    def _section_html(section_obj, doc, cls: str) -> str:
        """Render a header or footer to HTML.

        Handles:
        - Multi-paragraph headers/footers
        - Tab-based left/center/right alignment converted to flex layout
        - Page-number (PAGE) fields rendered as placeholder spans
        - Paragraphs that contain only field chars but no literal text
        """
        try:
            paras = section_obj.paragraphs
            if not paras:
                return ""

            def _has_content(p) -> bool:
                """True if the paragraph has visible text OR field chars."""
                if p.text.strip():
                    return True
                xml = p._element.xml if hasattr(p._element, "xml") else ""
                return ("fldChar" in xml or "instrText" in xml)

            def _inject_page_field(p_html: str, p) -> str:
                """Replace runs that are PAGE fields with a styled placeholder."""
                try:
                    instr = " ".join(
                        e.text or ""
                        for e in p._element.iter("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}instrText")
                    ).strip()
                    if "PAGE" in instr.upper():
                        # Insert page number span before closing </p>
                        if p_html.endswith("</p>"):
                            inner = p_html[:-4]
                            p_html = inner + '<span class="koto-hdr-page-num" style="font-size:inherit;color:inherit">1</span></p>'
                except Exception:
                    pass
                return p_html

            def _tabs_to_flex(p_html: str) -> str:
                """Convert tab-based header/footer to a flex row.

                Word headers commonly use tabs for left/center/right alignment.
                We detect <span> elements that contain only tab characters and
                use them as separators to create a flex layout.
                """
                import re as _re
                # Check if it contains tab characters at all
                if '\t' not in p_html:
                    return p_html
                m = _re.match(r'(<p[^>]*>)(.*)(</p>)$', p_html, _re.DOTALL)
                if not m:
                    return p_html
                open_tag, inner, close_tag = m.group(1), m.group(2), m.group(3)
                # Replace <span ...>\t</span> with a special marker
                _TAB_MARKER = '\x00TAB\x00'
                # Match spans containing only a tab (with optional trailing whitespace)
                inner = _re.sub(r'<span[^>]*>\s*\t\s*</span>', _TAB_MARKER, inner)
                # Also handle bare tab characters between spans
                inner = inner.replace('\t', _TAB_MARKER)
                parts = inner.split(_TAB_MARKER)
                parts = [p.strip() for p in parts]
                if len(parts) < 2:
                    return p_html
                # Build flex container
                flex_style = 'display:flex;justify-content:space-between;align-items:baseline;width:100%'
                if 'style="' in open_tag:
                    open_tag = open_tag.replace('style="', f'style="{flex_style};', 1)
                else:
                    open_tag = open_tag[:-1] + f' style="{flex_style}">'
                spans = []
                for i, part in enumerate(parts):
                    if not part:
                        part = '&nbsp;'
                    align = 'left' if i == 0 else ('right' if i == len(parts) - 1 else 'center')
                    spans.append(f'<span class="koto-hdr-col" style="text-align:{align};flex-shrink:0">{part}</span>')
                return open_tag + ''.join(spans) + close_tag

            texts = []
            for p in paras:
                if not _has_content(p):
                    continue
                p_html = _para_html(p, doc, "p")
                p_html = _inject_page_field(p_html, p)
                # _run_html replaces w:tab with 6 × &nbsp; for non-TOC runs,
                # but we need actual \t characters for _tabs_to_flex detection.
                # Restore tab markers from the raw XML w:tab elements.
                _nbsp6 = '\u00a0' * 6
                p_html = p_html.replace(_nbsp6, '\t')
                # Inject class
                if p_html.startswith("<p>"):
                    p_html = f'<p class="{cls}">' + p_html[3:]
                elif p_html.startswith("<p "):
                    p_html = f'<p class="{cls}" ' + p_html[3:]
                # Convert tab-based alignment to flex layout
                p_html = _tabs_to_flex(p_html)
                texts.append(p_html)
            if not texts:
                return ""
            return "".join(texts)
        except Exception:
            return ""

    def _detect_outline_level(para, p_el, style_ref=None) -> str | None:
        """Detect heading level from w:outlineLvl in paragraph or its style.

        Word uses outlineLvl (0-based) as the authoritative heading level.
        Many Chinese documents (especially WPS) set outlineLvl on custom
        styles rather than using the standard "Heading 1" style names.

        Returns "h1"-"h6" if a valid outline level is found, else None.
        """
        # 1) Check paragraph-level pPr/outlineLvl (most authoritative)
        try:
            pPr = p_el.find(qn("w:pPr"))
            if pPr is not None:
                olvl = pPr.find(qn("w:outlineLvl"))
                if olvl is not None:
                    val = olvl.get(qn("w:val"))
                    if val is not None:
                        lvl = int(val)
                        if lvl == 9:
                            return None  # 9 means "body text" (no heading)
                        if 0 <= lvl <= 5:
                            # Guard against WPS/Word documents that write
                            # outlineLvl on long body-text paragraphs for
                            # navigation-pane purposes.  Real headings are
                            # always concise; skip very long paragraphs.
                            try:
                                para_text = para.text.strip() if para is not None else ""
                                if len(para_text) > 20:
                                    return None
                            except Exception:
                                pass
                            return f"h{lvl + 1}"
        except Exception:
            pass

        # 2) Style-based fallback: walk the basedOn chain to find a recognised
        #    heading style ancestor.  This is the Word-native approach:
        #    a paragraph is a structural heading only if its style (or one of its
        #    base styles) is a named heading style ("Heading N" / "\u6807\u9898N",
        #    etc.).  We do NOT trust outlineLvl written on arbitrary body-text
        #    styles, which WPS inserts for its own navigation-pane purposes.
        _sr = style_ref if style_ref is not None else _resolve_para_style_ref(para)
        _visited_ids: set[int] = set()
        _style_iter = _sr
        while _style_iter is not None:
            _eid = id(getattr(_style_iter, "_element", None))
            if _eid in _visited_ids:
                break
            _visited_ids.add(_eid)
            try:
                _sname = getattr(_style_iter, "name", "") or ""
                _sid2  = getattr(_style_iter, "style_id", "") or ""
                _h = _is_heading_style_key(_sname) or _is_heading_style_key(_sid2)
                if _h:
                    return _h
            except Exception:
                pass
            try:
                _style_iter = _style_iter.base_style
            except Exception:
                break
        return None

    # ── List tracking ────────────────────────────────────────────────────────
    # We process the document body as a flat list of block elements.
    # When we encounter paragraphs with numPr (list items), we group
    # consecutive same-numId items into <ul>/<ol>.

    from docx.oxml.ns import qn as _qn  # already imported as qn above

    doc = Document(file_path)
    body_parts: list[str] = []
    preview_units_limit = (
        _DOCX_PREVIEW_TARGET_PAGES * _DOCX_PREVIEW_UNITS_PER_PAGE
        if progressive_preview
        else None
    )
    preview_units_used = 0
    preview_truncated = False

    # NOTE: Header no longer prepended to body — rendered per-page by NodeView.

    # ── Iterate body elements ─────────────────────────────────────────────
    # Use document.element.body children to interleave paragraphs and tables
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    # Buffer for list grouping
    current_list_id: int | None = None
    current_list_tag: str = "ul"
    current_list_items: list[str] = []

    def _record_preview_units(units: int) -> None:
        nonlocal preview_units_used
        if preview_units_limit is not None:
            preview_units_used += max(1, units)

    def _would_exceed_preview(units: int) -> bool:
        if preview_units_limit is None:
            return False
        if not body_parts:
            return False
        return (preview_units_used + max(1, units)) > preview_units_limit

    def _estimate_paragraph_units(para) -> int:
        text = (getattr(para, "text", "") or "").strip()
        if not text:
            lines = 1
        else:
            lines = max(1, math.ceil(len(text) / 42))
        try:
            xml = para._element.xml if hasattr(para, "_element") else ""
            if "w:drawing" in xml or "pic:pic" in xml:
                lines += 6
        except Exception:
            pass
        return max(1, min(18, lines + (1 if len(text) > 90 else 0)))

    def _estimate_table_units(tbl_el) -> int:
        row_count = len(tbl_el.findall(qn("w:tr")))
        cell_count = len(tbl_el.findall(".//" + qn("w:tc")))
        return max(6, min(220, row_count * 2 + math.ceil(cell_count / 24)))

    def _flush_list() -> bool:
        nonlocal current_list_id, current_list_tag, current_list_items
        if current_list_items:
            _list_units = max(2, min(24, len(current_list_items) * 2))
            if _would_exceed_preview(_list_units):
                return False
            body_parts.append(
                f"<{current_list_tag}>"
                + "".join(f"<li>{item}</li>" for item in current_list_items)
                + f"</{current_list_tag}>"
            )
            _record_preview_units(_list_units)
        current_list_id = None
        current_list_items = []
        return True

    stop_render = False
    body_elem = doc.element.body
    for child_elem in body_elem:
        if stop_render:
            break
        tag_local = child_elem.tag.split("}")[-1] if "}" in child_elem.tag else child_elem.tag

        if tag_local == "p":
            para = Paragraph(child_elem, doc)
            style_ref = _resolve_para_style_ref(para)
            style_defaults = _resolve_style_defaults(style_ref)

            # Check for list (numPr)
            numPr = child_elem.find(qn("w:pPr") + "/" + qn("w:numPr"))
            if numPr is None:
                pPr = child_elem.find(qn("w:pPr"))
                if pPr is not None:
                    numPr = pPr.find(qn("w:numPr"))

            if numPr is not None:
                numId_elem = numPr.find(qn("w:numId"))
                ilvl_elem = numPr.find(qn("w:ilvl"))
                num_id = int(numId_elem.get(qn("w:val"), 0)) if numId_elem is not None else 0
                ilvl = int(ilvl_elem.get(qn("w:val"), 0)) if ilvl_elem is not None else 0

                # Determine list type from numbering.xml (best-effort)
                list_tag = "ul"
                try:
                    numbering_part = doc.part.numbering_part
                    if numbering_part is not None:
                        num_elem = numbering_part._element.find(
                            f'.//{qn("w:num")}[@{qn("w:numId")}="{num_id}"]'
                        )
                        if num_elem is not None:
                            abstractNumId_elem = num_elem.find(qn("w:abstractNumId"))
                            if abstractNumId_elem is not None:
                                abs_id = abstractNumId_elem.get(qn("w:val"), "")
                                abs_num = numbering_part._element.find(
                                    f'.//{qn("w:abstractNum")}[@{qn("w:abstractNumId")}="{abs_id}"]'
                                )
                                if abs_num is not None:
                                    lvl = abs_num.find(
                                        f'.//{qn("w:lvl")}[@{qn("w:ilvl")}="{ilvl}"]'
                                    )
                                    if lvl is not None:
                                        numFmt = lvl.find(qn("w:numFmt"))
                                        if numFmt is not None:
                                            fmt_val = numFmt.get(qn("w:val"), "")
                                            if fmt_val in (
                                                "decimal", "lowerLetter", "upperLetter",
                                                "lowerRoman", "upperRoman", "chineseCounting",
                                            ):
                                                list_tag = "ol"
                except Exception:
                    pass

                if num_id != current_list_id or list_tag != current_list_tag:
                    if not _flush_list():
                        preview_truncated = True
                        stop_render = True
                        break
                    current_list_id = num_id
                    current_list_tag = list_tag

                # Render list item content (without block tag wrapper)
                inner_parts: list[str] = []
                for lc in child_elem:
                    lt = lc.tag.split("}")[-1] if "}" in lc.tag else lc.tag
                    if lt == "hyperlink":
                        rId = lc.get(qn("r:id"))
                        url = ""
                        if rId:
                            try:
                                url = para.part.relationships[rId].target_ref
                            except Exception:
                                pass
                        link_inner = ""
                        for r_el in lc.findall(qn("w:r")):
                            from docx.text.run import Run
                            link_inner += _run_html(Run(r_el, para), doc)
                        if link_inner:
                            esc_url = url.replace('"', "&quot;")
                            inner_parts.append(
                                f'<a href="{esc_url}" target="_blank" '
                                f'style="color:#1155CC;text-decoration:underline;">'
                                f'{link_inner}</a>'
                            )
                    elif lt == "r":
                        from docx.text.run import Run
                        inner_parts.append(_run_html(Run(lc, para), doc))
                current_list_items.append("".join(inner_parts) or "&nbsp;")
                continue

            # Not a list item — flush pending list
            if not _flush_list():
                preview_truncated = True
                stop_render = True
                break

            # Check for an explicit page break (<w:br w:type="page"/>) inside
            # any run of this paragraph.  These are emitted as a block-level
            # separator *after* the paragraph so the frontend can position
            # the page-break overlay exactly at the right spot.
            _has_explicit_pb = any(
                br.get(qn("w:type"), "") == "page"
                for br in child_elem.findall(".//" + qn("w:br"))
            )

            # Also detect section breaks embedded in this paragraph's
            # properties (<w:pPr><w:sectPr>...<w:type w:val="nextPage"/>).
            # Word uses these to end a section on a new page (e.g. TOC boundary).
            # Sections with type="continuous" do NOT force a new page.
            _has_section_pb = False
            try:
                _pPr = child_elem.find(qn("w:pPr"))
                if _pPr is not None:
                    _sectPr = _pPr.find(qn("w:sectPr"))
                    if _sectPr is not None:
                        # OOXML: section type is a CHILD element, not an attribute:
                        #   <w:sectPr><w:type w:val="nextPage"/></w:sectPr>
                        _type_el = _sectPr.find(qn("w:type"))
                        _stype = _type_el.get(qn("w:val"), "nextPage") if _type_el is not None else "nextPage"
                        if _stype != "continuous":
                            _has_section_pb = True
            except Exception:
                pass

            # Determine tag (heading vs paragraph)
            block_tag = "p"
            _is_toc_para, _ = _detect_toc_info(
                para=para,
                p_el=child_elem,
                style_ref=style_ref,
            )
            try:
                _style_name = style_defaults.get("style_name") or ""
            except Exception:
                _style_name = ""
            try:
                _style_id = style_defaults.get("style_id") or ""
            except Exception:
                _style_id = ""
            _h_tag = None if _is_toc_para else (
                _is_heading_style_key(_style_name) or _is_heading_style_key(_style_id)
            )

            # Fallback: check w:outlineLvl in paragraph properties (and
            # inherited style pPr).  Word uses outlineLvl as the
            # authoritative heading level for the navigation pane, even
            # when the style name is non-standard (e.g. WPS documents).
            if not _is_toc_para and not _h_tag:
                _h_tag = _detect_outline_level(para, child_elem, style_ref=style_ref)

            if _h_tag and not _p_elem_text_content(child_elem):
                _h_tag = None

            # Guard: h2–h6 paragraphs whose text is longer than 25 chars are
            # almost certainly body text that was *styled* as a heading for
            # visual emphasis (e.g. a company-name paragraph styled as "标题2").
            # Real subheadings in Chinese documents are concise (typically
            # under 15 chars); 25 chars gives a safe margin without picking
            # up body sentences (which are 30-100+ chars).
            # h1 (document-level title) is exempt — it can be longer.
            if _h_tag and _h_tag != "h1":
                try:
                    if len((para.text or "").strip()) > 25:
                        _h_tag = None
                except Exception:
                    pass

            if _h_tag:
                block_tag = _h_tag

            _para_units = _estimate_paragraph_units(para)
            if _would_exceed_preview(_para_units):
                preview_truncated = True
                stop_render = True
                break

            body_parts.append(_para_html(para, doc, block_tag, style_ref=style_ref))
            _record_preview_units(_para_units)

            # Emit page-break marker (hard break or section break).
            if _has_explicit_pb or _has_section_pb:
                if _would_exceed_preview(1):
                    preview_truncated = True
                    stop_render = True
                    break
                body_parts.append('<div data-page-break="true" '
                                  'class="koto-page-break" '
                                  'contenteditable="false"></div>')
                _record_preview_units(1)

        elif tag_local == "tbl":
            if not _flush_list():
                preview_truncated = True
                stop_render = True
                break
            _table_units = _estimate_table_units(child_elem)
            _table_row_limit = None
            if preview_units_limit is not None:
                _remaining_units = max(0, preview_units_limit - preview_units_used)
                if body_parts and _remaining_units < _table_units:
                    preview_truncated = True
                    stop_render = True
                    break
                if not body_parts and _remaining_units <= _table_units:
                    _table_row_limit = max(4, min(_DOCX_PREVIEW_MAX_TABLE_ROWS, max(4, _remaining_units // 2 or 4)))
                    preview_truncated = True
            tbl = Table(child_elem, doc)
            body_parts.append(_table_html(tbl, doc, max_rows=_table_row_limit))
            if _table_row_limit is not None:
                _record_preview_units(max(6, _table_row_limit * 2))
                stop_render = True
            else:
                _record_preview_units(_table_units)

        elif tag_local == "sdt":
            # Structured Document Tag — render inner content
            if not _flush_list():
                preview_truncated = True
                stop_render = True
                break
            # Detect TOC SDT for semantic wrapper
            _is_toc_sdt = False
            try:
                _sdtPr = child_elem.find(qn("w:sdtPr"))
                if _sdtPr is not None:
                    _dpg = _sdtPr.find(".//" + qn("w:docPartGallery"))
                    if _dpg is not None:
                        _gval = _dpg.get(qn("w:val"), "")
                        if "Table of Contents" in _gval or "目录" in _gval:
                            _is_toc_sdt = True
            except Exception:
                pass
            if _is_toc_sdt:
                body_parts.append('<div class="koto-toc">')
                _record_preview_units(1)
            sdtContent = child_elem.find(qn("w:sdtContent"))
            if sdtContent is not None:
                for sdt_child in sdtContent:
                    sdt_tag = sdt_child.tag.split("}")[-1] if "}" in sdt_child.tag else sdt_child.tag
                    if sdt_tag == "p":
                        para = Paragraph(sdt_child, doc)
                        style_ref = _resolve_para_style_ref(para)
                        style_defaults = _resolve_style_defaults(style_ref)
                        # Use heading tag if the style declares it
                        _btag = "p"
                        _is_toc_para, _ = _detect_toc_info(
                            para=para,
                            p_el=sdt_child,
                            style_ref=style_ref,
                        )
                        try:
                            _style_name = style_defaults.get("style_name") or ""
                        except Exception:
                            _style_name = ""
                        try:
                            _style_id = style_defaults.get("style_id") or ""
                        except Exception:
                            _style_id = ""
                        _h_tag = None if _is_toc_para else (
                            _is_heading_style_key(_style_name) or _is_heading_style_key(_style_id)
                        )
                        if not _is_toc_para and not _h_tag:
                            _h_tag = _detect_outline_level(para, sdt_child, style_ref=style_ref)
                        if _h_tag and not _p_elem_text_content(sdt_child):
                            _h_tag = None
                        if _h_tag:
                            _btag = _h_tag
                        _para_units = _estimate_paragraph_units(para)
                        if _would_exceed_preview(_para_units):
                            preview_truncated = True
                            stop_render = True
                            break
                        body_parts.append(_para_html(para, doc, _btag, style_ref=style_ref))
                        _record_preview_units(_para_units)
                    elif sdt_tag == "tbl":
                        _table_units = _estimate_table_units(sdt_child)
                        _table_row_limit = None
                        if preview_units_limit is not None:
                            _remaining_units = max(0, preview_units_limit - preview_units_used)
                            if body_parts and _remaining_units < _table_units:
                                preview_truncated = True
                                stop_render = True
                                break
                            if not body_parts and _remaining_units <= _table_units:
                                _table_row_limit = max(4, min(_DOCX_PREVIEW_MAX_TABLE_ROWS, max(4, _remaining_units // 2 or 4)))
                                preview_truncated = True
                        tbl = Table(sdt_child, doc)
                        body_parts.append(_table_html(tbl, doc, max_rows=_table_row_limit))
                        if _table_row_limit is not None:
                            _record_preview_units(max(6, _table_row_limit * 2))
                            stop_render = True
                            break
                        _record_preview_units(_table_units)
            if _is_toc_sdt:
                body_parts.append('</div>')
            if stop_render:
                break

        elif tag_local == "sectPr":
            # Direct body-level <w:sectPr> — the final section of the document.
            # In DOCX this is always the last child of <w:body> and represents
            # the document-level defaults (not a break). Do NOT emit a page
            # break here; it would add an extra blank page at the end.
            pass

    if not stop_render and not _flush_list():
        preview_truncated = True
        stop_render = True

    # NOTE: Footer no longer appended to body — rendered per-page by NodeView.

    # ── Extract per-section header/footer and page dimension metadata ─────
    # _section_html depends on _para_html (local function), so we extract
    # section data HERE (still inside _docx_to_rich_html) to avoid scope issues.
    sections_data: list[dict] = []
    try:
        _emu_px = lambda e: round(e / 914400 * 96) if e else 0
        for sec in doc.sections:
            sec_info: dict = {
                "page_width_px":  _emu_px(sec.page_width),
                "page_height_px": _emu_px(sec.page_height),
                "margin_top_px":  _emu_px(sec.top_margin),
                "margin_bottom_px": _emu_px(sec.bottom_margin),
                "margin_left_px": _emu_px(sec.left_margin),
                "margin_right_px": _emu_px(sec.right_margin),
            }
            # Default header / footer
            try:
                sec_info["header_html"] = _section_html(sec.header, doc, "koto-header")
            except Exception:
                sec_info["header_html"] = ""
            try:
                sec_info["footer_html"] = _section_html(sec.footer, doc, "koto-footer")
            except Exception:
                sec_info["footer_html"] = ""
            # First-page header / footer (Word's "Different First Page")
            try:
                sec_info["first_header_html"] = _section_html(sec.first_page_header, doc, "koto-header") if sec.different_first_page_header_footer else ""
            except Exception:
                sec_info["first_header_html"] = ""
            try:
                sec_info["first_footer_html"] = _section_html(sec.first_page_footer, doc, "koto-footer") if sec.different_first_page_header_footer else ""
            except Exception:
                sec_info["first_footer_html"] = ""
            # Even-page header / footer
            try:
                sec_info["even_header_html"] = _section_html(sec.even_page_header, doc, "koto-header")
            except Exception:
                sec_info["even_header_html"] = ""
            try:
                sec_info["even_footer_html"] = _section_html(sec.even_page_footer, doc, "koto-footer")
            except Exception:
                sec_info["even_footer_html"] = ""
            sections_data.append(sec_info)
    except Exception as sec_exc:
        import logging as _log
        _log.getLogger(__name__).debug("[DocxParser] 节段信息提取失败 (非致命): %s", sec_exc)

    preview_meta = {
        "pending": bool(progressive_preview and preview_truncated),
        "target_pages": _DOCX_PREVIEW_TARGET_PAGES if progressive_preview else None,
    }
    return "\n".join(body_parts), sections_data, preview_meta


def _extract_images_from_paragraphs(html: str) -> str:
    """将 <p> 标签内的 <img> 元素移到段落外部，避免编辑器路径错误。

    TipTap/ProseMirror 不允许图片 void 元素嵌套在段落块元素中，否则在
    规范化时会抛出路径错误。

    规则：
    - 若段落只包含一个 <img>（忽略空白），则整个 <p>...</p> 替换为裸 <img>。
    - 若段落混合了文字和图片，则将 <img> 提取出来放在段落之后。
    """
    import re as _re

    # <img ... > or <img ... /> — both have no '>' inside attributes
    _IMG = r'<img[^>]*>'

    # Match any <p> that contains at least one <img> anywhere inside.
    # Group 1 = p opening-tag attributes, Group 2 = full inner content.
    any_p_with_img = _re.compile(
        r'<p([^>]*)>((?:(?!</p>).)*' + _IMG + r'(?:(?!</p>).)*)</p>',
        _re.IGNORECASE | _re.DOTALL,
    )

    def _process(m: '_re.Match') -> str:
        p_attrs = m.group(1)
        p_inner = m.group(2)
        imgs = _re.findall(_IMG, p_inner, _re.IGNORECASE | _re.DOTALL)
        clean_inner = _re.sub(_IMG, '', p_inner, flags=_re.IGNORECASE | _re.DOTALL).strip()
        result = (f'<p{p_attrs}>{clean_inner}</p>' if clean_inner else '') + ''.join(imgs)
        return result

    return any_p_with_img.sub(_process, html)


def _extract_headings_from_html(html: str) -> list[dict]:
    """Extract a flat list of headings from the rendered HTML for the outline panel.

    First tries <h1>–<h6> tags. If none are found, falls back to extracting
    headings from TOC entries (koto-toc-N paragraphs with anchor links), which
    is reliable for Chinese documents that use custom style names.

    Returns a list like [{"level": 1, "text": "Title", "id": "anchorId"}, ...].
    """
    import re as _re
    headings: list[dict] = []
    _id_re = _re.compile(r'id="([^"]*)"', _re.IGNORECASE)
    _href_re = _re.compile(r'href="#([^"]*)"', _re.IGNORECASE)
    _strip_tags = _re.compile(r'<[^>]+>')
    _strip_toc_tab = _re.compile(r'<span[^>]*class="koto-toc-tab"[^>]*>.*?</span>', _re.IGNORECASE | _re.DOTALL)

    # Primary: look for <h1>-<h6>
    _heading_re = _re.compile(
        r'<(h[1-6])([^>]*)>(.*?)</\1>',
        _re.IGNORECASE | _re.DOTALL,
    )
    for m in _heading_re.finditer(html):
        level = int(m.group(1)[1])
        attrs = m.group(2)
        inner_html = m.group(3)
        id_m = _id_re.search(attrs)
        hid = id_m.group(1) if id_m else ""
        text = _strip_tags.sub('', inner_html).strip()
        if text:
            headings.append({"level": level, "text": text, "id": hid})

    # Fallback: extract headings from TOC entries (koto-toc-N paragraphs)
    # These are reliable for Chinese documents using custom role styles.
    if not headings:
        _toc_re = _re.compile(
            r'<p[^>]*class="koto-toc-(\d+)"[^>]*>(.*?)</p>',
            _re.IGNORECASE | _re.DOTALL,
        )
        for m in _toc_re.finditer(html):
            level = int(m.group(1))
            inner = m.group(2)
            # Extract anchor ID from href="#..." in the inner hyperlink
            href_m = _href_re.search(inner)
            hid = href_m.group(1) if href_m else ""
            # Remove the koto-toc-tab spacer span before stripping tags
            inner = _strip_toc_tab.sub('', inner)
            text = _strip_tags.sub('', inner).strip()
            # Remove trailing page number (digits at end, optionally with leading whitespace)
            text = _re.sub(r'\s*\d+\s*$', '', text).strip()
            if text:
                headings.append({"level": level, "text": text, "id": hid})

    return headings


def _extract_docx_comments(file_path: str) -> list[dict[str, Any]]:
    """
    从 DOCX 的 word/comments.xml 提取批注信息。

    Returns:
        [{id, author, date, text, anchor_text}]
        若无批注或解析失败返回空列表。
    """
    import zipfile
    from xml.etree import ElementTree as ET

    ns = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "w14": "http://schemas.microsoft.com/office/word/2010/wordml",
    }

    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            # ── 解析 comments.xml ────────────────────────────────────
            if "word/comments.xml" not in zf.namelist():
                return []
            comments_xml = zf.read("word/comments.xml")
            comments_tree = ET.fromstring(comments_xml)
            comments_map: dict[str, dict] = {}
            for c in comments_tree.findall(".//w:comment", ns):
                cid = c.get(f"{{{ns['w']}}}id", "")
                author = c.get(f"{{{ns['w']}}}author", "")
                date = c.get(f"{{{ns['w']}}}date", "")
                # 拼接所有 <w:t> 文本
                texts = [t.text or "" for t in c.findall(".//w:t", ns)]
                comments_map[cid] = {
                    "id": cid,
                    "author": author,
                    "date": date,
                    "text": "".join(texts).strip(),
                    "anchor_text": "",
                }

            if not comments_map:
                return []

            # ── 从 document.xml 提取批注锚定原文 ────────────────────
            if "word/document.xml" in zf.namelist():
                doc_xml = zf.read("word/document.xml")
                doc_tree = ET.fromstring(doc_xml)
                # 构建 commentRangeStart id → 对应的 body 元素范围
                body = doc_tree.find(".//w:body", ns)
                if body is not None:
                    _extract_anchor_texts(body, comments_map, ns)

            return list(comments_map.values())
    except Exception as exc:
        logger.debug("[DocxParser] 批注提取失败 (非致命): %s", exc)
        return []


def _extract_anchor_texts(
    body_el: Any, comments_map: dict[str, dict], ns: dict[str, str]
) -> None:
    """遍历 document.xml body，提取 commentRangeStart/End 之间的文本。"""
    from xml.etree import ElementTree as ET

    # 收集所有文本节点和 comment range 标记的顺序
    events: list[tuple[str, str]] = []  # (type, value_or_id)

    def _walk(el: Any) -> None:
        tag = el.tag.split("}")[-1] if "}" in el.tag else el.tag
        if tag == "commentRangeStart":
            cid = el.get(f"{{{ns['w']}}}id", "")
            events.append(("start", cid))
        elif tag == "commentRangeEnd":
            cid = el.get(f"{{{ns['w']}}}id", "")
            events.append(("end", cid))
        elif tag == "t" and el.text:
            events.append(("text", el.text))
        for child in el:
            _walk(child)

    _walk(body_el)

    # 按顺序扫描，收集每个 comment id 对应的 anchor 文本
    active_ids: set[str] = set()
    for etype, val in events:
        if etype == "start":
            active_ids.add(val)
        elif etype == "end":
            active_ids.discard(val)
        elif etype == "text" and active_ids:
            for cid in active_ids:
                if cid in comments_map:
                    comments_map[cid]["anchor_text"] += val


def parse_docx(file_path: str, *, progressive_preview: bool = False) -> dict[str, Any]:
    """
    将 DOCX 转换为保留完整 Word 格式的内联样式 HTML。

    首先尝试 _docx_to_rich_html()（python-docx 全格式管线）；
    若失败则回退到 mammoth 语义化管线（含后处理 Pass）。

    Returns:
        {"html": str, "messages": list[str]}

    When ``progressive_preview`` is true, only the initial preview chunk is
    rendered and the result includes ``progressive.pending`` so the frontend can
    fetch the full document in the background before enabling editing.
    """
    messages_out: list[str] = []

    # ── Primary path: rich python-docx renderer ───────────────────────────
    try:
        rich_html, sections_data, preview_meta = _docx_to_rich_html(
            file_path,
            progressive_preview=progressive_preview,
        )

        # strip <style>/<script> just in case
        rich_html = re.sub(
            r"<style[^>]*>.*?</style>|<script[^>]*>.*?</script>",
            "",
            rich_html,
            flags=re.DOTALL | re.IGNORECASE,
        )

        # Unwrap single-column layout tables that have no visible borders
        try:
            rich_html = _unwrap_layout_tables(rich_html)
        except Exception as exc:
            logger.warning("[DocxParser] 布局表格展开失败 (非致命): %s", exc)

        # Deduplicate images
        try:
            rich_html = _deduplicate_images(rich_html)
        except Exception as exc:
            logger.warning("[DocxParser] 图片去重失败 (非致命): %s", exc)

        # Extract <img> from <p> tags to prevent editor path errors
        try:
            rich_html = _extract_images_from_paragraphs(rich_html)
        except Exception as exc:
            logger.warning("[DocxParser] 图片段落提取失败 (非致命): %s", exc)

        # ── Extract page dimensions for frontend page-break overlay ───
        page_meta: dict[str, Any] = {}
        try:
            if sections_data:
                _sec0 = sections_data[0]
                page_meta = {
                    "page_width_px":    _sec0.get("page_width_px", 0),
                    "page_height_px":   _sec0.get("page_height_px", 0),
                    "margin_top_px":    _sec0.get("margin_top_px", 0),
                    "margin_bottom_px": _sec0.get("margin_bottom_px", 0),
                    "margin_left_px":   _sec0.get("margin_left_px", 0),
                    "margin_right_px":  _sec0.get("margin_right_px", 0),
                    "header_html":      _sec0.get("header_html", ""),
                    "footer_html":      _sec0.get("footer_html", ""),
                }
                page_meta["sections"] = sections_data
        except Exception as meta_exc:
            logger.debug("[DocxParser] 页面元数据提取失败 (非致命): %s", meta_exc)

        # ── Extract headings for the outline/navigation panel ─────────
        headings: list[dict] = []
        try:
            headings = _extract_headings_from_html(rich_html)
        except Exception:
            pass

        result = {"html": rich_html, "messages": messages_out, "headings": headings}
        result.update(page_meta)
        if progressive_preview:
            result["progressive"] = {
                "pending": bool(preview_meta.get("pending")),
                "target_pages": preview_meta.get("target_pages") or _DOCX_PREVIEW_TARGET_PAGES,
            }
        # ── 提取 DOCX 批注 ──────────────────────────────────────────
        try:
            comments = _extract_docx_comments(file_path)
            if comments:
                result["comments"] = comments
        except Exception:
            pass
        return result

    except Exception as primary_exc:
        logger.warning(
            "[DocxParser] 富格式渲染失败，回退到 mammoth: %s", primary_exc
        )

    # ── Fallback path: mammoth semantic renderer ──────────────────────────
    try:
        import mammoth
    except ImportError:
        raise RuntimeError("python-docx 和 mammoth 均未安装，请执行: pip install python-docx mammoth")

    def _img_handler(image: Any) -> dict[str, str]:
        """将图片转换为内联 base64 data URI，自动压缩大图。"""
        try:
            with image.open() as f:
                img_bytes = f.read()
            content_type = image.content_type or "image/png"
            img_bytes, content_type = _compress_image_bytes(img_bytes, content_type)
            b64 = base64.b64encode(img_bytes).decode("ascii")
            return {"src": f"data:{content_type};base64,{b64}"}
        except Exception as e:
            logger.warning(f"[DocxParser] 图片内联失败: {e}")
            return {"src": ""}

    try:
        with open(file_path, "rb") as f:
            result = mammoth.convert_to_html(
                f,
                convert_image=mammoth.images.img_element(_img_handler),
            )
        for msg in result.messages:
            messages_out.append(str(msg))

        clean_html = re.sub(
            r"<style[^>]*>.*?</style>|<script[^>]*>.*?</script>",
            "",
            result.value,
            flags=re.DOTALL | re.IGNORECASE,
        )
        clean_html = re.sub(
            r"(<p[^>]*>)(?:(?:body|h[1-6]|blockquote|strong|em|code|pre|table|td|th|ul|ol|li)\s*\{[^{}]*\})+",
            r"\1",
            clean_html,
            flags=re.IGNORECASE,
        )
        clean_html = re.sub(r"<p[^>]*>\s*</p>", "", clean_html)

        try:
            tbl_styles = _extract_table_styles(file_path)
            clean_html = _inject_table_styles(clean_html, tbl_styles)
        except Exception as exc:
            logger.warning("[DocxParser] 表格样式注入失败 (非致命): %s", exc)

        try:
            clean_html = _unwrap_layout_tables(clean_html)
        except Exception as exc:
            logger.warning("[DocxParser] 布局表格展开失败 (非致命): %s", exc)

        try:
            floating_srcs = _get_floating_image_srcs(file_path)
            if floating_srcs:
                imgs_html = "".join(
                    f'<img src="{src}" '
                    'style="float:right;max-width:28%;max-height:180px;'
                    'margin:0 0 10px 14px;object-fit:contain;border-radius:2px;" '
                    'alt="" />'
                    for src in floating_srcs
                )
                first_block = re.search(
                    r"<(?:p|h[1-6]|table|ul|ol|blockquote)\b", clean_html
                )
                if first_block:
                    clean_html = (
                        clean_html[: first_block.start()]
                        + imgs_html
                        + clean_html[first_block.start() :]
                    )
                else:
                    clean_html = imgs_html + clean_html
        except Exception as exc:
            logger.warning("[DocxParser] 浮动图片注入失败 (非致命): %s", exc)

        try:
            clean_html = _deduplicate_images(clean_html)
        except Exception as exc:
            logger.warning("[DocxParser] 图片去重失败 (非致命): %s", exc)

        fallback_result = {"html": clean_html, "messages": messages_out}
        try:
            comments = _extract_docx_comments(file_path)
            if comments:
                fallback_result["comments"] = comments
        except Exception:
            pass
        return fallback_result
    except Exception as e:
        logger.error(f"[DocxParser] 解析失败: {e}")
        raise


# ─────────────────────────────────────────────────────────────────────────────
# XLSX → Univer Sheets IWorkbookData JSON
# ─────────────────────────────────────────────────────────────────────────────

# Univer CellValueType constants (must match @univerjs/core CellValueType enum)
_UNIVER_TYPE_STRING = 1
_UNIVER_TYPE_NUMBER = 2
_UNIVER_TYPE_BOOLEAN = 3

# Horizontal alignment: Univer HorizontalAlign enum
_ALIGN_H_MAP = {
    "general": 0,
    "left": 1,
    "center": 2,
    "right": 3,
    "justify": 6,
    "distributed": 7,
}

# Vertical alignment: Univer VerticalAlign enum
_ALIGN_V_MAP = {
    "top": 1,
    "middle": 2,
    "center": 2,
    "bottom": 3,
    "justify": 4,
    "distributed": 5,
}


def _openpyxl_cell_to_univer(cell: Any) -> dict[str, Any] | None:
    """将单个 openpyxl Cell 转换为 Univer ICellData 对象。"""
    v = cell.value
    if v is None:
        return None

    cell_data: dict[str, Any] = {}

    # ── Value & type ──────────────────────────────────────────────────────────
    if isinstance(v, bool):
        cell_data["v"] = int(v)
        cell_data["t"] = _UNIVER_TYPE_BOOLEAN
    elif isinstance(v, (int, float)):
        cell_data["v"] = v
        cell_data["t"] = _UNIVER_TYPE_NUMBER
    else:
        cell_data["v"] = str(v)
        cell_data["t"] = _UNIVER_TYPE_STRING

    # ── Style ─────────────────────────────────────────────────────────────────
    style: dict[str, Any] = {}
    try:
        font = cell.font
        if font:
            if font.bold:
                style["bl"] = 1
            if font.italic:
                style["it"] = 1
            if font.size:
                style["fs"] = int(font.size)
            if font.color and font.color.type == "rgb" and font.color.rgb not in (
                "00000000", "FF000000"
            ):
                # Univer expects { rgb: "#RRGGBB" }
                style["cl"] = {"rgb": "#" + font.color.rgb[2:]}
        fill = cell.fill
        if fill and fill.fill_type not in (None, "none") and fill.fgColor:
            if fill.fgColor.type == "rgb" and fill.fgColor.rgb not in (
                "00000000",
                "FFFFFFFF",
            ):
                style["bg"] = {"rgb": "#" + fill.fgColor.rgb[2:]}
        ali = cell.alignment
        if ali:
            if ali.horizontal and ali.horizontal in _ALIGN_H_MAP:
                style["ht"] = _ALIGN_H_MAP[ali.horizontal]
            if ali.vertical and ali.vertical in _ALIGN_V_MAP:
                style["vt"] = _ALIGN_V_MAP[ali.vertical]
    except Exception:
        pass  # 样式提取失败不影响数据

    if style:
        cell_data["s"] = style

    return cell_data


def parse_xlsx(file_path: str, original_name: str | None = None) -> dict[str, Any]:
    """
    使用 openpyxl 将 XLSX 转换为 Univer Sheets IWorkbookData 快照格式。

    Args:
        file_path: 临时文件路径（可能是 UUID 命名）
        original_name: 用户上传的原始文件名（用于设置 workbook name）

    Returns:
        {
          "id": str,
          "name": str,
          "sheetOrder": [str, ...],
          "sheets": {
            "<sheetId>": { ... }
          },
          "_warnings": ["..."]   # 非空时应由前端显示提示
        }
    """
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl 未安装，请执行: pip install openpyxl")

    # ── 公式检测：通过 ZIP 级 sheet XML 快速扫描，避免额外完整加载一遍 workbook ──
    _warnings: list[str] = []
    def _xlsx_contains_formula_fast(path: str) -> bool:
        import zipfile

        try:
            with zipfile.ZipFile(path, "r") as zf:
                for info in zf.infolist():
                    name = info.filename
                    if not name.startswith("xl/worksheets/") or not name.endswith(".xml"):
                        continue
                    with zf.open(info) as fh:
                        tail = b""
                        for chunk in iter(lambda: fh.read(65536), b""):
                            data = tail + chunk
                            if b"<f" in data:
                                return True
                            tail = data[-8:]
        except Exception:
            return False
        return False

    if _xlsx_contains_formula_fast(file_path):
        _warnings.append(
            "此表格包含公式（如 =SUM(...)）。Koto 目前以「静态值」模式读取 Excel，"
            "公式已转换为计算结果，保存导出后公式将永久丢失。如需保留公式，请下载原始文件。"
        )

    wb = openpyxl.load_workbook(file_path, data_only=True)

    workbook_id = str(uuid.uuid4())
    if original_name:
        workbook_name = os.path.splitext(os.path.basename(original_name))[0]
    else:
        workbook_name = os.path.splitext(os.path.basename(file_path))[0]
    sheet_order: list[str] = []
    sheets: dict[str, Any] = {}

    # ── Shared styles registry ────────────────────────────────────────────────
    # Univer v0.5.x expects cellData["s"] to be a *string* style-ID that keys
    # into the top-level "styles" map.  Inline IStyleData objects (the old
    # approach) are silently ignored by Univer's createUnit(), which is why
    # all cells appeared blank even though the data was correctly parsed.
    _style_hash_to_id: dict[str, str] = {}
    _styles_registry: dict[str, Any] = {}

    def _get_style_id(style_obj: dict[str, Any]) -> str:
        """Return a stable string key for *style_obj*, registering it if new."""
        import json as _json
        h = _json.dumps(style_obj, sort_keys=True, ensure_ascii=False)
        if h not in _style_hash_to_id:
            sid = str(len(_style_hash_to_id))
            _style_hash_to_id[h] = sid
            _styles_registry[sid] = style_obj
        return _style_hash_to_id[h]

    for idx, ws in enumerate(wb.worksheets):
        sheet_id = f"sheet{idx + 1}"
        sheet_order.append(sheet_id)

        # ── Cell data: nested dict {row: {col: ICellData}} ────────────────────
        cell_data: dict[int, dict[int, Any]] = {}
        for row in ws.iter_rows():
            for cell in row:
                cd = _openpyxl_cell_to_univer(cell)
                if cd is not None:
                    # Convert inline style object → style ID string
                    if "s" in cd and isinstance(cd["s"], dict):
                        cd["s"] = _get_style_id(cd["s"])
                    r = cell.row - 1
                    c = cell.column - 1
                    if r not in cell_data:
                        cell_data[r] = {}
                    cell_data[r][c] = cd

        # ── Merge data ────────────────────────────────────────────────────────
        merge_data: list[dict[str, int]] = []
        for merge_range in ws.merged_cells.ranges:
            merge_data.append({
                "startRow": merge_range.min_row - 1,
                "startColumn": merge_range.min_col - 1,
                "endRow": merge_range.max_row - 1,
                "endColumn": merge_range.max_col - 1,
            })

        # ── Column widths ─────────────────────────────────────────────────────
        # openpyxl width is in Excel "character units"; Univer IColumnData.w is pixels.
        # Conversion: pixels = chars * 7 + 5  (96 DPI, default font ~11pt)
        col_data: dict[int, dict] = {}
        for col_letter, col_dim in ws.column_dimensions.items():
            if col_dim.width and col_dim.width > 0:
                col_idx = openpyxl.utils.column_index_from_string(col_letter) - 1
                px = max(4, round(col_dim.width * 7 + 5))
                col_data[col_idx] = {"w": px}

        # ── Row heights ───────────────────────────────────────────────────────
        # openpyxl height is in points (1pt = 1/72 inch); Univer IRowData.h is pixels.
        # Conversion: px = pt * 96 / 72 = pt * 4/3
        row_data: dict[int, dict] = {}
        for row_num, row_dim in ws.row_dimensions.items():
            if row_dim.height and row_dim.height > 0:
                row_idx = row_num - 1
                px = max(4, round(row_dim.height * 96 / 72))
                row_data[row_idx] = {"h": px}

        sheets[sheet_id] = {
            "id": sheet_id,
            "name": ws.title,
            "rowCount": max(ws.max_row or 30, 30),
            "columnCount": max(ws.max_column or 10, 10),
            "cellData": cell_data,
            "mergeData": merge_data,
            "columnData": col_data,
            "rowData": row_data,
        }

    wb.close()

    return {
        "id": workbook_id,
        "name": workbook_name,
        "appVersion": "0.5.0",   # required by IWorkbookData; Univer fails silently without it
        "locale": "zh-CN",       # required by IWorkbookData; used for cell formatting
        "sheetOrder": sheet_order,
        "sheets": sheets,
        "styles": _styles_registry,  # style-id → IStyleData (populated during cell scan)
        "resources": [],  # Univer plugin resources (e.g. conditional formatting)
        "_warnings": _warnings,  # 前端应在非空时展示用户提示
    }


# ─────────────────────────────────────────────────────────────────────────────
# PPTX → 结构化文本 JSON (卡片编辑模型)
# ─────────────────────────────────────────────────────────────────────────────


def parse_pptx(file_path: str) -> list[dict[str, Any]]:
    """
    使用 python-pptx 提取每个 Slide 的文本框内容。
    保留 shape_id 以便后端导出时可回写原文件。

    Returns:
        [{"slide_id": int, "slide_index": int,
          "texts": [{"shape_id": int, "shape_name": str, "text": str, "is_title": bool}]}]
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError:
        raise RuntimeError("python-pptx 未安装，请执行: pip install python-pptx")

    prs = Presentation(file_path)
    slides_data: list[dict[str, Any]] = []

    for slide_idx, slide in enumerate(prs.slides):
        texts: list[dict[str, Any]] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_content = "\n".join(
                para.text for para in shape.text_frame.paragraphs
            ).strip()

            if not text_content:
                continue

            is_title = False
            try:
                ph = shape.placeholder_format
                if ph is not None:
                    is_title = ph.type in (
                        PP_PLACEHOLDER.TITLE,
                        PP_PLACEHOLDER.CENTER_TITLE,
                    )
            except Exception:
                pass

            texts.append(
                {
                    "shape_id": shape.shape_id,
                    "shape_name": shape.name,
                    "text": text_content,
                    "is_title": is_title,
                }
            )

        slides_data.append(
            {
                "slide_id": slide_idx + 1,
                "slide_index": slide_idx,
                "texts": texts,
            }
        )

    return slides_data


# ─────────────────────────────────────────────────────────────────────────────
# PPTX → 几何画布 JSON (含图片/表格/备注)
# ─────────────────────────────────────────────────────────────────────────────


def parse_pptx_geometry(file_path: Any) -> dict[str, Any]:
    """
    使用 python-pptx 提取每个 Slide 的完整几何数据，供前端画布编辑器渲染。
    包含文本框 (含字体样式)、图片 (base64 data URI)、表格 (单元格文本)、备注。
    支持 GROUP 形状递归、从 layout/master 继承背景色。

    Returns:
        {
          "slide_width_emu": int,
          "slide_height_emu": int,
          "slides": [
            {
              "slide_index": int, "slide_id": int,
              "background": "#xxxxxx", "notes": str,
              "shapes": [
                {
                  "id": int, "name": str,
                  "_type": "TEXT" | "PICTURE" | "TABLE",
                  "left": int, "top": int, "width": int, "height": int,
                  "z_order": int, "fill": str | null,
                  # TEXT:    "has_text": true, "paragraphs": [{"align": str, "runs": [...]}]
                  # PICTURE: "image_b64": "data:image/...;base64,..."
                  # TABLE:   "table_rows": int, "table_cols": int,
                  #          "cells": [{"row": int, "col": int, "text": str}]
                }
              ]
            }
          ]
        }
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        raise RuntimeError("python-pptx 未安装，请执行: pip install python-pptx")

    prs = Presentation(file_path)
    slide_w: int = prs.slide_width or 9144000
    slide_h: int = prs.slide_height or 6858000
    slides_data: list[dict[str, Any]] = []

    # ── Build theme-color lookup table (schemeClr name → "#rrggbb") ──────
    # Many PPTX borders/fills use scheme colors (dk1, lt1, accent1, etc.) instead
    # of explicit RGB, causing line.color.rgb to throw. We pre-resolve them once.
    _theme_colors: dict[str, str] = {}
    try:
        _NS_T = "http://schemas.openxmlformats.org/drawingml/2006/main"
        _prs_part = prs.part
        for _rel in _prs_part.rels.values():
            if "slideMaster" in _rel.reltype:
                _mst_part = _rel.target_part
                for _r2 in _mst_part.rels.values():
                    if "theme" in _r2.reltype:
                        import lxml.etree as _ET
                        _theme_el = _ET.fromstring(_r2.target_part.blob)
                        _cs = _theme_el.find(f".//{{{_NS_T}}}clrScheme")
                        if _cs is not None:
                            for _c in _cs:
                                _name = _c.tag.split("}")[1]  # e.g. "dk1", "lt1", "accent1"
                                _srgb = _c.find(f"{{{_NS_T}}}srgbClr")
                                _sys  = _c.find(f"{{{_NS_T}}}sysClr")
                                if _srgb is not None:
                                    _theme_colors[_name] = "#" + _srgb.get("val", "000000").lower()
                                elif _sys is not None:
                                    _last = _sys.get("lastClr", "")
                                    if _last:
                                        _theme_colors[_name] = "#" + _last.lower()
                        break
                break
    except Exception:
        pass

    # ── Extract presentation-level default font size (defaultTextStyle) ──
    # Non-placeholder shapes inherit font size from this style chain.
    # Without it the frontend must guess, often resulting in inflated text.
    _default_font_size_pt: float = 18.0  # OOXML spec hardcoded fallback
    _default_title_font_size_pt: float = 36.0  # OOXML spec default for titles
    try:
        _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        _dts = prs.element.find(f"{{{_NS_P}}}defaultTextStyle")
        if _dts is not None:
            _lvl1 = _dts.find(f"{{{_NS_T}}}lvl1pPr")
            if _lvl1 is not None:
                _drp = _lvl1.find(f"{{{_NS_T}}}defRPr")
                if _drp is not None and _drp.get("sz"):
                    _default_font_size_pt = int(_drp.get("sz")) / 100.0
    except Exception:
        pass
    # ── Extract title font size from slideMaster txStyles/titleStyle ──
    try:
        _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        for _mst in prs.slide_masters:
            _txS = _mst._element.find(f"{{{_NS_P}}}txStyles")
            if _txS is not None:
                _titS = _txS.find(f"{{{_NS_P}}}titleStyle")
                if _titS is not None:
                    _lvl1t = _titS.find(f"{{{_NS_T}}}lvl1pPr")
                    if _lvl1t is not None:
                        _drpt = _lvl1t.find(f"{{{_NS_T}}}defRPr")
                        if _drpt is not None and _drpt.get("sz"):
                            _default_title_font_size_pt = int(_drpt.get("sz")) / 100.0
            break  # only need first master
    except Exception:
        pass

    def _resolve_color(color_obj: Any) -> str | None:
        """Return "#rrggbb" from a python-pptx color object, resolving scheme colors via theme."""
        try:
            return "#" + str(color_obj.rgb).lower()
        except Exception:
            pass
        try:
            from pptx.enum.dml import MSO_THEME_COLOR
            # python-pptx theme_color names like DARK_1, LIGHT_1, ACCENT_1 etc.
            # Map to OOXML scheme names: DARK_1→dk1, LIGHT_1→lt1, ACCENT_N→accentN, etc.
            _tc = color_obj.theme_color
            _name_map = {
                1: "dk1", 2: "lt1", 3: "dk2", 4: "lt2",
                5: "accent1", 6: "accent2", 7: "accent3", 8: "accent4",
                9: "accent5", 10: "accent6", 11: "hlink", 12: "folHlink",
            }
            _hex = _theme_colors.get(_name_map.get(int(_tc), ""))
            if _hex:
                return _hex
        except Exception:
            pass
        return None

    # ── Inner helpers ────────────────────────────────────────────────────

    _NS_DML = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _extract_grad_css(fill: Any) -> str | None:
        """Try to convert a python-pptx gradient fill into a CSS linear-gradient string."""
        try:
            _NS = _NS_DML
            gfill = fill._fill.find(f"{{{_NS}}}gradFill")
            if gfill is None:
                return None
            # Rotation angle is stored in <a:lin ang="…"/> where ang is in 60000ths of a degree
            ang_el = gfill.find(f"{{{_NS}}}lin")
            angle_css = 0
            if ang_el is not None:
                ang_raw = ang_el.get("ang", "0")
                angle_css = round(int(ang_raw) / 60000.0)
            gs_lst = gfill.find(f"{{{_NS}}}gsLst")
            if gs_lst is None:
                return None
            stops = []
            for gs in gs_lst.findall(f"{{{_NS}}}gs"):
                pos_pct = round(int(gs.get("pos", "0")) / 1000.0)
                srgb = gs.find(f"{{{_NS}}}srgbClr")
                if srgb is not None and len(srgb.get("val", "")) == 6:
                    stops.append(f"#{srgb.get('val').lower()} {pos_pct}%")
                else:
                    # sysClr or schemeClr — skip; gradient may be partial but better than nothing
                    lum_mod = gs.find(f".//{{{_NS}}}lumMod")
                    if lum_mod is None:
                        return None  # can't resolve theme colour, bail
                    return None
            if len(stops) < 2:
                return None
            return f"linear-gradient({angle_css}deg, {', '.join(stops)})"
        except Exception:
            return None

    def _extract_bg(slide: Any) -> dict[str, Any]:
        """Walk slide → layout → master for the first extractable fill.
        Returns dict: {"color": "#rrggbb"} for solid, {"gradient": "css"} for gradient,
        {"image": "data:…"} for picture fill, or {} (fallback white).

        IMPORTANT: slide.background._element returns <p:cSld> (the whole slide),
        NOT just the <p:bg> element.  We must narrow the search to <p:bg><p:bgPr>
        to avoid picking up shape images as backgrounds.
        """
        from pptx.oxml.ns import qn as _qn_bg
        _R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        _P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

        def _resolve_blip_image(bg_pr_element: Any, part: Any) -> dict:
            """Find <a:blip r:embed="rIdN"/> in <p:bgPr> XML, resolve to compressed data URI."""
            try:
                blip = bg_pr_element.find('.//' + _qn_bg('a:blip'))
                if blip is not None:
                    rId = blip.get(f'{{{_R_NS}}}embed')
                    if rId and hasattr(part, 'rels') and rId in part.rels:
                        img_part = part.rels[rId].target_part
                        img_bytes, mime = _compress_image_bytes(
                            img_part.blob, img_part.content_type or 'image/png'
                        )
                        b64 = base64.b64encode(img_bytes).decode('ascii')
                        return {"image": f"data:{mime};base64,{b64}"}
            except Exception:
                pass
            return {}

        def _resolve_solid_color(bg_pr_element: Any) -> dict:
            """Find <a:solidFill><a:srgbClr val="..."/> in <p:bgPr> XML."""
            try:
                solid = bg_pr_element.find('.//' + _qn_bg('a:solidFill'))
                if solid is not None:
                    srgb = solid.find(_qn_bg('a:srgbClr'))
                    if srgb is not None:
                        val = srgb.get('val', '')
                        if len(val) == 6:
                            return {"color": "#" + val.lower()}
            except Exception:
                pass
            return {}

        for src in (slide, getattr(slide, "slide_layout", None), getattr(slide, "slide_master", None)):
            if src is None:
                continue
            try:
                # Locate the <p:bg> element within <p:cSld>.
                # slide.background._element is <p:cSld>; searching it directly
                # would also find <a:blip> inside shape tree → wrong image!
                cSld = src.background._element
                bg_el = cSld.find(f'{{{_P_NS}}}bg')
                if bg_el is None:
                    continue  # this source has no background definition

                bgPr = bg_el.find(f'{{{_P_NS}}}bgPr')
                src_part = getattr(src, 'part', None)

                if bgPr is not None and src_part is not None:
                    # Check for image fill (blipFill)
                    result = _resolve_blip_image(bgPr, src_part)
                    if result:
                        return result
                    # Check for solid fill
                    result = _resolve_solid_color(bgPr)
                    if result:
                        return result
                    # Check for gradient fill
                    f = src.background.fill
                    if getattr(f.type, 'name', '') == 'GRADIENT':
                        css = _extract_grad_css(f)
                        if css:
                            return {"gradient": css}

                # bgRef (theme reference) — use python-pptx fill API
                bgRef = bg_el.find(f'{{{_P_NS}}}bgRef')
                if bgRef is not None:
                    f = src.background.fill
                    fill_name = getattr(f.type, 'name', '') if f.type is not None else ''
                    if fill_name == 'SOLID':
                        try:
                            return {"color": "#" + str(f.fore_color.rgb).lower()}
                        except Exception:
                            pass
                    if fill_name == 'GRADIENT':
                        css = _extract_grad_css(f)
                        if css:
                            return {"gradient": css}
            except Exception:
                pass
        return {"color": "#FFFFFF"}

    def _parse_tf(tf: Any, layout_defaults: dict[str, Any] | None = None) -> list[dict[str, Any]]:
        """Convert a python-pptx TextFrame into our [{align, runs:[...]}] format.

        Handles three common cases that cause text to disappear:
        1. Paragraph-level default run properties (defRPr) for font size/color
           — most real PPTs store the font size here, not on individual runs.
        2. Empty para.runs when text lives in <a:fld> field elements
           — use para.text as a fallback single run.
        3. Hard line-breaks (<a:br>) within a paragraph
           — emitted as a run with text='\n'.
        """
        # XML namespace used by DrawingML
        _NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

        # ── Read body-level lstStyle defaults (lvl1pPr.defRPr) ──────────────
        # Many PPTX files store the inherited font size here rather than on
        # individual runs or paragraphs.  Serves as lowest-priority fallback.
        _body_defaults: dict[str, Any] = {}
        try:
            lst_style = tf._txBody.find(f"{{{_NS}}}lstStyle")
            if lst_style is not None:
                lvl1 = lst_style.find(f"{{{_NS}}}lvl1pPr")
                if lvl1 is not None:
                    _def_rpr = lvl1.find(f"{{{_NS}}}defRPr")
                    if _def_rpr is not None:
                        # We need _read_rpr but it's defined below; inline a minimal version.
                        _sz = _def_rpr.get("sz")
                        if _sz:
                            try:
                                _v = int(_sz)
                                if _v > 0:
                                    _body_defaults["size"] = round(_v / 100.0, 1)
                            except Exception:
                                pass
                        for _key, _attr in (("bold", "b"), ("italic", "i")):
                            _av = _def_rpr.get(_attr)
                            if _av and _av.lower() not in ("0", "false"):
                                _body_defaults[_key] = True
                        _lat = _def_rpr.find(f"{{{_NS}}}latin")
                        if _lat is not None:
                            _tf = _lat.get("typeface", "")
                            if _tf and not _tf.startswith("+"):
                                _body_defaults["fontName"] = _tf
                        _sol = _def_rpr.find(f"{{{_NS}}}solidFill")
                        if _sol is not None:
                            _srgb = _sol.find(f"{{{_NS}}}srgbClr")
                            if _srgb is not None and len(_srgb.get("val", "")) == 6:
                                _body_defaults["color"] = "#" + _srgb.get("val", "").lower()
        except Exception:
            pass

        def _pt_from_emu_hundredths(val: Any) -> float | None:
            """python-pptx stores font size as 100ths of a point (e.g. 2400 = 24pt)."""
            try:
                v = int(val)
                return round(v / 100.0, 1) if v > 0 else None
            except Exception:
                return None

        def _read_rpr(rpr_el: Any) -> dict[str, Any]:
            """Extract font attributes from an <a:rPr> or <a:defRPr> XML element."""
            out: dict[str, Any] = {}
            if rpr_el is None:
                return out
            try:
                sz = rpr_el.get("sz")
                if sz:
                    pt = _pt_from_emu_hundredths(sz)
                    if pt:
                        out["size"] = pt
            except Exception:
                pass
            try:
                b = rpr_el.get("b")
                if b and b.lower() not in ("0", "false"):
                    out["bold"] = True
            except Exception:
                pass
            try:
                i = rpr_el.get("i")
                if i and i.lower() not in ("0", "false"):
                    out["italic"] = True
            except Exception:
                pass
            try:
                u = rpr_el.get("u")
                if u and u != "none":
                    out["underline"] = True
            except Exception:
                pass
            # Font name from <a:latin typeface="..."> child (Latin / default font)
            try:
                latin = rpr_el.find(f"{{{_NS}}}latin")
                if latin is not None:
                    tf_val = latin.get("typeface", "")
                    if tf_val and not tf_val.startswith("+"):
                        out["fontName"] = tf_val
            except Exception:
                pass
            # East Asian font from <a:ea typeface="..."> child (CJK-specific override)
            try:
                ea = rpr_el.find(f"{{{_NS}}}ea")
                if ea is not None:
                    ea_tf = ea.get("typeface", "")
                    if ea_tf and not ea_tf.startswith("+") and ea_tf != "+mj-ea" and ea_tf != "+mn-ea":
                        out["eaFontName"] = ea_tf
            except Exception:
                pass
            # Solid colour from <a:solidFill><a:srgbClr val="rrggbb">
            try:
                solid = rpr_el.find(f"{{{_NS}}}solidFill")
                if solid is not None:
                    srgb = solid.find(f"{{{_NS}}}srgbClr")
                    if srgb is not None:
                        val = srgb.get("val", "")
                        if len(val) == 6:
                            out["color"] = "#" + val.lower()
            except Exception:
                pass
            # Character spacing: <a:rPr spc="N"/> — hundredths of a point (can be negative)
            try:
                spc = rpr_el.get("spc")
                if spc:
                    out["charSpacing"] = int(spc)
            except Exception:
                pass
            return out

        paras: list[dict[str, Any]] = []
        for para in tf.paragraphs:
            align_name = "LEFT"
            try:
                if para.alignment:
                    align_name = para.alignment.name
            except Exception:
                pass

            # ── Read paragraph-level default run properties (defRPr) ──────
            # These serve as the fallback when a run has no explicit rPr attrs.
            para_defaults: dict[str, Any] = {}
            try:
                pPr = para._p.find(f"{{{_NS}}}pPr")
                if pPr is not None:
                    defRPr = pPr.find(f"{{{_NS}}}defRPr")
                    para_defaults = _read_rpr(defRPr)
            except Exception:
                pass

            p_obj: dict[str, Any] = {"align": align_name, "runs": []}

            # ── Extract paragraph spacing properties ─────────────────
            # <a:lnSpc> → line spacing, <a:spcBef> → space before, <a:spcAft> → space after
            try:
                pPr = para._p.find(f"{{{_NS}}}pPr")
                if pPr is not None:
                    # Line spacing: <a:lnSpc><a:spcPct val="150000"/></a:lnSpc>
                    lnSpc = pPr.find(f"{{{_NS}}}lnSpc")
                    if lnSpc is not None:
                        spcPct = lnSpc.find(f"{{{_NS}}}spcPct")
                        if spcPct is not None:
                            val = int(spcPct.get("val", "0"))
                            if val > 0:
                                p_obj["lineSpacing"] = round(val / 100000.0, 2)
                        else:
                            spcPts = lnSpc.find(f"{{{_NS}}}spcPts")
                            if spcPts is not None:
                                val = int(spcPts.get("val", "0"))
                                if val > 0:
                                    p_obj["lineSpacingPt"] = round(val / 100.0, 1)
                    # Space before: <a:spcBef><a:spcPts val="600"/></a:spcBef>
                    spcBef = pPr.find(f"{{{_NS}}}spcBef")
                    if spcBef is not None:
                        pts = spcBef.find(f"{{{_NS}}}spcPts")
                        if pts is not None:
                            val = int(pts.get("val", "0"))
                            if val > 0:
                                p_obj["spaceBefore"] = round(val / 100.0, 1)
                        else:
                            pct = spcBef.find(f"{{{_NS}}}spcPct")
                            if pct is not None:
                                val = int(pct.get("val", "0"))
                                if val > 0:
                                    p_obj["spaceBeforePct"] = round(val / 100000.0, 2)
                    # Space after: <a:spcAft><a:spcPts val="600"/></a:spcAft>
                    spcAft = pPr.find(f"{{{_NS}}}spcAft")
                    if spcAft is not None:
                        pts = spcAft.find(f"{{{_NS}}}spcPts")
                        if pts is not None:
                            val = int(pts.get("val", "0"))
                            if val > 0:
                                p_obj["spaceAfter"] = round(val / 100.0, 1)
                        else:
                            pct = spcAft.find(f"{{{_NS}}}spcPct")
                            if pct is not None:
                                val = int(pct.get("val", "0"))
                                if val > 0:
                                    p_obj["spaceAfterPct"] = round(val / 100000.0, 2)
            except Exception:
                pass

            # ── Iterate raw XML to capture <a:r> AND <a:br> (hard line-break) ──
            # python-pptx's para.runs only yields <a:r> elements and silently
            # skips <a:fld> field elements (e.g. slide numbers, date fields).
            # We walk _p directly so nothing is lost.
            try:
                for child in para._p:
                    tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

                    if tag == "br":
                        # Hard line-break: emit a newline run so rendering looks correct
                        p_obj["runs"].append({"text": "\n"})
                        continue

                    if tag not in ("r", "fld"):
                        continue

                    # Text node: <a:t> child
                    t_el = child.find(f"{{{_NS}}}t")
                    text_val = (t_el.text or "") if t_el is not None else ""

                    # Run-level properties (<a:rPr>)
                    rPr = child.find(f"{{{_NS}}}rPr")
                    run_attrs = _read_rpr(rPr)

                    # Merge: run → paragraph → layout placeholder → body (lstStyle) defaults
                    _layout_defs = layout_defaults or {}
                    r: dict[str, Any] = {"text": text_val}
                    for key in ("size", "bold", "italic", "underline", "fontName", "eaFontName", "color", "charSpacing"):
                        if key in run_attrs:
                            r[key] = run_attrs[key]
                        elif key in para_defaults:
                            r[key] = para_defaults[key]
                        elif key in _layout_defs:
                            r[key] = _layout_defs[key]
                        elif key in _body_defaults:
                            r[key] = _body_defaults[key]

                    p_obj["runs"].append(r)

            except Exception:
                # Fallback: use python-pptx's high-level .runs API
                for run in para.runs:
                    r_fb: dict[str, Any] = {"text": run.text}
                    try:
                        if run.font.size:
                            r_fb["size"] = round(run.font.size.pt, 1)
                    except Exception:
                        pass
                    try:
                        if run.font.bold:
                            r_fb["bold"] = True
                    except Exception:
                        pass
                    try:
                        if run.font.italic:
                            r_fb["italic"] = True
                    except Exception:
                        pass
                    try:
                        if run.font.underline:
                            r_fb["underline"] = True
                    except Exception:
                        pass
                    try:
                        if run.font.name:
                            r_fb["fontName"] = run.font.name
                    except Exception:
                        pass
                    try:
                        if run.font.color and run.font.color.type is not None:
                            r_fb["color"] = "#" + str(run.font.color.rgb).lower()
                    except Exception:
                        pass
                    p_obj["runs"].append(r_fb)

            # ── Fallback: if we still have no content, use para.text directly ──
            # This catches edge cases where XML walk also found nothing.
            if not any(r.get("text") for r in p_obj["runs"]):
                raw = ""
                try:
                    raw = para.text
                except Exception:
                    pass
                if raw.strip():
                    fallback_run: dict[str, Any] = {"text": raw}
                    # Apply body defaults first, layout, then para defaults (higher priority)
                    _layout_defs = layout_defaults or {}
                    for key, val in {**_body_defaults, **_layout_defs, **para_defaults}.items():
                        fallback_run.setdefault(key, val)
                    p_obj["runs"] = [fallback_run]

            # ── Sibling-size inheritance ─────────────────────────────────
            # In OOXML, a run with no explicit sz inherits the "effective" size
            # of the paragraph.  When defRPr & lstStyle also have no size, the
            # best heuristic is the most-common explicit size of neighbouring
            # runs in the same paragraph (e.g.  sz=2100 | sz=None | sz=2100 →
            # the middle run should also be 2100, not the global default).
            _runs = p_obj["runs"]
            if _runs and any("size" not in r for r in _runs if r.get("text", "").strip()):
                _sibling_sizes = [r["size"] for r in _runs if "size" in r]
                if _sibling_sizes:
                    # Use the most common size among siblings
                    from collections import Counter
                    _best_size = Counter(_sibling_sizes).most_common(1)[0][0]
                    for r in _runs:
                        if "size" not in r and r.get("text", "").strip():
                            r["size"] = _best_size

            paras.append(p_obj)
        return paras

    def _collect_shapes(
        shapes_iter: Any,
        out: list[dict[str, Any]],
        z_base: int = 0,
        off_left: int = 0,
        off_top: int = 0,
        grp_ch_off_x: int = 0,
        grp_ch_off_y: int = 0,
        grp_scale_x: float = 1.0,
        grp_scale_y: float = 1.0,
    ) -> None:
        """
        Recursively parse shapes into `out`.
        GROUP shapes are unwrapped and children are appended with their
        slide-absolute coordinates (group offset + internal coordinate scaling).
        """
        for z_idx, shape in enumerate(shapes_iter):
            # ── Resolve effective geometry ───────────────────────────────────
            # python-pptx returns None for left/top/width/height when a
            # placeholder shape inherits its position from the slide layout.
            # (TITLE and BODY placeholders are almost always in this state.)
            # Falling back to 0 makes every text shape invisible because their
            # div has zero width/height — TABLE and PICTURE shapes are
            # always explicitly sized, which is why only they were visible.
            # Fix: pull the real geometry from the matching layout placeholder.
            eff_left = shape.left
            eff_top  = shape.top
            eff_w    = shape.width
            eff_h    = shape.height

            if None in (eff_left, eff_top, eff_w, eff_h):
                try:
                    ph_fmt = shape.placeholder_format
                    if ph_fmt is not None:
                        slide_layout = getattr(
                            getattr(shape, "part", None), "slide_layout", None
                        )
                        if slide_layout is not None:
                            for lph in slide_layout.placeholders:
                                try:
                                    if lph.placeholder_format.idx == ph_fmt.idx:
                                        if eff_left is None:
                                            eff_left = lph.left
                                        if eff_top is None:
                                            eff_top = lph.top
                                        if eff_w is None:
                                            eff_w = lph.width
                                        if eff_h is None:
                                            eff_h = lph.height
                                        break
                                except Exception:
                                    pass
                        # Walk slide master if layout didn't resolve everything
                        if None in (eff_left, eff_top, eff_w, eff_h):
                            slide_master = getattr(slide_layout, "slide_master", None) if slide_layout else None
                            if slide_master is not None:
                                for mph in slide_master.placeholders:
                                    try:
                                        if mph.placeholder_format.idx == ph_fmt.idx:
                                            if eff_left is None:
                                                eff_left = mph.left
                                            if eff_top is None:
                                                eff_top = mph.top
                                            if eff_w is None:
                                                eff_w = mph.width
                                            if eff_h is None:
                                                eff_h = mph.height
                                            break
                                    except Exception:
                                        pass
                        # If still None after layout+master, use standard widescreen EMU defaults
                        if None in (eff_left, eff_top, eff_w, eff_h):
                            _ph_idx = getattr(ph_fmt, "idx", -1)
                            _prs_shape = getattr(getattr(shape, "part", None), "presentation", None)
                            _sw = int(getattr(_prs_shape, "slide_width",  None) or 9144000)
                            _sh = int(getattr(_prs_shape, "slide_height", None) or 6858000)
                            if _ph_idx == 0:   # Title
                                if eff_left is None: eff_left = 457200
                                if eff_top  is None: eff_top  = 274638
                                if eff_w    is None: eff_w    = 8229600
                                if eff_h    is None: eff_h    = 1143000
                            elif _ph_idx == 1:  # Body / Content
                                if eff_left is None: eff_left = 457200
                                if eff_top  is None: eff_top  = 1600200
                                if eff_w    is None: eff_w    = 8229600
                                if eff_h    is None: eff_h    = 4525963
                            else:
                                if eff_left is None: eff_left = 0
                                if eff_top  is None: eff_top  = 0
                                if eff_w    is None: eff_w    = _sw
                                if eff_h    is None: eff_h    = _sh
                except Exception:
                    pass

            abs_left = off_left + round((( eff_left or 0) - grp_ch_off_x) * grp_scale_x)
            abs_top  = off_top  + round((( eff_top  or 0) - grp_ch_off_y) * grp_scale_y)
            # Also scale the shape's own width/height when inside a group
            if grp_scale_x != 1.0 and eff_w is not None:
                eff_w = round(eff_w * grp_scale_x)
            if grp_scale_y != 1.0 and eff_h is not None:
                eff_h = round(eff_h * grp_scale_y)

            # ── Group: recurse with absolute-coordinate offset + internal scaling ──
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    # OOXML groups have an internal coordinate system.
                    # <p:grpSpPr><a:xfrm> carries:
                    #   off  x/y  : group's position on slide (same as eff_left/eff_top)
                    #   ext  cx/cy: group's rendered size on slide
                    #   chOff x/y : origin of the internal coordinate system
                    #   chExt cx/cy: size of the internal coordinate system
                    # Child shapes' left/top are in internal coordinates; we map them to slide coords.
                    _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
                    _NS_A2 = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    xfrm = shape.element.find(
                        f"{{{_NS_P}}}grpSpPr/{{{_NS_A2}}}xfrm"
                    ) or shape.element.find(
                        f"{{{_NS_A2}}}grpSpPr/{{{_NS_A2}}}xfrm"
                    )
                    ch_off_x = ch_off_y = 0
                    ch_ext_cx = eff_w or 1
                    ch_ext_cy = eff_h or 1
                    if xfrm is not None:
                        chOff = xfrm.find(f"{{{_NS_A2}}}chOff")
                        chExt = xfrm.find(f"{{{_NS_A2}}}chExt")
                        if chOff is not None:
                            ch_off_x = int(chOff.get("x", 0))
                            ch_off_y = int(chOff.get("y", 0))
                        if chExt is not None:
                            ch_ext_cx = int(chExt.get("cx", 1)) or 1
                            ch_ext_cy = int(chExt.get("cy", 1)) or 1
                    grp_w = eff_w or ch_ext_cx
                    grp_h = eff_h or ch_ext_cy
                    scale_x = grp_w / ch_ext_cx
                    scale_y = grp_h / ch_ext_cy
                    _collect_shapes(
                        shape.shapes,
                        out,
                        z_base=z_base + z_idx * 100,
                        off_left=abs_left,
                        off_top=abs_top,
                        grp_ch_off_x=ch_off_x,
                        grp_ch_off_y=ch_off_y,
                        grp_scale_x=scale_x,
                        grp_scale_y=scale_y,
                    )
                except Exception:
                    _collect_shapes(
                        shape.shapes,
                        out,
                        z_base=z_base + z_idx * 100,
                        off_left=abs_left,
                        off_top=abs_top,
                    )
                continue

            # ── Media (video/audio embedded shapes): skip entirely ───────────
            # MSO_SHAPE_TYPE.MEDIA = 16.  python-pptx cannot render video/audio
            # and attempting to read their blobs would base64-encode hundreds of MB.
            try:
                if getattr(MSO_SHAPE_TYPE, 'MEDIA', None) is not None and shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    continue
            except Exception:
                pass

            s: dict[str, Any] = {
                "id": shape.shape_id,
                "name": shape.name,
                "left": abs_left,
                "top": abs_top,
                "width":  eff_w or 0,
                "height": eff_h or 0,
                "z_order": z_base + z_idx,
                "fill": None,
            }

            # Shape fill — solid, gradient, or picture
            try:
                fill = shape.fill
                fill_name = getattr(fill.type, 'name', '') if fill.type is not None else ''
                if fill_name == 'SOLID':
                    _fc = _resolve_color(fill.fore_color)
                    if _fc:
                        s["fill"] = _fc
                elif fill_name == 'GRADIENT':
                    css = _extract_grad_css(fill)
                    if css:
                        s["fillGradient"] = css
                elif fill_name == 'PICTURE':
                    try:
                        from pptx.oxml.ns import qn as _qn_sh
                        blip = shape.element.find('.//' + _qn_sh('a:blip'))
                        if blip is not None:
                            rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                            if rId:
                                img_part = shape.part.related_parts[rId]
                                _raw_blob = img_part.blob
                                _raw_mime = img_part.content_type or 'image/png'
                                # Skip non-image blobs (e.g. video used as fill) or oversized media
                                if _raw_mime.startswith('image/') and len(_raw_blob) <= _MAX_BLOB_BYTES:
                                    img_bytes, mime = _compress_image_bytes(_raw_blob, _raw_mime)
                                    b64 = base64.b64encode(img_bytes).decode('ascii')
                                    s["fillImage"] = f"data:{mime};base64,{b64}"
                    except Exception:
                        pass
            except Exception:
                pass

            # Shape outline/border
            try:
                line = shape.line
                if line and line.width is not None and line.width > 0:
                    _lc = None
                    try:
                        _lc = _resolve_color(line.color)
                    except Exception:
                        pass
                    # line.width is in EMU; store as-is, frontend scales
                    s["border"] = {"widthEmu": int(line.width), "color": _lc or "#000000"}
            except Exception:
                pass

            # Shape rotation (degrees, clockwise-positive; python-pptx returns float or None)
            try:
                _rot = shape.rotation
                if _rot is not None and _rot != 0.0:
                    s["rotation"] = round(float(_rot), 2)
            except Exception:
                pass

            # ── Picture ───────────────────────────────────────────────
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_blob = shape.image.blob
                    img_mime = shape.image.content_type or "image/png"
                    # Guard: skip if non-image MIME (video poster frame) or oversized blob
                    if not img_mime.startswith("image/") or len(img_blob) > _MAX_BLOB_BYTES:
                        logger.warning(
                            "[parse_pptx] skipping oversized/non-image blob: mime=%s size=%.1f MB",
                            img_mime, len(img_blob) / 1048576,
                        )
                        s["_type"] = "PICTURE"
                        s["image_b64"] = ""  # placeholder — too large or non-image media
                        out.append(s)
                        continue
                    img_blob, img_mime = _compress_image_bytes(img_blob, img_mime)
                    b64 = base64.b64encode(img_blob).decode("ascii")
                    s["_type"] = "PICTURE"
                    s["image_b64"] = f"data:{img_mime};base64,{b64}"
                    out.append(s)
                    continue
            except Exception:
                pass

            # ── Table ─────────────────────────────────────────────────
            try:
                if shape.has_table:
                    tbl = shape.table
                    # Column widths and row heights (EMU)
                    col_widths = []
                    row_heights = []
                    try:
                        col_widths = [int(col.width) for col in tbl.columns]
                    except Exception:
                        pass
                    try:
                        row_heights = [int(row.height) for row in tbl.rows]
                    except Exception:
                        pass
                    cells: list[dict[str, Any]] = []
                    for r_idx, row in enumerate(tbl.rows):
                        for c_idx, cell in enumerate(row.cells):
                            cell_text = ""
                            try:
                                cell_text = (
                                    cell.text_frame.text if cell.text_frame else ""
                                )
                            except Exception:
                                pass
                            cell_d: dict[str, Any] = {"row": r_idx, "col": c_idx, "text": cell_text}
                            # Cell fill
                            try:
                                cfill = cell.fill
                                if cfill.type is not None and getattr(cfill.type, 'name', '') == 'SOLID':
                                    _cf = _resolve_color(cfill.fore_color)
                                    if _cf:
                                        cell_d["fill"] = _cf
                            except Exception:
                                pass
                            # Per-cell text formatting (first run of first paragraph)
                            try:
                                if cell.text_frame and cell.text_frame.paragraphs:
                                    fp = cell.text_frame.paragraphs[0]
                                    if fp.runs:
                                        fr = fp.runs[0]
                                        if fr.font.size:   cell_d["fontSize"] = round(fr.font.size.pt, 1)
                                        if fr.font.bold:   cell_d["bold"] = True
                                        if fr.font.color and fr.font.color.type is not None:
                                            cell_d["color"] = "#" + str(fr.font.color.rgb).lower()
                                    cell_d["align"] = (
                                        fp.alignment.name if fp.alignment else "LEFT"
                                    )
                            except Exception:
                                pass
                            cells.append(cell_d)
                    s["_type"] = "TABLE"
                    s["table_rows"] = len(tbl.rows)
                    s["table_cols"] = len(tbl.columns)
                    s["col_widths"] = col_widths
                    s["row_heights"] = row_heights
                    s["cells"] = cells
                    out.append(s)
                    continue
            except Exception:
                pass

            # ── Generic Shapes / Icons / SVGs / Charts / Connectors ────────
            try:
                _st = getattr(shape, "shape_type", None)
                if _st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM,
                           MSO_SHAPE_TYPE.GRAPHIC_FRAME):
                    s["_type"] = "SHAPE"
                    # Rounded rectangle: extract corner radius from XML
                    try:
                        from pptx.oxml.ns import qn as _qn_r
                        prstGeom = shape.element.find('.//' + _qn_r('a:prstGeom'))
                        if prstGeom is not None:
                            s["autoShapeType"] = prstGeom.get("prst", "")
                            avLst = prstGeom.find('.//' + _qn_r('a:gd'))
                            if avLst is not None and avLst.get("fmla", "").startswith("val "):
                                s["cornerRadiusEmu"] = int(avLst.get("fmla").split()[1])
                    except Exception:
                        pass
                elif _st in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREE_FORM):
                    s["_type"] = "LINE"
            except Exception:
                pass

            try:
                if getattr(shape, "has_chart", False):
                    s["_type"] = "CHART"
            except Exception:
                pass

            # ── Text frame ────────────────────────────────────────────
            try:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    s["_type"] = "TEXT"
                    s["has_text"] = True
                    is_title = False
                    try:
                        from pptx.enum.shapes import PP_PLACEHOLDER
                        ph = shape.placeholder_format
                        if ph is not None:
                            is_title = ph.type in (
                                PP_PLACEHOLDER.TITLE,
                                PP_PLACEHOLDER.CENTER_TITLE,
                            )
                    except Exception:
                        pass
                    s["is_title"] = is_title

                    # ── Extract text body properties (<a:bodyPr>) ─────────
                    _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    try:
                        bodyPr = shape.text_frame._txBody.find(f"{{{_NS_A}}}bodyPr")
                        if bodyPr is not None:
                            anchor = bodyPr.get("anchor")
                            # If not set at slide level, inherit from layout → master placeholder
                            if not anchor:
                                try:
                                    _ph_fmt = getattr(shape, "placeholder_format", None)
                                    if _ph_fmt is not None:
                                        _ph_idx = _ph_fmt.idx
                                        _sl_layout = getattr(getattr(shape, "part", None), "slide_layout", None)
                                        if _sl_layout is not None:
                                            for _lph in _sl_layout.placeholders:
                                                try:
                                                    if _lph.placeholder_format.idx == _ph_idx:
                                                        _lbPr = _lph.text_frame._txBody.find(f"{{{_NS_A}}}bodyPr")
                                                        if _lbPr is not None:
                                                            anchor = _lbPr.get("anchor") or anchor
                                                        break
                                                except Exception:
                                                    pass
                                        if not anchor:
                                            _sl_master = getattr(_sl_layout, "slide_master", None) if _sl_layout else None
                                            if _sl_master is not None:
                                                for _mph in _sl_master.placeholders:
                                                    try:
                                                        if _mph.placeholder_format.idx == _ph_idx:
                                                            _mbPr = _mph.text_frame._txBody.find(f"{{{_NS_A}}}bodyPr")
                                                            if _mbPr is not None:
                                                                anchor = _mbPr.get("anchor") or anchor
                                                            break
                                                    except Exception:
                                                        pass
                                except Exception:
                                    pass
                            if anchor:
                                s["textAnchor"] = anchor  # t | ctr | b | just | dist
                            # Text insets (EMU). OOXML defaults: l/r=91440, t/b=45720
                            lIns = bodyPr.get("lIns")
                            tIns = bodyPr.get("tIns")
                            rIns = bodyPr.get("rIns")
                            bIns = bodyPr.get("bIns")
                            s["textInsets"] = {
                                "l": int(lIns) if lIns is not None else 91440,
                                "t": int(tIns) if tIns is not None else 45720,
                                "r": int(rIns) if rIns is not None else 91440,
                                "b": int(bIns) if bIns is not None else 45720,
                            }
                            # Text autofit mode: spAutoFit | normAutofit (with fontScale) | noAutofit
                            _af_sp   = bodyPr.find(f"{{{_NS_A}}}spAutoFit")
                            _af_norm = bodyPr.find(f"{{{_NS_A}}}normAutofit")
                            if _af_sp is not None:
                                s["autoFit"] = "sp"
                            elif _af_norm is not None:
                                s["autoFit"] = "norm"
                                _fs = _af_norm.get("fontScale")
                                if _fs:
                                    # fontScale: 100000 = 100%, 75000 = 75% of original pt size
                                    s["fontScale"] = round(int(_fs) / 1000.0, 1)
                            # Word wrap: "square" (default, wraps) or "none" (no wrap)
                            _wrap = bodyPr.get("wrap")
                            if _wrap == "none":
                                s["wordWrap"] = "none"
                    except Exception:
                        pass

                    # ── Layout placeholder lstStyle inheritance ──────────
                    # OOXML font size resolution: run rPr → para defRPr →
                    # shape lstStyle → layout placeholder lstStyle →
                    # master titleStyle/bodyStyle → defaultTextStyle.
                    # Extract layout placeholder's lstStyle defaults so
                    # _parse_tf can insert them in the merge chain.
                    _layout_defaults: dict[str, Any] = {}
                    try:
                        _ph_fmt2 = getattr(shape, "placeholder_format", None)
                        if _ph_fmt2 is not None:
                            _ph_idx2 = _ph_fmt2.idx
                            _sl_layout2 = getattr(getattr(shape, "part", None), "slide_layout", None)
                            if _sl_layout2 is not None:
                                for _lph2 in _sl_layout2.placeholders:
                                    try:
                                        if _lph2.placeholder_format.idx == _ph_idx2:
                                            _ltxBody = _lph2.text_frame._txBody
                                            _llstStyle = _ltxBody.find(f"{{{_NS_A}}}lstStyle")
                                            if _llstStyle is not None:
                                                _llvl1 = _llstStyle.find(f"{{{_NS_A}}}lvl1pPr")
                                                if _llvl1 is not None:
                                                    _ldefRPr = _llvl1.find(f"{{{_NS_A}}}defRPr")
                                                    if _ldefRPr is not None:
                                                        _lsz = _ldefRPr.get("sz")
                                                        if _lsz:
                                                            try:
                                                                _lv = int(_lsz)
                                                                if _lv > 0:
                                                                    _layout_defaults["size"] = round(_lv / 100.0, 1)
                                                            except Exception:
                                                                pass
                                                        for _lk, _la in (("bold", "b"), ("italic", "i")):
                                                            _lav = _ldefRPr.get(_la)
                                                            if _lav and _lav.lower() not in ("0", "false"):
                                                                _layout_defaults[_lk] = True
                                                        _llat = _ldefRPr.find(f"{{{_NS_A}}}latin")
                                                        if _llat is not None:
                                                            _ltf = _llat.get("typeface", "")
                                                            if _ltf and not _ltf.startswith("+"):
                                                                _layout_defaults["fontName"] = _ltf
                                                        _lea = _ldefRPr.find(f"{{{_NS_A}}}ea")
                                                        if _lea is not None:
                                                            _leaf = _lea.get("typeface", "")
                                                            if _leaf and not _leaf.startswith("+"):
                                                                _layout_defaults["eaFontName"] = _leaf
                                                        _lsol = _ldefRPr.find(f"{{{_NS_A}}}solidFill")
                                                        if _lsol is not None:
                                                            _lsrgb = _lsol.find(f"{{{_NS_A}}}srgbClr")
                                                            if _lsrgb is not None and len(_lsrgb.get("val", "")) == 6:
                                                                _layout_defaults["color"] = "#" + _lsrgb.get("val", "").lower()
                                            break
                                    except Exception:
                                        pass
                    except Exception:
                        pass

                    s["paragraphs"] = _parse_tf(shape.text_frame, layout_defaults=_layout_defaults)
                    out.append(s)
            except Exception:
                pass

            # SHAPE/CHART types set their _type above but don't self-append (unlike
            # TEXT, PICTURE, TABLE which already called out.append+continue).
            if s.get("_type") in ("SHAPE", "CHART"):
                out.append(s)


    # ── Main loop ────────────────────────────────────────────────────────

    for slide_idx, slide in enumerate(prs.slides):
        bg_info = _extract_bg(slide)
        # bg_info is now a dict: {"color":"#..."} | {"gradient":"css"} | {"image":"data:..."} | {}

        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        shapes_data: list[dict[str, Any]] = []
        # ── Collect decorative shapes from slide master + layout first (lowest z-order).
        # These form the visual theme/background template visible behind slide content.
        try:
            _slide_layout = getattr(slide, "slide_layout", None)
            _slide_master = getattr(_slide_layout, "slide_master", None) if _slide_layout else None
            # Slide master shapes (z_base=-2000, bottom-most layer)
            if _slide_master is not None:
                _mst_decos = [_s for _s in _slide_master.shapes
                              if not getattr(_s, "is_placeholder", False)]
                _mst_out: list[dict[str, Any]] = []
                _collect_shapes(_mst_decos, _mst_out, z_base=-2000)
                for _ms in _mst_out:
                    _ms["editable"] = False
                shapes_data.extend(_mst_out)
            # Slide layout shapes (z_base=-1000, above master, below slide content)
            if _slide_layout is not None:
                _lay_decos = [_s for _s in _slide_layout.shapes
                              if not getattr(_s, "is_placeholder", False)]
                _lay_out: list[dict[str, Any]] = []
                _collect_shapes(_lay_decos, _lay_out, z_base=-1000)
                for _ls in _lay_out:
                    _ls["editable"] = False
                shapes_data.extend(_lay_out)
        except Exception:
            pass
        # Slide's own shapes (z_base=0, top layer)
        _collect_shapes(slide.shapes, shapes_data)

        # ── Text exclusion zones ─────────────────────────────────────────────
        # When a lower-z-order text shape physically overlaps a higher-z-order
        # text shape, expand the lower shape's right/left textInset so its text
        # content doesn't flow under the upper shape.  This reproduces
        # PowerPoint's implicit text-runaround behaviour for overlapping shapes.
        _EXCL_BUFFER = 91440  # 1 EMU point of extra breathing room
        _text_shapes = [
            s for s in shapes_data
            if s.get("has_text") and s.get("textInsets") is not None
        ]
        for _a in _text_shapes:
            _al  = _a["left"];  _at  = _a["top"]
            _ar  = _al + _a["width"];  _ab  = _at + _a["height"]
            _acx = _al + _a["width"] / 2
            for _b in shapes_data:
                if _b.get("z_order", 0) <= _a.get("z_order", 0):
                    continue  # only consider shapes rendered ON TOP of _a
                if not _b.get("has_text"):
                    continue  # only exclude space for text-bearing shapes
                _bl = _b["left"];  _bt = _b["top"]
                _br = _bl + _b["width"];  _bb = _bt + _b["height"]
                # Must overlap in BOTH axes
                if _bb <= _at or _bt >= _ab:
                    continue
                if _br <= _al or _bl >= _ar:
                    continue
                # Expand inset on the side where _b sits relative to _a's centre
                _ins = _a["textInsets"]
                if (_bl + _b["width"] / 2) > _acx:
                    # _b is to the right of _a's centre → expand right inset
                    _need_r = (_ar - _bl) + _EXCL_BUFFER
                    if _need_r > _ins.get("r", 91440):
                        _ins["r"] = _need_r
                else:
                    # _b is to the left of _a's centre → expand left inset
                    _need_l = (_br - _al) + _EXCL_BUFFER
                    if _need_l > _ins.get("l", 91440):
                        _ins["l"] = _need_l

        slides_data.append(
            {
                "slide_index": slide_idx,
                "slide_id": slide_idx + 1,
                "background": bg_info.get("color", "#FFFFFF"),
                "backgroundGradient": bg_info.get("gradient"),
                "backgroundImage": bg_info.get("image"),
                "notes": notes_text,
                "shapes": shapes_data,
            }
        )

    return {
        "slide_width_emu": int(slide_w),
        "slide_height_emu": int(slide_h),
        "default_font_size_pt": _default_font_size_pt,
        "default_title_font_size_pt": _default_title_font_size_pt,
        "slides": slides_data,
    }



# ─────────────────────────────────────────────────────────────────────────────
# PDF → 文本提取 + 原始 URL
# ─────────────────────────────────────────────────────────────────────────────

# OCR: minimum character count to consider a page as "text-bearing" (not scanned)
_PDF_OCR_THRESHOLD = 50

# Common Tesseract-OCR install paths on Windows
_TESSERACT_WIN_PATHS = [
    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
]


def _configure_tesseract() -> bool:
    """Set pytesseract.tesseract_cmd to the first existing Windows path. Returns True if found."""
    try:
        import pytesseract  # type: ignore
    except ImportError:
        return False
    for path in _TESSERACT_WIN_PATHS:
        if os.path.isfile(path):
            pytesseract.pytesseract.tesseract_cmd = path
            return True
    # If already on PATH, keep default
    import shutil
    return shutil.which("tesseract") is not None


def _ocr_pdf_pages(
    file_path: str, page_indices: list[int]
) -> dict[int, str]:
    """
    Render the given 0-based page indices with PyMuPDF and run Tesseract OCR.
    Returns {page_index: text}. Silently returns {} if any dependency is missing.
    """
    try:
        import fitz  # PyMuPDF
        import pytesseract
        from PIL import Image  # type: ignore
    except ImportError as exc:
        logger.info(f"[PdfOCR] 依赖缺失，跳过 OCR: {exc}")
        return {}

    if not _configure_tesseract():
        logger.warning(
            "[PdfOCR] 未找到 Tesseract 可执行文件。"
            " 请安装: https://github.com/UB-Mannheim/tesseract/wiki"
        )
        return {}

    results: dict[int, str] = {}
    try:
        doc = fitz.open(file_path)
        # Determine available language packs; prefer Chinese + English
        try:
            langs = pytesseract.get_languages()
            lang = "+".join(
                lc for lc in ("chi_sim", "chi_tra", "eng") if lc in langs
            ) or "eng"
        except Exception:
            lang = "eng"

        for idx in page_indices:
            if idx >= len(doc):
                continue
            page = doc[idx]
            # Render at 2× scale (150 DPI → 300 DPI) for better OCR accuracy
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            try:
                text = pytesseract.image_to_string(img, lang=lang).strip()
            except pytesseract.pytesseract.TesseractError as e:
                # If Chinese pack unavailable, retry with English only
                logger.warning(f"[PdfOCR] lang={lang} 失败，回退 eng: {e}")
                text = pytesseract.image_to_string(img, lang="eng").strip()
            results[idx] = text
        doc.close()
    except Exception as e:
        logger.warning(f"[PdfOCR] OCR 执行失败: {e}")
    return results


def _flatten_pdf_outline(reader: Any, items: list, depth: int = 0) -> list:
    """Recursively convert a pypdf outline into a JSON-serializable list."""
    result = []
    for item in items:
        if isinstance(item, list):
            result.extend(_flatten_pdf_outline(reader, item, depth + 1))
        else:
            page_num = None
            try:
                page_num = reader.get_destination_page_number(item) + 1
            except Exception:
                pass
            children = _flatten_pdf_outline(reader, getattr(item, "children", []), depth + 1)
            result.append({
                "title": str(getattr(item, "title", "") or ""),
                "page": page_num,
                "depth": depth,
                "children": children,
            })
    return result


def _get_pdf_meta(file_path: str) -> tuple:
    """Extract bookmarks (outline) and metadata from a PDF using pypdf."""
    outline: list = []
    meta: dict = {}
    try:
        from pypdf import PdfReader as _PdfReader  # type: ignore
        reader = _PdfReader(file_path)
        if reader.metadata:
            rm = reader.metadata
            meta = {
                "title": str(rm.get("/Title") or ""),
                "author": str(rm.get("/Author") or ""),
                "created": str(rm.get("/CreationDate") or ""),
                "modified": str(rm.get("/ModDate") or ""),
            }
        outline = _flatten_pdf_outline(reader, reader.outline)
    except Exception as e:
        logger.debug(f"[PdfParser] 书签/元数据提取失败（非致命）: {e}")
    return outline, meta


def parse_pdf(file_path: str, file_id: str) -> dict[str, Any]:
    """
    提取 PDF 全量文本，供 AI RAG 使用。
    同时返回原始文件的 raw URL，供前端 PDF.js 渲染。

    文字提取回退链：pdfplumber → pypdf → PyPDF2。
    若提取文本过少（扫描件），自动对空页运行 OCR（PyMuPDF + Tesseract）。
    若三个库均不可用，仍返回含 raw_url 的结果（PDF.js 视觉渲染不依赖文字提取）。

    Returns:
        {"text": str, "page_count": int, "raw_url": str,
         "pages": [{"page": int, "text": str}], "ocr_applied": bool,
         "outline": list, "metadata": dict}
    """
    raw_url = f"/api/v1/workspace/raw/{file_id}"
    outline, meta = _get_pdf_meta(file_path)
    pages_text: list[dict[str, Any]] = []
    full_text_parts: list[str] = []
    page_count = 0

    # ── 1. pdfplumber（首选：支持表格提取） ─────────────────────────────────
    try:
        import pdfplumber  # type: ignore

        with pdfplumber.open(file_path) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                try:
                    page_text = page.extract_text() or ""
                except Exception as e:
                    logger.warning(f"[PdfParser/pdfplumber] 第 {i+1} 页文本提取失败: {e}")
                    page_text = ""
                pages_text.append({"page": i + 1, "text": page_text})
                if page_text:
                    full_text_parts.append(page_text)

        return _apply_ocr_fallback(file_path, pages_text, full_text_parts, page_count, raw_url, outline, meta)
    except ImportError:
        logger.info("[PdfParser] pdfplumber 未安装，尝试 pypdf/PyPDF2")
    except Exception as e:
        logger.warning(f"[PdfParser] pdfplumber 解析失败: {e}，尝试下一库")

    # ── 2. pypdf / PyPDF2 回退 ───────────────────────────────────────────────
    for pkg_name, mod_name in [("pypdf", "pypdf"), ("PyPDF2", "PyPDF2")]:
        try:
            mod = __import__(mod_name)
            PdfReader = getattr(mod, "PdfReader")
            with open(file_path, "rb") as fh:
                reader = PdfReader(fh)
                page_count = len(reader.pages)
                for i, pg in enumerate(reader.pages):
                    try:
                        page_text = pg.extract_text() or ""
                    except Exception as e:
                        logger.warning(f"[PdfParser/{pkg_name}] 第 {i+1} 页文本提取失败: {e}")
                        page_text = ""
                    pages_text.append({"page": i + 1, "text": page_text})
                    if page_text:
                        full_text_parts.append(page_text)

            return _apply_ocr_fallback(file_path, pages_text, full_text_parts, page_count, raw_url, outline, meta)
        except ImportError:
            logger.info(f"[PdfParser] {pkg_name} 未安装，尝试下一库")
            continue
        except Exception as e:
            logger.warning(f"[PdfParser] {pkg_name} 解析失败: {e}，尝试下一库")
            pages_text = []
            full_text_parts = []
            continue

    # ── 3. 所有文字提取库均不可用 — 仍尝试全页 OCR ─────────────────────────
    logger.warning(
        "[PdfParser] pdfplumber / pypdf / PyPDF2 均不可用，尝试纯 OCR。"
        " 建议执行: pip install pdfplumber"
    )
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        page_count = len(doc)
        doc.close()
    except Exception:
        page_count = 0

    if page_count:
        ocr_results = _ocr_pdf_pages(file_path, list(range(page_count)))
        for i in range(page_count):
            t = ocr_results.get(i, "")
            pages_text.append({"page": i + 1, "text": t})
            if t:
                full_text_parts.append(t)
        return {
            "text": "\n\n".join(full_text_parts),
            "page_count": page_count,
            "raw_url": raw_url,
            "pages": pages_text,
            "ocr_applied": bool(ocr_results),
            "outline": outline,
            "metadata": meta,
        }

    return {
        "text": "",
        "page_count": 0,
        "raw_url": raw_url,
        "pages": [],
        "ocr_applied": False,
        "outline": outline,
        "metadata": meta,
    }


def _apply_ocr_fallback(
    file_path: str,
    pages_text: list[dict[str, Any]],
    full_text_parts: list[str],
    page_count: int,
    raw_url: str,
    outline: list | None = None,
    meta: dict | None = None,
) -> dict[str, Any]:
    """
    After normal text extraction: find pages with < _PDF_OCR_THRESHOLD chars,
    run OCR on those pages, and merge the results back.
    """
    scanned_indices = [
        i for i, p in enumerate(pages_text)
        if len(p["text"].strip()) < _PDF_OCR_THRESHOLD
    ]
    ocr_applied = False
    if scanned_indices:
        logger.info(
            f"[PdfParser] 发现 {len(scanned_indices)} 页文本稀少，尝试 OCR: "
            f"页码 {[i+1 for i in scanned_indices]}"
        )
        ocr_results = _ocr_pdf_pages(file_path, scanned_indices)
        for idx, text in ocr_results.items():
            if text:
                pages_text[idx]["text"] = text
                full_text_parts.append(text)
                ocr_applied = True

    return {
        "text": "\n\n".join(full_text_parts),
        "page_count": page_count,
        "raw_url": raw_url,
        "pages": pages_text,
        "ocr_applied": ocr_applied,
        "outline": outline or [],
        "metadata": meta or {},
    }


# ─────────────────────────────────────────────────────────────────────────────
# DOCX 导出: 修改后 HTML → .docx
# ─────────────────────────────────────────────────────────────────────────────


def _css_color_to_hex(value: str) -> str | None:
    """Convert a CSS color string like '#aabbcc' or 'rgb(r,g,b)' to a 6-char hex string."""
    if not value:
        return None
    value = value.strip()
    if value.startswith("#"):
        v = value.lstrip("#")
        if len(v) == 3:
            v = "".join(c * 2 for c in v)
        return v.upper() if len(v) == 6 else None
    if value.startswith("rgb"):
        nums = re.findall(r"\d+", value)
        if len(nums) >= 3:
            return "{:02X}{:02X}{:02X}".format(int(nums[0]), int(nums[1]), int(nums[2]))
    return None


def _parse_css_inline(style_str: str) -> dict[str, str]:
    """Parse a CSS inline style string into a {property: value} dict."""
    result: dict[str, str] = {}
    for part in style_str.split(";"):
        part = part.strip()
        if ":" not in part:
            continue
        prop, _, val = part.partition(":")
        result[prop.strip().lower()] = val.strip()
    return result


def _css_length_to_pt(value: str) -> float | None:
    """Convert a CSS length string (pt, px, em, cm, mm, in) to points.

    Returns None if the value cannot be parsed.
    """
    if not value:
        return None
    value = value.strip()
    try:
        if value.endswith("pt"):
            return float(value[:-2])
        if value.endswith("px"):
            return float(value[:-2]) * 0.75
        if value.endswith("em"):
            return float(value[:-2]) * 12  # 1em ≈ 12pt default
        if value.endswith("cm"):
            return float(value[:-2]) * 28.3465
        if value.endswith("mm"):
            return float(value[:-2]) * 2.83465
        if value.endswith("in"):
            return float(value[:-2]) * 72
    except (ValueError, TypeError):
        return None
    return None


def _set_run_east_asian_font(run: Any, font_name: str) -> None:
    """Set East Asian font on a run via direct XML manipulation."""
    try:
        from docx.oxml.ns import qn as _qn
        rPr = run._element.get_or_add_rPr()
        rFonts = rPr.find(_qn("w:rFonts"))
        if rFonts is None:
            from lxml import etree
            rFonts = etree.SubElement(rPr, _qn("w:rFonts"))
        rFonts.set(_qn("w:eastAsia"), font_name)
    except Exception:
        pass


def _apply_run_inline(run: Any, css: dict[str, str]) -> None:
    """Apply inline CSS to a python-docx Run (font, bold, italic, color, size, etc)."""
    from docx.shared import Pt, RGBColor

    # ── Font family ──────────────────────────────────────────────────────
    ff = css.get("font-family", "")
    if ff:
        # Parse font-family: may be "'Calibri','SimSun'" or "SimSun" etc.
        fonts = [f.strip().strip("'\"") for f in ff.split(",") if f.strip()]
        if fonts:
            run.font.name = fonts[0]
            # If there's a second font, treat it as East Asian font
            if len(fonts) >= 2:
                _set_run_east_asian_font(run, fonts[1])
            else:
                # Check if the single font is a CJK font — also set as eastAsia
                _CJK_FONTS = {
                    "SimHei", "SimSun", "KaiTi", "FangSong",
                    "Microsoft YaHei", "STZhongsong", "STSong", "STHeiti",
                    "STKaiti", "STFangsong", "FZShuSong-Z01", "FZHei-B01",
                    "NSimSun", "DengXian", "YouYuan", "LiSu",
                }
                if fonts[0] in _CJK_FONTS:
                    _set_run_east_asian_font(run, fonts[0])

    # ── Bold / Italic / Underline ────────────────────────────────────────
    if css.get("font-weight") in ("bold", "700", "800", "900"):
        run.bold = True
    if css.get("font-style") == "italic":
        run.italic = True
    td = css.get("text-decoration", "")
    if "underline" in td:
        run.underline = True
    if "line-through" in td:
        run.font.strike = True

    # ── Superscript / Subscript ──────────────────────────────────────────
    va = css.get("vertical-align", "").lower()
    if va == "super":
        run.font.superscript = True
    elif va == "sub":
        run.font.subscript = True

    # ── Text color ───────────────────────────────────────────────────────
    color_hex = _css_color_to_hex(css.get("color", ""))
    if color_hex:
        r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
        run.font.color.rgb = RGBColor(r, g, b)

    # ── Highlight / background color ─────────────────────────────────────
    bg = css.get("background-color", "") or css.get("background", "")
    bg_hex = _css_color_to_hex(bg)
    if bg_hex:
        try:
            from docx.oxml.ns import qn as _qn
            from lxml import etree
            rPr = run._element.get_or_add_rPr()
            shd = rPr.find(_qn("w:shd"))
            if shd is None:
                shd = etree.SubElement(rPr, _qn("w:shd"))
            shd.set(_qn("w:val"), "clear")
            shd.set(_qn("w:color"), "auto")
            shd.set(_qn("w:fill"), bg_hex.upper())
        except Exception:
            pass

    # ── Font size ────────────────────────────────────────────────────────
    fs_pt = _css_length_to_pt(css.get("font-size", ""))
    if fs_pt is not None and fs_pt > 0:
        run.font.size = Pt(fs_pt)

    # ── Letter spacing ───────────────────────────────────────────────────
    ls = css.get("letter-spacing", "")
    ls_pt = _css_length_to_pt(ls)
    if ls_pt is not None:
        try:
            from docx.oxml.ns import qn as _qn
            rPr = run._element.get_or_add_rPr()
            spacing = rPr.find(_qn("w:spacing"))
            if spacing is None:
                from lxml import etree
                spacing = etree.SubElement(rPr, _qn("w:spacing"))
            # OOXML spacing is in half-points (1/144 inch)
            spacing.set(_qn("w:val"), str(int(ls_pt * 20)))
        except Exception:
            pass


def _apply_para_format(para: Any, pcss: dict[str, str]) -> None:
    """Apply paragraph-level CSS properties to a python-docx Paragraph."""
    from docx.shared import Pt, Twips
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    # Text alignment
    align = pcss.get("text-align", "")
    if align == "center":
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    elif align == "right":
        para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    elif align == "justify":
        para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    elif align == "left":
        para.alignment = WD_ALIGN_PARAGRAPH.LEFT

    pf = para.paragraph_format

    # Space before (margin-top)
    mt_pt = _css_length_to_pt(pcss.get("margin-top", ""))
    if mt_pt is not None and mt_pt >= 0:
        pf.space_before = Pt(mt_pt)

    # Space after (margin-bottom)
    mb_pt = _css_length_to_pt(pcss.get("margin-bottom", ""))
    if mb_pt is not None and mb_pt >= 0:
        pf.space_after = Pt(mb_pt)

    # Line spacing
    lh = pcss.get("line-height", "")
    if lh:
        try:
            # Only use _css_length_to_pt if the value has a unit suffix
            has_unit = any(lh.rstrip().endswith(u) for u in ("pt", "px", "em", "cm", "mm", "in"))
            if has_unit:
                lh_pt = _css_length_to_pt(lh)
                if lh_pt is not None and lh_pt > 0:
                    from docx.enum.text import WD_LINE_SPACING
                    pf.line_spacing_rule = WD_LINE_SPACING.EXACTLY
                    pf.line_spacing = Pt(lh_pt)
            else:
                # Unitless multiplier (e.g. "1.15", "1.5")
                pf.line_spacing = float(lh)
        except (ValueError, TypeError):
            pass

    # First-line indent (text-indent)
    ti_pt = _css_length_to_pt(pcss.get("text-indent", ""))
    if ti_pt is not None and ti_pt > 0:
        pf.first_line_indent = Pt(ti_pt)

    # Left indent (padding-left or margin-left)
    pl = pcss.get("padding-left", "") or pcss.get("margin-left", "")
    pl_pt = _css_length_to_pt(pl)
    if pl_pt is not None and pl_pt > 0:
        pf.left_indent = Pt(pl_pt)

    # Right indent (padding-right or margin-right)
    pr = pcss.get("padding-right", "") or pcss.get("margin-right", "")
    pr_pt = _css_length_to_pt(pr)
    if pr_pt is not None and pr_pt > 0:
        pf.right_indent = Pt(pr_pt)


def _add_paragraph_from_tag(doc: Any, tag: Any) -> None:
    """Convert a single <p>/<h1>-<h6>/<li> BS4 tag into a python-docx paragraph."""
    from docx.shared import Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    tag_name = tag.name.lower()
    if tag_name in ("h1", "h2", "h3", "h4", "h5", "h6"):
        level = int(tag_name[1])
        para = doc.add_heading("", level=level)
    else:
        para = doc.add_paragraph()

    # Paragraph-level formatting from style attribute
    style_str = tag.get("style", "")
    pcss: dict[str, str] = {}
    if style_str:
        pcss = _parse_css_inline(style_str)
        _apply_para_format(para, pcss)

    # Walk inline children (text nodes + <strong>, <em>, <span>, <a>, <u>)
    for child in tag.children:
        if hasattr(child, "name") and child.name is None:
            # NavigableString
            text = str(child)
            if text:
                run = para.add_run(text)
                # Apply paragraph-level font info to bare text nodes
                _apply_run_inline(run, pcss)
        elif hasattr(child, "name"):
            name = child.name.lower() if child.name else ""
            # Inline image — void element with no text; insert as inline picture
            if name == "img":
                src = child.get("src", "")
                img_bytes = _resolve_img_src(src)
                if img_bytes:
                    try:
                        from docx.shared import Inches, Pt
                        run = para.add_run()
                        
                        img_css = _parse_css_inline(child.get("style", ""))
                        w_str = img_css.get("width", child.get("width", ""))
                        w_val = None
                        if w_str:
                            try:
                                if w_str.endswith("px"): w_val = Pt(float(w_str.replace("px", "")) * 0.75)
                                elif w_str.endswith("in"): w_val = Inches(float(w_str.replace("in", "")))
                                elif w_str.isdigit(): w_val = Pt(float(w_str) * 0.75)
                            except: pass
                        
                        if w_val:
                            run.add_picture(io.BytesIO(img_bytes), width=w_val)
                        else:
                            run.add_picture(io.BytesIO(img_bytes), width=Inches(3))
                    except Exception:
                        pass
                continue
            text = child.get_text()
            if not text:
                continue
            run = para.add_run(text)
            child_css: dict[str, str] = _parse_css_inline(child.get("style", ""))
            # Merge parent paragraph CSS as fallback for font properties
            merged_css = dict(pcss)
            merged_css.update(child_css)
            if name in ("strong", "b"):
                run.bold = True
            if name in ("em", "i"):
                run.italic = True
            if name == "u":
                run.underline = True
            if name == "s":
                run.font.strike = True
            _apply_run_inline(run, merged_css)


def _set_cell_shading(cell: Any, fill_hex: str) -> None:
    """Set the background shading of a python-docx table cell via direct XML."""
    from lxml import etree

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove any existing shd element
    for shd in tcPr.findall(
        "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}shd"
    ):
        tcPr.remove(shd)
    WNS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    shd = etree.SubElement(tcPr, f"{{{WNS}}}shd")
    shd.set(f"{{{WNS}}}val", "clear")
    shd.set(f"{{{WNS}}}color", "auto")
    shd.set(f"{{{WNS}}}fill", fill_hex.upper())


def _apply_border_xml(tcBorders: Any, side: str, val: str, sz: str, color: str) -> None:
    """Helper to add/update an individual border to tcBorders xml."""
    from lxml import etree
    WNS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    el = tcBorders.find(f"{{{WNS}}}{side}")
    if el is None:
        el = etree.SubElement(tcBorders, f"{{{WNS}}}{side}")
    if val == "none":
        el.set(f"{{{WNS}}}val", "none")
    else:
        el.set(f"{{{WNS}}}val", "single")
        el.set(f"{{{WNS}}}sz", sz)
        if color:
            el.set(f"{{{WNS}}}color", color.upper())

def _set_cell_borders(cell: Any, css: dict[str, str]) -> None:
    """Apply borders to a cell by parsing border/border-* CSS."""
    from lxml import etree
    import re

    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    WNS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    tcBorders = tcPr.find(f"{{{WNS}}}tcBorders")
    if tcBorders is None:
        tcBorders = etree.SubElement(tcPr, f"{{{WNS}}}tcBorders")

    # Helper to parse a shorthand border value like "1px solid #000000"
    def parse_border(val_str: str) -> tuple[str, str, str]:
        if not val_str or val_str == "none":
            return "none", "0", ""
        # defaults
        sz, color = "4", "auto" 
        parts = val_str.split()
        for p in parts:
            if "px" in p:
                try: # 1px is approx 4 eighths-of-a-point
                    sz = str(int(float(p.replace("px", "")) * 4))
                except Exception: pass
            elif "pt" in p:
                try: # 1pt is 8 eighths-of-a-point
                    sz = str(int(float(p.replace("pt", "")) * 8))
                except Exception: pass
            elif p.startswith("#") or p.startswith("rgb"):
                h = _css_color_to_hex(p)
                if h: color = h.upper()
        return "single", sz, color

    # Apply general border
    if "border" in css:
        val, sz, color = parse_border(css["border"])
        for side in ("top", "left", "bottom", "right", "insideH", "insideV"):
            _apply_border_xml(tcBorders, side, val, sz, color)
    
    # Override sides
    # We map CSS border sides to Word sides
    for css_prop, w_side in [
        ("border-top", "top"),
        ("border-bottom", "bottom"),
        ("border-left", "left"),
        ("border-right", "right"),
    ]:
        if css_prop in css:
            val, sz, color = parse_border(css[css_prop])
            _apply_border_xml(tcBorders, w_side, val, sz, color)



def _resolve_img_src(src: str) -> bytes | None:
    """Decode a base64 data URI or read a /api/v1/workspace/tmp_image/ temp file.

    Handles:
    - ``data:image/png;base64,<b64>`` — decodes the base64 payload
    - ``/api/v1/workspace/tmp_image/<session_id>/<filename>`` — reads the file
      from the session-scoped tmp directory on disk (written by upload_image)

    Returns the raw image bytes, or *None* on any failure.
    """
    if not src:
        return None
    # Case 1: inline base64 data URI
    if src.startswith("data:image/"):
        try:
            _, _, b64_str = src.partition(",")
            return base64.b64decode(b64_str)
        except Exception:
            return None
    # Case 2: server-side temp image URL (from /api/v1/workspace/upload_image)
    if src.startswith("/api/v1/workspace/tmp_image/"):
        parts = src.split("/")
        # expected: ['', 'api', 'v1', 'workspace', 'tmp_image', '{sid}', '{fname}']
        if len(parts) == 7:
            session_id, filename = parts[5], parts[6]
            if (len(session_id) == 32
                    and all(c in "0123456789abcdef" for c in session_id)
                    and "/" not in filename
                    and "\\" not in filename
                    and ".." not in filename):
                img_path = Path("workspace") / "tmp" / session_id / "images" / filename
                if img_path.is_file():
                    try:
                        return img_path.read_bytes()
                    except Exception:
                        pass
    return None


def _insert_block_image(doc: Any, tag: Any) -> None:
    """Insert an image into the document from an ``<img>`` or ``<figure>`` tag.

    Finds the ``<img>`` element, resolves its ``src`` attribute via
    :func:`_resolve_img_src`, and calls ``doc.add_picture()`` to embed it at
    5 inches wide (constrained by the page width with aspect ratio preserved).
    Silently skips unresolvable images.
    """
    from docx.shared import Inches

    img_el = (
        tag
        if (hasattr(tag, "name") and tag.name and tag.name.lower() == "img")
        else (tag.find("img") if hasattr(tag, "find") else None)
    )
    if not img_el:
        return
    src = img_el.get("src", "")
    img_bytes = _resolve_img_src(src)
    if not img_bytes:
        return
    
    img_css = _parse_css_inline(img_el.get("style", ""))
    w_str = img_css.get("width", img_el.get("width", ""))
    from docx.shared import Pt
    w_val = None
    if w_str:
        try:
            if w_str.endswith("px"): w_val = Pt(float(w_str.replace("px", "")) * 0.75)
            elif w_str.endswith("in"): w_val = Inches(float(w_str.replace("in", "")))
            elif w_str.isdigit(): w_val = Pt(float(w_str) * 0.75)
        except: pass

    try:
        if w_val:
            doc.add_picture(io.BytesIO(img_bytes), width=w_val)
        else:
            doc.add_picture(io.BytesIO(img_bytes), width=Inches(5))
    except Exception as exc:
        logger.debug("[_insert_block_image] 图片插入跳过: %s", exc)


def _setup_blank_doc_defaults(doc: Any) -> None:
    """Configure default styles for a blank python-docx Document."""
    from docx.shared import Pt
    from docx.oxml.ns import qn as _qn

    doc_style = doc.styles['Normal']
    doc_style.font.name = 'Calibri'
    doc_style.font.size = Pt(10.5)  # 五号 — standard Chinese document size
    try:
        rPr = doc_style.element.get_or_add_rPr()
        rFonts = rPr.find(_qn('w:rFonts'))
        if rFonts is None:
            from lxml import etree
            rFonts = etree.SubElement(rPr, _qn('w:rFonts'))
        rFonts.set(_qn('w:eastAsia'), 'DengXian')
    except Exception:
        pass
    # Remove default empty paragraph
    for para in list(doc.paragraphs):
        p = para._element
        p.getparent().remove(p)


def _extract_docx_save_parts(docx_input: Any) -> tuple[str, dict[str, Any]]:
    if isinstance(docx_input, dict):
        html_content = docx_input.get("html")
        if html_content is None:
            html_content = docx_input.get("body_html", "")
        payload = {
            "header_html": str(docx_input.get("header_html") or ""),
            "footer_html": str(docx_input.get("footer_html") or ""),
            "sections": docx_input.get("sections") if isinstance(docx_input.get("sections"), list) else [],
        }
        return str(html_content or ""), payload
    return str(docx_input or ""), {"header_html": "", "footer_html": "", "sections": []}


def _clear_block_container(container: Any) -> None:
    element = getattr(container, "_element", None)
    if element is None:
        return
    for child in list(element):
        element.remove(child)


def _apply_run_marks(run: Any, marks: dict[str, bool] | None) -> None:
    marks = marks or {}
    if marks.get("bold"):
        run.bold = True
    if marks.get("italic"):
        run.italic = True
    if marks.get("underline"):
        run.underline = True
    if marks.get("strike"):
        run.font.strike = True
    if marks.get("superscript"):
        run.font.superscript = True
    if marks.get("subscript"):
        run.font.subscript = True


def _add_page_number_field(para: Any, css: dict[str, str] | None = None, marks: dict[str, bool] | None = None) -> None:
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    def _styled_run(text: str = ""):
        run = para.add_run(text)
        _apply_run_marks(run, marks)
        _apply_run_inline(run, css or {})
        return run

    begin_run = _styled_run()
    fld_begin = OxmlElement("w:fldChar")
    fld_begin.set(qn("w:fldCharType"), "begin")
    begin_run._r.append(fld_begin)

    instr_run = _styled_run()
    instr = OxmlElement("w:instrText")
    instr.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")
    instr.text = " PAGE "
    instr_run._r.append(instr)

    sep_run = _styled_run()
    fld_sep = OxmlElement("w:fldChar")
    fld_sep.set(qn("w:fldCharType"), "separate")
    sep_run._r.append(fld_sep)

    _styled_run("1")

    end_run = _styled_run()
    fld_end = OxmlElement("w:fldChar")
    fld_end.set(qn("w:fldCharType"), "end")
    end_run._r.append(fld_end)


def _append_inline_html_to_paragraph(para: Any, node: Any, inherited_css: dict[str, str] | None = None, marks: dict[str, bool] | None = None) -> None:
    from bs4.element import NavigableString, Tag
    from docx.shared import Inches, Pt

    inherited_css = dict(inherited_css or {})
    marks = dict(marks or {})

    for child in list(getattr(node, "children", [])):
        if isinstance(child, NavigableString):
            text = str(child)
            if not text:
                continue
            run = para.add_run(text)
            _apply_run_marks(run, marks)
            _apply_run_inline(run, inherited_css)
            continue

        if not isinstance(child, Tag):
            continue

        name = child.name.lower() if child.name else ""
        child_css = dict(inherited_css)
        child_css.update(_parse_css_inline(child.get("style", "")))
        child_marks = dict(marks)

        if name in ("strong", "b"):
            child_marks["bold"] = True
        if name in ("em", "i"):
            child_marks["italic"] = True
        if name == "u":
            child_marks["underline"] = True
        if name == "s":
            child_marks["strike"] = True
        if name == "sup":
            child_marks["superscript"] = True
        if name == "sub":
            child_marks["subscript"] = True

        if name == "br":
            para.add_run().add_break()
            continue

        if "koto-hdr-page-num" in (child.get("class") or []):
            _add_page_number_field(para, child_css, child_marks)
            continue

        if name == "img":
            src = child.get("src", "")
            img_bytes = _resolve_img_src(src)
            if img_bytes:
                try:
                    run = para.add_run()
                    _apply_run_marks(run, child_marks)
                    img_css = _parse_css_inline(child.get("style", ""))
                    w_str = img_css.get("width", child.get("width", ""))
                    w_val = None
                    if w_str:
                        try:
                            if w_str.endswith("px"):
                                w_val = Pt(float(w_str.replace("px", "")) * 0.75)
                            elif w_str.endswith("in"):
                                w_val = Inches(float(w_str.replace("in", "")))
                            elif w_str.isdigit():
                                w_val = Pt(float(w_str) * 0.75)
                        except Exception:
                            pass
                    if w_val:
                        run.add_picture(io.BytesIO(img_bytes), width=w_val)
                    else:
                        run.add_picture(io.BytesIO(img_bytes), width=Inches(2.5))
                except Exception:
                    pass
            continue

        _append_inline_html_to_paragraph(para, child, child_css, child_marks)


def _configure_header_footer_tabs(para: Any, section: Any) -> None:
    try:
        from docx.enum.text import WD_TAB_ALIGNMENT, WD_TAB_LEADER
        from docx.shared import Emu

        usable_width = int(section.page_width) - int(section.left_margin) - int(section.right_margin)
        if usable_width <= 0:
            return
        tab_stops = para.paragraph_format.tab_stops
        tab_stops.add_tab_stop(Emu(usable_width // 2), WD_TAB_ALIGNMENT.CENTER, WD_TAB_LEADER.SPACES)
        tab_stops.add_tab_stop(Emu(usable_width), WD_TAB_ALIGNMENT.RIGHT, WD_TAB_LEADER.SPACES)
    except Exception:
        pass


def _add_header_footer_paragraph(container: Any, tag: Any, section: Any) -> None:
    para = container.add_paragraph()
    pcss = _parse_css_inline(tag.get("style", ""))
    if pcss:
        _apply_para_format(para, pcss)

    direct_cols = [
        child for child in list(getattr(tag, "children", []))
        if getattr(child, "name", None) and "koto-hdr-col" in (child.get("class") or [])
    ]
    if direct_cols:
        _configure_header_footer_tabs(para, section)
        for idx, col in enumerate(direct_cols):
            col_css = dict(pcss)
            col_css.update(_parse_css_inline(col.get("style", "")))
            _append_inline_html_to_paragraph(para, col, col_css, {})
            if idx < len(direct_cols) - 1:
                tab_run = para.add_run("\t")
                _apply_run_inline(tab_run, col_css)
        return

    _append_inline_html_to_paragraph(para, tag, pcss, {})


def _write_docx_header_footer_html(container: Any, html_content: str, section: Any) -> None:
    from bs4 import BeautifulSoup

    _clear_block_container(container)
    soup = BeautifulSoup(html_content or "", "html.parser")
    body = soup.find("body") or soup
    wrote_any = False

    for top in body.children:
        if not hasattr(top, "name") or top.name is None:
            text = str(top).strip()
            if text:
                para = container.add_paragraph()
                run = para.add_run(text)
                _apply_run_inline(run, {})
                wrote_any = True
            continue

        tag_name = top.name.lower()
        if tag_name in ("p", "li", "h1", "h2", "h3", "h4", "h5", "h6"):
            _add_header_footer_paragraph(container, top, section)
            wrote_any = True
            continue
        if tag_name in ("ul", "ol"):
            for li in top.find_all("li", recursive=False):
                _add_header_footer_paragraph(container, li, section)
                wrote_any = True
            continue

        text = top.get_text(separator=" ").strip()
        if text:
            para = container.add_paragraph()
            run = para.add_run(text)
            _apply_run_inline(run, {})
            wrote_any = True

    if not wrote_any:
        container.add_paragraph("")


def _payload_value(section_payload: dict[str, Any] | None, key: str, fallback: str = "") -> str:
    if isinstance(section_payload, dict) and key in section_payload:
        return str(section_payload.get(key) or "")
    return fallback


def _export_docx_python(docx_input: Any, original_path: str | None = None) -> bytes:
    """
    Build a .docx from editor HTML using python-docx directly.

    If *original_path* points to a valid .docx file, it is used as a
    template so that the original styles, theme, fonts, headers/footers,
    and section properties are preserved.  The body content is replaced
    with what was edited in the browser.
    """

    try:
        from bs4 import BeautifulSoup
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_ALIGN_VERTICAL
        from docx.oxml.ns import qn
    except ImportError as exc:
        raise RuntimeError(f"python-docx 或 beautifulsoup4 未安装: {exc}") from exc

    html_content, docx_payload = _extract_docx_save_parts(docx_input)
    default_header_html = docx_payload.get("header_html", "")
    default_footer_html = docx_payload.get("footer_html", "")
    sections_payload = docx_payload.get("sections") if isinstance(docx_payload.get("sections"), list) else []

    # ── Open original as template, or create blank ────────────────────────
    if original_path and os.path.isfile(original_path):
        try:
            doc = Document(original_path)
            # Clear body content but keep section properties
            body = doc.element.body
            sect_pr = body.findall(qn('w:sectPr'))
            for child in list(body):
                if child.tag != qn('w:sectPr'):
                    body.remove(child)
            logger.info("[export_docx] using original DOCX as template: %s", original_path)
        except Exception as exc:
            logger.warning("[export_docx] failed to open original (%s), creating blank", exc)
            doc = Document()
            _setup_blank_doc_defaults(doc)
    else:
        doc = Document()
        _setup_blank_doc_defaults(doc)

    soup = BeautifulSoup(html_content, "html.parser")
    body = soup.find("body") or soup

    for top in body.children:
        if not hasattr(top, "name") or top.name is None:
            # bare text node
            text = str(top).strip()
            if text:
                doc.add_paragraph(text)
            continue

        tag_name = top.name.lower()

        # ── Images (standalone <img> or <figure>) ────────────────────────
        if tag_name in ("img", "figure"):
            _insert_block_image(doc, top)
            continue

        # ── Block-level text elements ──────────────────────────────────────
        if tag_name in ("h1", "h2", "h3", "h4", "h5", "h6", "li"):
            _add_paragraph_from_tag(doc, top)
            continue

        if tag_name == "p":
            # <p><img .../></p> — image-only paragraph → embed as picture block
            if not top.get_text().strip() and top.find("img"):
                _insert_block_image(doc, top.find("img"))
            else:
                _add_paragraph_from_tag(doc, top)
            continue

        if tag_name in ("ul", "ol"):
            for li in top.find_all("li", recursive=False):
                _add_paragraph_from_tag(doc, li)
            continue

        # ── Tables ─────────────────────────────────────────────────────────
        if tag_name == "table":
            rows_tags = top.find_all("tr")
            if not rows_tags:
                continue

            # Determine table dimensions
            col_count = 0
            for r in rows_tags:
                c = 0
                for cell in r.find_all(["td", "th"], recursive=False):
                    c += int(cell.get("colspan", 1))
                col_count = max(col_count, c)
            if col_count == 0:
                continue

            tbl = doc.add_table(rows=len(rows_tags), cols=col_count)
            tbl.style = "Table Grid"

            # Track which (row, col) positions are already consumed by a
            # rowspan merge so we skip writing to them again.
            occupied: dict[tuple[int, int], bool] = {}

            for ri, row_tag in enumerate(rows_tags):
                cell_tags = row_tag.find_all(["td", "th"], recursive=False)
                ci = 0  # logical col cursor
                for cell_tag in cell_tags:
                    # Advance past cells already filled by a previous rowspan
                    while occupied.get((ri, ci)):
                        ci += 1

                    colspan = int(cell_tag.get("colspan", 1))
                    rowspan = int(cell_tag.get("rowspan", 1))

                    # Guard against out-of-range
                    if ci >= col_count:
                        break

                    # Mark all cells this span covers as occupied
                    for dr in range(rowspan):
                        for dc in range(colspan):
                            if dr == 0 and dc == 0:
                                continue
                            occupied[(ri + dr, ci + dc)] = True

                    tbl_cell = tbl.cell(ri, ci)

                    # ── Merge ────────────────────────────────────────────
                    end_r = min(ri + rowspan - 1, len(rows_tags) - 1)
                    end_c = min(ci + colspan - 1, col_count - 1)
                    if end_r > ri or end_c > ci:
                        tbl_cell.merge(tbl.cell(end_r, end_c))

                    # ── Cell text ────────────────────────────────────────
                    # Clear the default empty paragraph python-docx adds
                    for p_el in list(tbl_cell.paragraphs):
                        p_el._element.getparent().remove(p_el._element)

                    cell_css = _parse_css_inline(cell_tag.get("style", ""))

                    # Vertical alignment
                    v_align = cell_css.get("vertical-align", "").lower()
                    if v_align == "middle":
                        tbl_cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
                    elif v_align == "bottom":
                        tbl_cell.vertical_alignment = WD_ALIGN_VERTICAL.BOTTOM

                    # Bold if <th>
                    is_header = cell_tag.name.lower() == "th"

                    # Check if cell contains <p> sub-elements (TipTap style)
                    p_tags = cell_tag.find_all(["p", "h1", "h2", "h3", "h4", "h5", "h6"], recursive=False)
                    if p_tags:
                        # Multi-paragraph cell: each <p> becomes a paragraph
                        for pi, p_tag in enumerate(p_tags):
                            cell_para = tbl_cell.add_paragraph()
                            p_css = _parse_css_inline(p_tag.get("style", ""))
                            # Merge cell_css as fallback
                            merged_p = dict(cell_css)
                            merged_p.update(p_css)
                            _apply_para_format(cell_para, merged_p)
                            for child in p_tag.children:
                                if not hasattr(child, "name") or child.name is None:
                                    text = str(child)
                                    if text:
                                        run = cell_para.add_run(text)
                                        if is_header:
                                            run.bold = True
                                        _apply_run_inline(run, merged_p)
                                else:
                                    child_name = child.name.lower() if child.name else ""
                                    text = child.get_text()
                                    if not text:
                                        continue
                                    run = cell_para.add_run(text)
                                    if is_header:
                                        run.bold = True
                                    child_css_inner = _parse_css_inline(child.get("style", ""))
                                    merged = dict(merged_p)
                                    merged.update(child_css_inner)
                                    if child_name in ("strong", "b"):
                                        run.bold = True
                                    if child_name in ("em", "i"):
                                        run.italic = True
                                    if child_name == "u":
                                        run.underline = True
                                    if child_name == "s":
                                        run.font.strike = True
                                    _apply_run_inline(run, merged)
                    else:
                        # Flat cell content (no <p> wrappers)
                        cell_para = tbl_cell.add_paragraph()

                        # Text alignment from cell style
                        align = cell_css.get("text-align", "")
                        if align == "center":
                            cell_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        elif align == "right":
                            cell_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

                        for child in cell_tag.children:
                            if not hasattr(child, "name") or child.name is None:
                                text = str(child).strip()
                                if text:
                                    run = cell_para.add_run(text)
                                    if is_header:
                                        run.bold = True
                                    _apply_run_inline(run, cell_css)
                            else:
                                child_name = child.name.lower() if child.name else ""
                                text = child.get_text()
                                if not text:
                                    continue
                                run = cell_para.add_run(text)
                                if is_header:
                                    run.bold = True
                                child_css_inner = _parse_css_inline(child.get("style", ""))
                                merged = dict(cell_css)
                                merged.update(child_css_inner)
                                if child_name in ("strong", "b"):
                                    run.bold = True
                                if child_name in ("em", "i"):
                                    run.italic = True
                                if child_name == "u":
                                    run.underline = True
                                if child_name == "s":
                                    run.font.strike = True
                                _apply_run_inline(run, merged)

                    # ── Cell background colour ───────────────────────────
                    bg = cell_css.get("background-color", "") or cell_css.get("background", "")
                    bg_hex = _css_color_to_hex(bg)
                    if not bg_hex and is_header:
                        bg_hex = "EEF1F8"  # default header shading
                    if bg_hex:
                        _set_cell_shading(tbl_cell, bg_hex)

                    # Default borders for all tables (often lost in HTML conversion)
                    if not any(k.startswith("border") for k in cell_css):
                        _set_cell_borders(tbl_cell, {"border": "1px solid #000000"})
                    else:
                        _set_cell_borders(tbl_cell, cell_css)
                    
                    # ── Border colour override (legacy) ──────────────────
                    border_raw = cell_css.get("border-color", "")
                    border_hex = _css_color_to_hex(border_raw)
                    if border_hex:
                        _set_cell_borders(tbl_cell, border_hex)

                    ci += colspan

            # ── Column widths from <col> or first row <td style="width:..."> ──
            col_tags = top.find_all("col")
            if not col_tags and rows_tags:
                col_tags = rows_tags[0].find_all(["td", "th"], recursive=False)
            
            if col_tags:
                # Force table to not auto-fit so widths are respected
                try:
                    tbl.allow_autofit = False
                    from docx.oxml.ns import qn
                    tblPr = tbl._element.xpath('w:tblPr')
                    if tblPr:
                        tblW = tblPr[0].xpath('w:tblW')
                        if tblW:
                            tblW[0].set(qn('w:type'), 'pct')
                            tblW[0].set(qn('w:w'), '5000') # 100%
                except Exception:
                    pass

                for idx, col_tag in enumerate(col_tags[:col_count]):
                    w_css = _parse_css_inline(col_tag.get("style", ""))
                    width_str = w_css.get("width", col_tag.get("width", ""))
                    if width_str:
                        try:
                            if width_str.endswith("px"):
                                emu = int(float(width_str[:-2]) * 9144)  # 1px ≈ 9144 EMU
                                for r in range(len(rows_tags)):
                                    tbl.cell(r, idx).width = emu
                            elif width_str.endswith("%"):
                                pct_val = float(width_str[:-1])
                                w_val = int(pct_val * 50)  # 50 = 1%
                                for r in range(len(rows_tags)):
                                    try:
                                        tc = tbl.cell(r, idx)._tc
                                        tcW = tc.get_or_add_tcPr().get_or_add_tcW()
                                        tcW.type = 'pct'
                                        tcW.w = w_val
                                    except Exception:
                                        pass
                        except (ValueError, IndexError):
                            pass

            continue  # done with this table

        # ── Any other block: just extract its text ─────────────────────────
        text = top.get_text(separator=" ").strip()
        if text:
            doc.add_paragraph(text)

    try:
        has_even = any(
            _payload_value(section_payload, "even_header_html") or _payload_value(section_payload, "even_footer_html")
            for section_payload in sections_payload
            if isinstance(section_payload, dict)
        )
        if hasattr(doc.settings, "odd_and_even_pages_header_footer"):
            doc.settings.odd_and_even_pages_header_footer = bool(has_even)
    except Exception:
        has_even = False

    for idx, section in enumerate(doc.sections):
        section_payload = sections_payload[idx] if idx < len(sections_payload) and isinstance(sections_payload[idx], dict) else None
        section_header_html = _payload_value(section_payload, "header_html", default_header_html)
        section_footer_html = _payload_value(section_payload, "footer_html", default_footer_html)
        first_header_html = _payload_value(section_payload, "first_header_html", "")
        first_footer_html = _payload_value(section_payload, "first_footer_html", "")
        even_header_html = _payload_value(section_payload, "even_header_html", "")
        even_footer_html = _payload_value(section_payload, "even_footer_html", "")

        try:
            section.header.is_linked_to_previous = False
        except Exception:
            pass
        try:
            section.footer.is_linked_to_previous = False
        except Exception:
            pass
        _write_docx_header_footer_html(section.header, section_header_html, section)
        _write_docx_header_footer_html(section.footer, section_footer_html, section)

        try:
            section.different_first_page_header_footer = bool(first_header_html or first_footer_html)
        except Exception:
            pass
        try:
            section.first_page_header.is_linked_to_previous = False
        except Exception:
            pass
        try:
            section.first_page_footer.is_linked_to_previous = False
        except Exception:
            pass
        _write_docx_header_footer_html(section.first_page_header, first_header_html, section)
        _write_docx_header_footer_html(section.first_page_footer, first_footer_html, section)

        if has_even:
            try:
                section.even_page_header.is_linked_to_previous = False
            except Exception:
                pass
            try:
                section.even_page_footer.is_linked_to_previous = False
            except Exception:
                pass
            _write_docx_header_footer_html(section.even_page_header, even_header_html, section)
            _write_docx_header_footer_html(section.even_page_footer, even_footer_html, section)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def export_docx(docx_input: Any, original_path: str | None = None) -> bytes:
    """
    将编辑器产出的 HTML 转换为 .docx 字节流。

    如果提供 original_path（原始 .docx 文件路径），则以该文件为模板，
    保留样式、主题、字体、页眉页脚等原始格式，仅替换正文内容。

    主路径：_export_docx_python — 基于 python-docx 直接构建，完整支持：
      - 合并单元格 (colspan/rowspan)
      - 单元格背景色、边框色
      - 列宽、对齐方式
      - 行内格式（加粗、斜体、颜色、字号）

    备用路径：html2docx（当 python-docx/bs4 导入失败时）

    Returns:
        bytes — .docx 文件内容
    """
    html_content, _ = _extract_docx_save_parts(docx_input)
    logger.info(
        "[export_docx] html_content length=%d original=%s preview=%.200s",
        len(html_content or ""),
        original_path or "(none)",
        (html_content or "")[:200],
    )

    # ── Primary: rich python-docx builder ──────────────────────────────────
    try:
        return _export_docx_python(docx_input, original_path=original_path)
    except Exception as exc:
        logger.warning("[export_docx] python-docx builder failed (%s), falling back to html2docx", exc)

    # ── Fallback: html2docx (handles edge cases the builder misses) ─────────
    try:
        from html2docx import html2docx

        wrapped_html = (
            '<html><head><meta charset="utf-8">'
            '<style>body{font-family:"Microsoft YaHei","PingFang SC","SimHei",Arial,sans-serif;'
            "font-size:11pt;line-height:1.6;}"
            "h1{font-size:20pt}h2{font-size:16pt}h3{font-size:14pt}"
            "</style></head><body>" + html_content + "</body></html>"
        )
        buf = html2docx(wrapped_html, title="Koto 导出文档")
        buf.seek(0)
        data = buf.read()
        if data:
            return data
    except ImportError:
        pass

    raise RuntimeError("export_docx: 所有路径均失败，请确认 python-docx 和 beautifulsoup4 已安装")


# ─────────────────────────────────────────────────────────────────────────────
# XLSX 导出: Luckysheet JSON → .xlsx
# ─────────────────────────────────────────────────────────────────────────────


def export_xlsx(
    sheets_json: Any, images: list[dict] | None = None
) -> bytes:
    """
    将编辑器序列化数据重建为 .xlsx 字节流。

    支持两种输入格式:
    - Univer IWorkbookData (dict): {sheetOrder:[...], sheets:{id:{name, cellData:{row:{col:{v,...}}}}}}
    - Luckysheet legacy (list):  [{name, celldata:[{r,c,v:{v:...}}]}]

    副加 images (前端 overlay) 嵌入到第一个 sheet。

    Returns:
        bytes — .xlsx 文件内容
    """
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl 未安装")

    wb = openpyxl.Workbook()
    wb.remove(wb.active)  # 删除默认空 sheet

    # ── Detect format ────────────────────────────────────────────────────────
    is_univer = isinstance(sheets_json, dict) and "sheetOrder" in sheets_json

    if is_univer:
        # Univer IWorkbookData format
        sheet_order = sheets_json.get("sheetOrder", [])
        sheets_map = sheets_json.get("sheets", {})
        for sheet_id in sheet_order:
            sheet_data = sheets_map.get(sheet_id, {})
            ws = wb.create_sheet(title=sheet_data.get("name", "Sheet"))
            cell_data = sheet_data.get("cellData", {})
            for row_key, row_cells in cell_data.items():
                r = int(row_key) + 1  # Univer 0-indexed → openpyxl 1-indexed
                for col_key, cell in row_cells.items():
                    c = int(col_key) + 1
                    if cell and "v" in cell:
                        ws.cell(row=r, column=c, value=cell["v"])
    else:
        # Luckysheet legacy format (list of sheets)
        if not isinstance(sheets_json, list):
            sheets_json = []
        for sheet_data in sheets_json:
            ws = wb.create_sheet(title=sheet_data.get("name", "Sheet"))
            for cell_entry in sheet_data.get("celldata", []):
                r = cell_entry.get("r", 0) + 1  # Luckysheet 0-indexed → openpyxl 1-indexed
                c = cell_entry.get("c", 0) + 1
                v = cell_entry.get("v", {})
                if v:
                    ws.cell(row=r, column=c, value=v.get("v"))

    # Embed overlay images into the first sheet (if any)
    if images and wb.worksheets:
        import base64
        import io as _io

        try:
            from openpyxl.drawing.image import Image as XlImage
        except ImportError:
            XlImage = None
        if XlImage:
            ws_first = wb.worksheets[0]
            for img_data in images:
                src = img_data.get("src", "")
                if not src or not src.startswith("data:image"):
                    continue
                try:
                    # data:image/png;base64,<data>
                    b64 = src.split(",", 1)[1]
                    raw = base64.b64decode(b64)
                    ximg = XlImage(_io.BytesIO(raw))
                    ws_first.add_image(ximg, "A1")
                except Exception:
                    pass  # skip unreadable images

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()


# ─────────────────────────────────────────────────────────────────────────────
# PPTX 导出: 仅替换文字 (保留主题/动画)
# ─────────────────────────────────────────────────────────────────────────────


def export_pptx(original_path: str, slides_json: Any) -> bytes:
    """
    在原始 PPTX 文件上就地替换文字/表格内容，不重建 PPT 结构。
    保留原 PPT 的主题背景、图片、动画等。

    Args:
        original_path: 暂存的原始 .pptx 文件路径
        slides_json: 画布编辑器序列化数据（两种格式）:
          - 新格式 (几何画布): {"slides": [{"slide_index": int, "shapes": [
              {"_type": "TEXT", "id": int, "text": str} |
              {"_type": "TABLE", "id": int, "cells": [{"row", "col", "text"}]}
            ]}]}
          - 旧格式 (文本卡片): [{"slide_index": int, "texts": [{"shape_id": int, "text": str}]}]

    Returns:
        bytes — 修改后的 .pptx 文件内容
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx 未安装")

    prs = Presentation(original_path)

    # 构建 shape lookup: slide_index → shape_id → shape (all shapes, not just text)
    slides_map: dict[int, dict[int, Any]] = {}
    for slide_idx, slide in enumerate(prs.slides):
        shape_map: dict[int, Any] = {shape.shape_id: shape for shape in slide.shapes}
        slides_map[slide_idx] = shape_map

    def _replace_text_frame(tf: Any, new_text: str) -> None:
        """Replace all text in a text frame preserving the first run's style."""
        if not tf.paragraphs:
            return
        first_para = tf.paragraphs[0]
        if first_para.runs:
            first_para.runs[0].text = new_text
            for run in first_para.runs[1:]:
                run.text = ""
        else:
            # No existing runs — add one rather than using para.text= (which destroys XML)
            run = first_para.add_run()
            run.text = new_text
        for para in tf.paragraphs[1:]:
            for run in para.runs:
                run.text = ""

    def _apply_shape_fill(shape: Any, fill_color: str | None) -> None:
        """Apply solid fill colour to a shape. Accepts '#rrggbb' or None."""
        if fill_color is None:
            return
        try:
            from pptx.util import Pt
            from pptx.dml.color import RGBColor
            hex_val = fill_color.lstrip("#")
            if len(hex_val) != 6:
                return
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(
                int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
            )
        except Exception:
            pass

    def _apply_shape_border(shape: Any, border_data: dict | None) -> None:
        """Apply border (line) properties to a shape."""
        if not border_data:
            return
        try:
            from pptx.util import Pt, Emu
            from pptx.dml.color import RGBColor
            width = border_data.get("width", 0)
            color = border_data.get("color", "#000000")
            if width and width > 0:
                shape.line.width = Pt(float(width))
                hex_val = color.lstrip("#")
                if len(hex_val) == 6:
                    shape.line.color.rgb = RGBColor(
                        int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
                    )
        except Exception:
            pass

    def _apply_paragraphs(tf: Any, paragraphs_data: list | None) -> None:
        """Replace text frame content from structured paragraphs with run-level formatting."""
        if not paragraphs_data:
            return
        try:
            from pptx.util import Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            _AL = {
                "LEFT": PP_ALIGN.LEFT, "CENTER": PP_ALIGN.CENTER,
                "RIGHT": PP_ALIGN.RIGHT, "JUSTIFY": PP_ALIGN.JUSTIFY,
            }

            # We iterate paragraphs in parallel: reuse existing para XML elements where possible
            existing_paras = list(tf.paragraphs)
            for pi, p_data in enumerate(paragraphs_data):
                if pi < len(existing_paras):
                    para = existing_paras[pi]
                else:
                    # Add new paragraph by adding a run first
                    para = tf.paragraphs[0]  # fallback, will be overwritten
                    try:
                        from lxml import etree
                        _NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
                        new_p = etree.SubElement(tf._txBody, f"{{{_NS}}}p")
                        # Re-read paragraphs
                        existing_paras = list(tf.paragraphs)
                        if pi < len(existing_paras):
                            para = existing_paras[pi]
                    except Exception:
                        break

                # Apply alignment
                align_str = (p_data.get("align") or "LEFT").upper()
                if align_str in _AL:
                    para.alignment = _AL[align_str]

                # Apply line spacing
                line_spacing = p_data.get("lineSpacing")
                if line_spacing and float(line_spacing) > 0:
                    try:
                        para.line_spacing = float(line_spacing)
                    except Exception:
                        pass

                # Apply runs: replace run text and formatting
                runs_data = p_data.get("runs", [])
                existing_runs = list(para.runs)
                for ri, r_data in enumerate(runs_data):
                    if ri < len(existing_runs):
                        run = existing_runs[ri]
                    else:
                        run = para.add_run()
                    run.text = r_data.get("text", "")
                    try:
                        if "bold" in r_data:
                            run.font.bold = bool(r_data["bold"])
                        if "italic" in r_data:
                            run.font.italic = bool(r_data["italic"])
                        if "underline" in r_data:
                            run.font.underline = bool(r_data["underline"])
                        if "size" in r_data and r_data["size"]:
                            run.font.size = Pt(float(r_data["size"]))
                        if "fontName" in r_data and r_data["fontName"]:
                            run.font.name = r_data["fontName"]
                        if "color" in r_data and r_data["color"]:
                            hex_val = str(r_data["color"]).lstrip("#")
                            if len(hex_val) == 6:
                                run.font.color.rgb = RGBColor(
                                    int(hex_val[0:2], 16), int(hex_val[2:4], 16), int(hex_val[4:6], 16)
                                )
                    except Exception:
                        pass

                # Clear excess existing runs
                for ri in range(len(runs_data), len(existing_runs)):
                    existing_runs[ri].text = ""

            # Clear excess existing paragraphs
            for pi in range(len(paragraphs_data), len(existing_paras)):
                for run in existing_paras[pi].runs:
                    run.text = ""
        except Exception:
            pass

    # Detect format: new geometry canvas dict vs legacy text-card list
    if isinstance(slides_json, dict) and "slides" in slides_json:
        # ── New geometry canvas format ──
        for slide_data in slides_json["slides"]:
            slide_idx = slide_data.get("slide_index", slide_data.get("index", 0))
            shape_map = slides_map.get(slide_idx, {})
            for shape_entry in slide_data.get("shapes", []):
                shape_id = shape_entry.get("id") or shape_entry.get("shape_id")
                shape = shape_map.get(shape_id)
                if shape is None:
                    continue

                # Apply shape fill colour
                if "fill" in shape_entry:
                    _apply_shape_fill(shape, shape_entry["fill"])

                # Apply shape border / line
                if "border" in shape_entry:
                    _apply_shape_border(shape, shape_entry["border"])

                stype = shape_entry.get("_type", "TEXT")
                if stype in ("TEXT", "TEXT_BOX"):
                    if shape.has_text_frame:
                        paras = shape_entry.get("paragraphs")
                        if paras:
                            _apply_paragraphs(shape.text_frame, paras)
                        else:
                            # Fallback: plain text replacement
                            _replace_text_frame(
                                shape.text_frame, shape_entry.get("text", "")
                            )
                elif stype == "TABLE":
                    if shape.has_table:
                        tbl = shape.table
                        for cell_data in shape_entry.get("cells", []):
                            r, c = cell_data.get("row", 0), cell_data.get("col", 0)
                            try:
                                cell = tbl.cell(r, c)
                                if cell.text_frame:
                                    _replace_text_frame(
                                        cell.text_frame, cell_data.get("text", "")
                                    )
                            except Exception:
                                pass
    else:
        # ── Legacy text-card format ──
        for slide_data in slides_json or []:
            slide_idx = slide_data.get("slide_index", 0)
            shape_map = slides_map.get(slide_idx, {})
            for text_entry in slide_data.get("texts", []):
                shape_id = text_entry.get("shape_id")
                shape = shape_map.get(shape_id)
                if shape is None or not shape.has_text_frame:
                    continue
                _replace_text_frame(shape.text_frame, text_entry.get("text", ""))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
