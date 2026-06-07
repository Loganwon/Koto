# -*- coding: utf-8 -*-
# Copyright (C) 2024-2026 Koto AI. All rights reserved.
# SPDX-License-Identifier: LicenseRef-Koto-Proprietary
"""
MultiFileCoordinator — Cross-file Operation Coordinator
========================================================

Manages multi-file operations with:
- File locking for concurrent safety
- Change tracking for undo/redo and frontend sync
- Cross-file data extraction and injection
- Multi-document comparison and analysis

Usage::

    from app.core.file.multi_file_coordinator import MultiFileCoordinator

    coordinator = MultiFileCoordinator()

    # Extract data from Excel and inject into Word
    change = await coordinator.extract_and_inject(
        source=FileHandle(path="data.xlsx"),
        target=FileHandle(path="report.docx"),
        query="从Sheet1提取销售数据",
    )

    # Compare multiple documents
    result = await coordinator.compare_documents(
        files=[file1, file2, file3],
        aspect="content",
    )
"""
from __future__ import annotations

import asyncio
import difflib
import hashlib
import json
import logging
import os
import shutil
import threading
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, TYPE_CHECKING

if TYPE_CHECKING:
    from app.core.agent.doc_agent import FileHandle

logger = logging.getLogger(__name__)


# ============================================================================
# Data Classes
# ============================================================================


@dataclass
class FileChange:
    """Represents a change made to a file."""
    file_path: str
    change_type: str           # add/modify/delete
    range_start: int           # Character offset start
    range_end: int             # Character offset end
    original: str              # Original content
    modified: str              # New content
    timestamp: float = field(default_factory=time.time)
    step_id: str = ""
    operation: str = ""        # e.g., "extract_inject", "compare", "annotate"

    def to_dict(self) -> Dict[str, Any]:
        return {
            "file_path": self.file_path,
            "change_type": self.change_type,
            "range": [self.range_start, self.range_end],
            "original": self.original[:500],
            "modified": self.modified[:500],
            "timestamp": self.timestamp,
            "step_id": self.step_id,
            "operation": self.operation,
        }


@dataclass
class FileSnapshot:
    """Point-in-time snapshot of a file's content."""
    path: str
    content: str
    content_hash: str
    timestamp: float = field(default_factory=time.time)

    @classmethod
    def from_file(cls, path: str) -> "FileSnapshot":
        """Create snapshot from file on disk."""
        content = ""
        if not os.path.exists(path):
            content_hash = hashlib.md5(content.encode()).hexdigest()
            return cls(path=path, content=content, content_hash=content_hash)
        try:
            if path.endswith((".xlsx", ".xls")):
                content = _read_excel_as_text(path)
            elif path.endswith(".docx"):
                content = _read_docx_as_text(path)
            elif path.endswith(".pptx"):
                content = _read_pptx_as_text(path)
            elif path.endswith(".pdf"):
                content = _read_pdf_as_text(path)
            else:
                with open(path, "r", encoding="utf-8", errors="replace") as f:
                    content = f.read()
        except Exception as e:
            logger.warning("[FileSnapshot] Failed to read %s: %s", path, e)

        content_hash = hashlib.md5(content.encode()).hexdigest()
        return cls(path=path, content=content, content_hash=content_hash)


@dataclass
class CompareResult:
    """Result of comparing multiple documents."""
    files: List[str]
    aspect: str                # content/structure/metadata
    similarities: Dict[str, float]  # pairwise similarity scores
    differences: List[Dict[str, Any]]  # list of specific differences
    summary: str
    timestamp: float = field(default_factory=time.time)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "files": self.files,
            "aspect": self.aspect,
            "similarities": self.similarities,
            "differences": self.differences,
            "summary": self.summary,
            "timestamp": self.timestamp,
        }


# ============================================================================
# MultiFileCoordinator
# ============================================================================


class MultiFileCoordinator:
    """
    Coordinates operations across multiple files.

    Thread-safe with file locking. Tracks all changes for
    potential rollback and frontend synchronization.
    """

    def __init__(self, workspace_root: Optional[str] = None):
        self._workspace_root = workspace_root or self._get_default_workspace()
        self._file_locks: Dict[str, asyncio.Lock] = {}
        self._sync_locks: Dict[str, threading.Lock] = {}
        self._change_log: List[FileChange] = []
        self._snapshots: Dict[str, FileSnapshot] = {}
        self._lock = threading.Lock()

    def _get_default_workspace(self) -> str:
        """Get default workspace directory."""
        project_root = Path(__file__).resolve().parent.parent.parent.parent
        return str(project_root / "workspace")

    # ── File Locking ───────────────────────────────────────────────────────

    def _get_async_lock(self, path: str) -> asyncio.Lock:
        """Get or create async lock for a file."""
        with self._lock:
            if path not in self._file_locks:
                self._file_locks[path] = asyncio.Lock()
            return self._file_locks[path]

    def _get_sync_lock(self, path: str) -> threading.Lock:
        """Get or create sync lock for a file."""
        with self._lock:
            if path not in self._sync_locks:
                self._sync_locks[path] = threading.Lock()
            return self._sync_locks[path]

    # ── Snapshot Management ────────────────────────────────────────────────

    def take_snapshot(self, path: str) -> FileSnapshot:
        """Take a snapshot of file's current content."""
        snapshot = FileSnapshot.from_file(path)
        with self._lock:
            self._snapshots[path] = snapshot
        return snapshot

    def get_snapshot(self, path: str) -> Optional[FileSnapshot]:
        """Get last snapshot for a file."""
        return self._snapshots.get(path)

    def has_changed(self, path: str) -> bool:
        """Check if file has changed since last snapshot."""
        old_snapshot = self.get_snapshot(path)
        if not old_snapshot:
            return True
        new_snapshot = FileSnapshot.from_file(path)
        return new_snapshot.content_hash != old_snapshot.content_hash

    # ── Change Tracking ────────────────────────────────────────────────────

    def track_change(
        self,
        path: str,
        original: str,
        modified: str,
        change_type: str = "modify",
        range_start: int = 0,
        range_end: int = 0,
        step_id: str = "",
        operation: str = "",
    ) -> FileChange:
        """Record a file change."""
        change = FileChange(
            file_path=path,
            change_type=change_type,
            range_start=range_start,
            range_end=range_end,
            original=original,
            modified=modified,
            step_id=step_id,
            operation=operation,
        )
        with self._lock:
            self._change_log.append(change)
        return change

    def get_changes(self, path: Optional[str] = None) -> List[FileChange]:
        """Get change log, optionally filtered by path."""
        with self._lock:
            if path:
                return [c for c in self._change_log if c.file_path == path]
            return list(self._change_log)

    def clear_changes(self):
        """Clear the change log."""
        with self._lock:
            self._change_log.clear()

    # ── Cross-file Operations ──────────────────────────────────────────────

    async def extract_and_inject(
        self,
        source: "FileHandle",
        target: "FileHandle",
        query: str,
        insert_position: str = "end",
    ) -> FileChange:
        """
        Extract data from source file and inject into target file.

        Args:
            source: Source file to extract from
            target: Target file to inject into
            query: Description of what to extract
            insert_position: Where to insert ("start", "end", "cursor")

        Returns:
            FileChange describing the modification
        """
        source_lock = self._get_async_lock(source.path)
        target_lock = self._get_async_lock(target.path)

        # Lock both files (ordered to prevent deadlock)
        locks = sorted([
            (source.path, source_lock),
            (target.path, target_lock),
        ], key=lambda x: x[0])

        for _, lock in locks:
            await lock.acquire()

        try:
            # Take snapshots
            source_snapshot = self.take_snapshot(source.path)
            target_snapshot = self.take_snapshot(target.path)

            # Extract data from source
            extracted_data = await self._extract_data(source, query, source_snapshot.content)

            if not extracted_data:
                return FileChange(
                    file_path=target.path,
                    change_type="none",
                    range_start=0,
                    range_end=0,
                    original="",
                    modified="",
                    operation="extract_inject",
                )

            # Inject into target
            new_content, range_start, range_end = self._inject_data(
                target_snapshot.content,
                extracted_data,
                insert_position,
                target.cursor_position,
            )

            # Write the modified content
            await self._write_file(target.path, new_content)

            # Track the change
            change = self.track_change(
                path=target.path,
                original=target_snapshot.content[range_start:range_end] if range_start < range_end else "",
                modified=extracted_data,
                change_type="modify",
                range_start=range_start,
                range_end=range_end,
                operation="extract_inject",
            )

            return change

        finally:
            for _, lock in reversed(locks):
                lock.release()

    async def _extract_data(
        self,
        source: "FileHandle",
        query: str,
        content: str,
    ) -> str:
        """Extract data from source using LLM."""
        try:
            from app.core.workflow_engine import call_llm

            prompt = f"""从以下文件内容中提取数据：

## 提取要求
{query}

## 文件类型
{source.file_type}

## 文件内容
{content[:8000]}

## 输出
直接输出提取的数据，不要添加额外解释。"""

            result = call_llm(prompt)
            return result.strip() if result else ""

        except Exception as e:
            logger.error("[MultiFileCoordinator] Extraction failed: %s", e)
            return ""

    def _inject_data(
        self,
        content: str,
        data: str,
        position: str,
        cursor: int,
    ) -> Tuple[str, int, int]:
        """Inject data into content at specified position."""
        if position == "start":
            return data + "\n" + content, 0, 0
        elif position == "cursor" and 0 <= cursor <= len(content):
            return content[:cursor] + data + content[cursor:], cursor, cursor
        else:  # end
            return content + "\n" + data, len(content), len(content)

    async def _write_file(self, path: str, content: str):
        """Write content to file with backup."""
        # Create backup
        if os.path.exists(path):
            backup_path = path + ".bak"
            shutil.copy2(path, backup_path)

        # Write based on file type
        if path.endswith(".docx"):
            await self._write_docx(path, content)
        elif path.endswith(".xlsx"):
            await self._write_xlsx(path, content)
        else:
            with open(path, "w", encoding="utf-8") as f:
                f.write(content)

    async def _write_docx(self, path: str, content: str):
        """Write content to DOCX file."""
        try:
            from docx import Document
            doc = Document()
            for para in content.split("\n"):
                doc.add_paragraph(para)
            doc.save(path)
        except ImportError:
            logger.warning("[MultiFileCoordinator] python-docx not installed")
            raise

    async def _write_xlsx(self, path: str, content: str):
        """Write content to XLSX file (simplified - appends to first sheet)."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path)
            ws = wb.active
            # Find next empty row
            next_row = ws.max_row + 1
            for i, line in enumerate(content.split("\n")):
                ws.cell(row=next_row + i, column=1, value=line)
            wb.save(path)
            wb.close()
        except ImportError:
            logger.warning("[MultiFileCoordinator] openpyxl not installed")
            raise

    # ── Document Comparison ────────────────────────────────────────────────

    async def compare_documents(
        self,
        files: List["FileHandle"],
        aspect: str = "content",
    ) -> CompareResult:
        """
        Compare multiple documents.

        Args:
            files: List of FileHandle objects to compare
            aspect: What to compare ("content", "structure", "metadata")

        Returns:
            CompareResult with similarities and differences
        """
        if len(files) < 2:
            return CompareResult(
                files=[f.path for f in files],
                aspect=aspect,
                similarities={},
                differences=[],
                summary="至少需要两个文件进行对比",
            )

        # Take snapshots of all files
        snapshots = {}
        for f in files:
            lock = self._get_async_lock(f.path)
            async with lock:
                snapshots[f.path] = self.take_snapshot(f.path)

        # Calculate pairwise similarities
        similarities = {}
        differences = []

        paths = list(snapshots.keys())
        for i in range(len(paths)):
            for j in range(i + 1, len(paths)):
                path1, path2 = paths[i], paths[j]
                content1 = snapshots[path1].content
                content2 = snapshots[path2].content

                # Calculate similarity
                similarity = self._calculate_similarity(content1, content2)
                key = f"{Path(path1).name} vs {Path(path2).name}"
                similarities[key] = similarity

                # Find specific differences
                diffs = self._find_differences(content1, content2, path1, path2)
                differences.extend(diffs)

        # Generate summary
        summary = self._generate_comparison_summary(similarities, differences)

        return CompareResult(
            files=[f.path for f in files],
            aspect=aspect,
            similarities=similarities,
            differences=differences[:20],  # Limit for transport
            summary=summary,
        )

    def _calculate_similarity(self, text1: str, text2: str) -> float:
        """Calculate text similarity ratio (0-1)."""
        if not text1 and not text2:
            return 1.0
        if not text1 or not text2:
            return 0.0

        # Use sequence matcher for similarity
        matcher = difflib.SequenceMatcher(None, text1, text2)
        return round(matcher.ratio(), 3)

    def _find_differences(
        self,
        text1: str,
        text2: str,
        path1: str,
        path2: str,
    ) -> List[Dict[str, Any]]:
        """Find specific differences between two texts."""
        differences = []

        # Use unified diff
        lines1 = text1.splitlines(keepends=True)
        lines2 = text2.splitlines(keepends=True)

        diff = difflib.unified_diff(
            lines1, lines2,
            fromfile=Path(path1).name,
            tofile=Path(path2).name,
            lineterm="",
        )

        current_diff = {"type": None, "lines": []}
        for line in diff:
            if line.startswith("---") or line.startswith("+++"):
                continue
            elif line.startswith("@@"):
                if current_diff["lines"]:
                    differences.append(current_diff)
                current_diff = {"type": "context", "header": line.strip(), "lines": []}
            elif line.startswith("-"):
                if current_diff["type"] != "remove":
                    if current_diff["lines"]:
                        differences.append(current_diff)
                    current_diff = {"type": "remove", "file": path1, "lines": []}
                current_diff["lines"].append(line[1:].strip())
            elif line.startswith("+"):
                if current_diff["type"] != "add":
                    if current_diff["lines"]:
                        differences.append(current_diff)
                    current_diff = {"type": "add", "file": path2, "lines": []}
                current_diff["lines"].append(line[1:].strip())

        if current_diff["lines"]:
            differences.append(current_diff)

        return differences[:50]  # Limit number of differences

    def _generate_comparison_summary(
        self,
        similarities: Dict[str, float],
        differences: List[Dict[str, Any]],
    ) -> str:
        """Generate a human-readable comparison summary."""
        if not similarities:
            return "无可比较的文件对"

        parts = []

        # Similarity summary
        avg_similarity = sum(similarities.values()) / len(similarities)
        parts.append(f"平均相似度: {avg_similarity:.1%}")

        # Most/least similar
        sorted_sims = sorted(similarities.items(), key=lambda x: x[1])
        if len(sorted_sims) > 1:
            parts.append(f"最不相似: {sorted_sims[0][0]} ({sorted_sims[0][1]:.1%})")
            parts.append(f"最相似: {sorted_sims[-1][0]} ({sorted_sims[-1][1]:.1%})")

        # Difference count
        add_count = sum(1 for d in differences if d.get("type") == "add")
        remove_count = sum(1 for d in differences if d.get("type") == "remove")
        if add_count or remove_count:
            parts.append(f"差异: +{add_count} -{remove_count} 处")

        return " | ".join(parts)

    # ── Annotation Support ─────────────────────────────────────────────────

    async def annotate_file(
        self,
        path: str,
        annotations: List[Dict[str, Any]],
    ) -> List[FileChange]:
        """
        Add annotations/comments to a file.

        Args:
            path: File path
            annotations: List of {range_start, range_end, comment, color}

        Returns:
            List of FileChange objects for each annotation
        """
        lock = self._get_async_lock(path)
        async with lock:
            changes = []
            for ann in annotations:
                change = self.track_change(
                    path=path,
                    original="",
                    modified=ann.get("comment", ""),
                    change_type="annotate",
                    range_start=ann.get("range_start", 0),
                    range_end=ann.get("range_end", 0),
                    operation="annotate",
                )
                changes.append(change)
            return changes


# ============================================================================
# File Reading Helpers
# ============================================================================


def _read_excel_as_text(path: str) -> str:
    """Read Excel file as plain text."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            parts.append(f"=== Sheet: {sheet_name} ===")
            ws = wb[sheet_name]
            for row in ws.iter_rows(max_row=200, values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                parts.append("\t".join(cells))
        wb.close()
        return "\n".join(parts)
    except Exception as e:
        logger.warning("[_read_excel_as_text] Error: %s", e)
        return ""


def _read_docx_as_text(path: str) -> str:
    """Read DOCX file as plain text."""
    try:
        from docx import Document
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs]
        return "\n".join(parts)
    except Exception as e:
        logger.warning("[_read_docx_as_text] Error: %s", e)
        return ""


def _read_pptx_as_text(path: str) -> str:
    """Read PPTX file as plain text."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"=== Slide {i+1} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
        return "\n".join(parts)
    except Exception as e:
        logger.warning("[_read_pptx_as_text] Error: %s", e)
        return ""


def _read_pdf_as_text(path: str) -> str:
    """Read PDF file as plain text."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        parts = []
        for page in doc:
            parts.append(page.get_text())
        doc.close()
        return "\n".join(parts)
    except Exception as e:
        logger.warning("[_read_pdf_as_text] Error: %s", e)
        return ""


# ============================================================================
# Singleton instance
# ============================================================================

_coordinator_instance: Optional[MultiFileCoordinator] = None
_coordinator_lock = threading.Lock()


def get_file_coordinator() -> MultiFileCoordinator:
    """Get the singleton MultiFileCoordinator instance."""
    global _coordinator_instance
    if _coordinator_instance is None:
        with _coordinator_lock:
            if _coordinator_instance is None:
                _coordinator_instance = MultiFileCoordinator()
    return _coordinator_instance
