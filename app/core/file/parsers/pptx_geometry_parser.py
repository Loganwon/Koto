# -*- coding: utf-8 -*-
# Copyright (C) 2024-2026 Koto AI. All rights reserved.
# SPDX-License-Identifier: LicenseRef-Koto-Proprietary
from __future__ import annotations

import base64
import io
import logging
import os
import re
from typing import Any

from app.core.file.image_utils import MAX_BLOB_BYTES as _MAX_BLOB_BYTES
from app.core.file.image_utils import compress_image_bytes as _compress_image_bytes

logger = logging.getLogger(__name__)


def parse_pptx_geometry(file_path: Any) -> dict[str, Any]:
    """
    使用 python-pptx 提取每个 Slide 的完整几何数据，供前端画布编辑器渲染。
    包含文本框 (含字体样式)、图片 (base64 data URI)、表格 (单元格文本)、备注。
    支持 GROUP 形状递归、从 layout/master 继承背景色。

    Returns:
        {
          "slide_width_emu": int,
          "slide_height_emu": int,
          "slides": [
            {
              "slide_index": int, "slide_id": int,
              "background": "#xxxxxx", "notes": str,
              "shapes": [
                {
                  "id": int, "name": str,
                  "_type": "TEXT" | "PICTURE" | "TABLE",
                  "left": int, "top": int, "width": int, "height": int,
                  "z_order": int, "fill": str | null,
                  # TEXT:    "has_text": true, "paragraphs": [{"align": str, "runs": [...]}]
                  # PICTURE: "image_b64": "data:image/...;base64,..."
                  # TABLE:   "table_rows": int, "table_cols": int,
                  #          "cells": [{"row": int, "col": int, "text": str}]
                }
              ]
            }
          ]
        }
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        raise RuntimeError("python-pptx 未安装，请执行: pip install python-pptx")

    prs = Presentation(file_path)
    slide_w: int = prs.slide_width or 9144000
    slide_h: int = prs.slide_height or 6858000
    slides_data: list[dict[str, Any]] = []

    # ── Build theme-color lookup table (schemeClr name → "#rrggbb") ──────
    # Many PPTX borders/fills use scheme colors (dk1, lt1, accent1, etc.) instead
    # of explicit RGB, causing line.color.rgb to throw. We pre-resolve them once.
    _theme_colors: dict[str, str] = {}
    try:
        _NS_T = "http://schemas.openxmlformats.org/drawingml/2006/main"
        _prs_part = prs.part
        for _rel in _prs_part.rels.values():
            if "slideMaster" in _rel.reltype:
                _mst_part = _rel.target_part
                for _r2 in _mst_part.rels.values():
                    if "theme" in _r2.reltype:
                        import lxml.etree as _ET

                        _theme_el = _ET.fromstring(_r2.target_part.blob)
                        _cs = _theme_el.find(f".//{{{_NS_T}}}clrScheme")
                        if _cs is not None:
                            for _c in _cs:
                                _name = _c.tag.split("}")[
                                    1
                                ]  # e.g. "dk1", "lt1", "accent1"
                                _srgb = _c.find(f"{{{_NS_T}}}srgbClr")
                                _sys = _c.find(f"{{{_NS_T}}}sysClr")
                                if _srgb is not None:
                                    _theme_colors[_name] = (
                                        "#" + _srgb.get("val", "000000").lower()
                                    )
                                elif _sys is not None:
                                    _last = _sys.get("lastClr", "")
                                    if _last:
                                        _theme_colors[_name] = "#" + _last.lower()
                        break
                break
    except Exception as exc:
        # Theme colors determine the default fill/border colors of many slides.
        # Keep rendering with explicit/default colors, but make this fidelity loss
        # diagnosable instead of silently returning a visually different deck.
        logger.warning(
            "[PptxGeometry] Failed to read the presentation theme; "
            "using explicit/default colors: %s",
            exc,
        )

    # ── Extract presentation-level default font size (defaultTextStyle) ──
    # Non-placeholder shapes inherit font size from this style chain.
    # Without it the frontend must guess, often resulting in inflated text.
    _default_font_size_pt: float = 18.0  # OOXML spec hardcoded fallback
    _default_title_font_size_pt: float = 36.0  # OOXML spec default for titles
    try:
        _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        _dts = prs.element.find(f"{{{_NS_P}}}defaultTextStyle")
        if _dts is not None:
            _lvl1 = _dts.find(f"{{{_NS_T}}}lvl1pPr")
            if _lvl1 is not None:
                _drp = _lvl1.find(f"{{{_NS_T}}}defRPr")
                if _drp is not None and _drp.get("sz"):
                    _default_font_size_pt = int(_drp.get("sz")) / 100.0
    except Exception as exc:
        logger.warning(
            "[PptxGeometry] Failed to read the presentation default font size; "
            "using 18pt fallback: %s",
            exc,
        )
    # ── Extract title font size from slideMaster txStyles/titleStyle ──
    try:
        _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        for _mst in prs.slide_masters:
            _txS = _mst._element.find(f"{{{_NS_P}}}txStyles")
            if _txS is not None:
                _titS = _txS.find(f"{{{_NS_P}}}titleStyle")
                if _titS is not None:
                    _lvl1t = _titS.find(f"{{{_NS_T}}}lvl1pPr")
                    if _lvl1t is not None:
                        _drpt = _lvl1t.find(f"{{{_NS_T}}}defRPr")
                        if _drpt is not None and _drpt.get("sz"):
                            _default_title_font_size_pt = int(_drpt.get("sz")) / 100.0
            break  # only need first master
    except Exception as exc:
        logger.warning(
            "[PptxGeometry] Failed to read the title font size; "
            "using 36pt fallback: %s",
            exc,
        )

    def _resolve_color(color_obj: Any) -> str | None:
        """Return "#rrggbb" from a python-pptx color object, resolving scheme colors via theme."""
        try:
            return "#" + str(color_obj.rgb).lower()
        except Exception:
            pass
        try:
            from pptx.enum.dml import MSO_THEME_COLOR

            # python-pptx theme_color names like DARK_1, LIGHT_1, ACCENT_1 etc.
            # Map to OOXML scheme names: DARK_1→dk1, LIGHT_1→lt1, ACCENT_N→accentN, etc.
            _tc = color_obj.theme_color
            _name_map = {
                1: "dk1",
                2: "lt1",
                3: "dk2",
                4: "lt2",
                5: "accent1",
                6: "accent2",
                7: "accent3",
                8: "accent4",
                9: "accent5",
                10: "accent6",
                11: "hlink",
                12: "folHlink",
            }
            _hex = _theme_colors.get(_name_map.get(int(_tc), ""))
            if _hex:
                return _hex
        except Exception:
            pass
        return None

    # ── Inner helpers ────────────────────────────────────────────────────

    _NS_DML = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _extract_grad_css(fill: Any) -> str | None:
        """Try to convert a python-pptx gradient fill into a CSS linear-gradient string."""
        try:
            _NS = _NS_DML
            gfill = fill._fill.find(f"{{{_NS}}}gradFill")
            if gfill is None:
                return None
            # Rotation angle is stored in <a:lin ang="…"/> where ang is in 60000ths of a degree
            ang_el = gfill.find(f"{{{_NS}}}lin")
            angle_css = 0
            if ang_el is not None:
                ang_raw = ang_el.get("ang", "0")
                angle_css = round(int(ang_raw) / 60000.0)
            gs_lst = gfill.find(f"{{{_NS}}}gsLst")
            if gs_lst is None:
                return None
            stops = []
            for gs in gs_lst.findall(f"{{{_NS}}}gs"):
                pos_pct = round(int(gs.get("pos", "0")) / 1000.0)
                srgb = gs.find(f"{{{_NS}}}srgbClr")
                if srgb is not None and len(srgb.get("val", "")) == 6:
                    stops.append(f"#{srgb.get('val').lower()} {pos_pct}%")
                else:
                    # sysClr or schemeClr — skip; gradient may be partial but better than nothing
                    lum_mod = gs.find(f".//{{{_NS}}}lumMod")
                    if lum_mod is None:
                        return None  # can't resolve theme colour, bail
                    return None
            if len(stops) < 2:
                return None
            return f"linear-gradient({angle_css}deg, {', '.join(stops)})"
        except Exception:
            return None

    def _extract_bg(slide: Any) -> dict[str, Any]:
        """Walk slide → layout → master for the first extractable fill.
        Returns dict: {"color": "#rrggbb"} for solid, {"gradient": "css"} for gradient,
        {"image": "data:…"} for picture fill, or {} (fallback white).

        IMPORTANT: slide.background._element returns <p:cSld> (the whole slide),
        NOT just the <p:bg> element.  We must narrow the search to <p:bg><p:bgPr>
        to avoid picking up shape images as backgrounds.
        """
        from pptx.oxml.ns import qn as _qn_bg

        _R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        _P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

        def _resolve_blip_image(bg_pr_element: Any, part: Any) -> dict:
            """Find <a:blip r:embed="rIdN"/> in <p:bgPr> XML, resolve to compressed data URI."""
            try:
                blip = bg_pr_element.find(".//" + _qn_bg("a:blip"))
                if blip is not None:
                    rId = blip.get(f"{{{_R_NS}}}embed")
                    if rId and hasattr(part, "rels") and rId in part.rels:
                        img_part = part.rels[rId].target_part
                        img_bytes, mime = _compress_image_bytes(
                            img_part.blob, img_part.content_type or "image/png"
                        )
                        b64 = base64.b64encode(img_bytes).decode("ascii")
                        return {"image": f"data:{mime};base64,{b64}"}
            except Exception:
                pass
            return {}

        def _resolve_solid_color(bg_pr_element: Any) -> dict:
            """Find <a:solidFill><a:srgbClr val="..."/> in <p:bgPr> XML."""
            try:
                solid = bg_pr_element.find(".//" + _qn_bg("a:solidFill"))
                if solid is not None:
                    srgb = solid.find(_qn_bg("a:srgbClr"))
                    if srgb is not None:
                        val = srgb.get("val", "")
                        if len(val) == 6:
                            return {"color": "#" + val.lower()}
            except Exception:
                pass
            return {}

        for src in (
            slide,
            getattr(slide, "slide_layout", None),
            getattr(slide, "slide_master", None),
        ):
            if src is None:
                continue
            try:
                # Locate the <p:bg> element within <p:cSld>.
                # slide.background._element is <p:cSld>; searching it directly
                # would also find <a:blip> inside shape tree → wrong image!
                cSld = src.background._element
                bg_el = cSld.find(f"{{{_P_NS}}}bg")
                if bg_el is None:
                    continue  # this source has no background definition

                bgPr = bg_el.find(f"{{{_P_NS}}}bgPr")
                src_part = getattr(src, "part", None)

                if bgPr is not None and src_part is not None:
                    # Check for image fill (blipFill)
                    result = _resolve_blip_image(bgPr, src_part)
                    if result:
                        return result
                    # Check for solid fill
                    result = _resolve_solid_color(bgPr)
                    if result:
                        return result
                    # Check for gradient fill
                    f = src.background.fill
                    if getattr(f.type, "name", "") == "GRADIENT":
                        css = _extract_grad_css(f)
                        if css:
                            return {"gradient": css}

                # bgRef (theme reference) — use python-pptx fill API
                bgRef = bg_el.find(f"{{{_P_NS}}}bgRef")
                if bgRef is not None:
                    f = src.background.fill
                    fill_name = (
                        getattr(f.type, "name", "") if f.type is not None else ""
                    )
                    if fill_name == "SOLID":
                        try:
                            return {"color": "#" + str(f.fore_color.rgb).lower()}
                        except Exception:
                            pass
                    if fill_name == "GRADIENT":
                        css = _extract_grad_css(f)
                        if css:
                            return {"gradient": css}
            except Exception:
                pass
        return {"color": "#FFFFFF"}

    def _parse_tf(
        tf: Any, layout_defaults: dict[str, Any] | None = None
    ) -> list[dict[str, Any]]:
        """Convert a python-pptx TextFrame into our [{align, runs:[...]}] format.

        Handles three common cases that cause text to disappear:
        1. Paragraph-level default run properties (defRPr) for font size/color
           — most real PPTs store the font size here, not on individual runs.
        2. Empty para.runs when text lives in <a:fld> field elements
           — use para.text as a fallback single run.
        3. Hard line-breaks (<a:br>) within a paragraph
           — emitted as a run with text='\n'.
        """
        # XML namespace used by DrawingML
        _NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

        # ── Read body-level lstStyle defaults (lvl1pPr.defRPr) ──────────────
        # Many PPTX files store the inherited font size here rather than on
        # individual runs or paragraphs.  Serves as lowest-priority fallback.
        _body_defaults: dict[str, Any] = {}
        try:
            lst_style = tf._txBody.find(f"{{{_NS}}}lstStyle")
            if lst_style is not None:
                lvl1 = lst_style.find(f"{{{_NS}}}lvl1pPr")
                if lvl1 is not None:
                    _def_rpr = lvl1.find(f"{{{_NS}}}defRPr")
                    if _def_rpr is not None:
                        # We need _read_rpr but it's defined below; inline a minimal version.
                        _sz = _def_rpr.get("sz")
                        if _sz:
                            try:
                                _v = int(_sz)
                                if _v > 0:
                                    _body_defaults["size"] = round(_v / 100.0, 1)
                            except Exception:
                                pass
                        for _key, _attr in (("bold", "b"), ("italic", "i")):
                            _av = _def_rpr.get(_attr)
                            if _av and _av.lower() not in ("0", "false"):
                                _body_defaults[_key] = True
                        _lat = _def_rpr.find(f"{{{_NS}}}latin")
                        if _lat is not None:
                            _tf = _lat.get("typeface", "")
                            if _tf and not _tf.startswith("+"):
                                _body_defaults["fontName"] = _tf
                        _sol = _def_rpr.find(f"{{{_NS}}}solidFill")
                        if _sol is not None:
                            _srgb = _sol.find(f"{{{_NS}}}srgbClr")
                            if _srgb is not None and len(_srgb.get("val", "")) == 6:
                                _body_defaults["color"] = (
                                    "#" + _srgb.get("val", "").lower()
                                )
        except Exception:
            pass

        def _pt_from_emu_hundredths(val: Any) -> float | None:
            """python-pptx stores font size as 100ths of a point (e.g. 2400 = 24pt)."""
            try:
                v = int(val)
                return round(v / 100.0, 1) if v > 0 else None
            except Exception:
                return None

        def _read_rpr(rpr_el: Any) -> dict[str, Any]:
            """Extract font attributes from an <a:rPr> or <a:defRPr> XML element."""
            out: dict[str, Any] = {}
            if rpr_el is None:
                return out
            try:
                sz = rpr_el.get("sz")
                if sz:
                    pt = _pt_from_emu_hundredths(sz)
                    if pt:
                        out["size"] = pt
            except Exception:
                pass
            try:
                b = rpr_el.get("b")
                if b and b.lower() not in ("0", "false"):
                    out["bold"] = True
            except Exception:
                pass
            try:
                i = rpr_el.get("i")
                if i and i.lower() not in ("0", "false"):
                    out["italic"] = True
            except Exception:
                pass
            try:
                u = rpr_el.get("u")
                if u and u != "none":
                    out["underline"] = True
            except Exception:
                pass
            # Font name from <a:latin typeface="..."> child (Latin / default font)
            try:
                latin = rpr_el.find(f"{{{_NS}}}latin")
                if latin is not None:
                    tf_val = latin.get("typeface", "")
                    if tf_val and not tf_val.startswith("+"):
                        out["fontName"] = tf_val
            except Exception:
                pass
            # East Asian font from <a:ea typeface="..."> child (CJK-specific override)
            try:
                ea = rpr_el.find(f"{{{_NS}}}ea")
                if ea is not None:
                    ea_tf = ea.get("typeface", "")
                    if (
                        ea_tf
                        and not ea_tf.startswith("+")
                        and ea_tf != "+mj-ea"
                        and ea_tf != "+mn-ea"
                    ):
                        out["eaFontName"] = ea_tf
            except Exception:
                pass
            # Solid colour from <a:solidFill><a:srgbClr val="rrggbb">
            try:
                solid = rpr_el.find(f"{{{_NS}}}solidFill")
                if solid is not None:
                    srgb = solid.find(f"{{{_NS}}}srgbClr")
                    if srgb is not None:
                        val = srgb.get("val", "")
                        if len(val) == 6:
                            out["color"] = "#" + val.lower()
            except Exception:
                pass
            # Character spacing: <a:rPr spc="N"/> — hundredths of a point (can be negative)
            try:
                spc = rpr_el.get("spc")
                if spc:
                    out["charSpacing"] = int(spc)
            except Exception:
                pass
            return out

        paras: list[dict[str, Any]] = []
        for para in tf.paragraphs:
            align_name = "LEFT"
            try:
                if para.alignment:
                    align_name = para.alignment.name
            except Exception:
                pass

            # ── Read paragraph-level default run properties (defRPr) ──────
            # These serve as the fallback when a run has no explicit rPr attrs.
            para_defaults: dict[str, Any] = {}
            try:
                pPr = para._p.find(f"{{{_NS}}}pPr")
                if pPr is not None:
                    defRPr = pPr.find(f"{{{_NS}}}defRPr")
                    para_defaults = _read_rpr(defRPr)
            except Exception:
                pass

            p_obj: dict[str, Any] = {"align": align_name, "runs": []}

            # ── Extract paragraph spacing properties ─────────────────
            # <a:lnSpc> → line spacing, <a:spcBef> → space before, <a:spcAft> → space after
            try:
                pPr = para._p.find(f"{{{_NS}}}pPr")
                if pPr is not None:
                    # Line spacing: <a:lnSpc><a:spcPct val="150000"/></a:lnSpc>
                    lnSpc = pPr.find(f"{{{_NS}}}lnSpc")
                    if lnSpc is not None:
                        spcPct = lnSpc.find(f"{{{_NS}}}spcPct")
                        if spcPct is not None:
                            val = int(spcPct.get("val", "0"))
                            if val > 0:
                                p_obj["lineSpacing"] = round(val / 100000.0, 2)
                        else:
                            spcPts = lnSpc.find(f"{{{_NS}}}spcPts")
                            if spcPts is not None:
                                val = int(spcPts.get("val", "0"))
                                if val > 0:
                                    p_obj["lineSpacingPt"] = round(val / 100.0, 1)
                    # Space before: <a:spcBef><a:spcPts val="600"/></a:spcBef>
                    spcBef = pPr.find(f"{{{_NS}}}spcBef")
                    if spcBef is not None:
                        pts = spcBef.find(f"{{{_NS}}}spcPts")
                        if pts is not None:
                            val = int(pts.get("val", "0"))
                            if val > 0:
                                p_obj["spaceBefore"] = round(val / 100.0, 1)
                        else:
                            pct = spcBef.find(f"{{{_NS}}}spcPct")
                            if pct is not None:
                                val = int(pct.get("val", "0"))
                                if val > 0:
                                    p_obj["spaceBeforePct"] = round(val / 100000.0, 2)
                    # Space after: <a:spcAft><a:spcPts val="600"/></a:spcAft>
                    spcAft = pPr.find(f"{{{_NS}}}spcAft")
                    if spcAft is not None:
                        pts = spcAft.find(f"{{{_NS}}}spcPts")
                        if pts is not None:
                            val = int(pts.get("val", "0"))
                            if val > 0:
                                p_obj["spaceAfter"] = round(val / 100.0, 1)
                        else:
                            pct = spcAft.find(f"{{{_NS}}}spcPct")
                            if pct is not None:
                                val = int(pct.get("val", "0"))
                                if val > 0:
                                    p_obj["spaceAfterPct"] = round(val / 100000.0, 2)
            except Exception:
                pass

            # ── Iterate raw XML to capture <a:r> AND <a:br> (hard line-break) ──
            # python-pptx's para.runs only yields <a:r> elements and silently
            # skips <a:fld> field elements (e.g. slide numbers, date fields).
            # We walk _p directly so nothing is lost.
            try:
                for child in para._p:
                    tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

                    if tag == "br":
                        # Hard line-break: emit a newline run so rendering looks correct
                        p_obj["runs"].append({"text": "\n"})
                        continue

                    if tag not in ("r", "fld"):
                        continue

                    # Text node: <a:t> child
                    t_el = child.find(f"{{{_NS}}}t")
                    text_val = (t_el.text or "") if t_el is not None else ""

                    # Run-level properties (<a:rPr>)
                    rPr = child.find(f"{{{_NS}}}rPr")
                    run_attrs = _read_rpr(rPr)

                    # Merge: run → paragraph → layout placeholder → body (lstStyle) defaults
                    _layout_defs = layout_defaults or {}
                    r: dict[str, Any] = {"text": text_val}
                    for key in (
                        "size",
                        "bold",
                        "italic",
                        "underline",
                        "fontName",
                        "eaFontName",
                        "color",
                        "charSpacing",
                    ):
                        if key in run_attrs:
                            r[key] = run_attrs[key]
                        elif key in para_defaults:
                            r[key] = para_defaults[key]
                        elif key in _layout_defs:
                            r[key] = _layout_defs[key]
                        elif key in _body_defaults:
                            r[key] = _body_defaults[key]

                    p_obj["runs"].append(r)

            except Exception:
                # Fallback: use python-pptx's high-level .runs API
                for run in para.runs:
                    r_fb: dict[str, Any] = {"text": run.text}
                    try:
                        if run.font.size:
                            r_fb["size"] = round(run.font.size.pt, 1)
                    except Exception:
                        pass
                    try:
                        if run.font.bold:
                            r_fb["bold"] = True
                    except Exception:
                        pass
                    try:
                        if run.font.italic:
                            r_fb["italic"] = True
                    except Exception:
                        pass
                    try:
                        if run.font.underline:
                            r_fb["underline"] = True
                    except Exception:
                        pass
                    try:
                        if run.font.name:
                            r_fb["fontName"] = run.font.name
                    except Exception:
                        pass
                    try:
                        if run.font.color and run.font.color.type is not None:
                            r_fb["color"] = "#" + str(run.font.color.rgb).lower()
                    except Exception:
                        pass
                    p_obj["runs"].append(r_fb)

            # ── Fallback: if we still have no content, use para.text directly ──
            # This catches edge cases where XML walk also found nothing.
            if not any(r.get("text") for r in p_obj["runs"]):
                raw = ""
                try:
                    raw = para.text
                except Exception:
                    pass
                if raw.strip():
                    fallback_run: dict[str, Any] = {"text": raw}
                    # Apply body defaults first, layout, then para defaults (higher priority)
                    _layout_defs = layout_defaults or {}
                    for key, val in {
                        **_body_defaults,
                        **_layout_defs,
                        **para_defaults,
                    }.items():
                        fallback_run.setdefault(key, val)
                    p_obj["runs"] = [fallback_run]

            # ── Sibling-size inheritance ─────────────────────────────────
            # In OOXML, a run with no explicit sz inherits the "effective" size
            # of the paragraph.  When defRPr & lstStyle also have no size, the
            # best heuristic is the most-common explicit size of neighbouring
            # runs in the same paragraph (e.g.  sz=2100 | sz=None | sz=2100 →
            # the middle run should also be 2100, not the global default).
            _runs = p_obj["runs"]
            if _runs and any(
                "size" not in r for r in _runs if r.get("text", "").strip()
            ):
                _sibling_sizes = [r["size"] for r in _runs if "size" in r]
                if _sibling_sizes:
                    # Use the most common size among siblings
                    from collections import Counter

                    _best_size = Counter(_sibling_sizes).most_common(1)[0][0]
                    for r in _runs:
                        if "size" not in r and r.get("text", "").strip():
                            r["size"] = _best_size

            paras.append(p_obj)
        return paras

    def _collect_shapes(
        shapes_iter: Any,
        out: list[dict[str, Any]],
        z_base: int = 0,
        off_left: int = 0,
        off_top: int = 0,
        grp_ch_off_x: int = 0,
        grp_ch_off_y: int = 0,
        grp_scale_x: float = 1.0,
        grp_scale_y: float = 1.0,
    ) -> None:
        """
        Recursively parse shapes into `out`.
        GROUP shapes are unwrapped and children are appended with their
        slide-absolute coordinates (group offset + internal coordinate scaling).
        """
        for z_idx, shape in enumerate(shapes_iter):
            # ── Resolve effective geometry ───────────────────────────────────
            # python-pptx returns None for left/top/width/height when a
            # placeholder shape inherits its position from the slide layout.
            # (TITLE and BODY placeholders are almost always in this state.)
            # Falling back to 0 makes every text shape invisible because their
            # div has zero width/height — TABLE and PICTURE shapes are
            # always explicitly sized, which is why only they were visible.
            # Fix: pull the real geometry from the matching layout placeholder.
            eff_left = shape.left
            eff_top = shape.top
            eff_w = shape.width
            eff_h = shape.height

            if None in (eff_left, eff_top, eff_w, eff_h):
                try:
                    ph_fmt = shape.placeholder_format
                    if ph_fmt is not None:
                        slide_layout = getattr(
                            getattr(shape, "part", None), "slide_layout", None
                        )
                        if slide_layout is not None:
                            for lph in slide_layout.placeholders:
                                try:
                                    if lph.placeholder_format.idx == ph_fmt.idx:
                                        if eff_left is None:
                                            eff_left = lph.left
                                        if eff_top is None:
                                            eff_top = lph.top
                                        if eff_w is None:
                                            eff_w = lph.width
                                        if eff_h is None:
                                            eff_h = lph.height
                                        break
                                except Exception:
                                    pass
                        # Walk slide master if layout didn't resolve everything
                        if None in (eff_left, eff_top, eff_w, eff_h):
                            slide_master = (
                                getattr(slide_layout, "slide_master", None)
                                if slide_layout
                                else None
                            )
                            if slide_master is not None:
                                for mph in slide_master.placeholders:
                                    try:
                                        if mph.placeholder_format.idx == ph_fmt.idx:
                                            if eff_left is None:
                                                eff_left = mph.left
                                            if eff_top is None:
                                                eff_top = mph.top
                                            if eff_w is None:
                                                eff_w = mph.width
                                            if eff_h is None:
                                                eff_h = mph.height
                                            break
                                    except Exception:
                                        pass
                        # If still None after layout+master, use standard widescreen EMU defaults
                        if None in (eff_left, eff_top, eff_w, eff_h):
                            _ph_idx = getattr(ph_fmt, "idx", -1)
                            _prs_shape = getattr(
                                getattr(shape, "part", None), "presentation", None
                            )
                            _sw = int(
                                getattr(_prs_shape, "slide_width", None) or 9144000
                            )
                            _sh = int(
                                getattr(_prs_shape, "slide_height", None) or 6858000
                            )
                            if _ph_idx == 0:  # Title
                                if eff_left is None:
                                    eff_left = 457200
                                if eff_top is None:
                                    eff_top = 274638
                                if eff_w is None:
                                    eff_w = 8229600
                                if eff_h is None:
                                    eff_h = 1143000
                            elif _ph_idx == 1:  # Body / Content
                                if eff_left is None:
                                    eff_left = 457200
                                if eff_top is None:
                                    eff_top = 1600200
                                if eff_w is None:
                                    eff_w = 8229600
                                if eff_h is None:
                                    eff_h = 4525963
                            else:
                                if eff_left is None:
                                    eff_left = 0
                                if eff_top is None:
                                    eff_top = 0
                                if eff_w is None:
                                    eff_w = _sw
                                if eff_h is None:
                                    eff_h = _sh
                except Exception:
                    pass

            abs_left = off_left + round(((eff_left or 0) - grp_ch_off_x) * grp_scale_x)
            abs_top = off_top + round(((eff_top or 0) - grp_ch_off_y) * grp_scale_y)
            # Also scale the shape's own width/height when inside a group
            if grp_scale_x != 1.0 and eff_w is not None:
                eff_w = round(eff_w * grp_scale_x)
            if grp_scale_y != 1.0 and eff_h is not None:
                eff_h = round(eff_h * grp_scale_y)

            # ── Group: recurse with absolute-coordinate offset + internal scaling ──
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    # OOXML groups have an internal coordinate system.
                    # <p:grpSpPr><a:xfrm> carries:
                    #   off  x/y  : group's position on slide (same as eff_left/eff_top)
                    #   ext  cx/cy: group's rendered size on slide
                    #   chOff x/y : origin of the internal coordinate system
                    #   chExt cx/cy: size of the internal coordinate system
                    # Child shapes' left/top are in internal coordinates; we map them to slide coords.
                    _NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
                    _NS_A2 = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    xfrm = shape.element.find(
                        f"{{{_NS_P}}}grpSpPr/{{{_NS_A2}}}xfrm"
                    ) or shape.element.find(f"{{{_NS_A2}}}grpSpPr/{{{_NS_A2}}}xfrm")
                    ch_off_x = ch_off_y = 0
                    ch_ext_cx = eff_w or 1
                    ch_ext_cy = eff_h or 1
                    if xfrm is not None:
                        chOff = xfrm.find(f"{{{_NS_A2}}}chOff")
                        chExt = xfrm.find(f"{{{_NS_A2}}}chExt")
                        if chOff is not None:
                            ch_off_x = int(chOff.get("x", 0))
                            ch_off_y = int(chOff.get("y", 0))
                        if chExt is not None:
                            ch_ext_cx = int(chExt.get("cx", 1)) or 1
                            ch_ext_cy = int(chExt.get("cy", 1)) or 1
                    grp_w = eff_w or ch_ext_cx
                    grp_h = eff_h or ch_ext_cy
                    scale_x = grp_w / ch_ext_cx
                    scale_y = grp_h / ch_ext_cy
                    _collect_shapes(
                        shape.shapes,
                        out,
                        z_base=z_base + z_idx * 100,
                        off_left=abs_left,
                        off_top=abs_top,
                        grp_ch_off_x=ch_off_x,
                        grp_ch_off_y=ch_off_y,
                        grp_scale_x=scale_x,
                        grp_scale_y=scale_y,
                    )
                except Exception:
                    _collect_shapes(
                        shape.shapes,
                        out,
                        z_base=z_base + z_idx * 100,
                        off_left=abs_left,
                        off_top=abs_top,
                    )
                continue

            # ── Media (video/audio embedded shapes): skip entirely ───────────
            # MSO_SHAPE_TYPE.MEDIA = 16.  python-pptx cannot render video/audio
            # and attempting to read their blobs would base64-encode hundreds of MB.
            try:
                if (
                    getattr(MSO_SHAPE_TYPE, "MEDIA", None) is not None
                    and shape.shape_type == MSO_SHAPE_TYPE.MEDIA
                ):
                    continue
            except Exception:
                pass

            s: dict[str, Any] = {
                "id": shape.shape_id,
                "name": shape.name,
                "left": abs_left,
                "top": abs_top,
                "width": eff_w or 0,
                "height": eff_h or 0,
                "z_order": z_base + z_idx,
                "fill": None,
            }

            # Shape fill — solid, gradient, or picture
            try:
                fill = shape.fill
                fill_name = (
                    getattr(fill.type, "name", "") if fill.type is not None else ""
                )
                if fill_name == "SOLID":
                    _fc = _resolve_color(fill.fore_color)
                    if _fc:
                        s["fill"] = _fc
                elif fill_name == "GRADIENT":
                    css = _extract_grad_css(fill)
                    if css:
                        s["fillGradient"] = css
                elif fill_name == "PICTURE":
                    try:
                        from pptx.oxml.ns import qn as _qn_sh

                        blip = shape.element.find(".//" + _qn_sh("a:blip"))
                        if blip is not None:
                            rId = blip.get(
                                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                            )
                            if rId:
                                img_part = shape.part.related_parts[rId]
                                _raw_blob = img_part.blob
                                _raw_mime = img_part.content_type or "image/png"
                                # Skip non-image blobs (e.g. video used as fill) or oversized media
                                if (
                                    _raw_mime.startswith("image/")
                                    and len(_raw_blob) <= _MAX_BLOB_BYTES
                                ):
                                    img_bytes, mime = _compress_image_bytes(
                                        _raw_blob, _raw_mime
                                    )
                                    b64 = base64.b64encode(img_bytes).decode("ascii")
                                    s["fillImage"] = f"data:{mime};base64,{b64}"
                    except Exception:
                        pass
            except Exception:
                pass

            # Shape outline/border
            try:
                line = shape.line
                if line and line.width is not None and line.width > 0:
                    _lc = None
                    try:
                        _lc = _resolve_color(line.color)
                    except Exception:
                        pass
                    # line.width is in EMU; store as-is, frontend scales
                    s["border"] = {
                        "widthEmu": int(line.width),
                        "color": _lc or "#000000",
                    }
            except Exception:
                pass

            # Shape rotation (degrees, clockwise-positive; python-pptx returns float or None)
            try:
                _rot = shape.rotation
                if _rot is not None and _rot != 0.0:
                    s["rotation"] = round(float(_rot), 2)
            except Exception:
                pass

            # ── Picture ───────────────────────────────────────────────
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_blob = shape.image.blob
                    img_mime = shape.image.content_type or "image/png"
                    # Guard: skip if non-image MIME (video poster frame) or oversized blob
                    if (
                        not img_mime.startswith("image/")
                        or len(img_blob) > _MAX_BLOB_BYTES
                    ):
                        logger.warning(
                            "[parse_pptx] skipping oversized/non-image blob: mime=%s size=%.1f MB",
                            img_mime,
                            len(img_blob) / 1048576,
                        )
                        s["_type"] = "PICTURE"
                        s["image_b64"] = (
                            ""  # placeholder — too large or non-image media
                        )
                        out.append(s)
                        continue
                    img_blob, img_mime = _compress_image_bytes(img_blob, img_mime)
                    b64 = base64.b64encode(img_blob).decode("ascii")
                    s["_type"] = "PICTURE"
                    s["image_b64"] = f"data:{img_mime};base64,{b64}"
                    out.append(s)
                    continue
            except Exception:
                pass

            # ── Table ─────────────────────────────────────────────────
            try:
                if shape.has_table:
                    tbl = shape.table
                    # Column widths and row heights (EMU)
                    col_widths = []
                    row_heights = []
                    try:
                        col_widths = [int(col.width) for col in tbl.columns]
                    except Exception:
                        pass
                    try:
                        row_heights = [int(row.height) for row in tbl.rows]
                    except Exception:
                        pass
                    cells: list[dict[str, Any]] = []
                    for r_idx, row in enumerate(tbl.rows):
                        for c_idx, cell in enumerate(row.cells):
                            cell_text = ""
                            try:
                                cell_text = (
                                    cell.text_frame.text if cell.text_frame else ""
                                )
                            except Exception:
                                pass
                            cell_d: dict[str, Any] = {
                                "row": r_idx,
                                "col": c_idx,
                                "text": cell_text,
                            }
                            # Cell fill
                            try:
                                cfill = cell.fill
                                if (
                                    cfill.type is not None
                                    and getattr(cfill.type, "name", "") == "SOLID"
                                ):
                                    _cf = _resolve_color(cfill.fore_color)
                                    if _cf:
                                        cell_d["fill"] = _cf
                            except Exception:
                                pass
                            # Per-cell text formatting (first run of first paragraph)
                            try:
                                if cell.text_frame and cell.text_frame.paragraphs:
                                    fp = cell.text_frame.paragraphs[0]
                                    if fp.runs:
                                        fr = fp.runs[0]
                                        if fr.font.size:
                                            cell_d["fontSize"] = round(
                                                fr.font.size.pt, 1
                                            )
                                        if fr.font.bold:
                                            cell_d["bold"] = True
                                        if (
                                            fr.font.color
                                            and fr.font.color.type is not None
                                        ):
                                            cell_d["color"] = (
                                                "#" + str(fr.font.color.rgb).lower()
                                            )
                                    cell_d["align"] = (
                                        fp.alignment.name if fp.alignment else "LEFT"
                                    )
                            except Exception:
                                pass
                            cells.append(cell_d)
                    s["_type"] = "TABLE"
                    s["table_rows"] = len(tbl.rows)
                    s["table_cols"] = len(tbl.columns)
                    s["col_widths"] = col_widths
                    s["row_heights"] = row_heights
                    s["cells"] = cells
                    out.append(s)
                    continue
            except Exception:
                pass

            # ── Generic Shapes / Icons / SVGs / Charts / Connectors ────────
            try:
                _st = getattr(shape, "shape_type", None)
                if _st in (
                    MSO_SHAPE_TYPE.AUTO_SHAPE,
                    MSO_SHAPE_TYPE.FREEFORM,
                    MSO_SHAPE_TYPE.GRAPHIC_FRAME,
                ):
                    s["_type"] = "SHAPE"
                    # Rounded rectangle: extract corner radius from XML
                    try:
                        from pptx.oxml.ns import qn as _qn_r

                        prstGeom = shape.element.find(".//" + _qn_r("a:prstGeom"))
                        if prstGeom is not None:
                            s["autoShapeType"] = prstGeom.get("prst", "")
                            avLst = prstGeom.find(".//" + _qn_r("a:gd"))
                            if avLst is not None and avLst.get("fmla", "").startswith(
                                "val "
                            ):
                                s["cornerRadiusEmu"] = int(avLst.get("fmla").split()[1])
                    except Exception:
                        pass
                elif _st in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREE_FORM):
                    s["_type"] = "LINE"
            except Exception:
                pass

            try:
                if getattr(shape, "has_chart", False):
                    s["_type"] = "CHART"
            except Exception:
                pass

            # ── Text frame ────────────────────────────────────────────
            try:
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    s["_type"] = "TEXT"
                    s["has_text"] = True
                    is_title = False
                    try:
                        from pptx.enum.shapes import PP_PLACEHOLDER

                        ph = shape.placeholder_format
                        if ph is not None:
                            is_title = ph.type in (
                                PP_PLACEHOLDER.TITLE,
                                PP_PLACEHOLDER.CENTER_TITLE,
                            )
                    except Exception:
                        pass
                    s["is_title"] = is_title

                    # ── Extract text body properties (<a:bodyPr>) ─────────
                    _NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    try:
                        bodyPr = shape.text_frame._txBody.find(f"{{{_NS_A}}}bodyPr")
                        if bodyPr is not None:
                            anchor = bodyPr.get("anchor")
                            # If not set at slide level, inherit from layout → master placeholder
                            if not anchor:
                                try:
                                    _ph_fmt = getattr(shape, "placeholder_format", None)
                                    if _ph_fmt is not None:
                                        _ph_idx = _ph_fmt.idx
                                        _sl_layout = getattr(
                                            getattr(shape, "part", None),
                                            "slide_layout",
                                            None,
                                        )
                                        if _sl_layout is not None:
                                            for _lph in _sl_layout.placeholders:
                                                try:
                                                    if (
                                                        _lph.placeholder_format.idx
                                                        == _ph_idx
                                                    ):
                                                        _lbPr = _lph.text_frame._txBody.find(
                                                            f"{{{_NS_A}}}bodyPr"
                                                        )
                                                        if _lbPr is not None:
                                                            anchor = (
                                                                _lbPr.get("anchor")
                                                                or anchor
                                                            )
                                                        break
                                                except Exception:
                                                    pass
                                        if not anchor:
                                            _sl_master = (
                                                getattr(
                                                    _sl_layout, "slide_master", None
                                                )
                                                if _sl_layout
                                                else None
                                            )
                                            if _sl_master is not None:
                                                for _mph in _sl_master.placeholders:
                                                    try:
                                                        if (
                                                            _mph.placeholder_format.idx
                                                            == _ph_idx
                                                        ):
                                                            _mbPr = _mph.text_frame._txBody.find(
                                                                f"{{{_NS_A}}}bodyPr"
                                                            )
                                                            if _mbPr is not None:
                                                                anchor = (
                                                                    _mbPr.get("anchor")
                                                                    or anchor
                                                                )
                                                            break
                                                    except Exception:
                                                        pass
                                except Exception:
                                    pass
                            if anchor:
                                s["textAnchor"] = anchor  # t | ctr | b | just | dist
                            # Text insets (EMU). OOXML defaults: l/r=91440, t/b=45720
                            lIns = bodyPr.get("lIns")
                            tIns = bodyPr.get("tIns")
                            rIns = bodyPr.get("rIns")
                            bIns = bodyPr.get("bIns")
                            s["textInsets"] = {
                                "l": int(lIns) if lIns is not None else 91440,
                                "t": int(tIns) if tIns is not None else 45720,
                                "r": int(rIns) if rIns is not None else 91440,
                                "b": int(bIns) if bIns is not None else 45720,
                            }
                            # Text autofit mode: spAutoFit | normAutofit (with fontScale) | noAutofit
                            _af_sp = bodyPr.find(f"{{{_NS_A}}}spAutoFit")
                            _af_norm = bodyPr.find(f"{{{_NS_A}}}normAutofit")
                            if _af_sp is not None:
                                s["autoFit"] = "sp"
                            elif _af_norm is not None:
                                s["autoFit"] = "norm"
                                _fs = _af_norm.get("fontScale")
                                if _fs:
                                    # fontScale: 100000 = 100%, 75000 = 75% of original pt size
                                    s["fontScale"] = round(int(_fs) / 1000.0, 1)
                            # Word wrap: "square" (default, wraps) or "none" (no wrap)
                            _wrap = bodyPr.get("wrap")
                            if _wrap == "none":
                                s["wordWrap"] = "none"
                    except Exception:
                        pass

                    # ── Layout placeholder lstStyle inheritance ──────────
                    # OOXML font size resolution: run rPr → para defRPr →
                    # shape lstStyle → layout placeholder lstStyle →
                    # master titleStyle/bodyStyle → defaultTextStyle.
                    # Extract layout placeholder's lstStyle defaults so
                    # _parse_tf can insert them in the merge chain.
                    _layout_defaults: dict[str, Any] = {}
                    try:
                        _ph_fmt2 = getattr(shape, "placeholder_format", None)
                        if _ph_fmt2 is not None:
                            _ph_idx2 = _ph_fmt2.idx
                            _sl_layout2 = getattr(
                                getattr(shape, "part", None), "slide_layout", None
                            )
                            if _sl_layout2 is not None:
                                for _lph2 in _sl_layout2.placeholders:
                                    try:
                                        if _lph2.placeholder_format.idx == _ph_idx2:
                                            _ltxBody = _lph2.text_frame._txBody
                                            _llstStyle = _ltxBody.find(
                                                f"{{{_NS_A}}}lstStyle"
                                            )
                                            if _llstStyle is not None:
                                                _llvl1 = _llstStyle.find(
                                                    f"{{{_NS_A}}}lvl1pPr"
                                                )
                                                if _llvl1 is not None:
                                                    _ldefRPr = _llvl1.find(
                                                        f"{{{_NS_A}}}defRPr"
                                                    )
                                                    if _ldefRPr is not None:
                                                        _lsz = _ldefRPr.get("sz")
                                                        if _lsz:
                                                            try:
                                                                _lv = int(_lsz)
                                                                if _lv > 0:
                                                                    _layout_defaults[
                                                                        "size"
                                                                    ] = round(
                                                                        _lv / 100.0, 1
                                                                    )
                                                            except Exception:
                                                                pass
                                                        for _lk, _la in (
                                                            ("bold", "b"),
                                                            ("italic", "i"),
                                                        ):
                                                            _lav = _ldefRPr.get(_la)
                                                            if (
                                                                _lav
                                                                and _lav.lower()
                                                                not in ("0", "false")
                                                            ):
                                                                _layout_defaults[
                                                                    _lk
                                                                ] = True
                                                        _llat = _ldefRPr.find(
                                                            f"{{{_NS_A}}}latin"
                                                        )
                                                        if _llat is not None:
                                                            _ltf = _llat.get(
                                                                "typeface", ""
                                                            )
                                                            if (
                                                                _ltf
                                                                and not _ltf.startswith(
                                                                    "+"
                                                                )
                                                            ):
                                                                _layout_defaults[
                                                                    "fontName"
                                                                ] = _ltf
                                                        _lea = _ldefRPr.find(
                                                            f"{{{_NS_A}}}ea"
                                                        )
                                                        if _lea is not None:
                                                            _leaf = _lea.get(
                                                                "typeface", ""
                                                            )
                                                            if (
                                                                _leaf
                                                                and not _leaf.startswith(
                                                                    "+"
                                                                )
                                                            ):
                                                                _layout_defaults[
                                                                    "eaFontName"
                                                                ] = _leaf
                                                        _lsol = _ldefRPr.find(
                                                            f"{{{_NS_A}}}solidFill"
                                                        )
                                                        if _lsol is not None:
                                                            _lsrgb = _lsol.find(
                                                                f"{{{_NS_A}}}srgbClr"
                                                            )
                                                            if (
                                                                _lsrgb is not None
                                                                and len(
                                                                    _lsrgb.get(
                                                                        "val", ""
                                                                    )
                                                                )
                                                                == 6
                                                            ):
                                                                _layout_defaults[
                                                                    "color"
                                                                ] = (
                                                                    "#"
                                                                    + _lsrgb.get(
                                                                        "val", ""
                                                                    ).lower()
                                                                )
                                            break
                                    except Exception:
                                        pass
                    except Exception:
                        pass

                    s["paragraphs"] = _parse_tf(
                        shape.text_frame, layout_defaults=_layout_defaults
                    )
                    out.append(s)
            except Exception:
                pass

            # SHAPE/CHART types set their _type above but don't self-append (unlike
            # TEXT, PICTURE, TABLE which already called out.append+continue).
            if s.get("_type") in ("SHAPE", "CHART"):
                out.append(s)

    # ── Main loop ────────────────────────────────────────────────────────

    for slide_idx, slide in enumerate(prs.slides):
        bg_info = _extract_bg(slide)
        # bg_info is now a dict: {"color":"#..."} | {"gradient":"css"} | {"image":"data:..."} | {}

        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass

        shapes_data: list[dict[str, Any]] = []
        # ── Collect decorative shapes from slide master + layout first (lowest z-order).
        # These form the visual theme/background template visible behind slide content.
        try:
            _slide_layout = getattr(slide, "slide_layout", None)
            _slide_master = (
                getattr(_slide_layout, "slide_master", None) if _slide_layout else None
            )
            # Slide master shapes (z_base=-2000, bottom-most layer)
            if _slide_master is not None:
                _mst_decos = [
                    _s
                    for _s in _slide_master.shapes
                    if not getattr(_s, "is_placeholder", False)
                ]
                _mst_out: list[dict[str, Any]] = []
                _collect_shapes(_mst_decos, _mst_out, z_base=-2000)
                for _ms in _mst_out:
                    _ms["editable"] = False
                shapes_data.extend(_mst_out)
            # Slide layout shapes (z_base=-1000, above master, below slide content)
            if _slide_layout is not None:
                _lay_decos = [
                    _s
                    for _s in _slide_layout.shapes
                    if not getattr(_s, "is_placeholder", False)
                ]
                _lay_out: list[dict[str, Any]] = []
                _collect_shapes(_lay_decos, _lay_out, z_base=-1000)
                for _ls in _lay_out:
                    _ls["editable"] = False
                shapes_data.extend(_lay_out)
        except Exception:
            pass
        # Slide's own shapes (z_base=0, top layer)
        _collect_shapes(slide.shapes, shapes_data)

        # ── Text exclusion zones ─────────────────────────────────────────────
        # When a lower-z-order text shape physically overlaps a higher-z-order
        # text shape, expand the lower shape's right/left textInset so its text
        # content doesn't flow under the upper shape.  This reproduces
        # PowerPoint's implicit text-runaround behaviour for overlapping shapes.
        _EXCL_BUFFER = 91440  # 1 EMU point of extra breathing room
        _MAX_EXCL_OVERLAP_RATIO = 0.6
        _text_shapes = [
            s
            for s in shapes_data
            if s.get("has_text") and s.get("textInsets") is not None
        ]
        for _a in _text_shapes:
            _al = _a["left"]
            _at = _a["top"]
            _ar = _al + _a["width"]
            _ab = _at + _a["height"]
            _acx = _al + _a["width"] / 2
            for _b in shapes_data:
                if _b.get("z_order", 0) <= _a.get("z_order", 0):
                    continue  # only consider shapes rendered ON TOP of _a
                if not _b.get("has_text"):
                    continue  # only exclude space for text-bearing shapes
                _bl = _b["left"]
                _bt = _b["top"]
                _br = _bl + _b["width"]
                _bb = _bt + _b["height"]
                # Must overlap in BOTH axes
                if _bb <= _at or _bt >= _ab:
                    continue
                if _br <= _al or _bl >= _ar:
                    continue
                _overlap_l = max(_al, _bl)
                _overlap_r = min(_ar, _br)
                _overlap_t = max(_at, _bt)
                _overlap_b = min(_ab, _bb)
                _overlap_w = _overlap_r - _overlap_l
                _overlap_h = _overlap_b - _overlap_t
                if _overlap_w <= 0 or _overlap_h <= 0:
                    continue
                # Stacked placeholders often overlap by a thin strip while spanning
                # most of the line width. That is not a side obstruction and should
                # not force a huge left/right inset on the lower text box.
                if (
                    _overlap_w >= _a["width"] * _MAX_EXCL_OVERLAP_RATIO
                    or _overlap_w >= _b["width"] * _MAX_EXCL_OVERLAP_RATIO
                    or (_bl < _acx < _br)
                ):
                    continue
                # Expand inset on the side where _b sits relative to _a's centre
                _ins = _a["textInsets"]
                if (_bl + _b["width"] / 2) > _acx:
                    # _b is to the right of _a's centre → expand right inset
                    _need_r = (_ar - _bl) + _EXCL_BUFFER
                    if _need_r > _ins.get("r", 91440):
                        _ins["r"] = _need_r
                else:
                    # _b is to the left of _a's centre → expand left inset
                    _need_l = (_br - _al) + _EXCL_BUFFER
                    if _need_l > _ins.get("l", 91440):
                        _ins["l"] = _need_l

        slides_data.append(
            {
                "slide_index": slide_idx,
                "slide_id": slide_idx + 1,
                "background": bg_info.get("color", "#FFFFFF"),
                "backgroundGradient": bg_info.get("gradient"),
                "backgroundImage": bg_info.get("image"),
                "notes": notes_text,
                "shapes": shapes_data,
            }
        )

    return {
        "slide_width_emu": int(slide_w),
        "slide_height_emu": int(slide_h),
        "default_font_size_pt": _default_font_size_pt,
        "default_title_font_size_pt": _default_title_font_size_pt,
        "slides": slides_data,
    }


# ─────────────────────────────────────────────────────────────────────────────
# PDF → 文本提取 + 原始 URL
# ─────────────────────────────────────────────────────────────────────────────
