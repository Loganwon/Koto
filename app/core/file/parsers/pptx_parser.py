# -*- coding: utf-8 -*-
# Copyright (C) 2024-2026 Koto AI. All rights reserved.
# SPDX-License-Identifier: LicenseRef-Koto-Proprietary
from __future__ import annotations

from typing import Any


def parse_pptx(file_path: str) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import PP_PLACEHOLDER
    except ImportError:
        raise RuntimeError("python-pptx 未安装，请执行: pip install python-pptx")

    prs = Presentation(file_path)
    slides_data: list[dict[str, Any]] = []

    for slide_idx, slide in enumerate(prs.slides):
        texts: list[dict[str, Any]] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_content = "\n".join(
                paragraph.text for paragraph in shape.text_frame.paragraphs
            ).strip()
            if not text_content:
                continue

            is_title = False
            try:
                placeholder = shape.placeholder_format
                if placeholder is not None:
                    is_title = placeholder.type in (
                        PP_PLACEHOLDER.TITLE,
                        PP_PLACEHOLDER.CENTER_TITLE,
                    )
            except Exception:
                pass

            texts.append(
                {
                    "shape_id": shape.shape_id,
                    "shape_name": shape.name,
                    "text": text_content,
                    "is_title": is_title,
                }
            )

        slides_data.append(
            {
                "slide_id": slide_idx + 1,
                "slide_index": slide_idx,
                "texts": texts,
            }
        )

    return slides_data


def parse_pptx_geometry(file_path: Any) -> dict[str, Any]:
    from app.core.file.parsers.pptx_geometry_parser import (
        parse_pptx_geometry as _parse_pptx_geometry,
    )

    return _parse_pptx_geometry(file_path)

__all__ = ["parse_pptx", "parse_pptx_geometry"]
