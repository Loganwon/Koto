# -*- coding: utf-8 -*-
"""
╔══════════════════════════════════════════════════════════════════╗
║   Koto  ─  Document Generation Service                          ║
║   支持：Word (.docx) / PDF (reportlab) / Excel (.xlsx) /        ║
║          PowerPoint (.pptx)                                     ║
╚══════════════════════════════════════════════════════════════════╝

设计原则：
  1. 四种格式共用同一套 DocGenRequest 数据结构，reduce duplication
  2. 支持「专业/学术/简约/彩色」四种预设风格
  3. 中英双语：字体选择自动检测系统可用中文字体
  4. 每个格式的生成器独立封装，失败不影响其他格式
  5. 与 TaskPlanner 联动：generate_* 方法均为同步（在线程池中执行）
  6. 输出文件默认保存到 workspace/doc_gen_outputs/，可通过 output_dir 覆盖

用法示例::

    from app.core.services.doc_gen_service import DocGenService, DocGenRequest, DocSection

    svc = DocGenService()

    req = DocGenRequest(
        title="2026年Q1市场分析报告",
        sections=[
            DocSection(content_type="heading", content="执行摘要", level=1),
            DocSection(content_type="text",    content="本季度整体收入增长18%..."),
            DocSection(content_type="table",   content=[
                ["指标", "Q1 2026", "Q1 2025", "同比"],
                ["收入", "¥1,200万", "¥1,017万", "+18%"],
                ["毛利率", "42%", "38%", "+4pp"],
            ]),
            DocSection(content_type="heading", content="竞争态势分析", level=1),
            DocSection(content_type="bullet",  content=["竞品A: 市占率下降3%", "竞品B: 新品线发布", "机会窗口: 中小企业市场"]),
        ],
        author="Koto AI",
        style_preset="professional",
    )

    word_path = svc.generate_word(req)
    pdf_path  = svc.generate_pdf(req)
    xlsx_path = svc.generate_excel(req)  # 需要 ExcelSheetRequest 格式
    pptx_path = svc.generate_presentation(req)
"""

from __future__ import annotations

import logging
import os
import sys
import uuid
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

# ── 输出目录 ─────────────────────────────────────────────────────────────────
_BASE_DIR = (
    Path(sys.executable).parent
    if getattr(sys, "frozen", False)
    else Path(__file__).resolve().parents[3]  # project root
)
_DEFAULT_OUTPUT_DIR = _BASE_DIR / "workspace" / "doc_gen_outputs"


# ══════════════════════════════════════════════════════════════════════════════
# 数据模型
# ══════════════════════════════════════════════════════════════════════════════


class ContentType(str, Enum):
    TEXT = "text"
    HEADING = "heading"  # level 1/2/3 → H1/H2/H3
    TABLE = "table"  # content = [[header…], [row1…], …]
    BULLET_LIST = "bullet"  # content = ["item1", "item2", …]
    NUMBERED_LIST = "numbered"  # content = ["item1", "item2", …]
    CODE = "code"  # content = str, style.get("language") for syntax hint
    IMAGE = "image"  # content = file path or URL
    PAGE_BREAK = "page_break"
    HORIZONTAL_RULE = "hr"
    QUOTE = "quote"  # block quote


@dataclass
class DocSection:
    """
    文档的一个内容块。
    content_type 决定如何渲染，level 只对 heading 有效。
    """

    content_type: str  # ContentType 或字符串
    content: Any  # str | list，取决于 content_type
    level: int = 1  # heading 层级 (1=H1, 2=H2, 3=H3)
    style: Optional[Dict[str, Any]] = None  # 可选额外样式提示
    caption: Optional[str] = None  # 图表/表格图注


@dataclass
class ExcelSheet:
    """Excel 中一个工作表的数据描述。"""

    name: str = "Sheet1"
    title: str = ""
    column_headers: List[str] = field(default_factory=list)
    rows: List[List[Any]] = field(default_factory=list)
    add_totals_row: bool = False  # 是否在末尾追加合计行
    column_widths: Optional[List[int]] = None  # 各列宽度（字符数），None=自适应
    freeze_header: bool = True


@dataclass
class DocGenRequest:
    """
    统一的文档生成请求。同一个请求对象可用于生成四种格式。
    """

    title: str
    sections: List[DocSection]

    # ── 基本元数据 ────────────────────────────────────────────────────────
    author: str = "Koto AI"
    subject: str = ""
    keywords: str = ""
    language: str = "zh"  # "zh" | "en"

    # ── 样式预设 ──────────────────────────────────────────────────────────
    style_preset: str = "professional"  # "professional"|"academic"|"minimal"|"colorful"
    theme_color: str = "#1F4E79"  # 主色调（十六进制）

    # ── 输出控制 ──────────────────────────────────────────────────────────
    output_dir: str = ""  # 留空则使用默认目录
    output_filename: str = ""  # 留空则自动生成（含时间戳）

    # ── Excel 专用（多 Sheet 支持） ───────────────────────────────────────
    excel_sheets: List[ExcelSheet] = field(default_factory=list)

    # ── PPT 专用 ──────────────────────────────────────────────────────────
    pptx_layout: str = "widescreen"  # "widescreen"(16:9) | "standard"(4:3)
    pptx_theme: str = "default"

    # ── 附加元数据 ────────────────────────────────────────────────────────
    meta: Dict[str, Any] = field(default_factory=dict)


# ══════════════════════════════════════════════════════════════════════════════
# 工具函数
# ══════════════════════════════════════════════════════════════════════════════

# 4种预设的配色方案
_STYLE_PRESETS: Dict[str, Dict[str, Any]] = {
    "professional": {
        "title_font_size": 28,
        "h1_font_size": 16,
        "h2_font_size": 14,
        "h3_font_size": 12,
        "body_font_size": 11,
        "primary_rgb": (31, 78, 121),  # #1F4E79 深蓝
        "secondary_rgb": (68, 114, 196),  # #4472C4 中蓝
        "accent_rgb": (237, 125, 49),  # #ED7D31 橙
        "header_bg_rgb": (31, 78, 121),
        "header_fg_rgb": (255, 255, 255),
        "stripe_bg_rgb": (217, 225, 242),
        "line_spacing": 1.2,
    },
    "academic": {
        "title_font_size": 24,
        "h1_font_size": 14,
        "h2_font_size": 13,
        "h3_font_size": 12,
        "body_font_size": 11,
        "primary_rgb": (0, 0, 0),
        "secondary_rgb": (50, 50, 50),
        "accent_rgb": (0, 70, 130),
        "header_bg_rgb": (240, 240, 240),
        "header_fg_rgb": (0, 0, 0),
        "stripe_bg_rgb": (250, 250, 250),
        "line_spacing": 1.5,
    },
    "minimal": {
        "title_font_size": 26,
        "h1_font_size": 15,
        "h2_font_size": 13,
        "h3_font_size": 12,
        "body_font_size": 11,
        "primary_rgb": (60, 60, 60),
        "secondary_rgb": (100, 100, 100),
        "accent_rgb": (180, 180, 180),
        "header_bg_rgb": (230, 230, 230),
        "header_fg_rgb": (50, 50, 50),
        "stripe_bg_rgb": (248, 248, 248),
        "line_spacing": 1.3,
    },
    "colorful": {
        "title_font_size": 30,
        "h1_font_size": 18,
        "h2_font_size": 15,
        "h3_font_size": 13,
        "body_font_size": 11,
        "primary_rgb": (112, 48, 160),  # 紫色
        "secondary_rgb": (0, 176, 240),  # 青色
        "accent_rgb": (255, 192, 0),  # 金黄
        "header_bg_rgb": (112, 48, 160),
        "header_fg_rgb": (255, 255, 255),
        "stripe_bg_rgb": (235, 220, 255),
        "line_spacing": 1.2,
    },
}


def _get_preset(style_preset: str) -> Dict[str, Any]:
    return _STYLE_PRESETS.get(style_preset, _STYLE_PRESETS["professional"])


def _ensure_output_dir(output_dir: str) -> Path:
    if output_dir:
        p = Path(output_dir)
    else:
        p = _DEFAULT_OUTPUT_DIR
    p.mkdir(parents=True, exist_ok=True)
    return p


def _build_output_path(output_dir: str, filename: str, extension: str) -> Path:
    out_dir = _ensure_output_dir(output_dir)
    if filename:
        stem = Path(filename).stem
    else:
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        stem = f"koto_{ts}_{uuid.uuid4().hex[:6]}"
    return out_dir / f"{stem}.{extension}"


def _hex_to_rgb(hex_color: str) -> Tuple[int, int, int]:
    """将 #RRGGBB 转换为 (R, G, B) 整数元组。"""
    h = hex_color.lstrip("#")
    if len(h) == 6:
        return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))  # type: ignore
    return (31, 78, 121)  # fallback: professional blue


def _find_chinese_font() -> Optional[str]:
    """
    在系统字体目录中查找可用的中文 TrueType 字体路径。
    优先顺序：微软雅黑 → 黑体 → 宋体 → 苹方 → 文泉驿。
    """
    candidates = [
        # Windows
        r"C:\Windows\Fonts\msyh.ttc",  # 微软雅黑
        r"C:\Windows\Fonts\msyhbd.ttc",
        r"C:\Windows\Fonts\simhei.ttf",  # 黑体
        r"C:\Windows\Fonts\simsun.ttc",  # 宋体
        r"C:\Windows\Fonts\simfang.ttf",  # 仿宋
        # macOS
        "/System/Library/Fonts/PingFang.ttc",
        "/Library/Fonts/Arial Unicode MS.ttf",
        # Linux
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    ]
    for path in candidates:
        if os.path.exists(path):
            return path
    return None


# ══════════════════════════════════════════════════════════════════════════════
# Word 生成器
# ══════════════════════════════════════════════════════════════════════════════


class _WordGenerator:
    """基于 python-docx 生成专业品质的 Word 文档。"""

    def generate(self, req: DocGenRequest, output_path: Path) -> None:
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor, Inches, Cm
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.enum.table import WD_TABLE_ALIGNMENT
            from docx.oxml.ns import qn
        except ImportError as e:
            raise ImportError(f"python-docx 未安装: {e}")

        preset = _get_preset(req.style_preset)
        doc = Document()

        # ── 页面设置 ──────────────────────────────────────────────────────
        for section in doc.sections:
            section.page_height = Cm(29.7)
            section.page_width = Cm(21.0)
            section.top_margin = Cm(2.5)
            section.bottom_margin = Cm(2.5)
            section.left_margin = Cm(3.0)
            section.right_margin = Cm(2.5)

        # ── 标题页 ────────────────────────────────────────────────────────
        title_para = doc.add_paragraph()
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_run = title_para.add_run(req.title)
        title_run.bold = True
        title_run.font.size = Pt(preset["title_font_size"])
        pr, pg, pb = preset["primary_rgb"]
        title_run.font.color.rgb = RGBColor(pr, pg, pb)
        title_run.font.name = "微软雅黑"
        title_run._r.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")

        # 副标题行（作者 + 日期）
        meta_para = doc.add_paragraph()
        meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        meta_text = f"{req.author}  |  {datetime.now().strftime('%Y年%m月%d日')}"
        if req.subject:
            meta_text = f"{req.subject}\n{meta_text}"
        meta_run = meta_para.add_run(meta_text)
        meta_run.font.size = Pt(11)
        sr, sg, sb = preset["secondary_rgb"]
        meta_run.font.color.rgb = RGBColor(sr, sg, sb)
        meta_run.font.name = "微软雅黑"
        meta_run._r.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")

        doc.add_paragraph("")  # 空行

        # ── 正文内容 ──────────────────────────────────────────────────────
        for s in req.sections:
            ct = (
                s.content_type.lower()
                if isinstance(s.content_type, str)
                else s.content_type.value
            )

            if ct == ContentType.HEADING:
                head_level = max(1, min(s.level, 3))
                h = doc.add_heading(str(s.content), level=head_level)
                h.alignment = WD_ALIGN_PARAGRAPH.LEFT
                # 应用样式色彩
                for run in h.runs:
                    run.font.color.rgb = RGBColor(pr, pg, pb)
                    run.font.name = "微软雅黑"
                    try:
                        run._r.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
                    except Exception:
                        pass

            elif ct == ContentType.TEXT:
                p = doc.add_paragraph(str(s.content))
                p.paragraph_format.space_after = Pt(6)
                for run in p.runs:
                    run.font.size = Pt(preset["body_font_size"])
                    run.font.name = "微软雅黑"
                    try:
                        run._r.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
                    except Exception:
                        pass

            elif ct == ContentType.TABLE:
                rows_data = s.content
                if not isinstance(rows_data, list) or len(rows_data) == 0:
                    continue
                table = doc.add_table(rows=len(rows_data), cols=len(rows_data[0]))
                table.style = "Table Grid"
                hbr, hbg, hbb = preset["header_bg_rgb"]
                sbr, sbg, sbb = preset["stripe_bg_rgb"]
                from docx.oxml import OxmlElement

                for row_idx, row_data in enumerate(rows_data):
                    row = table.rows[row_idx]
                    for col_idx, cell_val in enumerate(row_data):
                        cell = row.cells[col_idx]
                        cell.text = str(cell_val)
                        run = (
                            cell.paragraphs[0].runs[0]
                            if cell.paragraphs[0].runs
                            else cell.paragraphs[0].add_run(str(cell_val))
                        )
                        run.font.size = Pt(10)
                        run.font.name = "微软雅黑"
                        try:
                            run._r.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
                        except Exception:
                            pass
                        if row_idx == 0:
                            run.bold = True
                            # 表头背景色
                            tc = cell._tc
                            tcPr = tc.get_or_add_tcPr()
                            shd = OxmlElement("w:shd")
                            shd.set(qn("w:fill"), f"{hbr:02X}{hbg:02X}{hbb:02X}")
                            shd.set(qn("w:val"), "clear")
                            tcPr.append(shd)
                            hfr, hfg, hfb = preset["header_fg_rgb"]
                            run.font.color.rgb = RGBColor(hfr, hfg, hfb)
                        elif row_idx % 2 == 0:
                            # 斑马条纹
                            tc = cell._tc
                            tcPr = tc.get_or_add_tcPr()
                            shd = OxmlElement("w:shd")
                            shd.set(qn("w:fill"), f"{sbr:02X}{sbg:02X}{sbb:02X}")
                            shd.set(qn("w:val"), "clear")
                            tcPr.append(shd)
                if s.caption:
                    cap = doc.add_paragraph(s.caption)
                    cap.style = "Caption"
                    cap.alignment = WD_ALIGN_PARAGRAPH.CENTER

            elif ct in (ContentType.BULLET_LIST, ContentType.NUMBERED_LIST):
                items = s.content if isinstance(s.content, list) else [str(s.content)]
                style = (
                    "List Bullet" if ct == ContentType.BULLET_LIST else "List Number"
                )
                for item in items:
                    p = doc.add_paragraph(str(item), style=style)
                    for run in p.runs:
                        run.font.size = Pt(preset["body_font_size"])
                        run.font.name = "微软雅黑"
                        try:
                            run._r.rPr.rFonts.set(qn("w:eastAsia"), "微软雅黑")
                        except Exception:
                            pass

            elif ct == ContentType.CODE:
                p = doc.add_paragraph()
                run = p.add_run(str(s.content))
                run.font.name = "Courier New"
                run.font.size = Pt(9)
                p.paragraph_format.left_indent = Inches(0.5)
                p.paragraph_format.space_before = Pt(6)
                p.paragraph_format.space_after = Pt(6)

            elif ct == ContentType.QUOTE:
                p = doc.add_paragraph(str(s.content))
                p.paragraph_format.left_indent = Inches(0.5)
                p.paragraph_format.right_indent = Inches(0.5)
                for run in p.runs:
                    ar, ag, ab = preset["accent_rgb"]
                    run.font.color.rgb = RGBColor(ar, ag, ab)
                    run.font.italic = True
                    run.font.size = Pt(preset["body_font_size"])

            elif ct == ContentType.IMAGE:
                try:
                    img_path = str(s.content)
                    if os.path.exists(img_path):
                        doc.add_picture(img_path, width=Inches(5.5))
                        last_para = doc.paragraphs[-1]
                        last_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    if s.caption:
                        cap = doc.add_paragraph(s.caption)
                        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                except Exception as img_err:
                    doc.add_paragraph(f"[图片加载失败: {img_err}]")

            elif ct == ContentType.PAGE_BREAK:
                doc.add_page_break()

            elif ct == ContentType.HORIZONTAL_RULE:
                doc.add_paragraph("─" * 40)

        # ── 文档属性 ──────────────────────────────────────────────────────
        core_props = doc.core_properties
        core_props.title = req.title
        core_props.author = req.author
        core_props.subject = req.subject
        core_props.keywords = req.keywords
        core_props.created = datetime.now()

        doc.save(str(output_path))
        logger.info(f"[DocGen] Word 生成完成: {output_path}")


# ══════════════════════════════════════════════════════════════════════════════
# PDF 生成器（基于 reportlab）
# ══════════════════════════════════════════════════════════════════════════════


class _PdfGenerator:
    """基于 reportlab 生成专业品质的 PDF 文档，支持中文字体自动检测。"""

    def _register_fonts(self) -> Tuple[str, str]:
        """
        注册系统中文字体并返回 (正文字体名, 粗体字体名)。
        若找不到中文字体则回退到 Helvetica。
        """
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        zh_path = _find_chinese_font()
        if zh_path:
            try:
                # .ttc 文件需指定子字体索引
                if zh_path.lower().endswith(".ttc"):
                    from reportlab.pdfbase.ttfonts import TTFont as _TTF

                    pdfmetrics.registerFont(_TTF("KotoCNBody", zh_path, subfontIndex=0))
                    pdfmetrics.registerFont(_TTF("KotoCNBold", zh_path, subfontIndex=0))
                else:
                    pdfmetrics.registerFont(TTFont("KotoCNBody", zh_path))
                    pdfmetrics.registerFont(TTFont("KotoCNBold", zh_path))
                return "KotoCNBody", "KotoCNBold"
            except Exception as fe:
                logger.warning(f"[DocGen] 中文字体注册失败 ({zh_path}): {fe}")
        logger.warning(
            "[DocGen] 未找到中文字体，PDF 正文可能显示为方框（推荐安装微软雅黑）"
        )
        return "Helvetica", "Helvetica-Bold"

    def generate(self, req: DocGenRequest, output_path: Path) -> None:
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import cm
            from reportlab.lib import colors
            from reportlab.platypus import (
                SimpleDocTemplate,
                Paragraph,
                Spacer,
                Table,
                TableStyle,
                PageBreak,
                HRFlowable,
                Preformatted,
                ListFlowable,
                ListItem,
            )
            from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_JUSTIFY
        except ImportError as e:
            raise ImportError(
                f"reportlab 未安装: {e}\n" "请运行: pip install reportlab>=4.0.0"
            )

        preset = _get_preset(req.style_preset)
        body_font, bold_font = self._register_fonts()

        pr, pg, pb = preset["primary_rgb"]
        sr, sg, sb = preset["secondary_rgb"]
        ar, ag, ab = preset["accent_rgb"]
        hbr, hbg, hbb = preset["header_bg_rgb"]
        hfr, hfg, hfb = preset["header_fg_rgb"]
        stripe_r, stripe_g, stripe_b = preset["stripe_bg_rgb"]

        primary_color = colors.Color(pr / 255, pg / 255, pb / 255)
        secondary_color = colors.Color(sr / 255, sg / 255, sb / 255)
        accent_color = colors.Color(ar / 255, ag / 255, ab / 255)
        header_bg_color = colors.Color(hbr / 255, hbg / 255, hbb / 255)
        header_fg_color = colors.Color(hfr / 255, hfg / 255, hfb / 255)
        stripe_color = colors.Color(stripe_r / 255, stripe_g / 255, stripe_b / 255)

        styles = getSampleStyleSheet()
        base_font_size = preset["body_font_size"]

        # 自定义样式
        style_title = ParagraphStyle(
            "KotoTitle",
            fontName=bold_font,
            fontSize=preset["title_font_size"],
            textColor=primary_color,
            spaceAfter=12,
            alignment=TA_CENTER,
        )
        style_meta = ParagraphStyle(
            "KotoMeta",
            fontName=body_font,
            fontSize=10,
            textColor=secondary_color,
            spaceAfter=20,
            alignment=TA_CENTER,
        )
        style_h1 = ParagraphStyle(
            "KotoH1",
            fontName=bold_font,
            fontSize=preset["h1_font_size"],
            textColor=primary_color,
            spaceBefore=14,
            spaceAfter=6,
            borderPad=(0, 0, 3, 0),
        )
        style_h2 = ParagraphStyle(
            "KotoH2",
            fontName=bold_font,
            fontSize=preset["h2_font_size"],
            textColor=secondary_color,
            spaceBefore=10,
            spaceAfter=4,
        )
        style_h3 = ParagraphStyle(
            "KotoH3",
            fontName=bold_font,
            fontSize=preset["h3_font_size"],
            textColor=accent_color,
            spaceBefore=6,
            spaceAfter=4,
        )
        style_body = ParagraphStyle(
            "KotoBody",
            fontName=body_font,
            fontSize=base_font_size,
            leading=base_font_size * preset["line_spacing"],
            spaceAfter=6,
            alignment=TA_JUSTIFY,
        )
        style_quote = ParagraphStyle(
            "KotoQuote",
            fontName=body_font,
            fontSize=base_font_size,
            textColor=accent_color,
            leftIndent=30,
            rightIndent=30,
            spaceAfter=8,
            fontStyle="italic" if body_font == "Helvetica" else "normal",
        )
        style_code = ParagraphStyle(
            "KotoCode",
            fontName="Courier",
            fontSize=9,
            leftIndent=20,
            backColor=colors.Color(0.95, 0.95, 0.95),
            spaceAfter=8,
        )
        style_caption = ParagraphStyle(
            "KotoCaption",
            fontName=body_font,
            fontSize=9,
            textColor=secondary_color,
            alignment=TA_CENTER,
            spaceAfter=8,
        )

        doc = SimpleDocTemplate(
            str(output_path),
            pagesize=A4,
            leftMargin=3.0 * cm,
            rightMargin=2.5 * cm,
            topMargin=2.5 * cm,
            bottomMargin=2.5 * cm,
            title=req.title,
            author=req.author,
            subject=req.subject,
        )

        story = []

        # ── 标题区 ────────────────────────────────────────────────────────
        story.append(Paragraph(req.title, style_title))
        meta_parts = [req.author, datetime.now().strftime("%Y年%m月%d日")]
        if req.subject:
            meta_parts.insert(0, req.subject)
        story.append(Paragraph("  |  ".join(meta_parts), style_meta))
        story.append(HRFlowable(width="100%", thickness=1.5, color=primary_color))
        story.append(Spacer(1, 12))

        # ── 正文 ──────────────────────────────────────────────────────────
        for s in req.sections:
            ct = (
                s.content_type.lower()
                if isinstance(s.content_type, str)
                else s.content_type.value
            )

            if ct == ContentType.HEADING:
                style_map = {1: style_h1, 2: style_h2, 3: style_h3}
                h_style = style_map.get(s.level, style_h1)
                # Add underline for H1
                if s.level == 1:
                    story.append(Paragraph(str(s.content), h_style))
                    story.append(
                        HRFlowable(
                            width="100%",
                            thickness=0.5,
                            color=primary_color,
                            spaceAfter=4,
                        )
                    )
                else:
                    story.append(Paragraph(str(s.content), h_style))

            elif ct == ContentType.TEXT:
                story.append(Paragraph(str(s.content), style_body))

            elif ct == ContentType.TABLE:
                rows_data = s.content
                if not isinstance(rows_data, list) or len(rows_data) == 0:
                    continue
                table_data = []
                for row in rows_data:
                    table_data.append([str(c) for c in row])

                col_count = max(len(row) for row in table_data)
                available_width = doc.width
                col_width = available_width / col_count

                tbl = Table(table_data, colWidths=[col_width] * col_count, repeatRows=1)
                tbl_style = TableStyle(
                    [
                        # 表头
                        ("BACKGROUND", (0, 0), (-1, 0), header_bg_color),
                        ("TEXTCOLOR", (0, 0), (-1, 0), header_fg_color),
                        ("FONTNAME", (0, 0), (-1, 0), bold_font),
                        ("FONTSIZE", (0, 0), (-1, 0), 10),
                        ("ALIGN", (0, 0), (-1, 0), "CENTER"),
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        # 边框
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.Color(0.7, 0.7, 0.7)),
                        ("BOX", (0, 0), (-1, -1), 1.0, primary_color),
                        # 正文
                        ("FONTNAME", (0, 1), (-1, -1), body_font),
                        ("FONTSIZE", (0, 1), (-1, -1), 9),
                        (
                            "ROWBACKGROUNDS",
                            (0, 1),
                            (-1, -1),
                            [colors.white, stripe_color],
                        ),
                        ("TOPPADDING", (0, 0), (-1, -1), 6),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                        ("LEFTPADDING", (0, 0), (-1, -1), 8),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                    ]
                )
                tbl.setStyle(tbl_style)
                story.append(tbl)
                if s.caption:
                    story.append(Paragraph(s.caption, style_caption))

            elif ct == ContentType.BULLET_LIST:
                items = s.content if isinstance(s.content, list) else [str(s.content)]
                story.append(
                    ListFlowable(
                        [
                            ListItem(
                                Paragraph(str(item), style_body),
                                bulletColor=primary_color,
                            )
                            for item in items
                        ],
                        bulletType="bullet",
                    )
                )

            elif ct == ContentType.NUMBERED_LIST:
                items = s.content if isinstance(s.content, list) else [str(s.content)]
                story.append(
                    ListFlowable(
                        [ListItem(Paragraph(str(item), style_body)) for item in items],
                        bulletType="1",
                    )
                )

            elif ct == ContentType.CODE:
                story.append(Preformatted(str(s.content), style_code))

            elif ct == ContentType.QUOTE:
                story.append(Paragraph(f'"{s.content}"', style_quote))

            elif ct == ContentType.IMAGE:
                try:
                    from reportlab.platypus import Image as RLImage

                    img_path = str(s.content)
                    if os.path.exists(img_path):
                        img = RLImage(img_path, width=14 * cm)
                        story.append(img)
                    if s.caption:
                        story.append(Paragraph(s.caption, style_caption))
                except Exception as img_err:
                    story.append(Paragraph(f"[图片加载失败: {img_err}]", style_body))

            elif ct == ContentType.PAGE_BREAK:
                story.append(PageBreak())

            elif ct == ContentType.HORIZONTAL_RULE:
                story.append(
                    HRFlowable(width="100%", thickness=0.5, color=accent_color)
                )
                story.append(Spacer(1, 6))

        doc.build(story)
        logger.info(f"[DocGen] PDF 生成完成: {output_path}")


# ══════════════════════════════════════════════════════════════════════════════
# Excel 生成器
# ══════════════════════════════════════════════════════════════════════════════


class _ExcelGenerator:
    """基于 openpyxl 生成专业品质的 Excel 报表。"""

    def _apply_header_style(
        self, ws, row_idx: int, col_count: int, preset: Dict, bg_rgb, fg_rgb
    ) -> None:
        from openpyxl.styles import PatternFill, Font, Alignment, Border, Side

        thick_border = Border(
            bottom=Side(style="medium"),
            top=Side(style="medium"),
        )
        bg_hex = f"{bg_rgb[0]:02X}{bg_rgb[1]:02X}{bg_rgb[2]:02X}"
        fg_hex = f"{fg_rgb[0]:02X}{fg_rgb[1]:02X}{fg_rgb[2]:02X}"

        for col in range(1, col_count + 1):
            cell = ws.cell(row=row_idx, column=col)
            cell.fill = PatternFill(
                start_color=bg_hex, end_color=bg_hex, fill_type="solid"
            )
            cell.font = Font(name="微软雅黑", bold=True, color=fg_hex, size=10)
            cell.alignment = Alignment(
                horizontal="center", vertical="center", wrap_text=True
            )
            cell.border = thick_border

    def generate(self, req: DocGenRequest, output_path: Path) -> None:
        try:
            import openpyxl
            from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
            from openpyxl.utils import get_column_letter
        except ImportError as e:
            raise ImportError(f"openpyxl 未安装: {e}")

        preset = _get_preset(req.style_preset)
        wb = openpyxl.Workbook()
        pr, pg, pb = preset["primary_rgb"]
        sr, sg, sb = preset["secondary_rgb"]
        hbr, hbg, hbb = preset["header_bg_rgb"]
        hfr, hfg, hfb = preset["header_fg_rgb"]
        stripe_r, stripe_g, stripe_b = preset["stripe_bg_rgb"]

        primary_hex = f"{pr:02X}{pg:02X}{pb:02X}"
        stripe_hex = f"{stripe_r:02X}{stripe_g:02X}{stripe_b:02X}"

        # ── 确定 Sheet 数据 ───────────────────────────────────────────────
        sheets_to_gen: List[ExcelSheet] = req.excel_sheets if req.excel_sheets else []

        # 如果没有提供 excel_sheets，从 sections 提取表格数据建立默认 Sheet
        if not sheets_to_gen:
            default_rows: List[List[Any]] = []
            default_headers: List[str] = []
            for s in req.sections:
                ct = (
                    s.content_type.lower()
                    if isinstance(s.content_type, str)
                    else s.content_type.value
                )
                if (
                    ct == ContentType.TABLE
                    and isinstance(s.content, list)
                    and len(s.content) > 0
                ):
                    if not default_headers:
                        default_headers = [str(c) for c in s.content[0]]
                        default_rows = [list(row) for row in s.content[1:]]
            if default_headers:
                sheets_to_gen = [
                    ExcelSheet(
                        name="数据",
                        title=req.title,
                        column_headers=default_headers,
                        rows=default_rows,
                    )
                ]
            else:
                # 如果完全没有表格数据，生成一个说明 Sheet
                sheets_to_gen = [
                    ExcelSheet(
                        name="说明",
                        title=req.title,
                        column_headers=["内容"],
                        rows=[
                            [s.content]
                            for s in req.sections
                            if isinstance(s.content, str)
                        ],
                    )
                ]

        # 删除默认空 Sheet
        if "Sheet" in wb.sheetnames:
            del wb["Sheet"]

        for sheet_def in sheets_to_gen:
            ws = wb.create_sheet(title=sheet_def.name[:31])  # Excel 限 31 字符

            row_cursor = 1

            # Sheet 大标题
            if sheet_def.title:
                ws.cell(row=row_cursor, column=1, value=sheet_def.title)
                ws.cell(row=row_cursor, column=1).font = Font(
                    name="微软雅黑",
                    bold=True,
                    size=14,
                    color=primary_hex,
                )
                ws.merge_cells(
                    start_row=row_cursor,
                    start_column=1,
                    end_row=row_cursor,
                    end_column=max(len(sheet_def.column_headers), 1),
                )
                ws.cell(row=row_cursor, column=1).alignment = Alignment(
                    horizontal="center"
                )
                row_cursor += 1
                # 生成日期行
                ws.cell(
                    row=row_cursor,
                    column=1,
                    value=f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}",
                )
                ws.cell(row=row_cursor, column=1).font = Font(
                    name="微软雅黑", size=9, color="888888"
                )
                row_cursor += 1

            # 表头
            headers = sheet_def.column_headers
            if headers:
                for col_idx, header in enumerate(headers, start=1):
                    ws.cell(row=row_cursor, column=col_idx, value=header)
                self._apply_header_style(
                    ws,
                    row_cursor,
                    len(headers),
                    preset,
                    (hbr, hbg, hbb),
                    (hfr, hfg, hfb),
                )
                row_cursor += 1

            # 数据行
            thin_border = Border(
                left=Side(style="thin", color="DDDDDD"),
                right=Side(style="thin", color="DDDDDD"),
                top=Side(style="thin", color="DDDDDD"),
                bottom=Side(style="thin", color="DDDDDD"),
            )
            totals: Dict[int, Any] = {}  # col_idx → running sum
            for data_row in sheet_def.rows:
                for col_idx, cell_val in enumerate(data_row, start=1):
                    c = ws.cell(row=row_cursor, column=col_idx, value=cell_val)
                    c.font = Font(name="微软雅黑", size=10)
                    c.border = thin_border
                    c.alignment = Alignment(vertical="center", wrap_text=True)
                    # 斑马条纹
                    if row_cursor % 2 == 0:
                        c.fill = PatternFill(
                            start_color=stripe_hex,
                            end_color=stripe_hex,
                            fill_type="solid",
                        )
                    # 尝试累加数值列
                    if sheet_def.add_totals_row:
                        try:
                            totals[col_idx] = totals.get(col_idx, 0) + float(
                                str(cell_val)
                                .replace(",", "")
                                .replace("¥", "")
                                .replace("%", "")
                            )
                        except ValueError:
                            totals.setdefault(col_idx, "─")
                row_cursor += 1

            # 合计行
            if sheet_def.add_totals_row and totals:
                ws.cell(row=row_cursor, column=1, value="合 计")
                ws.cell(row=row_cursor, column=1).font = Font(
                    name="微软雅黑", bold=True, size=10
                )
                for col_idx, val in totals.items():
                    if col_idx == 1:
                        continue
                    c = ws.cell(
                        row=row_cursor,
                        column=col_idx,
                        value=val if isinstance(val, str) else round(val, 2),
                    )
                    c.font = Font(
                        name="微软雅黑", bold=True, size=10, color=primary_hex
                    )
                    c.fill = PatternFill(
                        start_color=stripe_hex, end_color=stripe_hex, fill_type="solid"
                    )
                row_cursor += 1

            # ── 列宽自适应 ────────────────────────────────────────────────
            for col_idx in range(1, (len(headers) or 1) + 1):
                col_letter = get_column_letter(col_idx)
                if sheet_def.column_widths and col_idx - 1 < len(
                    sheet_def.column_widths
                ):
                    ws.column_dimensions[col_letter].width = sheet_def.column_widths[
                        col_idx - 1
                    ]
                else:
                    max_len = 0
                    for row_cells in ws.iter_rows(min_col=col_idx, max_col=col_idx):
                        for c in row_cells:
                            if c.value:
                                # 中文字符计2个单位
                                val_str = str(c.value)
                                char_width = sum(
                                    2 if ord(ch) > 127 else 1 for ch in val_str
                                )
                                max_len = max(max_len, char_width)
                    ws.column_dimensions[col_letter].width = min(max_len + 4, 50)

            # 冻结表头
            if sheet_def.freeze_header and headers:
                freeze_row = (3 if sheet_def.title else 1) + 1
                ws.freeze_panes = f"A{freeze_row}"

            # 行高
            ws.row_dimensions[1].height = 30 if sheet_def.title else 20

        # ── 元数据 ────────────────────────────────────────────────────────
        wb.properties.title = req.title
        wb.properties.creator = req.author
        wb.properties.subject = req.subject
        wb.save(str(output_path))
        logger.info(f"[DocGen] Excel 生成完成: {output_path}")


# ══════════════════════════════════════════════════════════════════════════════
# PowerPoint 生成器
# ══════════════════════════════════════════════════════════════════════════════


class _PptxGenerator:
    """基于 python-pptx 生成专业品质的 PowerPoint 演示文稿。"""

    def generate(self, req: DocGenRequest, output_path: Path) -> None:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
            from pptx.util import Inches, Pt
        except ImportError as e:
            raise ImportError(f"python-pptx 未安装: {e}")

        preset = _get_preset(req.style_preset)
        prs = Presentation()

        # 幻灯片尺寸
        if req.pptx_layout == "widescreen":
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
        else:
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

        slide_layouts = prs.slide_layouts
        # 0 = Title Slide, 1 = Title & Content, 5 = Title Only, 6 = Blank
        BLANK_LAYOUT = 6
        TITLE_LAYOUT = 0
        CONTENT_LAYOUT = 1

        pr, pg, pb = preset["primary_rgb"]
        sr, sg, sb = preset["secondary_rgb"]
        ar, ag, ab = preset["accent_rgb"]
        hbr, hbg, hbb = preset["header_bg_rgb"]
        hfr, hfg, hfb = preset["header_fg_rgb"]

        def _rgb(t):
            from pptx.dml.color import RGBColor

            return RGBColor(t[0], t[1], t[2])

        def _set_font(
            run_or_para, font_name="微软雅黑", size_pt=None, bold=False, color=None
        ):
            try:
                tf = run_or_para.font
                tf.name = font_name
                if size_pt:
                    tf.size = Pt(size_pt)
                tf.bold = bold
                if color:
                    tf.color.rgb = _rgb(color)
            except Exception:
                pass

        # ── 封面幻灯片 ─────────────────────────────────────────────────────
        title_slide = prs.slides.add_slide(slide_layouts[TITLE_LAYOUT])
        title_ph = title_slide.shapes.title
        subtitle_ph = (
            title_slide.placeholders[1] if len(title_slide.placeholders) > 1 else None
        )

        title_ph.text = req.title
        for para in title_ph.text_frame.paragraphs:
            for run in para.runs:
                _set_font(
                    run,
                    size_pt=preset["title_font_size"],
                    bold=True,
                    color=(pr, pg, pb),
                )

        if subtitle_ph:
            sub_parts = [req.author, datetime.now().strftime("%Y年%m月%d日")]
            if req.subject:
                sub_parts.insert(0, req.subject)
            subtitle_ph.text = "  |  ".join(sub_parts)
            for para in subtitle_ph.text_frame.paragraphs:
                for run in para.runs:
                    _set_font(run, size_pt=14, color=(sr, sg, sb))

        # ── 正文幻灯片 ─────────────────────────────────────────────────────
        # 将 sections 分组：H1 创建新幻灯片，其余内容追加到当前幻灯片
        current_slide = None
        current_tf = None

        def _new_content_slide(heading_text: str):
            nonlocal current_slide, current_tf
            slide = prs.slides.add_slide(slide_layouts[CONTENT_LAYOUT])
            title_sp = slide.shapes.title
            if title_sp:
                title_sp.text = heading_text
                for para in title_sp.text_frame.paragraphs:
                    for run in para.runs:
                        _set_font(
                            run,
                            size_pt=preset["h1_font_size"],
                            bold=True,
                            color=(pr, pg, pb),
                        )
            # Content placeholder
            content_sp = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    content_sp = ph
                    break
            if content_sp:
                current_tf = content_sp.text_frame
                current_tf.clear()
            else:
                current_tf = None
            current_slide = slide
            return slide

        def _add_text_to_tf(
            text: str, bold=False, size_pt=None, color=None, bullet=False
        ):
            if current_tf is None:
                return
            para = current_tf.add_paragraph()
            run = para.add_run()
            run.text = text
            _set_font(
                run, size_pt=size_pt or preset["body_font_size"], bold=bold, color=color
            )
            if bullet:
                para.level = 1

        for s in req.sections:
            ct = (
                s.content_type.lower()
                if isinstance(s.content_type, str)
                else s.content_type.value
            )

            if ct == ContentType.HEADING and s.level == 1:
                _new_content_slide(str(s.content))

            elif ct == ContentType.HEADING and s.level >= 2:
                _add_text_to_tf(
                    str(s.content),
                    bold=True,
                    size_pt=preset["h2_font_size"],
                    color=(sr, sg, sb),
                )

            elif ct == ContentType.TEXT:
                if current_slide is None:
                    _new_content_slide(req.title)
                _add_text_to_tf(str(s.content), size_pt=preset["body_font_size"])

            elif ct in (ContentType.BULLET_LIST, ContentType.NUMBERED_LIST):
                if current_slide is None:
                    _new_content_slide("内容")
                items = s.content if isinstance(s.content, list) else [str(s.content)]
                for item in items:
                    _add_text_to_tf(
                        f"• {item}" if ct == ContentType.BULLET_LIST else item,
                        size_pt=preset["body_font_size"],
                    )

            elif ct == ContentType.TABLE:
                # PPT 中的表格：新建一张独立幻灯片
                rows_data = s.content
                if not isinstance(rows_data, list) or len(rows_data) == 0:
                    continue
                table_slide = prs.slides.add_slide(slide_layouts[5])  # Title Only
                if table_slide.shapes.title and s.caption:
                    table_slide.shapes.title.text = s.caption or "数据表格"

                row_count = len(rows_data)
                col_count = max(len(r) for r in rows_data)
                left = Inches(0.5)
                top = Inches(1.5)
                width = prs.slide_width - Inches(1.0)
                height = min(
                    Inches(0.35 * row_count + 0.3), prs.slide_height - Inches(2.0)
                )

                tbl_shape = table_slide.shapes.add_table(
                    row_count, col_count, left, top, width, height
                )
                tbl = tbl_shape.table
                for row_idx, row_data in enumerate(rows_data):
                    for col_idx, cell_val in enumerate(row_data):
                        cell = tbl.cell(row_idx, col_idx)
                        cell.text = str(cell_val)
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                is_header = row_idx == 0
                                _set_font(
                                    run,
                                    size_pt=9,
                                    bold=is_header,
                                    color=(hfr, hfg, hfb) if is_header else None,
                                )
                                if is_header:
                                    from pptx.dml.color import RGBColor as PPTXRGB

                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = PPTXRGB(hbr, hbg, hbb)

            elif ct == ContentType.PAGE_BREAK:
                current_slide = None  # 强制下一次内容到新幻灯片

        prs.save(str(output_path))
        logger.info(f"[DocGen] PPT 生成完成: {output_path}")


# ══════════════════════════════════════════════════════════════════════════════
# 对外 API（单例服务）
# ══════════════════════════════════════════════════════════════════════════════


class DocGenService:
    """
    Koto 文档生成服务入口。
    提供 Word / PDF / Excel / PowerPoint 四种格式的生成方法。
    """

    def __init__(self):
        self._word_gen = _WordGenerator()
        self._pdf_gen = _PdfGenerator()
        self._excel_gen = _ExcelGenerator()
        self._pptx_gen = _PptxGenerator()

    # ── Word ──────────────────────────────────────────────────────────────────

    def generate_word(self, req: DocGenRequest) -> str:
        """生成 Word (.docx) 文档，返回输出文件路径。"""
        output_path = _build_output_path(req.output_dir, req.output_filename, "docx")
        self._word_gen.generate(req, output_path)
        return str(output_path)

    # ── PDF ───────────────────────────────────────────────────────────────────

    def generate_pdf(self, req: DocGenRequest) -> str:
        """生成 PDF 文档，返回输出文件路径。需要 reportlab。"""
        output_path = _build_output_path(req.output_dir, req.output_filename, "pdf")
        self._pdf_gen.generate(req, output_path)
        return str(output_path)

    # ── Excel ─────────────────────────────────────────────────────────────────

    def generate_excel(self, req: DocGenRequest) -> str:
        """生成 Excel (.xlsx) 报表，返回输出文件路径。需要 openpyxl。"""
        output_path = _build_output_path(req.output_dir, req.output_filename, "xlsx")
        self._excel_gen.generate(req, output_path)
        return str(output_path)

    # ── PPT ───────────────────────────────────────────────────────────────────

    def generate_presentation(self, req: DocGenRequest) -> str:
        """生成 PowerPoint (.pptx)，返回输出文件路径。需要 python-pptx。"""
        output_path = _build_output_path(req.output_dir, req.output_filename, "pptx")
        self._pptx_gen.generate(req, output_path)
        return str(output_path)

    # ── 便捷方法：从字典构建请求 ──────────────────────────────────────────────

    @staticmethod
    def build_request(data: Dict[str, Any]) -> DocGenRequest:
        """
        从 LLM 输出的字典结构构建 DocGenRequest。

        期望 data 格式::

            {
                "title": "报告标题",
                "author": "作者",
                "subject": "主题",
                "style_preset": "professional",
                "sections": [
                    {"content_type": "heading", "content": "第一章", "level": 1},
                    {"content_type": "text",    "content": "正文内容..."},
                    {"content_type": "table",   "content": [["列1","列2"],["A","B"]]},
                    {"content_type": "bullet",  "content": ["要点1","要点2"]},
                ],
                "excel_sheets": [
                    {
                        "name": "销售数据",
                        "title": "2026年Q1销售汇总",
                        "column_headers": ["地区","销售额","增长率"],
                        "rows": [["华东","¥520万","+22%"],["华南","¥380万","+15%"]],
                        "add_totals_row": true
                    }
                ]
            }
        """
        sections = [
            DocSection(
                content_type=sec.get("content_type", "text"),
                content=sec.get("content", ""),
                level=sec.get("level", 1),
                style=sec.get("style"),
                caption=sec.get("caption"),
            )
            for sec in data.get("sections", [])
        ]
        excel_sheets = [
            ExcelSheet(
                name=sh.get("name", "Sheet1"),
                title=sh.get("title", ""),
                column_headers=sh.get("column_headers", []),
                rows=sh.get("rows", []),
                add_totals_row=sh.get("add_totals_row", False),
                column_widths=sh.get("column_widths"),
                freeze_header=sh.get("freeze_header", True),
            )
            for sh in data.get("excel_sheets", [])
        ]
        return DocGenRequest(
            title=data.get("title", "未命名文档"),
            sections=sections,
            author=data.get("author", "Koto AI"),
            subject=data.get("subject", ""),
            keywords=data.get("keywords", ""),
            language=data.get("language", "zh"),
            style_preset=data.get("style_preset", "professional"),
            theme_color=data.get("theme_color", "#1F4E79"),
            output_dir=data.get("output_dir", ""),
            output_filename=data.get("output_filename", ""),
            excel_sheets=excel_sheets,
            pptx_layout=data.get("pptx_layout", "widescreen"),
            meta=data.get("meta", {}),
        )
