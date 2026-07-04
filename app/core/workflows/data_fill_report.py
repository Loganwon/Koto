# ══════════════════════════════════════════════════════════════
# data_fill_report.py — 数据填报（Excel → Word/PPT）
#
# 用户场景：
#   有一个 Excel 数据表，需要把数据填入 Word 报告或 PPT 模板。
#   AI 智能匹配占位符与数据列，批量替换生成成品文档。
# ══════════════════════════════════════════════════════════════

from __future__ import annotations

import json
import logging
import re
import shutil
from pathlib import Path
from typing import Any

from app.core.workflow_engine import (
    WorkflowExecutor,
    sse_error,
    sse_output,
    sse_step_done,
    sse_step_start,
)

logger = logging.getLogger(__name__)

_MAPPING_SYSTEM = """你是一个数据映射助手。
请将模板中的占位符与数据表的列名进行匹配。

模板中的占位符列表：
{placeholders_json}

数据表的列名和前几行样例：
{data_preview}

{instruction}

输出 JSON 数组，每个元素：
{{
  "placeholder": "模板中的占位符（精确原文）",
  "data_col": "数据表中对应的列名",
  "row_selector": "first | last | sum | average | max | min | concat"
}}

规则：
1. placeholder 必须与输入列表完全一致
2. 如果无法匹配某个占位符，设 data_col 为 null
3. row_selector 说明从数据中取哪个值（通常是 first，除非是汇总场景）
4. 只输出 JSON 数组"""

# 支持的占位符格式
_PLACEHOLDER_PATTERNS = [
    r"\{\{(.+?)\}\}",  # {{字段名}}
    r"<<(.+?)>>",  # <<字段名>>
    r"\[(.+?)\]",  # [字段名]
    r"__(.+?)__",  # __字段名__
    r"\$\{(.+?)\}",  # ${字段名}
]


class DataFillReport(WorkflowExecutor):
    """
    数据填报工作流。

    params 期望字段:
        data_file:     str — Excel/CSV 数据源路径
        template_file: str — Word/PPT 模板文件路径
        instruction:   str — 填写说明（可选）
        model_mode:    str — "auto" | "local"
    """

    WORKFLOW_ID = "data_fill_report"
    WORKFLOW_NAME = "数据填报"

    def execute(self, params: dict, yield_event) -> Any:
        data_file: str = params.get("data_file") or ""
        template_file: str = params.get("template_file") or ""
        instruction: str = params.get("instruction") or ""
        model_mode: str = params.get("model_mode") or "auto"

        if not data_file or not template_file:
            yield sse_error(
                "请同时提供数据文件（data_file）和模板文件（template_file）"
            )
            return

        template_ext = Path(template_file).suffix.lower()
        if template_ext not in (".docx", ".pptx"):
            yield sse_error("模板文件必须是 .docx 或 .pptx 格式")
            return

        # ── Step 1: 读取数据 ────────────────────────────────────────
        yield sse_step_start("data", "📊 读取数据源…")
        headers, rows = self._load_data(data_file, model_mode)
        if not headers:
            yield sse_error("数据文件为空或无法解析")
            return
        yield sse_step_done("data", f"📊 {len(rows)} 行 × {len(headers)} 列")

        # ── Step 2: 扫描模板占位符 ──────────────────────────────────
        yield sse_step_start("scan", "🔍 扫描模板占位符…")
        if template_ext == ".docx":
            placeholders, template_texts = self._scan_docx_placeholders(template_file)
        else:
            placeholders, template_texts = self._scan_pptx_placeholders(template_file)

        if not placeholders:
            yield sse_error(
                "模板中未找到占位符（支持格式：{{字段}}、<<字段>>、[字段]、__字段__、${字段}）"
            )
            return
        yield sse_step_done("scan", f"🔍 找到 {len(placeholders)} 个占位符")

        # ── Step 3: LLM 智能映射 ───────────────────────────────────
        yield sse_step_start("mapping", "🤖 AI 字段映射…")
        data_preview = self._build_data_preview(headers, rows, model_mode)
        mappings = self._map_fields(placeholders, data_preview, instruction, model_mode)
        valid_mappings = [m for m in mappings if m.get("data_col")]
        yield sse_step_done(
            "mapping", f"🤖 成功映射 {len(valid_mappings)}/{len(placeholders)} 个字段"
        )

        if not valid_mappings:
            yield sse_error("无法将任何占位符映射到数据列")
            return

        # ── Step 4: 构建替换字典 ────────────────────────────────────
        yield sse_step_start("fill", "📝 填入数据…")
        replacements = self._build_replacements(mappings, headers, rows)

        output_suffix = template_ext
        output_path = self.save_output_file(output_suffix)
        shutil.copy2(template_file, str(output_path))

        if template_ext == ".docx":
            filled_count = self._fill_docx(str(output_path), replacements)
        else:
            filled_count = self._fill_pptx(str(output_path), replacements)

        yield sse_step_done("fill", f"📝 已填入 {filled_count} 处")

        # ── 输出 ─────────────────────────────────────────────────────
        mapping_report = self._build_mapping_report(mappings, replacements)
        output_type = "docx_file" if template_ext == ".docx" else "pptx_file"
        yield sse_output(
            output_type,
            {"path": str(output_path), "filename": f"填报_{output_path.name}"},
            f"填报结果（{filled_count} 处替换）",
        )
        yield sse_output("markdown", mapping_report, "字段映射报告")

    # ── 数据加载 ──────────────────────────────────────────────────────

    def _load_data(
        self, file_path: str, model_mode: str
    ) -> tuple[list[str], list[list[str]]]:
        ext = Path(file_path).suffix.lower()
        try:
            if ext in (".xlsx", ".xls"):
                import openpyxl

                wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
                ws = wb.active
                all_rows = []
                for row in ws.iter_rows(values_only=True):
                    all_rows.append([str(c) if c is not None else "" for c in row])
                wb.close()
                if not all_rows:
                    return [], []
                return all_rows[0], all_rows[1:]
            elif ext == ".csv":
                import csv

                with open(file_path, "r", encoding="utf-8-sig", errors="replace") as f:
                    reader = csv.reader(f)
                    all_rows = list(reader)
                if not all_rows:
                    return [], []
                return all_rows[0], all_rows[1:]
        except Exception as e:
            logger.warning("[DataFill] 数据加载失败: %s", e)
        return [], []

    # ── 模板扫描 ──────────────────────────────────────────────────────

    def _scan_docx_placeholders(self, path: str) -> tuple[list[str], list[str]]:
        """扫描 DOCX 中的占位符。"""
        from docx import Document

        doc = Document(path)
        all_text = []
        for para in doc.paragraphs:
            all_text.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    all_text.append(cell.text)
        full_text = "\n".join(all_text)
        placeholders = self._find_placeholders(full_text)
        return placeholders, all_text

    def _scan_pptx_placeholders(self, path: str) -> tuple[list[str], list[str]]:
        """扫描 PPTX 中的占位符。"""
        from pptx import Presentation

        prs = Presentation(path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text
                        if text.strip():
                            all_text.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            all_text.append(cell.text)
        full_text = "\n".join(all_text)
        placeholders = self._find_placeholders(full_text)
        return placeholders, all_text

    def _find_placeholders(self, text: str) -> list[str]:
        """从文本中提取所有占位符（保留原始格式）。"""
        found = set()
        for pattern in _PLACEHOLDER_PATTERNS:
            for m in re.finditer(pattern, text):
                found.add(m.group(0))  # 完整匹配，包括分隔符
        return sorted(found)

    # ── 字段映射 ──────────────────────────────────────────────────────

    def _build_data_preview(
        self, headers: list[str], rows: list[list[str]], model_mode: str
    ) -> str:
        max_rows = 5 if model_mode == "local" else 20
        max_cols = 8 if model_mode == "local" else len(headers)
        preview_headers = headers[:max_cols]
        preview_rows = rows[:max_rows]

        lines = ["\t".join(preview_headers)]
        for r in preview_rows:
            lines.append("\t".join(r[:max_cols]))
        return "\n".join(lines)

    def _map_fields(
        self,
        placeholders: list[str],
        data_preview: str,
        instruction: str,
        model_mode: str,
    ) -> list[dict]:
        """LLM 映射占位符到数据列。"""
        inst_text = f"\n用户说明: {instruction}" if instruction else ""
        system = _MAPPING_SYSTEM.format(
            placeholders_json=json.dumps(placeholders, ensure_ascii=False),
            data_preview=data_preview,
            instruction=inst_text,
        )
        prompt = "请进行字段映射。"
        try:
            result = self.llm_json(prompt, system=system, model_mode=model_mode)
            if isinstance(result, list):
                return result
        except Exception as e:
            logger.warning("[DataFill] LLM 映射失败: %s", e)
        return []

    def _build_replacements(
        self, mappings: list[dict], headers: list[str], rows: list[list[str]]
    ) -> dict[str, str]:
        """构建 {占位符: 替换值} 字典。"""
        col_index = {h: i for i, h in enumerate(headers)}
        result: dict[str, str] = {}

        for m in mappings:
            placeholder = m.get("placeholder", "")
            data_col = m.get("data_col")
            selector = m.get("row_selector", "first")

            if not placeholder or not data_col or data_col not in col_index:
                continue

            ci = col_index[data_col]
            values = [r[ci] for r in rows if ci < len(r) and str(r[ci]).strip()]

            if not values:
                result[placeholder] = ""
                continue

            if selector == "first":
                result[placeholder] = str(values[0])
            elif selector == "last":
                result[placeholder] = str(values[-1])
            elif selector == "sum":
                try:
                    result[placeholder] = str(
                        sum(float(v.replace(",", "")) for v in values)
                    )
                except (ValueError, TypeError):
                    result[placeholder] = str(values[0])
            elif selector == "average":
                try:
                    nums = [float(v.replace(",", "")) for v in values]
                    result[placeholder] = f"{sum(nums)/len(nums):.2f}"
                except (ValueError, TypeError):
                    result[placeholder] = str(values[0])
            elif selector == "max":
                try:
                    result[placeholder] = str(
                        max(float(v.replace(",", "")) for v in values)
                    )
                except (ValueError, TypeError):
                    result[placeholder] = str(values[0])
            elif selector == "min":
                try:
                    result[placeholder] = str(
                        min(float(v.replace(",", "")) for v in values)
                    )
                except (ValueError, TypeError):
                    result[placeholder] = str(values[0])
            elif selector == "concat":
                result[placeholder] = "、".join(str(v) for v in values)
            else:
                result[placeholder] = str(values[0])

        return result

    # ── 文档填充 ──────────────────────────────────────────────────────

    def _fill_docx(self, path: str, replacements: dict[str, str]) -> int:
        """在 DOCX 中执行替换（保留格式）。"""
        from docx import Document

        doc = Document(path)
        count = 0

        # 段落替换
        for para in doc.paragraphs:
            count += self._replace_in_paragraph(para, replacements)

        # 表格替换
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        count += self._replace_in_paragraph(para, replacements)

        # Header/Footer
        for section in doc.sections:
            for header_para in section.header.paragraphs:
                count += self._replace_in_paragraph(header_para, replacements)
            for footer_para in section.footer.paragraphs:
                count += self._replace_in_paragraph(footer_para, replacements)

        doc.save(path)
        return count

    def _replace_in_paragraph(self, para: Any, replacements: dict[str, str]) -> int:
        """在段落的 runs 中进行替换，保留格式。"""
        full_text = para.text
        count = 0
        for placeholder, value in replacements.items():
            if placeholder in full_text:
                # 尝试在单个 run 中替换
                for run in para.runs:
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, value)
                        count += 1
                # 如果占位符跨 run 分割，回退到重建整段
                if placeholder in para.text:
                    # 占位符可能跨 run，需要合并再替换
                    new_text = para.text.replace(placeholder, value)
                    if new_text != para.text and para.runs:
                        # 保留第一个 run 的格式，清空其余
                        para.runs[0].text = new_text
                        for run in para.runs[1:]:
                            run.text = ""
                        count += 1
        return count

    def _fill_pptx(self, path: str, replacements: dict[str, str]) -> int:
        """在 PPTX 中执行替换。"""
        from pptx import Presentation

        prs = Presentation(path)
        count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            for placeholder, value in replacements.items():
                                if placeholder in run.text:
                                    run.text = run.text.replace(placeholder, value)
                                    count += 1
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    for placeholder, value in replacements.items():
                                        if placeholder in run.text:
                                            run.text = run.text.replace(
                                                placeholder, value
                                            )
                                            count += 1

        prs.save(path)
        return count

    # ── 报告 ──────────────────────────────────────────────────────────

    def _build_mapping_report(
        self, mappings: list[dict], replacements: dict[str, str]
    ) -> str:
        lines = ["# 数据填报映射报告\n"]
        lines.append("| 占位符 | 数据列 | 聚合方式 | 填入值 |")
        lines.append("|--------|--------|----------|--------|")
        for m in mappings:
            ph = m.get("placeholder", "")
            col = m.get("data_col") or "未匹配"
            sel = m.get("row_selector", "first")
            val = replacements.get(ph, "—")
            if len(val) > 40:
                val = val[:37] + "…"
            lines.append(f"| `{ph}` | {col} | {sel} | {val} |")
        return "\n".join(lines)
