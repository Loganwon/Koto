# ══════════════════════════════════════════════════════════════
# pptx_data_refresh.py — PPT 数据刷新
#
# 用户场景：
#   月报/周报 PPT 中的数据需要用新 Excel 数据更新，
#   AI 识别 PPT 中的数值/文本 → 匹配新数据 → 替换并保留格式。
# ══════════════════════════════════════════════════════════════

from __future__ import annotations

import json
import logging
import shutil
from pathlib import Path
from typing import Any

from app.core.workflow_engine import (
    WorkflowExecutor,
    sse_error,
    sse_output,
    sse_step_done,
    sse_step_start,
)

logger = logging.getLogger(__name__)

_MAPPING_SYSTEM = """你是一个数据匹配助手。
请将 PPT 中的旧数据与新 Excel 数据进行匹配，确定哪些值需要更新。

PPT 中提取的文本（按 slide 分组）：
{pptx_texts}

新数据表（列名 + 数据）：
{data_preview}

{instruction}

输出 JSON 数组，每个元素：
{{
  "slide_idx": 幻灯片序号（从0开始）,
  "old_value": "PPT 中的旧值（精确文本）",
  "new_value": "替换后的新值",
  "reason": "替换原因（简短说明）"
}}

规则：
1. old_value 必须是 PPT 文本中实际存在的子串
2. 只替换数据值，不要修改标题、标签等固定文本
3. 如果新旧值相同则不输出
4. 只输出 JSON 数组"""


class PptxDataRefresh(WorkflowExecutor):
    """
    PPT 数据刷新工作流。

    params 期望字段:
        pptx_file:   str — 现有 PPT 文件路径
        data_file:   str — 新数据 Excel/CSV 路径
        instruction: str — 更新说明（可选）
        model_mode:  str — "auto" | "local"
    """

    WORKFLOW_ID = "pptx_data_refresh"
    WORKFLOW_NAME = "PPT 数据刷新"

    def execute(self, params: dict, yield_event) -> Any:
        pptx_file: str = params.get("pptx_file") or ""
        data_file: str = params.get("data_file") or ""
        instruction: str = params.get("instruction") or ""
        model_mode: str = params.get("model_mode") or "auto"

        if not pptx_file or not data_file:
            yield sse_error("请同时提供 PPT 文件（pptx_file）和数据文件（data_file）")
            return

        # ── Step 1: 读取新数据 ──────────────────────────────────────
        yield sse_step_start("data", "📊 读取新数据…")
        headers, rows = self._load_data(data_file)
        if not headers:
            yield sse_error("数据文件为空或无法解析")
            return
        yield sse_step_done("data", f"📊 {len(rows)} 行 × {len(headers)} 列")

        # ── Step 2: 提取 PPT 文本 ──────────────────────────────────
        yield sse_step_start("scan", "📑 扫描 PPT 内容…")
        slide_texts = self._extract_pptx_texts(pptx_file)
        if not slide_texts:
            yield sse_error("PPT 内容为空或无法解析")
            return
        total_shapes = sum(len(s["texts"]) for s in slide_texts)
        yield sse_step_done("scan", f"📑 {len(slide_texts)} 页，{total_shapes} 个文本块")

        # ── Step 3: LLM 匹配映射 ───────────────────────────────────
        yield sse_step_start("mapping", "🤖 AI 数据匹配…")
        pptx_texts_str = json.dumps(slide_texts, ensure_ascii=False, indent=2)
        data_preview = self._build_data_preview(headers, rows)
        inst_text = f"\n用户说明: {instruction}" if instruction else ""

        system = _MAPPING_SYSTEM.format(
            pptx_texts=pptx_texts_str[:6000],
            data_preview=data_preview,
            instruction=inst_text,
        )
        replacements = self._get_replacements(system, model_mode)
        yield sse_step_done("mapping", f"🤖 找到 {len(replacements)} 处需更新")

        if not replacements:
            yield sse_output("markdown", "# 数据刷新结果\n\n未发现需要更新的数据。", "无需更新")
            return

        # ── Step 4: 执行替换 ────────────────────────────────────────
        yield sse_step_start("replace", "📝 更新 PPT 数据…")
        output_path = self.save_output_file(".pptx")
        shutil.copy2(pptx_file, str(output_path))
        applied = self._apply_replacements(str(output_path), replacements)
        yield sse_step_done("replace", f"📝 已更新 {applied} 处")

        # ── 输出 ─────────────────────────────────────────────────────
        change_report = self._build_change_report(replacements, applied)
        yield sse_output(
            "pptx_file",
            {"path": str(output_path), "filename": f"更新_{output_path.name}"},
            f"更新后 PPT（{applied} 处变更）",
        )
        yield sse_output("markdown", change_report, "变更清单")

    # ── 辅助方法 ──────────────────────────────────────────────────────

    def _load_data(self, file_path: str) -> tuple[list[str], list[list[str]]]:
        ext = Path(file_path).suffix.lower()
        try:
            if ext in (".xlsx", ".xls"):
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
                ws = wb.active
                all_rows = []
                for row in ws.iter_rows(values_only=True):
                    all_rows.append([str(c) if c is not None else "" for c in row])
                wb.close()
                if not all_rows:
                    return [], []
                return all_rows[0], all_rows[1:]
            elif ext == ".csv":
                import csv
                with open(file_path, "r", encoding="utf-8-sig", errors="replace") as f:
                    reader = csv.reader(f)
                    all_rows = list(reader)
                if not all_rows:
                    return [], []
                return all_rows[0], all_rows[1:]
        except Exception as e:
            logger.warning("[PptxRefresh] 数据加载失败: %s", e)
        return [], []

    def _extract_pptx_texts(self, path: str) -> list[dict]:
        """提取每页 slide 的文本内容。"""
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for si, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                texts.append(text)
            if texts:
                slides.append({"slide": si, "texts": texts})
        return slides

    def _build_data_preview(self, headers: list[str], rows: list[list[str]]) -> str:
        lines = ["\t".join(headers)]
        for r in rows[:20]:
            lines.append("\t".join(r))
        return "\n".join(lines)

    def _get_replacements(self, system: str, model_mode: str) -> list[dict]:
        prompt = "请进行数据匹配并输出替换列表。"
        try:
            result = self.llm_json(prompt, system=system, model_mode=model_mode)
            if isinstance(result, list):
                return [r for r in result if r.get("old_value") and r.get("new_value")
                        and r["old_value"] != r["new_value"]]
        except Exception as e:
            logger.warning("[PptxRefresh] LLM 匹配失败: %s", e)
        return []

    def _apply_replacements(self, path: str, replacements: list[dict]) -> int:
        """在 PPTX 中执行替换。"""
        from pptx import Presentation
        prs = Presentation(path)
        count = 0

        # 按 slide_idx 分组
        by_slide: dict[int, list[dict]] = {}
        for r in replacements:
            si = r.get("slide_idx", 0)
            by_slide.setdefault(si, []).append(r)

        for si, slide in enumerate(prs.slides):
            reps = by_slide.get(si, [])
            if not reps:
                continue
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            for r in reps:
                                old_val = r["old_value"]
                                new_val = r["new_value"]
                                if old_val in run.text:
                                    run.text = run.text.replace(old_val, new_val)
                                    count += 1
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    for r in reps:
                                        old_val = r["old_value"]
                                        new_val = r["new_value"]
                                        if old_val in run.text:
                                            run.text = run.text.replace(old_val, new_val)
                                            count += 1

        prs.save(path)
        return count

    def _build_change_report(self, replacements: list[dict], applied: int) -> str:
        lines = [
            "# PPT 数据刷新清单\n",
            f"计划更新 {len(replacements)} 处，实际替换 {applied} 处\n",
            "| 页码 | 旧值 | 新值 | 说明 |",
            "|------|------|------|------|",
        ]
        for r in replacements:
            si = r.get("slide_idx", 0) + 1
            old_v = str(r.get("old_value", ""))[:30]
            new_v = str(r.get("new_value", ""))[:30]
            reason = r.get("reason", "")
            lines.append(f"| P{si} | {old_v} | {new_v} | {reason} |")
        return "\n".join(lines)
