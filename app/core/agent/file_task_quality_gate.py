# -*- coding: utf-8 -*-
# Copyright (C) 2024-2026 Koto AI. All rights reserved.
# SPDX-License-Identifier: LicenseRef-Koto-Proprietary
from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any, Dict, List

from app.core.agent.file_task_contract import FileTaskFile, FileTaskRequest
from app.core.agent.file_task_recipes import (
    request_file_types,
    request_target_file_type,
    select_task_recipe,
    semantic_markers,
)
from app.core.agent.file_task_targeting import files_explicitly_mentioned_in_task

_SOURCE_CONTENT_REQUIRED_PATTERNS = (
    re.compile(
        r"\b(?:include|contain|copy|preserve|keep|write|put)\b.{0,80}\b(?:original|source|input)\b.{0,30}\b(?:content|text)\b",
        re.IGNORECASE,
    ),
    re.compile(
        r"(?:包含|包括|保留|写入|放入|复制).{0,40}(?:原文|原始|源文件|原文件|输入).{0,12}(?:内容|文本)",
        re.IGNORECASE,
    ),
)
_TEXT_SOURCE_SUFFIXES = {
    "csv",
    "htm",
    "html",
    "json",
    "log",
    "md",
    "rtf",
    "tsv",
    "txt",
    "xml",
    "yaml",
    "yml",
}
_QUALITY_TEXT_LIMIT = 24_000
_TOP_WORDS = {
    "one": 1,
    "two": 2,
    "three": 3,
    "four": 4,
    "five": 5,
    "six": 6,
    "seven": 7,
    "eight": 8,
    "nine": 9,
    "ten": 10,
}
_ZH_COUNT_WORDS = {
    "一": 1,
    "二": 2,
    "两": 2,
    "三": 3,
    "四": 4,
    "五": 5,
    "六": 6,
    "七": 7,
    "八": 8,
    "九": 9,
    "十": 10,
}
_IGNORED_REQUIRED_PHRASES = {
    "bullet",
    "bullets",
    "markdown",
    "md",
    "word",
    "docx",
    "excel",
    "xlsx",
    "pdf",
}
_DOCX_TABLE_OUTPUT_PATTERN = re.compile(
    r"(?:创建|新增|生成|插入|添加|加入|追加|写入|放入|包含|包括|附上).{0,40}"
    r"(?:真实\s*)?(?:word|docx|文档)?\s*(?:表格|table)"
    r"|(?:真实\s*)?(?:word|docx|文档)?\s*(?:表格|table).{0,40}"
    r"(?:创建|新增|生成|插入|添加|加入|追加|写入|放入|包含|包括|附上)"
    r"|(?:create|add|insert|append|write|include|put).{0,40}"
    r"(?:real\s*)?(?:word|docx|document)?\s*table"
    r"|(?:风险登记表|风险表|risk\s+register\s+table)",
    re.IGNORECASE,
)
_SPREADSHEET_FORMULA_REQUIRED_PATTERN = re.compile(
    r"(?:公式(?:列|单元格|计算)?|保留公式|使用公式|写入公式|"
    r"formula(?:s|\s+column|\s+cells?)?)",
    re.IGNORECASE,
)


def should_attempt_repair(
    check_payload: Dict[str, Any] | None,
    *,
    round_index: int,
    repair_attempts: int,
    max_rounds: int,
    max_repair_attempts: int,
) -> bool:
    if not isinstance(check_payload, dict):
        return False
    if repair_attempts >= max_repair_attempts:
        return False
    if round_index >= max_rounds:
        return False
    if bool(check_payload.get("passed")):
        return False
    status = str(check_payload.get("status") or "").strip().lower()
    return status in {
        "write_not_performed",
        "quality_gate_failed",
    }


def change_operations(file_changes: List[Dict[str, Any]]) -> set[str]:
    return {
        str(change.get("operation") or "").strip()
        for change in file_changes
        if str(change.get("operation") or "").strip()
    }


def change_sum_int(file_changes: List[Dict[str, Any]], key: str) -> int:
    total = 0
    for change in file_changes:
        try:
            total += int(change.get(key) or 0)
        except Exception:
            continue
    return total


def target_or_request_type(
    request: FileTaskRequest,
    file_changes: List[Dict[str, Any]],
) -> str:
    target_type = Path(str(request.target_path or "")).suffix.lstrip(".").lower()
    if target_type:
        return target_type
    for change in file_changes:
        candidate = (
            str(
                change.get("file_type")
                or Path(str(change.get("path") or "")).suffix.lstrip(".")
            )
            .lower()
            .strip()
        )
        if candidate:
            return candidate
    for file_info in request.files or []:
        if file_info.target:
            candidate = (
                str(
                    file_info.type
                    or Path(str(file_info.path or file_info.name)).suffix.lstrip(".")
                )
                .lower()
                .strip()
            )
            if candidate:
                return candidate
    return ""


def quality_gate_result(
    *,
    criterion: str,
    passed: bool,
    detail: str,
    priority: str = "high",
) -> Dict[str, Any]:
    return {
        "criterion": criterion,
        "passed": bool(passed),
        "detail": detail,
        "priority": priority,
    }


def evaluate_task_quality_gate(
    request: FileTaskRequest,
    file_changes: List[Dict[str, Any]],
    *,
    write_intent: bool,
    output_mode: str,
) -> Dict[str, Any]:
    del output_mode
    if not write_intent:
        return {"passed": True, "criteria_results": [], "remaining": []}

    operations = change_operations(file_changes)
    target_type = target_or_request_type(request, file_changes)
    paragraphs_written = change_sum_int(file_changes, "paragraphs_written")
    images_inserted = change_sum_int(file_changes, "images_inserted")
    rows_written = change_sum_int(file_changes, "rows_written")
    cells_written = change_sum_int(file_changes, "cells_written")
    slides_updated = change_sum_int(file_changes, "slides_updated")
    slides_added = change_sum_int(file_changes, "slides_added")
    slides_designed = change_sum_int(file_changes, "slides_designed")
    text_shapes_styled = change_sum_int(file_changes, "text_shapes_styled")
    annotations_added = change_sum_int(file_changes, "annotations_added")
    differences_detected = change_sum_int(file_changes, "differences_detected")
    comments_removed = change_sum_int(file_changes, "comments_removed")
    revisions_accepted = change_sum_int(file_changes, "revisions_accepted")
    paragraphs_rewritten = change_sum_int(file_changes, "paragraphs_rewritten")
    placeholders_replaced = change_sum_int(file_changes, "placeholders_replaced")
    replacements_made = change_sum_int(file_changes, "replacements_made")
    task_text = str(request.task or "")
    local_docx_replace_request = target_type in {
        "docx",
        "doc",
    } and _looks_like_local_docx_replace_request(task_text)
    local_docx_edit_request = (
        target_type in {"docx", "doc"}
        and (
            _looks_like_local_docx_edit_request(task_text) or local_docx_replace_request
        )
        and not _looks_like_table_request(task_text)
    )
    actual_docx_paragraphs = (
        _target_docx_quality_paragraph_count(request, file_changes)
        if target_type in {"docx", "doc"} and file_changes
        else 0
    )

    criteria: List[Dict[str, Any]] = []
    metric_values = {
        "paragraphs_written": paragraphs_written,
        "images_inserted": images_inserted,
        "rows_written": rows_written,
        "slides_updated": slides_updated,
        "slides_added": slides_added,
        "slides_designed": slides_designed,
        "text_shapes_styled": text_shapes_styled,
        "annotations_added": annotations_added,
        "differences_detected": differences_detected,
        "comments_removed": comments_removed,
        "revisions_accepted": revisions_accepted,
        "paragraphs_rewritten": paragraphs_rewritten,
        "cells_written": cells_written,
        "placeholders_replaced": placeholders_replaced,
        "replacements_made": replacements_made,
    }
    recipe_match = select_task_recipe(
        request, request.files or [], write_intent=write_intent
    )
    explicit_source_content_gate = _source_content_in_output_gate(request, file_changes)
    explicit_top_table_gate = _top_table_requirement_gate(request, file_changes)
    explicit_section_gates = _explicit_docx_section_gates(request, file_changes)
    explicit_phrase_gate = _explicit_required_phrase_gate(request, file_changes)
    explicit_bullet_count_gate = _explicit_bullet_count_gate(request, file_changes)
    explicit_slide_count_gate = _explicit_slide_count_gate(
        request, file_changes, target_type
    )
    explicit_pptx_native_visual_gate = _pptx_native_visual_requirement_gate(
        request,
        file_changes,
        target_type,
    )
    explicit_spreadsheet_chart_gate = _spreadsheet_chart_requirement_gate(
        request,
        file_changes,
        target_type,
    )
    explicit_spreadsheet_formula_gate = _spreadsheet_formula_requirement_gate(
        request,
        file_changes,
        target_type,
    )
    explicit_docx_table_gate = _docx_table_requirement_gate(
        request,
        file_changes,
        target_type,
        operations,
        rows_written,
    )
    seen_recipe_criteria: set[str] = set()
    if recipe_match and not local_docx_edit_request:
        for gate in recipe_match.recipe.quality_gates:
            criterion = str(gate.get("criterion") or "").strip()
            if not criterion or criterion in seen_recipe_criteria:
                continue
            seen_recipe_criteria.add(criterion)
            operation = str(gate.get("operation") or "").strip()
            any_operation = {
                str(item).strip()
                for item in gate.get("any_operation") or []
                if str(item).strip()
            }
            metric_name = str(gate.get("metric") or "").strip()
            actual = int(metric_values.get(metric_name, 0) or 0)
            if metric_name == "paragraphs_written" and actual_docx_paragraphs:
                actual = max(actual, actual_docx_paragraphs)
            minimum = int(gate.get("minimum") or 0)
            if any_operation:
                operation_present = bool(operations.intersection(any_operation))
                if (
                    not operation_present
                    and "write_docx_content" in any_operation
                    and "run_python_code" in operations
                    and actual_docx_paragraphs
                ):
                    operation_present = True
                passed = operation_present and actual >= minimum
                detail = str(gate.get("detail") or "").format(
                    operations=", ".join(sorted(operations)) or "无", actual=actual
                )
            else:
                operation_present = not operation or operation in operations
                if (
                    not operation_present
                    and operation == "write_docx_content"
                    and "run_python_code" in operations
                    and actual_docx_paragraphs
                ):
                    operation_present = True
                passed = operation_present and actual >= minimum
                detail = str(gate.get("detail") or "").format(
                    operations=", ".join(sorted(operations)) or "无",
                    actual=actual,
                    minimum=minimum,
                )
            criteria.append(
                quality_gate_result(
                    criterion=criterion,
                    passed=passed,
                    detail=detail or criterion,
                    priority=str(gate.get("priority") or "high"),
                )
            )

    if _looks_like_financial_xlsx_docx_chart_report_task(request, request.files or []):
        table_changes = [
            change
            for change in file_changes
            if isinstance(change, dict)
            and str(change.get("operation") or "").strip()
            == "insert_excel_as_docx_table"
        ]
        readable_table = False
        table_details: List[str] = []
        for change in table_changes:
            rows = int(change.get("rows_written") or 0)
            columns = int(change.get("columns_written") or 0)
            table_details.append(f"{rows} 行 × {columns or '未知'} 列")
            if 1 <= rows <= 24 and (columns == 0 or 2 <= columns <= 8):
                readable_table = True
        criteria.append(
            quality_gate_result(
                criterion="financial_report_has_readable_key_data_table",
                passed=readable_table,
                detail=(
                    "财务报告中的关键数据表必须为 1-24 行、2-8 列的可读摘要表；"
                    f"当前表格：{'；'.join(table_details) or '无'}。"
                ),
                priority="critical",
            )
        )
        duplicate_report_gate = _financial_report_duplicate_content_gate(
            request, file_changes
        )
        if duplicate_report_gate:
            criteria.append(duplicate_report_gate)

    if criteria:
        table_narrative_gate = _table_narrative_requirement_gate(
            task_text,
            target_type,
            operations,
            paragraphs_written,
        )
        if table_narrative_gate:
            criteria.append(table_narrative_gate)
        if explicit_source_content_gate:
            criteria.append(explicit_source_content_gate)
        if explicit_top_table_gate:
            criteria.append(explicit_top_table_gate)
        criteria.extend(explicit_section_gates)
        if explicit_phrase_gate:
            criteria.append(explicit_phrase_gate)
        if explicit_bullet_count_gate:
            criteria.append(explicit_bullet_count_gate)
        if explicit_slide_count_gate:
            criteria.append(explicit_slide_count_gate)
        if explicit_pptx_native_visual_gate:
            criteria.append(explicit_pptx_native_visual_gate)
        if explicit_spreadsheet_chart_gate:
            criteria.append(explicit_spreadsheet_chart_gate)
        if explicit_spreadsheet_formula_gate:
            criteria.append(explicit_spreadsheet_formula_gate)
        if explicit_docx_table_gate:
            criteria.append(explicit_docx_table_gate)
        failed = [item for item in criteria if not item.get("passed")]
        return {
            "passed": not failed,
            "criteria_results": criteria,
            "remaining": [
                str(item.get("detail") or item.get("criterion")) for item in failed
            ],
        }

    if target_type in {"docx", "doc"} and _looks_like_chart_request(task_text):
        criteria.append(
            quality_gate_result(
                criterion="docx_chart_request_has_image",
                passed="insert_image_into_docx" in operations and images_inserted >= 1,
                detail=f"用户要求图表/图片进入 Word；当前图片写入数：{images_inserted}。",
                priority="critical",
            )
        )

    if local_docx_edit_request:
        if local_docx_replace_request:
            targeted_replace_ops = {"run_python_code", "rewrite_docx_paragraph_window"}
            criteria.append(
                quality_gate_result(
                    criterion="docx_local_replace_has_targeted_write",
                    passed=bool(file_changes)
                    and bool(operations.intersection(targeted_replace_ops)),
                    detail=(
                        "DOCX 定位替换应在保留原结构的前提下完成一次目标写回；"
                        f"当前操作：{', '.join(sorted(operations)) or '无'}，"
                        f"目标变更数：{len(file_changes)}。"
                    ),
                    priority="high",
                )
            )
        else:
            inserted_text = any(
                str(change.get("inserted_text") or "").strip()
                for change in file_changes
                if str(change.get("operation") or "").strip() == "insert_docx_paragraph"
            )
            criteria.append(
                quality_gate_result(
                    criterion="docx_local_edit_has_paragraph_insert",
                    passed="insert_docx_paragraph" in operations
                    and (paragraphs_written >= 1 or inserted_text),
                    detail=(
                        "局部 Word 编辑任务应使用保留原结构的段落插入工具；"
                        f"当前操作：{', '.join(sorted(operations)) or '无'}，"
                        f"段落插入数：{paragraphs_written}。"
                    ),
                    priority="high",
                )
            )
        if explicit_source_content_gate:
            criteria.append(explicit_source_content_gate)
        criteria.extend(explicit_section_gates)
        if explicit_phrase_gate:
            criteria.append(explicit_phrase_gate)
        if explicit_bullet_count_gate:
            criteria.append(explicit_bullet_count_gate)
        if explicit_slide_count_gate:
            criteria.append(explicit_slide_count_gate)
        failed = [item for item in criteria if not item.get("passed")]
        return {
            "passed": not failed,
            "criteria_results": criteria,
            "remaining": [
                str(item.get("detail") or item.get("criterion")) for item in failed
            ],
        }

    if target_type in {"docx", "doc"} and _looks_like_docx_report_request(
        request, request.files or []
    ):
        narrative_minimum = 2 if _looks_like_table_request(task_text) else 3
        criteria.append(
            quality_gate_result(
                criterion="docx_report_has_narrative",
                passed=(
                    "write_docx_content" in operations
                    and paragraphs_written >= narrative_minimum
                )
                or paragraphs_written >= narrative_minimum
                or (bool(file_changes) and actual_docx_paragraphs >= narrative_minimum),
                detail=(
                    "DOCX 报告/分析任务应写入可读文本结构；"
                    f"当前段落写入数：{paragraphs_written}，"
                    f"目标文档实际段落数：{actual_docx_paragraphs}，"
                    f"最低要求：{narrative_minimum}。"
                ),
                priority="high",
            )
        )

    if explicit_docx_table_gate:
        criteria.append(explicit_docx_table_gate)
        table_narrative_gate = _table_narrative_requirement_gate(
            task_text,
            target_type,
            operations,
            paragraphs_written,
        )
        if table_narrative_gate:
            criteria.append(table_narrative_gate)

    if target_type in {"docx", "doc"} and operations.intersection(
        {"compare_docx_and_annotate", "write_docx_comments"}
    ):
        criteria.append(
            quality_gate_result(
                criterion="docx_compare_has_difference_annotations",
                passed=annotations_added > 0,
                detail=f"DOCX 对比标注任务必须写入真实差异批注；当前批注数：{annotations_added}。",
                priority="critical",
            )
        )

    if _looks_like_ppt_slide_write_request(request, request.files or []):
        criteria.append(
            quality_gate_result(
                criterion="ppt_request_has_slide_write",
                passed=bool(
                    operations.intersection(
                        {
                            "add_pptx_slides",
                            "write_pptx_slides",
                            "design_pptx_theme_layout",
                        }
                    )
                ),
                detail=f"PPT 任务应产生幻灯片写入/更新操作；当前操作：{', '.join(sorted(operations)) or '无'}。",
                priority="critical",
            )
        )

    if target_type in {"docx", "doc"} and not criteria:
        docx_write_ops = {
            "write_docx_content",
            "insert_excel_as_docx_table",
            "insert_image_into_docx",
            "annotate_file",
            "compare_docx_and_annotate",
            "clear_docx_review_marks",
            "insert_docx_paragraph",
            "fill_docx_template",
            "rewrite_docx_paragraph_window",
        }
        docx_metric_total = (
            paragraphs_written
            + images_inserted
            + rows_written
            + annotations_added
            + differences_detected
            + comments_removed
            + revisions_accepted
            + paragraphs_rewritten
            + placeholders_replaced
        )
        run_python_docx_writeback = (
            "run_python_code" in operations
            and bool(file_changes)
            and (
                _looks_like_polish_request(task_text)
                or _looks_like_translation_request(task_text)
            )
        )
        criteria.append(
            quality_gate_result(
                criterion="generic_docx_has_native_write",
                passed=bool(operations.intersection(docx_write_ops))
                and docx_metric_total > 0
                or run_python_docx_writeback,
                detail=(
                    "DOCX 写入任务必须产生可核验的 Word 原生写入指标；"
                    f"当前操作：{', '.join(sorted(operations)) or '无'}，"
                    f"段落/图片/表格/批注/修订指标合计：{docx_metric_total}。"
                ),
                priority="high",
            )
        )

    if target_type in {"pptx", "ppt"} and not criteria:
        pptx_write_ops = {
            "add_pptx_slides",
            "write_pptx_slides",
            "design_pptx_theme_layout",
        }
        pptx_metric_total = (
            slides_updated + slides_added + slides_designed + text_shapes_styled
        )
        criteria.append(
            quality_gate_result(
                criterion="generic_pptx_has_native_write",
                passed=bool(operations.intersection(pptx_write_ops))
                and pptx_metric_total > 0,
                detail=(
                    "PPTX 写入任务必须产生可核验的幻灯片写入、更新或设计指标；"
                    f"当前操作：{', '.join(sorted(operations)) or '无'}，"
                    f"幻灯片/文本样式指标合计：{pptx_metric_total}。"
                ),
                priority="high",
            )
        )

    if target_type in {"xlsx", "xlsm"} and not criteria:
        spreadsheet_metric_total = rows_written + cells_written
        spreadsheet_write_ops = {"write_sheet_data", "run_python_code"}
        criteria.append(
            quality_gate_result(
                criterion="generic_spreadsheet_has_native_write",
                passed=bool(operations.intersection(spreadsheet_write_ops))
                and spreadsheet_metric_total > 0,
                detail=(
                    "表格写入任务必须产生可核验的单元格/行写入指标；"
                    f"当前操作：{', '.join(sorted(operations)) or '无'}，"
                    f"行/单元格指标合计：{spreadsheet_metric_total}。"
                ),
                priority="high",
            )
        )

    if target_type == "csv" and not criteria:
        csv_metric_total = rows_written + cells_written
        csv_write_ops = {
            "create_file",
            "extract_to_file",
            "run_python_code",
            "write_sheet_data",
        }
        csv_output_text = _target_text_for_quality_gate(request, file_changes)
        csv_has_content = bool(str(csv_output_text or "").strip())
        criteria.append(
            quality_gate_result(
                criterion="generic_csv_has_written_content",
                passed=(
                    bool(operations.intersection(csv_write_ops))
                    and (csv_metric_total > 0 or csv_has_content)
                ),
                detail=(
                    "CSV 写入任务必须产生非空 CSV 内容或可核验的行/单元格指标；"
                    f"当前操作：{', '.join(sorted(operations)) or '无'}，"
                    f"行/单元格指标合计：{csv_metric_total}，"
                    f"CSV 内容：{'已生成' if csv_has_content else '未检测到'}。"
                ),
                priority="high",
            )
        )

    if explicit_source_content_gate:
        criteria.append(explicit_source_content_gate)
    if explicit_top_table_gate:
        criteria.append(explicit_top_table_gate)
    criteria.extend(explicit_section_gates)
    if explicit_phrase_gate:
        criteria.append(explicit_phrase_gate)
    if explicit_bullet_count_gate:
        criteria.append(explicit_bullet_count_gate)
    if explicit_slide_count_gate:
        criteria.append(explicit_slide_count_gate)
    if explicit_pptx_native_visual_gate:
        criteria.append(explicit_pptx_native_visual_gate)
    if explicit_spreadsheet_chart_gate:
        criteria.append(explicit_spreadsheet_chart_gate)
    if explicit_spreadsheet_formula_gate:
        criteria.append(explicit_spreadsheet_formula_gate)

    failed = [item for item in criteria if not item.get("passed")]
    return {
        "passed": not failed,
        "criteria_results": criteria,
        "remaining": [
            str(item.get("detail") or item.get("criterion")) for item in failed
        ],
    }


def _source_content_in_output_gate(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> Dict[str, Any] | None:
    if not _task_requires_source_content_in_output(request.task):
        return None
    candidates = _source_content_candidates(request, file_changes)
    if not candidates:
        return None
    output_text = _normalize_quality_text(
        _output_text_for_quality_gate(request, file_changes)
    )
    missing: List[str] = []
    for candidate in candidates:
        anchors = _source_content_anchors(candidate["text"])
        if not anchors:
            continue
        if not output_text or not all(anchor in output_text for anchor in anchors):
            missing.append(candidate["label"])
    if not missing:
        return quality_gate_result(
            criterion="source_content_included",
            passed=True,
            detail="用户要求产物包含源文件原文，已在输出内容中找到源内容片段。",
            priority="critical",
        )
    return quality_gate_result(
        criterion="source_content_included",
        passed=False,
        detail=(
            "用户要求产物包含源文件原文，但输出内容缺少这些源内容："
            + "、".join(missing[:3])
            + "。"
        ),
        priority="critical",
    )


def _task_requires_source_content_in_output(task: str) -> bool:
    task_text = str(task or "")
    return any(
        pattern.search(task_text) for pattern in _SOURCE_CONTENT_REQUIRED_PATTERNS
    )


def _source_content_candidates(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> List[Dict[str, str]]:
    target_paths = _quality_target_paths(request, file_changes)
    candidates: List[Dict[str, str]] = []
    if str(request.selection or "").strip():
        _append_source_candidate(candidates, "选区内容", str(request.selection or ""))
    if request.current_file:
        _append_file_source_candidate(candidates, request.current_file, target_paths)
    for file_info in request.files or []:
        _append_file_source_candidate(candidates, file_info, target_paths)
    return candidates


def _append_file_source_candidate(
    candidates: List[Dict[str, str]],
    file_info: FileTaskFile,
    target_paths: set[str],
) -> None:
    label = str(file_info.name or file_info.path or "源文件").strip()
    if _quality_file_is_target(file_info, target_paths):
        return
    content = str(file_info.content or "")
    if not content:
        content = _read_quality_text_from_path(file_info.path or file_info.name)
    _append_source_candidate(candidates, label, content)


def _append_source_candidate(
    candidates: List[Dict[str, str]], label: str, content: str
) -> None:
    clean = str(content or "").replace("\x00", "").strip()
    if len(_normalize_quality_text(clean)) < 4:
        return
    normalized_label = str(label or "源内容").strip() or "源内容"
    seen = {
        (_normalize_quality_text(item["label"]), _normalize_quality_text(item["text"]))
        for item in candidates
    }
    marker = (_normalize_quality_text(normalized_label), _normalize_quality_text(clean))
    if marker in seen:
        return
    candidates.append({"label": normalized_label, "text": clean[:_QUALITY_TEXT_LIMIT]})


def _source_content_anchors(text: str) -> List[str]:
    normalized = _normalize_quality_text(text)
    if len(normalized) < 4:
        return []
    if len(normalized) <= 600:
        return [normalized]
    starts = [0, max(0, len(normalized) // 2 - 80), max(0, len(normalized) - 160)]
    anchors: List[str] = []
    for start in starts:
        anchor = normalized[start : start + 160].strip()
        if len(anchor) >= 20 and anchor not in anchors:
            anchors.append(anchor)
    return anchors


def _output_text_for_quality_gate(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> str:
    parts: List[str] = []
    target_text = _target_text_for_quality_gate(request, file_changes)
    if target_text:
        parts.append(target_text)
    for change in file_changes:
        if not isinstance(change, dict):
            continue
        for key in (
            "preview",
            "content",
            "text",
            "after",
            "summary",
            "result_preview",
        ):
            value = change.get(key)
            if value:
                parts.append(str(value))
        diff = change.get("diff")
        if diff:
            parts.append(_quality_stringify(diff))
    return "\n".join(part for part in parts if part)


def _target_text_for_quality_gate(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> str:
    parts: List[str] = []
    for path_text in _quality_target_paths(request, file_changes):
        text = _read_quality_text_from_path(path_text)
        if text:
            parts.append(text)
    return "\n".join(parts)


def _read_quality_text_from_path(path_text: str) -> str:
    path_value = str(path_text or "").strip()
    if not path_value:
        return ""
    try:
        path = _resolve_quality_path(path_value)
        if not path:
            return ""
        suffix = path.suffix.lower().lstrip(".")
        if suffix in {"docx", "doc"}:
            return _read_docx_quality_text(path)
        if suffix in {"pptx", "ppt"}:
            return _read_pptx_quality_text(path)
        if suffix not in _TEXT_SOURCE_SUFFIXES:
            return ""
        return path.read_text(encoding="utf-8", errors="replace")[:_QUALITY_TEXT_LIMIT]
    except Exception:
        return ""


def _read_docx_quality_text(path: Path) -> str:
    try:
        from docx import Document

        document = Document(str(path))
    except Exception:
        return ""
    parts: List[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append("\t".join(cells))
    return "\n".join(parts)[:_QUALITY_TEXT_LIMIT]


def _read_pptx_quality_text(path: Path) -> str:
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
    except Exception:
        return ""
    parts: List[str] = []
    for index, slide in enumerate(presentation.slides, 1):
        slide_parts: List[str] = []
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                slide_parts.append(text)
        if slide_parts:
            parts.append(f"Slide {index}\n" + "\n".join(slide_parts))
    return "\n".join(parts)[:_QUALITY_TEXT_LIMIT]


def _explicit_required_phrase_gate(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> Dict[str, Any] | None:
    phrases = _explicit_required_phrases(str(request.task or ""))
    if not phrases:
        return None
    output_text = _output_text_for_quality_gate(request, file_changes)
    normalized_output = _normalize_quality_phrase(output_text)
    missing = [
        phrase
        for phrase in phrases
        if _normalize_quality_phrase(phrase) not in normalized_output
    ]
    return quality_gate_result(
        criterion="explicit_required_phrases_present",
        passed=not missing,
        detail=(
            "用户显式要求输出包含这些短语；"
            + (
                "目标产物中均已找到。"
                if not missing
                else "目标产物缺少：" + "、".join(missing[:5])
            )
        ),
        priority="critical",
    )


def _explicit_slide_count_gate(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]], target_type: str
) -> Dict[str, Any] | None:
    if target_type not in {"pptx", "ppt"}:
        return None
    expected = _explicit_slide_count(str(request.task or ""))
    if expected is None:
        return None
    actual = _target_presentation_slide_count(request, file_changes)
    return quality_gate_result(
        criterion="explicit_pptx_slide_count",
        passed=actual == expected,
        detail=(
            f"用户显式要求幻灯片页数为 {expected}；"
            f"目标演示文稿可核验页数：{actual}。"
        ),
        priority="critical",
    )


def _explicit_slide_count(task: str) -> int | None:
    text = str(task or "")
    count_token = r"\d+|one|two|three|four|five|six|seven|eight|nine|ten|一|二|两|三|四|五|六|七|八|九|十"
    patterns = (
        rf"(?:正好|恰好|exactly|必须正好)\s*(?P<count>{count_token})\s*(?:页|张|个)?\s*(?:幻灯片|slides?)",
        rf"(?P<count>{count_token})\s*(?:页|张|个)\s*(?:幻灯片|slides?)",
        rf"(?P<count>{count_token})\s*slides?",
    )
    for pattern in patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if not match:
            continue
        value = _count_word_to_int(match.group("count"))
        if value is not None:
            return value
    return None


def _target_presentation_slide_count(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> int:
    counts = [
        _presentation_slide_count(path_text)
        for path_text in _quality_target_paths(request, file_changes)
        if str(path_text or "").lower().endswith((".pptx", ".ppt"))
    ]
    return max(counts, default=0)


def _pptx_native_visual_requirement_gate(
    request: FileTaskRequest,
    file_changes: List[Dict[str, Any]],
    target_type: str,
) -> Dict[str, Any] | None:
    if target_type not in {"pptx", "ppt"}:
        return None
    task = str(request.task or "")
    native_requested = bool(re.search(r"(?:原生|native)", task, re.IGNORECASE))
    chart_requested = _looks_like_chart_request(task)
    timeline_requested = bool(
        re.search(r"(?:时间线|里程碑(?:图|时间线)?|timeline)", task, re.IGNORECASE)
    )
    if not native_requested or not (chart_requested or timeline_requested):
        return None

    chart_count = 0
    timeline_visual_count = 0
    for path_text in _quality_target_paths(request, file_changes):
        counts = _presentation_native_visual_counts(path_text)
        chart_count = max(chart_count, counts["charts"])
        timeline_visual_count = max(
            timeline_visual_count, counts["timeline_visual_shapes"]
        )

    chart_passed = chart_requested and chart_count >= 1
    timeline_passed = timeline_requested and timeline_visual_count >= 3
    return quality_gate_result(
        criterion="pptx_native_chart_or_timeline_present",
        passed=chart_passed or timeline_passed,
        detail=(
            "用户显式要求原生图表或原生时间线；"
            f"目标演示文稿中的原生图表数：{chart_count}，"
            f"时间线/里程碑页的原生视觉形状数：{timeline_visual_count}。"
        ),
        priority="critical",
    )


def _presentation_native_visual_counts(path_text: str) -> Dict[str, int]:
    resolved = _resolve_quality_path(path_text)
    if not resolved or resolved.suffix.lower() not in {".pptx", ".ppt"}:
        return {"charts": 0, "timeline_visual_shapes": 0}
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        presentation = Presentation(str(resolved))
    except Exception:
        return {"charts": 0, "timeline_visual_shapes": 0}

    charts = 0
    timeline_visual_shapes = 0
    visual_types = {
        MSO_SHAPE_TYPE.AUTO_SHAPE,
        MSO_SHAPE_TYPE.GROUP,
        MSO_SHAPE_TYPE.LINE,
    }
    for slide in presentation.slides:
        slide_text = "\n".join(
            str(getattr(shape, "text", "") or "") for shape in slide.shapes
        )
        charts += sum(
            1 for shape in slide.shapes if bool(getattr(shape, "has_chart", False))
        )
        if re.search(r"(?:时间线|里程碑|timeline)", slide_text, re.IGNORECASE):
            timeline_visual_shapes = max(
                timeline_visual_shapes,
                sum(
                    1
                    for shape in slide.shapes
                    if not bool(getattr(shape, "is_placeholder", False))
                    and getattr(shape, "shape_type", None) in visual_types
                ),
            )
    return {
        "charts": charts,
        "timeline_visual_shapes": timeline_visual_shapes,
    }


def _presentation_slide_count(path_text: str) -> int:
    resolved = _resolve_quality_path(path_text)
    if not resolved:
        return 0
    try:
        from pptx import Presentation

        presentation = Presentation(str(resolved))
        return len(presentation.slides)
    except Exception:
        return 0


def _explicit_bullet_count_gate(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> Dict[str, Any] | None:
    expected = _explicit_bullet_count(str(request.task or ""))
    if expected is None:
        return None
    output_text = _target_text_for_quality_gate(
        request, file_changes
    ) or _output_text_for_quality_gate(request, file_changes)
    actual = _quality_bullet_count(output_text)
    return quality_gate_result(
        criterion="explicit_bullet_count",
        passed=actual == expected,
        detail=(
            f"用户显式要求 bullet 数量为 {expected}；"
            f"目标产物中可核验 bullet 数量：{actual}。"
        ),
        priority="critical",
    )


def _explicit_bullet_count(task: str) -> int | None:
    text = str(task or "")
    count_token = r"\d+|one|two|three|four|five|six|seven|eight|nine|ten|一|二|两|三|四|五|六|七|八|九|十"
    patterns = (
        rf"(?:正好|恰好|exactly)\s*(?P<count>{count_token})\s*(?:条|个|项)?\s*(?:bullet|bullets|要点|项目)",
        rf"(?P<count>{count_token})\s*(?:条|个|项)?\s*(?:bullet|bullets)",
        rf"(?:bullet|bullets|要点|项目).{{0,12}}(?:正好|恰好|exactly)\s*(?P<count>{count_token})",
    )
    for pattern in patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if not match:
            continue
        value = _count_word_to_int(match.group("count"))
        if value is not None:
            return value
    return None


def _count_word_to_int(value: str) -> int | None:
    token = str(value or "").strip().casefold()
    if not token:
        return None
    if token.isdigit():
        number = int(token)
        return number if 0 <= number <= 50 else None
    if token in _TOP_WORDS:
        return _TOP_WORDS[token]
    return _ZH_COUNT_WORDS.get(token)


def _quality_bullet_count(text: str) -> int:
    return len(
        re.findall(
            r"(?m)^\s*(?:[-*+]|(?:\d+[.)])|(?:[一二三四五六七八九十]+[、.]))\s+\S",
            str(text or ""),
        )
    )


def _explicit_required_phrases(task: str) -> List[str]:
    text = str(task or "")
    phrases: List[str] = []
    patterns = (
        (r"(?:标题分别为|标题依次为)\s*([^。；;\n]+)", True, False),
        (r"(?:包含标题|标题为|标题是)\s*([^。；;\n，,]+)", False, False),
        (r"(?:包含客户|保留客户|客户(?:为|是)?)\s*([^。；;\n，,]+)", False, False),
        (r"(?:包含日期|日期(?:为|是)?)\s*([^。；;\n，]+)", False, False),
        (
            r"(?:最后包含一句|包含一句|包含短语|包含字样|包括字样)\s*([^。；;\n，,]+)",
            False,
            False,
        ),
        (r"(?:每条|每项).{0,16}(?:包含|包括)\s*([^。；;\n]+)", True, False),
        (
            r"([A-Za-z][A-Za-z0-9 /:&.+-]{2,80})\s*(?:小节|章节|section)",
            False,
            True,
        ),
    )
    for pattern, split_segment, section_name in patterns:
        for match in re.finditer(pattern, text, re.IGNORECASE):
            segment = match.group(1)
            extracted = (
                _english_phrases_from_segment(segment)
                if split_segment
                else [_clean_required_phrase(segment)]
            )
            for phrase in extracted:
                keep_phrase = (
                    _should_keep_required_section_phrase(phrase)
                    if section_name
                    else _should_keep_required_phrase(phrase)
                )
                if keep_phrase and phrase.casefold() not in {
                    item.casefold() for item in phrases
                }:
                    phrases.append(phrase)
    return phrases


def _english_phrases_from_segment(segment: str) -> List[str]:
    pieces = re.split(r"\s*(?:和|及|与|、|，|;|；)\s*", str(segment or ""))
    phrases: List[str] = []
    for piece in pieces:
        for match in re.finditer(
            r"[A-Za-z][A-Za-z0-9]*(?:[ /:&.+-]+[A-Za-z0-9][A-Za-z0-9]*){0,8}",
            piece,
        ):
            phrase = _clean_required_phrase(match.group(0))
            if _should_keep_required_phrase(phrase) and phrase.casefold() not in {
                item.casefold() for item in phrases
            }:
                phrases.append(phrase)
    return phrases


def _clean_required_phrase(text: str) -> str:
    phrase = re.sub(r"\s+", " ", str(text or "")).strip(" .:-：")
    phrase = re.sub(r"\s*(?:小节|章节|section)\s*$", "", phrase, flags=re.IGNORECASE)
    return phrase.strip(" .:-：")


def _should_keep_required_phrase(phrase: str) -> bool:
    normalized = _normalize_quality_phrase(phrase)
    return len(phrase) >= 3 and normalized not in _IGNORED_REQUIRED_PHRASES


def _should_keep_required_section_phrase(phrase: str) -> bool:
    normalized = _normalize_quality_phrase(phrase)
    if not _should_keep_required_phrase(phrase):
        return False
    if re.search(
        r"\b(?:create|generate|write|include|includes|must|should|the|this|that|report|document|docx|file)\b",
        normalized,
    ):
        return False
    generic_sections = {
        "risk",
        "risks",
        "action",
        "actions",
        "next action",
        "next actions",
        "summary",
        "summaries",
        "section",
        "sections",
    }
    return normalized not in generic_sections


def _normalize_quality_phrase(text: str) -> str:
    return re.sub(r"\s+", " ", str(text or "")).strip().casefold()


def _spreadsheet_chart_requirement_gate(
    request: FileTaskRequest,
    file_changes: List[Dict[str, Any]],
    target_type: str,
) -> Dict[str, Any] | None:
    if target_type not in {"xlsx", "xlsm"}:
        return None
    if not _looks_like_chart_request(str(request.task or "")):
        return None
    chart_count = _target_workbook_chart_count(request, file_changes)
    return quality_gate_result(
        criterion="spreadsheet_chart_request_has_workbook_chart",
        passed=chart_count > 0,
        detail=(
            "用户要求 Excel 产物包含真实图表；"
            f"目标工作簿中可核验图表数量：{chart_count}。"
        ),
        priority="critical",
    )


def _spreadsheet_formula_requirement_gate(
    request: FileTaskRequest,
    file_changes: List[Dict[str, Any]],
    target_type: str,
) -> Dict[str, Any] | None:
    if target_type not in {"xlsx", "xlsm"}:
        return None
    if not _SPREADSHEET_FORMULA_REQUIRED_PATTERN.search(str(request.task or "")):
        return None
    formula_count = _target_workbook_formula_count(request, file_changes)
    return quality_gate_result(
        criterion="spreadsheet_formula_request_has_workbook_formula",
        passed=formula_count > 0,
        detail=(
            "用户明确要求 Excel 产物保留公式；"
            f"目标工作簿中可核验公式单元格数量：{formula_count}。"
        ),
        priority="critical",
    )


def _target_workbook_chart_count(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> int:
    paths = list(_quality_target_paths(request, file_changes))
    for file_info in request.files or []:
        file_path = str(file_info.path or "").strip()
        file_name = str(file_info.name or file_path).strip()
        if file_info.target:
            paths.extend([file_path, file_name])
    counts = [_workbook_chart_count(path_text) for path_text in paths if path_text]
    return max(counts, default=0)


def _target_workbook_formula_count(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> int:
    paths = list(_quality_target_paths(request, file_changes))
    for file_info in request.files or []:
        file_path = str(file_info.path or "").strip()
        file_name = str(file_info.name or file_path).strip()
        if file_info.target:
            paths.extend([file_path, file_name])
    counts = [_workbook_formula_count(path_text) for path_text in paths if path_text]
    return max(counts, default=0)


def _workbook_formula_count(path_text: str) -> int:
    resolved = _resolve_quality_path(path_text)
    if not resolved or resolved.suffix.lower() not in {".xlsx", ".xlsm"}:
        return 0
    try:
        import openpyxl

        workbook = openpyxl.load_workbook(
            str(resolved), data_only=False, read_only=True
        )
    except Exception:
        return 0
    try:
        return sum(
            1
            for sheet in workbook.worksheets
            for row in sheet.iter_rows()
            for cell in row
            if getattr(cell, "data_type", "") == "f"
            or (
                isinstance(cell.value, str) and str(cell.value).lstrip().startswith("=")
            )
        )
    finally:
        workbook.close()


def _workbook_chart_count(path_text: str) -> int:
    resolved = _resolve_quality_path(path_text)
    if not resolved or resolved.suffix.lower() not in {".xlsx", ".xlsm"}:
        return 0
    try:
        import openpyxl

        workbook = openpyxl.load_workbook(str(resolved), data_only=False)
    except Exception:
        return 0
    try:
        return sum(
            len(getattr(sheet, "_charts", []) or []) for sheet in workbook.worksheets
        )
    finally:
        workbook.close()


def _resolve_quality_path(path_text: str) -> Path | None:
    raw = str(path_text or "").strip()
    if not raw:
        return None
    candidates = [Path(raw)]
    if not Path(raw).is_absolute():
        project_root = Path(__file__).resolve().parents[3]
        candidates.extend(
            [
                Path.cwd() / raw,
                project_root / raw,
                project_root / "workspace" / raw,
            ]
        )
    for candidate in candidates:
        try:
            if candidate.exists() and candidate.is_file():
                return candidate
        except Exception:
            continue
    return None


def _quality_target_paths(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> set[str]:
    paths: set[str] = set()
    if str(request.target_path or "").strip():
        paths.add(str(request.target_path or "").strip())
    for change in file_changes:
        if not isinstance(change, dict):
            continue
        for key in ("path", "file_path", "target_path"):
            value = str(change.get(key) or "").strip()
            if value:
                paths.add(value)
    return paths


def _target_docx_quality_paragraph_count(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> int:
    paths = set(_quality_target_paths(request, file_changes))
    target_names = {
        Path(path_text).name.casefold()
        for path_text in paths
        if str(path_text or "").strip()
    }
    for file_info in request.files or []:
        file_path = str(file_info.path or "").strip()
        file_name = str(file_info.name or file_path).strip()
        if file_info.target or Path(file_name).name.casefold() in target_names:
            if file_path:
                paths.add(file_path)
            if file_name:
                paths.add(file_name)
    paragraph_counts = [
        len(_read_docx_quality_paragraphs(path_text))
        for path_text in paths
        if str(path_text or "").lower().endswith((".docx", ".doc"))
    ]
    return max(paragraph_counts, default=0)


def _quality_file_is_target(file_info: FileTaskFile, target_paths: set[str]) -> bool:
    file_paths = [
        str(file_info.path or "").strip(),
        str(file_info.name or "").strip(),
    ]
    if target_paths:
        return any(
            _quality_same_path(candidate, target)
            for candidate in file_paths
            for target in target_paths
            if candidate and target
        )
    return bool(file_info.target)


def _quality_same_path(left: str, right: str) -> bool:
    left_text = str(left or "").replace("\\", "/").strip().rstrip("/").casefold()
    right_text = str(right or "").replace("\\", "/").strip().rstrip("/").casefold()
    if not left_text or not right_text:
        return False
    if left_text == right_text:
        return True
    left_path = Path(left_text)
    right_path = Path(right_text)
    return left_path.name == right_path.name and bool(left_path.suffix)


def _quality_stringify(value: Any) -> str:
    if isinstance(value, str):
        return value
    try:
        return json.dumps(value, ensure_ascii=False)
    except Exception:
        return str(value or "")


def _normalize_quality_text(value: Any) -> str:
    return re.sub(r"\s+", " ", str(value or "")).strip().casefold()


def _top_table_requirement_gate(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> Dict[str, Any] | None:
    requirement = _extract_top_table_requirement(request.task)
    if not requirement:
        return None
    source_path = _first_source_spreadsheet_path(request, file_changes)
    target_path = _first_target_docx_path(request, file_changes)
    if not source_path or not target_path:
        return None
    expected_labels = _expected_top_row_labels(
        source_path, requirement["sort_by"], requirement["count"]
    )
    if not expected_labels:
        return None
    tables = _read_docx_quality_tables(target_path)
    rows_passed, columns_passed = _docx_tables_match_top_requirement(
        tables,
        expected_labels,
        requirement.get("columns") or [],
    )
    duplicate_rows = _top_table_rows_duplicated_as_paragraphs(
        _read_docx_quality_paragraphs(target_path), expected_labels
    )
    passed = rows_passed and columns_passed and not duplicate_rows
    detail = (
        f"用户要求按 {requirement['sort_by']} 取 Top {requirement['count']} 表格；"
        f"期望顺序：{', '.join(expected_labels)}。"
    )
    if requirement.get("columns"):
        detail += f"期望列：{', '.join(requirement['columns'])}。"
    details = []
    details.append(
        "排序顺序已匹配。" if rows_passed else "目标 Word 表格未包含该排序结果。"
    )
    if requirement.get("columns"):
        details.append(
            "表格列已匹配。"
            if columns_passed
            else "目标 Word 表格列未严格匹配用户要求。"
        )
    if duplicate_rows:
        details.append("Top 表格行被重复写成了段落清单。")
    detail += "".join(details)
    return quality_gate_result(
        criterion="top_table_sorted_by_requested_metric",
        passed=passed,
        detail=detail,
        priority="critical",
    )


def _extract_top_table_requirement(task: str) -> Dict[str, Any] | None:
    task_text = str(task or "")
    match = re.search(
        r"\btop\s*(?P<count>\d+|one|two|three|four|five|six|seven|eight|nine|ten)\b.{0,80}?\bby\s+(?P<metric>[A-Za-z_][\w -]{0,40})",
        task_text,
        re.IGNORECASE,
    )
    if match:
        raw_count = match.group("count").casefold()
        count = int(raw_count) if raw_count.isdigit() else _TOP_WORDS.get(raw_count, 0)
        metric = re.split(
            r"\b(?:table|with|including|include|and|section)\b|[,.;，。；]",
            match.group("metric").strip(),
            maxsplit=1,
            flags=re.IGNORECASE,
        )[0].strip()
        if count > 0 and metric:
            return {
                "count": count,
                "sort_by": metric,
                "columns": _extract_requested_table_columns(task_text),
            }
    zh_match = re.search(
        r"(?:前|最高|最大|排名前)\s*(?P<count>\d+|一|二|三|四|五|六|七|八|九|十).{0,40}?(?:按|根据|以)(?P<metric>[\u4e00-\u9fffA-Za-z_][\w\u4e00-\u9fff -]{0,30})(?:排序|排名|取|筛选)",
        task_text,
        re.IGNORECASE,
    )
    if zh_match:
        count = _parse_top_count(zh_match.group("count"))
        metric = zh_match.group("metric").strip()
        if count > 0 and metric:
            return {
                "count": count,
                "sort_by": metric,
                "columns": _extract_requested_table_columns(task_text),
            }
    return None


def _extract_requested_table_columns(task: str) -> List[str]:
    match = re.search(
        r"\btable\s+with\s+(?P<columns>.+?)(?=(?:,\s*\d+\)|[.;；。]|,\s*(?:and\s+)?(?:a\s+)?(?:risk|next|summary|executive|section)\b|$))",
        str(task or ""),
        re.IGNORECASE,
    )
    if not match:
        return []
    columns_text = re.sub(
        r"\b(?:columns?|fields?)\b", "", match.group("columns"), flags=re.IGNORECASE
    )
    columns_text = re.sub(r"\s+", " ", columns_text).strip(" ,，、")
    raw_parts = re.split(
        r"\s*(?:,|，|、|\band\b)\s*", columns_text, flags=re.IGNORECASE
    )
    columns: List[str] = []
    for part in raw_parts:
        clean = re.sub(r"^(?:the|a|an)\s+", "", part.strip(), flags=re.IGNORECASE)
        clean = clean.strip(" ,，、")
        if clean and _normalize_quality_text(clean) not in {
            _normalize_quality_text(item) for item in columns
        }:
            columns.append(clean)
    return columns if len(columns) >= 2 else []


def _parse_top_count(value: str) -> int:
    text = str(value or "").strip()
    if text.isdigit():
        return int(text)
    return {
        "一": 1,
        "二": 2,
        "三": 3,
        "四": 4,
        "五": 5,
        "六": 6,
        "七": 7,
        "八": 8,
        "九": 9,
        "十": 10,
    }.get(text, 0)


def _first_source_spreadsheet_path(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> str:
    target_paths = _quality_target_paths(request, file_changes)
    for change in file_changes:
        if not isinstance(change, dict):
            continue
        source_path = str(change.get("source_path") or "").strip()
        if (
            source_path.lower().endswith((".xlsx", ".xlsm", ".csv"))
            and Path(source_path).exists()
        ):
            return source_path
    for file_info in request.files or []:
        path = str(file_info.path or file_info.name or "").strip()
        if path.lower().endswith(
            (".xlsx", ".xlsm", ".csv")
        ) and not _quality_file_is_target(file_info, target_paths):
            return path
    return ""


def _first_target_docx_path(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> str:
    if str(request.target_path or "").lower().endswith((".docx", ".doc")):
        return str(request.target_path or "").strip()
    for change in file_changes:
        if not isinstance(change, dict):
            continue
        path = str(change.get("path") or change.get("file_path") or "").strip()
        if path.lower().endswith((".docx", ".doc")):
            return path
    return ""


def _expected_top_row_labels(source_path: str, sort_by: str, count: int) -> List[str]:
    try:
        import openpyxl

        workbook = openpyxl.load_workbook(source_path, read_only=True, data_only=True)
        worksheet = workbook[workbook.sheetnames[0]]
        rows = [
            ["" if value is None else str(value) for value in row]
            for row in worksheet.iter_rows(values_only=True)
            if any(str(value or "").strip() for value in row)
        ]
        workbook.close()
    except Exception:
        return []
    if len(rows) < 2:
        return []
    headers = rows[0]
    sort_index = _quality_header_index(headers, sort_by)
    if sort_index is None:
        return []
    data_rows = rows[1:]
    sorted_rows = sorted(
        data_rows,
        key=lambda row: _quality_sort_value(
            row[sort_index] if sort_index < len(row) else ""
        ),
        reverse=True,
    )
    return [
        str(row[0]).strip()
        for row in sorted_rows[:count]
        if row and str(row[0]).strip()
    ]


def _read_docx_quality_tables(path_text: str) -> List[List[List[str]]]:
    try:
        from docx import Document

        document = Document(str(path_text))
    except Exception:
        return []
    tables: List[List[List[str]]] = []
    for table in document.tables:
        table_rows: List[List[str]] = []
        for row in table.rows:
            table_rows.append([cell.text.strip() for cell in row.cells])
        tables.append(table_rows)
    return tables


def _docx_tables_include_ordered_labels(
    tables: List[List[List[str]]], expected_labels: List[str]
) -> bool:
    expected = [_normalize_quality_text(label) for label in expected_labels if label]
    if not expected:
        return False
    for table in tables:
        row_labels = [
            _normalize_quality_text(row[0])
            for row in table[1:]
            if row and _normalize_quality_text(row[0])
        ]
        if len(row_labels) < len(expected):
            continue
        for start in range(0, len(row_labels) - len(expected) + 1):
            if row_labels[start : start + len(expected)] == expected:
                return True
    return False


def _docx_tables_match_top_requirement(
    tables: List[List[List[str]]],
    expected_labels: List[str],
    expected_columns: List[str],
) -> tuple[bool, bool]:
    expected = [_normalize_quality_text(label) for label in expected_labels if label]
    expected_headers = [_normalize_quality_text(column) for column in expected_columns]
    if not expected:
        return False, not expected_headers
    any_rows_match = False
    any_columns_match = not expected_headers
    for table in tables:
        if not table:
            continue
        row_labels = [
            _normalize_quality_text(row[0])
            for row in table[1:]
            if row and _normalize_quality_text(row[0])
        ]
        rows_match = False
        for start in range(0, max(0, len(row_labels) - len(expected) + 1)):
            if row_labels[start : start + len(expected)] == expected:
                rows_match = True
                break
        if not rows_match:
            continue
        any_rows_match = True
        if not expected_headers:
            return True, True
        headers = [
            _normalize_quality_text(cell)
            for cell in table[0]
            if _normalize_quality_text(cell)
        ]
        columns_match = headers == expected_headers
        if columns_match:
            return True, True
    return any_rows_match, any_columns_match


def _top_table_rows_duplicated_as_paragraphs(
    paragraphs: List[str], expected_labels: List[str]
) -> bool:
    expected = [_normalize_quality_text(label) for label in expected_labels if label]
    if not expected:
        return False
    duplicate_like_count = 0
    for paragraph in paragraphs:
        normalized = _normalize_quality_text(paragraph)
        if not any(label in normalized for label in expected):
            continue
        if re.search(r"^\s*(?:customer|客户)\s*[:：]", paragraph, re.IGNORECASE) or (
            "|" in paragraph
            and re.search(r"\b(?:region|revenue|margin)\b", paragraph, re.IGNORECASE)
        ):
            duplicate_like_count += 1
    return duplicate_like_count >= min(2, len(expected))


def _quality_header_index(headers: List[str], wanted: str) -> int | None:
    wanted_text = _normalize_quality_text(wanted)
    normalized_headers = [_normalize_quality_text(header) for header in headers]
    for index, header in enumerate(normalized_headers):
        if header == wanted_text:
            return index
    for index, header in enumerate(normalized_headers):
        if wanted_text in header or header in wanted_text:
            return index
    return None


def _quality_sort_value(value: Any) -> tuple[int, Any]:
    text = str(value or "").strip()
    if not text:
        return (0, 0)
    numeric_text = re.sub(r"[,$%￥¥\s]", "", text)
    try:
        return (1, float(numeric_text))
    except ValueError:
        return (1, text.casefold())


def _explicit_docx_section_gates(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> List[Dict[str, Any]]:
    target_path = _first_target_docx_path(request, file_changes)
    if not target_path:
        return []
    task_text = str(request.task or "")
    paragraphs = _read_docx_quality_paragraphs(target_path)
    if not paragraphs:
        return []
    gates: List[Dict[str, Any]] = []
    if _task_requires_risk_section(task_text):
        passed = _has_heading_like_paragraph(paragraphs, ("risk", "risks", "风险"))
        gates.append(
            quality_gate_result(
                criterion="required_risk_section_present",
                passed=passed,
                detail=(
                    "用户要求独立风险 section；"
                    + (
                        "目标 Word 中已找到风险章节。"
                        if passed
                        else "目标 Word 中未找到独立风险章节。"
                    )
                ),
                priority="critical",
            )
        )
    next_action_count = _required_next_action_count(task_text)
    if next_action_count:
        passed = _has_next_actions(paragraphs, next_action_count)
        gates.append(
            quality_gate_result(
                criterion="required_next_actions_present",
                passed=passed,
                detail=(
                    f"用户要求 {next_action_count} 条具体 next actions；"
                    + (
                        "目标 Word 中已找到足够的行动项。"
                        if passed
                        else "目标 Word 中未找到足够的行动项。"
                    )
                ),
                priority="critical",
            )
        )
    return gates


def _task_requires_risk_section(task: str) -> bool:
    return bool(
        re.search(r"\brisk\s+section\b", task, re.IGNORECASE)
        or re.search(r"风险.{0,8}(?:章节|部分|小节|板块|section)", task)
    )


def _required_next_action_count(task: str) -> int:
    match = re.search(
        r"\b(?P<count>\d+|one|two|three|four|five|six|seven|eight|nine|ten)\s+(?:concrete\s+)?next actions?\b",
        task,
        re.IGNORECASE,
    )
    if match:
        raw_count = match.group("count").casefold()
        return int(raw_count) if raw_count.isdigit() else _TOP_WORDS.get(raw_count, 0)
    zh_match = re.search(
        r"(?:至少|正好|共)?\s*(?P<count>\d+|一|二|三|四|五|六|七|八|九|十)"
        r"\s*(?:条|项|个)\s*(?:具体)?(?:下一步|行动项|行动|动作)",
        task,
    )
    if zh_match:
        return _parse_top_count(zh_match.group("count"))
    return 0


def _read_docx_quality_paragraphs(path_text: str) -> List[str]:
    try:
        from docx import Document

        document = Document(str(path_text))
    except Exception:
        return []
    return [
        paragraph.text.strip()
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]


def _has_heading_like_paragraph(
    paragraphs: List[str], keywords: tuple[str, ...]
) -> bool:
    for paragraph in paragraphs:
        normalized = _normalize_quality_text(paragraph)
        if len(normalized) > 80:
            continue
        if len(normalized) > 40 and re.search(r"[.!?。！？]", paragraph):
            continue
        if any(keyword in normalized for keyword in keywords):
            return True
    return False


def _has_next_actions(paragraphs: List[str], required_count: int) -> bool:
    combined = "\n".join(paragraphs)
    markers = re.findall(r"(?:^|\s)(?:\d+[.)、]|[-•])\s+\S+", combined)
    if len(markers) >= required_count:
        return True
    has_heading = _has_heading_like_paragraph(
        paragraphs, ("next action", "next actions", "下一步", "行动项")
    )
    action_like = [
        paragraph
        for paragraph in paragraphs
        if re.search(
            r"\b(?:prioritize|initiate|monitor|track|assign|schedule|review|follow up|推进|启动|跟进|监控|安排|复盘)\b",
            paragraph,
            re.IGNORECASE,
        )
    ]
    return has_heading and len(action_like) >= required_count


def repair_retry_message(
    request: FileTaskRequest,
    check_payload: Dict[str, Any],
    file_changes: List[Dict[str, Any]],
) -> str:
    lines = [
        "核验未通过，当前任务还不能结束。下一轮必须修复目标文件，而不是重复上一轮完全相同的调用。",
    ]
    status = str(check_payload.get("status") or "").strip()
    summary = str(check_payload.get("summary") or "").strip()
    if status:
        lines.append(f"当前核验状态：{status}")
    if summary:
        lines.append(f"核验摘要：{summary}")
    if request.target_path:
        lines.append(f"目标文件：{request.target_path}")

    request_files = list(getattr(request, "files", []) or [])
    if not request_files:
        request_files.extend(
            files_explicitly_mentioned_in_task(
                workspace_root=Path.cwd() / "workspace",
                task=request.task,
            )
        )
    recipe_match = select_task_recipe(request, request_files, write_intent=True)
    if recipe_match:
        lines.append(f"当前任务路线：{recipe_match.recipe.id}")
        if recipe_match.recipe.success_criteria:
            lines.append("本路线验收标准：")
            for criterion in recipe_match.recipe.success_criteria[:5]:
                text = str(criterion or "").strip()
                if text:
                    lines.append(f"- {text}")
    if _looks_like_financial_xlsx_docx_chart_report_task(request, request_files):
        lines.append(
            "财务预测图表写入修复要求：本任务不能只完成 Python 计算或打印 stdout。"
            "必须产生写入工具事件：write_docx_content 写入问题清单/分析结论，"
            "insert_excel_as_docx_table 写入关键数据表，insert_image_into_docx 插入至少两张真实 PNG/JPG 图表。"
        )
        lines.append(
            "Excel 解析要求：如果 pandas 读出的列名是 Unnamed，不要用 df.columns 判断年份列；"
            "应扫描表格行，找到包含 2025E/2026E/2027E/2028E 等年份标签的 header row，"
            "再根据这些列抽取“收入合计、毛利合计、费用合计、净利润、销量”等指标。"
        )

    remaining = (
        check_payload.get("remaining")
        if isinstance(check_payload.get("remaining"), list)
        else []
    )
    if remaining:
        lines.append("仍需满足：")
        for index, item in enumerate(remaining[:5], start=1):
            text = str(item or "").strip()
            if text:
                lines.append(f"{index}. {text}")

    if any(
        "表格" in str(item or "") and "Word" in str(item or "") for item in remaining
    ):
        target_path = str(request.target_path or "").strip()
        xlsx_sources = [
            str(file_info.path or file_info.name or "").strip()
            for file_info in request_files
            if str(
                file_info.type
                or Path(str(file_info.path or file_info.name)).suffix.lstrip(".")
            )
            .lower()
            .strip()
            in {"xlsx", "xlsm", "xls", "csv"}
            and str(file_info.path or file_info.name or "").strip()
        ]
        if xlsx_sources:
            lines.append(
                "表格修复要求：调用 insert_excel_as_docx_table，把源表格作为真实 Word 表格写入目标 DOCX；"
                "只调用 write_docx_content 写文字段落不能通过。"
            )
            lines.append(f"推荐 source_path：{xlsx_sources[0]}")
        else:
            lines.append(
                "表格修复要求：当前没有表格源文件。请使用可用的原生 Word 表格写入能力，"
                "或通过 run_python_code 使用 python-docx 打开并保存 TASK_TARGET_PATH，"
                "在该沙箱目标中创建真实 table；不要直接写宿主绝对路径。"
                "只写标题和文字段落不能通过。"
            )
        if target_path:
            lines.append(f"推荐 target_path：{target_path}")

    if file_changes:
        lines.append("已观察到的文件变更：")
        for change in file_changes[-3:]:
            if not isinstance(change, dict):
                continue
            change_summary = str(change.get("summary") or "").strip()
            path_text = str(change.get("path") or change.get("file_path") or "").strip()
            if change_summary and path_text:
                lines.append(f"- {path_text}: {change_summary}")
            elif change_summary:
                lines.append(f"- {change_summary}")
            elif path_text:
                lines.append(f"- {path_text}")

    lines.append(
        "要求：先理解核验失败原因；只有当参数、代码、工具选择或写入位置已经改变时，才允许再次调用工具；修复后再结束。"
    )
    return "\n".join(lines)


def success_criteria(
    request: FileTaskRequest,
    *,
    write_intent: bool,
    output_mode: str,
) -> List[str]:
    criteria = [
        "每个步骤都产生 typed event，可被前端时间线渲染",
        "所有上下文来源都来自显式输入",
    ]
    recipe_match = select_task_recipe(
        request, request.files or [], write_intent=write_intent
    )
    if recipe_match and recipe_match.recipe.success_criteria:
        criteria.extend(
            str(item)
            for item in recipe_match.recipe.success_criteria
            if str(item or "").strip()
        )
        return criteria
    if write_intent:
        criteria.extend(
            [
                "写入工具必须产生 file.changed 事件",
                "最终 checker 必须确认目标文件已更新",
            ]
        )
    elif output_mode == "hybrid":
        criteria.append("最终摘要必须给出明确建议，且当前轮不默认直接写入原文件")
    else:
        criteria.append("最终摘要说明已使用的上下文和未完成项")
    return criteria


def _file_types(files: List[FileTaskFile]) -> set[str]:
    return request_file_types(files)


def _looks_like_chart_request(task: str) -> bool:
    return semantic_markers(task).get("chart_request", False)


def _looks_like_problem_analysis_request(task: str) -> bool:
    return semantic_markers(task).get("problem_analysis_request", False)


def _looks_like_table_request(task: str) -> bool:
    return semantic_markers(task).get("table_request", False)


def _looks_like_docx_table_output_request(task: str) -> bool:
    return bool(_DOCX_TABLE_OUTPUT_PATTERN.search(str(task or "")))


def _docx_table_requirement_gate(
    request: FileTaskRequest,
    file_changes: List[Dict[str, Any]],
    target_type: str,
    operations: set[str],
    reported_rows: int,
) -> Dict[str, Any] | None:
    task = str(request.task or "")
    if target_type not in {"docx", "doc"}:
        return None
    if not (
        _looks_like_docx_table_output_request(task)
        or (
            _looks_like_table_request(task)
            and not _looks_like_problem_analysis_request(task)
        )
    ):
        return None
    table_count, actual_rows, max_columns = _target_docx_quality_table_metrics(
        request, file_changes
    )
    passed = ("insert_excel_as_docx_table" in operations and reported_rows > 0) or (
        table_count >= 1 and actual_rows >= 2 and max_columns >= 2
    )
    return quality_gate_result(
        criterion="docx_table_request_has_table",
        passed=passed,
        detail=(
            "用户要求表格数据进入 Word；"
            f"工具报告写入行数：{reported_rows}，"
            f"目标文档实际表格数：{table_count}，"
            f"实际表格总行数：{actual_rows}，最大列数：{max_columns}。"
        ),
        priority="critical",
    )


def _target_docx_quality_table_metrics(
    request: FileTaskRequest, file_changes: List[Dict[str, Any]]
) -> tuple[int, int, int]:
    best = (0, 0, 0)
    for path_text in _quality_target_paths(request, file_changes):
        if not str(path_text or "").lower().endswith((".docx", ".doc")):
            continue
        resolved = _resolve_quality_path(path_text)
        if not resolved:
            continue
        try:
            from docx import Document

            document = Document(str(resolved))
        except Exception:
            continue
        table_count = len(document.tables)
        total_rows = sum(len(table.rows) for table in document.tables)
        max_columns = max((len(table.columns) for table in document.tables), default=0)
        best = max(best, (table_count, total_rows, max_columns))
    return best


def _financial_report_duplicate_content_gate(
    request: FileTaskRequest,
    file_changes: List[Dict[str, Any]],
) -> Dict[str, Any] | None:
    inspected = False
    repeated_title_count = 0
    duplicate_table_count = 0
    for path_text in _quality_target_paths(request, file_changes):
        if not str(path_text or "").lower().endswith((".docx", ".doc")):
            continue
        resolved = _resolve_quality_path(path_text)
        if not resolved:
            continue
        try:
            from collections import Counter

            from docx import Document

            document = Document(str(resolved))
        except Exception:
            continue

        inspected = True
        title_counts = Counter(
            " ".join(paragraph.text.split()).casefold()
            for paragraph in document.paragraphs
            if paragraph.style is not None
            and str(paragraph.style.name or "").strip().casefold() == "title"
            and " ".join(paragraph.text.split())
        )
        repeated_title_count = max(
            repeated_title_count,
            sum(count - 1 for count in title_counts.values() if count > 1),
        )

        table_fingerprints = Counter(
            tuple(
                tuple(" ".join(cell.text.split()).casefold() for cell in row.cells)
                for row in table.rows
            )
            for table in document.tables
            if table.rows
        )
        duplicate_table_count = max(
            duplicate_table_count,
            sum(count - 1 for count in table_fingerprints.values() if count > 1),
        )

    if not inspected:
        return None

    passed = repeated_title_count == 0 and duplicate_table_count == 0
    return quality_gate_result(
        criterion="financial_report_has_no_duplicate_report_body",
        passed=passed,
        detail=(
            "新建财务报告不能包含重复的报告标题或完全相同的数据表；"
            f"重复 Title 段落数：{repeated_title_count}，"
            f"重复表格数：{duplicate_table_count}。"
        ),
        priority="critical",
    )


def _table_narrative_requirement_gate(
    task: str,
    target_type: str,
    operations: set[str],
    paragraphs_written: int,
) -> Dict[str, Any] | None:
    if target_type not in {"docx", "doc"} or not _looks_like_table_request(task):
        return None
    if not re.search(
        r"(?:表格.{0,24}(?:说明|摘要|总结|结论|分析|要点)|"
        r"(?:说明|摘要|总结|结论|分析|要点|写一句|一句简短说明).{0,40}表格|"
        r"\b(?:summary|narrative|explanation|analysis|conclusion)\b.{0,40}\btable\b|"
        r"\btable\b.{0,40}\b(?:summary|narrative|explanation|analysis|conclusion)\b)",
        task,
        re.IGNORECASE,
    ):
        return None
    return quality_gate_result(
        criterion="docx_table_request_has_narrative",
        passed="write_docx_content" in operations and paragraphs_written >= 1,
        detail=(
            "用户要求 Word 表格配套文字说明/摘要；"
            f"当前说明段落写入数：{paragraphs_written}。"
        ),
        priority="high",
    )


def _looks_like_translation_request(task: str) -> bool:
    return semantic_markers(task).get("translation_request", False)


def _looks_like_polish_request(task: str) -> bool:
    return semantic_markers(task).get("polish_request", False)


def _looks_like_ppt_slide_write_request(
    request: FileTaskRequest, files: List[FileTaskFile]
) -> bool:
    return semantic_markers(request.task, file_types=_file_types(files)).get(
        "ppt_slide_write_request", False
    )


def _looks_like_docx_report_request(
    request: FileTaskRequest, files: List[FileTaskFile]
) -> bool:
    return semantic_markers(
        request.task,
        file_types=_file_types(files),
        target_file_type=request_target_file_type(request, files),
    ).get("docx_report_request", False)


def _looks_like_local_docx_edit_request(task: str) -> bool:
    text = str(task or "")
    if not text.strip():
        return False
    return bool(
        re.search(
            r"(?:只|仅|只是|仅仅)?(?:追加|添加|插入|加上|补充|写入).{0,30}"
            r"(?:一句|一段|一条|少量|这句|这段|文本|段落)",
            text,
            re.IGNORECASE,
        )
        or re.search(
            r"(?:保留|保持).{0,30}(?:已有|现有|原有).{0,20}"
            r"(?:表格|结构|内容|格式).{0,12}(?:不变|不修改|不要改|原样)",
            text,
            re.IGNORECASE,
        )
        or re.search(
            r"\b(?:append|insert|add)\b.{0,50}\b(?:sentence|paragraph|line)\b",
            text,
            re.IGNORECASE,
        )
        or re.search(
            r"\b(?:keep|preserve)\b.{0,40}\b(?:existing|current|original)\b"
            r".{0,30}\b(?:table|structure|content|format)\b",
            text,
            re.IGNORECASE,
        )
    )


def _looks_like_local_docx_replace_request(task: str) -> bool:
    text = str(task or "")
    if not text.strip():
        return False
    return bool(
        re.search(
            r"(?:只替换|仅替换|替换.{0,16}(?:一处|那一处|这一处|第\s*\d+\s*处|第\s*\d+\s*段)|"
            r"(?:第\s*\d+\s*段|一处|那一处|这一处).{0,24}(?:替换|删除|改成|换成)|"
            r"local(?:ized)?\s+(?:replace|edit)|replace\s+only)",
            text,
            re.IGNORECASE,
        )
    )


def _looks_like_financial_xlsx_docx_chart_report_task(
    request: FileTaskRequest, files: List[FileTaskFile]
) -> bool:
    return semantic_markers(
        request.task,
        file_types=_file_types(files),
        target_file_type=request_target_file_type(request, files),
    ).get("financial_xlsx_docx_chart_report", False)
