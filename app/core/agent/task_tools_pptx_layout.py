# -*- coding: utf-8 -*-
# Copyright (C) 2024-2026 Koto AI. All rights reserved.
# SPDX-License-Identifier: LicenseRef-Koto-Proprietary
from __future__ import annotations

import json
from typing import Any, Dict, List, Optional

from app.core.agent.task_tools_pptx_theme import _hex_to_rgb_color


def _parse_jsonish_list(value: Any, field_name: str) -> tuple[List[Any], Optional[str]]:
    if isinstance(value, str):
        try:
            parsed = json.loads(value)
        except json.JSONDecodeError as exc:
            return [], f"Invalid {field_name} JSON: {exc}"
    else:
        parsed = value
    if parsed is None or parsed == "":
        return [], None
    if isinstance(parsed, list):
        return parsed, None
    if isinstance(parsed, dict):
        return [parsed], None
    return [{"content": str(parsed)}], None


def _pptx_text_lines(value: Any) -> List[str]:
    if value is None:
        return []
    if isinstance(value, str):
        return [line.strip() for line in value.splitlines() if line.strip()]
    if isinstance(value, list):
        lines: List[str] = []
        for item in value:
            lines.extend(_pptx_text_lines(item))
        return lines
    if isinstance(value, dict):
        for key in ("text", "content", "body", "bullet", "point", "summary"):
            if key in value:
                return _pptx_text_lines(value.get(key))
        for key in ("bullets", "bullet_points", "points", "items", "lines"):
            if key in value:
                return _pptx_text_lines(value.get(key))
        return [
            "：".join(str(part).strip() for part in (key, val) if str(part).strip())
            for key, val in value.items()
        ]
    return [str(value).strip()] if str(value).strip() else []


def _pptx_first_text(value: Any) -> str:
    lines = _pptx_text_lines(value)
    return lines[0] if lines else ""


def _remove_koto_theme_shapes(slide: Any) -> None:
    for shape in list(slide.shapes):
        if str(getattr(shape, "name", "") or "").startswith("KOTO_THEME_"):
            element = shape._element
            element.getparent().remove(element)


def _is_title_shape(slide: Any, shape: Any) -> bool:
    if getattr(slide.shapes, "title", None) is shape:
        return True
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        return shape.placeholder_format.type in {
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        }
    except Exception:
        return False


def _is_body_placeholder(shape: Any) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER

        return shape.placeholder_format.type in {
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.SUBTITLE,
        }
    except Exception:
        return False


def _apply_text_style(
    text_frame: Any, *, font_family: str, size_pt: float, color: Any, bold: bool = False
) -> None:
    from pptx.util import Pt

    for paragraph in text_frame.paragraphs:
        paragraph.font.name = font_family
        paragraph.font.size = Pt(size_pt)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = color
        for run in paragraph.runs:
            run.font.name = font_family
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = color


def _set_slide_background(slide: Any, color: Any) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_theme_background_shape(
    slide: Any, slide_width: int, slide_height: int, color: Any
) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height
    )
    background.name = "KOTO_THEME_BACKGROUND"
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    element = background._element
    tree = element.getparent()
    tree.remove(element)
    tree.insert(2, element)


def _add_theme_accent_shapes(
    slide: Any,
    slide_width: int,
    slide_height: int,
    theme: Dict[str, Any],
    slide_number: int,
) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    accent = _hex_to_rgb_color(theme["accent"], "0F766E")
    accent2 = _hex_to_rgb_color(theme["accent2"], "D97706")
    muted = _hex_to_rgb_color(theme["muted"], "E6DED2")
    footer_text = _hex_to_rgb_color(
        theme["body_text"] if not theme.get("is_dark") else theme["inverse_text"],
        "25313B",
    )

    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, slide_width, Inches(0.08)
    )
    top_bar.name = "KOTO_THEME_ACCENT_BAR"
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent
    top_bar.line.fill.background()

    corner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, slide_height - Inches(0.09), slide_width, Inches(0.09)
    )
    corner.name = "KOTO_THEME_FOOTER_BAR"
    corner.fill.solid()
    corner.fill.fore_color.rgb = muted
    corner.line.fill.background()

    marker = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(0.94), Inches(0.46), Inches(0.06)
    )
    marker.name = "KOTO_THEME_TITLE_MARKER"
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent2
    marker.line.fill.background()

    footer = slide.shapes.add_textbox(
        slide_width - Inches(1.2),
        slide_height - Inches(0.36),
        Inches(0.72),
        Inches(0.18),
    )
    footer.name = "KOTO_THEME_SLIDE_NUMBER"
    footer.text_frame.clear()
    footer.text_frame.paragraphs[0].text = f"{slide_number:02d}"
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    footer.text_frame.paragraphs[0].font.name = str(
        theme.get("font_family") or "Microsoft YaHei"
    )
    footer.text_frame.paragraphs[0].font.size = Pt(8)
    footer.text_frame.paragraphs[0].font.color.rgb = footer_text
