# =============================================================================
# Koto Task Tools — Public API Index
# =============================================================================
# Lines: 5395 (guarded) | Public: 33 | Private helpers: 100+
#
# QUICK NAVIGATION:
#   XLSX Tools ........... read_sheet_data → write_sheet_data
#   DOCX Read ............ read_docx_content
#   File Parse ........... parse_file_to_text
#   Workspace Files ...... list_workspace_files, open_file_in_editor
#   LLM Tools ............ llm_extract, llm_transform
#   File Compare ......... compare_files
#   DOCX Annotate ........ compare_docx_and_annotate → write_docx_comments
#   File Create .......... create_file, extract_to_file
#   DOCX Edit ............ write_docx_content, fill_docx_template
#   PDF Convert .......... convert_docx_to_pdf, convert_file
#   PPTX Tools ........... write_pptx_slides → add_pptx_slides
#   TaskToolsPlugin ...... Agent plugin class
#
# PRIVATE HELPERS (internal, called by public functions):
#   General I/O .... _resolve_path, _result_path, _success_result, _blocked_write_result
#   XLSX helpers ... _select_workbook_sheet, _build_workbook_structure_payload, ...
#   DOCX helpers ... _docx_nonempty_paragraph_texts, _build_docx_compare_annotations, ...
#   Sandbox ........ _stage_task_files_for_sandbox, _sync_staged_files_to_source, ...
#   PDF helpers .... _read_pdf_excerpt, _read_pdf_letter_window, ...
# =============================================================================
# ══════════════════════════════════════════════════════════════
# task_tools.py — Composable file-operation tools for FileTaskRuntime
#
# These tools are the building blocks the AI orchestrates freely
# to accomplish user tasks on workspace files.  Each tool is
# self-contained: read → process → write, with no hardcoded
# workflow assumptions.
# ══════════════════════════════════════════════════════════════
from __future__ import annotations

import base64
import importlib
import io
import json
import logging
import os
import filecmp
import hashlib
import re
import shutil
import sys
import tempfile
import stat
import time
import types
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional

from app.core.agent.base import AgentPlugin
from app.core.agent.file_task_contract import FileTaskToolStreamResult
from app.core.agent.file_task_completion_verifier import verify_task_completion
from app.core.agent.file_task_result_markers import (
    KOTO_CREATED_FALLBACK_KEY,
    KOTO_CREATED_RESULT_KEY,
    KOTO_CREATED_RESULT_MARKER,
    KOTO_MODIFIED_FALLBACK_KEY,
    KOTO_MODIFIED_RESULT_KEY,
    KOTO_MODIFIED_RESULT_MARKER,
)
from app.core.agent.path_utils import (
    default_search_roots,
    has_parent_path_segment,
    is_within_roots,
    resolve_existing_path,
)
from app.core.agent.task_tools_docx_minimal import (
    _coerce_docx_paragraphs_for_write,
    _minimal_docx_package_bytes,
    _normalize_docx_paragraphs,
    _plain_text_to_docx_paragraphs,
)
from app.core.agent.task_tools_docx_style import apply_docx_style as _apply_docx_style
from app.core.agent.task_tools_docx_compare import (
    _contract_risk_summary_from_annotations,
    _build_docx_compare_annotations,
    _docx_compare_annotation_candidates,
    _docx_diff_comment,
    _docx_diff_key,
    _docx_nonempty_paragraph_texts,
    _short_docx_diff_text,
)
from app.core.agent.task_tools_pptx_theme import (
    _hex_to_rgb_color,
    _pptx_density_settings,
    _select_pptx_theme,
)
from app.core.agent.task_tools_pptx_layout import (
    _add_theme_accent_shapes,
    _add_theme_background_shape,
    _apply_text_style,
    _is_body_placeholder,
    _is_title_shape,
    _parse_jsonish_list,
    _pptx_first_text,
    _pptx_text_lines,
    _remove_koto_theme_shapes,
    _set_slide_background,
)
from app.core.agent.task_tools_office_create import (
    create_docx_file as _office_create_docx_file,
    create_pptx_file as _office_create_pptx_file,
    create_xlsx_file as _office_create_xlsx_file,
    plain_text_to_pptx_slides as _office_plain_text_to_pptx_slides,
)
from app.core.agent.task_tools_conversion import (
    convert_docx_to_pdf as _conversion_convert_docx_to_pdf,
    convert_docx_to_pdf_with_docx2pdf as _conversion_docx2pdf,
    convert_docx_to_pdf_with_libreoffice as _conversion_libreoffice,
    convert_docx_to_pdf_with_word as _conversion_word,
    convert_file as _conversion_convert_file,
    list_conversions as _conversion_list_conversions,
    normalize_conversion_extension as _conversion_normalize_extension,
)
from app.core.agent.task_tools_docx_template import (
    replace_docx_placeholders_in_paragraph as _replace_docx_placeholders_in_paragraph,
)
from app.core.agent.task_tools_docx_review_cleanup import (
    DOCX_COMMENT_MARKUP_TAGS as _DOCX_COMMENT_MARKUP_TAGS,
    DOCX_W_NS as _DOCX_W_NS,
    accept_docx_revision_markup as _accept_docx_revision_markup,
    build_docx_review_clear_summary as _build_docx_review_clear_summary,
    normalize_docx_review_clear_scope as _review_normalize_docx_review_clear_scope,
    remove_comments_content_type_override as _remove_comments_content_type_override,
    remove_comments_relationships_xml as _remove_comments_relationships_xml,
    remove_docx_comment_markup as _remove_docx_comment_markup,
    _serialize_xml_root,
)
from app.core.agent.task_tools_pdf_window import (
    int_to_chinese_letter_number as _int_to_chinese_letter_number,
    int_to_pdf_roman as _int_to_pdf_roman,
    pdf_letter_heading_terms as _pdf_letter_heading_terms,
    pdf_page_has_letter_heading as _pdf_page_has_letter_heading,
    read_pdf_excerpt as _pdf_read_excerpt,
    read_pdf_letter_window as _pdf_read_letter_window,
)
from app.core.agent.task_tools_xlsx_sheet_selection import (
    select_workbook_sheet as _select_workbook_sheet,
    sheet_matches_statement as _sheet_matches_statement,
)
from app.core.agent.task_tools_xlsx_structure import (
    collect_formula_examples as _collect_formula_examples,
    detect_year_header as _detect_year_header,
    extract_external_link_targets as _extract_external_link_targets,
    sample_sheet_rows as _sample_sheet_rows,
    row_label_for_year_series as _row_label_for_year_series,
    display_series_value as _display_series_value,
    severity_for_financial_label as _xlsx_severity_for_financial_label,
)
from app.core.agent.task_tool_operation_bindings import build_task_tool_operations
from app.core.agent.task_tools_registry import build_task_tool_definitions
from app.core.services.file_service import FileService

logger = logging.getLogger(__name__)

_TEXT_LIMIT_MIN = 1_000
_TEXT_LIMIT_DEFAULT = 60_000
_TEXT_LIMIT_DOCX_DEFAULT = 24_000
_TEXT_LIMIT_MAX = 200_000
_TASK_TOOL_LLM_CALL_TIMEOUT = float(os.getenv("KOTO_FILE_TASK_LLM_TIMEOUT", "45"))
_SANDBOX_CLEANUP_RETRIES = int(os.getenv("KOTO_SANDBOX_CLEANUP_RETRIES", "3"))
_SANDBOX_CLEANUP_RETRY_DELAY_SECONDS = float(
    os.getenv("KOTO_SANDBOX_CLEANUP_RETRY_DELAY", "0.1")
)

# ── Workspace root (same resolver as WorkspaceEditorPlugin) ──────────────────
_WORKSPACE_ROOT: Optional[str] = None


def _get_workspace_root() -> str:
    global _WORKSPACE_ROOT
    if _WORKSPACE_ROOT is None:
        configured_root = str(os.getenv("KOTO_WORKSPACE_DIR") or "").strip()
        if configured_root:
            workspace_root = Path(configured_root).expanduser().resolve()
        elif getattr(sys, "frozen", False):
            workspace_root = Path(sys.executable).resolve().parent / "workspace"
        else:
            project_root = Path(__file__).resolve().parent.parent.parent.parent
            workspace_root = project_root / "workspace"
        _WORKSPACE_ROOT = str(workspace_root)
    return _WORKSPACE_ROOT


def _get_project_root() -> str:
    """Return the project root (parent of workspace)."""
    return os.path.dirname(_get_workspace_root())


def _file_service(*, backup_enabled: bool = False) -> FileService:
    return FileService(
        workspace_dir=_get_workspace_root(),
        backup_enabled=backup_enabled,
    )


def _safe_resolve(relative_path: str) -> Optional[str]:
    """Resolve a user path inside workspace root or project root.

    Tries workspace root first, then project root as fallback.
    Returns None on traversal or if neither resolves.
    """
    root = _get_workspace_root()
    project_root = _get_project_root()
    # Strip leading "workspace/" prefix — the model sometimes includes it even
    # though paths are already relative to the workspace root.
    stripped = relative_path.replace("\\", "/")
    if has_parent_path_segment(stripped):
        return None
    if stripped.startswith("workspace/"):
        stripped = stripped[len("workspace/") :]
    try:
        resolved = os.path.normpath(os.path.join(root, stripped))
        if not is_within_roots(resolved, [root]):
            try:
                project_resolved = os.path.normpath(os.path.join(project_root, stripped))
                if is_within_roots(project_resolved, [project_root]):
                    return project_resolved
            except (ValueError, TypeError):
                pass
            return None
        if not os.path.exists(resolved):
            try:
                project_resolved = os.path.normpath(os.path.join(project_root, stripped))
                if (is_within_roots(project_resolved, [project_root])
                        and os.path.exists(project_resolved)):
                    return project_resolved
            except (ValueError, TypeError):
                pass
        return resolved
    except (ValueError, TypeError):
        return None


def _resolve_path(path: str, *, must_exist: bool = True) -> Optional[str]:
    """Accept both absolute and relative-to-workspace paths.

    Args:
        path: File path, absolute or relative.
        must_exist: If True (default), returns None for non-existent files.
                    If False, returns the resolved path even if the file does not exist yet.
    """
    if os.path.isabs(path):
        normalized = os.path.normpath(path)
        if must_exist and not os.path.exists(normalized):
            return None
        return normalized
    if has_parent_path_segment(path):
        return None

    # Keep workspace-relative priority for backward compatibility.
    ws_candidate = _safe_resolve(path)
    if ws_candidate and os.path.exists(ws_candidate):
        return ws_candidate

    # Try project root as well
    project_root = _get_project_root()
    try:
        project_candidate = os.path.normpath(os.path.join(project_root, path.replace("\\", "/")))
        if (is_within_roots(project_candidate, [project_root])
                and os.path.exists(project_candidate)):
            return project_candidate
    except (ValueError, TypeError):
        pass

    # For write targets, return the workspace candidate even if it does not exist
    if not must_exist and ws_candidate:
        return ws_candidate

    roots = [_get_workspace_root(), *default_search_roots()]
    resolved, _ = resolve_existing_path(path, roots=roots)
    if resolved and (
        is_within_roots(resolved, [_get_workspace_root()])
        or is_within_roots(resolved, [project_root])
    ):
        return resolved
    # For write targets, try project root as a last resort
    if not must_exist:
        try:
            fallback = os.path.normpath(os.path.join(project_root, path.replace("\\", "/")))
            if is_within_roots(fallback, [project_root]):
                return fallback
        except (ValueError, TypeError):
            pass
    return resolved


def _result_path(raw_path: str, resolved_path: str) -> str:
    """Return a workspace-relative path so the frontend can match it to wsSourcePath.

    Falls back to absolute path when the file is outside the workspace root.
    Frontend refresh (_refreshChangedFile) compares this to state.wsSourcePath which
    is always workspace-relative, so returning absolute paths breaks the match.
    """
    target = resolved_path or raw_path
    if not target:
        return raw_path
    ws_root = _get_workspace_root()
    try:
        rel = os.path.relpath(target, ws_root).replace("\\", "/")
        if not rel.startswith(".."):
            return rel
    except (ValueError, TypeError):
        pass
    return target


def _best_effort_backup(path: str) -> str:
    target = str(path or "").strip()
    if not target:
        return ""
    backup_path = target + ".bak"
    if os.path.exists(backup_path):
        _clear_readonly_attribute(backup_path)
    try:
        shutil.copy2(target, backup_path)
        return ""
    except OSError as primary_exc:
        try:
            _create_alternate_backup(target)
            return ""
        except OSError as fallback_exc:
            logger.warning(
                "[TaskTools] backup skipped for %s: primary=%s; fallback=%s",
                target,
                primary_exc,
                fallback_exc,
            )
            return f"无法创建备份 {os.path.basename(backup_path)}（{primary_exc}），已继续直接写入原文件。"


def _create_alternate_backup(target: str) -> str:
    directory = os.path.dirname(target) or os.getcwd()
    prefix = os.path.basename(target) + ".koto-backup-"
    handle: Optional[int] = None
    fallback_path = ""
    try:
        handle, fallback_path = tempfile.mkstemp(
            prefix=prefix, suffix=".bak", dir=directory
        )
        os.close(handle)
        handle = None
        shutil.copy2(target, fallback_path)
        return fallback_path
    except OSError:
        if handle is not None:
            try:
                os.close(handle)
            except OSError:
                pass
        if fallback_path and os.path.exists(fallback_path):
            try:
                os.unlink(fallback_path)
            except OSError:
                pass
        raise


def _clear_readonly_attribute(path: str) -> bool:
    target = str(path or "").strip()
    if not target or not os.path.exists(target):
        return False

    try:
        stats = os.stat(target)
    except OSError:
        return False

    readonly_attr = getattr(stat, "FILE_ATTRIBUTE_READONLY", 0)
    file_attrs = getattr(stats, "st_file_attributes", 0)
    is_readonly = bool(readonly_attr and file_attrs & readonly_attr) or not bool(
        stats.st_mode & stat.S_IWRITE
    )
    if not is_readonly:
        return False

    try:
        os.chmod(target, stats.st_mode | stat.S_IWRITE)
        return True
    except OSError as exc:
        logger.warning(
            "[TaskTools] failed to clear readonly attribute for %s: %s", target, exc
        )
        return False


def _nonwritable_target_message(target_path: str) -> str:
    name = (
        os.path.basename(str(target_path or "").strip())
        or str(target_path or "").strip()
        or "目标文件"
    )
    return (
        f"目标文件 {name} 当前不可写，无法写回原文件。"
        "Koto 已尝试清除只读属性；如果仍失败，请检查文件权限，"
        "或关闭可能占用该文件的 Koto 页签及其他程序后重试。"
    )


def _nonwritable_target_next_step(target_path: str) -> str:
    name = (
        os.path.basename(str(target_path or "").strip())
        or str(target_path or "").strip()
        or "目标文件"
    )
    return f"检查 {name} 的文件权限；如果文件正在被占用，" "请关闭相关 Koto 页签或其他程序后重新执行写回原文件。"


def _ensure_existing_file_writable(path: str, *, label: str = "目标文件") -> str:
    target = str(path or "").strip()
    if not target or not os.path.exists(target):
        return ""
    if _clear_readonly_attribute(target):
        name = os.path.basename(target) or target
        return f"检测到{label} {name} 为只读，已自动移除只读属性后继续写入。"
    return ""


def _merge_warnings(*parts: Any) -> str:
    cleaned = [str(part).strip() for part in parts if str(part or "").strip()]
    return "；".join(cleaned)


def _save_via_temp_file(save_func: Any, target_path: str, *, suffix: str) -> None:
    target = str(target_path or "").strip()
    if not target:
        raise ValueError("target_path is required")

    directory = os.path.dirname(target) or os.getcwd()
    temp_handle = tempfile.NamedTemporaryFile(
        delete=False, suffix=suffix or ".tmp", dir=directory
    )
    temp_path = temp_handle.name
    temp_handle.close()
    try:
        save_func(temp_path)
        if os.path.exists(target):
            _clear_readonly_attribute(target)
        try:
            os.replace(temp_path, target)
        except PermissionError as exc:
            raise PermissionError(_nonwritable_target_message(target)) from exc
    finally:
        if os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except OSError:
                pass


def _save_docx_via_temp_file(document: Any, target_path: str) -> None:
    _save_via_temp_file(document.save, target_path, suffix=".docx")


def _write_bytes_via_temp_file(data: bytes, target_path: str, *, suffix: str) -> None:
    def _write_temp(temp_path: str) -> None:
        with open(temp_path, "wb") as handle:
            handle.write(data)

    _save_via_temp_file(_write_temp, target_path, suffix=suffix)


def _save_workbook_via_temp_file(workbook: Any, target_path: str) -> None:
    suffix = Path(str(target_path or "")).suffix or ".xlsx"
    _save_via_temp_file(workbook.save, target_path, suffix=suffix)


def _save_pptx_via_temp_file(presentation: Any, target_path: str) -> None:
    suffix = Path(str(target_path or "")).suffix or ".pptx"
    _save_via_temp_file(presentation.save, target_path, suffix=suffix)


def _next_available_docx_copy_path(target_path: str) -> str:
    target = str(target_path or "").strip()
    if not target:
        raise ValueError("target_path is required")

    directory = os.path.dirname(target) or os.getcwd()
    target_path_obj = Path(target)
    stem = target_path_obj.stem or "document"
    suffix = target_path_obj.suffix or ".docx"
    for index in range(1, 1000):
        copy_suffix = ".koto-copy" if index == 1 else f".koto-copy-{index}"
        candidate = os.path.join(directory, f"{stem}{copy_suffix}{suffix}")
        if not os.path.exists(candidate):
            return candidate
    raise RuntimeError(f"Unable to allocate fallback copy path for {target}")


def _normalize_text_limit(max_chars: Any, default: int) -> int:
    try:
        value = int(max_chars)
    except (TypeError, ValueError):
        value = default
    return min(max(_TEXT_LIMIT_MIN, value), _TEXT_LIMIT_MAX)


def _normalize_positive_int(value: Any, *, default: int, upper: int) -> int:
    try:
        if isinstance(value, str):
            value = value.strip()
        parsed = int(float(value))
    except (TypeError, ValueError):
        parsed = default
    if parsed < 1:
        parsed = default
    return min(parsed, upper)


def _normalize_positive_float(value: Any, *, default: float, upper: float) -> float:
    try:
        if isinstance(value, str):
            value = value.strip()
        parsed = float(value)
    except (TypeError, ValueError):
        parsed = default
    if parsed <= 0:
        parsed = default
    return min(parsed, upper)


_GENERIC_SHEET_NAME_GUESSES = {"sheet", "sheet1", "工作表", "工作表1"}
_YEAR_HEADER_RE = re.compile(
    r"(?<!\d)(20\d{2}(?:\s*[A-Za-z]{0,4})?)(?!\d)", re.IGNORECASE
)
_FINANCIAL_STATEMENT_PATTERNS = {
    "profit_and_loss": (
        re.compile(
            r"(?:^|\b)(p\s*&\s*l|profit\s*&?\s*loss|income\s*statement)(?:\b|$)",
            re.IGNORECASE,
        ),
        re.compile(r"利润|损益|收入成本", re.IGNORECASE),
    ),
    "balance_sheet": (
        re.compile(
            r"(?:^|\b)(balance\s*sheet|statement\s*of\s*financial\s*position|bs)(?:\b|$)",
            re.IGNORECASE,
        ),
        re.compile(r"资产负债", re.IGNORECASE),
    ),
    "cash_flow": (
        re.compile(r"(?:^|\b)(cash\s*flow|cashflow|cf)(?:\b|$)", re.IGNORECASE),
        re.compile(r"现金流", re.IGNORECASE),
    ),
}
_HIGH_PRIORITY_FINANCIAL_LABEL_HINTS = (
    "税",
    "所得税",
    "净利润",
    "现金",
    "资产",
    "负债",
    "收入",
    "成本",
    "毛利",
)


def _build_workbook_structure_payload(
    resolved_path: str,
    *,
    sample_rows_per_sheet: int,
    max_formula_examples_per_sheet: int,
) -> Dict[str, Any]:
    import openpyxl

    wb_formula = openpyxl.load_workbook(resolved_path, data_only=False)
    wb_values = openpyxl.load_workbook(resolved_path, data_only=True)
    try:
        value_sheets = {sheet.title: sheet for sheet in wb_values.worksheets}
        external_link_targets = _extract_external_link_targets(wb_formula)
        sheet_summaries: List[Dict[str, Any]] = []
        total_formulas = 0
        total_external_formula_refs = 0
        statement_presence = {
            "profit_and_loss": {"present": False, "matched_sheet": ""},
            "balance_sheet": {"present": False, "matched_sheet": ""},
            "cash_flow": {"present": False, "matched_sheet": ""},
        }

        for sheet in wb_formula.worksheets:
            value_sheet = value_sheets.get(sheet.title, sheet)
            year_header = _detect_year_header(value_sheet)
            formula_info = _collect_formula_examples(
                sheet, max_formula_examples=max_formula_examples_per_sheet
            )
            total_formulas += int(formula_info["formula_count"])
            total_external_formula_refs += int(formula_info["external_formula_count"])

            for statement_key in statement_presence:
                if not statement_presence[statement_key][
                    "present"
                ] and _sheet_matches_statement(sheet.title, statement_key):
                    statement_presence[statement_key] = {
                        "present": True,
                        "matched_sheet": sheet.title,
                    }

            summary: Dict[str, Any] = {
                "name": sheet.title,
                "max_row": int(getattr(sheet, "max_row", 0) or 0),
                "max_column": int(getattr(sheet, "max_column", 0) or 0),
                "hidden": str(getattr(sheet, "sheet_state", "visible") or "visible")
                != "visible",
                "sample_rows": _sample_sheet_rows(
                    value_sheet, max_rows=sample_rows_per_sheet
                ),
                "formula_count": int(formula_info["formula_count"]),
                "external_formula_count": int(formula_info["external_formula_count"]),
            }
            if formula_info["formula_examples"]:
                summary["formula_examples"] = formula_info["formula_examples"]
            if formula_info["external_formula_examples"]:
                summary["external_formula_examples"] = formula_info[
                    "external_formula_examples"
                ]
            if year_header:
                summary["year_header"] = year_header
            sheet_summaries.append(summary)

        summary_parts = [f"已检查 {len(sheet_summaries)} 个工作表"]
        if external_link_targets:
            summary_parts.append(f"发现 {len(external_link_targets)} 个外部链接")
        if total_formulas:
            summary_parts.append(f"共扫描到 {total_formulas} 个公式单元格")
        return {
            "sheet_count": len(sheet_summaries),
            "sheet_names": [sheet["name"] for sheet in sheet_summaries],
            "sheets": sheet_summaries,
            "external_link_targets": external_link_targets,
            "external_link_count": len(external_link_targets),
            "total_formula_cells": total_formulas,
            "total_external_formula_refs": total_external_formula_refs,
            "statement_presence": statement_presence,
            "summary": "；".join(summary_parts) + "。",
        }
    finally:
        wb_formula.close()
        wb_values.close()


def _severity_for_financial_label(label: str) -> str:
    return _xlsx_severity_for_financial_label(
        label,
        _HIGH_PRIORITY_FINANCIAL_LABEL_HINTS,
    )


def _detect_financial_series_gap_findings(
    resolved_path: str,
    *,
    max_findings: int,
    max_rows_per_sheet: int = 240,
) -> List[Dict[str, Any]]:
    import openpyxl

    try:
        from openpyxl.utils import get_column_letter
    except Exception:  # pragma: no cover - openpyxl is already required by caller
        return []

    wb_formula = openpyxl.load_workbook(resolved_path, data_only=False)
    wb_values = openpyxl.load_workbook(resolved_path, data_only=True)
    findings: List[Dict[str, Any]] = []
    try:
        value_sheets = {sheet.title: sheet for sheet in wb_values.worksheets}
        for formula_sheet in wb_formula.worksheets:
            value_sheet = value_sheets.get(formula_sheet.title, formula_sheet)
            year_header = _detect_year_header(value_sheet)
            if not year_header:
                continue
            year_columns = [int(item["index"]) for item in year_header["columns"]]
            headers = [str(item["header"]) for item in year_header["columns"]]
            first_year_column = min(year_columns)
            row_limit = min(
                int(getattr(formula_sheet, "max_row", 0) or 0),
                year_header["row"] + max_rows_per_sheet,
            )
            for row_index in range(int(year_header["row"]) + 1, row_limit + 1):
                label = _row_label_for_year_series(
                    value_sheet, row_index, before_column=first_year_column
                )
                if not label:
                    continue
                populated_positions: List[int] = []
                series: List[Dict[str, Any]] = []
                for offset, column_index in enumerate(year_columns):
                    formula_cell = formula_sheet.cell(
                        row=row_index, column=column_index
                    )
                    value_cell = value_sheet.cell(row=row_index, column=column_index)
                    formula_text = (
                        str(formula_cell.value or "")
                        if getattr(formula_cell, "data_type", "") == "f"
                        else ""
                    )
                    display_value = _display_series_value(
                        value_cell.value, formula_text
                    )
                    populated = display_value not in (None, "")
                    if populated:
                        populated_positions.append(offset)
                    series.append(
                        {
                            "header": headers[offset],
                            "column": get_column_letter(column_index),
                            "value": display_value,
                        }
                    )
                if len(populated_positions) < 2:
                    continue
                first_populated = min(populated_positions)
                last_populated = max(populated_positions)
                gap_positions = [
                    index
                    for index in range(first_populated, last_populated + 1)
                    if index not in populated_positions
                ]
                if not gap_positions:
                    continue
                findings.append(
                    {
                        "severity": _severity_for_financial_label(label),
                        "type": "year_series_gap",
                        "sheet": formula_sheet.title,
                        "row": row_index,
                        "label": label,
                        "message": f"行“{label}”在年份序列中存在空档：{', '.join(headers[index] for index in gap_positions)} 为空，但前后年份仍有数据或公式。",
                        "evidence": {
                            "header_row": year_header["row"],
                            "series": series,
                        },
                    }
                )
                if len(findings) >= max_findings:
                    return findings
        return findings
    finally:
        wb_formula.close()
        wb_values.close()


def _read_pdf_excerpt(
    path: str,
    *,
    max_chars: int,
    start_page: int = 1,
    end_page: int = 0,
    allow_full_fallback: bool = True,
) -> str:
    return _pdf_read_excerpt(
        path,
        max_chars=max_chars,
        start_page=start_page,
        end_page=end_page,
        allow_full_fallback=allow_full_fallback,
        logger=logger,
    )


def _read_pdf_letter_window(
    path: str,
    *,
    max_chars: int,
    start_letter: int,
    end_letter: int,
) -> str:
    return _pdf_read_letter_window(
        path,
        max_chars=max_chars,
        start_letter=start_letter,
        end_letter=end_letter,
        read_excerpt=_read_pdf_excerpt,
    )


def _success_result(
    path: str,
    *,
    operation: str,
    summary: str,
    file_type: str = "",
    change_type: str = "modify",
    preview: str = "",
    focus: bool = False,
    **extra: Any,
) -> str:
    summary_code = str(
        extra.pop(
            "summary_code",
            "CREATE_OK" if str(change_type or "").lower() == "create" else "WRITE_OK",
        )
        or ""
    )
    payload: Dict[str, Any] = {
        "success": True,
        "path": path,
        "file_type": file_type or Path(str(path)).suffix.lstrip(".").lower(),
        "change_type": change_type,
        "operation": operation,
        "summary_code": summary_code,
        "summary": summary,
        "preview": preview[:400],
        "focus": focus,
    }
    payload.update(extra)
    return json.dumps(payload, ensure_ascii=False, default=str)


def _file_task_diff(
    kind: str,
    items: List[Dict[str, Any]],
    *,
    limit: int = 80,
) -> Dict[str, Any]:
    normalized = [dict(item) for item in items if isinstance(item, dict)]
    return {
        "kind": kind,
        "items": normalized[:limit],
        "changed_count": len(normalized),
        "truncated": len(normalized) > limit,
    }


def _blocked_write_result(
    path: str,
    *,
    summary: str,
    suggested_next_step: str = "",
    **extra: Any,
) -> str:
    payload: Dict[str, Any] = {
        "success": False,
        "path": path,
        "status": "write_blocked",
        "summary_code": "WRITE_BLOCKED",
        "summary": summary,
        "error": summary,
    }
    if suggested_next_step:
        payload["suggested_next_step"] = suggested_next_step
    payload.update(extra)
    return json.dumps(payload, ensure_ascii=False, default=str)


# ══════════════════════════════════════════════════════════════
# Tool implementations
# ══════════════════════════════════════════════════════════════


def read_docx_content(path: str, max_chars: int = _TEXT_LIMIT_DOCX_DEFAULT) -> str:
    """Read DOCX paragraphs as structured JSON.

    Returns JSON with paragraphs and tables.
    """
    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)

    max_chars = _normalize_text_limit(max_chars, _TEXT_LIMIT_DOCX_DEFAULT)

    try:
        from docx import Document as DocxDocument

        doc = DocxDocument(resolved)
        paragraphs = []
        tables = []
        total = 0
        for p in doc.paragraphs:
            if total >= max_chars:
                break
            paragraphs.append(
                {
                    "text": p.text,
                    "style": p.style.name if p.style else "",
                }
            )
            total += len(p.text)

        for table_index, table in enumerate(doc.tables):
            if total >= max_chars:
                break
            table_rows = []
            for row_index, row in enumerate(table.rows):
                if row_index >= 20 or total >= max_chars:
                    break
                values = [cell.text for cell in row.cells]
                table_rows.append(values)
                total += sum(len(v) for v in values)
            tables.append(
                {
                    "index": table_index,
                    "rows": table_rows,
                    "row_count": len(table.rows),
                    "column_count": len(table.columns),
                }
            )

        return json.dumps(
            {
                "paragraphs": paragraphs,
                "tables": tables,
                "total_paragraphs": len(doc.paragraphs),
                "total_tables": len(doc.tables),
            },
            ensure_ascii=False,
        )
    except ImportError:
        return json.dumps({"error": "python-docx not installed"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def _read_docx_paragraph_window(
    path: str,
    *,
    max_chars: int,
    start: int,
    end: int,
) -> str:
    from docx import Document as DocxDocument

    doc = DocxDocument(path)
    start = max(1, int(start or 1))
    end = max(start, int(end or start))
    parts: list[str] = []
    total = 0
    for index, paragraph in enumerate(doc.paragraphs, start=1):
        if index < start:
            continue
        if index > end:
            break
        text = str(paragraph.text or "").strip()
        if not text:
            continue
        block = f"[Paragraph {index}]\n{text}"
        parts.append(block)
        total += len(block)
        if total >= max_chars:
            break
    return "\n\n".join(parts)


def _read_pptx_slide_window(
    path: str,
    *,
    max_chars: int,
    start: int,
    end: int,
) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    start = max(1, int(start or 1))
    end = max(start, int(end or start))
    parts: list[str] = []
    total = 0
    for index, slide in enumerate(presentation.slides, start=1):
        if index < start:
            continue
        if index > end:
            break
        lines: list[str] = []
        for shape in slide.shapes:
            text = str(getattr(shape, "text", "") or "").strip()
            if text:
                lines.append(text)
        if not lines:
            continue
        block = f"[Slide {index}]\n" + "\n".join(lines)
        parts.append(block)
        total += len(block)
        if total >= max_chars:
            break
    return "\n\n".join(parts)


def _read_xlsx_sheet_window(
    path: str,
    *,
    max_chars: int,
    sheet_index: int,
) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        sheet_names = list(wb.sheetnames or [])
        if not sheet_names:
            return ""
        safe_index = max(0, min(int(sheet_index or 0), len(sheet_names) - 1))
        ws = wb[sheet_names[safe_index]]
        parts = [f"[Sheet {safe_index + 1}: {ws.title}]"]
        total = len(parts[0])
        for row_index, row in enumerate(ws.iter_rows(values_only=True), start=1):
            values = ["" if value is None else str(value) for value in row]
            if not any(value.strip() for value in values):
                continue
            line = f"R{row_index}: " + " | ".join(values)
            parts.append(line)
            total += len(line)
            if total >= max_chars:
                break
        return "\n".join(parts)
    finally:
        wb.close()


def parse_file_to_text(
    path: str,
    max_chars: int = _TEXT_LIMIT_DEFAULT,
    start_page: int = 1,
    end_page: int = 0,
    window_unit: str = "",
    start: int = 0,
    end: int = 0,
    sheet_index: int = -1,
) -> str:
    """Parse any supported file to plain text (DOCX/XLSX/PPTX/PDF/TXT/CSV)."""
    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)

    try:
        from app.core.workflow_engine import parse_source_file

        suffix = Path(resolved).suffix.lower()
        if suffix == ".pdf" and window_unit == "pdf_letter":
            text = _read_pdf_letter_window(
                resolved,
                max_chars=max_chars,
                start_letter=start or 1,
                end_letter=end or start or 1,
            )
        elif suffix == ".pdf":
            text = _read_pdf_excerpt(
                resolved,
                max_chars=max_chars,
                start_page=start_page,
                end_page=end_page,
            )
        elif suffix in {".doc", ".docx"} and window_unit == "paragraph":
            text = _read_docx_paragraph_window(
                resolved,
                max_chars=max_chars,
                start=start or 1,
                end=end or start or 1,
            )
        elif suffix in {".ppt", ".pptx"} and window_unit == "slide":
            text = _read_pptx_slide_window(
                resolved,
                max_chars=max_chars,
                start=start or 1,
                end=end or start or 1,
            )
        elif suffix in {".xls", ".xlsx", ".xlsm", ".csv"} and window_unit == "sheet":
            text = _read_xlsx_sheet_window(
                resolved,
                max_chars=max_chars,
                sheet_index=sheet_index,
            )
        else:
            text = parse_source_file(resolved)
        if not text.strip():
            return f"(File parsed but no text content: {path})"
        return text[:max_chars]
    except Exception as e:
        return f"Error parsing file: {e}"


def _resolve_task_file_entries(
    task_files: Optional[List[Dict[str, str]]],
) -> List[Dict[str, str]]:
    """Resolve task file metadata into concrete files that can be staged."""
    resolved_entries: List[Dict[str, str]] = []
    seen: set[str] = set()

    for item in task_files or []:
        if not isinstance(item, dict):
            continue

        raw_path = str(item.get("path") or "").strip()
        raw_name = str(item.get("name") or "").strip()
        resolved_path = _resolve_path(raw_path or raw_name)
        if not resolved_path or not os.path.isfile(resolved_path):
            continue

        key = os.path.normcase(os.path.abspath(resolved_path))
        if key in seen:
            continue
        seen.add(key)

        resolved_entries.append(
            {
                "display_name": raw_name or os.path.basename(resolved_path),
                "source_path": resolved_path,
                "source_fingerprint_initial": _fingerprint_file(resolved_path),
            }
        )

    return resolved_entries


def _fingerprint_file(path: str) -> Dict[str, Any]:
    """Capture a stable fingerprint for later change detection."""
    try:
        stat = os.stat(path)
        digest = hashlib.sha1(usedforsecurity=False)
        with open(path, "rb") as handle:
            while True:
                chunk = handle.read(1024 * 1024)
                if not chunk:
                    break
                digest.update(chunk)
        return {
            "size": stat.st_size,
            "mtime_ns": getattr(
                stat, "st_mtime_ns", int(stat.st_mtime * 1_000_000_000)
            ),
            "sha1": digest.hexdigest(),
        }
    except OSError:
        return {}


def _fingerprint_changed(path: str, fingerprint: Dict[str, Any]) -> bool:
    """Return True when the file differs from its captured fingerprint."""
    if not path or not fingerprint or not os.path.isfile(path):
        return False
    return _fingerprint_file(path) != fingerprint


def _unique_staged_name(name: str, used_names: set[str]) -> str:
    """Return a unique basename for files mirrored into the sandbox workdir."""
    candidate = os.path.basename(name or "") or "task_file"
    candidate = (
        re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", candidate).strip(" .") or "task_file"
    )
    stem, ext = os.path.splitext(candidate)
    normalized = candidate.lower()
    index = 2
    while normalized in used_names:
        candidate = f"{stem}_{index}{ext}"
        normalized = candidate.lower()
        index += 1
    used_names.add(normalized)
    return candidate


def _stage_task_files_for_sandbox(
    resolved_entries: List[Dict[str, str]], sandbox_dir: str
) -> List[Dict[str, str]]:
    """Copy resolved task files into the sandbox workdir for basename-based access."""
    staged_entries: List[Dict[str, str]] = []
    used_names: set[str] = set()

    for entry in resolved_entries:
        staged_name = _unique_staged_name(entry["display_name"], used_names)
        staged_path = os.path.join(sandbox_dir, staged_name)
        shutil.copy2(entry["source_path"], staged_path)
        staged_entries.append(
            {
                **entry,
                "staged_name": staged_name,
                "staged_path": staged_path,
                "staged_mtime_initial": os.stat(staged_path).st_mtime,
                "staged_fingerprint_initial": _fingerprint_file(staged_path),
            }
        )

    return staged_entries


def _prepend_task_file_context(
    code: str,
    staged_entries: List[Dict[str, str]],
    *,
    output_dir: str = "",
) -> str:
    """Expose task file paths to sandbox code and keep basename access working."""
    absolute_paths = {
        entry["display_name"]: entry["source_path"] for entry in staged_entries
    }
    staged_paths = {
        entry["display_name"]: entry["staged_path"] for entry in staged_entries
    }
    staged_names = [entry["staged_name"] for entry in staged_entries]

    workspace_root = _get_workspace_root()
    normalized_output_dir = _normalize_workspace_relative_path(output_dir)
    output_dir_abs = ""
    if normalized_output_dir:
        output_dir_abs = str(
            (Path(workspace_root) / Path(*normalized_output_dir.split("/"))).resolve()
        )
    preamble = (
        "# Attached task files are mirrored into the sandbox working directory.\n"
        f"TASK_WORKSPACE_ROOT = {json.dumps(workspace_root, ensure_ascii=False)}\n"
        f"TASK_OUTPUT_DIR = {json.dumps(output_dir_abs, ensure_ascii=False)}\n"
        f"TASK_SANDBOX_FILE_PATHS = {json.dumps(staged_paths, ensure_ascii=False)}\n"
        f"TASK_FILE_PATHS = {json.dumps(absolute_paths, ensure_ascii=False)}\n"
        f"TASK_SANDBOX_FILES = {json.dumps(staged_names, ensure_ascii=False)}\n"
        "# Prefer TASK_SANDBOX_FILE_PATHS[...] for opening and editing attached files.\n"
        "# If TASK_OUTPUT_DIR is not empty, create new relative output files there.\n"
        "# After modifying an attached file, print: KOTO_MODIFIED:<sandbox_absolute_path>\n"
        "# Koto will sync the staged edit back to the source file automatically.\n"
        "# After creating a file in the workspace, print: KOTO_CREATED:<absolute_path>\n"
        "# TASK_FILE_PATHS is retained only as a compatibility fallback for legacy flows.\n"
        "# e.g. print('KOTO_MODIFIED:' + TASK_SANDBOX_FILE_PATHS['report.docx'])\n"
        "# e.g. print('KOTO_CREATED:' + output_path)\n\n"
    )
    return preamble + code


def _normalize_workspace_relative_path(path: str) -> str:
    candidate = str(path or "").strip().replace("\\", "/")
    if not candidate:
        return ""
    if os.path.isabs(candidate):
        try:
            candidate = str(
                Path(candidate).resolve().relative_to(Path(_get_workspace_root()).resolve())
            ).replace("\\", "/")
        except (OSError, ValueError):
            return ""
    candidate = candidate.lstrip("/")
    if candidate.startswith("workspace/"):
        candidate = candidate[len("workspace/") :]
    parts = [part for part in candidate.split("/") if part and part not in {".", ".."}]
    if not parts:
        return ""
    return "/".join(parts)


def _target_output_candidates(sandbox_dir: str, target_path: str) -> List[Path]:
    rel_target = _normalize_workspace_relative_path(target_path)
    if not rel_target:
        return []
    sandbox_root = Path(sandbox_dir)
    target_rel_path = Path(*rel_target.split("/"))
    candidates = [
        sandbox_root / target_rel_path,
        sandbox_root / "workspace" / target_rel_path,
        sandbox_root / target_rel_path.name,
    ]
    unique: List[Path] = []
    seen: set[str] = set()
    for candidate in candidates:
        key = os.path.normcase(str(candidate.resolve(strict=False)))
        if key not in seen:
            seen.add(key)
            unique.append(candidate)
    return unique


def _workspace_target_path(target_path: str) -> Optional[Path]:
    rel_target = _normalize_workspace_relative_path(target_path)
    if not rel_target:
        return None
    workspace_root = Path(_get_workspace_root()).resolve()
    real_target = (workspace_root / Path(*rel_target.split("/"))).resolve()
    try:
        real_target.relative_to(workspace_root)
    except ValueError:
        return None
    return real_target


def _prepare_sandbox_target_paths(sandbox_dir: str, target_path: str) -> None:
    for candidate in _target_output_candidates(sandbox_dir, target_path):
        try:
            candidate.parent.mkdir(parents=True, exist_ok=True)
        except OSError:
            continue


def _sync_target_outputs_from_sandbox(
    sandbox_dir: str,
    result: Dict[str, Any],
    *,
    target_path: str,
    target_existed: bool,
) -> Dict[str, Any]:
    rel_target = _normalize_workspace_relative_path(target_path)
    if not rel_target:
        return result
    real_target = _workspace_target_path(rel_target)
    if real_target is None:
        return result

    source: Optional[Path] = None
    for candidate in _target_output_candidates(sandbox_dir, rel_target):
        try:
            if candidate.is_file():
                source = candidate
                break
        except OSError:
            continue
    if source is None:
        return result

    try:
        if real_target.exists() and filecmp.cmp(str(source), str(real_target), shallow=False):
            copied = False
        else:
            real_target.parent.mkdir(parents=True, exist_ok=True)
            if real_target.exists():
                _clear_readonly_attribute(str(real_target))
            shutil.copy2(source, real_target)
            copied = True
    except OSError:
        return result

    updated = dict(result)
    stdout = str(updated.get("stdout") or "")
    markers = _parse_koto_file_markers(stdout)
    marker_kind = "modified" if target_existed else "created"
    existing = {
        os.path.normcase(os.path.abspath(path))
        for path in markers.get(marker_kind, [])
    }
    norm_target = os.path.normcase(os.path.abspath(str(real_target)))
    if copied and norm_target not in existing:
        marker_name = "KOTO_MODIFIED" if target_existed else "KOTO_CREATED"
        if stdout and not stdout.endswith("\n"):
            stdout += "\n"
        stdout += f"{marker_name}:{real_target}"
        updated["stdout"] = stdout
    return updated


_WORKSPACE_CREATED_FILE_EXTS = {
    ".csv",
    ".doc",
    ".docx",
    ".gif",
    ".jpeg",
    ".jpg",
    ".json",
    ".md",
    ".pdf",
    ".png",
    ".ppt",
    ".pptx",
    ".svg",
    ".txt",
    ".webp",
    ".xls",
    ".xlsm",
    ".xlsx",
}


def _copy_sandbox_workspace_file_to_workspace(path: Path, workspace_tmp: Path) -> str:
    """Copy sandbox/workspace/... output back to the real workspace root."""
    try:
        source = path.resolve()
        rel_path = source.relative_to(workspace_tmp.resolve())
    except (OSError, ValueError):
        return ""
    if not rel_path.parts or ".." in rel_path.parts:
        return ""
    target = (Path(_get_workspace_root()) / rel_path).resolve()
    try:
        target.relative_to(Path(_get_workspace_root()).resolve())
    except ValueError:
        return ""
    target.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, target)
    return str(target)


def _sync_created_workspace_files_from_sandbox(
    sandbox_dir: str,
    result: Dict[str, Any],
) -> Dict[str, Any]:
    """Preserve files created under sandbox/workspace/... before temp cleanup."""
    workspace_tmp = Path(sandbox_dir) / "workspace"
    if not workspace_tmp.is_dir():
        return result

    discovered: List[str] = []
    for path in workspace_tmp.rglob("*"):
        if not path.is_file():
            continue
        if path.suffix.lower() not in _WORKSPACE_CREATED_FILE_EXTS:
            continue
        copied = _copy_sandbox_workspace_file_to_workspace(path, workspace_tmp)
        if copied:
            discovered.append(copied)

    if not discovered:
        return result

    updated = dict(result)
    stdout = str(updated.get("stdout") or "")
    existing = {
        os.path.normcase(os.path.abspath(path))
        for path in _parse_koto_file_markers(stdout).get("created", [])
    }
    marker_lines = []
    for path in discovered:
        key = os.path.normcase(os.path.abspath(path))
        if key not in existing:
            marker_lines.append(f"KOTO_CREATED:{path}")
            existing.add(key)
    if marker_lines:
        updated["stdout"] = "\n".join(
            part for part in [stdout.rstrip(), *marker_lines] if part
        )
    return updated


def _workspace_output_dir_path(output_dir: str) -> Optional[Path]:
    rel_dir = _normalize_workspace_relative_path(output_dir)
    if not rel_dir:
        return None
    workspace_root = Path(_get_workspace_root()).resolve()
    target_dir = (workspace_root / Path(*rel_dir.split("/"))).resolve()
    try:
        target_dir.relative_to(workspace_root)
    except ValueError:
        return None
    return target_dir


def _relocate_root_created_files_to_output_dir(
    result: Dict[str, Any],
    *,
    output_dir: str,
) -> Dict[str, Any]:
    target_dir = _workspace_output_dir_path(output_dir)
    if target_dir is None:
        return result
    workspace_root = Path(_get_workspace_root()).resolve()
    markers = _parse_koto_file_markers(str(result.get("stdout") or ""))
    replacements: Dict[str, str] = {}
    for raw_path in markers.get("created", []):
        try:
            source = Path(raw_path).resolve()
            source.relative_to(workspace_root)
        except (OSError, ValueError):
            continue
        if source.parent != workspace_root:
            continue
        if source.suffix.lower() not in _WORKSPACE_CREATED_FILE_EXTS:
            continue
        target = (target_dir / source.name).resolve()
        try:
            target.relative_to(workspace_root)
        except ValueError:
            continue
        if source == target or target.exists() or not source.is_file():
            continue
        try:
            target.parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(source), str(target))
        except OSError:
            continue
        replacements[str(source)] = str(target)

    if not replacements:
        return result
    updated = dict(result)
    stdout = str(updated.get("stdout") or "")
    for source, target in replacements.items():
        stdout = stdout.replace(source, target)
    updated["stdout"] = stdout
    return updated


def _parse_koto_file_markers(stdout: str) -> Dict[str, List[str]]:
    """Extract KOTO_CREATED/KOTO_MODIFIED markers printed by sandbox code."""
    created: List[str] = []
    modified: List[str] = []
    for line in (stdout or "").splitlines():
        line = line.strip()
        if line.startswith("KOTO_CREATED:"):
            candidate = line[len("KOTO_CREATED:") :].strip()
            if candidate and os.path.isabs(candidate):
                created.append(candidate)
        elif line.startswith("KOTO_MODIFIED:"):
            candidate = line[len("KOTO_MODIFIED:") :].strip()
            if candidate and os.path.isabs(candidate):
                modified.append(candidate)
    return {"created": created, "modified": modified}


def _canonicalize_koto_markers(
    stdout: str, staged_entries: List[Dict[str, str]]
) -> str:
    """Rewrite staged-file KOTO_MODIFIED markers to their source paths and dedupe them."""
    if not stdout:
        return stdout or ""

    staged_to_source = {
        os.path.normcase(os.path.abspath(str(entry.get("staged_path") or ""))): str(
            entry.get("source_path") or ""
        )
        for entry in staged_entries
        if entry.get("staged_path") and entry.get("source_path")
    }
    seen_created: set[str] = set()
    seen_modified: set[str] = set()
    normalized_lines: List[str] = []

    for raw_line in stdout.splitlines():
        stripped = raw_line.strip()
        if stripped.startswith("KOTO_CREATED:"):
            candidate = stripped[len("KOTO_CREATED:") :].strip()
            if candidate and os.path.isabs(candidate):
                normalized = os.path.normcase(os.path.abspath(candidate))
                if normalized not in seen_created:
                    seen_created.add(normalized)
                    normalized_lines.append(
                        f"KOTO_CREATED:{os.path.abspath(candidate)}"
                    )
            continue

        if stripped.startswith("KOTO_MODIFIED:"):
            candidate = stripped[len("KOTO_MODIFIED:") :].strip()
            if candidate and os.path.isabs(candidate):
                normalized_candidate = os.path.normcase(os.path.abspath(candidate))
                rewritten = staged_to_source.get(normalized_candidate, candidate)
                rewritten_abs = os.path.abspath(rewritten)
                normalized_rewritten = os.path.normcase(rewritten_abs)
                if normalized_rewritten not in seen_modified:
                    seen_modified.add(normalized_rewritten)
                    normalized_lines.append(f"KOTO_MODIFIED:{rewritten_abs}")
            continue

        normalized_lines.append(raw_line)

    return "\n".join(normalized_lines)


def _sync_staged_files_to_source(
    staged_entries: List[Dict[str, str]],
    result: Dict[str, Any],
) -> Dict[str, Any]:
    """Copy staged-file edits back to source files and emit KOTO_MODIFIED markers."""
    stdout = _canonicalize_koto_markers(str(result.get("stdout") or ""), staged_entries)
    existing_markers = _parse_koto_file_markers(stdout)
    already_reported = {
        os.path.normcase(os.path.abspath(path))
        for path in existing_markers.get("modified", [])
    }
    extra_modified: List[str] = []

    for entry in staged_entries:
        staged_path = str(entry.get("staged_path") or "")
        source_path = str(entry.get("source_path") or "")
        if not staged_path or not source_path or not os.path.isfile(staged_path):
            continue

        staged_changed = _fingerprint_changed(
            staged_path,
            entry.get("staged_fingerprint_initial") or {},
        )
        source_changed = _fingerprint_changed(
            source_path,
            entry.get("source_fingerprint_initial") or {},
        )
        if not staged_changed and not source_changed:
            continue

        files_match_after_run = False
        if staged_changed:
            try:
                files_match_after_run = filecmp.cmp(
                    staged_path, source_path, shallow=False
                )
            except OSError:
                files_match_after_run = False

        if staged_changed and not source_changed:
            try:
                if os.path.exists(source_path):
                    _clear_readonly_attribute(source_path)
                shutil.copy2(staged_path, source_path)
                logger.info(
                    "[sandbox] Auto-synced modified staged file -> %s", source_path
                )
                source_changed = True
            except Exception as exc:
                logger.warning(
                    "[sandbox] Sync failed %s -> %s: %s", staged_path, source_path, exc
                )
                continue
        elif staged_changed and source_changed and not files_match_after_run:
            logger.warning(
                "[sandbox] Source and staged file both changed; keeping direct source version for %s",
                source_path,
            )

        norm_source = os.path.normcase(os.path.abspath(source_path))
        if source_changed and norm_source not in already_reported:
            extra_modified.append(source_path)

    if not extra_modified:
        merged_result = dict(result)
        merged_result["stdout"] = stdout
        return merged_result

    extra_lines = "\n".join(f"KOTO_MODIFIED:{path}" for path in extra_modified)
    merged_stdout = stdout
    if merged_stdout and not merged_stdout.endswith("\n"):
        merged_stdout += "\n"
    merged_stdout += extra_lines

    merged_result = dict(result)
    merged_result["stdout"] = merged_stdout
    return merged_result


def _format_sandbox_result(result: Dict[str, Any]) -> str:
    """Normalize sandbox execution output into a single text payload."""
    parts = []
    if result.get("stdout"):
        parts.append(result["stdout"])
    if result.get("stderr"):
        parts.append(f"[stderr] {result['stderr']}")
    generated_files = result.get("files") or result.get("images") or {}
    if generated_files:
        parts.append(f"[{len(generated_files)} image(s) generated]")
    if result.get("error"):
        parts.append(f"[error] {result['error']}")
    return "\n".join(parts) if parts else "(no output)"


def _load_sandbox_run_python():
    module_name = "app.core.sandbox"
    module = sys.modules.get(module_name)
    if module is not None and not isinstance(module, types.ModuleType):
        sys.modules.pop(module_name, None)

    module = importlib.import_module(module_name)
    run_python = getattr(module, "run_python", None)
    if not callable(run_python) or getattr(run_python, "__module__", module_name) != module_name:
        sys.modules.pop(module_name, None)
        module = importlib.import_module(module_name)
        run_python = getattr(module, "run_python", None)

    if not callable(run_python):
        raise RuntimeError("app.core.sandbox.run_python is unavailable")
    return run_python


class SandboxRunResult(dict):
    """Structured sandbox result with a marker-text compatibility view."""

    def _marker_text(self) -> str:
        parts: List[str] = []
        created = self.get(KOTO_CREATED_RESULT_KEY) or self.get(KOTO_CREATED_FALLBACK_KEY) or []
        modified = self.get(KOTO_MODIFIED_RESULT_KEY) or self.get(KOTO_MODIFIED_FALLBACK_KEY) or []

        if isinstance(created, list) and created:
            parts.append(
                KOTO_CREATED_RESULT_MARKER + json.dumps(created, ensure_ascii=False)
            )
        if isinstance(modified, list) and modified:
            parts.append(
                KOTO_MODIFIED_RESULT_MARKER + json.dumps(modified, ensure_ascii=False)
            )
        return "\n".join(parts)

    def as_text(self) -> str:
        text = str(self.get("summary") or "")
        marker_text = self._marker_text()
        if marker_text:
            if text and not text.endswith("\n"):
                text += "\n"
            text += marker_text
        return text

    def _legacy_marker_text(self) -> str:
        return self._marker_text()

    def as_legacy_text(self) -> str:
        return self.as_text()

    def __contains__(self, item: object) -> bool:
        if dict.__contains__(self, item):
            return True
        if isinstance(item, str):
            return item in self.as_text()
        return False

    def __str__(self) -> str:
        return self.as_text()


def run_python_in_sandbox(
    code: str,
    timeout: int = 30,
    task_files: Optional[List[Dict[str, str]]] = None,
    target_path: str = "",
    output_dir: str = "",
) -> Dict[str, Any]:
    """Execute Python code in the sandbox. Returns structured stdout/stderr/files.

    If the code prints ``KOTO_CREATED:<absolute_path>`` or
    ``KOTO_MODIFIED:<absolute_path>`` lines, those paths are returned as hidden
    structured fields so the file-task runtime can emit file_change events.
    """
    tmpdir: str | None = None
    normalized_timeout = _normalize_positive_int(timeout, default=30, upper=120)
    try:
        run_python = _load_sandbox_run_python()
        resolved_task_files = _resolve_task_file_entries(task_files)
        tmpdir = tempfile.mkdtemp(prefix="koto-task-")
        staged_entries: List[Dict[str, str]] = []
        if resolved_task_files:
            staged_entries = _stage_task_files_for_sandbox(resolved_task_files, tmpdir)
        resolved_target = _workspace_target_path(target_path)
        target_existed = bool(resolved_target and resolved_target.is_file())
        _prepare_sandbox_target_paths(tmpdir, target_path)
        prepared_code = _prepend_task_file_context(
            code,
            staged_entries,
            output_dir=output_dir,
        )
        result = run_python(prepared_code, timeout=normalized_timeout, work_dir=tmpdir)
        result = _sync_target_outputs_from_sandbox(
            tmpdir,
            result,
            target_path=target_path,
            target_existed=target_existed,
        )
        result = _sync_created_workspace_files_from_sandbox(tmpdir, result)
        result = _relocate_root_created_files_to_output_dir(
            result,
            output_dir=output_dir,
        )
        if staged_entries:
            result = _sync_staged_files_to_source(staged_entries, result)
        return _wrap_sandbox_result(result, output_dir=output_dir)
    except Exception as e:
        return SandboxRunResult(
            {
                "summary": f"Sandbox error: {e}",
                "stdout": "",
                "stderr": "",
                "error": str(e),
                "files": {},
                "_koto_created": [],
                "_koto_modified": [],
                "__koto_created__": [],
                "__koto_modified__": [],
            }
        )
    finally:
        if tmpdir and os.path.isdir(tmpdir):
            _cleanup_sandbox_tmpdir(tmpdir)


def _cleanup_sandbox_tmpdir(tmpdir: str) -> None:
    last_error: OSError | None = None
    for attempt in range(1, _SANDBOX_CLEANUP_RETRIES + 1):
        try:
            shutil.rmtree(tmpdir)
            return
        except OSError as exc:
            last_error = exc
            if attempt >= _SANDBOX_CLEANUP_RETRIES:
                break
            time.sleep(_SANDBOX_CLEANUP_RETRY_DELAY_SECONDS)
    if last_error is not None:
        # Try to quarantine the dir so it is not re-used accidentally (Windows
        # file-locks are a common cause of rmtree failure).
        try:
            dead_name = f".dead.{int(time.time())}.{os.path.basename(tmpdir)}"
            dead_path = os.path.join(os.path.dirname(tmpdir), dead_name)
            os.rename(tmpdir, dead_path)
            logger.warning(
                "[task_tools] sandbox tmpdir cleanup failed; quarantined as %s: %s",
                dead_path,
                last_error,
            )
        except OSError:
            logger.warning(
                "[task_tools] temp sandbox cleanup skipped for %s: %s",
                tmpdir,
                last_error,
            )


def _relocate_materialized_files_to_output_dir(
    materialized_files: Dict[str, str],
    *,
    output_dir: str,
) -> Dict[str, str]:
    target_dir = _workspace_output_dir_path(output_dir)
    if target_dir is None or not materialized_files:
        return materialized_files
    workspace_root = Path(_get_workspace_root()).resolve()
    relocated: Dict[str, str] = {}
    changed = False
    for name, raw_path in materialized_files.items():
        path_text = str(raw_path or "").strip()
        if not path_text:
            continue
        try:
            source = Path(path_text).resolve()
        except OSError:
            relocated[name] = path_text
            continue
        if (
            source.parent != workspace_root
            or source.suffix.lower() not in _WORKSPACE_CREATED_FILE_EXTS
        ):
            relocated[name] = path_text
            continue
        target = (target_dir / source.name).resolve()
        try:
            target.relative_to(workspace_root)
        except ValueError:
            relocated[name] = path_text
            continue
        if source == target:
            relocated[name] = str(source)
            continue
        if target.exists() or not source.is_file():
            relocated[name] = path_text
            continue
        try:
            target.parent.mkdir(parents=True, exist_ok=True)
            shutil.move(str(source), str(target))
        except OSError:
            relocated[name] = path_text
            continue
        relocated[name] = str(target)
        changed = True
    return relocated if changed else materialized_files


def _wrap_sandbox_result(
    result: Dict[str, Any],
    *,
    output_dir: str = "",
) -> Dict[str, Any]:
    """Normalize sandbox result into structured runtime payload."""
    text = _format_sandbox_result(result)
    markers = _parse_koto_file_markers(str(result.get("stdout", "")))
    created = markers.get("created", [])
    modified = markers.get("modified", [])
    generated_files = result.get("files") or result.get("images") or {}
    if not isinstance(generated_files, dict):
        generated_files = {}
    materialized_files = _materialize_sandbox_files(generated_files)
    materialized_files = _relocate_materialized_files_to_output_dir(
        materialized_files,
        output_dir=output_dir,
    )
    # fallback: when stdout has no KOTO_CREATED markers, use sandbox-captured files
    if not created and materialized_files:
        created = list(materialized_files.values())
    return SandboxRunResult(
        {
            "summary": text,
            "stdout": str(result.get("stdout") or ""),
            "stderr": str(result.get("stderr") or ""),
            "error": str(result.get("error") or ""),
            "files": dict(generated_files),
            "generated_file_paths": materialized_files,
            "generated_files": [
                {
                    "name": name,
                    "path": path,
                    "file_type": Path(path).suffix.lstrip(".").lower(),
                }
                for name, path in materialized_files.items()
            ],
            "_koto_created": created,
            "_koto_modified": modified,
            "__koto_created__": created,
            "__koto_modified__": modified,
        }
    )


def _materialize_sandbox_files(files: Dict[str, Any]) -> Dict[str, str]:
    """Persist sandbox-captured files so later tools can consume real paths."""
    materialized: Dict[str, str] = {}
    if not files:
        return materialized
    artifact_dir = tempfile.mkdtemp(prefix="koto-task-artifacts-")
    for raw_name, raw_data in files.items():
        filename = _safe_artifact_filename(str(raw_name or "artifact").strip())
        if not filename:
            filename = "artifact.bin"
        try:
            data = base64.b64decode(str(raw_data or ""), validate=False)
        except Exception:
            continue
        if not data:
            continue
        target = os.path.join(artifact_dir, filename)
        try:
            with open(target, "wb") as fh:
                fh.write(data)
        except OSError:
            continue
        materialized[str(raw_name)] = target
    if not materialized:
        try:
            shutil.rmtree(artifact_dir)
        except OSError:
            pass
    return materialized


def _safe_artifact_filename(filename: str) -> str:
    name = os.path.basename(filename.replace("\\", "/")).strip()
    if not name or name in {".", ".."}:
        return ""
    stem = re.sub(r"[^A-Za-z0-9._ \-\u4e00-\u9fff]", "_", name)
    return stem[:120] or "artifact.bin"


def list_workspace_files(path: str = "", recursive: bool = False) -> str:
    """List files in workspace directory. Returns JSON array of file info."""
    root = _get_workspace_root()
    target = _safe_resolve(path) if path else root
    if target is None or not os.path.isdir(target):
        return json.dumps({"error": f"Directory not found: {path}"}, ensure_ascii=False)

    entries = []
    try:
        if recursive:
            for dirpath, dirnames, filenames in os.walk(target):
                dirnames[:] = [d for d in dirnames if not d.startswith(".")]
                for fname in filenames:
                    if fname.startswith("."):
                        continue
                    fpath = os.path.join(dirpath, fname)
                    rel = os.path.relpath(fpath, root).replace("\\", "/")
                    try:
                        st = os.stat(fpath)
                        entries.append({"name": rel, "size": st.st_size})
                    except OSError:
                        pass
                    if len(entries) >= 300:
                        break
        else:
            for item in sorted(os.listdir(target)):
                if item.startswith("."):
                    continue
                fpath = os.path.join(target, item)
                rel = os.path.relpath(fpath, root).replace("\\", "/")
                is_dir = os.path.isdir(fpath)
                entries.append(
                    {
                        "name": rel,
                        "type": "dir" if is_dir else "file",
                        "size": 0 if is_dir else os.path.getsize(fpath),
                    }
                )
    except PermissionError:
        return json.dumps({"error": "Permission denied"}, ensure_ascii=False)

    return json.dumps(entries, ensure_ascii=False)


def open_file_in_editor(path: str) -> str:
    """Open a file in the frontend editor so the user can view it.
    Use this when the user wants to open, view, or navigate to a file
    — NOT for reading content. The file will be brought into focus in the UI.
    """
    resolved = _resolve_path(path)
    if not resolved or not os.path.isfile(resolved):
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)
    fname = os.path.basename(resolved)
    root = _get_workspace_root()
    rel = os.path.relpath(resolved, root).replace("\\", "/") if root else resolved
    return _success_result(
        rel,
        operation="open_file",
        summary=f"已打开 {fname}",
        file_type=Path(resolved).suffix.lstrip(".").lower(),
        change_type="open",
        focus=True,
    )


def copy_file(source: str, destination: str) -> str:
    """Copy a file within the workspace."""
    src = _resolve_path(source)
    if not src:
        return json.dumps({"error": f"Source not found: {source}"}, ensure_ascii=False)
    dst = _safe_resolve(destination)
    if not dst:
        return json.dumps(
            {"error": f"Invalid destination: {destination}"}, ensure_ascii=False
        )
    try:
        write_warning = (
            _ensure_existing_file_writable(dst) if os.path.exists(dst) else ""
        )
        result = _file_service().copy_file(src, dst, overwrite=True)
        if not result.get("success"):
            error_text = str(result.get("error") or "复制失败")
            if (
                any(
                    token in error_text.lower()
                    for token in ("permission", "denied")
                )
                or "权限" in error_text
            ):
                return _blocked_write_result(
                    _result_path(destination, dst),
                    summary=error_text,
                    suggested_next_step=_nonwritable_target_next_step(dst),
                    operation="copy_file",
                    file_type=Path(dst).suffix.lstrip(".").lower(),
                )
            return json.dumps({"error": error_text}, ensure_ascii=False)
        copied_path = str(result.get("destination") or dst)
        return _success_result(
            _result_path(destination, copied_path),
            operation="copy_file",
            summary=f"已复制文件到 {os.path.basename(copied_path)}",
            change_type="create",
            file_type=Path(copied_path).suffix.lstrip(".").lower(),
            focus=True,
            warning=write_warning,
        )
    except PermissionError as exc:
        return _blocked_write_result(
            _result_path(destination, dst),
            summary=str(exc).strip() or _nonwritable_target_message(dst),
            suggested_next_step=_nonwritable_target_next_step(dst),
            operation="copy_file",
            file_type=Path(dst).suffix.lstrip(".").lower(),
        )
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def _write_docx_content_without_python_docx(
    path: str, resolved: str, para_list: Any
) -> str:
    paragraphs = _normalize_docx_paragraphs(para_list)
    if not paragraphs:
        paragraphs = [{"text": "", "style": "Normal"}]
    file_exists = os.path.exists(resolved)
    backup_warning = _best_effort_backup(resolved) if file_exists else ""
    write_warning = _ensure_existing_file_writable(resolved) if file_exists else ""
    os.makedirs(os.path.dirname(resolved), exist_ok=True)
    docx_bytes = _minimal_docx_package_bytes(
        paragraphs, existing_path=resolved if file_exists else ""
    )
    _write_bytes_via_temp_file(docx_bytes, resolved, suffix=".docx")
    preview = "\n".join(str(p.get("text", "")) for p in paragraphs[:3])
    diff_items = [
        {
            "paragraph_index": index,
            "before": "",
            "after": str(item.get("text") or ""),
            "style": str(item.get("style") or ""),
        }
        for index, item in enumerate(paragraphs, start=1)
    ]
    return _success_result(
        _result_path(path, resolved),
        operation="write_docx_content",
        summary=f"已写入 {len(paragraphs)} 个段落到 Word 文档",
        file_type="docx",
        change_type="modify" if file_exists else "create",
        preview=preview,
        focus=True,
        summary_code="WRITE_OK" if file_exists else "CREATE_OK",
        diff=_file_task_diff("docx_paragraphs", diff_items),
        paragraphs_written=len(paragraphs),
        warning=_merge_warnings(
            backup_warning,
            write_warning,
            "python-docx 未安装，已使用基础 DOCX 写入兜底。",
        ),
    )


def _create_docx_file(path: str, resolved: str, content: str) -> str:
    return _office_create_docx_file(
        path,
        resolved,
        content,
        plain_text_to_paragraphs=_plain_text_to_docx_paragraphs,
        save_document=_save_docx_via_temp_file,
        fallback_writer=_write_docx_content_without_python_docx,
        success_result=_success_result,
        result_path=_result_path,
        file_task_diff=_file_task_diff,
    )


def _create_xlsx_file(path: str, resolved: str, content: str) -> str:
    return _office_create_xlsx_file(
        path,
        resolved,
        content,
        save_workbook=_save_workbook_via_temp_file,
        success_result=_success_result,
        result_path=_result_path,
        file_task_diff=_file_task_diff,
    )


def _plain_text_to_pptx_slides(content: str) -> List[Dict[str, Any]]:
    return _office_plain_text_to_pptx_slides(content)


def _create_pptx_file(path: str, resolved: str, content: str) -> str:
    return _office_create_pptx_file(
        path,
        resolved,
        content,
        text_lines=_pptx_text_lines,
        save_presentation=_save_pptx_via_temp_file,
        success_result=_success_result,
        result_path=_result_path,
    )


def create_file(path: str, content: str = "") -> str:
    """Create a new file in the workspace.

    Office targets are created as real package files and emit the same write
    metrics as their specialized tools, so task quality gates can verify them.
    """
    resolved = _safe_resolve(path)
    if not resolved:
        return json.dumps({"error": f"Invalid path: {path}"}, ensure_ascii=False)
    if os.path.exists(resolved):
        return json.dumps({"error": "File already exists"}, ensure_ascii=False)
    try:
        suffix = Path(resolved).suffix.lower()
        if suffix in {".docx", ".doc"}:
            return _create_docx_file(path, resolved, content)
        if suffix in {".xlsx", ".xlsm"}:
            return _create_xlsx_file(path, resolved, content)
        if suffix in {".pptx", ".ppt"}:
            return _create_pptx_file(path, resolved, content)
        os.makedirs(os.path.dirname(resolved), exist_ok=True)
        with open(resolved, "w", encoding="utf-8") as f:
            f.write(content)
        return _success_result(
            _result_path(path, resolved),
            operation="create_file",
            summary=f"已创建文件 {os.path.basename(resolved)}",
            change_type="create",
            preview=content,
            focus=True,
        )
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def llm_extract(text: str, fields: str, instructions: str = "") -> str:
    """Use LLM to extract structured data from text.

    Args:
        text: The source text to extract from.
        fields: Comma-separated field names to extract.
        instructions: Optional additional instructions for the LLM.

    Returns: JSON object with field names as keys.
    """
    from app.core.workflow_engine import call_llm_json

    field_list = [f.strip() for f in fields.split(",") if f.strip()]
    prompt = (
        f"从以下文本中提取指定字段的值。\n\n"
        f"字段列表: {json.dumps(field_list, ensure_ascii=False)}\n\n"
        f"文本:\n---\n{text[:6000]}\n---\n\n"
    )
    if instructions:
        prompt += f"额外要求: {instructions}\n\n"
    prompt += "以 JSON 对象输出，key 为字段名，value 为提取到的值。" "找不到的字段值设为 null。只输出 JSON。"
    try:
        result = call_llm_json(prompt, call_timeout=_TASK_TOOL_LLM_CALL_TIMEOUT)
        if isinstance(result, dict):
            return json.dumps(result, ensure_ascii=False)
        return json.dumps({"raw": str(result)}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def llm_transform(text: str, instruction: str) -> str:
    """Use LLM to transform text according to instruction.

    Args:
        text: The source text to transform.
        instruction: What transformation to apply.

    Returns: The transformed text.
    """
    from app.core.workflow_engine import call_llm

    prompt = f"{instruction}\n\n原文:\n---\n{text[:6000]}\n---"
    try:
        return call_llm(prompt, call_timeout=_TASK_TOOL_LLM_CALL_TIMEOUT)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


# ══════════════════════════════════════════════════════════════
# New tools for DocAgent - Multi-file operations
# ══════════════════════════════════════════════════════════════


def compare_files(file_paths: str, aspect: str = "content") -> str:
    """Compare multiple files and return similarity analysis.

    Args:
        file_paths: Comma-separated list of file paths to compare.
        aspect: What to compare ("content", "structure").

    Returns: JSON with similarity scores and differences.
    """
    import asyncio
    from app.core.file.multi_file_coordinator import get_file_coordinator
    from app.core.agent.doc_agent import FileHandle

    paths = [p.strip() for p in file_paths.split(",") if p.strip()]
    if len(paths) < 2:
        return json.dumps({"error": "至少需要两个文件进行对比"}, ensure_ascii=False)

    resolved_paths = []
    for p in paths:
        resolved = _resolve_path(p)
        if not resolved:
            return json.dumps({"error": f"文件不存在: {p}"}, ensure_ascii=False)
        resolved_paths.append(resolved)

    try:
        files = [FileHandle(path=p) for p in resolved_paths]
        coordinator = get_file_coordinator()

        # Run async comparison in sync context
        loop = asyncio.new_event_loop()
        try:
            result = loop.run_until_complete(
                coordinator.compare_documents(files, aspect)
            )
        finally:
            loop.close()

        return json.dumps(result.to_dict(), ensure_ascii=False, default=str)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def plan_docx_compare_annotations(
    original_path: str,
    revised_path: str,
    target_path: str = "",
    max_differences: int = 80,
) -> str:
    """Return DOCX difference anchors so the model can write comment text.

    This tool does not modify files. It finds target-document anchors and
    provides default "另一份为/本文件为" comments that the model may refine before
    calling write_docx_comments.
    """
    original_resolved = _resolve_path(original_path)
    revised_resolved = _resolve_path(revised_path)
    if not original_resolved:
        return json.dumps({"error": f"原始文件不存在: {original_path}"}, ensure_ascii=False)
    if not revised_resolved:
        return json.dumps({"error": f"对比文件不存在: {revised_path}"}, ensure_ascii=False)
    if Path(original_resolved).suffix.lstrip(".").lower() != "docx":
        return json.dumps(
            {"error": f"只支持 DOCX 原始文件: {original_path}"}, ensure_ascii=False
        )
    if Path(revised_resolved).suffix.lstrip(".").lower() != "docx":
        return json.dumps(
            {"error": f"只支持 DOCX 对比文件: {revised_path}"}, ensure_ascii=False
        )

    target_raw = str(target_path or original_path or "").strip()
    target_resolved = _resolve_path(target_raw) if target_raw else original_resolved
    if not target_resolved:
        return json.dumps({"error": f"目标 DOCX 不存在: {target_raw}"}, ensure_ascii=False)
    if Path(target_resolved).suffix.lstrip(".").lower() != "docx":
        return json.dumps({"error": f"目标文件必须是 DOCX: {target_raw}"}, ensure_ascii=False)

    try:
        max_items = _normalize_positive_int(max_differences, default=80, upper=200)
        annotations, differences_detected = _docx_compare_annotation_candidates(
            original_resolved,
            revised_resolved,
            target_resolved,
            max_differences=max_items,
        )
        risks = _contract_risk_summary_from_annotations(annotations)
        return json.dumps(
            {
                "success": True,
                "operation": "plan_docx_compare_annotations",
                "original_path": _result_path(original_path, original_resolved),
                "revised_path": _result_path(revised_path, revised_resolved),
                "target_path": _result_path(target_raw, target_resolved),
                "differences_detected": differences_detected,
                "annotation_candidates": annotations,
                "contract_risk_summary": risks,
                "summary": f"已定位 {len(annotations)} 处可写入目标 DOCX 的差异批注候选。",
            },
            ensure_ascii=False,
        )
    except Exception as exc:
        return json.dumps({"error": str(exc)}, ensure_ascii=False)


def write_docx_comments(
    path: str,
    comments_json: Any = "[]",
    source_path: str = "",
    compare_path: str = "",
    differences_detected: int = 0,
) -> str:
    """Write model-authored Word comments into an existing DOCX in place."""
    target_resolved = _resolve_path(path)
    if not target_resolved:
        return json.dumps({"error": f"目标 DOCX 不存在: {path}"}, ensure_ascii=False)
    if Path(target_resolved).suffix.lstrip(".").lower() != "docx":
        return json.dumps({"error": f"目标文件必须是 DOCX: {path}"}, ensure_ascii=False)

    try:
        annotations = _parse_annotations_payload(comments_json)
    except ValueError as exc:
        return json.dumps({"error": str(exc)}, ensure_ascii=False)
    normalized: List[Dict[str, str]] = []
    for item in annotations:
        if not isinstance(item, dict):
            continue
        anchor = str(
            item.get("原文片段")
            or item.get("anchor_text")
            or item.get("anchor")
            or item.get("text")
            or item.get("target_text")
            or ""
        ).strip()
        comment = str(
            item.get("批注内容")
            or item.get("comment")
            or item.get("note")
            or item.get("修改后文本")
            or ""
        ).strip()
        if not anchor or not comment:
            continue
        normalized.append(
            {
                "原文片段": anchor,
                "批注内容": comment,
                "批注标签": str(item.get("批注标签") or item.get("label") or "").strip(),
                "修改原因": str(item.get("修改原因") or item.get("reason") or "").strip(),
            }
        )

    if not normalized:
        return json.dumps({"error": "没有可写入的 DOCX 批注"}, ensure_ascii=False)

    try:
        from web.track_changes_editor import TrackChangesEditor
    except ImportError:
        return json.dumps({"error": "缺少 Word 批注写入组件"}, ensure_ascii=False)

    try:
        backup_warning = _best_effort_backup(target_resolved)
        write_warning = _ensure_existing_file_writable(target_resolved)
        editor = TrackChangesEditor(author="Koto AI")
        applied = editor.apply_comment_changes(target_resolved, normalized)
        if not applied.get("success"):
            return json.dumps(
                {"error": str(applied.get("error") or "DOCX 批注写入失败")},
                ensure_ascii=False,
            )
        annotations_added = int(applied.get("applied") or 0)
        failed = int(applied.get("failed") or 0)
        risks = _contract_risk_summary_from_annotations(normalized)
        detected = _normalize_positive_int(
            differences_detected, default=len(normalized), upper=10000
        )
        return _success_result(
            _result_path(path, target_resolved),
            operation="write_docx_comments",
            summary=f"已在目标 DOCX 原文上写入 {annotations_added} 条 Word 批注。",
            file_type="docx",
            change_type="annotate",
            preview="\n".join(
                [
                    f"差异候选：{detected}",
                    f"已写入批注：{annotations_added}",
                    *([f"未能定位：{failed}"] if failed else []),
                ]
            ),
            focus=True,
            annotations_added=annotations_added,
            differences_detected=detected,
            comments_failed=failed,
            source_path=source_path,
            compare_path=compare_path,
            contract_risk_summary=risks,
            warning=_merge_warnings(backup_warning, write_warning),
        )
    except PermissionError as exc:
        return _blocked_write_result(
            _result_path(path, target_resolved),
            summary=str(exc).strip() or _nonwritable_target_message(target_resolved),
            suggested_next_step=_nonwritable_target_next_step(target_resolved),
            operation="write_docx_comments",
            file_type="docx",
        )
    except Exception as exc:
        return json.dumps({"error": str(exc)}, ensure_ascii=False)


def compare_docx_and_annotate(
    original_path: str,
    revised_path: str,
    target_path: str = "",
    max_differences: int = 80,
) -> str:
    """Compare two DOCX files and mark differences with Word comments.

    The comments are written to target_path when provided, otherwise to revised_path.
    The annotated text always comes from the actual target document; the other
    DOCX is used only as the comparison source.
    """
    original_resolved = _resolve_path(original_path)
    revised_resolved = _resolve_path(revised_path)
    if not original_resolved:
        return json.dumps({"error": f"原始文件不存在: {original_path}"}, ensure_ascii=False)
    if not revised_resolved:
        return json.dumps({"error": f"对比文件不存在: {revised_path}"}, ensure_ascii=False)
    if Path(original_resolved).suffix.lstrip(".").lower() != "docx":
        return json.dumps(
            {"error": f"只支持 DOCX 原始文件: {original_path}"}, ensure_ascii=False
        )
    if Path(revised_resolved).suffix.lstrip(".").lower() != "docx":
        return json.dumps(
            {"error": f"只支持 DOCX 对比文件: {revised_path}"}, ensure_ascii=False
        )

    target_raw = str(target_path or revised_path or "").strip()
    target_resolved = _resolve_path(target_raw) if target_raw else revised_resolved
    if not target_resolved:
        return json.dumps({"error": f"目标 DOCX 不存在: {target_raw}"}, ensure_ascii=False)
    if Path(target_resolved).suffix.lstrip(".").lower() != "docx":
        return json.dumps({"error": f"目标文件必须是 DOCX: {target_raw}"}, ensure_ascii=False)

    try:
        from web.track_changes_editor import TrackChangesEditor
    except ImportError:
        return json.dumps({"error": "缺少 Word 批注写入组件"}, ensure_ascii=False)

    try:
        max_items = _normalize_positive_int(max_differences, default=80, upper=200)
        annotations, differences_detected = _docx_compare_annotation_candidates(
            original_resolved,
            revised_resolved,
            target_resolved,
            max_differences=max_items,
        )
        contract_risk_summary = _contract_risk_summary_from_annotations(annotations)
        if not annotations:
            return _success_result(
                _result_path(target_raw, target_resolved),
                operation="compare_docx_and_annotate",
                summary="两份 DOCX 未发现可标注的正文差异。",
                file_type="docx",
                change_type="annotate",
                preview="未发现正文差异",
                focus=True,
                annotations_added=0,
                differences_detected=0,
                source_path=_result_path(original_path, original_resolved),
                compare_path=_result_path(revised_path, revised_resolved),
            )

        backup_warning = _best_effort_backup(target_resolved)
        write_warning = _ensure_existing_file_writable(target_resolved)
        editor = TrackChangesEditor(author="Koto AI")
        applied = editor.apply_comment_changes(target_resolved, annotations)
        if not applied.get("success"):
            return json.dumps(
                {"error": str(applied.get("error") or "DOCX 差异批注写入失败")},
                ensure_ascii=False,
            )
        annotations_added = int(applied.get("applied") or 0)
        failed = int(applied.get("failed") or 0)
        preview_lines = [
            f"发现差异：{differences_detected}",
            f"已标注批注：{annotations_added}",
        ]
        if failed:
            preview_lines.append(f"未能定位：{failed}")
        if annotations:
            preview_lines.append(
                str(annotations[0].get("批注内容") or annotations[0].get("修改后文本") or "")[
                    :160
                ]
            )
        return _success_result(
            _result_path(target_raw, target_resolved),
            operation="compare_docx_and_annotate",
            summary=f"已对比两份 DOCX，并在目标文档标注 {annotations_added} 处差异。",
            file_type="docx",
            change_type="annotate",
            preview="\n".join(preview_lines),
            focus=True,
            annotations_added=annotations_added,
            differences_detected=differences_detected,
            comments_failed=failed,
            source_path=_result_path(original_path, original_resolved),
            compare_path=_result_path(revised_path, revised_resolved),
            contract_risk_summary=contract_risk_summary,
            warning=_merge_warnings(backup_warning, write_warning),
        )
    except PermissionError as exc:
        return _blocked_write_result(
            _result_path(target_raw, target_resolved),
            summary=str(exc).strip() or _nonwritable_target_message(target_resolved),
            suggested_next_step=_nonwritable_target_next_step(target_resolved),
            operation="compare_docx_and_annotate",
            file_type="docx",
        )
    except Exception as exc:
        return json.dumps({"error": str(exc)}, ensure_ascii=False)


def extract_to_file(
    source_path: str,
    target_path: str,
    extract_query: str,
    insert_position: str = "end",
) -> str:
    """Extract data from source file and inject into target file.

    Args:
        source_path: Path to source file (e.g., Excel with data).
        target_path: Path to target file (e.g., Word document).
        extract_query: Description of what to extract.
        insert_position: Where to insert ("start", "end", "cursor").

    Returns: JSON with operation result and change details.
    """
    import asyncio
    from app.core.file.multi_file_coordinator import get_file_coordinator
    from app.core.agent.doc_agent import FileHandle

    src = _resolve_path(source_path)
    if not src:
        return json.dumps({"error": f"源文件不存在: {source_path}"}, ensure_ascii=False)

    tgt = _resolve_path(target_path)
    if not tgt:
        # Target can be new file
        tgt = _safe_resolve(target_path)
        if not tgt:
            return json.dumps({"error": f"目标路径无效: {target_path}"}, ensure_ascii=False)

    try:
        source = FileHandle(path=src)
        target = FileHandle(path=tgt)
        coordinator = get_file_coordinator()

        loop = asyncio.new_event_loop()
        try:
            change = loop.run_until_complete(
                coordinator.extract_and_inject(
                    source, target, extract_query, insert_position
                )
            )
        finally:
            loop.close()

        return json.dumps(
            {
                "success": change.change_type != "none",
                "change": change.to_dict(),
            },
            ensure_ascii=False,
        )
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def _parse_annotations_payload(annotations: Any) -> List[Any]:
    try:
        ann_list = (
            json.loads(annotations) if isinstance(annotations, str) else annotations
        )
    except json.JSONDecodeError as e:
        raise ValueError(f"无效的 annotations JSON: {e}") from e
    if ann_list in (None, ""):
        return []
    if not isinstance(ann_list, list):
        raise ValueError("annotations 必须是 JSON 数组")
    return ann_list


def _annotation_preview(annotations: List[Any]) -> str:
    return "\n".join(
        str(item.get("comment") or "").strip()
        for item in annotations[:3]
        if isinstance(item, dict) and str(item.get("comment") or "").strip()
    )


def _annotation_result_payload(
    raw_path: str,
    resolved_path: str,
    *,
    annotations_added: int,
    summary: str,
    preview: str,
    **extra: Any,
) -> Dict[str, Any]:
    display_path = _result_path(raw_path, resolved_path)
    payload: Dict[str, Any] = {
        "success": True,
        "path": display_path,
        "file_path": display_path,
        "file_type": Path(resolved_path).suffix.lstrip(".").lower(),
        "change_type": "annotate",
        "operation": "annotate_file",
        "summary": str(summary or "").strip(),
        "preview": str(preview or "").strip(),
        "focus": True,
        "annotations_added": int(annotations_added or 0),
    }
    for key, value in extra.items():
        if value not in (None, ""):
            payload[key] = value
    return payload


def _build_docx_annotation_request(
    path: str,
    *,
    requirement: str,
    model_id: str = "",
    task_files: Optional[List[Dict[str, Any]]] = None,
    request_context: Optional[Dict[str, Any]] = None,
):
    from app.core.agent.file_task_contract import FileTaskFile, FileTaskRequest

    context = dict(request_context or {})
    target_path = (
        str(context.get("target_path") or path or "").strip() or str(path or "").strip()
    )
    task_text = str(context.get("task") or requirement or "").strip()
    options = (
        dict(context.get("options") or {})
        if isinstance(context.get("options"), dict)
        else {}
    )

    files: List[FileTaskFile] = []
    seen_paths: set[str] = set()
    target_norm = os.path.normcase(str(path or target_path or "").strip())
    for item in task_files or []:
        if not isinstance(item, dict):
            continue
        raw_path = str(item.get("path") or item.get("name") or "").strip()
        normalized = os.path.normcase(raw_path)
        if raw_path and normalized in seen_paths:
            continue
        parsed = FileTaskFile.from_mapping(item)
        if raw_path and normalized == target_norm:
            parsed.target = True
        if parsed.path or parsed.name or parsed.content:
            files.append(parsed)
            if raw_path:
                seen_paths.add(normalized)

    if not any(
        os.path.normcase(str(file_info.path or "")) == target_norm
        for file_info in files
    ):
        files.append(
            FileTaskFile(
                path=str(path or target_path or "").strip(),
                name=Path(str(path or target_path or "").strip()).name,
                type=Path(str(path or target_path or "").strip())
                .suffix.lstrip(".")
                .lower(),
                target=True,
            )
        )

    return FileTaskRequest(
        task=task_text,
        files=files,
        target_path=target_path,
        model_mode=str(context.get("model_mode") or "deepseek").strip() or "deepseek",
        model_id=str(model_id or context.get("model_id") or "").strip(),
        options=options,
    )


def _stream_docx_annotation_tool_result(
    path: str,
    *,
    requirement: str,
    model_id: str = "",
    gemini_client: Any = None,
    workspace_root: str = "",
    task_files: Optional[List[Dict[str, Any]]] = None,
    request_context: Optional[Dict[str, Any]] = None,
) -> FileTaskToolStreamResult:
    from app.core.agent import file_task_doc_annotate_boundary

    annotation_request = _build_docx_annotation_request(
        path,
        requirement=requirement,
        model_id=model_id,
        task_files=task_files,
        request_context=request_context,
    )
    return file_task_doc_annotate_boundary.stream_bridge_request_as_tool(
        annotation_request,
        workspace_root=workspace_root,
        gemini_client=gemini_client,
    )


def annotate_file(
    path: str,
    annotations: Any = "[]",
    requirement: str = "",
    model_id: str = "",
    gemini_client: Any = None,
) -> Any:
    """Add annotations/highlights to a file.

    Args:
        path: Path to the file to annotate.
        annotations: JSON array of annotations:
            [{"range_start": 0, "range_end": 100, "comment": "...", "color": "yellow"}]
        requirement: Optional AI review requirement for DOCX native comment generation.
        model_id: Optional model ID for DOCX native comment generation.

    Returns: JSON with annotation results.
    """
    import asyncio
    from app.core.file.multi_file_coordinator import get_file_coordinator

    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"文件不存在: {path}"}, ensure_ascii=False)

    try:
        ann_list = _parse_annotations_payload(annotations)
    except ValueError as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)

    requirement_text = str(requirement or "").strip()
    if (
        Path(resolved).suffix.lstrip(".").lower() == "docx"
        and requirement_text
        and not ann_list
    ):
        return _stream_docx_annotation_tool_result(
            path,
            requirement=requirement_text,
            model_id=model_id,
            gemini_client=gemini_client,
            workspace_root=_WORKSPACE_ROOT,
        )

    try:
        coordinator = get_file_coordinator()

        loop = asyncio.new_event_loop()
        try:
            changes = loop.run_until_complete(
                coordinator.annotate_file(resolved, ann_list)
            )
        finally:
            loop.close()

        return json.dumps(
            {
                **_annotation_result_payload(
                    path,
                    resolved,
                    annotations_added=len(changes),
                    summary=f"已添加 {len(changes)} 条批注到 {os.path.basename(resolved)}",
                    preview=_annotation_preview(ann_list),
                ),
                "changes": [c.to_dict() for c in changes],
            },
            ensure_ascii=False,
        )
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


_DOCX_REVIEW_CLEAR_SCOPE_ALIASES = {
    "": "comments",
    "comment": "comments",
    "comments": "comments",
    "annotation": "comments",
    "annotations": "comments",
    "批注": "comments",
    "标注": "comments",
    "评论": "comments",
    "注释": "comments",
    "review": "all",
    "review_marks": "all",
    "all": "all",
    "全部": "all",
    "所有": "all",
    "修订": "revisions",
    "revision": "revisions",
    "revisions": "revisions",
    "tracked_changes": "revisions",
}

def _normalize_docx_review_clear_scope(scope: str) -> str:
    return _review_normalize_docx_review_clear_scope(
        scope,
        _DOCX_REVIEW_CLEAR_SCOPE_ALIASES,
    )


def clear_docx_review_marks(path: str, scope: str = "comments") -> str:
    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"文件不存在: {path}"}, ensure_ascii=False)
    if Path(resolved).suffix.lstrip(".").lower() != "docx":
        return json.dumps({"error": f"只支持 DOCX 文件: {path}"}, ensure_ascii=False)

    try:
        normalized_scope = _normalize_docx_review_clear_scope(scope)
    except ValueError as exc:
        return json.dumps({"error": str(exc)}, ensure_ascii=False)

    try:
        from lxml import etree
    except ImportError:
        return json.dumps({"error": "lxml not installed"}, ensure_ascii=False)

    backup_warning = _best_effort_backup(resolved)
    write_warning = _ensure_existing_file_writable(resolved)

    try:
        with zipfile.ZipFile(resolved, "r") as archive:
            entries = {
                info.filename: archive.read(info.filename)
                for info in archive.infolist()
            }
    except zipfile.BadZipFile as exc:
        return json.dumps({"error": f"无法读取 DOCX 压缩包: {exc}"}, ensure_ascii=False)
    except Exception as exc:
        return json.dumps({"error": str(exc)}, ensure_ascii=False)

    document_xml = entries.get("word/document.xml")
    if not document_xml:
        return json.dumps({"error": "DOCX 缺少 word/document.xml"}, ensure_ascii=False)

    try:
        document_root = etree.fromstring(document_xml)
    except Exception as exc:
        return json.dumps({"error": f"无法解析 DOCX 正文 XML: {exc}"}, ensure_ascii=False)

    changed = False
    comments_removed = 0
    comment_markup_removed = 0
    revisions_accepted = 0

    if normalized_scope in {"comments", "all"}:
        comments_xml = entries.pop("word/comments.xml", None)
        if comments_xml is not None:
            try:
                comments_root = etree.fromstring(comments_xml)
                comments_removed = len(
                    comments_root.xpath(".//w:comment", namespaces={"w": _DOCX_W_NS})
                )
            except Exception:
                comments_removed = 0
            changed = True

        comment_markup_removed = _remove_docx_comment_markup(document_root)
        if comment_markup_removed:
            changed = True

        rels_xml = entries.get("word/_rels/document.xml.rels")
        if rels_xml is not None:
            updated_rels, rels_removed = _remove_comments_relationships_xml(rels_xml)
            if rels_removed:
                entries["word/_rels/document.xml.rels"] = updated_rels
                changed = True

        content_types_xml = entries.get("[Content_Types].xml")
        if content_types_xml is not None:
            (
                updated_content_types,
                overrides_removed,
            ) = _remove_comments_content_type_override(content_types_xml)
            if overrides_removed:
                entries["[Content_Types].xml"] = updated_content_types
                changed = True

        approx_comments = comment_markup_removed // len(_DOCX_COMMENT_MARKUP_TAGS)
        if comment_markup_removed and not approx_comments:
            approx_comments = 1
        comments_removed = max(comments_removed, approx_comments)

    if normalized_scope in {"revisions", "all"}:
        revisions_accepted = _accept_docx_revision_markup(document_root)
        if revisions_accepted:
            changed = True

    if changed:
        entries["word/document.xml"] = _serialize_xml_root(document_root)

    summary = _build_docx_review_clear_summary(
        normalized_scope,
        comments_removed,
        revisions_accepted,
        changed=changed,
    )
    preview_lines: List[str] = []
    if normalized_scope in {"comments", "all"}:
        preview_lines.append(f"批注: {comments_removed}")
    if normalized_scope in {"revisions", "all"}:
        preview_lines.append(f"修订: {revisions_accepted}")

    if not changed:
        result_warning = _merge_warnings(backup_warning, write_warning)
        return _success_result(
            _result_path(path, resolved),
            operation="clear_docx_review_marks",
            summary=summary,
            file_type="docx",
            change_type="modify",
            preview="\n".join(preview_lines),
            focus=False,
            changed=False,
            scope=normalized_scope,
            comments_removed=comments_removed,
            comment_markup_removed=comment_markup_removed,
            revisions_accepted=revisions_accepted,
            warning=result_warning,
        )

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        for name, data in entries.items():
            archive.writestr(name, data)
    docx_bytes = buffer.getvalue()

    try:
        _write_bytes_via_temp_file(docx_bytes, resolved, suffix=".docx")
    except PermissionError as exc:
        locked_message = str(exc).strip() or _nonwritable_target_message(resolved)
        try:
            fallback_resolved = _next_available_docx_copy_path(resolved)
            _write_bytes_via_temp_file(docx_bytes, fallback_resolved, suffix=".docx")
        except Exception:
            return _blocked_write_result(
                _result_path(path, resolved),
                summary=locked_message,
                suggested_next_step=_nonwritable_target_next_step(resolved),
                scope=normalized_scope,
                comments_removed=comments_removed,
                comment_markup_removed=comment_markup_removed,
                revisions_accepted=revisions_accepted,
            )

        result_warning = _merge_warnings(
            backup_warning,
            write_warning,
            f"原目标文件 {os.path.basename(resolved)} 当前不可写，结果已写入同目录副本 {os.path.basename(fallback_resolved)}。",
        )
        return _blocked_write_result(
            _result_path(fallback_resolved, fallback_resolved),
            summary=(
                f"原目标文件 {os.path.basename(resolved)} 当前不可写，尚未写回原文件；"
                f"已将清理结果写入恢复副本 {os.path.basename(fallback_resolved)}"
            ),
            suggested_next_step=_nonwritable_target_next_step(resolved),
            operation="clear_docx_review_marks",
            file_type="docx",
            change_type="create",
            preview="\n".join(preview_lines),
            focus=True,
            warning=result_warning,
            original_target_path=_result_path(path, resolved),
            blocked_target=True,
            blocked_reason=locked_message,
            fallback_copy=True,
            scope=normalized_scope,
            comments_removed=comments_removed,
            comment_markup_removed=comment_markup_removed,
            revisions_accepted=revisions_accepted,
        )

    result_warning = _merge_warnings(backup_warning, write_warning)
    return _success_result(
        _result_path(path, resolved),
        operation="clear_docx_review_marks",
        summary=summary,
        file_type="docx",
        change_type="modify",
        preview="\n".join(preview_lines),
        focus=True,
        changed=True,
        scope=normalized_scope,
        comments_removed=comments_removed,
        comment_markup_removed=comment_markup_removed,
        revisions_accepted=revisions_accepted,
        warning=result_warning,
    )


def _normalize_table_columns(columns: Any) -> List[str]:
    if not columns:
        return []
    value = columns
    if isinstance(columns, str):
        text = columns.strip()
        if not text:
            return []
        try:
            value = json.loads(text)
        except Exception:
            value = re.split(r"[,，、|]", text)
    if not isinstance(value, list):
        return []
    normalized: List[str] = []
    for item in value:
        text = str(item or "").strip()
        if text and text not in normalized:
            normalized.append(text)
    return normalized


def _match_header_index(headers: List[str], wanted: str) -> Optional[int]:
    wanted_text = str(wanted or "").strip().casefold()
    if not wanted_text:
        return None
    normalized_headers = [str(header or "").strip().casefold() for header in headers]
    for index, header in enumerate(normalized_headers):
        if header == wanted_text:
            return index
    for index, header in enumerate(normalized_headers):
        if wanted_text in header or header in wanted_text:
            return index
    return None


def _table_sort_value(value: Any) -> tuple[int, Any]:
    text = str(value or "").strip()
    if not text:
        return (0, 0)
    numeric_text = re.sub(r"[,$%￥¥\s]", "", text)
    try:
        return (1, float(numeric_text))
    except ValueError:
        return (1, text.casefold())


def insert_excel_as_docx_table(
    source_path: str,
    target_path: str,
    sheet_name: str = "",
    table_title: str = "",
    max_rows: int = 200,
    sort_by: str = "",
    sort_order: str = "desc",
    columns: Any = "",
) -> str:
    """Insert spreadsheet data into a DOCX file as a real Word table."""
    max_rows = _normalize_positive_int(max_rows, default=200, upper=5_000)
    sort_by_text = str(sort_by or "").strip()
    sort_descending = str(sort_order or "desc").strip().lower() not in {
        "asc",
        "ascending",
        "smallest",
        "lowest",
        "正序",
        "升序",
    }
    selected_columns = _normalize_table_columns(columns)
    source_resolved = _resolve_path(source_path)
    if not source_resolved:
        return json.dumps(
            {"error": f"File not found: {source_path}"}, ensure_ascii=False
        )

    target_resolved = _resolve_path(target_path)
    if not target_resolved:
        if os.path.isabs(target_path):
            target_resolved = os.path.normpath(target_path)
        else:
            target_resolved = _safe_resolve(target_path)
    if not target_resolved:
        return json.dumps({"error": f"Invalid path: {target_path}"}, ensure_ascii=False)

    try:
        from io import BytesIO

        import openpyxl
        from docx import Document

        workbook = openpyxl.load_workbook(
            source_resolved, read_only=True, data_only=True
        )
        target_sheet, requested_sheet, sheet_warning = _select_workbook_sheet(
            workbook, sheet_name
        )
        if not target_sheet:
            workbook.close()
            return json.dumps(
                {"error": sheet_warning},
                ensure_ascii=False,
            )

        worksheet = workbook[target_sheet]
        raw_rows: List[List[str]] = []
        read_row_limit = 5_000 if sort_by_text else max_rows + 1
        for row in worksheet.iter_rows(values_only=True):
            if len(raw_rows) >= read_row_limit + 1:
                break
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                raw_rows.append(values)
        workbook.close()

        if not raw_rows:
            return json.dumps(
                {"error": f"Sheet '{target_sheet}' has no data"}, ensure_ascii=False
            )

        column_count = max(len(row) for row in raw_rows)
        normalized_rows = [row + [""] * (column_count - len(row)) for row in raw_rows]
        headers = normalized_rows[0]
        data_rows = normalized_rows[1:]
        sort_warning = ""
        selected_warning = ""

        if sort_by_text:
            sort_index = _match_header_index(headers, sort_by_text)
            if sort_index is None:
                sort_warning = f"未找到排序列“{sort_by_text}”，已保留原始行顺序。"
            else:
                data_rows = sorted(
                    data_rows,
                    key=lambda row_values: _table_sort_value(
                        row_values[sort_index] if sort_index < len(row_values) else ""
                    ),
                    reverse=sort_descending,
                )

        data_rows = data_rows[:max_rows]

        if selected_columns:
            selected_indexes: List[int] = []
            missing_columns: List[str] = []
            for column_name in selected_columns:
                matched_index = _match_header_index(headers, column_name)
                if matched_index is None:
                    missing_columns.append(column_name)
                    continue
                if matched_index not in selected_indexes:
                    selected_indexes.append(matched_index)
            if selected_indexes:
                headers = [headers[index] for index in selected_indexes]
                data_rows = [
                    [
                        row_values[index] if index < len(row_values) else ""
                        for index in selected_indexes
                    ]
                    for row_values in data_rows
                ]
                column_count = len(headers)
            if missing_columns:
                selected_warning = "未找到列：" + "、".join(missing_columns[:6])

        normalized_rows = [headers, *data_rows]

        os.makedirs(os.path.dirname(target_resolved), exist_ok=True)
        target_exists = os.path.exists(target_resolved)
        backup_warning = ""
        if target_exists:
            backup_warning = _best_effort_backup(target_resolved)
            try:
                with open(target_resolved, "rb") as existing_doc:
                    document = Document(BytesIO(existing_doc.read()))
            except Exception as open_err:
                return json.dumps(
                    {
                        "error": (
                            f"无法用 python-docx 打开 {os.path.basename(target_resolved)}"
                            f"（{open_err}）。"
                            "请改用 run_python_code + python-docx 手动追加表格，"
                            "或用 openpyxl 读取数据后用脚本写入 docx。"
                        )
                    },
                    ensure_ascii=False,
                )
        else:
            document = Document()

        if document.paragraphs and any(p.text.strip() for p in document.paragraphs):
            document.add_paragraph("")

        style_warning = ""
        if table_title:
            title_paragraph = document.add_paragraph(table_title)
            style_warning = _apply_docx_style(title_paragraph, "Heading 2")

        table = document.add_table(rows=len(normalized_rows), cols=column_count)
        style_warning = _merge_warnings(
            style_warning,
            _apply_docx_style(table, "Table Grid"),
        )

        for row_index, row_values in enumerate(normalized_rows):
            for column_index, value in enumerate(row_values):
                table.cell(row_index, column_index).text = value

        preview_lines = []
        if headers:
            preview_lines.append(" | ".join(headers[:6]))
        for row_values in data_rows[:3]:
            preview_lines.append(" | ".join(row_values[:6]))

        try:
            _save_docx_via_temp_file(document, target_resolved)
        except PermissionError as exc:
            locked_message = str(exc).strip() or _nonwritable_target_message(
                target_resolved
            )
            try:
                fallback_resolved = _next_available_docx_copy_path(target_resolved)
                _save_docx_via_temp_file(document, fallback_resolved)
            except Exception:
                return _blocked_write_result(
                    _result_path(target_path, target_resolved),
                    summary=locked_message,
                    suggested_next_step=_nonwritable_target_next_step(target_resolved),
                    source_path=_result_path(source_path, source_resolved),
                    sheet=target_sheet,
                    requested_sheet=requested_sheet,
                )

            result_warning = "；".join(
                part
                for part in (
                    sheet_warning,
                    backup_warning,
                    f"原目标文件 {os.path.basename(target_resolved)} 当前不可写，结果已写入同目录副本 {os.path.basename(fallback_resolved)}。",
                )
                if part
            )
            return _blocked_write_result(
                _result_path(fallback_resolved, fallback_resolved),
                summary=(
                    f"原目标文件 {os.path.basename(target_resolved)} 当前不可写，尚未写回原文件；"
                    f"已将工作表“{target_sheet}”的 {len(data_rows)} 行数据写入恢复副本 {os.path.basename(fallback_resolved)}"
                ),
                suggested_next_step=_nonwritable_target_next_step(target_resolved),
                operation="insert_excel_as_docx_table",
                file_type="docx",
                change_type="create",
                preview="\n".join(preview_lines),
                focus=True,
                source_path=_result_path(source_path, source_resolved),
                sheet=target_sheet,
                requested_sheet=requested_sheet,
                warning=result_warning,
                rows_written=len(data_rows),
                columns_written=column_count,
                table_title=table_title,
                table_count=1,
                sort_by=sort_by_text,
                sort_order="desc" if sort_descending else "asc",
                selected_columns=selected_columns,
                original_target_path=_result_path(target_path, target_resolved),
                blocked_target=True,
                blocked_reason=locked_message,
                fallback_copy=True,
            )

        result_warning = "；".join(
            part
            for part in (sheet_warning, sort_warning, selected_warning, style_warning)
            if part
        )
        if backup_warning:
            result_warning = "；".join(
                part for part in (result_warning, backup_warning) if part
            )

        return _success_result(
            _result_path(target_path, target_resolved),
            operation="insert_excel_as_docx_table",
            summary=f"已将工作表“{target_sheet}”的 {len(data_rows)} 行数据写入 Word 表格",
            file_type="docx",
            change_type="modify" if target_exists else "create",
            preview="\n".join(preview_lines),
            focus=True,
            source_path=_result_path(source_path, source_resolved),
            sheet=target_sheet,
            requested_sheet=requested_sheet,
            warning=result_warning,
            rows_written=len(data_rows),
            columns_written=column_count,
            table_title=table_title,
            table_count=1,
            sort_by=sort_by_text,
            sort_order="desc" if sort_descending else "asc",
            selected_columns=selected_columns,
        )
    except ImportError as exc:
        return json.dumps({"error": f"Missing dependency: {exc}"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def insert_image_into_docx(
    path: str,
    image_path: str,
    title: str = "",
    caption: str = "",
    width_inches: Any = 6.5,
) -> str:
    """Append an image or chart into a DOCX file as a real inline picture."""
    resolved = _resolve_path(path)
    if not resolved:
        if os.path.isabs(path):
            resolved = os.path.normpath(path)
        else:
            resolved = _safe_resolve(path)
    if not resolved:
        return json.dumps({"error": f"Invalid path: {path}"}, ensure_ascii=False)

    image_resolved = _resolve_path(image_path)
    if not image_resolved:
        return json.dumps(
            {"error": f"File not found: {image_path}"}, ensure_ascii=False
        )

    title_text = str(title or "").strip()
    caption_text = str(caption or "").strip()
    image_width = _normalize_positive_float(width_inches, default=6.5, upper=10.0)

    try:
        from io import BytesIO

        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches

        os.makedirs(os.path.dirname(resolved), exist_ok=True)
        target_exists = os.path.exists(resolved)
        backup_warning = ""
        write_warning = ""
        if target_exists:
            backup_warning = _best_effort_backup(resolved)
            write_warning = _ensure_existing_file_writable(resolved)
            try:
                with open(resolved, "rb") as existing_doc:
                    document = Document(BytesIO(existing_doc.read()))
            except Exception as open_err:
                return json.dumps(
                    {
                        "error": (
                            f"无法用 python-docx 打开 {os.path.basename(resolved)}"
                            f"（{open_err}）。"
                            "请改用 run_python_code + python-docx 手动插图，"
                            "或检查目标文件是否损坏。"
                        )
                    },
                    ensure_ascii=False,
                )
        else:
            document = Document()

        if document.paragraphs and any(
            paragraph.text.strip() for paragraph in document.paragraphs
        ):
            document.add_paragraph("")

        style_warning = ""
        if title_text:
            title_paragraph = document.add_paragraph(title_text)
            style_warning = _apply_docx_style(title_paragraph, "Heading 2")

        picture_paragraph = document.add_paragraph()
        picture_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        picture_run = picture_paragraph.add_run()
        picture_run.add_picture(image_resolved, width=Inches(image_width))

        if caption_text:
            caption_paragraph = document.add_paragraph(caption_text)
            caption_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            style_warning = _merge_warnings(
                style_warning,
                _apply_docx_style(caption_paragraph, "Caption"),
            )

        image_name = os.path.basename(image_resolved)
        preview_lines: List[str] = []
        if title_text:
            preview_lines.append(title_text)
        preview_lines.append(image_name)
        if caption_text:
            preview_lines.append(caption_text)

        try:
            _save_docx_via_temp_file(document, resolved)
        except PermissionError as exc:
            locked_message = str(exc).strip() or _nonwritable_target_message(resolved)
            try:
                fallback_resolved = _next_available_docx_copy_path(resolved)
                _save_docx_via_temp_file(document, fallback_resolved)
            except Exception:
                return _blocked_write_result(
                    _result_path(path, resolved),
                    summary=locked_message,
                    suggested_next_step=_nonwritable_target_next_step(resolved),
                    image_path=_result_path(image_path, image_resolved),
                    image_name=image_name,
                    images_inserted=1,
                    title=title_text,
                    caption=caption_text,
                )

            result_warning = "；".join(
                part
                for part in (
                    backup_warning,
                    write_warning,
                    style_warning,
                    f"原目标文件 {os.path.basename(resolved)} 当前不可写，结果已写入同目录副本 {os.path.basename(fallback_resolved)}。",
                )
                if part
            )
            return _blocked_write_result(
                _result_path(fallback_resolved, fallback_resolved),
                summary=(
                    f"原目标文件 {os.path.basename(resolved)} 当前不可写，尚未写回原文件；"
                    f"已将图片“{image_name}”插入恢复副本 {os.path.basename(fallback_resolved)}"
                ),
                suggested_next_step=_nonwritable_target_next_step(resolved),
                operation="insert_image_into_docx",
                file_type="docx",
                change_type="create",
                preview="\n".join(preview_lines),
                focus=True,
                image_path=_result_path(image_path, image_resolved),
                image_name=image_name,
                images_inserted=1,
                title=title_text,
                caption=caption_text,
                warning=result_warning,
                original_target_path=_result_path(path, resolved),
                blocked_target=True,
                blocked_reason=locked_message,
                fallback_copy=True,
            )

        return _success_result(
            _result_path(path, resolved),
            operation="insert_image_into_docx",
            summary=f"已将图片“{image_name}”插入 Word 文档",
            file_type="docx",
            change_type="modify" if target_exists else "create",
            preview="\n".join(preview_lines),
            focus=True,
            image_path=_result_path(image_path, image_resolved),
            image_name=image_name,
            images_inserted=1,
            title=title_text,
            caption=caption_text,
            warning=_merge_warnings(backup_warning, write_warning, style_warning),
        )
    except ImportError:
        return json.dumps({"error": "python-docx not installed"}, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def _coerce_read_line_number(value: Any, default: int) -> int:
    if value is None:
        return default
    if isinstance(value, bool):
        return default
    try:
        if isinstance(value, str):
            stripped = value.strip()
            if not stripped:
                return default
            value = float(stripped) if "." in stripped else int(stripped)
        coerced = int(value)
    except (TypeError, ValueError):
        return default
    return max(1, coerced)


def read_file_range(path: str, start_line: int = 1, end_line: int = 100) -> str:
    """Read a specific range of lines from a file.

    Args:
        path: File path.
        start_line: Starting line number (1-indexed).
        end_line: Ending line number (inclusive).

    Returns: The specified lines as text.
    """
    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)

    try:
        start_line = _coerce_read_line_number(start_line, 1)
        end_line = _coerce_read_line_number(end_line, 100)
        if end_line < start_line:
            end_line = start_line

        with open(resolved, "r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()

        start_idx = max(0, start_line - 1)
        end_idx = min(len(lines), end_line)

        selected = lines[start_idx:end_idx]
        return "".join(selected)
    except Exception as e:
        return f"Error: {e}"


def replace_file_selection(
    path: str,
    original_selection: str = "",
    new_content: str = "",
    occurrence: int = 1,
) -> str:
    """Replace a selected text span in an existing text-like file.

    Args:
        path: Target text-like file path.
        original_selection: Exact selected text to replace.
        new_content: Replacement text.
        occurrence: 1-indexed occurrence to replace when the same text appears multiple times.

    Returns: JSON with standard file-change metadata.
    """
    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)

    suffix = Path(resolved).suffix.lstrip(".").lower()
    text_like = {
        "txt",
        "md",
        "markdown",
        "text",
        "csv",
        "json",
        "py",
        "js",
        "html",
        "css",
    }
    if suffix not in text_like:
        return json.dumps(
            {
                "error": f"replace_file_selection only supports text-like files, got: {suffix or 'unknown'}",
                "changed": False,
            },
            ensure_ascii=False,
        )

    original = str(original_selection or "")
    replacement = str(new_content or "")
    if not original:
        return json.dumps(
            {"error": "original_selection is required", "changed": False},
            ensure_ascii=False,
        )

    try:
        index_to_replace = max(1, int(occurrence or 1))
    except (TypeError, ValueError):
        index_to_replace = 1

    try:
        with open(
            resolved, "r", encoding="utf-8", errors="replace", newline=""
        ) as handle:
            content = handle.read()
    except Exception as exc:
        return json.dumps({"error": f"Read failed: {exc}"}, ensure_ascii=False)

    starts = [match.start() for match in re.finditer(re.escape(original), content)]
    if not starts:
        return json.dumps(
            {
                "success": False,
                "changed": False,
                "path": _result_path(path, resolved),
                "operation": "replace_file_selection",
                "summary": "未在目标文件中找到选区原文，未写入。",
                "error": "original_selection_not_found",
            },
            ensure_ascii=False,
        )
    if index_to_replace > len(starts):
        return json.dumps(
            {
                "success": False,
                "changed": False,
                "path": _result_path(path, resolved),
                "operation": "replace_file_selection",
                "summary": f"目标文件中只找到 {len(starts)} 处匹配，无法替换第 {index_to_replace} 处。",
                "error": "occurrence_out_of_range",
                "matches_found": len(starts),
            },
            ensure_ascii=False,
        )

    start = starts[index_to_replace - 1]
    end = start + len(original)
    updated = content[:start] + replacement + content[end:]

    backup_warning = _best_effort_backup(resolved)
    write_warning = _ensure_existing_file_writable(resolved)
    try:
        _write_bytes_via_temp_file(
            updated.encode("utf-8"), resolved, suffix=f".{suffix}" if suffix else ".txt"
        )
    except PermissionError as exc:
        return _blocked_write_result(
            _result_path(path, resolved),
            summary=str(exc).strip() or _nonwritable_target_message(resolved),
            suggested_next_step=_nonwritable_target_next_step(resolved),
            operation="replace_file_selection",
            file_type=suffix or "text",
            preview=replacement,
            replacements_made=0,
        )
    except Exception as exc:
        return json.dumps({"error": f"Write failed: {exc}"}, ensure_ascii=False)

    return _success_result(
        _result_path(path, resolved),
        operation="replace_file_selection",
        summary="已替换目标文件中的选区内容。",
        file_type=suffix or "text",
        change_type="modify",
        preview=replacement,
        focus=True,
        replacements_made=1,
        matches_found=len(starts),
        occurrence=index_to_replace,
        original_selection=original[:240],
        warning=_merge_warnings(backup_warning, write_warning),
    )


def write_docx_content(path: str, paragraphs: str = "[]") -> str:
    """Write paragraphs to a DOCX file.

    Args:
        path: Path to the DOCX file (will be created if not exists).
        paragraphs: JSON array of paragraph objects:
            [{"text": "...", "style": "Heading 1"}, {"text": "..."}]

    Returns: JSON with operation result.
    """
    resolved = _resolve_path(path, must_exist=False)
    if not resolved:
        return json.dumps({"error": f"无效路径: {path}"}, ensure_ascii=False)

    para_list = _coerce_docx_paragraphs_for_write(paragraphs)

    try:
        from docx import Document

        # Create or load document
        file_exists = os.path.exists(resolved)
        backup_warning = ""
        write_warning = ""
        if file_exists:
            doc = Document(resolved)
            original_paragraph_count = len(doc.paragraphs)
            # Backup
            backup_warning = _best_effort_backup(resolved)
            write_warning = _ensure_existing_file_writable(resolved)
        else:
            doc = Document()
            original_paragraph_count = 0
            os.makedirs(os.path.dirname(resolved), exist_ok=True)

        if not para_list:
            para_list = [{"text": "", "style": "Normal"}]

        style_warnings: List[str] = []
        for p in para_list:
            text = p.get("text", "")
            style = p.get("style")
            para = doc.add_paragraph(text)
            if style:
                style_warnings.append(_apply_docx_style(para, style))

        _save_docx_via_temp_file(doc, resolved)
        preview = "\n".join(str(p.get("text", "")) for p in para_list[:3])
        diff_items = [
            {
                "paragraph_index": original_paragraph_count + index,
                "before": "",
                "after": str(item.get("text", "")),
                "style": str(item.get("style") or ""),
            }
            for index, item in enumerate(para_list, start=1)
        ]
        return _success_result(
            _result_path(path, resolved),
            operation="write_docx_content",
            summary=f"已写入 {len(para_list)} 个段落到 Word 文档",
            file_type="docx",
            change_type="modify" if file_exists else "create",
            preview=preview,
            focus=True,
            summary_code="WRITE_OK" if file_exists else "CREATE_OK",
            diff=_file_task_diff("docx_paragraphs", diff_items),
            paragraphs_written=len(para_list),
            warning=_merge_warnings(backup_warning, write_warning, *style_warnings),
        )
    except ImportError:
        try:
            return _write_docx_content_without_python_docx(path, resolved, para_list)
        except PermissionError as exc:
            preview = "\n".join(
                str(p.get("text", ""))
                for p in _normalize_docx_paragraphs(para_list)[:3]
            )
            return _blocked_write_result(
                _result_path(path, resolved),
                summary=str(exc).strip() or _nonwritable_target_message(resolved),
                suggested_next_step=_nonwritable_target_next_step(resolved),
                operation="write_docx_content",
                file_type="docx",
                preview=preview,
                paragraphs_written=len(_normalize_docx_paragraphs(para_list)),
            )
        except Exception as exc:
            return json.dumps({"error": str(exc)}, ensure_ascii=False)
    except PermissionError as exc:
        preview = "\n".join(str(p.get("text", "")) for p in para_list[:3])
        return _blocked_write_result(
            _result_path(path, resolved),
            summary=str(exc).strip() or _nonwritable_target_message(resolved),
            suggested_next_step=_nonwritable_target_next_step(resolved),
            operation="write_docx_content",
            file_type="docx",
            preview=preview,
            paragraphs_written=len(para_list),
        )
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def insert_docx_paragraph(
    path: str,
    text: str,
    after_heading: str = "",
    before_heading: str = "",
    style: str = "Normal",
) -> str:
    """Insert a single paragraph into an existing DOCX without rewriting the document."""
    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)
    clean_text = str(text or "").strip()
    if not clean_text:
        return json.dumps({"error": "text is required"}, ensure_ascii=False)

    try:
        from docx import Document
        from docx.oxml import OxmlElement
        from docx.text.paragraph import Paragraph

        doc = Document(resolved)
        backup_warning = _best_effort_backup(resolved)
        write_warning = _ensure_existing_file_writable(resolved)
        after = str(after_heading or "").strip().casefold()
        before = str(before_heading or "").strip().casefold()
        insertion_index = len(doc.paragraphs)

        if before:
            for index, paragraph in enumerate(doc.paragraphs):
                if str(paragraph.text or "").strip().casefold() == before:
                    insertion_index = index
                    break
        elif after:
            for index, paragraph in enumerate(doc.paragraphs):
                if str(paragraph.text or "").strip().casefold() == after:
                    insertion_index = index + 1
                    break

        if insertion_index >= len(doc.paragraphs):
            paragraph = doc.add_paragraph(clean_text)
        else:
            anchor = doc.paragraphs[insertion_index]
            new_p = OxmlElement("w:p")
            anchor._p.addprevious(new_p)
            paragraph = Paragraph(new_p, anchor._parent)
            paragraph.add_run(clean_text)

        style_warning = _apply_docx_style(paragraph, style)

        _save_docx_via_temp_file(doc, resolved)
        return _success_result(
            _result_path(path, resolved),
            operation="insert_docx_paragraph",
            summary="已向 Word 文档插入 1 个段落",
            file_type="docx",
            change_type="modify",
            preview=clean_text,
            paragraphs_written=1,
            focus=True,
            inserted_text=clean_text,
            after_heading=after_heading,
            before_heading=before_heading,
            warning=_merge_warnings(backup_warning, write_warning, style_warning),
            diff=_file_task_diff(
                "docx_paragraphs",
                [
                    {
                        "paragraph_index": insertion_index + 1,
                        "before": "",
                        "after": clean_text,
                        "style": style or "",
                    }
                ],
            ),
        )
    except ImportError:
        return json.dumps({"error": "python-docx not installed"}, ensure_ascii=False)
    except PermissionError as exc:
        return _blocked_write_result(
            _result_path(path, resolved),
            summary=str(exc).strip() or _nonwritable_target_message(resolved),
            suggested_next_step=_nonwritable_target_next_step(resolved),
            operation="insert_docx_paragraph",
            file_type="docx",
            preview=clean_text,
        )
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def _resolve_output_path(
    source_path: str,
    source_resolved: str,
    target_path: str,
    suffix: str,
) -> Optional[str]:
    requested = str(target_path or "").strip()
    if requested:
        if os.path.isabs(requested):
            return os.path.normpath(requested)
        resolved = _safe_resolve(requested)
        if not resolved:
            return None
        return resolved
    base = Path(source_resolved or source_path).with_suffix(suffix)
    return str(base)


def _convert_docx_to_pdf_with_docx2pdf(source: str, target: str) -> str:
    return _conversion_docx2pdf(source, target)


def _convert_docx_to_pdf_with_word(source: str, target: str) -> str:
    return _conversion_word(source, target)


def _convert_docx_to_pdf_with_libreoffice(source: str, target: str) -> str:
    return _conversion_libreoffice(source, target)


def convert_docx_to_pdf(path: str, target_path: str = "") -> str:
    return _conversion_convert_docx_to_pdf(
        path,
        target_path,
        resolve_path=_resolve_path,
        resolve_output_path=_resolve_output_path,
        result_path=_result_path,
        success_result=_success_result,
        blocked_write_result=_blocked_write_result,
        converters=(
            _convert_docx_to_pdf_with_docx2pdf,
            _convert_docx_to_pdf_with_word,
            _convert_docx_to_pdf_with_libreoffice,
        ),
    )


def _normalize_conversion_extension(target_format: str) -> str:
    return _conversion_normalize_extension(target_format)


def convert_file(file_path: str, target_format: str, output_path: str = "") -> str:
    return _conversion_convert_file(
        file_path,
        target_format,
        output_path,
        resolve_path=_resolve_path,
        resolve_output_path=_resolve_output_path,
        result_path=_result_path,
        success_result=_success_result,
        blocked_write_result=_blocked_write_result,
    )


def list_conversions(file_ext: str = "") -> str:
    return _conversion_list_conversions(file_ext)


def fill_docx_template(
    path: str,
    data: str = "{}",
    target_path: str = "",
    placeholder_style: str = "braces",
) -> str:
    """Fill simple placeholders in a DOCX template.

    `data` is a JSON object. With placeholder_style=braces, field `name`
    replaces both `{{name}}` and `{name}`.
    """
    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)
    try:
        values = json.loads(data) if isinstance(data, str) else data
    except json.JSONDecodeError as exc:
        return json.dumps({"error": f"Invalid data JSON: {exc}"}, ensure_ascii=False)
    if not isinstance(values, dict):
        return json.dumps({"error": "Template data must be a JSON object"}, ensure_ascii=False)

    try:
        from docx import Document

        target = _resolve_output_path(path, resolved, target_path, ".docx")
        if not target:
            return json.dumps({"error": f"Invalid target path: {target_path}"}, ensure_ascii=False)
        file_exists = os.path.exists(target)
        if target != resolved:
            os.makedirs(os.path.dirname(target), exist_ok=True)
            shutil.copy2(resolved, target)

        backup_warning = _best_effort_backup(target) if file_exists else ""
        write_warning = _ensure_existing_file_writable(target) if file_exists else ""
        doc = Document(target)
        replacements: Dict[str, str] = {}
        for raw_key, raw_value in values.items():
            key = str(raw_key)
            value = "" if raw_value is None else str(raw_value)
            replacements[f"{{{{{key}}}}}"] = value
            if str(placeholder_style or "").strip().lower() in {"braces", "single"}:
                replacements[f"{{{key}}}"] = value

        diff_items: List[Dict[str, Any]] = []
        fields_seen: set[str] = set()
        for index, paragraph in enumerate(doc.paragraphs, start=1):
            changed, before, after = _replace_docx_placeholders_in_paragraph(
                paragraph, replacements
            )
            if changed:
                diff_items.append(
                    {
                        "location": "paragraph",
                        "paragraph_index": index,
                        "before": before,
                        "after": after,
                    }
                )
        for table_index, table in enumerate(doc.tables, start=1):
            for row_index, row in enumerate(table.rows, start=1):
                for col_index, cell in enumerate(row.cells, start=1):
                    for paragraph_index, paragraph in enumerate(cell.paragraphs, start=1):
                        changed, before, after = _replace_docx_placeholders_in_paragraph(
                            paragraph, replacements
                        )
                        if changed:
                            diff_items.append(
                                {
                                    "location": "table_cell",
                                    "table_index": table_index,
                                    "row": row_index,
                                    "col": col_index,
                                    "paragraph_index": paragraph_index,
                                    "before": before,
                                    "after": after,
                                }
                            )
        for item in diff_items:
            before = str(item.get("before") or "")
            for raw_key in values:
                key = str(raw_key)
                if f"{{{{{key}}}}}" in before or f"{{{key}}}" in before:
                    fields_seen.add(key)

        _save_docx_via_temp_file(doc, target)
        return _success_result(
            _result_path(target_path or path, target),
            operation="fill_docx_template",
            summary=f"已填充 Word 模板中的 {len(fields_seen)} 个字段",
            file_type="docx",
            change_type="modify" if file_exists else "create",
            summary_code="WRITE_OK",
            diff=_file_task_diff("docx_template_fields", diff_items),
            source_path=_result_path(path, resolved),
            fields_filled=sorted(fields_seen),
            placeholders_replaced=len(diff_items),
            warning=_merge_warnings(backup_warning, write_warning),
            focus=True,
        )
    except ImportError:
        return json.dumps({"error": "python-docx not installed"}, ensure_ascii=False)
    except PermissionError as exc:
        target = _resolve_output_path(path, resolved, target_path, ".docx") or resolved
        return _blocked_write_result(
            _result_path(target_path or path, target),
            summary=str(exc).strip() or _nonwritable_target_message(target),
            suggested_next_step=_nonwritable_target_next_step(target),
            operation="fill_docx_template",
            file_type="docx",
        )
    except Exception as exc:
        return json.dumps({"error": str(exc)}, ensure_ascii=False)


def write_pptx_slides(path: str, updates: str = "[]") -> str:
    """Update text in existing PPTX slides.

    updates: JSON array of:
      [{"slide_index": 0, "shape_name": "标题 1", "text": "新内容"}, ...]
    or
      [{"slide_index": 0, "shape_index": 0, "text": "新内容"}, ...]
    slide_index is 0-based.
    """
    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)

    try:
        upd_list = json.loads(updates) if isinstance(updates, str) else updates
    except json.JSONDecodeError as e:
        return json.dumps({"error": f"Invalid updates JSON: {e}"}, ensure_ascii=False)

    try:
        from pptx import Presentation
        from pptx.util import Pt

        backup_warning = _best_effort_backup(resolved)
        write_warning = _ensure_existing_file_writable(resolved)
        prs = Presentation(resolved)
        slides_updated = 0

        for upd in upd_list:
            slide_idx = int(upd.get("slide_index", 0))
            if slide_idx < 0 or slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]
            new_text = str(upd.get("text", ""))

            shape_name = upd.get("shape_name")
            shape_index = upd.get("shape_index")
            target_shape = None

            for i, shape in enumerate(slide.shapes):
                if shape_name and shape.name == shape_name:
                    target_shape = shape
                    break
                if shape_index is not None and i == int(shape_index):
                    target_shape = shape
                    break

            if target_shape and target_shape.has_text_frame:
                # Preserve first run's font settings; replace paragraph text
                tf = target_shape.text_frame
                if tf.paragraphs:
                    para = tf.paragraphs[0]
                    # Keep existing run formatting, just replace text
                    if para.runs:
                        para.runs[0].text = new_text
                        for run in para.runs[1:]:
                            run.text = ""
                    else:
                        para.text = new_text
                slides_updated += 1

        _save_pptx_via_temp_file(prs, resolved)
        return _success_result(
            _result_path(path, resolved),
            operation="write_pptx_slides",
            summary=f"已更新 {slides_updated} 个形状的文字内容",
            file_type="pptx",
            change_type="modify",
            slides_updated=slides_updated,
            warning=_merge_warnings(backup_warning, write_warning),
        )
    except ImportError:
        return json.dumps(
            {
                "error": "python-pptx not installed. Use run_python_code with python-pptx instead."
            },
            ensure_ascii=False,
        )
    except PermissionError as exc:
        return _blocked_write_result(
            _result_path(path, resolved),
            summary=str(exc).strip() or _nonwritable_target_message(resolved),
            suggested_next_step=_nonwritable_target_next_step(resolved),
            operation="write_pptx_slides",
            file_type="pptx",
            slides_updated=slides_updated,
        )
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def design_pptx_theme_layout(
    path: str,
    style_brief: str = "",
    theme: str = "",
    palette: str = "",
    typography: str = "",
    density: str = "balanced",
    preserve_content: bool = True,
) -> str:
    """Apply a conservative visual theme and placeholder layout pass to a PPTX file."""
    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)

    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches

        selected_theme = _select_pptx_theme(style_brief, theme, palette, typography)
        density_settings = _pptx_density_settings(density)
        backup_warning = _best_effort_backup(resolved)
        write_warning = _ensure_existing_file_writable(resolved)

        prs = Presentation(resolved)
        original_slide_count = len(prs.slides)
        if original_slide_count == 0:
            return json.dumps(
                {"error": "PPTX has no slides to design"}, ensure_ascii=False
            )

        slide_width = int(prs.slide_width)
        slide_height = int(prs.slide_height)
        margin_x = Inches(density_settings["margin_x"])
        title_top = Inches(density_settings["title_top"])
        body_top = Inches(density_settings["body_top"])
        title_height = Inches(0.72)
        body_height = max(slide_height - body_top - Inches(0.62), Inches(1.2))

        title_color = _hex_to_rgb_color(
            (
                selected_theme["inverse_text"]
                if selected_theme.get("is_dark")
                else selected_theme["primary"]
            ),
            "17324D",
        )
        body_color = _hex_to_rgb_color(selected_theme["body_text"], "25313B")
        background_color = _hex_to_rgb_color(selected_theme["background"], "F7F3EA")
        font_family = str(selected_theme.get("font_family") or "Microsoft YaHei")
        layout_warnings: List[str] = []
        text_shapes_styled = 0
        title_shapes_styled = 0
        body_placeholders_adjusted = 0

        for slide_index, slide in enumerate(prs.slides, start=1):
            _remove_koto_theme_shapes(slide)
            _set_slide_background(slide, background_color)
            _add_theme_background_shape(
                slide, slide_width, slide_height, background_color
            )
            _add_theme_accent_shapes(
                slide, slide_width, slide_height, selected_theme, slide_index
            )

            slide_text_chars = 0
            for shape in slide.shapes:
                if str(getattr(shape, "name", "") or "").startswith("KOTO_THEME_"):
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = str(getattr(shape, "text", "") or "")
                slide_text_chars += len(text)
                if _is_title_shape(slide, shape):
                    if preserve_content:
                        shape.left = margin_x
                        shape.top = title_top
                        shape.width = max(slide_width - (margin_x * 2), Inches(4.0))
                        shape.height = title_height
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.LEFT
                    _apply_text_style(
                        shape.text_frame,
                        font_family=font_family,
                        size_pt=density_settings["title_size"],
                        color=title_color,
                        bold=True,
                    )
                    title_shapes_styled += 1
                else:
                    if preserve_content and _is_body_placeholder(shape):
                        shape.left = margin_x
                        shape.top = body_top
                        shape.width = max(slide_width - (margin_x * 2), Inches(4.0))
                        shape.height = body_height
                        body_placeholders_adjusted += 1
                    _apply_text_style(
                        shape.text_frame,
                        font_family=font_family,
                        size_pt=density_settings["body_size"],
                        color=body_color,
                        bold=False,
                    )
                text_shapes_styled += 1

            if slide_text_chars > 950:
                layout_warnings.append(f"第 {slide_index} 页文本较多，建议人工复核拥挤度")

        _save_pptx_via_temp_file(prs, resolved)
        reopened = Presentation(resolved)
        if len(reopened.slides) != original_slide_count:
            return json.dumps(
                {"error": "PPTX validation failed: slide count changed unexpectedly"},
                ensure_ascii=False,
            )

        preview_lines = [
            f"主题：{selected_theme.get('display_name')}",
            f"字体：{font_family}",
            f"已处理幻灯片：{original_slide_count}",
        ]
        if layout_warnings:
            preview_lines.append("检查提示：" + "；".join(layout_warnings[:3]))

        return _success_result(
            _result_path(path, resolved),
            operation="design_pptx_theme_layout",
            summary=f"已为 {original_slide_count} 张幻灯片应用统一主题、字体和安全版式",
            file_type="pptx",
            change_type="modify",
            preview="\n".join(preview_lines),
            focus=True,
            slides_designed=original_slide_count,
            total_slides=original_slide_count,
            theme_name=selected_theme.get("display_name") or selected_theme.get("name"),
            layout_strategy="safe_placeholder_grid",
            font_family=font_family,
            title_shapes_styled=title_shapes_styled,
            text_shapes_styled=text_shapes_styled,
            body_placeholders_adjusted=body_placeholders_adjusted,
            layout_warnings=layout_warnings,
            warning=_merge_warnings(backup_warning, write_warning),
        )
    except ImportError:
        return json.dumps(
            {
                "error": "python-pptx not installed. Install python-pptx before designing PPTX layout."
            },
            ensure_ascii=False,
        )
    except PermissionError as exc:
        return _blocked_write_result(
            _result_path(path, resolved),
            summary=str(exc).strip() or _nonwritable_target_message(resolved),
            suggested_next_step=_nonwritable_target_next_step(resolved),
            operation="design_pptx_theme_layout",
            file_type="pptx",
        )
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


def add_pptx_slides(path: str, slides: str = "[]") -> str:
    """Append new slides to an existing PPTX file.

    slides: JSON array of:
      [{"title": "幻灯片标题", "content": "第一行\\n第二行\\n第三行", "layout_index": 1}, ...]
    content may also be a list of strings or bullet dictionaries.
    layout_index: 0=空白, 1=标题+内容(默认), 2=章节标题, etc.
    """
    resolved = _resolve_path(path)
    if not resolved:
        return json.dumps({"error": f"File not found: {path}"}, ensure_ascii=False)

    slides_list, parse_error = _parse_jsonish_list(slides, "slides")
    if parse_error:
        return json.dumps({"error": parse_error}, ensure_ascii=False)

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        backup_warning = _best_effort_backup(resolved)
        write_warning = _ensure_existing_file_writable(resolved)
        prs = Presentation(resolved)
        slides_added = 0

        for slide_data in slides_list:
            if not isinstance(slide_data, dict):
                slide_data = {"content": slide_data}
            layout_idx = int(slide_data.get("layout_index", 1))
            layout_idx = min(layout_idx, len(prs.slide_layouts) - 1)
            layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(layout)

            title_text = _pptx_first_text(slide_data.get("title", ""))
            content_value = slide_data.get("content")
            if content_value in (None, ""):
                for fallback_key in (
                    "bullets",
                    "bullet_points",
                    "points",
                    "items",
                    "lines",
                ):
                    if fallback_key in slide_data:
                        content_value = slide_data.get(fallback_key)
                        break
            content_lines = _pptx_text_lines(content_value)

            # Set title placeholder if available
            if slide.shapes.title and title_text:
                slide.shapes.title.text = title_text

            # Set body/content placeholder if available
            body_written = False
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1 and content_lines:
                    tf = ph.text_frame
                    tf.clear()
                    for i, line in enumerate(content_lines):
                        if i == 0:
                            tf.paragraphs[0].text = line
                        else:
                            p = tf.add_paragraph()
                            p.text = line
                    body_written = True
                    break

            if content_lines and not body_written:
                box = slide.shapes.add_textbox(
                    Inches(0.9), Inches(1.55), Inches(8.2), Inches(4.8)
                )
                tf = box.text_frame
                tf.clear()
                for i, line in enumerate(content_lines):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(20 if len(content_lines) <= 4 else 16)

            slides_added += 1

        total_slides = len(prs.slides)
        _save_pptx_via_temp_file(prs, resolved)
        return _success_result(
            _result_path(path, resolved),
            operation="add_pptx_slides",
            summary=f"已新增 {slides_added} 张幻灯片，当前共 {total_slides} 张",
            file_type="pptx",
            change_type="modify",
            slides_added=slides_added,
            total_slides=total_slides,
            warning=_merge_warnings(backup_warning, write_warning),
        )
    except ImportError:
        return json.dumps(
            {"error": "python-pptx not installed. Use run_python_code instead."},
            ensure_ascii=False,
        )
    except PermissionError as exc:
        return _blocked_write_result(
            _result_path(path, resolved),
            summary=str(exc).strip() or _nonwritable_target_message(resolved),
            suggested_next_step=_nonwritable_target_next_step(resolved),
            operation="add_pptx_slides",
            file_type="pptx",
            slides_added=slides_added if "slides_added" in locals() else 0,
            total_slides=len(prs.slides) if "prs" in locals() else 0,
        )
    except Exception as e:
        return json.dumps({"error": str(e)}, ensure_ascii=False)


# ══════════════════════════════════════════════════════════════
# TaskToolsPlugin — registers all tools above into ToolRegistry
# ══════════════════════════════════════════════════════════════


class TaskToolsPlugin(AgentPlugin):
    """File-focused tools for the Koto file-task runtime."""

    def __init__(
        self,
        socketio=None,
        task_files: Optional[List[Dict[str, str]]] = None,
        gemini_client: Any = None,
        workspace_root: str = "",
        request_context: Optional[Dict[str, Any]] = None,
    ):
        self._socketio = socketio
        self._task_files = list(task_files or [])
        self._gemini_client = gemini_client
        self._workspace_root = str(workspace_root or "").strip()
        self._request_context = dict(request_context or {})

    def _run_python_code(self, code: str, timeout: int = 30) -> str:
        return run_python_in_sandbox(
            code,
            timeout=timeout,
            task_files=self._task_files,
            target_path=str(self._request_context.get("target_path") or ""),
            output_dir=self._contextual_output_directory(),
        )

    def _contextual_output_directory(self) -> str:
        target_path = str(self._request_context.get("target_path") or "").strip()
        if target_path:
            normalized_target = target_path.replace("\\", "/").strip("/")
            if "/" in normalized_target:
                parent = str(Path(normalized_target).parent).replace("\\", "/")
                if parent and parent != ".":
                    return parent

        file_dirs: set[str] = set()
        for item in self._task_files:
            if not isinstance(item, dict):
                continue
            path_text = str(item.get("path") or "").strip()
            if not path_text:
                continue
            parent = str(Path(path_text.replace("\\", "/")).parent).replace("\\", "/")
            if parent not in {"", "."}:
                file_dirs.add(parent)
        if len(file_dirs) == 1:
            return next(iter(file_dirs))

        task_text = str(self._request_context.get("task") or "").casefold()
        if not task_text:
            return ""
        task_dir = self._contextual_output_directory_from_task_text(task_text)
        if task_dir:
            return task_dir
        workspace_root = Path(self._workspace_root or _get_workspace_root())
        try:
            matches = [
                child.name
                for child in workspace_root.iterdir()
                if child.is_dir() and child.name.casefold() in task_text
            ]
        except Exception:
            matches = []
        unique_matches = sorted(set(matches))
        if len(unique_matches) == 1:
            return unique_matches[0]
        return ""

    def _contextual_output_directory_from_task_text(self, task_text: str) -> str:
        candidates: set[str] = set()
        for match in re.finditer(
            r"(?P<dir>[A-Za-z0-9_. -]+)/(?:[A-Za-z0-9_. -]+)\.(?:csv|xlsx|xlsm|docx|md|txt|json|pptx|pdf)",
            task_text,
            re.IGNORECASE,
        ):
            directory = str(match.group("dir") or "").strip().strip("/")
            if directory and not Path(directory).suffix:
                candidates.add(directory)
        workspace_root = Path(self._workspace_root or _get_workspace_root())
        existing = []
        for candidate in candidates:
            try:
                target = (workspace_root / candidate).resolve()
                target.relative_to(workspace_root.resolve())
            except Exception:
                continue
            if target.is_dir():
                existing.append(candidate.replace("\\", "/"))
        unique_existing = sorted(set(existing))
        return unique_existing[0] if len(unique_existing) == 1 else ""

    def _contextual_create_file_path(self, path: str) -> str:
        raw_path = str(path or "").strip()
        if not raw_path:
            return raw_path
        normalized = raw_path.replace("\\", "/").strip()
        if os.path.isabs(raw_path) or "/" in normalized.strip("/"):
            return raw_path
        output_dir = self._contextual_output_directory()
        if not output_dir:
            return raw_path
        return str(Path(output_dir.replace("\\", "/")) / normalized).replace("\\", "/")

    def _create_file(self, path: str, content: str = "") -> str:
        return create_file(self._contextual_create_file_path(path), content)

    def _annotate_file(
        self,
        path: str,
        annotations: Any = "[]",
        requirement: str = "",
        model_id: str = "",
    ) -> Any:
        requirement_text = str(requirement or "").strip()
        try:
            ann_list = _parse_annotations_payload(annotations)
        except ValueError:
            ann_list = None
        if (
            Path(str(path or "")).suffix.lstrip(".").lower() == "docx"
            and requirement_text
            and not ann_list
        ):
            return _stream_docx_annotation_tool_result(
                path,
                requirement=requirement_text,
                model_id=model_id,
                gemini_client=self._gemini_client,
                workspace_root=self._workspace_root,
                task_files=self._task_files,
                request_context=self._request_context,
            )
        return annotate_file(
            path,
            annotations,
            requirement=requirement_text,
            model_id=model_id,
            gemini_client=self._gemini_client,
        )

    @property
    def name(self) -> str:
        return "TaskTools"

    @property
    def description(self) -> str:
        return "Composable file-operation tools for dynamic task execution."

    def get_tools(self) -> List[Dict[str, Any]]:
        return build_task_tool_definitions(self, build_task_tool_operations())

    def _editor_live_update(self, type: str, **kwargs) -> str:
        """Push a change to the live editor via WebSocket."""
        if not self._socketio:
            return "Error: no WebSocket connection"
        payload = {"type": type, **kwargs}
        try:
            self._socketio.emit(
                "agent_execute_command",
                {"action": "editor_apply", "payload": payload},
                namespace="/doc",
            )
            return f"Applied: {type}"
        except Exception as e:
            return json.dumps({"error": str(e)}, ensure_ascii=False)

# XLSX functions extracted to task_tools_xlsx.py
from .task_tools_xlsx import (  # noqa: E402, F401
    read_sheet_data,
    inspect_workbook_structure,
    audit_financial_workbook,
    write_sheet_data,
)
