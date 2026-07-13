"""Office-package creation operations used by :mod:`task_tools`.

The public tool module owns path validation and the stable JSON response
contract.  This module owns only the format-specific package construction so
that DOCX, XLSX, and PPTX dependencies do not keep accumulating in the tool
registry module.
"""

from __future__ import annotations

import csv
import io
import json
import os
from typing import Any, Callable, Dict, List


def create_docx_file(
    path: str,
    resolved: str,
    content: str,
    *,
    plain_text_to_paragraphs: Callable[[str], List[Dict[str, Any]]],
    save_document: Callable[[Any, str], None],
    fallback_writer: Callable[[str, str, Any], str],
    success_result: Callable[..., str],
    result_path: Callable[[str, str], str],
    file_task_diff: Callable[[str, List[Dict[str, Any]]], Dict[str, Any]],
) -> str:
    """Create a DOCX package, preserving the task-result response contract."""
    paragraphs = plain_text_to_paragraphs(content)
    if not paragraphs:
        paragraphs = [{"text": "", "style": "Normal"}]
    try:
        from docx import Document

        doc = Document()
        os.makedirs(os.path.dirname(resolved), exist_ok=True)
        for item in paragraphs:
            paragraph = doc.add_paragraph(str(item.get("text") or ""))
            style = str(item.get("style") or "").strip()
            if style:
                try:
                    paragraph.style = style
                except Exception:
                    pass
        save_document(doc, resolved)
        preview = "\n".join(str(item.get("text") or "") for item in paragraphs[:3])
        diff_items = [
            {
                "paragraph_index": index,
                "before": "",
                "after": str(item.get("text") or ""),
                "style": str(item.get("style") or ""),
            }
            for index, item in enumerate(paragraphs, start=1)
        ]
        return success_result(
            result_path(path, resolved),
            operation="write_docx_content",
            summary=f"已创建并写入 {len(paragraphs)} 个段落到 Word 文档",
            file_type="docx",
            change_type="create",
            preview=preview,
            focus=True,
            summary_code="CREATE_OK",
            diff=file_task_diff("docx_paragraphs", diff_items),
            paragraphs_written=len(paragraphs),
        )
    except ImportError:
        return fallback_writer(path, resolved, paragraphs)


def create_xlsx_file(
    path: str,
    resolved: str,
    content: str,
    *,
    save_workbook: Callable[[Any, str], None],
    success_result: Callable[..., str],
    result_path: Callable[[str, str], str],
    file_task_diff: Callable[[str, List[Dict[str, Any]]], Dict[str, Any]],
) -> str:
    """Create an XLSX package from CSV-like tool content."""
    try:
        import openpyxl

        workbook = openpyxl.Workbook()
        worksheet = workbook.active
        worksheet.title = "Sheet1"
        rows_written = 0
        columns_written = 0
        text = str(content or "").strip()
        rows = list(csv.reader(io.StringIO(text))) if text else []
        if not rows:
            rows = [[""]]
        diff_items: List[Dict[str, Any]] = []
        for row_index, row in enumerate(rows, start=1):
            columns_written = max(columns_written, len(row))
            for col_index, value in enumerate(row, start=1):
                worksheet.cell(row=row_index, column=col_index, value=value)
                diff_items.append(
                    {
                        "sheet": worksheet.title,
                        "cell": worksheet.cell(
                            row=row_index, column=col_index
                        ).coordinate,
                        "row": row_index,
                        "col": col_index,
                        "before": None,
                        "after": value,
                    }
                )
            rows_written += 1
        os.makedirs(os.path.dirname(resolved), exist_ok=True)
        save_workbook(workbook, resolved)
        workbook.close()
        cells_written = rows_written * max(columns_written, 1)
        return success_result(
            result_path(path, resolved),
            operation="write_sheet_data",
            summary=f"已创建工作簿并写入 {rows_written} 行、{columns_written} 列",
            file_type="xlsx",
            change_type="create",
            preview=text,
            focus=True,
            summary_code="CREATE_OK",
            diff=file_task_diff("xlsx_cells", diff_items),
            rows_written=rows_written,
            columns_written=columns_written,
            cells_written=cells_written,
        )
    except ImportError:
        return json.dumps({"error": "openpyxl not installed"}, ensure_ascii=False)


def plain_text_to_pptx_slides(content: str) -> List[Dict[str, Any]]:
    """Turn the lightweight markdown accepted by ``create_file`` into slides."""
    import re

    slides: List[Dict[str, Any]] = []
    current: Dict[str, Any] | None = None
    for raw_line in str(content or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        heading_match = re.match(r"^#{1,6}\s+(.+)$", line)
        if heading_match:
            if current:
                slides.append(current)
            current = {"title": heading_match.group(1).strip(), "content": []}
            continue
        bullet_match = re.match(r"^[-*•]\s+(.+)$", line)
        text = bullet_match.group(1).strip() if bullet_match else line
        if current is None:
            current = {"title": text[:56] or "新幻灯片", "content": []}
            if text:
                continue
        current.setdefault("content", []).append(text)
    if current:
        slides.append(current)
    return slides or [{"title": "新幻灯片", "content": []}]


def create_pptx_file(
    path: str,
    resolved: str,
    content: str,
    *,
    text_lines: Callable[[Any], List[str]],
    save_presentation: Callable[[Any, str], None],
    success_result: Callable[..., str],
    result_path: Callable[[str, str], str],
) -> str:
    """Create a PPTX package while keeping presentation concerns isolated."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        presentation = Presentation()
        slides = plain_text_to_pptx_slides(content)
        first_blank = len(presentation.slides) == 0
        for index, slide_data in enumerate(slides):
            layout_index = (
                0
                if index == 0 and first_blank
                else min(1, len(presentation.slide_layouts) - 1)
            )
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[layout_index]
            )
            title_text = str(slide_data.get("title") or "新幻灯片").strip()
            content_lines = text_lines(slide_data.get("content"))
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            else:
                title_box = slide.shapes.add_textbox(
                    Inches(0.6), Inches(0.4), Inches(8.4), Inches(0.7)
                )
                title_frame = title_box.text_frame
                title_frame.clear()
                title_frame.paragraphs[0].text = title_text
                title_frame.paragraphs[0].font.size = Pt(30)
            if content_lines:
                body_box = slide.shapes.add_textbox(
                    Inches(0.8), Inches(1.4), Inches(8.2), Inches(4.8)
                )
                body_frame = body_box.text_frame
                body_frame.clear()
                for line_index, line in enumerate(content_lines):
                    paragraph = (
                        body_frame.paragraphs[0]
                        if line_index == 0
                        else body_frame.add_paragraph()
                    )
                    paragraph.text = line
                    paragraph.font.size = Pt(18)
        os.makedirs(os.path.dirname(resolved), exist_ok=True)
        save_presentation(presentation, resolved)
        preview = "\n".join(str(slide.get("title") or "") for slide in slides[:3])
        return success_result(
            result_path(path, resolved),
            operation="add_pptx_slides",
            summary=f"已创建 PPT 并新增 {len(slides)} 张幻灯片",
            file_type="pptx",
            change_type="create",
            preview=preview,
            focus=True,
            slides_added=len(slides),
            total_slides=len(presentation.slides),
        )
    except ImportError:
        return json.dumps({"error": "python-pptx not installed"}, ensure_ascii=False)
