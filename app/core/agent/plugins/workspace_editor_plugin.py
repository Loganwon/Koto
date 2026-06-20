# Copyright (C) 2024-2026 Koto AI. All rights reserved.
# SPDX-License-Identifier: LicenseRef-Koto-Proprietary
"""
WorkspaceEditorPlugin — Agent ↔ 前端编辑器桥梁

Provides tools that let the Agent:
  1. List/read/create/save workspace files (via FileRegistry)
  2. Push edits to the live frontend editor (via SocketIO → /doc)
  3. Open files in new editor tabs

The plugin talks to the frontend by emitting WebSocket events on the
``/doc`` namespace.  The frontend SocketBridge listens for these events
and routes them to the appropriate DocController.
"""

import json
import logging
import os
from pathlib import Path
from typing import Any, Dict, List, Optional

from app.core.agent.base import AgentPlugin

logger = logging.getLogger(__name__)

# Workspace root — set at app startup; defaults to <project>/workspace
_WORKSPACE_ROOT: Optional[str] = None


def set_workspace_root(path: str):
    """Call once at app startup to set the workspace directory."""
    global _WORKSPACE_ROOT
    _WORKSPACE_ROOT = os.path.abspath(path)


def _get_workspace_root() -> str:
    """Resolve workspace root lazily."""
    global _WORKSPACE_ROOT
    if _WORKSPACE_ROOT is None:
        project_root = Path(__file__).parent.parent.parent.parent
        _WORKSPACE_ROOT = str(project_root / "workspace")
    return _WORKSPACE_ROOT


def _safe_resolve(relative_path: str) -> Optional[str]:
    """Resolve a user-provided path within the workspace root.

    Returns None if the resolved path escapes the workspace (directory
    traversal protection).
    """
    root = _get_workspace_root()
    try:
        resolved = os.path.normpath(os.path.join(root, relative_path))
        if not resolved.startswith(os.path.normpath(root)):
            return None
        return resolved
    except (ValueError, TypeError):
        return None


class WorkspaceEditorPlugin(AgentPlugin):
    """Agent tools bridging the workspace file system and the live editor UI."""

    def __init__(self, socketio=None):
        # socketio instance is injected so we can emit events to the frontend.
        # If None, editor_apply / editor_open_file become no-ops (unit-test safe).
        self._socketio = socketio

    @property
    def name(self) -> str:
        return "WorkspaceEditor"

    @property
    def description(self) -> str:
        return (
            "Tools to interact with the Koto workspace file system and "
            "push edits to the live document editor in the browser."
        )

    def get_tools(self) -> List[Dict[str, Any]]:
        return [
            {
                "name": "workspace_list_files",
                "func": self.workspace_list_files,
                "description": (
                    "List files in the workspace directory.  Returns a JSON array "
                    "of {name, type, size, mtime} objects.  Optionally accepts a "
                    "relative sub-directory path."
                ),
            },
            {
                "name": "workspace_read_file",
                "func": self.workspace_read_file,
                "description": (
                    "Read the text content of a workspace file.  Accepts a "
                    "relative path within the workspace.  Returns up to 8000 "
                    "characters of the file content.  "
                    "For structured files (docx/xlsx/pptx/pdf) the extracted "
                    "text representation is returned."
                ),
            },
            {
                "name": "workspace_create_file",
                "func": self.workspace_create_file,
                "description": (
                    "Create a new file in the workspace with the given content.  "
                    "Returns the relative path of the created file.  "
                    "The frontend is notified to refresh the file tree and open "
                    "the new file in an editor tab."
                ),
            },
            {
                "name": "workspace_save_file",
                "func": self.workspace_save_file,
                "description": (
                    "Overwrite an existing workspace file with new content.  "
                    "Use for saving processed results.  A backup is created "
                    "automatically before overwriting."
                ),
            },
            {
                "name": "editor_apply",
                "func": self.editor_apply,
                "description": (
                    "Push a change directly into the currently-open editor tab "
                    "in the browser.  Supported operations:\n"
                    '  - type="set_html", content="<p>..." → replace document body (Word/TipTap)\n'
                    '  - type="set_cell", cell="A1", value="..." → set a single spreadsheet cell\n'
                    '  - type="set_cells", cells=[{cell,value},...] → batch set cells\n'
                    '  - type="insert_text", position="end", content="..." → append text\n'
                    "Returns a confirmation string."
                ),
            },
            {
                "name": "editor_open_file",
                "func": self.editor_open_file,
                "description": (
                    "Ask the frontend to open a workspace file in a new editor "
                    "tab.  Accepts a relative path.  The file must already exist."
                ),
            },
        ]

    # ── Tool implementations ─────────────────────────────────────────────

    def workspace_list_files(self, path: str = "", recursive: bool = False) -> str:
        """List files under workspace/[path]."""
        root = _get_workspace_root()
        target = _safe_resolve(path) if path else root
        if target is None:
            return "Error: path escapes workspace boundary"
        if not os.path.isdir(target):
            return f"Error: '{path}' is not a directory"

        entries = []
        try:
            if recursive:
                for dirpath, dirnames, filenames in os.walk(target):
                    # Skip hidden dirs
                    dirnames[:] = [d for d in dirnames if not d.startswith(".")]
                    for fname in filenames:
                        if fname.startswith("."):
                            continue
                        fpath = os.path.join(dirpath, fname)
                        rel = os.path.relpath(fpath, root)
                        try:
                            st = os.stat(fpath)
                            entries.append(
                                {
                                    "name": rel.replace("\\", "/"),
                                    "type": "file",
                                    "size": st.st_size,
                                    "mtime": int(st.st_mtime),
                                }
                            )
                        except OSError:
                            pass
                    if len(entries) >= 500:
                        break
            else:
                for item in sorted(os.listdir(target)):
                    if item.startswith("."):
                        continue
                    fpath = os.path.join(target, item)
                    rel = os.path.relpath(fpath, root).replace("\\", "/")
                    try:
                        st = os.stat(fpath)
                        entries.append(
                            {
                                "name": rel,
                                "type": "dir" if os.path.isdir(fpath) else "file",
                                "size": st.st_size if os.path.isfile(fpath) else 0,
                                "mtime": int(st.st_mtime),
                            }
                        )
                    except OSError:
                        pass
        except PermissionError:
            return f"Error: permission denied for '{path}'"

        return json.dumps(entries, ensure_ascii=False)

    def workspace_read_file(self, path: str, max_chars: int = 8000) -> str:
        """Read text content of a workspace file."""
        resolved = _safe_resolve(path)
        if resolved is None:
            return "Error: path escapes workspace boundary"
        if not os.path.isfile(resolved):
            return f"Error: file not found — '{path}'"

        ext = os.path.splitext(resolved)[1].lower()

        # For structured file types, try to extract text via existing parsers
        if ext in (".docx", ".xlsx", ".pptx", ".pdf"):
            return self._read_structured(resolved, ext, max_chars)

        # Plain text / code / csv / json / markdown
        try:
            with open(resolved, "r", encoding="utf-8", errors="replace") as f:
                content = f.read(max_chars)
            truncated = os.path.getsize(resolved) > max_chars
            header = f"[{path}] ({os.path.getsize(resolved)} bytes)"
            if truncated:
                header += " (truncated)"
            return f"{header}\n{content}"
        except Exception as e:
            return f"Error reading file: {e}"

    def workspace_create_file(self, path: str, content: str = "") -> str:
        """Create a new file in the workspace."""
        resolved = _safe_resolve(path)
        if resolved is None:
            return "Error: path escapes workspace boundary"
        if os.path.exists(resolved):
            return f"Error: file already exists — '{path}'. Use workspace_save_file to overwrite."

        try:
            os.makedirs(os.path.dirname(resolved), exist_ok=True)
            with open(resolved, "w", encoding="utf-8") as f:
                f.write(content)
        except Exception as e:
            return f"Error creating file: {e}"

        # Notify frontend to refresh file tree and open the file
        self._emit("workspace_file_created", {"path": path.replace("\\", "/")})
        return f"Created: {path}"

    def workspace_save_file(self, path: str, content: str) -> str:
        """Overwrite a workspace file (auto-backup)."""
        resolved = _safe_resolve(path)
        if resolved is None:
            return "Error: path escapes workspace boundary"
        if not os.path.isfile(resolved):
            return f"Error: file not found — '{path}'. Use workspace_create_file for new files."

        # Create backup
        try:
            import shutil

            backup_path = resolved + ".bak"
            shutil.copy2(resolved, backup_path)
        except Exception:
            pass  # backup is best-effort

        try:
            with open(resolved, "w", encoding="utf-8") as f:
                f.write(content)
        except Exception as e:
            return f"Error saving file: {e}"

        self._emit("workspace_file_updated", {"path": path.replace("\\", "/")})
        return f"Saved: {path}"

    def editor_apply(self, type: str, **kwargs) -> str:
        """Push a change to the live editor via WebSocket."""
        if not self._socketio:
            return "Error: no WebSocket connection (editor_apply unavailable)"

        payload = {"type": type, **kwargs}

        # Validate known operation types
        valid_types = {
            "set_html",
            "set_cell",
            "set_cells",
            "insert_text",
            "set_pptx_text",
        }
        if type not in valid_types:
            return f"Error: unknown editor_apply type '{type}'. Valid: {', '.join(sorted(valid_types))}"

        try:
            self._socketio.emit(
                "agent_execute_command",
                {"action": "editor_apply", "payload": payload},
                namespace="/doc",
            )
            return f"Applied: {type}"
        except Exception as e:
            return f"Error emitting editor_apply: {e}"

    def editor_open_file(self, path: str) -> str:
        """Ask the frontend to open a workspace file."""
        resolved = _safe_resolve(path)
        if resolved is None:
            return "Error: path escapes workspace boundary"
        if not os.path.exists(resolved):
            return f"Error: file not found — '{path}'"

        self._emit("workspace_open_file", {"path": path.replace("\\", "/")})
        return f"Opened: {path}"

    # ── Internal helpers ─────────────────────────────────────────────────

    def _emit(self, event: str, data: dict):
        """Emit a WebSocket event on /doc namespace (best-effort)."""
        if self._socketio:
            try:
                self._socketio.emit(event, data, namespace="/doc")
            except Exception as e:
                logger.warning("[WorkspaceEditorPlugin] emit %s failed: %s", event, e)

    @staticmethod
    def _read_structured(fpath: str, ext: str, max_chars: int) -> str:
        """Extract text from structured file formats."""
        text = ""
        try:
            if ext == ".docx":
                try:
                    from docx import Document as DocxDocument

                    doc = DocxDocument(fpath)
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    text = "\n".join(paragraphs)
                except ImportError:
                    text = "(python-docx not installed — cannot read .docx)"

            elif ext == ".xlsx":
                try:
                    import openpyxl

                    wb = openpyxl.load_workbook(fpath, read_only=True, data_only=True)
                    parts = []
                    for sheet_name in wb.sheetnames[:5]:
                        ws = wb[sheet_name]
                        rows = []
                        for row in ws.iter_rows(max_row=200, values_only=True):
                            cells = [str(c) if c is not None else "" for c in row]
                            rows.append("\t".join(cells))
                        parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
                    wb.close()
                    text = "\n\n".join(parts)
                except ImportError:
                    text = "(openpyxl not installed — cannot read .xlsx)"

            elif ext == ".pptx":
                try:
                    from pptx import Presentation

                    prs = Presentation(fpath)
                    slides = []
                    for i, slide in enumerate(prs.slides[:30], 1):
                        slide_texts = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_texts.append(shape.text.strip())
                        if slide_texts:
                            slides.append(f"[Slide {i}]\n" + "\n".join(slide_texts))
                    text = "\n\n".join(slides)
                except ImportError:
                    text = "(python-pptx not installed — cannot read .pptx)"

            elif ext == ".pdf":
                try:
                    import fitz  # PyMuPDF

                    doc = fitz.open(fpath)
                    pages = []
                    for i, page in enumerate(doc[:20], 1):
                        page_text = page.get_text().strip()
                        if page_text:
                            pages.append(f"[Page {i}]\n{page_text}")
                    doc.close()
                    text = "\n\n".join(pages)
                except ImportError:
                    text = "(PyMuPDF not installed — cannot read .pdf)"

        except Exception as e:
            text = f"Error extracting text: {e}"

        name = os.path.basename(fpath)
        if len(text) > max_chars:
            text = text[:max_chars] + "\n...(truncated)"
        return (
            f"[{name}] ({ext})\n{text}" if text else f"[{name}] (empty or unreadable)"
        )
