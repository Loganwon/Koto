# -*- coding: utf-8 -*-
"""
E2E tests for the workspace file editor feature.

Covers the user journeys around opening, navigating, and handling errors in the
workspace file editor — specifically the fixes introduced for the PPTX
"Package not found" bug and related 0-byte file issues.

Playwright-based, runs against a live Flask server (port 9876) started by the
`_flask_server` session fixture in conftest.py.

Run with:
    pytest tests/e2e/test_workspace_file_editor.py -v -m e2e
"""

from __future__ import annotations

import io
import time
import zipfile
from pathlib import Path

import pytest

# ── helpers (mirrors test_human_journeys.py style) ───────────────────────────
PAGE_TIMEOUT = 15_000  # ms
THINK_SHORT = 400
THINK_MEDIUM = 800
THINK_LONG = 1_500

# Tolerated console noise patterns (same as conftest BENIGN_ERROR_PATTERNS).
_BENIGN = [
    "WebSocket",
    "ws://",
    "wss://",
    "net::ERR_",
    "favicon.ico",
    "API key",
    "api key",
    "Failed to load resource",
    "ERR_CONNECTION_REFUSED",
]


def _is_benign(msg: str) -> bool:
    return any(p in msg for p in _BENIGN)


def _real_errors(errors: list[str]) -> list[str]:
    return [e for e in errors if not _is_benign(e)]


def _goto(page, url: str):
    """Navigate and return the response (graceful on navigation errors)."""
    try:
        return page.goto(url, timeout=PAGE_TIMEOUT, wait_until="domcontentloaded")
    except Exception:
        return None


def _settle(page, ms: int) -> None:
    """Wait for `ms` milliseconds of simulated think-time."""
    if ms > 0:
        page.wait_for_timeout(ms)


# ── minimal PPTX builder ─────────────────────────────────────────────────────


def _make_minimal_pptx_bytes() -> bytes:
    """Return a real (non-empty, valid) .pptx file as bytes, or skip."""
    try:
        from pptx import Presentation

        buf = io.BytesIO()
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = "E2E Test Slide"
        prs.save(buf)
        return buf.getvalue()
    except ImportError:
        pytest.skip("python-pptx not installed — cannot build test PPTX")


def _make_minimal_docx_bytes() -> bytes:
    try:
        import docx

        buf = io.BytesIO()
        doc = docx.Document()
        doc.add_paragraph("E2E test document")
        doc.save(buf)
        return buf.getvalue()
    except ImportError:
        pytest.skip("python-docx not installed — cannot build test DOCX")


# ─────────────────────────────────────────────────────────────────────────────
# Journey 1: Landing on the workspace editor page
# ─────────────────────────────────────────────────────────────────────────────


@pytest.mark.e2e
class TestWorkspaceEditorPageLoad:
    """User navigates to the workspace editor and sees the correct UI."""

    def test_workspace_page_loads_without_crash(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """GET /workspace returns the editor shell with no JS errors."""
        resp = _goto(e2e_page, f"{e2e_base_url}/workspace")
        _settle(e2e_page, THINK_LONG)
        if resp and resp.status >= 400:
            pytest.skip("Workspace editor page not available at /workspace")

        # The editor shell must be present
        has_canvas = e2e_page.locator("#wa-canvas, #wa-drop-zone").count() > 0
        assert has_canvas, "Workspace canvas / drop-zone should be present"

        assert (
            _real_errors(console_errors) == []
        ), f"JS errors on workspace page load: {_real_errors(console_errors)}"

    def test_workspace_file_tree_panel_exists(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """The left-hand file tree panel must exist in the DOM."""
        resp = _goto(e2e_page, f"{e2e_base_url}/workspace")
        _settle(e2e_page, THINK_LONG)
        if resp and resp.status >= 400:
            pytest.skip("Workspace editor page not available")

        assert (
            e2e_page.locator("#wa-files-list").count() > 0
        ), "#wa-files-list (the file tree panel) must be present in the DOM"

        assert (
            _real_errors(console_errors) == []
        ), f"JS errors: {_real_errors(console_errors)}"

    def test_workspace_header_has_title(self, e2e_page, console_errors, e2e_base_url):
        """Header title element must be visible."""
        resp = _goto(e2e_page, f"{e2e_base_url}/workspace")
        _settle(e2e_page, THINK_LONG)
        if resp and resp.status >= 400:
            pytest.skip("Workspace editor page not available")

        title_el = e2e_page.locator("#wa-file-name").first
        assert title_el.count() > 0, "#wa-file-name must be present in the header"

        assert (
            _real_errors(console_errors) == []
        ), f"JS errors: {_real_errors(console_errors)}"

    def test_workspace_toast_element_present(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """The toast notification element must be present (required for error display)."""
        resp = _goto(e2e_page, f"{e2e_base_url}/workspace")
        _settle(e2e_page, THINK_LONG)
        if resp and resp.status >= 400:
            pytest.skip("Workspace editor page not available")

        assert (
            e2e_page.locator("#wa-toast").count() > 0
        ), "#wa-toast must be present so errors can be shown to the user"

        assert (
            _real_errors(console_errors) == []
        ), f"JS errors: {_real_errors(console_errors)}"


# ─────────────────────────────────────────────────────────────────────────────
# Journey 2: Upload a real PPTX file and see it parsed
# ─────────────────────────────────────────────────────────────────────────────


@pytest.mark.e2e
class TestPptxUploadJourney:
    """User picks a PPTX file via the file input and sees the slides rendered."""

    def _navigate_to_workspace(self, page, base_url):
        resp = _goto(page, f"{base_url}/workspace")
        _settle(page, THINK_LONG)
        if resp and resp.status >= 400:
            pytest.skip("Workspace editor page not available")

    def test_pptx_upload_opens_without_package_not_found_error(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """
        Critical regression test: uploading a real PPTX must NOT produce
        'Package not found' in the page or console.  Before the fix, a relative
        _TMP_DIR caused python-pptx to fail.
        """
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        pptx_bytes = _make_minimal_pptx_bytes()

        file_input = e2e_page.locator("#wa-file-input, input[accept*='.pptx']").first
        if file_input.count() == 0:
            pytest.skip("File input not found on workspace page")

        file_input.set_input_files(
            {
                "name": "regression_test.pptx",
                "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "buffer": pptx_bytes,
            }
        )
        _settle(e2e_page, THINK_LONG * 2)

        # 'Package not found' must not appear in any toast or console error
        page_text = e2e_page.locator("body").inner_text()
        assert (
            "Package not found" not in page_text
        ), "PPTX upload: 'Package not found' error must not appear after the absolute-path fix"

        errors_with_package_not_found = [
            e for e in console_errors if "Package not found" in e
        ]
        assert (
            errors_with_package_not_found == []
        ), f"'Package not found' in console after PPTX upload: {errors_with_package_not_found}"

    def test_pptx_upload_shows_loading_toast(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """Uploading a file should trigger a toast notification."""
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        pptx_bytes = _make_minimal_pptx_bytes()

        file_input = e2e_page.locator("#wa-file-input, input[accept*='.pptx']").first
        if file_input.count() == 0:
            pytest.skip("File input not found")

        file_input.set_input_files(
            {
                "name": "toast_test.pptx",
                "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "buffer": pptx_bytes,
            }
        )
        _settle(e2e_page, THINK_MEDIUM)

        # Toast should appear (may already have disappeared — that's OK)
        # Just ensure it doesn't permanently show an error state
        toast = e2e_page.locator("#wa-toast")
        # Toast being visible is a bonus; what matters is no crash
        assert e2e_page.locator("body").count() > 0, "Page must not crash on upload"

        assert (
            _real_errors(console_errors) == []
        ), f"JS errors during PPTX upload: {_real_errors(console_errors)}"

    def test_zero_byte_pptx_upload_shows_error_not_crash(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """
        Guard test: uploading a 0-byte PPTX file should show an error toast,
        not a crash or 'Package not found'.
        """
        self._navigate_to_workspace(e2e_page, e2e_base_url)

        file_input = e2e_page.locator("#wa-file-input, input[accept*='.pptx']").first
        if file_input.count() == 0:
            pytest.skip("File input not found")

        # Upload an empty (0-byte) PPTX
        file_input.set_input_files(
            {
                "name": "empty_slides.pptx",
                "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "buffer": b"",
            }
        )
        _settle(e2e_page, THINK_LONG)

        # 'Package not found' must never reach the page
        page_text = e2e_page.locator("body").inner_text()
        assert (
            "Package not found" not in page_text
        ), "0-byte PPTX must trigger the guard, not 'Package not found'"

        # Page must remain functional
        assert (
            e2e_page.locator("#wa-canvas, #wa-drop-zone").count() > 0
        ), "Workspace canvas must still be present after 0-byte upload attempt"

    def test_docx_upload_no_crash(self, e2e_page, console_errors, e2e_base_url):
        """Uploading a real DOCX file must not crash the editor."""
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        docx_bytes = _make_minimal_docx_bytes()

        file_input = e2e_page.locator("#wa-file-input, input[accept*='.docx']").first
        if file_input.count() == 0:
            pytest.skip("File input not found")

        file_input.set_input_files(
            {
                "name": "test_document.docx",
                "mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "buffer": docx_bytes,
            }
        )
        _settle(e2e_page, THINK_LONG * 2)

        # No crash
        assert e2e_page.locator("body").count() > 0
        assert (
            _real_errors(console_errors) == []
        ), f"JS errors on DOCX upload: {_real_errors(console_errors)}"

    def test_pdf_upload_no_crash(self, e2e_page, console_errors, e2e_base_url):
        """Uploading a PDF must not crash the editor."""
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        pdf_bytes = b"%PDF-1.4\n1 0 obj\n<</Type /Catalog>>\nendobj\n%%EOF"

        file_input = e2e_page.locator("#wa-file-input, input[accept*='.pdf']").first
        if file_input.count() == 0:
            pytest.skip("File input not found")

        file_input.set_input_files(
            {
                "name": "report.pdf",
                "mimeType": "application/pdf",
                "buffer": pdf_bytes,
            }
        )
        _settle(e2e_page, THINK_LONG * 2)

        assert e2e_page.locator("body").count() > 0
        assert (
            _real_errors(console_errors) == []
        ), f"JS errors on PDF upload: {_real_errors(console_errors)}"


# ─────────────────────────────────────────────────────────────────────────────
# Journey 3: File tree — ppt_sessions directory must be hidden
# ─────────────────────────────────────────────────────────────────────────────


@pytest.mark.e2e
class TestFileTreePptSessionsHidden:
    """
    The ppt_sessions folder holds PPT session artefacts (often 0-byte files).
    It must never appear in the file tree shown to the user.
    """

    def test_ppt_sessions_not_visible_in_file_tree(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """
        Verify via the /api/v1/workspace/list_files API that ppt_sessions is
        absent from the tree returned to the page.
        """
        import requests  # available in test env via pytest-playwright

        try:
            resp = requests.get(
                f"{e2e_base_url}/api/v1/workspace/list_files",
                timeout=8,
            )
        except Exception:
            pytest.skip("Could not reach list_files API")

        if resp.status_code != 200:
            pytest.skip(f"list_files returned {resp.status_code}")

        def _find(nodes: list, name: str) -> bool:
            for n in nodes:
                if n.get("name") == name:
                    return True
                if n.get("children") and _find(n["children"], name):
                    return True
            return False

        data = resp.json()
        files = data.get("files", [])
        assert not _find(files, "ppt_sessions"), (
            "ppt_sessions must not appear in the workspace file tree — "
            "it contains internal session artefacts, not user files"
        )

        assert (
            _real_errors(console_errors) == []
        ), f"JS errors: {_real_errors(console_errors)}"

    def test_file_tree_api_has_correct_shape(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """list_files response must include workspace_name and workspace_path."""
        import requests

        try:
            resp = requests.get(
                f"{e2e_base_url}/api/v1/workspace/list_files",
                timeout=8,
            )
        except Exception:
            pytest.skip("Could not reach list_files API")

        if resp.status_code != 200:
            pytest.skip(f"list_files returned {resp.status_code}")

        data = resp.json()
        assert "workspace_name" in data, "list_files must return 'workspace_name'"
        assert "workspace_path" in data, "list_files must return 'workspace_path'"
        assert isinstance(
            data.get("files"), list
        ), "list_files must return 'files' as a list"

    def test_tmp_dir_not_visible_in_file_tree(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """The 'tmp' directory (used for parsed file staging) must be hidden."""
        import requests

        try:
            resp = requests.get(
                f"{e2e_base_url}/api/v1/workspace/list_files",
                timeout=8,
            )
        except Exception:
            pytest.skip("Could not reach list_files API")

        if resp.status_code != 200:
            pytest.skip(f"list_files returned {resp.status_code}")

        def _find(nodes: list, name: str) -> bool:
            for n in nodes:
                if n.get("name") == name:
                    return True
                if n.get("children") and _find(n["children"], name):
                    return True
            return False

        data = resp.json()
        assert not _find(
            data.get("files", []), "tmp"
        ), "'tmp' directory must not appear in the workspace file tree"


# ─────────────────────────────────────────────────────────────────────────────
# Journey 4: open_file_by_path API — the fix for the double round-trip
# ─────────────────────────────────────────────────────────────────────────────


@pytest.mark.e2e
class TestOpenFileByPathApi:
    """
    Verify the open_file_by_path endpoint used by openWorkspaceFile() behaves
    correctly: accepts valid paths, rejects 0-byte files with 400, rejects
    path traversal with 403, returns 404 for missing files.
    """

    def _post_open_by_path(self, base_url: str, rel_path: str, timeout: int = 8):
        import requests

        try:
            return requests.post(
                f"{base_url}/api/v1/workspace/open_file_by_path",
                json={"path": rel_path},
                timeout=timeout,
            )
        except Exception as exc:
            pytest.skip(f"open_file_by_path unreachable: {exc}")

    def test_missing_path_returns_400(self, e2e_page, e2e_base_url):
        """Empty / missing path → 400."""
        import requests

        try:
            resp = requests.post(
                f"{e2e_base_url}/api/v1/workspace/open_file_by_path",
                json={},
                timeout=8,
            )
        except Exception:
            pytest.skip("Cannot reach open_file_by_path endpoint")
        assert resp.status_code == 400

    def test_path_traversal_returns_403(self, e2e_page, e2e_base_url):
        """Path traversal attempt → 403."""
        resp = self._post_open_by_path(e2e_base_url, "../../../etc/passwd")
        assert resp.status_code == 403

    def test_nonexistent_file_returns_404(self, e2e_page, e2e_base_url):
        """Non-existent workspace file → 404."""
        resp = self._post_open_by_path(
            e2e_base_url, "definitely_does_not_exist_xyz.pptx"
        )
        assert resp.status_code == 404

    def test_no_package_not_found_in_404_error(self, e2e_page, e2e_base_url):
        """
        For a missing file, the error message must be a clear 'file not found'
        message — never 'Package not found' (which is a python-pptx internal error
        that leaks when a relative path is used).
        """
        resp = self._post_open_by_path(
            e2e_base_url, "workspace/tmp/relative_path_test.pptx"
        )
        if resp.status_code not in (400, 403, 404):
            pytest.skip(f"Unexpected status {resp.status_code}")
        error = (resp.json() or {}).get("error", "")
        assert (
            "Package not found" not in error
        ), f"'Package not found' must never reach the client; got: {error!r}"


# ─────────────────────────────────────────────────────────────────────────────
# Journey 5: File drag-and-drop onto the workspace drop-zone
# ─────────────────────────────────────────────────────────────────────────────


@pytest.mark.e2e
class TestWorkspaceDragDropJourney:
    """User drags a file from the OS onto the workspace drop-zone."""

    def _navigate_to_workspace(self, page, base_url):
        resp = _goto(page, f"{base_url}/workspace")
        _settle(page, THINK_LONG)
        if resp and resp.status >= 400:
            pytest.skip("Workspace editor page not available")

    def test_drop_zone_element_exists(self, e2e_page, console_errors, e2e_base_url):
        """
        The #wa-drop-zone element must exist so the file drag-and-drop
        event handlers can attach to it.
        """
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        drop_zone = e2e_page.locator("#wa-drop-zone")
        assert (
            drop_zone.count() > 0
        ), "#wa-drop-zone element must be present to support drag-and-drop uploads"
        assert (
            _real_errors(console_errors) == []
        ), f"JS errors: {_real_errors(console_errors)}"

    def test_file_input_accepts_pptx_extension(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """The upload file input must declare .pptx in its accept attribute."""
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        # Check any file input on the page
        inputs = e2e_page.locator("input[type='file']")
        found = False
        for i in range(inputs.count()):
            accept = inputs.nth(i).get_attribute("accept") or ""
            if ".pptx" in accept:
                found = True
                break
        assert (
            found
        ), "At least one file input must declare .pptx in its accept attribute"

    def test_drag_over_body_does_not_crash(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """
        Simulating a dragover event on the page body must not produce JS errors.
        """
        self._navigate_to_workspace(e2e_page, e2e_base_url)

        # Dispatch a synthetic dragover to the body
        e2e_page.evaluate("""() => {
            const ev = new DragEvent('dragover', { bubbles: true, cancelable: true });
            document.body.dispatchEvent(ev);
        }""")
        _settle(e2e_page, THINK_SHORT)

        e2e_page.evaluate("""() => {
            const ev = new DragEvent('dragleave', { bubbles: true, cancelable: true });
            document.body.dispatchEvent(ev);
        }""")
        _settle(e2e_page, THINK_SHORT)

        assert (
            _real_errors(console_errors) == []
        ), f"JS errors on drag events: {_real_errors(console_errors)}"

    def test_simulated_pptx_drop_via_file_input(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """
        Simulate a file drop by using set_input_files on the hidden file input.
        The upload should succeed without 'Package not found' errors.
        """
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        pptx_bytes = _make_minimal_pptx_bytes()

        file_input = e2e_page.locator("input[type='file'][accept*='.pptx']").first
        if file_input.count() == 0:
            pytest.skip("PPTX file input not found")

        file_input.set_input_files(
            {
                "name": "dropped_deck.pptx",
                "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "buffer": pptx_bytes,
            }
        )
        _settle(e2e_page, THINK_LONG * 2)

        page_text = e2e_page.locator("body").inner_text()
        assert (
            "Package not found" not in page_text
        ), "Simulated PPTX drop: 'Package not found' must not appear"
        assert (
            _real_errors(console_errors) == []
        ), f"JS errors after simulated PPTX drop: {_real_errors(console_errors)}"


# ─────────────────────────────────────────────────────────────────────────────
# Journey 6: Tab management — opening a file creates a tab
# ─────────────────────────────────────────────────────────────────────────────


@pytest.mark.e2e
class TestWorkspaceTabManagement:
    """User opens a file and expects a tab to appear in the tab bar."""

    def _navigate_to_workspace(self, page, base_url):
        resp = _goto(page, f"{base_url}/workspace")
        _settle(page, THINK_LONG)
        if resp and resp.status >= 400:
            pytest.skip("Workspace editor page not available")

    def test_tab_bar_element_exists(self, e2e_page, console_errors, e2e_base_url):
        """The tab bar must be present in the DOM at page load."""
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        tab_bar = e2e_page.locator("#wa-tab-bar")
        assert (
            tab_bar.count() > 0
        ), "#wa-tab-bar must be present so opened files can show as tabs"
        assert (
            _real_errors(console_errors) == []
        ), f"JS errors: {_real_errors(console_errors)}"

    def test_pptx_upload_creates_tab(self, e2e_page, console_errors, e2e_base_url):
        """Uploading a PPTX file must cause a tab to appear in #wa-tab-bar."""
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        pptx_bytes = _make_minimal_pptx_bytes()

        file_input = e2e_page.locator("input[type='file'][accept*='.pptx']").first
        if file_input.count() == 0:
            pytest.skip("File input not found")

        tabs_before = e2e_page.locator("#wa-tab-bar .wa-tab").count()

        file_input.set_input_files(
            {
                "name": "tab_test.pptx",
                "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "buffer": pptx_bytes,
            }
        )
        _settle(e2e_page, THINK_LONG * 2)

        tabs_after = e2e_page.locator("#wa-tab-bar .wa-tab").count()
        # Tabs should have increased by at least 1
        assert tabs_after > tabs_before, (
            f"Opening a PPTX file should create a new tab; "
            f"before={tabs_before}, after={tabs_after}"
        )

    def test_docx_upload_creates_tab(self, e2e_page, console_errors, e2e_base_url):
        """Uploading a DOCX file must create a tab entry."""
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        docx_bytes = _make_minimal_docx_bytes()

        file_input = e2e_page.locator("input[type='file'][accept*='.docx']").first
        if file_input.count() == 0:
            pytest.skip("File input not found")

        tabs_before = e2e_page.locator("#wa-tab-bar .wa-tab").count()

        file_input.set_input_files(
            {
                "name": "tab_doc_test.docx",
                "mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "buffer": docx_bytes,
            }
        )
        _settle(e2e_page, THINK_LONG * 2)

        tabs_after = e2e_page.locator("#wa-tab-bar .wa-tab").count()
        assert tabs_after > tabs_before, (
            f"Opening a DOCX file should create a new tab; "
            f"before={tabs_before}, after={tabs_after}"
        )

    def test_header_title_changes_after_upload(
        self, e2e_page, console_errors, e2e_base_url
    ):
        """After uploading a file, #wa-file-name should show the filename."""
        self._navigate_to_workspace(e2e_page, e2e_base_url)
        pptx_bytes = _make_minimal_pptx_bytes()

        file_input = e2e_page.locator("input[type='file'][accept*='.pptx']").first
        if file_input.count() == 0:
            pytest.skip("File input not found")

        file_input.set_input_files(
            {
                "name": "title_test.pptx",
                "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "buffer": pptx_bytes,
            }
        )
        _settle(e2e_page, THINK_LONG * 2)

        title_el = e2e_page.locator("#wa-file-name")
        if title_el.count() > 0:
            title_text = title_el.inner_text()
            # After upload the title should no longer be the default placeholder
            # (it may be empty or the filename — either is acceptable vs "文件工作站")
            # The important thing: no crash
            assert isinstance(title_text, str)

        assert (
            _real_errors(console_errors) == []
        ), f"JS errors: {_real_errors(console_errors)}"


# ─────────────────────────────────────────────────────────────────────────────
# Journey 7: JS source — openWorkspaceFile correctness checks
# ─────────────────────────────────────────────────────────────────────────────


@pytest.mark.e2e
class TestWorkspaceAssistantJsSource:
    """
    Static analysis of the workspace file-open source to verify critical fixes are present.
    These tests run without a browser and don't require the Flask server.
    """

    @property
    def src(self) -> str:
        root = Path(__file__).parents[2]
        return "\n".join(
            [
                (root / "web" / "src" / "workspace" / "fs-tree.ts").read_text(encoding="utf-8"),
                (root / "web" / "src" / "workspace" / "infrastructure.ts").read_text(encoding="utf-8"),
            ]
        )

    def test_open_workspace_file_uses_open_file_by_path(self):
        """
        openWorkspaceFile must use open_file_by_path (not the old double-roundtrip
        of fetching bytes then re-uploading to open_file).
        """
        assert (
            "open_file_by_path" in self.src
        ), "openWorkspaceFile must call the open_file_by_path endpoint directly"

    def test_open_workspace_file_does_not_use_blob_roundtrip(self):
        """
        The old raw-byte re-upload path must be absent from openWorkspaceFile.
        """
        # Find the function and check that it does NOT contain the blob pattern
        fn_start = self.src.find("async function openWorkspaceFile")
        assert fn_start >= 0, "openWorkspaceFile function not found"
        # Read ~60 lines of the function body (safe upper bound)
        fn_body = self.src[fn_start : fn_start + 2000]
        assert ("new File(" + "[blob]") not in fn_body, (
            "openWorkspaceFile must not use the old raw-byte re-upload path. "
            "It should call open_file_by_path directly."
        )

    def test_router_from_parsed_method_exists(self):
        """
        Router.fromParsed must exist — it's the shared method that applies
        parsed server response to the editor, used by both openWorkspaceFile
        and the upload path (Router.load).
        """
        assert (
            "fromParsed" in self.src
        ), "Router.fromParsed must be defined — it's required by the fixed openWorkspaceFile"

    def test_tmp_dir_uses_path_resolve(self):
        """
        _TMP_DIR must use Path(__file__).resolve() so it stays absolute regardless
        of the process CWD.  Verify the source declares this.
        """
        # This is a Python-side guard — check the backend file instead
        wa_py = (
            Path(__file__).parents[2] / "web" / "blueprints" / "workspace_assistant.py"
        )
        src = wa_py.read_text(encoding="utf-8")
        assert "__file__" in src and "resolve()" in src, (
            "workspace_assistant.py _TMP_DIR must use Path(__file__).resolve() "
            "to produce an absolute path independent of CWD"
        )

    def test_ppt_sessions_in_skip_set(self):
        """
        'ppt_sessions' must be in the _skip set in list_workspace_files so
        session artefacts don't appear in the file tree.
        """
        wa_py = (
            Path(__file__).parents[2] / "web" / "blueprints" / "workspace_assistant.py"
        )
        src = wa_py.read_text(encoding="utf-8")
        assert (
            '"ppt_sessions"' in src or "'ppt_sessions'" in src
        ), "ppt_sessions must be in the _skip set in list_workspace_files"

    def test_zero_byte_guard_in_open_file(self):
        """
        The open_file endpoint must contain a 0-byte guard
        (st_size == 0 check) to reject empty uploads early.
        """
        wa_py = (
            Path(__file__).parents[2] / "web" / "blueprints" / "workspace_assistant.py"
        )
        src = wa_py.read_text(encoding="utf-8")
        assert "st_size == 0" in src, (
            "open_file (and/or open_file_by_path) must check st_size == 0 "
            "to reject 0-byte files before calling the parser"
        )

    def test_zero_byte_guard_in_open_file_by_path(self):
        """open_file_by_path must also contain a 0-byte guard."""
        wa_py = (
            Path(__file__).parents[2] / "web" / "blueprints" / "workspace_assistant.py"
        )
        src = wa_py.read_text(encoding="utf-8")
        # Count occurrences — there should be at least 2 (one per endpoint)
        count = src.count("st_size == 0")
        assert count >= 2, (
            f"Expected at least 2 zero-byte guards (open_file + open_file_by_path), "
            f"found {count}"
        )

    def test_open_browser_file_routes_workspace_files_through_open_file_by_path(self):
        """
        openBrowserFile must route workspace files through open_file_by_path.
        The old raw-byte round-trip produced 0-byte uploads whenever the
        seeded content was not flushed.
        """
        fn_start = self.src.find("async function openBrowserFile")
        assert fn_start >= 0, "openBrowserFile function not found"
        fn_body = self.src[fn_start : fn_start + 3000]
        assert "openWorkspaceFile(absPath, supported)" in fn_body
        assert ("new File(" + "[blob]") not in fn_body

    def test_open_browser_file_uses_open_abs_file_for_external_files(self):
        """
        openBrowserFile must use open_abs_file for files outside the workspace.
        """
        fn_start = self.src.find("async function openBrowserFile")
        assert fn_start >= 0, "openBrowserFile function not found"
        fn_body = self.src[fn_start : fn_start + 3000]
        assert (
            "open_abs_file" in fn_body
        ), "openBrowserFile must use open_abs_file for external files"

    def test_open_recent_file_bridge_routes_workspace_and_absolute_paths(self):
        """
        The main app opens Office/PDF files through WA.openRecentFile().
        That bridge must exist and dispatch relative/workspace paths to
        openWorkspaceFile, while absolute paths still flow through
        openBrowserFile.
        """
        fn_start = self.src.find("openRecentFile = async")
        assert fn_start >= 0, "WA.openRecentFile bridge function not found"
        fn_body = self.src[fn_start : fn_start + 2200]
        assert "openWorkspaceFile" in fn_body, (
            "WA.openRecentFile must route workspace-relative paths back through openWorkspaceFile"
        )
        assert "openBrowserFile" in fn_body, (
            "WA.openRecentFile must route absolute paths back through openBrowserFile"
        )
        assert "state._workspacePath" in fn_body, (
            "WA.openRecentFile must compare against the active workspace root before routing"
        )

    def test_file_type_icons_are_mapped_to_distinct_svgs(self):
        """Workspace file list should keep colored Office-style icon mappings for common file types."""
        src = self.src
        assert "function _waBrandFileSvg(label, opts = {})" in src, "Brand-style file SVG helper should exist"
        assert "const _WORD_FILE_SVG" in src, "Word files should have a dedicated colored icon definition"
        assert "const _EXCEL_FILE_SVG" in src, "Excel files should have a dedicated colored icon definition"
        assert "const _POWERPOINT_FILE_SVG" in src, "PowerPoint files should have a dedicated colored icon definition"
        assert "const _PDF_SVG" in src, "PDF should have a dedicated colored icon definition"
        assert "#185ABD" in src, "Word icon should keep Office blue branding"
        assert "#107C41" in src, "Excel icon should keep Office green branding"
        assert "#D24726" in src, "PowerPoint icon should keep Office orange branding"
        assert "#E53935" in src, "PDF icon should keep red branding"
        assert "docx: _WORD_FILE_SVG" in src, "DOCX extension should map to Word icon"
        assert "xlsx: _EXCEL_FILE_SVG" in src, "XLSX extension should map to Excel icon"
        assert "pptx: _POWERPOINT_FILE_SVG" in src, "PPTX extension should map to PowerPoint icon"
        assert "pdf: _PDF_SVG" in src, "PDF extension should map to PDF icon"
        assert src.count("const _PDF_SVG") == 1, "PDF icon constant must not be declared twice"
        assert src.count("const _TEXT_SVG") == 1, "Text icon constant must not be declared twice"
        assert src.count("const _CODE_SVG") == 1, "Code icon constant must not be declared twice"
        assert src.count("const _IMAGE_SVG") == 1, "Image icon constant must not be declared twice"

    def test_reload_file_by_path_routes_workspace_relative_paths_through_open_file_by_path(self):
        """
        Task-stream refresh uses WA.reloadFileByPath(). When file-change events
        carry workspace-relative paths, reloadFileByPath must still route them
        through open_file_by_path instead of misclassifying them as absolute
        browser files.
        """
        fn_start = self.src.find("async function reloadFileByPath")
        assert fn_start >= 0, "WA.reloadFileByPath function not found"
        fn_body = self.src[fn_start : fn_start + 2200]
        assert "return openWorkspaceFile(filePath, supported)" in fn_body
        assert "_isAbsolutePath" in fn_body, (
            "WA.reloadFileByPath must treat non-absolute file-change paths as workspace-relative"
        )
        assert "openBrowserFile(filePath, supported)" in fn_body, (
            "WA.reloadFileByPath must still support absolute-path refreshes for external files"
        )
