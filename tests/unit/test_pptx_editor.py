"""
Unit tests for web/blueprints/pptx_editor.py

Tests the parsing, editing, and export functions in isolation using
an in-memory PPTX file built with python-pptx.
"""

from __future__ import annotations

import io
import json
import os

import pytest

# ── Helpers ─────────────────────────────────────────────────────────────────


def _make_pptx(slides: list[list[str]]) -> bytes:
    """Build a minimal .pptx in memory. `slides` is a list of lists of text strings."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    for texts in slides:
        slide_layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(slide_layout)
        for i, text in enumerate(texts):
            if i == 0:
                slide.shapes.title.text = text
            else:
                # Add an extra text box for subsequent entries
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(2 + i * 0.5), Inches(8), Inches(0.5)
                )
                txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── _parse_slides ─────────────────────────────────────────────────────────


class TestParseSlides:
    def test_slide_count(self):
        from web.blueprints.pptx_editor import _parse_slides

        raw = _make_pptx([["Title 1", "Body 1"], ["Title 2"]])
        result = _parse_slides(raw)
        assert len(result["slides"]) == 2

    def test_returns_width_height(self):
        from web.blueprints.pptx_editor import _parse_slides

        raw = _make_pptx([["Hello"]])
        result = _parse_slides(raw)
        assert result["slide_width_emu"] > 0
        assert result["slide_height_emu"] > 0

    def test_shapes_have_required_keys(self):
        from web.blueprints.pptx_editor import _parse_slides

        raw = _make_pptx([["Title", "Body"]])
        result = _parse_slides(raw)
        slide = result["slides"][0]
        assert "shapes" in slide
        assert "background" in slide
        assert "index" in slide
        for shape in slide["shapes"]:
            for key in (
                "id",
                "name",
                "type",
                "left",
                "top",
                "width",
                "height",
                "has_text",
            ):
                assert key in shape, f"missing key {key!r} in shape"

    def test_text_extracted(self):
        from web.blueprints.pptx_editor import _parse_slides

        raw = _make_pptx([["My Title"]])
        result = _parse_slides(raw)
        texts = []
        for shape in result["slides"][0]["shapes"]:
            for para in shape.get("paragraphs", []):
                for run in para.get("runs", []):
                    texts.append(run["text"])
        assert any("My Title" in t for t in texts)

    def test_empty_pptx(self):
        from pptx import Presentation
        from web.blueprints.pptx_editor import _parse_slides

        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        result = _parse_slides(buf.getvalue())
        assert result["slides"] == []

    def test_invalid_bytes_raises(self):
        from web.blueprints.pptx_editor import _parse_slides

        with pytest.raises(Exception):
            _parse_slides(b"not a pptx file")


# ── _apply_edits ─────────────────────────────────────────────────────────


class TestApplyEdits:
    def test_text_is_updated(self):
        from web.blueprints.pptx_editor import _apply_edits, _parse_slides
        from pptx import Presentation

        raw = _make_pptx([["Original Title"]])
        parsed = _parse_slides(raw)

        # Change the first text run in the title shape
        for shape in parsed["slides"][0]["shapes"]:
            if shape["has_text"]:
                for para in shape["paragraphs"]:
                    if para["runs"]:
                        para["runs"][0]["text"] = "Updated Title"
                        break
                break

        edited_bytes = _apply_edits(raw, parsed["slides"])

        # Re-parse to verify
        prs = Presentation(io.BytesIO(edited_bytes))
        texts = []
        for shape in prs.slides[0].shapes:
            if hasattr(shape, "text_frame"):
                texts.append(shape.text_frame.text)
        assert any("Updated Title" in t for t in texts)

    def test_original_not_mutated(self):
        from web.blueprints.pptx_editor import _apply_edits, _parse_slides

        raw = _make_pptx([["Stable"]])
        parsed = _parse_slides(raw)
        original_copy = bytes(raw)

        for shape in parsed["slides"][0]["shapes"]:
            if shape["has_text"]:
                for para in shape["paragraphs"]:
                    if para["runs"]:
                        para["runs"][0]["text"] = "Changed"
                        break

        _apply_edits(raw, parsed["slides"])
        assert raw == original_copy, "original bytes must not be modified"

    def test_missing_shape_id_is_skipped(self):
        """Edits referencing a non-existent shape_id must not crash."""
        from web.blueprints.pptx_editor import _apply_edits, _parse_slides

        raw = _make_pptx([["Text"]])
        parsed = _parse_slides(raw)

        # Corrupt shape id
        parsed["slides"][0]["shapes"][0]["id"] = 99999
        # Should not raise
        _apply_edits(raw, parsed["slides"])

    def test_returns_bytes(self):
        from web.blueprints.pptx_editor import _apply_edits, _parse_slides

        raw = _make_pptx([["Hello"]])
        parsed = _parse_slides(raw)
        result = _apply_edits(raw, parsed["slides"])
        assert isinstance(result, bytes)
        assert len(result) > 0

    def test_multiple_slides_edited_independently(self):
        from web.blueprints.pptx_editor import _apply_edits, _parse_slides
        from pptx import Presentation

        raw = _make_pptx([["Slide One"], ["Slide Two"]])
        parsed = _parse_slides(raw)

        # Only edit slide 0
        for shape in parsed["slides"][0]["shapes"]:
            if shape["has_text"]:
                for para in shape["paragraphs"]:
                    if para["runs"]:
                        para["runs"][0]["text"] = "EDITED"
                        break

        edited = _apply_edits(raw, parsed["slides"])
        prs = Presentation(io.BytesIO(edited))
        all_texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    all_texts.append(shape.text_frame.text)

        joined = " ".join(all_texts)
        assert "EDITED" in joined
        assert "Slide Two" in joined
