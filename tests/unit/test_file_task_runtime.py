import json
import base64
from pathlib import Path

import pytest

from app.core.agent.file_task_contract import (
    FileTaskFile,
    FileTaskIntentPlan,
    FileTaskLedger,
    FileTaskRequest,
    FileTaskToolStreamChunk,
    FileTaskToolStreamResult,
    event_to_sse,
)
from app.core.agent.file_task_runtime import FileTaskRuntime
from app.core.agent.file_task_model import FileTaskModelClient
from app.core.agent.file_task_tool_catalog import (
    file_task_tool_specs,
    supported_file_workflows,
)
from app.core.agent.file_task_tool_gateway import (
    FileTaskToolContext,
    FileTaskToolGateway,
)
from app.core.agent.file_task_workflow_state import build_workflow_state



def test_file_task_runtime_relays_streaming_tool_events_before_tool_finished():
    def fake_executor(tool_name, args):
        assert tool_name == "write_docx_content"
        return FileTaskToolStreamResult(
            chunks=[
                FileTaskToolStreamChunk(
                    kind="event",
                    event_type="step_progress",
                    payload={
                        "detail": "已写入 1/3 段落",
                        "progress": 33,
                        "level": "progress",
                        "file_updated": True,
                        "path": "draft.docx",
                        "file_path": "draft.docx",
                        "supported": True,
                    },
                ),
                FileTaskToolStreamChunk(
                    kind="result",
                    payload={
                        "success": True,
                        "path": "draft.docx",
                        "file_path": "draft.docx",
                        "file_type": "docx",
                        "operation": "write_docx_content",
                        "summary": "已将 3 段改写写回 draft.docx。",
                        "preview": "已写入 3 段",
                        "change_type": "rewrite",
                        "focus": True,
                    },
                ),
            ]
        )

    def fake_model(**kwargs):
        return {
            "content": "开始改写",
            "tool_calls": [
                {
                    "id": "write_demo",
                    "name": "write_docx_content",
                    "args": {
                        "path": "draft.docx",
                        "paragraphs": '[{"text":"改写后的内容"}]',
                    },
                }
            ],
        }

    request = FileTaskRequest(
        task="将文档的第一段改写后写回",
        run_id="stream_tool_demo",
        target_path="draft.docx",
        files=[
            FileTaskFile(path="draft.docx", name="draft.docx", type="docx", target=True)
        ],
    )

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=fake_model, max_rounds=1
        ).run(request)
    )

    progress_index = next(
        i for i, event in enumerate(events) if event.type == "step_progress"
    )
    tool_finished_index = next(
        i
        for i, event in enumerate(events)
        if event.type == "tool.finished"
        and event.payload.get("tool_name") == "write_docx_content"
    )
    file_changed = next(event for event in events if event.type == "file.changed")

    assert progress_index < tool_finished_index
    assert events[progress_index].payload["file_updated"] is True
    assert events[progress_index].payload["path"] == "draft.docx"
    assert file_changed.payload["path"] == "draft.docx"
    step_verified = next(
        event
        for event in events
        if event.type == "supervisor.step_verified"
        and event.payload.get("tool_name") == "write_docx_content"
    )
    assert step_verified.payload["passed"] is True
    assert step_verified.payload["outcome"] == "succeeded"
    assert step_verified.payload["file_changes"][0]["path"] == "draft.docx"



def test_file_task_runtime_merges_current_file_content_into_target_context():
    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task="总结当前文档",
        run_id="current_file_context_merge",
        target_path="workspace/interview.docx",
        current_file=FileTaskFile(
            path="workspace/interview.docx",
            name="interview.docx",
            type="docx",
            content="当前打开文档正文",
        ),
    )

    files = runtime._context_files(request)

    assert len(files) == 1
    assert files[0].path == "workspace/interview.docx"
    assert files[0].target is True
    assert files[0].content == "当前打开文档正文"






def test_context_files_resolves_workspace_relative_path_from_task(tmp_path):
    workspace = tmp_path / "workspace"
    target = workspace / "ui_smoke_tests" / "koto_ui_smoke_sales.xlsx"
    target.parent.mkdir(parents=True)
    target.write_bytes(b"fake xlsx")
    runtime = FileTaskRuntime(workspace_root=str(workspace))
    request = FileTaskRequest(
        task="请读取工作区 ui_smoke_tests/koto_ui_smoke_sales.xlsx，分析 Revenue 最高的客户。",
        run_id="task_text_relative_path",
    )

    files = runtime._context_files(request)

    assert [Path(file.path).resolve() for file in files] == [target.resolve()]
    assert files[0].name == "koto_ui_smoke_sales.xlsx"
    assert files[0].type == "xlsx"


def test_context_files_deduplicates_relative_and_absolute_workspace_paths(tmp_path, monkeypatch):
    workspace = tmp_path / "workspace"
    target = workspace / "sales.xlsx"
    workspace.mkdir()
    target.write_bytes(b"fake xlsx")
    monkeypatch.chdir(tmp_path)
    runtime = FileTaskRuntime(workspace_root=str(workspace))
    request = FileTaskRequest(
        task="请读取 workspace/sales.xlsx 并生成报告。",
        run_id="dedupe_relative_absolute_context",
        files=[
            FileTaskFile(
                path="workspace/sales.xlsx",
                name="sales.xlsx",
                type="xlsx",
            )
        ],
    )

    files = runtime._context_files(request)

    assert [file_info.name for file_info in files].count("sales.xlsx") == 1


def test_context_files_deduplicates_current_file_basename_with_explicit_workspace_match(tmp_path):
    workspace = tmp_path / "workspace"
    target = workspace / "_test_integration_workspace.txt"
    workspace.mkdir()
    target.write_text("workspace file content", encoding="utf-8")
    runtime = FileTaskRuntime(workspace_root=str(workspace))
    request = FileTaskRequest(
        task="请分析当前打开的 _test_integration_workspace.txt，用一句话总结内容。",
        run_id="dedupe_current_file_and_explicit_reference",
        files=[
            FileTaskFile(
                path="_test_integration_workspace.txt",
                name="_test_integration_workspace.txt",
                type="txt",
                content="workspace file content",
            )
        ],
        current_file=FileTaskFile(
            path="_test_integration_workspace.txt",
            name="_test_integration_workspace.txt",
            type="txt",
            content="workspace file content",
        ),
    )

    files = runtime._context_files(request)

    assert [file_info.name for file_info in files] == ["_test_integration_workspace.txt"]
    assert runtime._plan_summary(request, files, write_intent=False) == "准备处理 1 个文件。"


def test_context_files_resolves_unique_workspace_basename_from_task(tmp_path):
    workspace = tmp_path / "workspace"
    target = workspace / "销售台账.xlsx"
    workspace.mkdir()
    target.write_bytes(b"fake xlsx")
    runtime = FileTaskRuntime(workspace_root=str(workspace))
    request = FileTaskRequest(
        task="请读取我的工作区里的销售台账.xlsx，告诉我哪个客户贡献最高。",
        run_id="task_text_basename",
    )

    files = runtime._context_files(request)

    assert [Path(file.path).resolve() for file in files] == [target.resolve()]


def test_context_files_ignores_ambiguous_workspace_basename_from_task(tmp_path):
    workspace = tmp_path / "workspace"
    (workspace / "a").mkdir(parents=True)
    (workspace / "b").mkdir(parents=True)
    (workspace / "a" / "sales.xlsx").write_bytes(b"a")
    (workspace / "b" / "sales.xlsx").write_bytes(b"b")
    runtime = FileTaskRuntime(workspace_root=str(workspace))
    request = FileTaskRequest(
        task="请读取 sales.xlsx 并总结。",
        run_id="task_text_ambiguous_basename",
    )

    assert runtime._context_files(request) == []


def test_readonly_missing_file_reference_does_not_complete(tmp_path):
    runtime = FileTaskRuntime(
        workspace_root=str(tmp_path / "workspace"),
        model_client=lambda **kwargs: {
            "content": "根据错误信息，文件 missing.xlsx 未找到。",
            "tool_calls": [],
        },
        max_rounds=1,
    )
    request = FileTaskRequest(
        task="请读取工作区 missing.xlsx，分析 Revenue 最高的客户。不要修改任何文件。",
        run_id="missing_file_not_completed",
    )

    events = list(runtime.run(request))

    run_started = next(event for event in events if event.type == "run.started")
    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = next(event for event in events if event.type == "run.finished")
    assert check_finished.payload["passed"] is False
    assert check_finished.payload["status"] == "needs_attention"
    assert "没有成功读取任何显式文件上下文" in check_finished.payload["summary"]
    assert run_finished.payload["completed_task"] is False
    assert run_started.payload["performance"]["classification_ms"] >= 0
    assert run_started.payload["performance"]["intent_adjudication_ms"] >= 0
    assert check_finished.payload["performance"]["total_ms"] >= 0
    assert run_finished.payload["runtime"]["performance"]["total_ms"] >= 0


def test_readonly_directory_listing_does_not_satisfy_missing_file_reference(tmp_path):
    workspace = tmp_path / "workspace"
    workspace.mkdir()
    (workspace / "existing.xlsx").write_bytes(b"fake xlsx")

    responses = iter(
        [
            {
                "content": "",
                "tool_calls": [
                    {
                        "name": "list_workspace_files",
                        "args": {"path": ".", "recursive": True},
                    }
                ],
            },
            {"content": "目录里没有 missing.xlsx。", "tool_calls": []},
        ]
    )

    def fake_model(**kwargs):
        return next(responses)

    def fake_executor(tool_name, args):
        if tool_name == "list_workspace_files":
            return json.dumps(
                [{"name": "existing.xlsx", "size": 9}], ensure_ascii=False
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    runtime = FileTaskRuntime(
        workspace_root=str(workspace),
        model_client=fake_model,
        tool_executor=fake_executor,
        max_rounds=2,
    )
    request = FileTaskRequest(
        task="请读取工作区 missing.xlsx，告诉我里面的客户列表。不要修改任何文件。",
        run_id="missing_file_listing_not_completed",
    )

    events = list(runtime.run(request))

    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = next(event for event in events if event.type == "run.finished")
    assert check_finished.payload["passed"] is False
    assert check_finished.payload["status"] == "needs_attention"
    assert "missing.xlsx" in check_finished.payload["summary"]
    assert run_finished.payload["completed_task"] is False




def test_explicit_answer_mode_is_not_overridden_by_docx_review_words():
    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task="帮我看看这份 DOCX 哪里需要修改，先只给建议，不要写入文件",
        run_id="docx_review_answer_mode",
        target_path="draft.docx",
        files=[
            FileTaskFile(
                path="draft.docx", name="draft.docx", type="docx", target=True
            ),
        ],
        options={"output_mode": "answer"},
    )

    classification = runtime._classify_request(request, runtime._context_files(request))

    assert classification.output_mode == "answer"
    assert classification.write_intent is False
    assert classification.docx_annotation_request is False
    assert classification.execution_mode == "generic_tool_loop"
    # decision_trace mechanism was removed from _classify_request
    # (simplification: trace was never consumed in production)


def test_soft_edit_keywords_do_not_control_file_task_mainline():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {
            "content": "建议先调整结构，但本轮不写入。",
            "tool_calls": [],
        },
        max_rounds=1,
    )
    request = FileTaskRequest(
        task="优化这个 docx 的结构，先给我建议",
        run_id="soft_keywords_are_not_authority",
        target_path="draft.docx",
        files=[
            FileTaskFile(path="draft.docx", name="draft.docx", type="docx", target=True)
        ],
    )

    events = list(runtime.run(request))
    run_started = next(event for event in events if event.type == "run.started")

    assert run_started.payload["write_intent"] is False
    assert run_started.payload["output_mode"] == "answer"
    assert run_started.payload["execution_mode"] == "generic_tool_loop"
    assert run_started.payload["docx_annotation_request"] is False
    assert "mainline_contract:v1" in run_started.payload["reason_codes"]


def test_explicit_writeback_still_controls_file_task_mainline():
    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task="润色这个 docx 并写回当前 docx",
        run_id="explicit_writeback_authority",
        target_path="draft.docx",
        files=[
            FileTaskFile(path="draft.docx", name="draft.docx", type="docx", target=True)
        ],
    )

    classification = runtime._normalize_mainline_contract(
        request,
        request.files,
        runtime._classify_request(request, request.files),
    )

    assert classification.write_intent is True
    assert classification.output_mode == "write"
    assert classification.selected_recipe == "docx_polish_writeback"
    assert "mainline_contract:keyword_write_demoted" not in classification.reason_codes



def test_readonly_financial_keywords_do_not_preempt_mainline_runner():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {
            "content": "只读分析，不写入报告。",
            "tool_calls": [],
        },
        max_rounds=1,
    )
    request = FileTaskRequest(
        task="分析这个财务模型的收入图表问题，只回答，不要修改文件",
        run_id="readonly_financial_keywords_no_runner",
        files=[
            FileTaskFile(path="model.xlsx", name="model.xlsx", type="xlsx"),
            FileTaskFile(
                path="report.docx", name="report.docx", type="docx", target=True
            ),
        ],
    )

    events = list(runtime.run(request))
    run_started = next(event for event in events if event.type == "run.started")

    assert run_started.payload["mode"] == "whitebox_v1"
    assert run_started.payload["write_intent"] is False
    assert run_started.payload["selected_recipe"] != "financial_xlsx_docx_report"
    assert "mainline_contract:readonly_guard" in run_started.payload["reason_codes"]







def test_file_task_runtime_classifies_two_docx_compare_annotation_as_compare_write():
    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task="对比这两份文件，找出他们有区别的地方标注出来",
        run_id="docx_compare_annotation_classification",
        files=[
            FileTaskFile(path="humanise!.docx", name="humanise!.docx", type="docx"),
            FileTaskFile(
                path="humanise!_revised.docx",
                name="humanise!_revised.docx",
                type="docx",
                target=True,
            ),
        ],
        target_path="humanise!_revised.docx",
    )

    classification = runtime._classify_request(request, request.files)

    assert classification.task_family == "compare"
    assert classification.operation_kind == "compare_annotate"
    assert classification.execution_mode == "generic_tool_loop"
    assert classification.output_mode == "write"
    assert classification.selected_recipe == "docx_compare_annotation"
    assert "compare_docx_and_annotate" in classification.matched_capabilities
    assert "annotate_file" not in classification.matched_capabilities


def test_file_task_runtime_keeps_named_docx_compare_annotation_in_compare_mode():
    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "对比 humanise!.docx 和 humanise!_revised.docx，"
            "并在 humanise!.docx 上标注不同之处"
        ),
        run_id="docx_compare_annotation_named_target",
        files=[
            FileTaskFile(path="humanise!.docx", name="humanise!.docx", type="docx"),
            FileTaskFile(
                path="humanise!_revised.docx",
                name="humanise!_revised.docx",
                type="docx",
            ),
        ],
        target_path="humanise!.docx",
    )

    classification = runtime._classify_request(request, request.files)

    assert classification.task_family == "compare"
    assert classification.operation_kind == "compare_annotate"
    assert classification.execution_mode == "generic_tool_loop"
    assert classification.output_mode == "write"
    assert classification.selected_recipe == "docx_compare_annotation"
    assert "write_docx_comments" in classification.matched_capabilities
    assert "annotate_file" not in classification.matched_capabilities


def test_file_task_runtime_limits_docx_compare_route_tools():
    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task="对比两份 DOCX，并在原文上标注不同之处",
        files=[
            FileTaskFile(path="old.docx", name="old.docx", type="docx", target=True),
            FileTaskFile(path="new.docx", name="new.docx", type="docx"),
        ],
        target_path="old.docx",
    )
    classification = runtime._classify_request(request, request.files)

    tool_defs = runtime._tool_defs_for_classification(
        [
            {"name": "parse_file_to_text"},
            {"name": "plan_docx_compare_annotations"},
            {"name": "write_docx_comments"},
            {"name": "compare_docx_and_annotate"},
            {"name": "verify_task_completion"},
            {"name": "annotate_file"},
            {"name": "write_docx_content"},
        ],
        classification,
    )
    tool_names = {str(item.get("name") or "") for item in tool_defs}

    assert classification.selected_recipe == "docx_compare_annotation"
    assert {
        "parse_file_to_text",
        "plan_docx_compare_annotations",
        "write_docx_comments",
        "compare_docx_and_annotate",
        "verify_task_completion",
    }.issubset(tool_names)
    assert "annotate_file" not in tool_names
    assert "write_docx_content" not in tool_names


def test_file_task_runtime_executes_two_docx_compare_annotation_through_model_loop():
    calls = []
    model_calls = []

    def fake_executor(tool_name, args):
        calls.append((tool_name, dict(args)))
        if tool_name == "parse_file_to_text":
            return "DOCX context"
        if tool_name == "plan_docx_compare_annotations":
            return json.dumps(
                {
                    "success": True,
                    "operation": "plan_docx_compare_annotations",
                    "target_path": args["target_path"],
                    "differences_detected": 2,
                    "annotation_candidates": [
                        {
                            "原文片段": "付款期限为30日",
                            "批注内容": "另一版为：付款期限为15日\n本版为：付款期限为30日",
                        }
                    ],
                },
                ensure_ascii=False,
            )
        if tool_name == "write_docx_comments":
            assert args["path"] == "old.docx"
            assert "付款期限为30日" in args["comments_json"]
            return json.dumps(
                {
                    "success": True,
                    "path": args["path"],
                    "file_path": args["path"],
                    "file_type": "docx",
                    "operation": "write_docx_comments",
                    "summary": "已在目标 DOCX 原文上写入 1 条 Word 批注。",
                    "change_type": "annotate",
                    "annotations_added": 1,
                    "differences_detected": 2,
                    "focus": True,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "summary": "已核验 DOCX 差异批注。"},
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool: {tool_name}")

    def fake_model(**kwargs):
        model_calls.append(kwargs)
        # Adjudicator call: no tools, adjudicator system prompt
        if not kwargs.get("tools") and "任务意图裁判" in str(kwargs.get("system", "")):
            return {
                "content": json.dumps({
                    "intent": "edit_file",
                    "confidence": 0.90,
                    "should_write": True,
                    "should_use_annotate_bridge": False,
                    "reason": "Two-DOCX comparison uses compare tool, not annotation bridge"
                }, ensure_ascii=False),
                "tool_calls": [],
            }
        if len(model_calls) == 2:
            return {
                "content": "我将比较两份 DOCX 并把差异写成批注。",
                "tool_calls": [
                    {
                        "id": "plan_docx",
                        "name": "plan_docx_compare_annotations",
                        "args": {
                            "original_path": "old.docx",
                            "revised_path": "new.docx",
                            "target_path": "old.docx",
                        },
                    }
                ],
            }
        if len(model_calls) == 3:
            return {
                "content": "我将把差异作为 Word 批注写回原文。",
                "tool_calls": [
                    {
                        "id": "write_comments",
                        "name": "write_docx_comments",
                        "args": {
                            "path": "old.docx",
                            "comments_json": json.dumps(
                                [
                                    {
                                        "原文片段": "付款期限为30日",
                                        "批注内容": "另一版为：付款期限为15日\n本版为：付款期限为30日",
                                    }
                                ],
                                ensure_ascii=False,
                            ),
                            "source_path": "old.docx",
                            "compare_path": "new.docx",
                            "differences_detected": 2,
                        },
                    }
                ],
            }
        return {"content": "已完成差异批注。", "tool_calls": []}

    request = FileTaskRequest(
        task="对比这两份文件，找出他们有区别的地方标注出来",
        run_id="docx_compare_annotation_model_loop",
        files=[
            FileTaskFile(path="old.docx", name="old.docx", type="docx", target=True),
            FileTaskFile(path="new.docx", name="new.docx", type="docx"),
        ],
        target_path="old.docx",
    )

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor,
            model_client=fake_model,
            max_rounds=3,
        ).run(request)
    )

    run_finished = next(
        event for event in reversed(events) if event.type == "run.finished"
    )
    file_changed = next(event for event in events if event.type == "file.changed")

    assert model_calls
    assert (
        "plan_docx_compare_annotations",
        {
            "original_path": "old.docx",
            "revised_path": "new.docx",
            "target_path": "old.docx",
        },
    ) in calls
    assert any(tool_name == "write_docx_comments" for tool_name, _args in calls)
    assert any(tool_name == "verify_task_completion" for tool_name, _args in calls)
    assert file_changed.payload["path"] == "old.docx"
    assert file_changed.payload["operation"] == "write_docx_comments"
    assert file_changed.payload["annotations_added"] == 1
    assert run_finished.payload["completed_task"] is True
    assert run_finished.payload["operation_kind"] == "compare_annotate"
    assert run_finished.payload["selected_recipe"] == "docx_compare_annotation"
    assert not any(
        event.payload.get("tool_name") == "annotate_file"
        for event in events
        if isinstance(event.payload, dict)
    )


def test_file_task_runtime_contract_compare_returns_risk_summary():
    model_calls = []

    def fake_executor(tool_name, args):
        if tool_name == "parse_file_to_text":
            return "contract context"
        if tool_name == "plan_docx_compare_annotations":
            return json.dumps(
                {
                    "success": True,
                    "operation": "plan_docx_compare_annotations",
                    "target_path": args["target_path"],
                    "differences_detected": 2,
                    "annotation_candidates": [
                        {
                            "原文片段": "付款期限为30日",
                            "批注内容": "另一版为：付款期限为15日\n本版为：付款期限为30日",
                        }
                    ],
                },
                ensure_ascii=False,
            )
        if tool_name == "write_docx_comments":
            assert args["path"] == "old_contract.docx"
            assert "风险" in args["comments_json"]
            return json.dumps(
                {
                    "success": True,
                    "path": args["path"],
                    "file_path": args["path"],
                    "file_type": "docx",
                    "operation": "write_docx_comments",
                    "summary": "已在目标 DOCX 原文上写入 1 条 Word 批注。",
                    "change_type": "annotate",
                    "annotations_added": 1,
                    "differences_detected": 2,
                    "contract_risk_summary": [
                        "付款/费用：付款或费用条款发生变化，需确认金额、期限、开票和逾期责任是否可接受。",
                    ],
                    "focus": True,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "summary": "已核验合同差异批注。"},
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool: {tool_name}")

    def fake_model(**kwargs):
        model_calls.append(kwargs)
        # Adjudicator call: no tools, adjudicator system prompt
        if not kwargs.get("tools") and "任务意图裁判" in str(kwargs.get("system", "")):
            return {
                "content": json.dumps({
                    "intent": "edit_file",
                    "confidence": 0.90,
                    "should_write": True,
                    "should_use_annotate_bridge": False,
                    "reason": "Two-DOCX comparison uses compare tool, not annotation bridge"
                }, ensure_ascii=False),
                "tool_calls": [],
            }
        if len(model_calls) == 1:
            return {
                "content": "我将对比两份合同并写入差异批注。",
                "tool_calls": [
                    {
                        "id": "plan_contract_docx",
                        "name": "plan_docx_compare_annotations",
                        "args": {
                            "original_path": "old_contract.docx",
                            "revised_path": "new_contract.docx",
                            "target_path": "old_contract.docx",
                        },
                    }
                ],
            }
        if len(model_calls) == 2:
            return {
                "content": "我将把合同差异、风险和建议写回原合同批注。",
                "tool_calls": [
                    {
                        "id": "write_contract_comments",
                        "name": "write_docx_comments",
                        "args": {
                            "path": "old_contract.docx",
                            "comments_json": json.dumps(
                                [
                                    {
                                        "原文片段": "付款期限为30日",
                                        "批注内容": (
                                            "另一版为：付款期限为15日\n"
                                            "本版为：付款期限为30日\n"
                                            "风险：付款周期缩短，可能增加现金流压力。\n"
                                            "建议：确认是否接受该付款安排。"
                                        ),
                                    }
                                ],
                                ensure_ascii=False,
                            ),
                            "source_path": "old_contract.docx",
                            "compare_path": "new_contract.docx",
                            "differences_detected": 2,
                        },
                    }
                ],
            }
        return {"content": "已完成合同差异批注，并整理风险点。", "tool_calls": []}

    request = FileTaskRequest(
        task="对比这两份合同，找出变化并标注出来，同时总结风险点",
        run_id="docx_contract_compare_review_demo",
        files=[
            FileTaskFile(
                path="old_contract.docx",
                name="old_contract.docx",
                type="docx",
                target=True,
            ),
            FileTaskFile(
                path="new_contract.docx",
                name="new_contract.docx",
                type="docx",
            ),
        ],
        target_path="old_contract.docx",
    )

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor,
            model_client=fake_model,
            max_rounds=3,
        ).run(request)
    )

    run_finished = next(
        event for event in reversed(events) if event.type == "run.finished"
    )
    file_changed = next(event for event in events if event.type == "file.changed")

    assert model_calls
    assert run_finished.payload["selected_recipe"] == "docx_contract_compare_review"
    assert run_finished.payload["task_family"] == "contract_review"
    assert run_finished.payload["completed_task"] is True
    assert "风险关注点" in run_finished.payload["summary"]
    assert "付款/费用" in run_finished.payload["summary"]
    assert file_changed.payload["path"] == "old_contract.docx"
    assert file_changed.payload["operation"] == "write_docx_comments"
    assert file_changed.payload["contract_risk_summary"]

    # decision_trace mechanism was removed from _classify_request




def test_file_task_runtime_generic_office_quality_gate_rejects_unstructured_docx_write():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="修改这个 Word 文档并保存",
        run_id="generic_docx_gate",
        target_path="draft.docx",
        files=[
            FileTaskFile(path="draft.docx", name="draft.docx", type="docx", target=True)
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [{"path": "draft.docx", "operation": "run_python_code", "file_type": "docx"}],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is False
    assert any(
        item["criterion"] == "generic_docx_has_native_write"
        for item in result["criteria_results"]
    )


def test_file_task_runtime_generic_office_quality_gate_accepts_native_docx_write():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="润色这个 Word 文档并写回",
        run_id="generic_docx_gate_pass",
        target_path="draft.docx",
        files=[
            FileTaskFile(path="draft.docx", name="draft.docx", type="docx", target=True)
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "draft.docx",
                "operation": "write_docx_content",
                "file_type": "docx",
                "paragraphs_written": 3,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True


def test_file_task_runtime_quality_gate_accepts_docx_template_fill():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="Fill the placeholders in this Word contract template.",
        run_id="docx_template_fill_gate",
        target_path="filled.docx",
        files=[
            FileTaskFile(path="template.docx", name="template.docx", type="docx")
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "filled.docx",
                "operation": "fill_docx_template",
                "file_type": "docx",
                "placeholders_replaced": 2,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "docx_template_fill_replaces_placeholders"
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_docx_pdf_export():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="Export the current Word document as PDF.",
        run_id="docx_pdf_export_gate",
        target_path="report.docx",
        files=[
            FileTaskFile(path="report.docx", name="report.docx", type="docx", target=True)
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "report.pdf",
                "operation": "convert_docx_to_pdf",
                "file_type": "pdf",
                "converter": "fake",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "docx_pdf_export_uses_converter"
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_generic_file_convert():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="Convert notes.txt to markdown and save it as notes.md.",
        run_id="file_format_convert_gate",
        target_path="notes.md",
        files=[
            FileTaskFile(path="notes.txt", name="notes.txt", type="txt")
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "notes.md",
                "operation": "convert_file",
                "file_type": "md",
                "target_format": "md",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "file_format_convert_uses_converter"
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_local_docx_paragraph_insert():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "请继续优化 workspace/report.docx：只追加一句"
            "“Overall risk level: Moderate.”，保留已有表格不变，保存同一个 DOCX。"
        ),
        run_id="local_docx_insert_gate",
        target_path="workspace/report.docx",
        files=[
            FileTaskFile(
                path="workspace/report.docx",
                name="report.docx",
                type="docx",
                target=True,
            )
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "workspace/report.docx",
                "operation": "insert_docx_paragraph",
                "file_type": "docx",
                "paragraphs_written": 1,
                "inserted_text": "Overall risk level: Moderate.",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    criteria = {item["criterion"] for item in result["criteria_results"]}
    assert "docx_local_edit_has_paragraph_insert" in criteria
    assert "docx_report_has_narrative" not in criteria
    assert "docx_table_request_has_table" not in criteria


def test_file_task_runtime_quality_gate_rejects_docx_missing_requested_source_content(tmp_path):
    from docx import Document

    source_path = tmp_path / "workspace" / "_test_integration_workspace.txt"
    source_path.parent.mkdir()
    source_path.write_text("workspace file content", encoding="utf-8")
    target_path = tmp_path / "workspace" / "koto_frontend_ai_local_test.docx"
    document = Document()
    document.add_paragraph("Koto Frontend Local AI Test")
    document.add_paragraph(
        "This was executed through the Koto workspace assistant frontend using the local model."
    )
    document.save(str(target_path))
    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "Read the currently opened _test_integration_workspace.txt and create a new Word "
            f"file at {target_path}. Include the title, the original content, and one "
            "sentence saying this was executed through the Koto workspace assistant frontend "
            "using the local model. Do not modify the original txt file."
        ),
        target_path=str(target_path),
        files=[
            FileTaskFile(
                path=str(source_path),
                name="_test_integration_workspace.txt",
                type="txt",
                content="workspace file content",
            )
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(target_path),
                "file_type": "docx",
                "operation": "write_docx_content",
                "paragraphs_written": 3,
                "summary": "已写入 3 个段落到 Word 文档",
                "preview": "Koto Frontend Local AI Test\nThis was executed through the Koto workspace assistant frontend using the local model.",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is False
    assert any(
        item["criterion"] == "source_content_included" and not item["passed"]
        for item in result["criteria_results"]
    )
    assert any("_test_integration_workspace.txt" in item for item in result["remaining"])


def test_file_task_runtime_quality_gate_accepts_docx_with_requested_source_content(tmp_path):
    from docx import Document

    source_path = tmp_path / "workspace" / "_test_integration_workspace.txt"
    source_path.parent.mkdir()
    source_path.write_text("workspace file content", encoding="utf-8")
    target_path = tmp_path / "workspace" / "koto_frontend_ai_local_test.docx"
    document = Document()
    document.add_paragraph("Koto Frontend Local AI Test")
    document.add_paragraph("workspace file content")
    document.add_paragraph(
        "This was executed through the Koto workspace assistant frontend using the local model."
    )
    document.save(str(target_path))
    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "Read the currently opened _test_integration_workspace.txt and create a new Word "
            f"file at {target_path}. Include the title, the original content, and one "
            "sentence saying this was executed through the Koto workspace assistant frontend "
            "using the local model. Do not modify the original txt file."
        ),
        target_path=str(target_path),
        files=[
            FileTaskFile(
                path=str(source_path),
                name="_test_integration_workspace.txt",
                type="txt",
                content="workspace file content",
            )
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(target_path),
                "file_type": "docx",
                "operation": "write_docx_content",
                "paragraphs_written": 3,
                "summary": "已写入 3 个段落到 Word 文档",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "source_content_included" and item["passed"]
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_rejects_unsorted_top_n_docx_table(tmp_path):
    import openpyxl
    from docx import Document

    source_path = tmp_path / "sales.xlsx"
    target_path = tmp_path / "report.docx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    sheet.append(["Customer", "Region", "Revenue", "Margin", "Risk"])
    sheet.append(["Northwind Labs", "NA", 128000, 0.34, "Security review"])
    sheet.append(["Aurora Retail", "EU", 96000, 0.28, "Payment terms"])
    sheet.append(["Blue Harbor", "APAC", 142000, 0.31, "Capacity"])
    sheet.append(["Delta Foods", "EU", 118000, 0.37, "Upsell"])
    workbook.save(source_path)

    document = Document()
    document.add_paragraph("Top 3 Customers by Revenue")
    table = document.add_table(rows=4, cols=4)
    for column, header in enumerate(["Customer", "Region", "Revenue", "Margin"]):
        table.cell(0, column).text = header
    for row_index, row_values in enumerate(
        [
            ["Northwind Labs", "NA", "128000", "0.34"],
            ["Aurora Retail", "EU", "96000", "0.28"],
            ["Blue Harbor", "APAC", "142000", "0.31"],
        ],
        start=1,
    ):
        for column, value in enumerate(row_values):
            table.cell(row_index, column).text = value
    document.save(target_path)

    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "Read sales.xlsx and create report.docx with a top three customers by "
            "Revenue table."
        ),
        target_path=str(target_path),
        files=[FileTaskFile(path=str(source_path), name="sales.xlsx", type="xlsx")],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(target_path),
                "source_path": str(source_path),
                "file_type": "docx",
                "operation": "insert_excel_as_docx_table",
                "rows_written": 3,
                "columns_written": 4,
                "table_title": "Top 3 Customers by Revenue",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is False
    assert any(
        item["criterion"] == "top_table_sorted_by_requested_metric"
        and not item["passed"]
        for item in result["criteria_results"]
    )
    assert any("Blue Harbor" in item and "Delta Foods" in item for item in result["remaining"])


def test_file_task_runtime_quality_gate_accepts_sorted_top_n_docx_table(tmp_path):
    import openpyxl
    from docx import Document

    source_path = tmp_path / "sales.xlsx"
    target_path = tmp_path / "report.docx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    sheet.append(["Customer", "Region", "Revenue", "Margin", "Risk"])
    sheet.append(["Northwind Labs", "NA", 128000, 0.34, "Security review"])
    sheet.append(["Aurora Retail", "EU", 96000, 0.28, "Payment terms"])
    sheet.append(["Blue Harbor", "APAC", 142000, 0.31, "Capacity"])
    sheet.append(["Delta Foods", "EU", 118000, 0.37, "Upsell"])
    workbook.save(source_path)

    document = Document()
    document.add_paragraph("Top 3 Customers by Revenue")
    table = document.add_table(rows=4, cols=4)
    for column, header in enumerate(["Customer", "Region", "Revenue", "Margin"]):
        table.cell(0, column).text = header
    for row_index, row_values in enumerate(
        [
            ["Blue Harbor", "APAC", "142000", "0.31"],
            ["Northwind Labs", "NA", "128000", "0.34"],
            ["Delta Foods", "EU", "118000", "0.37"],
        ],
        start=1,
    ):
        for column, value in enumerate(row_values):
            table.cell(row_index, column).text = value
    document.save(target_path)

    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "Read sales.xlsx and create report.docx with a top three customers by "
            "Revenue table."
        ),
        target_path=str(target_path),
        files=[FileTaskFile(path=str(source_path), name="sales.xlsx", type="xlsx")],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(target_path),
                "source_path": str(source_path),
                "file_type": "docx",
                "operation": "insert_excel_as_docx_table",
                "rows_written": 3,
                "columns_written": 4,
                "table_title": "Top 3 Customers by Revenue",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "top_table_sorted_by_requested_metric"
        and item["passed"]
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_rejects_extra_columns_when_table_columns_requested(tmp_path):
    import openpyxl
    from docx import Document

    source_path = tmp_path / "sales.xlsx"
    target_path = tmp_path / "report.docx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    sheet.append(["Customer", "Region", "Revenue", "Margin", "Risk"])
    sheet.append(["Northwind Labs", "NA", 128000, 0.34, "Security review"])
    sheet.append(["Aurora Retail", "EU", 96000, 0.28, "Payment terms"])
    sheet.append(["Blue Harbor", "APAC", 142000, 0.31, "Capacity"])
    sheet.append(["Delta Foods", "EU", 118000, 0.37, "Upsell"])
    workbook.save(source_path)

    document = Document()
    document.add_paragraph("Top 3 Customers by Revenue")
    table = document.add_table(rows=4, cols=5)
    for column, header in enumerate(["Customer", "Region", "Revenue", "Margin", "Risk"]):
        table.cell(0, column).text = header
    for row_index, row_values in enumerate(
        [
            ["Blue Harbor", "APAC", "142000", "0.31", "Capacity"],
            ["Northwind Labs", "NA", "128000", "0.34", "Security review"],
            ["Delta Foods", "EU", "118000", "0.37", "Upsell"],
        ],
        start=1,
    ):
        for column, value in enumerate(row_values):
            table.cell(row_index, column).text = value
    document.save(target_path)

    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "Read sales.xlsx and create report.docx with a top three customers by "
            "Revenue table with Customer, Region, Revenue and Margin."
        ),
        target_path=str(target_path),
        files=[FileTaskFile(path=str(source_path), name="sales.xlsx", type="xlsx")],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(target_path),
                "source_path": str(source_path),
                "file_type": "docx",
                "operation": "insert_excel_as_docx_table",
                "rows_written": 3,
                "columns_written": 5,
                "table_title": "Top 3 Customers by Revenue",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is False
    assert any(
        item["criterion"] == "top_table_sorted_by_requested_metric"
        and "列未严格匹配" in item["detail"]
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_uses_request_source_when_change_has_short_source_path(tmp_path):
    import openpyxl
    from docx import Document

    source_path = tmp_path / "workspace" / "sales.xlsx"
    target_path = tmp_path / "workspace" / "report.docx"
    source_path.parent.mkdir()
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    sheet.append(["Customer", "Region", "Revenue", "Margin"])
    sheet.append(["Northwind Labs", "NA", 128000, 0.34])
    sheet.append(["Blue Harbor", "APAC", 142000, 0.31])
    sheet.append(["Delta Foods", "EU", 118000, 0.37])
    workbook.save(source_path)

    document = Document()
    table = document.add_table(rows=4, cols=4)
    for column, header in enumerate(["Customer", "Region", "Revenue", "Margin"]):
        table.cell(0, column).text = header
    for row_index, row_values in enumerate(
        [
            ["Blue Harbor", "APAC", "142000", "0.31"],
            ["Northwind Labs", "NA", "128000", "0.34"],
            ["Delta Foods", "EU", "118000", "0.37"],
        ],
        start=1,
    ):
        for column, value in enumerate(row_values):
            table.cell(row_index, column).text = value
    document.save(target_path)

    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "Read workspace/sales.xlsx and create report.docx with a top three "
            "customers by Revenue table with Customer, Region, Revenue and Margin."
        ),
        target_path=str(target_path),
        files=[FileTaskFile(path=str(source_path), name="sales.xlsx", type="xlsx")],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(target_path),
                "source_path": "sales.xlsx",
                "file_type": "docx",
                "operation": "insert_excel_as_docx_table",
                "rows_written": 3,
                "columns_written": 4,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "top_table_sorted_by_requested_metric"
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_rejects_duplicate_top_table_rows_in_paragraphs(tmp_path):
    import openpyxl
    from docx import Document

    source_path = tmp_path / "sales.xlsx"
    target_path = tmp_path / "report.docx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    sheet.append(["Customer", "Region", "Revenue", "Margin"])
    sheet.append(["Northwind Labs", "NA", 128000, 0.34])
    sheet.append(["Blue Harbor", "APAC", 142000, 0.31])
    sheet.append(["Delta Foods", "EU", 118000, 0.37])
    workbook.save(source_path)

    document = Document()
    document.add_paragraph("Top Three Customers by Revenue")
    document.add_paragraph("Customer: Blue Harbor | Region: APAC | Revenue: 142000 | Margin: 0.31")
    document.add_paragraph("Customer: Northwind Labs | Region: NA | Revenue: 128000 | Margin: 0.34")
    table = document.add_table(rows=4, cols=4)
    for column, header in enumerate(["Customer", "Region", "Revenue", "Margin"]):
        table.cell(0, column).text = header
    for row_index, row_values in enumerate(
        [
            ["Blue Harbor", "APAC", "142000", "0.31"],
            ["Northwind Labs", "NA", "128000", "0.34"],
            ["Delta Foods", "EU", "118000", "0.37"],
        ],
        start=1,
    ):
        for column, value in enumerate(row_values):
            table.cell(row_index, column).text = value
    document.save(target_path)

    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "Read sales.xlsx and create report.docx with a top three customers by "
            "Revenue table with Customer, Region, Revenue and Margin."
        ),
        target_path=str(target_path),
        files=[FileTaskFile(path=str(source_path), name="sales.xlsx", type="xlsx")],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(target_path),
                "source_path": str(source_path),
                "file_type": "docx",
                "operation": "insert_excel_as_docx_table",
                "rows_written": 3,
                "columns_written": 4,
                "table_title": "Top 3 Customers by Revenue",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is False
    assert any(
        item["criterion"] == "top_table_sorted_by_requested_metric"
        and "重复写成了段落清单" in item["detail"]
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_rejects_missing_required_sections(tmp_path):
    from docx import Document

    target_path = tmp_path / "report.docx"
    document = Document()
    document.add_paragraph("Koto Complex Multi File Report")
    document.add_paragraph("Executive Summary")
    document.add_paragraph(
        "This report mentions risk profiles and next actions in passing, but omits the required sections."
    )
    document.add_paragraph("Top Three Customers by Revenue")
    document.save(target_path)

    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "Create report.docx. The report must include a risk section combining risks "
            "from both files and three concrete next actions."
        ),
        target_path=str(target_path),
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(target_path),
                "file_type": "docx",
                "operation": "write_docx_content",
                "paragraphs_written": 4,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is False
    assert any(
        item["criterion"] == "required_risk_section_present" and not item["passed"]
        for item in result["criteria_results"]
    )
    assert any(
        item["criterion"] == "required_next_actions_present" and not item["passed"]
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_required_sections_and_actions(tmp_path):
    from docx import Document

    target_path = tmp_path / "report.docx"
    document = Document()
    document.add_paragraph("Koto Complex Multi File Report")
    document.add_paragraph("Executive Summary")
    document.add_paragraph("Top Three Customers by Revenue")
    document.add_paragraph("Risk Section")
    document.add_paragraph("Blue Harbor has implementation capacity risk.")
    document.add_paragraph("Next Actions")
    document.add_paragraph("1. Prioritize Blue Harbor delivery staffing.")
    document.add_paragraph("2. Initiate Northwind Labs security review.")
    document.add_paragraph("3. Monitor Aurora Retail payment terms.")
    document.save(target_path)

    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "Create report.docx. The report must include a risk section combining risks "
            "from both files and three concrete next actions."
        ),
        target_path=str(target_path),
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(target_path),
                "file_type": "docx",
                "operation": "write_docx_content",
                "paragraphs_written": 9,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "required_risk_section_present" and item["passed"]
        for item in result["criteria_results"]
    )
    assert any(
        item["criterion"] == "required_next_actions_present" and item["passed"]
        for item in result["criteria_results"]
    )


def test_file_task_runtime_generic_office_quality_gate_accepts_native_pptx_write():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="编辑这个 PPT 并保存",
        run_id="generic_pptx_gate_pass",
        target_path="deck.pptx",
        files=[
            FileTaskFile(path="deck.pptx", name="deck.pptx", type="pptx", target=True)
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "deck.pptx",
                "operation": "design_pptx_theme_layout",
                "file_type": "pptx",
                "slides_designed": 6,
                "text_shapes_styled": 14,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True


def test_file_task_runtime_quality_gate_accepts_pptx_explicit_phrases_and_slide_count(
    tmp_path,
):
    from pptx import Presentation

    pptx_path = tmp_path / "operations_update_deck.pptx"
    presentation = Presentation()
    slides = [
        (
            "Executive Summary",
            "Northwind Retail depends on manual reconciliation. Decision date: August 15, 2026.",
        ),
        (
            "Current Risk",
            "Client: Northwind Retail. Current risk: manual reconciliation delays reporting.",
        ),
        (
            "Pilot Decision",
            "Approve the automation pilot by August 15, 2026 for Northwind Retail.",
        ),
    ]
    for title, body in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    presentation.save(pptx_path)

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "新演示文稿必须正好 3 页幻灯片，三页标题分别为 "
            "Executive Summary、Current Risk、Pilot Decision；"
            "内容必须包含客户 Northwind Retail，包含短语 manual reconciliation，"
            "包含日期 August 15, 2026。"
        ),
        run_id="pptx_explicit_phrases_and_slide_count",
        target_path=str(pptx_path),
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(pptx_path),
                "operation": "write_pptx_slides",
                "file_type": "pptx",
                "slides_added": 3,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "explicit_pptx_slide_count" and item["passed"] is True
        for item in result["criteria_results"]
    )
    assert any(
        item["criterion"] == "explicit_required_phrases_present"
        and item["passed"] is True
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_rejects_wrong_pptx_slide_count(tmp_path):
    from pptx import Presentation

    pptx_path = tmp_path / "operations_update_deck.pptx"
    presentation = Presentation()
    for title in ("Executive Summary", "Current Risk"):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = (
            "Northwind Retail uses manual reconciliation. August 15, 2026."
        )
    presentation.save(pptx_path)

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "新演示文稿必须正好 3 页幻灯片，三页标题分别为 "
            "Executive Summary、Current Risk、Pilot Decision；"
            "内容必须包含客户 Northwind Retail，包含短语 manual reconciliation，"
            "包含日期 August 15, 2026。"
        ),
        run_id="pptx_wrong_slide_count",
        target_path=str(pptx_path),
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(pptx_path),
                "operation": "write_pptx_slides",
                "file_type": "pptx",
                "slides_added": 2,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is False
    assert any(
        item["criterion"] == "explicit_pptx_slide_count" and item["passed"] is False
        for item in result["criteria_results"]
    )


def test_file_task_runtime_generic_office_quality_gate_accepts_native_xlsx_write():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="创建一个 Excel 表格",
        run_id="generic_xlsx_gate",
        target_path="analysis.xlsx",
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "analysis.xlsx",
                "operation": "write_sheet_data",
                "file_type": "xlsx",
                "rows_written": 4,
                "cells_written": 12,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True


def test_file_task_runtime_generic_office_quality_gate_accepts_python_xlsx_metrics():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="创建一个 Excel 表格",
        run_id="generic_python_xlsx_gate",
        target_path="analysis.xlsx",
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "analysis.xlsx",
                "operation": "run_python_code",
                "file_type": "xlsx",
                "rows_written": 4,
                "cells_written": 20,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True


def test_file_task_runtime_quality_gate_rejects_xlsx_chart_request_without_chart(tmp_path):
    import openpyxl

    workbook_path = tmp_path / "sales_profit_report.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    sheet.append(["month", "revenue", "cost", "profit"])
    sheet.append(["Jan", 12000, 7100, 4900])
    workbook.save(workbook_path)

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="生成一个包含 profit 列和月度 profit 折线图的 Excel 文件",
        run_id="xlsx_chart_gate_rejects_missing_chart",
        target_path=str(workbook_path),
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(workbook_path),
                "operation": "write_sheet_data",
                "file_type": "xlsx",
                "cells_written": 7,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is False
    assert any(
        item["criterion"] == "spreadsheet_chart_request_has_workbook_chart"
        and item["passed"] is False
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_rejects_missing_explicit_docx_phrase(tmp_path):
    from docx import Document

    workbook_path = tmp_path / "service_note_revised.docx"
    document = Document()
    document.add_heading("Operations Action Plan", level=1)
    document.add_paragraph("Client: Northwind Retail")
    document.add_paragraph(
        "1. Automate reporting - Owner: Sarah Chen, Due: 21 July 2026."
    )
    document.add_paragraph("Ready for operations leadership.")
    document.save(workbook_path)

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "生成 service_note_revised.docx，包含标题 Operations Action Plan，"
            "保留客户 Northwind Retail，每条包含 Owner 和 Due date，"
            "最后包含一句 Ready for operations leadership。"
        ),
        run_id="docx_missing_explicit_phrase",
        target_path=str(workbook_path),
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(workbook_path),
                "operation": "write_docx_content",
                "file_type": "docx",
                "paragraphs_written": 4,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is False
    assert any(
        item["criterion"] == "explicit_required_phrases_present"
        and item["passed"] is False
        and "Due date" in item["detail"]
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_explicit_docx_phrases(tmp_path):
    from docx import Document

    docx_path = tmp_path / "service_note_revised.docx"
    document = Document()
    document.add_heading("Operations Action Plan", level=1)
    document.add_paragraph("Client: Northwind Retail")
    document.add_paragraph(
        "1. Automate reporting - Owner: Sarah Chen, Due date: 21 July 2026."
    )
    document.add_paragraph("Ready for operations leadership.")
    document.save(docx_path)

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "生成 service_note_revised.docx，包含标题 Operations Action Plan，"
            "保留客户 Northwind Retail，每条包含 Owner 和 Due date，"
            "最后包含一句 Ready for operations leadership。"
        ),
        run_id="docx_has_explicit_phrases",
        target_path=str(docx_path),
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(docx_path),
                "operation": "write_docx_content",
                "file_type": "docx",
                "paragraphs_written": 4,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True


def test_file_task_runtime_quality_gate_accepts_pdf_markdown_key_facts_bullets(tmp_path):
    markdown_path = tmp_path / "operations_brief_summary.md"
    markdown_path.write_text(
        "\n".join(
            [
                "# PDF Operations Summary",
                "",
                "**Client:** Northwind Retail",
                "",
                "The brief depends on manual reconciliation.",
                "",
                "Decision date: August 15, 2026",
                "",
                "## Key facts",
                "",
                "- Northwind Retail relies on manual reconciliation.",
                "- Leadership receives the packet five business days late.",
                "- The automation pilot decision is due by August 15, 2026.",
            ]
        ),
        encoding="utf-8",
    )

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "新文件需要包含标题 PDF Operations Summary，包含客户 Northwind Retail，"
            "包含短语 manual reconciliation，包含日期 August 15, 2026，"
            "并在 Key facts 小节下写出正好三条 bullet。"
        ),
        run_id="pdf_markdown_has_key_facts_bullets",
        target_path=str(markdown_path),
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(markdown_path),
                "operation": "create_file",
                "file_type": "md",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True


def test_file_task_runtime_quality_gate_reads_workspace_relative_markdown_target(
    tmp_path,
):
    project_root = Path(__file__).resolve().parents[2]
    relative_path = (
        Path("_pytest_quality_gate_relative")
        / f"operations_brief_summary_{tmp_path.name}.md"
    )
    markdown_path = project_root / "workspace" / relative_path
    markdown_path.parent.mkdir(parents=True, exist_ok=True)
    markdown_path.write_text(
        "\n".join(
            [
                "# PDF Operations Summary",
                "",
                "**Client:** Northwind Retail",
                "",
                "manual reconciliation",
                "",
                "August 15, 2026",
                "",
                "## Key facts",
                "",
                "- One",
                "- Two",
                "- Three",
            ]
        ),
        encoding="utf-8",
    )

    try:
        runtime = FileTaskRuntime(
            tool_executor=lambda name, args: "",
            model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
        )
        request = FileTaskRequest(
            task=(
                "新文件需要包含标题 PDF Operations Summary，包含客户 Northwind Retail，"
                "包含短语 manual reconciliation，包含日期 August 15, 2026，"
                "并在 Key facts 小节下写出正好三条 bullet。"
            ),
            run_id="pdf_markdown_workspace_relative_target",
            target_path=str(relative_path),
            files=[],
        )

        result = runtime._evaluate_task_quality_gate(
            request,
            [
                {
                    "path": str(relative_path),
                    "operation": "create_file",
                    "file_type": "md",
                }
            ],
            write_intent=True,
            output_mode="write",
        )
    finally:
        markdown_path.unlink(missing_ok=True)

    assert result["passed"] is True


def test_file_task_runtime_quality_gate_does_not_double_count_markdown_preview(
    tmp_path,
):
    markdown_text = "\n".join(
        [
            "# PDF Operations Summary",
            "",
            "**Client:** Northwind Retail",
            "",
            "manual reconciliation",
            "",
            "August 15, 2026",
            "",
            "## Key facts",
            "",
            "- One",
            "- Two",
            "- Three",
        ]
    )
    markdown_path = tmp_path / "operations_brief_summary.md"
    markdown_path.write_text(markdown_text, encoding="utf-8")

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "新文件需要包含标题 PDF Operations Summary，包含客户 Northwind Retail，"
            "包含短语 manual reconciliation，包含日期 August 15, 2026，"
            "并在 Key facts 小节下写出正好三条 bullet。"
        ),
        run_id="pdf_markdown_preview_not_double_counted",
        target_path=str(markdown_path),
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(markdown_path),
                "operation": "create_file",
                "file_type": "md",
                "content": markdown_text,
                "preview": markdown_text,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True


def test_file_task_runtime_quality_gate_rejects_wrong_pdf_markdown_bullet_count(
    tmp_path,
):
    markdown_path = tmp_path / "operations_brief_summary.md"
    markdown_path.write_text(
        "\n".join(
            [
                "# PDF Operations Summary",
                "",
                "**Client:** Northwind Retail",
                "",
                "The brief depends on manual reconciliation.",
                "",
                "Decision date: August 15, 2026",
                "",
                "## Key facts",
                "",
                "- Northwind Retail relies on manual reconciliation.",
                "- The automation pilot decision is due by August 15, 2026.",
            ]
        ),
        encoding="utf-8",
    )

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "新文件需要包含标题 PDF Operations Summary，包含客户 Northwind Retail，"
            "包含短语 manual reconciliation，包含日期 August 15, 2026，"
            "并在 Key facts 小节下写出正好三条 bullet。"
        ),
        run_id="pdf_markdown_wrong_bullet_count",
        target_path=str(markdown_path),
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(markdown_path),
                "operation": "create_file",
                "file_type": "md",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is False
    assert any(
        item["criterion"] == "explicit_bullet_count" and item["passed"] is False
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_xlsx_chart_request_with_chart(tmp_path):
    import openpyxl
    from openpyxl.chart import LineChart, Reference

    workbook_path = tmp_path / "sales_profit_report.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sales"
    rows = [
        ["month", "revenue", "cost", "profit"],
        ["Jan", 12000, 7100, 4900],
        ["Feb", 15800, 8900, 6900],
        ["Mar", 14350, 8200, 6150],
    ]
    for row in rows:
        sheet.append(row)
    chart = LineChart()
    chart.title = "Monthly Profit"
    chart.add_data(Reference(sheet, min_col=4, min_row=1, max_row=4), titles_from_data=True)
    chart.set_categories(Reference(sheet, min_col=1, min_row=2, max_row=4))
    sheet.add_chart(chart, "A7")
    workbook.save(workbook_path)

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="生成一个包含 profit 列和月度 profit 折线图的 Excel 文件",
        run_id="xlsx_chart_gate_accepts_real_chart",
        target_path=str(workbook_path),
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(workbook_path),
                "operation": "write_sheet_data",
                "file_type": "xlsx",
                "cells_written": 7,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "spreadsheet_chart_request_has_workbook_chart"
        and item["passed"] is True
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_spreadsheet_cell_write():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="Update cell B2 in the Excel worksheet with the sales amount.",
        run_id="spreadsheet_cell_write_gate",
        target_path="sales.xlsx",
        files=[
            FileTaskFile(path="sales.xlsx", name="sales.xlsx", type="xlsx", target=True)
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "sales.xlsx",
                "operation": "write_sheet_data",
                "file_type": "xlsx",
                "cells_written": 1,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "spreadsheet_write_has_cells"
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_created_csv_content(tmp_path):
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    csv_path = tmp_path / "restock_plan.csv"
    csv_path.write_text(
        "sku,restock_quantity,priority,lead_time_days\nA100,30,normal,7\n",
        encoding="utf-8",
    )
    request = FileTaskRequest(
        task="Create restock_plan.csv with restock rows.",
        run_id="csv_create_gate",
        target_path=str(csv_path),
        files=[],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": str(csv_path),
                "operation": "run_python_code",
                "file_type": "csv",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "generic_csv_has_written_content"
        and item["passed"] is True
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_text_selection_replace():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="Rewrite the selected text and write it back to this Markdown file.",
        run_id="text_selection_replace_gate",
        target_path="notes.md",
        files=[
            FileTaskFile(path="notes.md", name="notes.md", type="md", target=True)
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "notes.md",
                "operation": "replace_file_selection",
                "file_type": "md",
                "replacements_made": 1,
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "text_selection_replace_has_replacement"
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_workspace_file_copy():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="Copy this PDF file to archive.pdf.",
        run_id="workspace_file_copy_gate",
        target_path="archive.pdf",
        files=[
            FileTaskFile(path="source.pdf", name="source.pdf", type="pdf")
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "archive.pdf",
                "operation": "copy_file",
                "file_type": "pdf",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "workspace_file_copy_uses_copy_tool"
        for item in result["criteria_results"]
    )


def test_file_task_runtime_quality_gate_accepts_cross_file_extract_to_file():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="Extract the action items from notes.pdf into action_items.md.",
        run_id="cross_file_extract_gate",
        target_path="action_items.md",
        files=[
            FileTaskFile(path="notes.pdf", name="notes.pdf", type="pdf")
        ],
    )

    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "action_items.md",
                "operation": "extract_to_file",
                "file_type": "md",
            }
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True
    assert any(
        item["criterion"] == "cross_file_extract_uses_write_tool"
        for item in result["criteria_results"]
    )


def test_file_task_runtime_does_not_treat_docx_organize_write_as_cross_file_extract():
    from app.core.agent.file_task_recipes import semantic_markers

    markers = semantic_markers(
        "整理当前文档并写入 report.docx",
        file_types={"docx"},
        target_file_type="docx",
    )

    assert markers["cross_file_extract_request"] is False


def test_file_task_runtime_emits_typed_event_sequence_with_monotonic_seq():
    def fake_model(**kwargs):
        return {"content": "已总结：alpha beta gamma", "tool_calls": []}

    def fake_executor(tool_name, args):
        assert tool_name == "parse_file_to_text"
        assert args["path"] == "notes.md"
        return "alpha beta gamma"

    request = FileTaskRequest(
        task="总结这个文件",
        run_id="run_demo",
        files=[FileTaskFile(path="notes.md", name="notes.md", type="md")],
    )
    events = list(
        FileTaskRuntime(tool_executor=fake_executor, model_client=fake_model).run(
            request
        )
    )
    event_types = [event.type for event in events]
    run_started = events[0]

    assert event_types[0] == "run.started"
    assert "task.classified" in event_types
    assert "plan.checked" in event_types
    assert "supervisor.status" in event_types
    assert event_types.index("task.classified") < event_types.index("plan.checked")
    assert event_types.index("plan.checked") < event_types.index("plan.created")
    assert event_types.index("plan.created") < event_types.index("step.started")
    assert "tool.started" in event_types
    assert "tool.finished" in event_types
    assert "step.result" in event_types
    assert event_types[-1] == "run.finished"
    assert [event.seq for event in events] == list(range(1, len(events) + 1))
    assert all(event.run_id == "run_demo" for event in events)
    assert run_started.payload["request_kind"] == "new_task"
    assert run_started.payload["task_family"] in {"summarize", "analyze"}
    assert run_started.payload["execution_mode"] == "generic_tool_loop"
    assert run_started.payload["workflow_state"]["task_plan"]["mainline_locked"] is True

    finished = next(event for event in events if event.type == "tool.finished")
    supervisor_events = [event for event in events if event.type == "supervisor.status"]
    step_result_ids = [event.step_id for event in events if event.type == "step.result"]
    execute_result = next(
        event
        for event in events
        if event.type == "step.result" and event.step_id == "execute"
    )
    check_result = next(
        event
        for event in events
        if event.type == "step.result" and event.step_id == "check"
    )

    assert finished.payload["success"] is True
    assert supervisor_events[0].payload["stage"] == "planned"
    assert supervisor_events[-1].payload["task_plan"]["steps"]
    assert supervisor_events[-1].payload["mainline_locked"] is True
    assert "alpha beta" in finished.payload["result_preview"]
    assert "context" in step_result_ids
    assert "execute" in step_result_ids
    assert "check" in step_result_ids
    assert execute_result.payload["summary"] == "已总结：alpha beta gamma"
    assert check_result.payload["passed"] is True


















def test_file_task_runtime_rolls_up_step_results_for_generic_write_tasks():
    def fake_executor(tool_name, args):
        if tool_name == "parse_file_to_text":
            return "第一段\n第二段"
        if tool_name == "write_docx_content":
            return FileTaskToolStreamResult(
                chunks=[
                    FileTaskToolStreamChunk(
                        kind="event",
                        event_type="step_progress",
                        payload={
                            "detail": "已写入 1/2 个段落",
                            "progress": 80,
                            "level": "progress",
                            "file_updated": True,
                            "path": "interview.docx",
                            "file_path": "interview.docx",
                            "supported": True,
                        },
                    ),
                    FileTaskToolStreamChunk(
                        kind="result",
                        payload={
                            "success": True,
                            "path": "interview.docx",
                            "file_path": "interview.docx",
                            "file_type": "docx",
                            "operation": "write_docx_content",
                            "summary": "已将 2 个段落写回 interview.docx。",
                            "preview": "重写了开头两段",
                            "change_type": "modify",
                            "focus": True,
                            "paragraphs_written": 2,
                            "updated_in_place": True,
                        },
                    ),
                ]
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "summary": "已更新 interview.docx。"},
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    def fake_model(**kwargs):
        if any(
            message.get("role") == "function"
            and message.get("name") == "write_docx_content"
            for message in kwargs["messages"]
        ):
            return {"content": "已完成写回", "tool_calls": []}
        return {
            "content": "开始改写并写回",
            "tool_calls": [
                {
                    "id": "rewrite_docx_demo",
                    "name": "write_docx_content",
                    "args": {
                        "path": "interview.docx",
                        "paragraphs": '[{"text":"重写后的开头第一段"},{"text":"重写后的开头第二段"}]',
                    },
                }
            ],
        }

    request = FileTaskRequest(
        task="把访谈文稿的开头改写后写回文档",
        run_id="step_result_write_demo",
        target_path="interview.docx",
        files=[
            FileTaskFile(
                path="interview.docx", name="interview.docx", type="docx", target=True
            )
        ],
    )

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=fake_model, max_rounds=2
        ).run(request)
    )
    execute_results = [
        event
        for event in events
        if event.type == "step.result" and event.step_id == "execute"
    ]
    check_result = next(
        event
        for event in events
        if event.type == "step.result" and event.step_id == "check"
    )

    assert execute_results
    assert execute_results[-1].payload["status"] == "completed"
    assert execute_results[-1].payload["file_change_count"] == 1
    assert execute_results[-1].payload["file_changes"][0]["path"] == "interview.docx"
    assert check_result.payload["passed"] is True
    assert check_result.payload["status"] == "completed"







def test_insert_docx_paragraph_append_request_clears_heading_anchors():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task="只追加一句到目标 DOCX 末尾，保留已有表格不变。",
        target_path="report.docx",
        files=[FileTaskFile(path="report.docx", name="report.docx", type="docx", target=True)],
    )

    repaired = runtime._repair_tool_args_for_context(
        "insert_docx_paragraph",
        {
            "text": "append me",
            "before_heading": "Report Title",
            "after_heading": "Intro",
        },
        request,
        request.files,
    )

    assert repaired["path"] == "report.docx"
    assert "before_heading" not in repaired
    assert "after_heading" not in repaired



def test_file_task_runtime_readonly_negation_overrides_write_word_in_task():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task="分析这个docx，只分析，不写入文件。请总结主要内容、指出访谈问题覆盖的重点和可能缺口。",
        run_id="readonly_negation_classification",
        files=[
            FileTaskFile(
                path="雷鸟访谈问题.docx",
                name="雷鸟访谈问题.docx",
                type="docx",
            )
        ],
    )

    classification = runtime._classify_request(request, request.files)

    assert classification.output_mode == "answer"
    assert classification.write_intent is False
    assert classification.task_family == "analyze"
    assert classification.operation_kind == "read"
    assert "readonly_write_negation" in classification.reason_codes
    assert "write_intent" not in classification.reason_codes


def test_file_task_runtime_frontend_readonly_summary_cannot_be_upgraded_to_write():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    task = (
        "请读取当前附加的 codex_frontend_task_flow_probe.txt，只做只读分析："
        "1）用三条 bullet 总结测试文件内容；"
        "2）指出一个当前工作流风险；"
        "3）说明你没有修改文件。请把任务过程用简洁步骤说明。"
    )
    request = FileTaskRequest(
        task=task,
        run_id="frontend_readonly_summary_probe",
        target_path="workspace/_codex_frontend_task_tests/codex_frontend_task_flow_probe.txt",
        files=[
            FileTaskFile(
                path="workspace/_codex_frontend_task_tests/codex_frontend_task_flow_probe.txt",
                name="codex_frontend_task_flow_probe.txt",
                type="txt",
                target=True,
            )
        ],
        options={"enable_ai_intent_adjudicator": True},
    )

    classification = runtime._classify_request(request, request.files)
    upgraded = runtime._apply_intent_adjudication(
        request,
        request.files,
        classification,
        {
            "status": "ok",
            "intent": "edit_file",
            "confidence": 0.95,
            "should_write": True,
        },
    )
    details = upgraded.public_dict()

    assert upgraded.output_mode == "answer", details
    assert upgraded.write_intent is False, details
    assert upgraded.operation_kind == "read", details
    assert "readonly_write_negation" in upgraded.reason_codes
    assert "ai_intent_adjudicator_readonly_guard" in upgraded.reason_codes
    assert "ai_intent_adjudicator_override" not in upgraded.reason_codes


def test_file_task_runtime_readonly_attached_filename_does_not_become_target():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    task = (
        "请读取当前附加的 codex_frontend_task_flow_probe.txt，只做只读分析："
        "1）用三条 bullet 总结测试文件内容；"
        "2）指出一个当前工作流风险；"
        "3）说明你没有修改文件。"
    )
    request = FileTaskRequest(
        task=task,
        run_id="frontend_readonly_attached_name_not_target",
        files=[
            FileTaskFile(
                path="workspace/_codex_frontend_task_tests/codex_frontend_task_flow_probe.txt",
                name="codex_frontend_task_flow_probe.txt",
                type="txt",
                content="Koto frontend task flow probe",
            )
        ],
    )

    normalized = runtime._request_with_inferred_target_path(request)
    context_files = runtime._context_files(normalized)
    classification = runtime._classify_request(normalized, context_files)
    details = classification.public_dict()

    assert normalized.target_path == ""
    assert len(context_files) == 1
    assert context_files[0].path.endswith("codex_frontend_task_flow_probe.txt")
    assert context_files[0].target is False
    assert classification.operation_kind == "read", details
    assert classification.write_intent is False, details


def test_file_task_runtime_create_new_docx_allows_source_file_protection():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "请读取已添加的 _test_integration_workspace.txt，并创建一个新的 Word 文件 "
            "workspace/koto_ai_assistant_eval_generated.docx。文件中包含标题、原文内容、"
            "以及一句用途判断。请不要修改原文件。"
        ),
        run_id="create_new_docx_with_source_protection",
        files=[
            FileTaskFile(
                path="_test_integration_workspace.txt",
                name="_test_integration_workspace.txt",
                type="txt",
            )
        ],
        target_path="workspace/koto_ai_assistant_eval_generated.docx",
    )

    classification = runtime._classify_request(request, request.files)
    details = classification.public_dict()

    assert classification.output_mode == "write", details
    assert classification.write_intent is True, details
    assert classification.operation_kind != "read", details
    assert "write_intent" in classification.reason_codes, details
    assert "readonly_write_negation" not in classification.reason_codes, details



def test_file_task_runtime_same_file_protection_still_blocks_write():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "请继续优化 workspace/report.docx，保存同一个 DOCX，"
            "但是不要修改 workspace/report.docx。"
        ),
        run_id="same_file_protection_blocks_write",
        target_path="workspace/report.docx",
        files=[
            FileTaskFile(
                path="workspace/report.docx",
                name="report.docx",
                type="docx",
                target=True,
            )
        ],
    )

    classification = runtime._classify_request(request, runtime._context_files(request))
    details = classification.public_dict()

    assert classification.output_mode == "answer", details
    assert classification.write_intent is False, details
    assert "readonly_write_negation" in classification.reason_codes, details


def test_file_task_runtime_blocks_multi_paragraph_write_for_local_docx_append():
    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task=(
            "请继续优化 workspace/report.docx：只追加一句风险声明，"
            "保留已有表格不变，保存同一个 DOCX。"
        ),
        target_path="workspace/report.docx",
    )

    block = runtime._local_docx_edit_block_message(
        request,
        "write_docx_content",
        {
            "path": "workspace/report.docx",
            "paragraphs": json.dumps(
                [{"text": "Risk Review"}, {"text": "Overall risk level: Moderate."}],
                ensure_ascii=False,
            ),
        },
    )

    assert "insert_docx_paragraph" in block


def test_file_task_runtime_english_create_new_docx_allows_source_file_protection():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=lambda **kwargs: {"content": "ok", "tool_calls": []},
    )
    request = FileTaskRequest(
        task=(
            "Read the attached _test_integration_workspace.txt and create a new Word "
            "file at workspace/koto_model_primary_intent_retest.docx. Include the "
            "original content. Do not modify the original txt file."
        ),
        run_id="english_create_new_docx_with_source_protection",
        files=[
            FileTaskFile(
                path="workspace/_test_integration_workspace.txt",
                name="_test_integration_workspace.txt",
                type="txt",
            )
        ],
        target_path="workspace/koto_model_primary_intent_retest.docx",
    )

    classification = runtime._classify_request(request, request.files)
    details = classification.public_dict()

    assert classification.output_mode == "write", details
    assert classification.write_intent is True, details
    assert classification.target_file_type == "docx", details
    assert "write_intent" in classification.reason_codes, details
    assert "readonly_write_negation" not in classification.reason_codes, details


def test_file_task_runtime_blocks_write_tool_when_task_is_readonly():
    responses = [
        {
            "content": "我准备误写入。",
            "tool_calls": [
                {
                    "name": "write_docx_content",
                    "args": {
                        "path": "雷鸟访谈问题_分析报告.docx",
                        "paragraphs": '[{"text":"bad"}]',
                    },
                }
            ],
        },
        {"content": "已改为只给分析答案。", "tool_calls": []},
    ]
    called_tools = []

    def fake_model(**kwargs):
        return responses.pop(0) if responses else {"content": "done", "tool_calls": []}

    def fake_executor(tool_name, args):
        called_tools.append(tool_name)
        if tool_name == "parse_file_to_text":
            return "访谈问题覆盖产品、市场、团队和融资。"
        raise AssertionError(
            f"readonly task should not execute write tool: {tool_name}"
        )

    request = FileTaskRequest(
        task="分析这个docx，只分析，不写入文件。",
        run_id="readonly_blocks_write_tool",
        files=[
            FileTaskFile(
                path="雷鸟访谈问题.docx",
                name="雷鸟访谈问题.docx",
                type="docx",
            )
        ],
    )

    events = list(
        FileTaskRuntime(tool_executor=fake_executor, model_client=fake_model).run(
            request
        )
    )
    blocked = next(
        event
        for event in events
        if event.payload.get("tool_name") == "write_docx_content"
    )
    run_finished = events[-1]

    assert called_tools == ["parse_file_to_text"]
    assert blocked.payload["blocked"] is True
    assert "用户没有授权写入文件" in blocked.payload["result_preview"]
    assert run_finished.payload["summary"] == "已改为只给分析答案。"
    assert run_finished.payload["file_changes"] == []
    assert run_finished.payload["completed_task"] is True


def test_file_task_runtime_blocks_python_file_writes_when_task_is_explicitly_readonly():
    responses = [
        {
            "content": "我准备用 Python 生成文件。",
            "tool_calls": [
                {
                    "name": "run_python_code",
                    "args": {
                        "code": "import matplotlib.pyplot as plt\nplt.plot([1,2,3])\nplt.savefig('bad.png')\nprint('KOTO_CREATED: bad.png')",
                    },
                }
            ],
        },
        {"content": "已改为只给只读分析。", "tool_calls": []},
    ]
    called_tools = []

    def fake_model(**kwargs):
        return responses.pop(0) if responses else {"content": "done", "tool_calls": []}

    def fake_executor(tool_name, args):
        called_tools.append(tool_name)
        if tool_name == "parse_file_to_text":
            return "访谈问题覆盖产品、市场、团队和融资。"
        raise AssertionError(
            f"readonly task should not execute python writer: {tool_name}"
        )

    request = FileTaskRequest(
        task="分析这个docx，只分析，不写入文件。",
        run_id="readonly_blocks_python_writer",
        files=[
            FileTaskFile(
                path="雷鸟访谈问题.docx",
                name="雷鸟访谈问题.docx",
                type="docx",
            )
        ],
    )

    events = list(
        FileTaskRuntime(tool_executor=fake_executor, model_client=fake_model).run(
            request
        )
    )
    blocked = next(
        event for event in events if event.payload.get("tool_name") == "run_python_code"
    )
    run_finished = events[-1]

    assert called_tools == ["parse_file_to_text"]
    assert blocked.payload["blocked"] is True
    assert "已拦截 run_python_code 中的文件写入/保存代码" in blocked.payload["result_preview"]
    assert run_finished.payload["summary"] == "已改为只给只读分析。"
    assert run_finished.payload["file_changes"] == []
    assert run_finished.payload["completed_task"] is True


def test_file_task_runtime_executes_model_planned_write_and_emits_file_change():
    responses = iter(
        [
            {
                "content": "准备写入 Word。",
                "tool_calls": [
                    {
                        "name": "write_docx_content",
                        "args": {
                            "path": "report.docx",
                            "paragraphs": '[{"text":"hello"}]',
                        },
                    }
                ],
            },
            {"content": "已完成。", "tool_calls": []},
        ]
    )

    def fake_model(**kwargs):
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        if tool_name == "write_docx_content":
            return json.dumps(
                {
                    "path": args["path"],
                    "operation": tool_name,
                    "summary": "已写入 1 个段落到 Word 文档",
                    "file_type": "docx",
                    "change_type": "modify",
                    "paragraphs_written": 1,
                    "focus": True,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "confidence": 0.9, "summary": "写入已核验"},
                ensure_ascii=False,
            )
        return ""

    request = FileTaskRequest(task="修改当前文件并保存", run_id="write_demo")
    events = list(
        FileTaskRuntime(tool_executor=fake_executor, model_client=fake_model).run(
            request
        )
    )

    check_finished = next(event for event in events if event.type == "check.finished")
    file_changed = next(event for event in events if event.type == "file.changed")
    run_finished = events[-1]

    assert file_changed.payload["path"] == "report.docx"
    assert check_finished.payload["passed"] is True
    assert check_finished.payload["status"] == "verified"
    assert run_finished.type == "run.finished"
    assert run_finished.payload["completed_task"] is True


def test_file_task_runtime_passes_structured_file_changes_to_checker():
    responses = iter(
        [
            {
                "content": "准备写入 Word。",
                "tool_calls": [
                    {
                        "name": "write_docx_content",
                        "args": {
                            "path": "report.docx",
                            "paragraphs": '[{"text":"hello"}]',
                        },
                    }
                ],
            },
            {"content": "已完成。", "tool_calls": []},
        ]
    )
    captured = {}

    def fake_model(**kwargs):
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        if tool_name == "write_docx_content":
            return json.dumps(
                {
                    "path": args["path"],
                    "operation": tool_name,
                    "summary": "已写入 1 个段落到 Word 文档",
                    "file_type": "docx",
                    "change_type": "modify",
                    "paragraphs_written": 1,
                    "focus": True,
                    "diff": {
                        "kind": "docx_paragraphs",
                        "items": [{"before": "", "after": "hello"}],
                        "changed_count": 1,
                    },
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            captured.update(args)
            return json.dumps(
                {"completed": True, "confidence": 0.9, "summary": "写入已核验"},
                ensure_ascii=False,
            )
        return ""

    request = FileTaskRequest(
        task="修改当前文件并保存",
        run_id="write_structured_check_demo",
        target_path="report.docx",
    )
    events = list(
        FileTaskRuntime(tool_executor=fake_executor, model_client=fake_model).run(
            request
        )
    )
    run_finished = next(event for event in events if event.type == "run.finished")

    assert captured["target_path"] == "report.docx"
    parsed_changes = json.loads(captured["file_changes"])
    assert parsed_changes[0]["path"] == "report.docx"
    assert parsed_changes[0]["operation"] == "write_docx_content"
    assert parsed_changes[0]["paragraphs_written"] == 1
    assert len(run_finished.payload.get("file_changes", [])) == 1


def test_file_task_runtime_ignores_planner_metadata_from_model_response():
    def fake_model(**kwargs):
        request = kwargs["request"]
        assert request.options.get("planner_policy") == "native_only"
        assert not request.options.get("planner_backend")
        return {
            "content": "已总结：alpha beta gamma",
            "tool_calls": [],
            "_planner": {
                "backend": "retired_external",
                "source": "external",
                "policy": "prefer_external",
                "transport": "embedded",
                "reason": "external_system_task",
                "fallback_from": "native",
            },
        }

    request = FileTaskRequest(
        task="访问网页并整理报告",
        run_id="planner_event_demo",
        options={
            "planner_backend": "retired_external",
            "planner_policy": "prefer_external",
        },
    )
    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=fake_model
        ).run(request)
    )

    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = next(event for event in events if event.type == "run.finished")
    expected_runtime = {
        "execution_path": "native",
        "terminal_status": "completed",
        "model_unavailable": False,
        "readonly_fallback_used": False,
        "planner": {
            "backend": "native",
            "source": "native",
            "policy": "native_only",
            "transport": "native",
            "reason": "file_task_native_only",
            "round": 1,
        },
    }

    assert not any(event.type == "planner.selected" for event in events)
    assert not any(event.type == "planner.fallback" for event in events)
    assert {
        key: value
        for key, value in check_finished.payload["runtime"].items()
        if key != "performance"
    } == expected_runtime
    assert check_finished.payload["runtime"]["performance"]["total_ms"] >= 0
    assert run_finished.payload["runtime"] == check_finished.payload["runtime"]


def test_file_task_runtime_emits_model_confirmed_plan_before_tools():
    responses = iter(
        [
            {
                "content": "我会先读取表格，再把表格写入 Word 并核验结果。",
                "tool_calls": [
                    {
                        "name": "read_sheet_data",
                        "args": {
                            "path": "sales.xlsx",
                            "sheet_name": "汇总表",
                            "max_rows": 200,
                        },
                    },
                    {
                        "name": "insert_excel_as_docx_table",
                        "args": {
                            "source_path": "sales.xlsx",
                            "target_path": "report.docx",
                            "sheet_name": "汇总表",
                            "table_title": "销售台账数据",
                        },
                    },
                ],
            },
            {"content": "已完成写入。", "tool_calls": []},
        ]
    )

    def fake_model(**kwargs):
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        if tool_name == "read_sheet_data":
            return json.dumps({"sheet": "汇总表", "row_count": 200}, ensure_ascii=False)
        if tool_name == "insert_excel_as_docx_table":
            return json.dumps(
                {
                    "path": args["target_path"],
                    "source_path": args["source_path"],
                    "summary": "已将工作表“汇总表”的 200 行数据写入 Word 表格",
                    "file_type": "docx",
                    "change_type": "modify",
                    "sheet": "汇总表",
                    "rows_written": 200,
                    "columns_written": 4,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "confidence": 0.95, "summary": "写入已核验"},
                ensure_ascii=False,
            )
        return ""

    events = list(
        FileTaskRuntime(tool_executor=fake_executor, model_client=fake_model).run(
            FileTaskRequest(
                task="将 xlsx 信息加入 docx",
                run_id="confirmed_plan_demo",
                target_path="report.docx",
                files=[
                    FileTaskFile(path="sales.xlsx", name="销售台账.xlsx", type="xlsx"),
                    FileTaskFile(
                        path="report.docx",
                        name="雷鸟访谈问题.docx",
                        type="docx",
                        target=True,
                    ),
                ],
            )
        )
    )

    event_types = [event.type for event in events]
    plan_index = event_types.index("plan.confirmed")
    first_tool_index = next(
        idx
        for idx, event in enumerate(events)
        if event.type == "tool.started" and event.step_id.startswith("tool_")
    )
    confirmed = events[plan_index]

    assert plan_index < first_tool_index
    assert confirmed.step_id == "execute"
    assert confirmed.payload["summary"] == "我会先读取表格，再把表格写入 Word 并核验结果。"
    assert [step["title"] for step in confirmed.payload["steps"]] == [
        "读取 Excel 表格",
        "写入 Word 表格",
        "核验结果",
    ]
    assert "sales.xlsx" in confirmed.payload["steps"][0]["description"]
    assert "report.docx" in confirmed.payload["steps"][1]["description"]




def test_file_task_runtime_explicit_write_beats_advisory_analysis_words():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")

    assert runtime._has_write_intent("分析问题并把结论写入 docx") is True
    assert runtime._has_write_intent("看看哪里需要修改，先不要写回文件") is False


def test_file_task_runtime_quality_gate_rejects_docx_table_task_without_real_table():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task="把 xlsx 表格加入 docx",
        target_path="report.docx",
        files=[
            FileTaskFile(path="sales.xlsx", name="sales.xlsx", type="xlsx"),
            FileTaskFile(
                path="report.docx", name="report.docx", type="docx", target=True
            ),
        ],
    )
    file_changes = [
        {
            "operation": "write_docx_content",
            "path": "report.docx",
            "file_type": "docx",
            "paragraphs_written": 4,
        }
    ]

    check = runtime._verify_task(
        request,
        lambda name, args: json.dumps(
            {"completed": True, "summary": "文件已更新。"}, ensure_ascii=False
        ),
        file_changes,
        write_intent=True,
        output_mode="write",
        model_failed=False,
    )

    assert check["passed"] is False
    assert check["status"] == "quality_gate_failed"
    assert any(
        item["criterion"] == "docx_table_request_has_table"
        for item in check["criteria_results"]
    )


def test_file_task_runtime_quality_gate_rejects_explicit_docx_table_output_without_table():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task=(
            "请读取 workspace/koto_complex_task_test/service_agreement_v1.docx、"
            "workspace/koto_complex_task_test/service_agreement_v2.docx、"
            "workspace/koto_complex_task_test/renewal_budget.xlsx，并在 "
            "workspace/koto_complex_task_test/service_agreement_full_test_20260628.docx "
            "末尾追加核验章节，章节里写三条核验结论，并创建一个真实 Word 表格。"
        ),
        target_path="workspace/koto_complex_task_test/service_agreement_full_test_20260628.docx",
    )
    file_changes = [
        {
            "operation": "write_docx_content",
            "path": request.target_path,
            "file_type": "docx",
            "paragraphs_written": 6,
        }
    ]

    check = runtime._verify_task(
        request,
        lambda name, args: json.dumps(
            {"completed": True, "summary": "文件已更新。"}, ensure_ascii=False
        ),
        file_changes,
        write_intent=True,
        output_mode="write",
        model_failed=False,
    )

    assert check["passed"] is False
    assert check["status"] == "quality_gate_failed"
    assert any(
        item["criterion"] == "docx_table_request_has_table"
        and item["passed"] is False
        for item in check["criteria_results"]
    )


def test_file_task_runtime_repair_message_guides_explicit_docx_table_tool():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task=(
            "请读取 workspace/koto_complex_task_test/renewal_budget.xlsx，并在 "
            "workspace/koto_complex_task_test/service_agreement_full_test_20260628.docx "
            "末尾新增一个真实 Word 表格。"
        ),
        target_path="workspace/koto_complex_task_test/service_agreement_full_test_20260628.docx",
    )

    message = runtime._repair_retry_message(
        request,
        {
            "status": "quality_gate_failed",
            "summary": "缺少真实 Word 表格。",
            "remaining": ["用户要求表格进入 Word；当前表格写入行数：0。"],
        },
        [
            {
                "operation": "write_docx_content",
                "path": request.target_path,
                "paragraphs_written": 6,
            }
        ],
    )

    assert "insert_excel_as_docx_table" in message
    assert "renewal_budget.xlsx" in message
    assert request.target_path in message


def test_file_task_runtime_quality_gate_accepts_verified_docx_narrative_after_python_write(tmp_path):
    from docx import Document

    target_path = tmp_path / "service_agreement_full_test_20260628.docx"
    document = Document()
    document.add_heading("服务协议全面审查报告", level=1)
    document.add_paragraph("执行摘要：本报告已补全关键风险和预算核验。")
    document.add_paragraph("预算分析：预算表与合同金额已完成交叉核验。")
    document.add_paragraph("谈判建议：优先处理责任上限、验收范围和付款期限。")
    document.save(target_path)

    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task="请修复目标 Word 报告，补全预算分析、风险矩阵和谈判建议。",
        target_path=target_path.name,
        files=[
            FileTaskFile(
                path=str(target_path),
                name=target_path.name,
                type="docx",
                target=True,
            )
        ],
    )
    file_changes = [
        {
            "operation": "run_python_code",
            "path": str(target_path),
            "file_type": "docx",
        }
    ]

    check = runtime._evaluate_task_quality_gate(
        request,
        file_changes,
        write_intent=True,
        output_mode="write",
    )

    assert check["passed"] is True
    assert any(
        item["criterion"] == "docx_report_has_narrative"
        and item["passed"] is True
        for item in check["criteria_results"]
    )


def test_file_task_runtime_classifies_ppt_page_write_as_presentation():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task="将docx里的访谈问题总结成3页并加入pptx",
        target_path="deck.pptx",
        files=[
            FileTaskFile(path="interview.docx", name="interview.docx", type="docx"),
            FileTaskFile(path="deck.pptx", name="deck.pptx", type="pptx", target=True),
        ],
    )

    classification = runtime._classify_request(request, request.files)

    assert classification.task_family == "presentation"
    assert classification.operation_kind == "write_slides"
    assert classification.output_mode == "write"
    assert classification.selected_recipe == "ppt_slide_write"
    assert "ppt_slide_write_request" in classification.reason_codes


def test_file_task_runtime_classifies_beautiful_ppt_as_high_quality_design():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task="把这个 PPT 编辑得好看一点，做成专业高级的汇报风格",
        target_path="deck.pptx",
        files=[
            FileTaskFile(path="deck.pptx", name="deck.pptx", type="pptx", target=True)
        ],
    )

    classification = runtime._classify_request(request, request.files)

    assert classification.task_family == "presentation"
    assert classification.operation_kind == "design_slides"
    assert classification.output_mode == "write"
    assert classification.selected_recipe == "pptx_design_edit_high_quality"
    assert "ppt_design_request" in classification.reason_codes
    assert "design_pptx_theme_layout" in classification.matched_capabilities


def test_file_task_runtime_ppt_light_theme_edit_overrides_answer_mode_and_writes():
    tool_calls = []

    def fake_model(**kwargs):
        return {
            "content": "开始应用浅色系主题。",
            "tool_calls": [
                {
                    "name": "design_pptx_theme_layout",
                    "args": {
                        "path": "AI Agent.pptx",
                        "theme": "light_blue_professional",
                        "style_prompt": "浅色系专业商务风格，柔和蓝白配色，保留原内容。",
                    },
                }
            ],
        }

    def fake_executor(tool_name, args):
        tool_calls.append((tool_name, dict(args or {})))
        if tool_name == "parse_file_to_text":
            return "AI Agent PPT：核心概念、价值主张、付费逻辑。"
        if tool_name == "design_pptx_theme_layout":
            return json.dumps(
                {
                    "success": True,
                    "path": args["path"],
                    "operation": "design_pptx_theme_layout",
                    "file_type": "pptx",
                    "summary": "已应用浅色系主题。",
                    "slides_designed": 4,
                    "text_shapes_styled": 12,
                    "theme_name": "light_blue_professional",
                    "layout_strategy": "preserve_content_refresh_theme",
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "summary": "PPT 已更新。"}, ensure_ascii=False
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    request = FileTaskRequest(
        task="我不喜欢这个ppt的风格，换一个浅色系的",
        run_id="ppt_light_theme_answer_mode_regression",
        target_path="AI Agent.pptx",
        files=[
            FileTaskFile(
                path="AI Agent.pptx", name="AI Agent.pptx", type="pptx", target=True
            )
        ],
        options={"output_mode": "answer"},
    )

    events = list(
        FileTaskRuntime(tool_executor=fake_executor, model_client=fake_model).run(
            request
        )
    )
    run_started = events[0]
    file_changed = next(event for event in events if event.type == "file.changed")
    check_finished = next(event for event in events if event.type == "check.finished")

    assert run_started.payload["output_mode"] == "write"
    assert run_started.payload["write_intent"] is True
    assert (
        "answer_mode_overridden_by_write_intent" in run_started.payload["reason_codes"]
    )
    assert any(name == "design_pptx_theme_layout" for name, _args in tool_calls)
    assert file_changed.payload["operation"] == "design_pptx_theme_layout"
    assert check_finished.payload["passed"] is True


def test_file_task_runtime_ai_intent_adjudicator_upgrades_ambiguous_ppt_design_to_write():
    model_calls = []
    tool_calls = []

    def fake_model(**kwargs):
        model_calls.append(kwargs)
        system = str(kwargs.get("system") or "")
        if "任务意图裁判" in system:
            return {
                "content": json.dumps(
                    {
                        "intent": "edit_file",
                        "confidence": 0.88,
                        "should_write": True,
                        "needs_clarification": False,
                        "target_file_type": "pptx",
                        "operation": "redesign_theme",
                        "reason": "用户要求把当前 PPT 风格改得更清爽，应直接修改目标演示文稿。",
                    },
                    ensure_ascii=False,
                )
            }
        return {
            "content": "开始应用清爽浅色主题。",
            "tool_calls": [
                {
                    "name": "design_pptx_theme_layout",
                    "args": {
                        "path": "deck.pptx",
                        "theme": "light_clean",
                        "style_prompt": "清爽浅色系 PPT 风格，保留原内容。",
                    },
                }
            ],
        }

    def fake_executor(tool_name, args):
        tool_calls.append((tool_name, dict(args or {})))
        if tool_name == "parse_file_to_text":
            return "PPT 内容包括 AI Agent 概念、产品价值和商业模式。"
        if tool_name == "design_pptx_theme_layout":
            return json.dumps(
                {
                    "success": True,
                    "path": args["path"],
                    "operation": "design_pptx_theme_layout",
                    "file_type": "pptx",
                    "summary": "已应用清爽浅色主题。",
                    "slides_designed": 5,
                    "text_shapes_styled": 16,
                    "theme_name": "light_clean",
                    "layout_strategy": "preserve_content_refresh_theme",
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "summary": "PPT 已更新。"}, ensure_ascii=False
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    request = FileTaskRequest(
        task="看看这个ppt整体感觉怎么样",
        run_id="ai_intent_adjudicator_ppt_upgrade",
        target_path="deck.pptx",
        files=[
            FileTaskFile(path="deck.pptx", name="deck.pptx", type="pptx", target=True)
        ],
        options={"enable_ai_intent_adjudicator": True},
    )

    events = list(
        FileTaskRuntime(tool_executor=fake_executor, model_client=fake_model).run(
            request
        )
    )
    run_started = events[0]
    task_classified = next(event for event in events if event.type == "task.classified")
    file_changed = next(event for event in events if event.type == "file.changed")

    assert len(model_calls) >= 2
    assert run_started.payload["output_mode"] == "write"
    assert run_started.payload["write_intent"] is True
    assert "ai_intent_adjudicator:edit_file" in run_started.payload["reason_codes"]
    assert "ai_intent_adjudicator_override" in run_started.payload["reason_codes"]
    assert task_classified.payload["intent_adjudication"]["intent"] == "edit_file"
    assert file_changed.payload["operation"] == "design_pptx_theme_layout"


def test_file_task_runtime_ai_intent_adjudicator_does_not_override_explicit_readonly():
    model_calls = []

    def fake_model(**kwargs):
        model_calls.append(kwargs)
        return {
            "content": json.dumps(
                {
                    "intent": "edit_file",
                    "confidence": 0.95,
                    "should_write": True,
                    "needs_clarification": False,
                    "target_file_type": "docx",
                    "operation": "rewrite",
                    "reason": "测试模型试图越权写入。",
                },
                ensure_ascii=False,
            )
        }

    request = FileTaskRequest(
        task="分析这个docx，只分析，不写入文件。",
        run_id="ai_intent_adjudicator_readonly_guard",
        target_path="draft.docx",
        files=[
            FileTaskFile(path="draft.docx", name="draft.docx", type="docx", target=True)
        ],
        options={"enable_ai_intent_adjudicator": True},
    )

    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=fake_model
        ).run(request)
    )
    run_started = events[0]

    assert run_started.payload["output_mode"] == "answer"
    assert run_started.payload["write_intent"] is False
    assert "readonly_write_negation" in run_started.payload["reason_codes"]
    assert not any(
        "ai_intent_adjudicator_override" == code
        for code in run_started.payload["reason_codes"]
    )


def test_file_task_runtime_intent_adjudicator_prompt_keeps_advice_readonly_by_default():
    prompt = FileTaskRuntime()._intent_adjudicator_system_prompt()

    assert "优化论点" in prompt
    assert "有什么风险和机会" in prompt
    assert "默认是 answer_only" in prompt
    assert "不等待确认" in prompt
    assert "不写入文件" in prompt
    assert "只有用户明确说" in prompt
    assert "应用到文件" in prompt
    assert "才使用 analyze_then_confirm" in prompt


def test_file_task_runtime_ai_intent_adjudicator_preserves_create_file_with_source_protection():
    model_calls = []

    def fake_model(**kwargs):
        model_calls.append(kwargs)
        system = str(kwargs.get("system") or "")
        if "任务意图裁判" in system:
            return {
                "content": json.dumps(
                    {
                        "intent": "analyze_then_confirm",
                        "confidence": 0.91,
                        "should_write": False,
                        "needs_clarification": False,
                        "target_file_type": "docx",
                        "operation": "summarize",
                        "reason": "模型误把源文件保护理解成只读分析。",
                    },
                    ensure_ascii=False,
                )
            }
        return {
            "content": "已创建文件。",
            "tool_calls": [
                {
                    "name": "write_docx_content",
                    "args": {
                        "path": "workspace/koto_model_primary_intent_retest.docx",
                        "paragraphs": json.dumps(
                            [{"text": "Koto Model Primary Real Task Test"}],
                            ensure_ascii=False,
                        ),
                    },
                }
            ],
        }

    tool_calls = []

    def fake_executor(tool_name, args):
        tool_calls.append((tool_name, dict(args or {})))
        if tool_name == "parse_file_to_text":
            return "workspace file content"
        if tool_name == "write_docx_content":
            return json.dumps(
                {
                    "success": True,
                    "path": args["path"],
                    "operation": "write_docx_content",
                    "file_type": "docx",
                    "summary": "已创建文件。",
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "summary": "DOCX 已创建。"},
                ensure_ascii=False,
            )
        return ""

    request = FileTaskRequest(
        task=(
            "Read the attached _test_integration_workspace.txt and create a new Word "
            "file at workspace/koto_model_primary_intent_retest.docx. Include the "
            "original content. Do not modify the original txt file."
        ),
        run_id="ai_intent_adjudicator_preserves_create_file_source_guard",
        target_path="workspace/koto_model_primary_intent_retest.docx",
        files=[
            FileTaskFile(
                path="workspace/_test_integration_workspace.txt",
                name="_test_integration_workspace.txt",
                type="txt",
            )
        ],
        options={"enable_ai_intent_adjudicator": True},
    )

    events = list(
        FileTaskRuntime(tool_executor=fake_executor, model_client=fake_model).run(
            request
        )
    )
    run_started = events[0]

    assert run_started.payload["output_mode"] == "write"
    assert run_started.payload["write_intent"] is True
    assert "readonly_write_negation" not in run_started.payload["reason_codes"]
    assert (
        "ai_intent_adjudicator_preserved_explicit_artifact_write"
        in run_started.payload["reason_codes"]
    )
    assert any(name == "write_docx_content" for name, _args in tool_calls)


def test_file_task_runtime_quality_gate_rejects_ppt_beautify_without_design_pass():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task="把这个 PPT 编辑得好看一点，做成专业高级的汇报风格",
        target_path="deck.pptx",
        files=[
            FileTaskFile(path="deck.pptx", name="deck.pptx", type="pptx", target=True)
        ],
    )
    file_changes = [
        {
            "operation": "write_pptx_slides",
            "path": "deck.pptx",
            "file_type": "pptx",
            "slides_updated": 2,
        }
    ]

    check = runtime._verify_task(
        request,
        lambda name, args: json.dumps(
            {"completed": True, "summary": "文件已更新。"}, ensure_ascii=False
        ),
        file_changes,
        write_intent=True,
        output_mode="write",
        model_failed=False,
    )

    assert check["passed"] is False
    assert check["status"] == "quality_gate_failed"
    assert any(
        item["criterion"] == "pptx_design_has_real_design_pass"
        for item in check["criteria_results"]
    )



def test_file_task_runtime_repairs_missing_docx_write_path_for_single_target(tmp_path):
    from docx import Document

    target = tmp_path / "draft.docx"
    doc = Document()
    doc.add_paragraph("原文")
    doc.save(target)

    responses = iter(
        [
            {
                "content": "直接写回润色结果。",
                "tool_calls": [
                    {
                        "name": "write_docx_content",
                        "args": {"paragraphs": "润色后的正文"},
                    }
                ],
            },
            {"content": "已完成。", "tool_calls": []},
        ]
    )

    request = FileTaskRequest(
        task="润色这个 docx 文档的表达，并直接写回文件",
        run_id="repair_missing_docx_write_path",
        target_path=str(target),
        files=[
            FileTaskFile(path=str(target), name="draft.docx", type="docx", target=True)
        ],
    )

    events = list(
        FileTaskRuntime(model_client=lambda **kwargs: next(responses, {"content": "", "tool_calls": []})).run(request)
    )

    file_changed = next(event for event in events if event.type == "file.changed")
    check_finished = next(event for event in events if event.type == "check.finished")
    saved = Document(str(target))
    text = "\n".join(p.text for p in saved.paragraphs)

    assert Path(file_changed.payload["path"]).resolve() == target.resolve()
    assert file_changed.payload["operation"] == "write_docx_content"
    assert "润色后的正文" in text
    assert check_finished.payload["status"] == "verified"


def test_file_task_runtime_accepts_execution_brief_before_tool_calls():
    seen_last_messages = []

    def fake_executor(tool_name, args):
        if tool_name == "parse_file_to_text":
            return "财务预测草稿。"
        if tool_name == "write_docx_content":
            return json.dumps(
                {
                    "success": True,
                    "path": args["path"],
                    "file_type": "docx",
                    "operation": "write_docx_content",
                    "summary": "已写入 2 个段落到 Word 文档",
                    "change_type": "modify",
                    "paragraphs_written": 2,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "summary": "report.docx 已完成更新。"},
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    def fake_model(**kwargs):
        seen_last_messages.append(str(kwargs["messages"][-1]["content"]))
        if any(
            message.get("role") == "function"
            and message.get("name") == "write_docx_content"
            for message in kwargs["messages"]
        ):
            return {"content": "已完成写入。", "tool_calls": []}
        if "已收到 execution_brief" in str(kwargs["messages"][-1]["content"]):
            return {
                "content": "现在开始写入文档。",
                "tool_calls": [
                    {
                        "id": "write_after_brief",
                        "name": "write_docx_content",
                        "args": {
                            "path": "report.docx",
                            "paragraphs": '[{"text":"财务预测整理摘要"},{"text":"收入增长主要来自新品放量。"}]',
                        },
                    }
                ],
            }
        return {
            "content": "",
            "execution_brief": {
                "title": "任务分析",
                "summary": "先归纳财务预测结论，再把摘要写入 report.docx。",
                "steps": [
                    {
                        "title": "整理关键结论",
                        "description": "基于显式上下文提炼财务预测的核心结论",
                    },
                    {
                        "title": "写入目标文档",
                        "description": "把整理后的摘要写回 report.docx",
                    },
                ],
                "planned_tools": ["write_docx_content"],
                "write_targets": ["report.docx"],
                "verification": "检查 report.docx 是否真的更新。",
            },
            "tool_calls": [],
        }

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=fake_model, max_rounds=4
        ).run(
            FileTaskRequest(
                task="整理当前财务预测并写入 report.docx",
                run_id="execution_brief_demo",
                target_path="report.docx",
                files=[
                    FileTaskFile(
                        path="report.docx", name="report.docx", type="docx", target=True
                    )
                ],
            )
        )
    )

    briefed = next(event for event in events if event.type == "plan.briefed")
    confirmed = next(event for event in events if event.type == "plan.confirmed")
    file_changed = next(event for event in events if event.type == "file.changed")
    run_finished = events[-1]

    assert briefed.payload["title"] == "任务分析"
    assert briefed.payload["planned_tools"] == ["write_docx_content"]
    assert confirmed.payload["steps"][0]["tool_name"] == "write_docx_content"
    assert any("已收到 execution_brief" in message for message in seen_last_messages)
    assert file_changed.payload["operation"] == "write_docx_content"
    assert run_finished.payload["completed_task"] is True


def test_file_task_runtime_accepts_execution_brief_as_tool_call_before_write():
    seen_last_messages = []

    def fake_executor(tool_name, args):
        if tool_name == "parse_file_to_text":
            return "财务预测草稿。"
        if tool_name == "write_docx_content":
            return json.dumps(
                {
                    "success": True,
                    "path": args["path"],
                    "file_type": "docx",
                    "operation": "write_docx_content",
                    "summary": "已写入 2 个段落到 Word 文档",
                    "change_type": "modify",
                    "paragraphs_written": 2,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "summary": "report.docx 已完成更新。"},
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    def fake_model(**kwargs):
        seen_last_messages.append(str(kwargs["messages"][-1]["content"]))
        if "已收到 execution_brief" in str(kwargs["messages"][-1]["content"]):
            return {
                "content": "现在开始写入文档。",
                "tool_calls": [
                    {
                        "name": "write_docx_content",
                        "args": {
                            "path": "report.docx",
                            "paragraphs": '[{"text":"财务预测整理摘要"},{"text":"收入增长主要来自新品放量。"}]',
                        },
                    }
                ],
            }
        return {
            "content": "",
            "tool_calls": [
                {
                    "name": "execution_brief",
                    "args": {
                        "summary": "先归纳财务预测结论，再把摘要写入 report.docx。",
                        "steps": [
                            {
                                "title": "整理关键结论",
                                "description": "基于显式上下文提炼核心结论",
                            },
                            {
                                "title": "写入目标文档",
                                "description": "把整理后的摘要写回 report.docx",
                            },
                        ],
                        "planned_tools": ["write_docx_content"],
                        "write_targets": ["report.docx"],
                        "verification": "检查 report.docx 是否真的更新。",
                    },
                }
            ],
        }

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=fake_model, max_rounds=4
        ).run(
            FileTaskRequest(
                task="整理当前财务预测并写入 report.docx",
                run_id="execution_brief_tool_call_demo",
                target_path="report.docx",
                files=[
                    FileTaskFile(
                        path="report.docx", name="report.docx", type="docx", target=True
                    )
                ],
            )
        )
    )

    briefed = next(event for event in events if event.type == "plan.briefed")
    file_changed = next(event for event in events if event.type == "file.changed")
    run_finished = events[-1]

    assert briefed.payload["planned_tools"] == ["write_docx_content"]
    assert any("已收到 execution_brief" in message for message in seen_last_messages)
    assert not any(
        event.type == "tool.finished"
        and event.payload.get("tool_name") == "execution_brief"
        and event.payload.get("success") is False
        for event in events
    )
    assert file_changed.payload["operation"] == "write_docx_content"
    assert run_finished.payload["completed_task"] is True


def test_file_task_runtime_execution_brief_ignores_legacy_delegated_planner_and_stays_native():
    class FakeModelClient:
        def __init__(self):
            self.options_seen = []
            self.tool_gap_seen = False

        def call(self, **kwargs):
            request = kwargs["request"]
            messages = kwargs["messages"]
            self.options_seen.append(dict(request.options or {}))

            if any(
                message.get("role") == "function"
                and message.get("name") == "write_docx_content"
                for message in messages
            ):
                assert request.options.get("planner_policy") == "native_only"
                assert not request.options.get("planner_backend")
                return {
                    "content": "已完成写入。",
                    "tool_calls": [],
                    "_planner": {
                        "backend": "native",
                        "source": "native",
                        "policy": "native_only",
                        "transport": "native",
                        "reason": "file_task_native_only",
                    },
                }

            if "已收到 execution_brief" in str(messages[-1]["content"]):
                assert request.options.get("planner_policy") == "native_only"
                assert not request.options.get("planner_backend")
                return {
                    "content": "继续原生执行写入。",
                    "tool_calls": [
                        {
                            "id": "write_from_retired_external",
                            "name": "write_docx_content",
                            "args": {
                                "path": "report.docx",
                                "paragraphs": '[{"text":"原生白盒执行。"},{"text":"第二段满足质量门。"},{"text":"第三段记录核验依据。"}]',
                            },
                        }
                    ],
                    "_planner": {
                        "backend": "native",
                        "source": "native",
                        "policy": "native_only",
                        "transport": "native",
                        "reason": "file_task_native_only",
                    },
                }

            assert request.options.get("planner_policy") == "native_only"
            assert not request.options.get("planner_backend")
            return {
                "content": "",
                "execution_brief": {
                    "title": "任务分析",
                    "summary": "先完成任务分析，再按白盒骨架继续原生执行。",
                    "delegated_planner": "retired_external",
                    "steps": [
                        {
                            "title": "继续原生执行",
                            "description": "在 Koto 工具骨架内完成写入",
                        }
                    ],
                },
                "tool_calls": [],
            }

    def fake_executor(tool_name, args):
        if tool_name == "parse_file_to_text":
            return "文档草稿。"
        if tool_name == "write_docx_content":
            return json.dumps(
                {
                    "success": True,
                    "path": args["path"],
                    "file_type": "docx",
                    "operation": "write_docx_content",
                    "summary": "已写入 3 个段落到 Word 文档",
                    "change_type": "modify",
                    "paragraphs_written": 3,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "summary": "report.docx 已完成更新。"},
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    model_client = FakeModelClient()
    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=model_client, max_rounds=4
        ).run(
            FileTaskRequest(
                task="整理当前文档并写入 report.docx",
                run_id="execution_brief_delegate_demo",
                target_path="report.docx",
                files=[
                    FileTaskFile(
                        path="report.docx", name="report.docx", type="docx", target=True
                    )
                ],
            )
        )
    )

    briefed = next(event for event in events if event.type == "plan.briefed")
    run_finished = events[-1]

    assert "delegated_planner" not in briefed.payload
    assert all(
        options.get("planner_policy") == "native_only"
        for options in model_client.options_seen
    )
    assert all(
        not options.get("planner_backend") for options in model_client.options_seen
    )
    assert not any(event.type == "planner.selected" for event in events)
    assert not any(event.type == "planner.fallback" for event in events)
    assert run_finished.payload["completed_task"] is True


def test_file_task_runtime_classification_defers_planner_without_explicit_override():
    def fake_model(**kwargs):
        return {"content": "已完成摘要。", "tool_calls": []}

    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=fake_model, max_rounds=2
        ).run(
            FileTaskRequest(
                task="总结当前文件内容",
                run_id="planner_deferred_classification_demo",
                files=[
                    FileTaskFile(
                        path="notes.txt",
                        name="notes.txt",
                        type="txt",
                        content="alpha beta",
                        target=True,
                    )
                ],
            )
        )
    )

    run_started = events[0]

    assert run_started.payload["planner_policy"] == "native_only"
    assert run_started.payload["planner_backend"] == "native"
    assert run_started.payload["planner_reason"] == "file_task_native_only"
    assert "planner_deferred:model_first" not in run_started.payload["reason_codes"]


def test_file_task_runtime_simple_quick_action_mode_skips_classification_and_plan_created():
    def fake_model(**kwargs):
        return {"content": "已总结当前文件重点。", "tool_calls": []}

    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=fake_model, max_rounds=2
        ).run(
            FileTaskRequest(
                task="请总结当前文件内容",
                run_id="quick_action_simple_demo",
                options={"quick_action_mode": "simple"},
                files=[
                    FileTaskFile(
                        path="notes.txt",
                        name="notes.txt",
                        type="txt",
                        content="alpha beta",
                        target=True,
                    )
                ],
            )
        )
    )

    event_types = [event.type for event in events]
    run_started = events[0]
    run_finished = events[-1]

    assert run_started.payload["quick_action_mode"] == "simple"
    assert "task.classified" not in event_types
    assert "plan.created" not in event_types
    assert "plan.checked" in event_types
    plan_checked = next(event for event in events if event.type == "plan.checked")
    assert plan_checked.payload["quick_action_bypass"] is True
    assert plan_checked.payload["passed"] is True
    assert run_finished.payload["completed_task"] is True


def test_file_task_runtime_system_prompt_mentions_execution_brief_protocol():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=lambda **kwargs: {}
    )
    request = FileTaskRequest(task="整理文件并写回目标文档", run_id="execution_brief_prompt_demo")

    system = runtime._build_system_prompt(request, [])

    assert "execution_brief" in system
    assert "execution_plan" in system
    assert "首轮协议" in system
    assert "返回 execution_plan 或 execution_brief 后" in system
    assert "白盒任务骨架" in system


def test_file_task_runtime_prompt_guides_answer_mode_without_writeback():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=lambda **kwargs: {}
    )
    request = FileTaskRequest(
        task="总结这个文档",
        files=[FileTaskFile(path="notes.docx", name="notes.docx", type="docx")],
    )

    messages = runtime._build_messages(request, [], request.files)
    system = runtime._build_system_prompt(request, request.files)
    content = messages[-1]["content"]

    assert "当前任务反馈模式：只给答案" in system
    assert "不要调用写入工具" in system
    assert '"task_feedback_mode"' in content
    assert '"output_mode": "answer"' in content
    assert '"should_write_this_round": false' in content


def test_file_task_runtime_prompt_guides_hybrid_mode_as_analysis_first():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=lambda **kwargs: {}
    )
    request = FileTaskRequest(
        task="分析这个投资建议书，看看有哪些大方向需要修改的地方",
        target_path="雷鸟创新-投资建议书.docx",
        files=[
            FileTaskFile(
                path="雷鸟创新-投资建议书.docx",
                name="雷鸟创新-投资建议书.docx",
                type="docx",
                target=True,
            )
        ],
    )

    messages = runtime._build_messages(request, [], request.files)
    system = runtime._build_system_prompt(request, request.files)
    content = messages[-1]["content"]

    assert "当前任务反馈模式：先分析后决定" in system
    assert "除非用户这轮已经明确要求直接应用到文件，否则不要直接调用写入工具" in system
    assert '"task_feedback_mode"' in content
    assert '"output_mode": "hybrid"' in content
    assert '"should_write_this_round": false' in content


def test_file_task_runtime_prompt_includes_structured_intent_plan_context():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=lambda **kwargs: {}
    )
    request = FileTaskRequest(
        task="分析这个投资建议书，看看有哪些大方向需要修改的地方",
        target_path="雷鸟创新-投资建议书.docx",
        files=[
            FileTaskFile(
                path="雷鸟创新-投资建议书.docx",
                name="雷鸟创新-投资建议书.docx",
                type="docx",
                target=True,
            )
        ],
    )

    messages = runtime._build_messages(request, [], request.files)
    system = runtime._build_system_prompt(request, request.files)
    content = messages[-1]["content"]

    assert "高阶意图规划：" in system
    assert "- 策略：analyze_then_optional_apply" in system
    assert '"intent_plan"' in content
    assert '"recommended_strategy": "analyze_then_optional_apply"' in content
    assert '"requires_confirmation": false' in content
    assert "当前轮会暂停等待确认：否" in system


def test_file_task_runtime_hybrid_requires_confirmation_only_when_explicit():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=lambda **kwargs: {}
    )
    request = FileTaskRequest(
        task="review this investment memo and suggest improvements, wait for my confirmation before applying to the file",
        target_path="雷鸟创新-投资建议书.docx",
        options={"output_mode": "hybrid"},
        files=[
            FileTaskFile(
                path="雷鸟创新-投资建议书.docx",
                name="雷鸟创新-投资建议书.docx",
                type="docx",
                target=True,
            )
        ],
    )

    messages = runtime._build_messages(request, [], request.files)
    system = runtime._build_system_prompt(request, request.files)
    content = messages[-1]["content"]

    assert "- 策略：analyze_then_confirm" in system
    assert '"recommended_strategy": "analyze_then_confirm"' in content
    assert '"requires_confirmation": true' in content
    assert "当前轮会暂停等待确认：是" in system


def test_file_task_runtime_reuses_execution_context_for_prompt_building():
    class StubIntentPlanner:
        def __init__(self):
            self.calls = 0

        def plan(self, request, files, classification, *, known_tool_gap=None):
            self.calls += 1
            return FileTaskIntentPlan(
                intent_type=classification.task_family,
                goal_statement="先分析，再等待确认。",
                output_mode=classification.output_mode,
                confidence=classification.confidence,
                write_intent=classification.write_intent,
                can_apply=False,
                requires_confirmation=False,
                recommended_strategy="analyze_then_optional_apply",
            )

    class StubModelClient:
        def planner_decision_for_request(self, request):
            class Decision:
                policy = "native_only"
                reason = "covered_by_koto_native"
                preferred_backend = ""

            return Decision()

        def call(self, **kwargs):
            return {"content": "ok", "tool_calls": []}

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=StubModelClient()
    )
    runtime._intent_planner = StubIntentPlanner()
    request = FileTaskRequest(
        task="分析这个投资建议书，看看有哪些大方向需要修改的地方",
        target_path="雷鸟创新-投资建议书.docx",
        files=[
            FileTaskFile(
                path="雷鸟创新-投资建议书.docx",
                name="雷鸟创新-投资建议书.docx",
                type="docx",
                target=True,
            )
        ],
    )

    execution_context = runtime._build_execution_context(request, request.files)
    messages = runtime._build_messages(
        request, [], request.files, execution_context=execution_context
    )
    system = runtime._build_system_prompt(
        request, request.files, execution_context=execution_context
    )

    assert runtime._intent_planner.calls == 1
    assert execution_context.effective_planner_policy == "native_only"
    assert execution_context.effective_planner_reason == "file_task_native_only"
    assert '"intent_plan"' in messages[-1]["content"]
    assert "高阶意图规划：" in system


def test_file_task_runtime_run_uses_custom_intent_planner_steps_and_payload():
    class StubIntentPlanner:
        def plan(self, request, files, classification, *, known_tool_gap=None):
            return FileTaskIntentPlan(
                intent_type=classification.task_family,
                goal_statement="先分析风险，再等待确认应用。",
                output_mode=classification.output_mode,
                confidence=0.84,
                write_intent=classification.write_intent,
                can_apply=True,
                requires_confirmation=True,
                recommended_strategy="analyze_then_confirm",
                dynamic_steps=[
                    {
                        "id": "context",
                        "title": "收集上下文",
                        "description": "先锁定目标文档与显式输入。",
                    },
                    {
                        "id": "execute",
                        "title": "生成建议",
                        "description": "先做局部分析，再等待确认。",
                    },
                    {
                        "id": "check",
                        "title": "确认出口",
                        "description": "确认当前轮不直接写回。",
                    },
                ],
                reason_codes=["stub_intent_plan"],
            )

    def fake_model(**kwargs):
        return {"content": "已完成分析。", "tool_calls": []}

    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "",
        model_client=fake_model,
        intent_planner=StubIntentPlanner(),
        max_rounds=1,
    )
    request = FileTaskRequest(
        task="分析这个投资建议书，看看有哪些大方向需要修改的地方",
        run_id="intent_plan_runtime_demo",
        target_path="雷鸟创新-投资建议书.docx",
        files=[
            FileTaskFile(
                path="雷鸟创新-投资建议书.docx",
                name="雷鸟创新-投资建议书.docx",
                type="docx",
                target=True,
            )
        ],
    )

    events = list(runtime.run(request))
    run_started = events[0]
    plan_created = next(event for event in events if event.type == "plan.created")

    assert run_started.payload["intent_plan"]["goal_statement"] == "先分析风险，再等待确认应用。"
    assert (
        run_started.payload["intent_plan"]["recommended_strategy"]
        == "analyze_then_confirm"
    )
    assert plan_created.payload["intent_plan"]["can_apply"] is True
    assert plan_created.payload["steps"][1]["description"] == "先做局部分析，再等待确认。"


def test_file_task_runtime_surfaces_tool_gap_without_retrying_write_guard():
    def fake_model(**kwargs):
        return {
            "content": "当前任务需要新的 Koto 工具。",
            "tool_calls": [],
            "tool_gap": {
                "summary": "当前缺少读取 CAD 文件的 Koto 原生工具。",
                "missing_capability": "read_cad_file",
                "why_missing": "allowlist 中没有可读取 dwg 的工具。",
                "suggested_next_step": "先实现只读 CAD 解析工具，再决定写入方案。",
                "proposed_tool": {
                    "name": "read_cad_file",
                    "description": "解析 DWG/DXF 为可检索的结构化文本。",
                    "parameters": {"type": "object"},
                },
            },
            "_planner": {
                "backend": "retired_external",
                "source": "external",
                "policy": "prefer_external",
                "transport": "embedded",
                "reason": "unsupported_file_types:dwg",
            },
        }

    request = FileTaskRequest(
        task="修改 CAD 文件并导出总结",
        run_id="tool_gap_demo",
        target_path="drawing.dwg",
    )
    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=fake_model
        ).run(request)
    )

    tool_missing = next(event for event in events if event.type == "tool.missing")
    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = next(event for event in events if event.type == "run.finished")
    expected_runtime = {
        "execution_path": "native",
        "terminal_status": "tool_gap",
        "model_unavailable": False,
        "readonly_fallback_used": False,
        "planner": {
            "backend": "native",
            "source": "native",
            "policy": "native_only",
            "transport": "native",
            "reason": "file_task_native_only",
            "round": 1,
        },
    }
    expected_artifact = {
        "artifact_type": "koto_next_action_v1",
        "category": "missing_native_tool",
        "tool_design_protocol": "koto_tool_design_v1",
        "tool_design_status": "draft",
        "generated_by": "koto_file_task_runtime",
        "external_planner_required": False,
        "title": "Koto 下一步：read_cad_file",
        "summary": "当前缺少读取 CAD 文件的 Koto 原生工具。",
        "source_task": "修改 CAD 文件并导出总结",
        "target_path": "drawing.dwg",
        "missing_capability": "read_cad_file",
        "why_missing": "allowlist 中没有可读取 dwg 的工具。",
        "suggested_next_step": "先实现只读 CAD 解析工具，再决定写入方案。",
        "implementation_scope": "smallest_next_capability",
        "acceptance_criteria": [
            "为 read_cad_file 提供稳定的 Koto 原生工具入口",
            "工具返回结构需要可被 file-task 规划器直接消费",
            "补齐能力后可重新执行当前任务而无需改写用户意图",
        ],
        "proposed_tool": {
            "name": "read_cad_file",
            "description": "解析 DWG/DXF 为可检索的结构化文本。",
            "parameters": {"type": "object"},
            "returns": "",
            "rationale": "",
        },
        "runtime_context": expected_runtime,
    }

    assert tool_missing.payload == {
        "summary": "当前缺少读取 CAD 文件的 Koto 原生工具。",
        "missing_capability": "read_cad_file",
        "why_missing": "allowlist 中没有可读取 dwg 的工具。",
        "suggested_next_step": "先实现只读 CAD 解析工具，再决定写入方案。",
        "proposed_tool": {
            "name": "read_cad_file",
            "description": "解析 DWG/DXF 为可检索的结构化文本。",
            "parameters": {"type": "object"},
        },
        "next_action_artifact": expected_artifact,
        "runtime": expected_runtime,
        "round": 1,
    }
    assert not any(
        event.type == "tool.finished"
        and event.payload.get("tool_name") == "write_guard"
        for event in events
    )
    assert check_finished.payload["status"] == "tool_gap"
    assert check_finished.payload["passed"] is False
    assert check_finished.payload["tool_gap"]["missing_capability"] == "read_cad_file"
    assert check_finished.payload["next_action_artifact"] == expected_artifact
    assert {
        key: value
        for key, value in check_finished.payload["runtime"].items()
        if key != "performance"
    } == {**expected_runtime, "terminal_status": "tool_gap"}
    assert check_finished.payload["runtime"]["performance"]["total_ms"] >= 0
    assert run_finished.payload["completed_task"] is False
    assert run_finished.payload["tool_gap"]["missing_capability"] == "read_cad_file"
    assert run_finished.payload["next_action_artifact"] == expected_artifact
    assert run_finished.payload["runtime"] == check_finished.payload["runtime"]


def test_file_task_runtime_does_not_external_fallback_after_tool_gap():
    class FakeModelClient:
        def __init__(self):
            self.options_seen = []
            self.tool_gap_seen = False

        def fallback_planner_backend_for_request(self, request):
            return "retired_external"

        def call(self, **kwargs):
            request = kwargs["request"]
            messages = kwargs["messages"]
            self.options_seen.append(dict(request.options or {}))

            if request.options.get("planner_backend") == "retired_external":
                self.tool_gap_seen = any(
                    isinstance(message.get("tool_gap"), dict)
                    for message in messages
                    if isinstance(message, dict)
                )
                return {
                    "content": "Retired planner 已接管并完成分析。",
                    "tool_calls": [],
                    "_planner": {
                        "backend": "retired_external",
                        "source": "external",
                        "policy": "explicit_backend",
                        "transport": "embedded",
                        "reason": str(
                            request.options.get("planner_runtime_reason") or ""
                        ),
                    },
                }

            assert request.options.get("planner_policy") == "native_only"
            assert not request.options.get("planner_backend")
            return {
                "content": "当前缺少 Koto 原生 CAD 工具。",
                "tool_calls": [],
                "tool_gap": {
                    "summary": "当前缺少读取 CAD 文件的 Koto 原生工具。",
                    "missing_capability": "read_cad_file",
                    "why_missing": "allowlist 中没有可读取 dwg 的工具。",
                    "suggested_next_step": "先实现只读 CAD 解析工具，再决定写入方案。",
                },
                "_planner": {
                    "backend": "native",
                    "source": "native",
                    "policy": "native_only",
                    "transport": "native",
                    "reason": "native_tool_design_required:read_cad_file",
                },
            }

    model_client = FakeModelClient()
    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=model_client, max_rounds=4
        ).run(
            FileTaskRequest(
                task="分析 CAD 文件并整理结论",
                run_id="tool_gap_external_fallback_demo",
                target_path="drawing.dwg",
            )
        )
    )

    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = next(event for event in events if event.type == "run.finished")

    assert model_client.options_seen[0].get("planner_policy") == "native_only"
    assert len(model_client.options_seen) == 1
    assert not any(event.type == "planner.selected" for event in events)
    assert not any(event.type == "planner.fallback" for event in events)
    assert model_client.tool_gap_seen is False
    assert check_finished.payload["status"] == "tool_gap"
    assert check_finished.payload["passed"] is False
    assert run_finished.payload["completed_task"] is False
    assert run_finished.payload["runtime"]["execution_path"] == "native"


def test_file_task_runtime_does_not_external_fallback_after_native_model_failure():
    class FakeModelClient:
        def __init__(self):
            self.options_seen = []

        def fallback_planner_backend_for_request(self, request):
            return "retired_external"

        def call(self, **kwargs):
            request = kwargs["request"]
            self.options_seen.append(dict(request.options or {}))

            if request.options.get("planner_backend") == "retired_external":
                return {
                    "content": "Retired planner 已生成最终摘要。",
                    "tool_calls": [],
                    "_planner": {
                        "backend": "retired_external",
                        "source": "external",
                        "policy": "explicit_backend",
                        "transport": "embedded",
                        "reason": str(
                            request.options.get("planner_runtime_reason") or ""
                        ),
                    },
                }

            raise RuntimeError("native provider offline")

    model_client = FakeModelClient()
    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=model_client, max_rounds=3
        ).run(
            FileTaskRequest(
                task="总结当前文件内容",
                run_id="native_model_failure_external_fallback_demo",
                files=[
                    FileTaskFile(
                        path="notes.txt",
                        name="notes.txt",
                        type="txt",
                        content="alpha beta",
                        target=True,
                    )
                ],
            )
        )
    )

    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = next(event for event in events if event.type == "run.finished")

    assert model_client.options_seen[0]["planner_policy"] == "native_only"
    assert len(model_client.options_seen) == 1
    assert not any(event.type == "planner.selected" for event in events)
    assert not any(event.type == "planner.fallback" for event in events)
    assert check_finished.payload["status"] == "needs_attention"
    assert check_finished.payload["passed"] is False
    assert run_finished.payload["completed_task"] is False
    assert run_finished.payload["runtime"]["execution_path"] == "readonly_fallback"
    assert "Retired planner" not in str(run_finished.payload["summary"])


def test_file_task_runtime_does_not_external_fallback_after_verify_error():
    class FakeModelClient:
        def __init__(self):
            self.options_seen = []

        def fallback_planner_backend_for_request(self, request):
            return "retired_external"

        def call(self, **kwargs):
            request = kwargs["request"]
            messages = kwargs["messages"]
            self.options_seen.append(dict(request.options or {}))

            if request.options.get("planner_backend") == "retired_external":
                if any(
                    message.get("role") == "function"
                    and message.get("name") == "write_docx_content"
                    for message in messages
                ):
                    return {
                        "content": "已完成修复。",
                        "tool_calls": [],
                        "_planner": {
                            "backend": "retired_external",
                            "source": "external",
                            "policy": "explicit_backend",
                            "transport": "embedded",
                            "reason": str(
                                request.options.get("planner_runtime_reason") or ""
                            ),
                        },
                    }
                return {
                    "content": "Retired planner 开始修复写入。",
                    "tool_calls": [
                        {
                            "id": "retired_external_write_docx",
                            "name": "write_docx_content",
                            "args": {
                                "path": "report.docx",
                                "paragraphs": '[{"text":"Retired planner 修正后的内容。"}]',
                            },
                        }
                    ],
                    "_planner": {
                        "backend": "retired_external",
                        "source": "external",
                        "policy": "explicit_backend",
                        "transport": "embedded",
                        "reason": str(
                            request.options.get("planner_runtime_reason") or ""
                        ),
                    },
                }

            if any(
                message.get("role") == "function"
                and message.get("name") == "write_docx_content"
                for message in messages
            ):
                return {
                    "content": "原生写入完成。",
                    "tool_calls": [],
                    "_planner": {
                        "backend": "native",
                        "source": "native",
                        "policy": "native_only",
                        "transport": "native",
                        "reason": "covered_by_koto_native",
                    },
                }

            return {
                "content": "先写入 report.docx。",
                "tool_calls": [
                    {
                        "id": "native_write_docx",
                        "name": "write_docx_content",
                        "args": {
                            "path": "report.docx",
                            "paragraphs": '[{"text":"原生初稿。"}]',
                        },
                    }
                ],
                "_planner": {
                    "backend": "native",
                    "source": "native",
                    "policy": "native_only",
                    "transport": "native",
                    "reason": "covered_by_koto_native",
                },
            }

    verify_calls = {"count": 0}

    def fake_executor(tool_name, args):
        if tool_name == "write_docx_content":
            return json.dumps(
                {
                    "success": True,
                    "path": args["path"],
                    "file_type": "docx",
                    "operation": "write_docx_content",
                    "summary": "已写入 1 个段落到 Word 文档",
                    "change_type": "modify",
                    "paragraphs_written": 1,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            verify_calls["count"] += 1
            if verify_calls["count"] == 1:
                return json.dumps({"error": "judge unavailable"}, ensure_ascii=False)
            return json.dumps(
                {"completed": True, "summary": "report.docx 已完成更新。"},
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    model_client = FakeModelClient()
    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=model_client, max_rounds=5
        ).run(
            FileTaskRequest(
                task="整理当前文档并写入 report.docx",
                run_id="verify_error_external_fallback_demo",
                target_path="report.docx",
                files=[
                    FileTaskFile(
                        path="report.docx", name="report.docx", type="docx", target=True
                    )
                ],
            )
        )
    )

    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = next(event for event in events if event.type == "run.finished")

    assert model_client.options_seen[0].get("planner_policy") == "native_only"
    assert all(
        not options.get("planner_backend") for options in model_client.options_seen
    )
    assert not any(event.type == "planner.selected" for event in events)
    assert not any(event.type == "planner.fallback" for event in events)
    assert verify_calls["count"] == 1
    assert check_finished.payload["status"] == "verify_error"
    assert check_finished.payload["passed"] is False
    assert run_finished.payload["completed_task"] is False


def test_file_task_runtime_pptx_design_retry_points_to_native_tool():
    def fake_model(**kwargs):
        return {
            "content": "很抱歉，我无法直接为 PPTX 文件设计主题和排版。",
            "tool_calls": [],
            "_planner": {
                "backend": "native",
                "source": "native",
                "policy": "native_only",
                "reason": "covered_by_koto_native",
            },
        }

    request = FileTaskRequest(
        task="目前这个 pptx 是没有风格设计的，请帮我设计主题和排版",
        run_id="pptx_design_native_retry",
        files=[
            FileTaskFile(
                path="deck.pptx",
                name="deck.pptx",
                type="pptx",
                content="PPT 文本上下文",
                target=True,
            )
        ],
        target_path="deck.pptx",
    )
    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=fake_model
        ).run(request)
    )

    write_guard = next(
        event
        for event in events
        if event.type == "tool.finished"
        and event.payload.get("tool_name") == "write_guard"
    )
    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = next(event for event in events if event.type == "run.finished")

    assert not any(event.type == "tool.missing" for event in events)
    assert "design_pptx_theme_layout" in write_guard.payload["result_preview"]
    assert check_finished.payload["status"] == "no_file_change"
    assert check_finished.payload["passed"] is False
    assert run_finished.payload["completed_task"] is False


def test_file_task_runtime_parses_native_tool_design_protocol_from_model_text():
    model_payload = {
        "tool_gap": {
            "summary": "需要一个 CAD 读取工具。",
            "missing_capability": "read_cad_file",
            "why_missing": "现有工具不能解析 DWG/DXF 文件。",
            "suggested_next_step": "生成并实现 Koto 原生 CAD 读取工具。",
            "proposed_tool": {
                "name": "read_cad_file",
                "description": "解析 DWG/DXF 为可检索的结构化文本。",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string"},
                        "max_entities": {"type": "integer"},
                    },
                    "required": ["path"],
                },
                "returns": "结构化 CAD 文本摘要。",
                "rationale": "CAD 文件需要格式感知解析。",
                "implementation_notes": ["第一版只读，不写回 CAD。"],
                "safety_constraints": ["不得修改源文件。"],
                "acceptance_tests": ["DWG/DXF 示例文件可以返回图层和实体摘要。"],
            },
        }
    }

    def fake_model(**kwargs):
        assert "koto_tool_design_v1" in kwargs["system"]
        assert "优先组合多个现有工具" in kwargs["system"]
        assert "tool_design_protocol" in kwargs["messages"][-1]["content"]
        return {
            "content": json.dumps(model_payload, ensure_ascii=False),
            "tool_calls": [],
            "_planner": {
                "backend": "native",
                "source": "native",
                "policy": "native_only",
                "reason": "covered_by_koto_native",
            },
        }

    request = FileTaskRequest(
        task="分析这个 CAD 文件",
        run_id="native_tool_design_protocol",
        files=[
            FileTaskFile(
                path="drawing.dwg",
                name="drawing.dwg",
                type="dwg",
                content="CAD 文件上下文",
                target=True,
            )
        ],
        target_path="drawing.dwg",
    )
    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=fake_model
        ).run(request)
    )

    tool_missing = next(event for event in events if event.type == "tool.missing")
    artifact = tool_missing.payload["next_action_artifact"]

    assert tool_missing.payload["missing_capability"] == "read_cad_file"
    assert tool_missing.payload["proposed_tool"]["implementation_notes"] == [
        "第一版只读，不写回 CAD。"
    ]
    assert artifact["tool_design_protocol"] == "koto_tool_design_v1"
    assert artifact["external_planner_required"] is False
    assert "DWG/DXF 示例文件可以返回图层和实体摘要。" in artifact["acceptance_criteria"]
    assert events[-1].payload["completed_task"] is False


def test_file_task_runtime_messages_include_capability_profiles():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=lambda **kwargs: {}
    )
    request = FileTaskRequest(
        task="把表格总结写进文档",
        run_id="capability_context_demo",
        current_file=FileTaskFile(
            path="metrics.xlsx", name="metrics.xlsx", type="xlsx"
        ),
        target_path="summary.docx",
    )

    messages = runtime._build_messages(request, [], [request.current_file])
    content = messages[-1]["content"]

    assert "file_capability_profiles" in content
    assert '"format": "xlsx"' in content
    assert '"format": "docx"' in content
    assert '"write_support": "native"' in content













def test_file_task_runtime_repairs_after_failed_verification(tmp_path):
    target_path = tmp_path / "report.docx"
    target_path.write_text("placeholder", encoding="utf-8")

    responses = iter(
        [
            {
                "content": "先写入第一版。",
                "tool_calls": [
                    {
                        "name": "write_docx_content",
                        "args": {
                            "path": str(target_path),
                            "paragraphs": '[{"text":"draft"}]',
                        },
                    }
                ],
            },
            {"content": "已完成初稿。", "tool_calls": []},
            {
                "content": "根据核验结果修复文档。",
                "tool_calls": [
                    {
                        "name": "write_docx_content",
                        "args": {
                            "path": str(target_path),
                            "paragraphs": '[{"text":"final"}]',
                        },
                    }
                ],
            },
            {"content": "修复完成。", "tool_calls": []},
        ]
    )
    seen_last_messages = []
    verify_calls = []
    write_calls = []

    def fake_model(**kwargs):
        seen_last_messages.append(kwargs["messages"][-1]["content"])
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        if tool_name == "write_docx_content":
            write_calls.append(dict(args))
            return json.dumps(
                {
                    "path": args["path"],
                    "operation": tool_name,
                    "summary": "已写入 Word 文档",
                    "file_type": "docx",
                    "change_type": "modify",
                    "paragraphs_written": 1,
                    "focus": True,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            verify_calls.append(dict(args))
            if len(verify_calls) == 1:
                return json.dumps(
                    {
                        "completed": False,
                        "summary": "正文还没有写到目标位置。",
                        "remaining_steps": ["把正文结论写到目标段落，而不是停留在草稿区"],
                    },
                    ensure_ascii=False,
                )
            return json.dumps(
                {
                    "completed": True,
                    "summary": "修复后核验通过。",
                    "confidence": 0.93,
                },
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    request = FileTaskRequest(
        task="修改当前文件并保存",
        run_id="repair_after_verify_demo",
        target_path=str(target_path),
        files=[
            FileTaskFile(
                path=str(target_path),
                name="report.docx",
                type="docx",
                content="现有 Word 文档",
                target=True,
            )
        ],
    )

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=fake_model, max_rounds=4
        ).run(request)
    )

    repair_guard = next(
        event
        for event in events
        if event.type == "tool.finished"
        and event.payload.get("tool_name") == "repair_guard"
    )
    check_finished_events = [
        event for event in events if event.type == "check.finished"
    ]
    run_finished = events[-1]

    assert len(write_calls) == 2
    assert len(verify_calls) == 2
    assert len(check_finished_events) == 2
    assert check_finished_events[0].payload["status"] == "needs_attention"
    assert check_finished_events[1].payload["status"] == "verified"
    assert "核验未通过" in repair_guard.payload["result_preview"]
    assert any("正文还没有写到目标位置" in message for message in seen_last_messages)
    assert any("把正文结论写到目标段落" in message for message in seen_last_messages)
    assert run_finished.payload["completed_task"] is True
    assert run_finished.payload["summary"] == "修复后核验通过。"


def test_file_task_runtime_uses_extra_repair_round_after_final_round_quality_failure(tmp_path):
    target_path = tmp_path / "sales_profit_report.xlsx"
    source_path = tmp_path / "sales_sample.xlsx"
    source_path.write_bytes(b"fake xlsx")

    responses = iter(
        [
            {
                "content": "先复制源文件建立目标。",
                "tool_calls": [
                    {
                        "name": "copy_file",
                        "args": {
                            "source_path": str(source_path),
                            "target_path": str(target_path),
                        },
                    }
                ],
            },
            {
                "content": "根据质量门补写表格数据。",
                "tool_calls": [
                    {
                        "name": "write_sheet_data",
                        "args": {
                            "path": str(target_path),
                            "sheet_name": "Sales",
                            "rows": [
                                ["month", "revenue", "cost", "profit"],
                                ["2026-01", 120000, 72000, 48000],
                            ],
                        },
                    }
                ],
            },
            {"content": "修复完成。", "tool_calls": []},
        ]
    )
    seen_last_messages = []
    tool_calls = []

    def fake_model(**kwargs):
        seen_last_messages.append(kwargs["messages"][-1]["content"])
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        tool_calls.append((tool_name, dict(args)))
        if tool_name == "copy_file":
            return json.dumps(
                {
                    "path": args["target_path"],
                    "operation": "copy_file",
                    "summary": "已复制文件",
                    "file_type": "xlsx",
                    "change_type": "create",
                },
                ensure_ascii=False,
            )
        if tool_name == "write_sheet_data":
            return json.dumps(
                {
                    "path": args["path"],
                    "operation": "write_sheet_data",
                    "summary": "已写入表格数据",
                    "file_type": "xlsx",
                    "change_type": "modify",
                    "rows_written": 2,
                    "cells_written": 8,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {
                    "completed": True,
                    "summary": "目标表格已核验通过。",
                    "confidence": 0.95,
                },
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    request = FileTaskRequest(
        task="读取 sales_sample.xlsx，生成新的 Excel 文件 sales_profit_report.xlsx，并新增 profit 列。",
        run_id="repair_after_final_round_quality_failure",
        target_path=str(target_path),
        files=[
            FileTaskFile(path=str(source_path), name="sales_sample.xlsx", type="xlsx"),
            FileTaskFile(
                path=str(target_path),
                name="sales_profit_report.xlsx",
                type="xlsx",
                target=True,
            ),
        ],
    )

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor,
            model_client=fake_model,
            max_rounds=1,
        ).run(request)
    )

    check_finished_events = [
        event for event in events if event.type == "check.finished"
    ]
    run_finished = events[-1]

    assert [name for name, _args in tool_calls if name in {"copy_file", "write_sheet_data"}] == [
        "copy_file",
        "write_sheet_data",
    ]
    assert check_finished_events[0].payload["status"] == "quality_gate_failed"
    assert check_finished_events[-1].payload["status"] == "verified"
    assert any("当前核验状态：quality_gate_failed" in message for message in seen_last_messages)
    assert run_finished.payload["completed_task"] is True


def test_file_task_runtime_accepts_python_docx_writeback_for_polish_task():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task="润色这个 docx 文档并写回文件",
        target_path="draft.docx",
        files=[
            FileTaskFile(path="draft.docx", name="draft.docx", type="docx", target=True)
        ],
    )
    file_changes = [
        {
            "path": "draft.docx",
            "file_type": "docx",
            "operation": "run_python_code",
            "summary": "Python 代码更新了 draft.docx",
            "change_type": "modify",
        }
    ]

    check = runtime._verify_task(
        request,
        lambda name, args: json.dumps(
            {"completed": True, "summary": "文件已更新。"}, ensure_ascii=False
        ),
        file_changes,
        write_intent=True,
        output_mode="write",
        model_failed=False,
    )

    assert check["passed"] is True
    assert check["status"] == "verified"


def test_file_task_runtime_preserves_write_blocked_status_in_immediate_verify(tmp_path):
    target_path = tmp_path / "locked.docx"
    target_path.write_text("placeholder", encoding="utf-8")

    responses = iter(
        [
            {
                "content": "尝试写入目标文档。",
                "tool_calls": [
                    {
                        "name": "write_docx_content",
                        "args": {
                            "path": str(target_path),
                            "paragraphs": '[{"text":"draft"}]',
                        },
                    }
                ],
            },
            {"content": "写入受阻，停止继续。", "tool_calls": []},
        ]
    )

    def fake_model(**kwargs):
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        if tool_name == "write_docx_content":
            return json.dumps(
                {
                    "success": False,
                    "status": "write_blocked",
                    "path": args["path"],
                    "operation": tool_name,
                    "summary": "目标文件当前不可写。",
                    "suggested_next_step": "关闭占用目标文件的程序或页签后重试。",
                    "file_type": "docx",
                },
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    request = FileTaskRequest(
        task="修改当前文件并保存",
        run_id="write_blocked_immediate_verify_demo",
        target_path=str(target_path),
        files=[
            FileTaskFile(
                path=str(target_path),
                name="locked.docx",
                type="docx",
                content="现有 Word 文档",
                target=True,
            )
        ],
    )

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=fake_model, max_rounds=2
        ).run(request)
    )

    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = events[-1]

    assert check_finished.payload["status"] == "write_blocked"
    assert check_finished.payload["passed"] is False
    assert check_finished.payload["remaining"] == ["关闭占用目标文件的程序或页签后重试。"]
    assert run_finished.payload["completed_task"] is False
    assert run_finished.payload["runtime"]["terminal_status"] == "write_blocked"


def test_file_task_runtime_packages_failed_python_feedback_for_next_model_turn():
    responses = iter(
        [
            {
                "content": "先运行 Python 脚本。",
                "tool_calls": [
                    {
                        "name": "run_python_code",
                        "args": {"code": "print(missing_name)"},
                    },
                ],
            },
            {"content": "收到错误后停止重复执行。", "tool_calls": []},
        ]
    )
    seen_last_messages = []

    def fake_model(**kwargs):
        seen_last_messages.append(kwargs["messages"][-1]["content"])
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        assert tool_name == "run_python_code"
        return {
            "summary": "Python 脚本执行失败。",
            "stdout": "",
            "stderr": "NameError: name 'missing_name' is not defined",
            "error": "NameError",
            "files": {},
            "_koto_created": [],
            "_koto_modified": [],
        }

    request = FileTaskRequest(
        task="用 Python 分析当前数据", run_id="python_failure_feedback_demo"
    )
    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=fake_model, max_rounds=2
        ).run(request)
    )

    failed_python = next(
        event
        for event in events
        if event.type == "tool.finished"
        and event.payload.get("tool_name") == "run_python_code"
    )
    feedback_payload = json.loads(seen_last_messages[1])

    assert failed_python.payload["success"] is False
    assert feedback_payload["tool_name"] == "run_python_code"
    assert feedback_payload["tool_args"] == {"code": "print(missing_name)"}
    assert feedback_payload["success"] is False
    assert feedback_payload["failure_reason"] == "execution_failed"
    assert feedback_payload["retry_same_call_allowed"] is False
    assert feedback_payload["result"]["error"] == "NameError"
    assert (
        feedback_payload["result"]["stderr"]
        == "NameError: name 'missing_name' is not defined"
    )
    assert "不要重复完全相同的调用" in feedback_payload["next_action"]


def test_file_task_runtime_allows_multiple_python_reads_without_file_markers():
    responses = iter(
        [
            {
                "content": "先读取 Excel。",
                "tool_calls": [
                    {
                        "name": "run_python_code",
                        "args": {"code": "print('first read')"},
                    },
                ],
            },
            {
                "content": "继续读取更多信息。",
                "tool_calls": [
                    {
                        "name": "run_python_code",
                        "args": {"code": "print('second read')"},
                    },
                ],
            },
            {"content": "读取完成。", "tool_calls": []},
        ]
    )
    calls = []

    def fake_model(**kwargs):
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        calls.append((tool_name, dict(args)))
        return "stdout only"

    request = FileTaskRequest(task="分析 Excel 数据", run_id="python_read_demo")
    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=fake_model, max_rounds=3
        ).run(request)
    )

    assert [name for name, _ in calls] == ["run_python_code", "run_python_code"]
    assert not any(
        event.type == "tool.finished" and event.payload.get("skipped")
        for event in events
    )


def test_file_task_runtime_blocks_python_pdf_text_extraction_and_guides_native_read():
    responses = iter(
        [
            {
                "content": "我先用 Python 读 PDF。",
                "tool_calls": [
                    {
                        "name": "run_python_code",
                        "args": {
                            "code": "from PyPDF2 import PdfReader\nreader = PdfReader('source.pdf')\nprint(reader.pages[0].extract_text())",
                        },
                    },
                ],
            },
            {
                "content": "改用原生 PDF 读取。",
                "tool_calls": [
                    {
                        "name": "parse_file_to_text",
                        "args": {
                            "path": "source.pdf",
                            "start_page": 1,
                            "end_page": 3,
                            "max_chars": 4000,
                        },
                    },
                ],
            },
            {"content": "已完成读取。", "tool_calls": []},
        ]
    )
    calls = []

    def fake_model(**kwargs):
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        calls.append((tool_name, dict(args)))
        if tool_name == "parse_file_to_text":
            return "[Page 1]\nThe global rules of art"
        raise AssertionError(f"unexpected tool execution: {tool_name}")

    request = FileTaskRequest(
        task="读取 PDF 原文并对照译稿",
        run_id="pdf_python_guard_demo",
        files=[
            FileTaskFile(path="source.pdf", name="source.pdf", type="pdf"),
            FileTaskFile(
                path="draft.docx", name="draft.docx", type="docx", target=True
            ),
        ],
    )
    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=fake_model, max_rounds=3
        ).run(request)
    )

    blocked = next(
        event
        for event in events
        if event.type == "tool.finished" and event.payload.get("blocked")
    )

    assert all(name != "run_python_code" for name, _ in calls)
    assert any(
        name == "parse_file_to_text"
        and args.get("path") == "source.pdf"
        and args.get("start_page") == 1
        and args.get("end_page") == 3
        and args.get("max_chars") == 4000
        for name, args in calls
    )
    assert blocked.payload["tool_name"] == "run_python_code"
    assert blocked.payload["success"] is False
    assert "不要用 run_python_code 直接读取 PDF 文本" in blocked.payload["result_preview"]
    assert "parse_file_to_text" in blocked.payload["result_preview"]


def test_file_task_runtime_surfaces_python_image_artifacts_in_tool_finished():
    responses = iter(
        [
            {
                "content": "先生成图表。",
                "tool_calls": [
                    {"name": "run_python_code", "args": {"code": "print('ready')"}},
                ],
            },
            {"content": "图表已生成。", "tool_calls": []},
        ]
    )

    def fake_model(**kwargs):
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        assert tool_name == "run_python_code"
        return {
            "summary": "ready\n[1 image(s) generated]",
            "stdout": "ready",
            "stderr": "",
            "error": "",
            "files": {"chart.png": "ZmFrZQ=="},
            "_koto_created": [],
            "_koto_modified": [],
        }

    request = FileTaskRequest(task="基于当前数据生成图表", run_id="python_chart_artifact_demo")
    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor, model_client=fake_model, max_rounds=2
        ).run(request)
    )

    tool_finished = next(
        event
        for event in events
        if event.type == "tool.finished"
        and event.payload.get("tool_name") == "run_python_code"
    )
    code_output = next(event for event in events if event.type == "code.output")

    assert tool_finished.payload["result_preview"] == "ready\n[1 image(s) generated]"
    assert tool_finished.payload["artifacts"] == [
        {
            "kind": "image",
            "name": "chart.png",
            "mime_type": "image/png",
            "data": "ZmFrZQ==",
        }
    ]
    assert code_output.payload["text"] == "ready"


def test_file_task_runtime_keeps_readonly_duplicate_guard_internal():
    repeated_call = {"name": "read_sheet_data", "args": {"path": "sales.xlsx"}}
    responses = iter(
        [
            {"content": "先读取。", "tool_calls": [repeated_call]},
            {"content": "再次读取。", "tool_calls": [repeated_call]},
        ]
    )

    def fake_model(**kwargs):
        return next(responses, {"content": "", "tool_calls": []})

    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: json.dumps(
                {"rows": [["sales", 1]]}, ensure_ascii=False
            ),
            model_client=fake_model,
            max_rounds=2,
        ).run(FileTaskRequest(task="分析 Excel 数据", run_id="duplicate_guard_demo"))
    )
    run_finished = events[-1]

    assert not any(
        event.type == "tool.finished"
        and event.payload.get("tool_name") == "duplicate_guard"
        for event in events
    )
    assert run_finished.payload["completed_task"] is True
    assert "避免重复写入" not in run_finished.payload["summary"]


def test_file_task_runtime_redirects_readonly_duplicate_read_to_final_answer():
    repeated_call = {"name": "parse_file_to_text", "args": {"path": "humanise.docx"}}
    model_calls = []
    responses = iter(
        [
            {"content": "先读取文档。", "tool_calls": [repeated_call]},
            {"content": "我再读取一次。", "tool_calls": [repeated_call]},
            {
                "content": "## 总结与回答\n\n这篇文章讨论艺术、技术与游戏中的操作性身体。",
                "tool_calls": [],
            },
        ]
    )

    def fake_model(**kwargs):
        model_calls.append(kwargs)
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        assert tool_name == "parse_file_to_text"
        return "文章主张：电子游戏通过艺术与技术共同生成一具可被玩家认领的操作性身体。"

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor,
            model_client=fake_model,
            max_rounds=3,
        ).run(
            FileTaskRequest(
                task="分析这个文章",
                run_id="readonly_duplicate_redirect_demo",
                files=[FileTaskFile(path="humanise.docx", name="humanise.docx", type="docx")],
            )
        )
    )

    run_finished = events[-1]

    assert not any(
        event.type == "tool.finished"
        and event.payload.get("tool_name") == "duplicate_guard"
        for event in events
    )
    assert len(model_calls) == 3
    assert model_calls[-1]["tools"] == []
    assert any(
        "不要再次调用任何工具" in str(message.get("content") or "")
        for message in model_calls[-1]["messages"]
    )
    assert "操作性身体" in run_finished.payload["summary"]
    assert "避免重复写入" not in run_finished.payload["summary"]
    assert run_finished.payload["completed_task"] is True


def test_file_task_runtime_disables_tools_after_readonly_answer_guard():
    read_call = {"name": "parse_file_to_text", "args": {"path": "humanise.docx"}}
    model_calls = []
    responses = iter(
        [
            {"content": "先读取文档。", "tool_calls": [read_call]},
            {"content": "", "tool_calls": []},
            {
                "content": "## 总结与回答\n\n这篇文章的核心是操作性身体与认领关系。",
                "tool_calls": [],
            },
        ]
    )

    def fake_model(**kwargs):
        model_calls.append(kwargs)
        return next(responses, {"content": "", "tool_calls": []})

    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "文章主张：操作性身体让玩家认领虚拟身体。",
            model_client=fake_model,
            max_rounds=3,
        ).run(
            FileTaskRequest(
                task="分析这个文章",
                run_id="readonly_answer_guard_answer_only_demo",
                files=[FileTaskFile(path="humanise.docx", name="humanise.docx", type="docx")],
            )
        )
    )

    run_finished = events[-1]
    started_events = [
        event for event in events if event.type == "model.call.started"
    ]

    assert len(model_calls) == 3
    assert model_calls[-1]["tools"] == []
    assert started_events[-1].payload["tool_count"] == 0
    assert started_events[-1].payload["answer_only"] is True
    assert "操作性身体" in run_finished.payload["summary"]
    assert "模型未返回完整自然语言答案" not in run_finished.payload["summary"]
    assert run_finished.payload["completed_task"] is True


def test_file_task_runtime_ignores_tool_calls_in_readonly_answer_only_round():
    read_call = {"name": "parse_file_to_text", "args": {"path": "humanise.docx"}}
    model_calls = []
    executed_tools = []
    responses = iter(
        [
            {"content": "先读取文档。", "tool_calls": [read_call]},
            {"content": "我再读取一次。", "tool_calls": [read_call]},
            {"content": "", "tool_calls": [read_call]},
        ]
    )

    def fake_model(**kwargs):
        model_calls.append(kwargs)
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        executed_tools.append(tool_name)
        return "文章主张：操作性身体让玩家认领虚拟身体。"

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor,
            model_client=fake_model,
            max_rounds=3,
        ).run(
            FileTaskRequest(
                task="分析这个文章",
                run_id="readonly_answer_only_ignores_tools_demo",
                files=[FileTaskFile(path="humanise.docx", name="humanise.docx", type="docx")],
            )
        )
    )

    run_finished = events[-1]
    check_finished = next(event for event in events if event.type == "check.finished")
    last_finished = [
        event for event in events if event.type == "model.call.finished"
    ][-1]

    assert len(model_calls) == 3
    assert model_calls[-1]["tools"] == []
    assert executed_tools == ["parse_file_to_text", "parse_file_to_text"]
    assert last_finished.payload["tool_call_count"] == 0
    assert last_finished.payload["discarded_tool_call_count"] == 1
    assert "操作性身体" in run_finished.payload["summary"]
    assert check_finished.payload["status"] == "needs_attention"
    assert check_finished.payload["passed"] is False
    assert run_finished.payload["completed_task"] is False


def test_file_task_runtime_supervisor_redirects_duplicate_read_before_write():
    model_calls = []

    def fake_model(**kwargs):
        model_calls.append(kwargs)
        raise RuntimeError("model unavailable for fallback test")

    def fake_executor(tool_name, args):
        if tool_name == "parse_file_to_text":
            return (
                "[Page 1] 报告介绍博物馆数字技术应用的年度研究背景，强调数字化保护、智慧服务和数据治理。"
                "当前页窗还列出多个案例方向，包括数字敦煌、知识图谱和藏品档案管理系统，构成后续分析的主题框架。"
                "[Page 2] 综述部分关注以观众为中心的可持续发展、数据驱动的观众分析和数字技术重塑策展逻辑。"
                "案例部分则围绕藏品信息资源管理、沉浸式展示、数字人文系统和公共服务数据体系展开，呈现博物馆数智化从基础数据建设走向场景应用的趋势。"
            )
        if tool_name == "write_docx_content":
            return json.dumps(
                {
                    "success": True,
                    "operation": "write_docx_content",
                    "path": args["path"],
                    "file_type": "docx",
                    "summary": "已写入 4 个段落到 Word 文档",
                    "paragraphs_written": 4,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "summary": "文件已写入。"}, ensure_ascii=False
            )
        raise AssertionError(f"unexpected tool call: {tool_name}")

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor,
            model_client=fake_model,
            max_rounds=3,
        ).run(
            FileTaskRequest(
                task="这是一篇非常长的pdf，请分步总结整篇文章，创建一个docx记录每一步发现，每完成一步等我继续。",
                run_id="supervisor_duplicate_read_demo",
                files=[FileTaskFile(path="source.pdf", name="source.pdf", type="pdf")],
            )
        )
    )

    assert len(model_calls) == 1
    assert not any(
        event.type == "tool.finished"
        and event.payload.get("tool_name") == "supervisor_guard"
        for event in events
    )
    assert not any(
        event.type == "tool.finished"
        and event.payload.get("tool_name") == "duplicate_guard"
        for event in events
    )
    assert any(
        event.type == "file.changed"
        and event.payload.get("operation") == "write_docx_content"
        for event in events
    )
    assert events[-1].payload["completed_task"] is False
    assert events[-1].payload["runtime"]["terminal_status"] == "awaiting_confirmation"



def test_file_task_runtime_treats_add_into_docx_as_write_intent():
    def fake_model(**kwargs):
        return {"content": "当前工具未写入。", "tool_calls": []}

    request = FileTaskRequest(task="将 xlsx 信息加入 docx", run_id="add_docx_demo")
    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=fake_model
        ).run(request)
    )

    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = events[-1]

    assert check_finished.payload["passed"] is False
    assert check_finished.payload["status"] == "no_file_change"
    assert run_finished.payload["completed_task"] is False


def test_file_task_runtime_treats_copy_table_into_target_docx_as_write_intent():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    task = "Copy the workspace spreadsheet table into the target DOCX file."
    request = FileTaskRequest(
        task=task,
        run_id="copy_table_into_target_docx_intent",
        target_path="report.docx",
        files=[
            FileTaskFile(path="sales.xlsx", name="sales.xlsx", type="xlsx"),
            FileTaskFile(
                path="report.docx", name="report.docx", type="docx", target=True
            ),
        ],
    )

    classification = runtime._classify_request(request, request.files)

    assert runtime._has_write_intent(task) is True
    assert classification.write_intent is True
    assert classification.output_mode == "write"
    assert classification.selected_recipe == "xlsx_table_to_docx"


def test_file_task_runtime_treats_file_copy_as_write_intent():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    task = "Copy this PDF file to archive.pdf."
    request = FileTaskRequest(
        task=task,
        run_id="copy_file_intent",
        target_path="archive.pdf",
        files=[
            FileTaskFile(path="source.pdf", name="source.pdf", type="pdf"),
        ],
    )

    classification = runtime._classify_request(request, request.files)

    assert runtime._has_write_intent(task) is True
    assert classification.write_intent is True
    assert classification.selected_recipe == "workspace_file_copy"


def test_file_task_runtime_treats_put_summary_into_new_slides_as_write_intent():
    def fake_model(**kwargs):
        return {"content": "我先读完了内容。", "tool_calls": []}

    request = FileTaskRequest(
        task="将内容总结并放到新的3页里",
        run_id="pptx_write_intent_demo",
        files=[
            FileTaskFile(
                path="AI Agent.pptx", name="AI Agent.pptx", type="pptx", target=True
            )
        ],
    )
    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=fake_model, max_rounds=2
        ).run(request)
    )

    write_guard = next(
        event
        for event in events
        if event.type == "tool.finished"
        and event.payload.get("tool_name") == "write_guard"
    )
    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = events[-1]

    assert "add_pptx_slides" in write_guard.payload["result_preview"]
    assert "read_docx_content" in write_guard.payload["result_preview"]
    assert check_finished.payload["passed"] is False
    assert check_finished.payload["status"] == "no_file_change"
    assert run_finished.payload["completed_task"] is False


def test_file_task_runtime_treats_pptx_page_content_supplement_as_write_intent():
    def fake_model(**kwargs):
        return {"content": "我先读取每一页现有内容。", "tool_calls": []}

    request = FileTaskRequest(
        task="我要你每一页做的内容补充呢？",
        run_id="pptx_page_content_supplement_demo",
        files=[
            FileTaskFile(
                path="AI Agent.pptx", name="AI Agent.pptx", type="pptx", target=True
            )
        ],
    )
    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "", model_client=fake_model, max_rounds=2
        ).run(request)
    )

    run_started = events[0]
    write_guard = next(
        event
        for event in events
        if event.type == "tool.finished"
        and event.payload.get("tool_name") == "write_guard"
    )
    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = events[-1]

    assert run_started.payload["output_mode"] == "write"
    assert run_started.payload["task_family"] == "presentation"
    assert run_started.payload["operation_kind"] == "write_slides"
    assert "write_intent" in run_started.payload["reason_codes"]
    assert "write_pptx_slides" in write_guard.payload["result_preview"]
    assert check_finished.payload["passed"] is False
    assert check_finished.payload["status"] == "no_file_change"
    assert run_finished.payload["completed_task"] is False


def test_file_task_runtime_treats_write_back_text_prompt_as_write_intent():
    def fake_model(**kwargs):
        return {"content": "我先理解目标内容。", "tool_calls": []}

    request = FileTaskRequest(
        task="请把我选中的内容润色后直接写回当前 txt 文件",
        run_id="txt_write_back_demo",
        target_path="copilot_whitebox_demo.txt",
        selection="原始段落",
        selection_source="copilot_whitebox_demo.txt",
        files=[
            FileTaskFile(
                path="copilot_whitebox_demo.txt",
                name="copilot_whitebox_demo.txt",
                type="txt",
                target=True,
            )
        ],
    )
    events = list(
        FileTaskRuntime(
            tool_executor=lambda name, args: "",
            model_client=fake_model,
            max_rounds=2,
        ).run(request)
    )

    run_started = events[0]
    write_guard = next(
        event
        for event in events
        if event.type == "tool.finished"
        and event.payload.get("tool_name") == "write_guard"
    )
    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = events[-1]

    assert run_started.payload["task_family"] == "polish"
    assert run_started.payload["operation_kind"] == "write"
    assert "write_intent" in run_started.payload["reason_codes"]
    assert "run_python_code" in write_guard.payload["result_preview"]
    assert check_finished.payload["passed"] is False
    assert check_finished.payload["status"] == "no_file_change"
    assert run_finished.payload["completed_task"] is False


def test_file_task_runtime_infers_text_target_write_step_without_ppt_mislabel():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task="请把我选中的内容润色后直接写回当前 txt 文件",
        run_id="txt_inferred_step_label",
        target_path="notes.txt",
        files=[
            FileTaskFile(path="notes.txt", name="notes.txt", type="txt", target=True)
        ],
    )

    step = runtime._inferred_write_plan_step(request, request.files)

    assert step["title"] == "写回文本文件"
    assert "notes.txt" in step["description"]
    assert "PPT" not in step["title"]
    assert "PPT" not in step["description"]


def test_file_task_runtime_adds_pptx_slides_from_list_content(tmp_path):
    from pptx import Presentation

    pptx_path = tmp_path / "AI Agent.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "原始页"
    presentation.save(pptx_path)

    responses = iter(
        [
            {
                "content": "我会把总结内容作为 3 页新 PPT 加入文件。",
                "tool_calls": [
                    {
                        "name": "add_pptx_slides",
                        "args": {
                            "path": str(pptx_path),
                            "slides": [
                                {
                                    "title": "总结一",
                                    "content": ["市场需求明确", "替代成本是关键"],
                                },
                                {
                                    "title": "总结二",
                                    "bullets": [
                                        {"text": "本地文件交付"},
                                        {"content": "高质量生成"},
                                    ],
                                },
                                {
                                    "title": "总结三",
                                    "content": {
                                        "points": [
                                            "下一步做规格核验",
                                            "确认客户使用场景",
                                        ]
                                    },
                                },
                            ],
                        },
                    }
                ],
            },
            {"content": "已新增 3 页总结。", "tool_calls": []},
        ]
    )

    def fake_model(**kwargs):
        return next(responses, {"content": "", "tool_calls": []})

    request = FileTaskRequest(
        task="将总结的内容作为3页新ppt加入",
        run_id="pptx_add_slides_demo",
        target_path=str(pptx_path),
        files=[
            FileTaskFile(
                path=str(pptx_path),
                name="AI Agent.pptx",
                type="pptx",
                content="原 PPT 内容",
                target=True,
            )
        ],
    )

    events = list(FileTaskRuntime(model_client=fake_model).run(request))
    add_finished = next(
        event
        for event in events
        if event.type == "tool.finished"
        and event.payload.get("tool_name") == "add_pptx_slides"
    )
    file_changed = next(event for event in events if event.type == "file.changed")
    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = events[-1]

    assert add_finished.payload["success"] is True
    assert Path(file_changed.payload["path"]).resolve() == pptx_path.resolve()
    assert check_finished.payload["status"] == "verified"
    assert run_finished.payload["completed_task"] is True

    saved = Presentation(str(pptx_path))
    all_text = "\n".join(
        shape.text
        for slide in saved.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert len(saved.slides) == 4
    assert "总结一" in all_text
    assert "本地文件交付" in all_text
    assert "确认客户使用场景" in all_text


def test_file_task_runtime_executes_pptx_theme_design_tool(tmp_path):
    from pptx import Presentation

    pptx_path = tmp_path / "AI Agent.pptx"
    presentation = Presentation()
    for title, body in [
        ("原始标题", "第一点\n第二点"),
        ("增长计划", "市场\n产品\n交付"),
    ]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    presentation.save(pptx_path)

    responses = iter(
        [
            {
                "content": "我会直接调用 Koto 原生 PPTX 设计工具做统一主题和安全版式。",
                "tool_calls": [
                    {
                        "name": "design_pptx_theme_layout",
                        "args": {
                            "path": str(pptx_path),
                            "style_brief": "科技感商业 BP",
                            "density": "balanced",
                        },
                    }
                ],
            },
            {"content": "已完成统一主题与版式设计。", "tool_calls": []},
        ]
    )

    def fake_model(**kwargs):
        return next(responses, {"content": "", "tool_calls": []})

    request = FileTaskRequest(
        task="帮这个 pptx 做一套统一视觉风格和排版",
        run_id="pptx_theme_design_demo",
        target_path=str(pptx_path),
        files=[
            FileTaskFile(
                path=str(pptx_path),
                name="AI Agent.pptx",
                type="pptx",
                content="原 PPT 内容",
                target=True,
            )
        ],
    )

    events = list(FileTaskRuntime(model_client=fake_model).run(request))
    design_finished = next(
        event
        for event in events
        if event.type == "tool.finished"
        and event.payload.get("tool_name") == "design_pptx_theme_layout"
    )
    file_changed = next(event for event in events if event.type == "file.changed")
    check_finished = next(event for event in events if event.type == "check.finished")
    run_finished = events[-1]

    assert design_finished.payload["success"] is True
    assert Path(file_changed.payload["path"]).resolve() == pptx_path.resolve()
    assert file_changed.payload["operation"] == "design_pptx_theme_layout"
    assert file_changed.payload["slides_designed"] == 2
    assert file_changed.payload["theme_name"] == "科技深色"
    assert check_finished.payload["status"] == "verified"
    assert run_finished.payload["completed_task"] is True

    saved = Presentation(str(pptx_path))
    all_text = "\n".join(
        shape.text
        for slide in saved.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    background_shapes = [
        shape
        for slide in saved.slides
        for shape in slide.shapes
        if getattr(shape, "name", "") == "KOTO_THEME_BACKGROUND"
    ]
    assert len(saved.slides) == 2
    assert len(background_shapes) == 2
    assert str(background_shapes[0].fill.fore_color.rgb) == "0F172A"
    assert "原始标题" in all_text
    assert "交付" in all_text


def test_design_pptx_theme_layout_keeps_light_theme_against_dark_palette(tmp_path):
    from pptx import Presentation

    from app.core.agent.task_tools import design_pptx_theme_layout

    pptx_path = tmp_path / "AI Agent.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "AI Agent"
    slide.placeholders[1].text = "核心概念\n价值主张"
    presentation.save(pptx_path)

    result = json.loads(
        design_pptx_theme_layout(
            str(pptx_path),
            style_brief="我不喜欢这个ppt的风格，换一个浅色系的",
            palette={"background": "666666", "body_text": "E5E7EB"},
            typography={"font_family": "serif"},
        )
    )

    saved = Presentation(str(pptx_path))
    background_shapes = [
        shape
        for slide in saved.slides
        for shape in slide.shapes
        if getattr(shape, "name", "") == "KOTO_THEME_BACKGROUND"
    ]
    assert result["success"] is True
    assert result["theme_name"] == "清爽简约"
    assert result["font_family"] == "Microsoft YaHei"
    assert len(background_shapes) == 1
    assert str(background_shapes[0].fill.fore_color.rgb) == "F8FAFC"


def test_file_task_runtime_prompt_tells_model_not_to_guess_sheet1():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=lambda **kwargs: {}
    )
    request = FileTaskRequest(
        task="将 xlsx 信息加入 docx",
        run_id="sheet_prompt_demo",
        files=[
            FileTaskFile(path="销售台账.xlsx", name="销售台账.xlsx", type="xlsx"),
            FileTaskFile(
                path="雷鸟访谈问题.docx",
                name="雷鸟访谈问题.docx",
                type="docx",
                target=True,
            ),
        ],
    )

    prompt = runtime._build_system_prompt(request, request.files)

    assert "不要猜 Sheet1" in prompt
    assert "省略 sheet_name" in prompt
    assert "available_sheets" in prompt


def test_file_task_runtime_prompt_routes_pptx_read_and_write_tools_correctly():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=lambda **kwargs: {}
    )
    request = FileTaskRequest(
        task="将内容总结并放到新的3页里",
        run_id="pptx_prompt_demo",
        files=[
            FileTaskFile(
                path="AI Agent.pptx", name="AI Agent.pptx", type="pptx", target=True
            )
        ],
    )

    prompt = runtime._build_system_prompt(request, request.files)

    assert "读取 PPTX 内容优先用 parse_file_to_text" in prompt
    assert "read_docx_content 只用于 DOCX" in prompt
    assert "新增 PPT 总结页时优先用 add_pptx_slides" in prompt


def test_file_task_runtime_prompt_forbids_python_pdf_text_reads_and_requires_windows():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=lambda **kwargs: {}
    )
    request = FileTaskRequest(
        task="根据 PDF 原文润色 docx 译稿",
        run_id="pdf_prompt_demo",
        files=[
            FileTaskFile(path="source.pdf", name="source.pdf", type="pdf"),
            FileTaskFile(
                path="draft.docx", name="draft.docx", type="docx", target=True
            ),
        ],
    )

    prompt = runtime._build_system_prompt(request, request.files)

    assert "读取 PDF 文本时只能使用 parse_file_to_text" in prompt
    assert "start_page/end_page" in prompt
    assert "不要用 run_python_code 调用 PyPDF2" in prompt
    assert "PDF 原文 + DOCX 译稿/润色/审校任务" in prompt


def test_file_task_runtime_has_specific_plan_copy_for_analysis_tools():
    runtime = FileTaskRuntime(
        tool_executor=lambda name, args: "", model_client=lambda **kwargs: {}
    )
    request = FileTaskRequest(task="对比文件并读取片段", run_id="analysis_plan_demo")

    assert runtime._tool_plan_title("compare_files") == "对比文件"
    assert runtime._tool_plan_title("read_file_range") == "读取文本片段"
    assert runtime._tool_plan_title("replace_file_selection") == "替换文本选区"
    assert "比较维度" in runtime._tool_plan_description(
        "compare_files",
        {"file_paths": "a.md, b.md", "aspect": "content"},
        [],
        request,
    )
    assert "第 2 到 5 行" in runtime._tool_plan_description(
        "read_file_range",
        {"path": "notes.md", "start_line": 2, "end_line": 5},
        [],
        request,
    )
    assert "写回" in runtime._tool_plan_description(
        "replace_file_selection",
        {
            "path": "notes.md",
            "original_selection": "旧句",
            "new_content": "新句",
        },
        [],
        request,
    )


def test_file_task_event_serializes_as_sse_json():
    def fake_model(**kwargs):
        return {"content": "收到", "tool_calls": []}

    request = FileTaskRequest(task="读取选区", run_id="sse_demo", selection="hello")
    event = next(
        iter(
            FileTaskRuntime(
                tool_executor=lambda name, args: "", model_client=fake_model
            ).run(request)
        )
    )

    raw = event_to_sse(event)
    assert raw.startswith("data: ")
    assert raw.endswith("\n\n")
    payload = json.loads(raw.removeprefix("data: ").strip())
    assert payload["run_id"] == "sse_demo"
    assert payload["type"] == "run.started"


def test_file_task_tool_catalog_covers_mainstream_office_workflows():
    tool_names = {spec.name for spec in file_task_tool_specs()}
    workflows = supported_file_workflows()

    assert {"docx", "xlsx", "pptx", "pdf", "text", "sandbox"}.issubset(workflows)
    assert {
        "read_docx_content",
        "write_docx_content",
        "fill_docx_template",
        "convert_docx_to_pdf",
        "convert_file",
        "list_conversions",
        "clear_docx_review_marks",
        "insert_image_into_docx",
        "read_sheet_data",
        "inspect_workbook_structure",
        "audit_financial_workbook",
        "write_sheet_data",
        "replace_file_selection",
    }.issubset(tool_names)
    assert {
        "design_pptx_theme_layout",
        "write_pptx_slides",
        "add_pptx_slides",
        "parse_file_to_text",
        "run_python_code",
    }.issubset(tool_names)
    assert any("append images/charts" in item for item in workflows["docx"])
    assert any("fill template placeholders" in item for item in workflows["docx"])
    assert any("convert DOCX to PDF" in item for item in workflows["docx"])
    assert any("clear review comments" in item for item in workflows["docx"])
    assert any("audit financial models" in item for item in workflows["xlsx"])
    assert any("replace exact selections" in item for item in workflows["text"])


def test_file_task_runtime_system_prompt_guides_text_selection_replacement():
    runtime = FileTaskRuntime(tool_executor=lambda name, args: "")
    request = FileTaskRequest(
        task="请把我选中的内容润色得更专业，直接写回当前 md 文件",
        selection="这个方案不错，但是说得不够专业。",
        files=[FileTaskFile(path="notes.md", name="notes.md", type="md", target=True)],
    )
    prompt = runtime._build_system_prompt(request, request.files)

    assert "replace_file_selection" in prompt
    assert "original_selection=用户选区原文" in prompt
    assert "不要为了单个选区改写去 run_python_code 整文件覆写" in prompt


def test_file_task_tool_gateway_is_the_extension_entry_and_filters_allowlist():
    class FakeProvider:
        def __init__(self):
            self.calls = []

        def definitions(self):
            return [
                {
                    "name": "parse_file_to_text",
                    "description": "fake read",
                    "parameters": {"type": "object"},
                },
                {
                    "name": "shell_exec",
                    "description": "must not leak",
                    "parameters": {"type": "object"},
                },
            ]

        def allowed_names(self):
            return {"parse_file_to_text", "shell_exec"}

        def execute(self, tool_name, tool_args):
            self.calls.append((tool_name, tool_args))
            return f"provider:{tool_args.get('path', '')}"

    provider = FakeProvider()
    gateway = FileTaskToolGateway(
        context=FileTaskToolContext(workspace_root="workspace"), providers=[provider]
    )

    assert gateway.allowed_names() == {"parse_file_to_text"}
    assert [definition["name"] for definition in gateway.definitions()] == [
        "parse_file_to_text"
    ]
    assert (
        gateway.execute("parse_file_to_text", {"path": "notes.md"})
        == "provider:notes.md"
    )
    assert provider.calls == [("parse_file_to_text", {"path": "notes.md"})]
    with pytest.raises(ValueError):
        gateway.execute("shell_exec", {})


def test_file_task_tool_gateway_filters_tools_by_task_file_type_context():
    class FakeProvider:
        def definitions(self):
            return [
                {
                    "name": "read_docx_content",
                    "description": "docx only",
                    "parameters": {"type": "object"},
                },
                {
                    "name": "parse_file_to_text",
                    "description": "generic read",
                    "parameters": {"type": "object"},
                },
                {
                    "name": "add_pptx_slides",
                    "description": "pptx write",
                    "parameters": {"type": "object"},
                },
                {
                    "name": "run_python_code",
                    "description": "sandbox",
                    "parameters": {"type": "object"},
                },
            ]

        def allowed_names(self):
            return {
                "read_docx_content",
                "parse_file_to_text",
                "add_pptx_slides",
                "run_python_code",
            }

        def execute(self, tool_name, tool_args):
            return "ok"

    gateway = FileTaskToolGateway(
        context=FileTaskToolContext(
            task_files=[{"path": "AI Agent.pptx", "type": "pptx"}]
        ),
        providers=[FakeProvider()],
    )

    assert gateway.allowed_names() == {
        "parse_file_to_text",
        "add_pptx_slides",
        "run_python_code",
    }
    assert [definition["name"] for definition in gateway.definitions()] == [
        "parse_file_to_text",
        "add_pptx_slides",
        "run_python_code",
    ]


def test_file_task_tool_gateway_infers_docx_output_tools_from_task_context():
    class FakeProvider:
        def definitions(self):
            return [
                {
                    "name": "parse_file_to_text",
                    "description": "pdf read",
                    "parameters": {"type": "object"},
                },
                {
                    "name": "write_docx_content",
                    "description": "docx write",
                    "parameters": {"type": "object"},
                },
                {
                    "name": "add_pptx_slides",
                    "description": "pptx write",
                    "parameters": {"type": "object"},
                },
                {
                    "name": "run_python_code",
                    "description": "sandbox",
                    "parameters": {"type": "object"},
                },
            ]

        def allowed_names(self):
            return {
                "parse_file_to_text",
                "write_docx_content",
                "add_pptx_slides",
                "run_python_code",
            }

        def execute(self, tool_name, tool_args):
            return "ok"

    gateway = FileTaskToolGateway(
        context=FileTaskToolContext(
            task_files=[{"path": "museum.pdf", "type": "pdf"}],
            request_context={
                "task": "总结这个 PDF 并创建一个 docx 记录要点",
                "target_path": "",
            },
        ),
        providers=[FakeProvider()],
    )

    assert gateway.allowed_names() == {
        "parse_file_to_text",
        "write_docx_content",
        "run_python_code",
    }
    assert "add_pptx_slides" not in gateway.allowed_names()


def test_file_task_runtime_uses_injected_tool_provider_boundary():
    model_tools = []

    class FakeProvider:
        def __init__(self):
            self.calls = []

        def definitions(self):
            return [
                {
                    "name": "parse_file_to_text",
                    "description": "fake read",
                    "parameters": {"type": "object"},
                }
            ]

        def allowed_names(self):
            return {"parse_file_to_text"}

        def execute(self, tool_name, tool_args):
            self.calls.append((tool_name, tool_args))
            return "provider context"

    def fake_model(**kwargs):
        model_tools.extend(kwargs["tools"])
        return {"content": "已读取 provider context", "tool_calls": []}

    provider = FakeProvider()
    request = FileTaskRequest(
        task="总结这个文件",
        run_id="provider_demo",
        files=[FileTaskFile(path="notes.md", name="notes.md", type="md")],
    )
    events = list(
        FileTaskRuntime(tool_provider=provider, model_client=fake_model).run(request)
    )

    assert provider.calls == [
        ("parse_file_to_text", {"path": "notes.md", "max_chars": 12000})
    ]
    assert [tool["name"] for tool in model_tools] == ["parse_file_to_text"]
    assert events[-1].type == "run.finished"


def test_file_task_model_client_routes_local_and_cloud():
    calls = []

    class FakeClient(FileTaskModelClient):
        def _call_cloud(self, **kwargs):
            calls.append("cloud")
            return {"content": "cloud ok", "tool_calls": []}

        def _call_local(self, **kwargs):
            calls.append("local")
            return {"content": "local ok", "tool_calls": []}

    client = FakeClient()
    client.call(
        request=FileTaskRequest(task="t", model_mode="cloud"),
        messages=[],
        system="",
        tools=[],
    )
    client.call(
        request=FileTaskRequest(task="t", model_mode="local"),
        messages=[],
        system="",
        tools=[],
    )

    assert calls == ["cloud", "local"]


def test_file_task_model_client_routes_deepseek_cloud_provider(monkeypatch):
    from app.core.agent import file_task_model as file_task_model_module
    import app.core.llm.model_fallback as fallback_module
    import app.core.llm.provider_factory as provider_factory

    captured = {}

    class FakeProvider:
        pass

    class FakeFallbackExecutor:
        def generate_with_fallback(self, **kwargs):
            captured["fallback"] = kwargs
            return {"content": "deepseek ok", "tool_calls": []}

    def fake_get_llm_provider(**kwargs):
        captured["provider_kwargs"] = kwargs
        return FakeProvider()

    monkeypatch.setattr(
        file_task_model_module,
        "get_configured_cloud_model",
        lambda **kwargs: "deepseek-v4-pro",
    )
    monkeypatch.setattr(provider_factory, "get_llm_provider", fake_get_llm_provider)
    monkeypatch.setattr(
        fallback_module, "get_fallback_executor", lambda: FakeFallbackExecutor()
    )

    client = FileTaskModelClient()
    response = client.call(
        request=FileTaskRequest(task="t", model_mode="deepseek", model_id="deepseek"),
        messages=[{"role": "user", "content": "hi"}],
        system="system",
        tools=[{"name": "parse_file_to_text"}],
    )

    assert response["content"] == "deepseek ok"
    assert captured["provider_kwargs"] == {
        "provider": "deepseek",
        "model": "deepseek-v4-pro",
    }
    assert captured["fallback"]["preferred_model"] == "deepseek-v4-pro"
    assert captured["fallback"]["task_type"] == "FILE_TASK"
    assert captured["fallback"]["system_instruction"] == "system"
    assert captured["fallback"]["tools"] == [{"name": "parse_file_to_text"}]


def test_file_task_model_client_passes_file_task_timeout_to_local_provider(monkeypatch):
    from app.core.agent import file_task_model as file_task_model_module
    import app.core.llm.ollama_llm_provider as ollama_module

    captured = {}

    class FakeOllamaProvider:
        def __init__(self, model=None):
            captured["model"] = model

        def generate_content(self, **kwargs):
            captured.update(kwargs)
            return {"content": "local ok", "tool_calls": []}

    monkeypatch.setattr(ollama_module, "OllamaLLMProvider", FakeOllamaProvider)

    client = FileTaskModelClient()
    monkeypatch.setattr(client, "_is_local_available", lambda: True)

    response = client.call(
        request=FileTaskRequest(task="t", model_mode="local"),
        messages=[{"role": "user", "content": "hi"}],
        system="system",
        tools=[{"name": "parse_file_to_text"}],
    )

    assert response["content"] == "local ok"
    assert captured["model"] is None
    assert (
        captured["call_timeout"] == file_task_model_module._FILE_TASK_LLM_CALL_TIMEOUT
    )
    assert captured["system_instruction"] == "system"
    assert captured["tools"] == [{"name": "parse_file_to_text"}]
    assert captured["stream"] is False
    assert captured["temperature"] == 0.2


def test_file_task_model_client_prefers_file_task_model_route(monkeypatch):
    from web import runtime_context

    monkeypatch.setattr(
        runtime_context,
        "get_model_map",
        lambda: {"CHAT": "chat-model", "FILE_TASK": "file-task-model"},
    )

    client = FileTaskModelClient()

    assert (
        client._cloud_model_id(FileTaskRequest(task="整理文件", model_mode="gemini"))
        == "file-task-model"
    )





def test_file_task_runtime_accepts_short_summary_plus_real_table_for_docx_report():
    runtime = FileTaskRuntime()
    request = FileTaskRequest(
        task="Create a DOCX report from this Excel table. Keep the real table in Word and write a short summary before the table.",
        target_path="report.docx",
        files=[
            FileTaskFile(path="sales.xlsx", name="sales.xlsx", type="xlsx"),
            FileTaskFile(
                path="report.docx", name="report.docx", type="docx", target=True
            ),
        ],
    )
    result = runtime._evaluate_task_quality_gate(
        request,
        [
            {
                "path": "report.docx",
                "file_type": "docx",
                "operation": "write_docx_content",
                "paragraphs_written": 2,
            },
            {
                "path": "report.docx",
                "file_type": "docx",
                "operation": "insert_excel_as_docx_table",
                "rows_written": 3,
            },
        ],
        write_intent=True,
        output_mode="write",
    )

    assert result["passed"] is True



def test_file_task_runtime_infers_new_artifact_target_without_marking_source_target(tmp_path):
    source_path = tmp_path / "workspace" / "_test_integration_workspace.txt"
    source_path.parent.mkdir()
    source_path.write_text("workspace file content", encoding="utf-8")
    runtime = FileTaskRuntime(workspace_root=str(tmp_path))
    request = FileTaskRequest(
        task=(
            "请读取已添加的 _test_integration_workspace.txt，并创建一个新的 Word 文件 "
            "workspace/koto_ai_assistant_eval_generated.docx。请不要修改原文件。"
        ),
        files=[
            FileTaskFile(
                path=str(source_path),
                name="_test_integration_workspace.txt",
                type="txt",
                content="workspace file content",
            )
        ],
    )

    normalized = runtime._request_with_inferred_target_path(request)
    context_files = runtime._context_files(normalized)

    assert normalized.target_path == "workspace/koto_ai_assistant_eval_generated.docx"
    assert any(
        file_info.path == "workspace/koto_ai_assistant_eval_generated.docx"
        and file_info.target
        for file_info in context_files
    )
    source_context = next(
        file_info
        for file_info in context_files
        if file_info.name == "_test_integration_workspace.txt"
    )
    assert source_context.target is False

    ui_request = FileTaskRequest(
        task=request.task,
        target_path=str(source_path),
        files=[
            FileTaskFile(
                path=str(source_path),
                name="_test_integration_workspace.txt",
                type="txt",
                content="workspace file content",
                target=True,
            )
        ],
    )
    normalized_ui_request = runtime._request_with_inferred_target_path(ui_request)
    assert normalized_ui_request.target_path == "workspace/koto_ai_assistant_eval_generated.docx"


def test_file_task_runtime_explicit_output_overrides_open_source_target(tmp_path):
    source_path = tmp_path / "workspace" / "_codex_frontend_task_tests" / "koto_task_smoke.txt"
    source_path.parent.mkdir(parents=True)
    source_path.write_text("Koto task smoke fixture", encoding="utf-8")
    output_path = "workspace/_codex_frontend_task_tests/koto_complex_task_report_20260617_1345.md"
    runtime = FileTaskRuntime(workspace_root=str(tmp_path))
    request = FileTaskRequest(
        task=(
            "请基于当前打开的 koto_task_smoke.txt 完成一个复杂文件任务，"
            f"将结果保存为 Markdown 文件：{output_path}"
        ),
        target_path=str(source_path),
        files=[
            FileTaskFile(
                path=str(source_path),
                name="koto_task_smoke.txt",
                type="txt",
                content="Koto task smoke fixture",
                target=True,
            )
        ],
    )

    normalized = runtime._request_with_inferred_target_path(request)
    context_files = runtime._context_files(normalized)
    summary = runtime._plan_summary(normalized, context_files, write_intent=True)
    verify_target = runtime._verification_target_path(
        normalized,
        [
            {
                "path": output_path,
                "operation": "write_markdown",
            }
        ],
    )

    assert normalized.target_path == output_path
    assert summary == "准备生成 koto_complex_task_report_20260617_1345.md。"
    assert verify_target == output_path
    assert any(file_info.path == output_path and file_info.target for file_info in context_files)
    source_context = next(file_info for file_info in context_files if file_info.name == "koto_task_smoke.txt")
    assert source_context.target is False


def test_file_task_runtime_infers_split_directory_and_filename_output_target(tmp_path):
    workspace = tmp_path / "workspace"
    current_dir = workspace / "codex_real_task_20260701"
    old_dir = workspace / "codex_real_task_20260630"
    current_dir.mkdir(parents=True)
    old_dir.mkdir(parents=True)
    source_path = current_dir / "sales_sample.xlsx"
    old_output_path = old_dir / "sales_profit_report.xlsx"
    source_path.write_bytes(b"fake xlsx")
    old_output_path.write_bytes(b"old xlsx")

    prompt = (
        "请读取当前打开的 sales_sample.xlsx，生成一个新的 Excel 文件，"
        "文件名为 sales_profit_report.xlsx，保存在 codex_real_task_20260701 目录下。"
        "新文件需要包含原始 month/revenue/cost 数据、新增 profit=revenue-cost 列，"
        "并插入一个月度 profit 折线图。请执行文件任务并在完成后汇报保存路径，"
        "不要只给建议。"
    )
    runtime = FileTaskRuntime(workspace_root=str(workspace))
    request = FileTaskRequest(
        task=prompt,
        run_id="split_directory_filename_output_target",
        target_path="sales_profit_report.xlsx",
        current_file=FileTaskFile(
            path="codex_real_task_20260701/sales_sample.xlsx",
            name="sales_sample.xlsx",
            type="xlsx",
        ),
    )

    normalized = runtime._request_with_inferred_target_path(request)
    context_files = runtime._context_files(normalized)

    assert normalized.target_path == "codex_real_task_20260701/sales_profit_report.xlsx"
    assert any(
        file_info.path == "codex_real_task_20260701/sales_profit_report.xlsx"
        and file_info.target
        for file_info in context_files
    )
    assert any(
        Path(file_info.path).name == "sales_sample.xlsx" and not file_info.target
        for file_info in context_files
    )
    assert all(
        Path(file_info.path).resolve() != old_output_path.resolve()
        for file_info in context_files
        if file_info.path
    )


def test_file_task_runtime_blocks_model_write_to_protected_source_when_creating_artifact(tmp_path):
    source_path = tmp_path / "workspace" / "_test_integration_workspace.txt"
    source_path.parent.mkdir()
    source_path.write_text("workspace file content", encoding="utf-8")
    target_path = "workspace/koto_ai_assistant_eval_generated.docx"
    responses = iter(
        [
            {
                "content": "读取源文件并创建 Word。",
                "tool_calls": [
                    {
                        "name": "parse_file_to_text",
                        "args": {"path": str(source_path), "max_chars": 12000},
                    },
                    {
                        "name": "write_docx_content",
                        "args": {
                            "path": target_path,
                            "paragraphs": '[{"text":"Koto AI 助手写入测试","style":"Heading 1"},{"text":"workspace file content"}]',
                        },
                    },
                    {
                        "name": "write_docx_content",
                        "args": {
                            "path": str(source_path),
                            "paragraphs": '[{"text":"should not touch source"}]',
                        },
                    },
                ],
            },
            {"content": "已完成。", "tool_calls": []},
        ]
    )
    called_tools = []

    def fake_model(**kwargs):
        assert kwargs["request"].target_path == target_path
        return next(responses, {"content": "", "tool_calls": []})

    def fake_executor(tool_name, args):
        called_tools.append((tool_name, dict(args)))
        if tool_name == "parse_file_to_text":
            return "workspace file content"
        if tool_name == "write_docx_content":
            assert args["path"] == target_path
            return json.dumps(
                {
                    "path": args["path"],
                    "operation": "write_docx_content",
                    "summary": "已写入 2 个段落到 Word 文档",
                    "file_type": "docx",
                    "change_type": "create",
                    "paragraphs_written": 2,
                    "focus": True,
                },
                ensure_ascii=False,
            )
        if tool_name == "verify_task_completion":
            return json.dumps(
                {"completed": True, "confidence": 0.95, "summary": "写入已核验"},
                ensure_ascii=False,
            )
        raise AssertionError(f"unexpected tool: {tool_name}")

    request = FileTaskRequest(
        task=(
            "请读取已添加的 _test_integration_workspace.txt，并创建一个新的 Word 文件 "
            f"{target_path}。文件中包含标题、原文内容、以及用途判断。请不要修改原文件。"
        ),
        run_id="protect_source_create_artifact",
        target_path=str(source_path),
        files=[
            FileTaskFile(
                path=str(source_path),
                name="_test_integration_workspace.txt",
                type="txt",
                content="workspace file content",
                target=True,
            )
        ],
    )

    events = list(
        FileTaskRuntime(
            tool_executor=fake_executor,
            model_client=fake_model,
            workspace_root=str(tmp_path),
            max_rounds=2,
        ).run(request)
    )

    blocked = next(
        event
        for event in events
        if event.type == "tool.finished"
        and event.payload.get("blocked")
        and event.payload.get("tool_name") == "write_docx_content"
    )
    changed = [event.payload for event in events if event.type == "file.changed"]
    run_started = next(event for event in events if event.type == "run.started")

    assert "不能写入 _test_integration_workspace.txt" in blocked.payload["result_preview"]
    write_calls = [
        args for name, args in called_tools if name == "write_docx_content"
    ]
    assert write_calls == [
        {
            "path": target_path,
            "paragraphs": '[{"text":"Koto AI 助手写入测试","style":"Heading 1"},{"text":"workspace file content"}]',
        }
    ]
    assert changed == [
        {
            "path": target_path,
            "operation": "write_docx_content",
            "summary": "已写入 2 个段落到 Word 文档",
            "preview": "",
            "file_type": "docx",
            "change_type": "create",
            "paragraphs_written": 2,
            "focus": True,
        }
    ]
    assert run_started.payload["target_path"] == target_path
    assert source_path.read_text(encoding="utf-8") == "workspace file content"
