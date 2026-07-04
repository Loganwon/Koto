import json
import zipfile
from pathlib import Path

import pytest


def test_parse_file_to_text_reads_office_windows(tmp_path):
    import openpyxl
    from docx import Document
    from pptx import Presentation

    from app.core.agent.task_tools import parse_file_to_text

    docx_path = tmp_path / "windowed.docx"
    document = Document()
    for index in range(1, 6):
        document.add_paragraph(f"Paragraph {index}")
    document.save(docx_path)

    xlsx_path = tmp_path / "windowed.xlsx"
    workbook = openpyxl.Workbook()
    workbook.active.title = "First"
    workbook.active.append(["first value"])
    second = workbook.create_sheet("Second")
    second.append(["second value"])
    workbook.save(xlsx_path)

    pptx_path = tmp_path / "windowed.pptx"
    presentation = Presentation()
    for index in range(1, 5):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"Slide {index}"
    presentation.save(pptx_path)

    docx_text = parse_file_to_text(
        str(docx_path),
        window_unit="paragraph",
        start=2,
        end=3,
    )
    xlsx_text = parse_file_to_text(
        str(xlsx_path),
        window_unit="sheet",
        sheet_index=1,
    )
    pptx_text = parse_file_to_text(
        str(pptx_path),
        window_unit="slide",
        start=2,
        end=3,
    )

    assert "Paragraph 2" in docx_text
    assert "Paragraph 4" not in docx_text
    assert "Second" in xlsx_text
    assert "second value" in xlsx_text
    assert "first value" not in xlsx_text
    assert "Slide 2" in pptx_text
    assert "Slide 4" not in pptx_text


def test_parse_file_to_text_pdf_letter_window_skips_table_of_contents(tmp_path, monkeypatch):
    from app.core.agent import task_tools

    pdf_path = tmp_path / "schiller.pdf"
    pdf_path.write_bytes(b"%PDF-1.4\n% fake test pdf\n")

    def fake_pdf_excerpt(
        path,
        *,
        max_chars,
        start_page=1,
        end_page=0,
        allow_full_fallback=True,
    ):
        if start_page == 6 and end_page == 6:
            return (
                "[Page 6]\n\u76ee \u5f55\n"
                "\u7b2c\u5341\u4e00\u5c01\u4fe1\n"
                "\u7b2c\u5341\u4e8c\u5c01\u4fe1\n"
                "\u7b2c\u5341\u4e09\u5c01\u4fe1\n"
                "\u7b2c\u5341\u56db\u5c01\u4fe1\n"
                "\u7b2c\u5341\u4e94\u5c01\u4fe1\n"
                "\u7b2c\u5341\u516d\u5c01\u4fe1"
            )
        if start_page == 63 and end_page == 63:
            return "[Page 63]\n\u7b2c\u5341\u4e00\u5c01\u4fe1\n\u4eba\u683c\u4e0e\u72b6\u6001\u3002"
        if start_page == 87 and end_page == 87:
            return "[Page 87]\n\u7b2c\u5341\u516d\u5c01\u4fe1\n\u540e\u7eed\u7ae0\u8282\u3002"
        if start_page == 63 and end_page == 86:
            return (
                "[Page 63]\n\u7b2c\u5341\u4e00\u5c01\u4fe1\n\u4eba\u683c\u4e0e\u72b6\u6001\u3002\n"
                "[Page 81]\n\u7b2c\u5341\u4e94\u5c01\u4fe1\n"
                "\u6e38\u620f\u51b2\u52a8\u7684\u5bf9\u8c61\u662f\u6d3b\u7684\u5f62\u8c61\u3002"
            )
        return ""

    monkeypatch.setattr(task_tools, "_read_pdf_excerpt", fake_pdf_excerpt)

    text = task_tools.parse_file_to_text(
        str(pdf_path),
        window_unit="pdf_letter",
        start=11,
        end=15,
        max_chars=4000,
    )

    assert "[PDF letter window: 11-15; resolved pages 63-86]" in text
    assert "\u7b2c\u5341\u4e00\u5c01\u4fe1" in text
    assert "\u7b2c\u5341\u4e94\u5c01\u4fe1" in text
    assert "\u6e38\u620f\u51b2\u52a8" in text


def test_file_task_event_schema_exposes_diff_contract():
    from app.core.agent.file_task_contract import file_task_event_schema

    schema = file_task_event_schema()

    assert schema["title"] == "FileTaskEvent"
    assert "payload" in schema["properties"]
    assert schema["properties"]["payload"]["properties"]["diff"]["$ref"].endswith(
        "FileTaskDiff"
    )
    assert schema["$defs"]["FileTaskDiff"]["required"] == [
        "kind",
        "items",
        "changed_count",
    ]


def test_extract_koto_paths_reads_primary_structured_marker_keys():
    from app.core.agent.file_task_result_markers import (
        KOTO_CREATED_RESULT_MARKER,
        KOTO_MODIFIED_RESULT_MARKER,
    )
    from app.core.agent.file_task_tool_catalog import extract_koto_paths

    result = {
        "__koto_created__": ["created.docx"],
        "__koto_modified__": ["modified.xlsx"],
    }

    assert extract_koto_paths(result, KOTO_CREATED_RESULT_MARKER) == ["created.docx"]
    assert extract_koto_paths(result, KOTO_MODIFIED_RESULT_MARKER) == ["modified.xlsx"]


def test_extract_file_changes_reads_run_python_spreadsheet_metrics():
    from app.core.agent.file_task_result_markers import (
        KOTO_CREATED_RESULT_MARKER,
        KOTO_MODIFIED_RESULT_MARKER,
    )
    from app.core.agent.file_task_tool_feedback import extract_file_changes

    result = {
        "stdout": (
            "Data rows written: 4 rows\n"
            "Total cells written: 20\n"
            "KOTO_CREATED:C:\\workspace\\sales_profit_report.xlsx"
        ),
        "__koto_created__": ["C:\\workspace\\sales_profit_report.xlsx"],
    }

    changes = extract_file_changes(
        "run_python_code",
        {"code": "write report"},
        result,
        created_marker=KOTO_CREATED_RESULT_MARKER,
        modified_marker=KOTO_MODIFIED_RESULT_MARKER,
    )

    assert changes[0]["operation"] == "run_python_code"
    assert changes[0]["rows_written"] == 4
    assert changes[0]["cells_written"] == 20


def test_task_tools_create_file_uses_context_directory_for_bare_output(
    tmp_path, monkeypatch
):
    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    source_dir = tmp_path / "codex_context_dir"
    source_dir.mkdir()
    (source_dir / "orders.csv").write_text("sku,units\nA100,1\n", encoding="utf-8")

    provider = task_tools.TaskToolsPlugin(
        workspace_root=str(tmp_path),
        request_context={
            "task": "请从 codex_context_dir 读取 orders.csv 并创建 restock_plan.csv",
        },
    )
    create_tool = next(
        tool for tool in provider.get_tools() if tool["name"] == "create_file"
    )

    payload = json.loads(
        create_tool["func"](
            "restock_plan.csv",
            "sku,restock_quantity\nA100,30\n",
        )
    )

    assert payload["success"] is True
    assert payload["path"] == "codex_context_dir/restock_plan.csv"
    assert (source_dir / "restock_plan.csv").exists()
    assert not (tmp_path / "restock_plan.csv").exists()


def test_run_python_relocates_root_created_files_to_context_directory(
    tmp_path, monkeypatch
):
    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    source_dir = tmp_path / "codex_context_dir"
    source_dir.mkdir()
    root_output = tmp_path / "restock_plan.csv"
    root_output.write_text("sku,restock_quantity\nA100,30\n", encoding="utf-8")

    result = task_tools._relocate_root_created_files_to_output_dir(
        {"stdout": f"KOTO_CREATED:{root_output}"},
        output_dir="codex_context_dir",
    )

    relocated = source_dir / "restock_plan.csv"
    assert relocated.exists()
    assert not root_output.exists()
    assert str(relocated) in result["stdout"]
    assert str(root_output) not in result["stdout"]


def test_task_tools_context_directory_uses_resolved_workspace_root(
    tmp_path, monkeypatch
):
    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    source_dir = tmp_path / "codex_context_dir"
    source_dir.mkdir()
    plugin = task_tools.TaskToolsPlugin(
        workspace_root=str(tmp_path),
        request_context={
            "task": "读取 codex_context_dir/orders.csv 和 codex_context_dir/rules.md，生成 restock_plan.csv",
        },
    )

    assert plugin._contextual_output_directory() == "codex_context_dir"


def test_create_file_docx_emits_docx_write_metrics_and_valid_package(
    tmp_path, monkeypatch
):
    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    result = json.loads(
        task_tools.create_file(
            "step-summary.docx",
            "# 第 0 步：封面与目录信息摘要\n\n## 来源页码：Page 4-8\n\n- 书名：The Global Rules of Art\n- 作者：Larissa Buchholz",
        )
    )

    assert result["success"] is True
    assert result["operation"] == "write_docx_content"
    assert result["file_type"] == "docx"
    assert result["change_type"] == "create"
    assert result["paragraphs_written"] >= 4

    from docx import Document

    doc = Document(tmp_path / "step-summary.docx")
    text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
    assert "第 0 步：封面与目录信息摘要" in text
    assert "Larissa Buchholz" in text
    assert "**" not in text


def test_create_file_docx_cleans_markdown_rule_and_bold_markers(tmp_path, monkeypatch):
    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    result = json.loads(
        task_tools.create_file(
            "clean-summary.docx",
            "# 当前页窗摘要（第 4-6 页）\n---\n**文档识别：** 年报目录。\n- 来源页码：第 4-6 页",
        )
    )

    assert result["success"] is True

    from docx import Document

    doc = Document(tmp_path / "clean-summary.docx")
    paragraphs = [paragraph.text for paragraph in doc.paragraphs]
    text = "\n".join(paragraphs)
    assert "---" not in text
    assert "**" not in text
    assert "文档识别：" in text
    assert paragraphs[0] == "当前页窗摘要（第 4-6 页）"


def test_write_docx_content_parses_loose_paragraph_objects_with_inner_quotes(
    tmp_path, monkeypatch
):
    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    loose = (
        '[{"text": "Section 1", "style": "Normal"}, '
        '{"text": "关于"身体"的讨论", "style": "Normal"}]'
    )

    result = json.loads(task_tools.write_docx_content("loose.docx", loose))

    from docx import Document

    doc = Document(tmp_path / "loose.docx")
    paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text]
    assert result["success"] is True
    assert result["paragraphs_written"] == 2
    assert paragraphs == ["Section 1", '关于"身体"的讨论']


def test_write_docx_content_returns_paragraph_diff_in_file_change(
    tmp_path, monkeypatch
):
    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_tool_catalog import parse_file_change

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    result_text = task_tools.write_docx_content(
        "diff-demo.docx",
        json.dumps([{"text": "第一段"}, {"text": "第二段"}], ensure_ascii=False),
    )
    result = json.loads(result_text)
    change = parse_file_change(
        "write_docx_content", {"path": "diff-demo.docx"}, result_text
    )

    assert result["diff"]["kind"] == "docx_paragraphs"
    assert result["diff"]["changed_count"] == 2
    assert result["diff"]["items"][0]["before"] == ""
    assert result["diff"]["items"][0]["after"] == "第一段"
    assert change["diff"]["kind"] == "docx_paragraphs"
    assert change["summary_code"] == "CREATE_OK"


def test_insert_docx_paragraph_preserves_existing_table(tmp_path, monkeypatch):
    monkeypatch.chdir(tmp_path)
    from docx import Document

    from app.core.agent import task_tools

    doc = Document()
    doc.add_heading("Risk Review", level=1)
    doc.add_paragraph("Existing risk text.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Customer"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "Blue Harbor"
    table.cell(1, 1).text = "128000"
    doc.add_heading("Next Actions", level=1)
    doc.save("report.docx")

    result = json.loads(
        task_tools.insert_docx_paragraph(
            "report.docx",
            "Overall risk level: Moderate.",
            before_heading="Next Actions",
        )
    )
    updated = Document("report.docx")
    texts = [paragraph.text for paragraph in updated.paragraphs]

    assert result["operation"] == "insert_docx_paragraph"
    assert result["paragraphs_written"] == 1
    assert "Overall risk level: Moderate." in texts
    assert texts.index("Overall risk level: Moderate.") < texts.index("Next Actions")
    assert len(updated.tables) == 1
    assert updated.tables[0].cell(1, 0).text == "Blue Harbor"


def test_create_file_xlsx_emits_sheet_metrics_and_valid_workbook(tmp_path, monkeypatch):
    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    result = json.loads(task_tools.create_file("analysis.xlsx", "指标,数值\n样本,42"))

    assert result["success"] is True
    assert result["operation"] == "write_sheet_data"
    assert result["file_type"] == "xlsx"
    assert result["rows_written"] == 2
    assert result["columns_written"] == 2

    import openpyxl

    workbook = openpyxl.load_workbook(tmp_path / "analysis.xlsx")
    try:
        sheet = workbook.active
        assert sheet.cell(row=1, column=1).value == "指标"
        assert sheet.cell(row=2, column=2).value == "42"
    finally:
        workbook.close()


def test_write_sheet_data_returns_cell_diff_in_file_change(tmp_path, monkeypatch):
    import openpyxl

    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_tool_catalog import parse_file_change

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    workbook_path = tmp_path / "budget.xlsx"
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "预算"
    sheet["A1"] = "旧值"
    workbook.save(workbook_path)
    workbook.close()

    result_text = task_tools.write_sheet_data(
        "budget.xlsx",
        sheet_name="预算",
        updates=json.dumps([{"row": 1, "col": 1, "value": "新值"}], ensure_ascii=False),
    )
    result = json.loads(result_text)
    change = parse_file_change("write_sheet_data", {"path": "budget.xlsx"}, result_text)

    assert result["diff"]["kind"] == "xlsx_cells"
    assert result["diff"]["changed_count"] == 1
    assert result["diff"]["items"][0]["cell"] == "A1"
    assert result["diff"]["items"][0]["before"] == "旧值"
    assert result["diff"]["items"][0]["after"] == "新值"
    assert change["diff"]["items"][0]["sheet"] == "预算"
    assert change["summary_code"] == "WRITE_OK"


def test_file_task_change_tracker_mirrors_list_appends_to_coordinator():
    from app.core.agent.file_task_change_tracker import FileTaskChangeTracker

    tracker = FileTaskChangeTracker()
    tracker.changes.append(
        {
            "path": "report.docx",
            "operation": "write_docx_content",
            "change_type": "modify",
            "diff": {
                "kind": "docx_paragraphs",
                "items": [{"before": "old", "after": "new"}],
                "changed_count": 1,
            },
        }
    )

    assert len(tracker.changes) == 1
    coordinator_changes = tracker.coordinator_changes()
    assert len(coordinator_changes) == 1
    assert coordinator_changes[0]["file_path"] == "report.docx"
    assert coordinator_changes[0]["original"] == "old"
    assert coordinator_changes[0]["modified"] == "new"


def test_fill_docx_template_replaces_placeholders_and_emits_diff(
    tmp_path, monkeypatch
):
    from docx import Document

    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_tool_catalog import parse_file_change

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    source = tmp_path / "template.docx"
    doc = Document()
    doc.add_paragraph("甲方：{{party_a}}")
    doc.add_paragraph("金额：{amount}")
    doc.save(source)

    result_text = task_tools.fill_docx_template(
        "template.docx",
        data=json.dumps({"party_a": "杭州公司", "amount": "100万元"}, ensure_ascii=False),
        target_path="filled.docx",
    )
    result = json.loads(result_text)
    change = parse_file_change(
        "fill_docx_template",
        {"path": "template.docx", "target_path": "filled.docx"},
        result_text,
    )

    assert result["success"] is True
    assert result["operation"] == "fill_docx_template"
    assert result["diff"]["kind"] == "docx_template_fields"
    assert result["diff"]["changed_count"] == 2
    assert result["fields_filled"] == ["amount", "party_a"]
    assert change["diff"]["items"][0]["before"] == "甲方：{{party_a}}"
    saved = Document(tmp_path / "filled.docx")
    assert [paragraph.text for paragraph in saved.paragraphs] == [
        "甲方：杭州公司",
        "金额：100万元",
    ]


def test_convert_docx_to_pdf_emits_file_change_with_converter(
    tmp_path, monkeypatch
):
    from docx import Document

    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_tool_catalog import parse_file_change

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    source = tmp_path / "report.docx"
    doc = Document()
    doc.add_paragraph("报告")
    doc.save(source)

    def fake_converter(source_path, target_path):
        assert source_path.endswith("report.docx")
        with open(target_path, "wb") as fh:
            fh.write(b"%PDF-1.4\n% koto test\n")
        return "fake_converter"

    monkeypatch.setattr(task_tools, "_convert_docx_to_pdf_with_docx2pdf", fake_converter)

    result_text = task_tools.convert_docx_to_pdf("report.docx", "report.pdf")
    result = json.loads(result_text)
    change = parse_file_change(
        "convert_docx_to_pdf",
        {"path": "report.docx", "target_path": "report.pdf"},
        result_text,
    )

    assert result["success"] is True
    assert result["operation"] == "convert_docx_to_pdf"
    assert result["converter"] == "fake_converter"
    assert result["summary_code"] == "CONVERT_OK"
    assert change["path"] == "report.pdf"
    assert change["converter"] == "fake_converter"
    assert (tmp_path / "report.pdf").exists()


def test_convert_file_emits_standard_file_change(tmp_path, monkeypatch):
    import app.core.agent.task_tools as task_tools
    import web.file_converter as file_converter
    from app.core.agent.file_task_tool_catalog import parse_file_change

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    source = tmp_path / "notes.txt"
    source.write_text("hello", encoding="utf-8")

    def fake_convert(source_path, target_format, output_path=None, output_dir=None):
        assert source_path.endswith("notes.txt")
        assert target_format == ".md"
        assert output_path is not None
        Path(output_path).write_text("# hello", encoding="utf-8")
        return {
            "success": True,
            "output_path": output_path,
            "from_format": "txt",
            "to_format": "md",
            "message": "converted",
            "warning": "",
            "error": "",
        }

    monkeypatch.setattr(file_converter, "convert", fake_convert)

    result_text = task_tools.convert_file("notes.txt", "md", "notes.md")
    result = json.loads(result_text)
    change = parse_file_change(
        "convert_file",
        {"file_path": "notes.txt", "target_format": "md", "output_path": "notes.md"},
        result_text,
    )

    assert result["success"] is True
    assert result["operation"] == "convert_file"
    assert result["summary_code"] == "CONVERT_OK"
    assert result["target_format"] == "md"
    assert change["path"] == "notes.md"
    assert change["target_format"] == "md"
    assert (tmp_path / "notes.md").read_text(encoding="utf-8") == "# hello"


def test_list_conversions_returns_structured_matrix():
    import app.core.agent.task_tools as task_tools

    result = json.loads(task_tools.list_conversions("txt"))

    assert result["success"] is True
    assert result["source_format"] == "txt"
    assert "md" in result["targets"]
    assert result["summary"]


def test_create_file_pptx_emits_slide_metrics_and_valid_deck(tmp_path, monkeypatch):
    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    result = json.loads(task_tools.create_file("deck.pptx", "# 汇报标题\n- 关键发现\n- 后续行动"))

    assert result["success"] is True
    assert result["operation"] == "add_pptx_slides"
    assert result["file_type"] == "pptx"
    assert result["slides_added"] >= 1

    from pptx import Presentation

    presentation = Presentation(tmp_path / "deck.pptx")
    assert len(presentation.slides) >= 1


def test_verify_task_completion_uses_structured_docx_table_metadata():
    from app.core.agent.task_tools import verify_task_completion

    result = json.loads(
        verify_task_completion(
            task_description="把 xlsx 表格加入 docx",
            file_states=json.dumps(
                [
                    {
                        "path": "report.docx",
                        "exists": True,
                        "modified": True,
                        "preview": "...",
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": "report.docx",
                        "operation": "insert_excel_as_docx_table",
                        "sheet": "汇总表",
                        "rows_written": 200,
                        "columns_written": 4,
                    }
                ],
                ensure_ascii=False,
            ),
            target_path="report.docx",
        )
    )

    assert result["completed"] is True
    assert "文件已成功修改：report.docx" in result["summary"]
    assert "工作表“汇总表”" in result["summary"]
    assert "200 行 × 4 列" in result["summary"]
    assert any(
        item["criterion"] == "all_tracked_files_modified" and item["passed"] is True
        for item in result["criteria_results"]
    )
    assert any(
        item["criterion"] == "target_file_hit" and item["passed"] is True
        for item in result["criteria_results"]
    )


def test_verify_task_completion_rejects_missing_explicit_output_file():
    from app.core.agent.task_tools import verify_task_completion

    result = json.loads(
        verify_task_completion(
            task_description=(
                "读取 codex_context_dir/orders.csv，生成 optimized_restock_plan.csv "
                "和 optimized_operations_report.md。"
            ),
            file_states=json.dumps(
                [
                    {
                        "path": "optimized_restock_plan.csv",
                        "exists": True,
                        "modified": True,
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": "optimized_restock_plan.csv",
                        "operation": "run_python_code",
                    }
                ],
                ensure_ascii=False,
            ),
        )
    )

    assert result["completed"] is False
    assert result["criteria_results"][0]["criterion"] == "explicit_output_files_present"
    assert "optimized_operations_report.md" in result["summary"]


def test_verify_task_completion_summarizes_multiple_docx_changes_on_target():
    from app.core.agent.task_tools import verify_task_completion

    result = json.loads(
        verify_task_completion(
            task_description="把 xlsx 表格加入 docx，并追加核验说明",
            file_states=json.dumps(
                [
                    {
                        "path": "report.docx",
                        "exists": True,
                        "modified": True,
                        "preview": "...",
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": "report.docx",
                        "operation": "write_docx_content",
                        "paragraphs_written": 2,
                    },
                    {
                        "path": "report.docx",
                        "operation": "insert_excel_as_docx_table",
                        "sheet": "Budget",
                        "rows_written": 4,
                        "columns_written": 5,
                    },
                ],
                ensure_ascii=False,
            ),
            target_path="report.docx",
        )
    )

    assert result["completed"] is True
    assert "文件已成功修改：report.docx" in result["summary"]
    assert "已写入 2 个段落" in result["summary"]
    assert "工作表“Budget”" in result["summary"]
    assert "4 行 × 5 列" in result["summary"]
    assert "相关文件变更" not in result["summary"]
    assert "其他文件变更" not in result["summary"]


def test_verify_task_completion_uses_structured_docx_image_metadata():
    from app.core.agent.task_tools import verify_task_completion

    result = json.loads(
        verify_task_completion(
            task_description="把图表加入 docx",
            file_states=json.dumps(
                [
                    {
                        "path": "report.docx",
                        "exists": True,
                        "modified": True,
                        "preview": "...",
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": "report.docx",
                        "operation": "insert_image_into_docx",
                        "image_name": "chart.png",
                        "images_inserted": 1,
                        "caption": "收入与利润趋势",
                    }
                ],
                ensure_ascii=False,
            ),
            target_path="report.docx",
        )
    )

    assert result["completed"] is True
    assert "文件已成功修改：report.docx" in result["summary"]
    assert "已插入 1 张图片" in result["summary"]
    assert "chart.png" in result["summary"]


def test_verify_task_completion_summarizes_multiple_docx_images_once():
    from app.core.agent.task_tools import verify_task_completion

    result = json.loads(
        verify_task_completion(
            task_description="把多张图表加入 docx",
            file_states=json.dumps(
                [
                    {
                        "path": "report.docx",
                        "exists": True,
                        "modified": True,
                        "preview": "...",
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": "report.docx",
                        "operation": "insert_image_into_docx",
                        "image_name": "chart1_revenue_profit_trend.png",
                        "images_inserted": 1,
                        "caption": "收入和利润趋势",
                    },
                    {
                        "path": "report.docx",
                        "operation": "insert_image_into_docx",
                        "image_name": "chart2_product_mix.png",
                        "images_inserted": 1,
                        "caption": "产品结构",
                    },
                    {
                        "path": "report.docx",
                        "operation": "insert_image_into_docx",
                        "image_name": "chart3_margin_analysis.png",
                        "images_inserted": 1,
                    },
                ],
                ensure_ascii=False,
            ),
            target_path="report.docx",
        )
    )

    assert result["completed"] is True
    assert "文件已成功修改：report.docx" in result["summary"]
    assert "已插入 3 张图片" in result["summary"]
    assert "chart1_revenue_profit_trend.png、chart2_product_mix.png、chart3_margin_analysis.png" in result["summary"]
    assert result["summary"].count("已插入") == 1


def test_verify_task_completion_ignores_intermediate_chart_artifacts_for_docx_target():
    from app.core.agent.task_tools import verify_task_completion

    result = json.loads(
        verify_task_completion(
            task_description="分析这个xlsx财务预测，找出其中的问题并将数据做成图，将数据和图都加入docx",
            file_states=json.dumps(
                [
                    {
                        "path": "chart1_revenue_profit_trend.png",
                        "exists": False,
                        "modified": False,
                        "preview": "生成收入利润趋势图。",
                    },
                    {
                        "path": "report.docx",
                        "exists": True,
                        "modified": True,
                        "preview": "财务预测问题分析。",
                    },
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": "chart1_revenue_profit_trend.png",
                        "operation": "run_python_code",
                        "file_type": "png",
                        "summary": "生成收入利润趋势图。",
                    },
                    {
                        "path": "report.docx",
                        "operation": "write_docx_content",
                        "file_type": "docx",
                        "paragraphs_written": 12,
                    },
                    {
                        "path": "report.docx",
                        "operation": "insert_image_into_docx",
                        "image_name": "chart1_revenue_profit_trend.png",
                        "images_inserted": 1,
                    },
                ],
                ensure_ascii=False,
            ),
            target_path="report.docx",
        )
    )

    assert result["completed"] is True
    assert "文件已成功修改：report.docx" in result["summary"]
    assert "以下文件尚未修改" not in result["summary"]


def test_verify_task_completion_rejects_table_only_result_when_task_requires_summary_text():
    from app.core.agent.task_tools import verify_task_completion

    result = json.loads(
        verify_task_completion(
            task_description="整理 xlsx 中的财务预测，并加入 docx",
            file_states=json.dumps(
                [
                    {
                        "path": "report.docx",
                        "exists": True,
                        "modified": True,
                        "preview": "...",
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": "report.docx",
                        "operation": "insert_excel_as_docx_table",
                        "sheet": "P&L",
                        "rows_written": 50,
                        "columns_written": 13,
                    }
                ],
                ensure_ascii=False,
            ),
            target_path="report.docx",
        )
    )

    assert result["completed"] is False
    assert "当前只写入了表格" in result["summary"]
    assert result["remaining_steps"] == [
        "先提炼关键结论，再用 write_docx_content 把摘要/说明写入目标 DOCX"
    ]
    assert any(
        item["criterion"] == "docx_narrative_write_present" and item["passed"] is False
        for item in result["criteria_results"]
    )


def test_verify_task_completion_detects_target_mismatch_from_structured_changes():
    from app.core.agent.task_tools import verify_task_completion

    result = json.loads(
        verify_task_completion(
            task_description="修改目标报告",
            file_states=json.dumps(
                [
                    {
                        "path": "other.docx",
                        "exists": True,
                        "modified": True,
                        "preview": "...",
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": "other.docx",
                        "operation": "write_docx_content",
                        "paragraphs_written": 1,
                    }
                ],
                ensure_ascii=False,
            ),
            target_path="report.docx",
        )
    )

    assert result["completed"] is False
    assert "未命中目标文件：report.docx" in result["summary"]
    assert result["remaining_steps"] == ["把结果写入 report.docx"]
    assert result["criteria_results"] == [
        {
            "criterion": "target_file_hit",
            "passed": False,
            "detail": "已修改 other.docx，但未命中目标文件：report.docx",
            "priority": "critical",
        }
    ]


def test_verify_task_completion_accepts_display_name_for_absolute_target(tmp_path):
    from app.core.agent.task_tools import verify_task_completion

    target = tmp_path / "report.docx"

    result = json.loads(
        verify_task_completion(
            task_description="创建 docx 总结",
            file_states=json.dumps(
                [
                    {
                        "path": target.name,
                        "exists": True,
                        "modified": True,
                        "preview": "已写入摘要",
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": target.name,
                        "operation": "write_docx_content",
                        "paragraphs_written": 3,
                    }
                ],
                ensure_ascii=False,
            ),
            target_path=str(target),
        )
    )

    assert result["completed"] is True
    assert "文件已成功修改：report.docx" in result["summary"]
    assert any(
        item["criterion"] == "target_file_hit" and item["passed"] is True
        for item in result["criteria_results"]
    )


def test_verify_task_completion_rejects_locked_target_fallback_copy_as_original_write():
    from app.core.agent.task_tools import verify_task_completion

    result = json.loads(
        verify_task_completion(
            task_description="把 xlsx 表格加入 docx",
            file_states=json.dumps(
                [
                    {
                        "path": "report.koto-copy.docx",
                        "exists": True,
                        "modified": True,
                        "preview": "...",
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": "report.koto-copy.docx",
                        "operation": "insert_excel_as_docx_table",
                        "sheet": "汇总表",
                        "rows_written": 200,
                        "columns_written": 4,
                        "original_target_path": "report.docx",
                        "fallback_copy": True,
                        "blocked_target": True,
                    }
                ],
                ensure_ascii=False,
            ),
            target_path="report.docx",
        )
    )

    assert result["completed"] is False
    assert "目标文件尚未完成修改：report.docx" in result["summary"]
    assert "report.koto-copy.docx" in result["summary"]
    assert result["remaining_steps"] == [
        "检查 report.docx 的文件权限；如果文件正在被占用，关闭相关程序后重新写回原文件"
    ]


def test_verify_task_completion_matches_workspace_relative_target_to_absolute_change(
    tmp_path, monkeypatch
):
    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    from app.core.agent.task_tools import verify_task_completion

    target_path = tmp_path / "notes.txt"
    target_path.write_text("updated", encoding="utf-8")

    result = json.loads(
        verify_task_completion(
            task_description="润色当前文本并写回原文件",
            file_states=json.dumps(
                [
                    {
                        "path": str(target_path),
                        "exists": True,
                        "modified": True,
                        "preview": "updated",
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": str(target_path),
                        "operation": "run_python_code",
                        "summary": "Python 代码更新了 notes.txt",
                    }
                ],
                ensure_ascii=False,
            ),
            target_path="notes.txt",
        )
    )

    assert result["completed"] is True
    assert "文件已成功修改：notes.txt" in result["summary"]


def test_verify_task_completion_rejects_same_basename_in_wrong_directory(tmp_path):
    from app.core.agent.task_tools import verify_task_completion

    expected_dir = tmp_path / "expected"
    actual_dir = tmp_path / "actual"
    expected_dir.mkdir()
    actual_dir.mkdir()
    expected_target = expected_dir / "report.docx"
    actual_target = actual_dir / "report.docx"
    expected_target.write_text("original", encoding="utf-8")
    actual_target.write_text("updated", encoding="utf-8")

    result = json.loads(
        verify_task_completion(
            task_description="修改指定目录里的目标报告",
            file_states=json.dumps(
                [
                    {
                        "path": str(actual_target),
                        "exists": True,
                        "modified": True,
                        "preview": "updated",
                    }
                ],
                ensure_ascii=False,
            ),
            file_changes=json.dumps(
                [
                    {
                        "path": str(actual_target),
                        "operation": "write_docx_content",
                        "paragraphs_written": 2,
                    }
                ],
                ensure_ascii=False,
            ),
            target_path=str(expected_target),
        )
    )

    assert result["completed"] is False
    assert f"未命中目标文件：{expected_target.name}" in result["summary"]
    assert result["remaining_steps"] == [f"把结果写入 {expected_target.name}"]


def test_task_file_sandbox_staging_sanitizes_invalid_display_names(tmp_path):
    from app.core.agent.task_tools import (
        _prepend_task_file_context,
        _stage_task_files_for_sandbox,
    )

    source = tmp_path / "source.xlsx"
    source.write_text("demo", encoding="utf-8")
    sandbox = tmp_path / "sandbox"
    sandbox.mkdir()

    staged = _stage_task_files_for_sandbox(
        [
            {
                "display_name": "????-financial model.xlsx",
                "source_path": str(source),
                "source_fingerprint_initial": {},
            }
        ],
        str(sandbox),
    )
    preamble = _prepend_task_file_context("print('ok')", staged)

    assert staged[0]["staged_name"] == "____-financial model.xlsx"
    assert "?" not in staged[0]["staged_path"]
    assert (sandbox / "____-financial model.xlsx").exists()
    assert '"????-financial model.xlsx"' in preamble
    assert "____-financial model.xlsx" in preamble


def test_copy_file_delegates_to_canonical_file_service(tmp_path, monkeypatch):
    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    source = tmp_path / "source.txt"
    source.write_text("hello", encoding="utf-8")
    calls = []

    class FakeFileService:
        def __init__(self, *, workspace_dir, backup_enabled):
            calls.append(
                {
                    "workspace_dir": workspace_dir,
                    "backup_enabled": backup_enabled,
                }
            )

        def copy_file(self, source_path, destination_path, overwrite=False):
            calls.append(
                {
                    "source": source_path,
                    "destination": destination_path,
                    "overwrite": overwrite,
                }
            )
            return {
                "success": True,
                "destination": destination_path,
            }

    monkeypatch.setattr(task_tools, "FileService", FakeFileService)

    payload = json.loads(task_tools.copy_file("source.txt", "copied.txt"))

    assert calls == [
        {
            "workspace_dir": str(tmp_path),
            "backup_enabled": False,
        },
        {
            "source": str(source),
            "destination": str(tmp_path / "copied.txt"),
            "overwrite": True,
        },
    ]
    assert payload["success"] is True
    assert payload["path"] == "copied.txt"
    assert payload["operation"] == "copy_file"
    assert payload["change_type"] == "create"


def test_annotate_file_returns_standard_file_change_payload(tmp_path, monkeypatch):
    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_tool_catalog import parse_file_change

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    notes_path = tmp_path / "notes.txt"
    notes_path.write_text("alpha beta", encoding="utf-8")

    result = task_tools.annotate_file(
        "notes.txt",
        [{"range_start": 0, "range_end": 5, "comment": "需要核对"}],
    )
    payload = json.loads(result)
    change = parse_file_change("annotate_file", {"path": "notes.txt"}, result)

    assert payload["success"] is True
    assert payload["path"] == "notes.txt"
    assert payload["operation"] == "annotate_file"
    assert payload["change_type"] == "annotate"
    assert payload["annotations_added"] == 1
    assert change["path"] == "notes.txt"
    assert change["operation"] == "annotate_file"
    assert change["annotations_added"] == 1


def test_replace_file_selection_writes_text_selection_and_registers_change(
    tmp_path, monkeypatch
):
    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_tool_catalog import parse_file_change

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    notes_path = tmp_path / "notes.md"
    notes_path.write_text("开场很啰嗦，需要改短。\n下一段保留。", encoding="utf-8")

    result = task_tools.replace_file_selection(
        "notes.md",
        original_selection="开场很啰嗦，需要改短。",
        new_content="开场应更短、更直接。",
    )
    payload = json.loads(result)
    change = parse_file_change(
        "replace_file_selection",
        {
            "path": "notes.md",
            "original_selection": "开场很啰嗦，需要改短。",
            "new_content": "开场应更短、更直接。",
        },
        result,
    )

    assert payload["success"] is True
    assert payload["path"] == "notes.md"
    assert payload["operation"] == "replace_file_selection"
    assert payload["replacements_made"] == 1
    assert change["operation"] == "replace_file_selection"
    assert change["replacements_made"] == 1
    assert "开场应更短、更直接。" in notes_path.read_text(encoding="utf-8")
    assert (tmp_path / "notes.md.bak").exists()


def test_replace_file_selection_requires_existing_selection(tmp_path, monkeypatch):
    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_tool_catalog import parse_file_change

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    notes_path = tmp_path / "notes.txt"
    notes_path.write_text("alpha beta", encoding="utf-8")

    result = task_tools.replace_file_selection(
        "notes.txt",
        original_selection="gamma",
        new_content="delta",
    )
    payload = json.loads(result)
    change = parse_file_change(
        "replace_file_selection",
        {"path": "notes.txt", "original_selection": "gamma", "new_content": "delta"},
        result,
    )

    assert payload["changed"] is False
    assert payload["error"] == "original_selection_not_found"
    assert change is None
    assert notes_path.read_text(encoding="utf-8") == "alpha beta"


def test_annotate_file_docx_requirement_returns_streaming_native_tool_result(
    tmp_path, monkeypatch
):
    import app.core.agent.file_task_doc_annotate_bridge as bridge
    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_contract import (
        FileTaskToolStreamChunk,
        FileTaskToolStreamResult,
    )
    from app.core.agent.file_task_tool_gateway import (
        FileTaskToolContext,
        FileTaskToolGateway,
    )

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))
    docx_path = tmp_path / "draft.docx"
    docx_path.write_bytes(b"PK\x03\x04")

    captured = {}

    def fake_stream_request_as_tool(request, *, workspace_root="", gemini_client=None):
        captured["request"] = request
        captured["workspace_root"] = workspace_root
        captured["gemini_client"] = gemini_client
        return FileTaskToolStreamResult(
            chunks=[
                FileTaskToolStreamChunk(
                    kind="event",
                    event_type="step_progress",
                    payload={
                        "detail": "已写入 1/2 条修订",
                        "file_updated": True,
                        "path": "draft.docx",
                        "file_path": "draft.docx",
                        "applied": 1,
                    },
                ),
                FileTaskToolStreamChunk(
                    kind="result",
                    payload={
                        "success": True,
                        "path": "draft.docx",
                        "file_path": "draft.docx",
                        "annotations_added": 2,
                        "updated_in_place": True,
                    },
                ),
            ]
        )

    monkeypatch.setattr(bridge, "stream_request_as_tool", fake_stream_request_as_tool)

    gateway = FileTaskToolGateway(
        context=FileTaskToolContext(
            task_files=[{"path": "draft.docx", "type": "docx", "target": True}],
            workspace_root=str(tmp_path),
            gemini_client="gemini-client",
            request_context={
                "task": "请批注不通顺的地方",
                "target_path": "draft.docx",
                "model_mode": "cloud",
                "model_id": "gemini-2.5-pro",
            },
        )
    )

    result = gateway.execute(
        "annotate_file",
        {"path": "draft.docx", "annotations": "[]", "requirement": "请批注不通顺的地方"},
    )

    assert isinstance(result, FileTaskToolStreamResult)

    chunks = list(result.chunks)
    progress_chunk = next(
        chunk
        for chunk in chunks
        if chunk.kind == "event"
        and chunk.event_type == "step_progress"
        and chunk.payload.get("file_updated")
    )
    final_chunk = chunks[-1]

    assert captured["gemini_client"] == "gemini-client"
    assert captured["workspace_root"] == str(tmp_path)
    assert captured["request"].target_path == "draft.docx"
    assert captured["request"].task == "请批注不通顺的地方"
    assert captured["request"].model_id == "gemini-2.5-pro"
    assert any(
        file_info.type == "docx" and file_info.target
        for file_info in captured["request"].files
    )
    assert progress_chunk.payload["path"] == "draft.docx"
    assert final_chunk.kind == "result"
    assert final_chunk.payload["path"] == "draft.docx"
    assert final_chunk.payload["annotations_added"] == 2
    assert final_chunk.payload["updated_in_place"] is True


def test_clear_docx_review_marks_removes_docx_comments_and_registers_file_change(
    tmp_path, monkeypatch
):
    docx_module = pytest.importorskip("docx")

    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_tool_catalog import parse_file_change
    from web.track_changes_editor import TrackChangesEditor

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    docx_path = tmp_path / "draft.docx"
    document = docx_module.Document()
    document.add_paragraph("第一段用于清除批注测试。")
    document.save(docx_path)

    editor = TrackChangesEditor(author="Koto Test")
    applied = editor.apply_comment_changes(
        str(docx_path),
        [{"原文片段": "第一段用于清除批注测试。", "修改后文本": "建议改写", "修改原因": "测试"}],
    )

    assert applied["applied"] == 1

    import zipfile

    with zipfile.ZipFile(docx_path) as archive:
        assert "word/comments.xml" in archive.namelist()
        assert "commentReference" in archive.read("word/document.xml").decode(
            "utf-8", errors="ignore"
        )

    result = task_tools.clear_docx_review_marks("draft.docx", scope="comments")
    payload = json.loads(result)
    change = parse_file_change(
        "clear_docx_review_marks", {"path": "draft.docx", "scope": "comments"}, result
    )

    assert payload["success"] is True
    assert payload["path"] == "draft.docx"
    assert payload["operation"] == "clear_docx_review_marks"
    assert payload["scope"] == "comments"
    assert payload["comments_removed"] >= 1
    assert change["path"] == "draft.docx"
    assert change["operation"] == "clear_docx_review_marks"
    assert change["comments_removed"] >= 1

    with zipfile.ZipFile(docx_path) as archive:
        names = archive.namelist()
        document_xml = archive.read("word/document.xml").decode(
            "utf-8", errors="ignore"
        )
    assert "word/comments.xml" not in names
    assert "commentRangeStart" not in document_xml
    assert "commentRangeEnd" not in document_xml
    assert "commentReference" not in document_xml


def test_task_tools_plugin_exposes_clear_docx_review_marks_tool():
    from app.core.agent.task_tools import TaskToolsPlugin

    tool_names = {
        tool["name"]
        for tool in TaskToolsPlugin(
            task_files=[{"path": "draft.docx", "type": "docx"}]
        ).get_tools()
    }

    assert "clear_docx_review_marks" in tool_names
    assert "compare_docx_and_annotate" in tool_names
    assert "plan_docx_compare_annotations" in tool_names
    assert "write_docx_comments" in tool_names
    assert "replace_file_selection" in tool_names
    assert "fill_docx_template" in tool_names
    assert "convert_docx_to_pdf" in tool_names


def test_contract_risk_summary_groups_common_clause_changes():
    import app.core.agent.task_tools as task_tools

    risks = task_tools._contract_risk_summary_from_annotations(
        [
            {
                "原文片段": "付款期限调整为十五日。",
                "批注内容": "另一份为：付款期限为三十日。本文件为：付款期限为十五日。",
            },
            {
                "原文片段": "任一方可提前终止合同。",
                "批注内容": "另一份为：提前三十日通知。本文件为：提前七日通知。",
            },
            {
                "原文片段": "违约金按合同金额百分之二十计算。",
                "批注内容": "本文件提高违约金比例。",
            },
        ]
    )

    assert any(item.startswith("付款/费用") for item in risks)
    assert any(item.startswith("终止/解除") for item in risks)
    assert any(item.startswith("违约责任") for item in risks)


def test_compare_docx_and_annotate_writes_word_comments(tmp_path, monkeypatch):
    docx_module = pytest.importorskip("docx")

    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_tool_catalog import parse_file_change

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    original_path = tmp_path / "humanise!.docx"
    revised_path = tmp_path / "humanise!_revised.docx"

    original = docx_module.Document()
    original.add_paragraph("AI 助手应该保持稳定的任务执行。")
    original.add_paragraph("旧版本只给出建议，没有写回文件。")
    original.save(original_path)

    revised = docx_module.Document()
    revised.add_paragraph("AI 助手应该保持稳定且可核验的任务执行。")
    revised.add_paragraph("新版本必须写回文件，并展示真实差异。")
    revised.save(revised_path)

    annotations, detected = task_tools._build_docx_compare_annotations(
        ["AI 助手应该保持稳定的任务执行。"],
        ["AI 助手应该保持稳定且可核验的任务执行。"],
        max_differences=5,
    )
    assert detected == 1
    assert annotations[0]["原文片段"] != "AI 助手应该保持稳定且可核验的任务执行。"
    assert "可核验" in annotations[0]["原文片段"]
    assert "批注内容" in annotations[0]
    assert "修改后文本" not in annotations[0]

    result = task_tools.compare_docx_and_annotate(
        "humanise!.docx",
        "humanise!_revised.docx",
    )
    payload = json.loads(result)
    change = parse_file_change(
        "compare_docx_and_annotate",
        {
            "original_path": "humanise!.docx",
            "revised_path": "humanise!_revised.docx",
        },
        result,
    )

    assert payload["success"] is True
    assert payload["operation"] == "compare_docx_and_annotate"
    assert payload["annotations_added"] >= 1
    assert payload["differences_detected"] >= 1
    assert change["operation"] == "compare_docx_and_annotate"
    assert change["annotations_added"] >= 1

    with zipfile.ZipFile(revised_path, "r") as archive:
        names = archive.namelist()
        document_xml = archive.read("word/document.xml").decode(
            "utf-8", errors="ignore"
        )
        comments_xml = archive.read("word/comments.xml").decode(
            "utf-8", errors="ignore"
        )
    assert "word/comments.xml" in names
    assert "commentRangeStart" in document_xml
    assert "差异" in comments_xml
    assert "另一份为" in comments_xml
    assert "本文件为" in comments_xml
    assert "差异类型" not in comments_xml
    assert "建议改为" not in comments_xml


def test_compare_docx_and_annotate_can_mark_original_document(tmp_path, monkeypatch):
    docx_module = pytest.importorskip("docx")

    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    original_path = tmp_path / "humanise!.docx"
    revised_path = tmp_path / "humanise!_revised.docx"

    original = docx_module.Document()
    original.add_paragraph("文件助手需要理解用户真正想要的结果，而不是只看关键词。")
    original.save(original_path)

    revised = docx_module.Document()
    revised.add_paragraph("文件助手需要理解用户真正想要的结果，而不是只看触发词。")
    revised.save(revised_path)

    result = task_tools.compare_docx_and_annotate(
        "humanise!.docx",
        "humanise!_revised.docx",
        target_path="humanise!.docx",
    )
    payload = json.loads(result)

    assert payload["success"] is True
    assert payload["path"] == "humanise!.docx"
    assert payload["annotations_added"] == 1

    with zipfile.ZipFile(original_path, "r") as archive:
        names = archive.namelist()
        original_document_xml = archive.read("word/document.xml").decode(
            "utf-8", errors="ignore"
        )
        original_comments_xml = archive.read("word/comments.xml").decode(
            "utf-8", errors="ignore"
        )
    with zipfile.ZipFile(revised_path, "r") as archive:
        revised_names = archive.namelist()

    assert "word/comments.xml" in names
    assert "word/comments.xml" not in revised_names
    assert "commentRangeStart" in original_document_xml
    assert "关键词" in original_document_xml
    assert "本文件为" in original_comments_xml
    assert "触发词" in original_comments_xml
    assert "建议改为" not in original_comments_xml


def test_plan_then_write_docx_comments_marks_original_document(tmp_path, monkeypatch):
    docx_module = pytest.importorskip("docx")

    import app.core.agent.task_tools as task_tools
    from app.core.agent.file_task_tool_catalog import parse_file_change

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    original_path = tmp_path / "contract_old.docx"
    revised_path = tmp_path / "contract_new.docx"

    original = docx_module.Document()
    original.add_paragraph("付款期限为30日。")
    original.add_paragraph("任一方提前30日通知可终止合同。")
    original.save(original_path)

    revised = docx_module.Document()
    revised.add_paragraph("付款期限为15日。")
    revised.add_paragraph("甲方提前7日通知可终止合同。")
    revised.save(revised_path)

    plan_payload = json.loads(
        task_tools.plan_docx_compare_annotations(
            "contract_old.docx",
            "contract_new.docx",
            target_path="contract_old.docx",
        )
    )
    assert plan_payload["success"] is True
    assert plan_payload["target_path"] == "contract_old.docx"
    assert plan_payload["annotation_candidates"]

    first_candidate = plan_payload["annotation_candidates"][0]
    comments = [
        {
            "原文片段": first_candidate["原文片段"],
            "批注内容": (
                "另一版为：付款期限为15日。\n"
                "本版为：付款期限为30日。\n"
                "风险：付款周期缩短，可能增加现金流压力。\n"
                "建议：确认是否接受该付款安排。"
            ),
        }
    ]
    result = task_tools.write_docx_comments(
        "contract_old.docx",
        comments_json=json.dumps(comments, ensure_ascii=False),
        source_path="contract_old.docx",
        compare_path="contract_new.docx",
        differences_detected=plan_payload["differences_detected"],
    )
    payload = json.loads(result)
    change = parse_file_change(
        "write_docx_comments",
        {
            "path": "contract_old.docx",
            "comments_json": json.dumps(comments, ensure_ascii=False),
        },
        result,
    )

    assert payload["success"] is True
    assert payload["operation"] == "write_docx_comments"
    assert payload["path"] == "contract_old.docx"
    assert payload["annotations_added"] == 1
    assert change["operation"] == "write_docx_comments"
    assert change["path"] == "contract_old.docx"

    with zipfile.ZipFile(original_path, "r") as archive:
        original_names = archive.namelist()
        original_comments_xml = archive.read("word/comments.xml").decode(
            "utf-8", errors="ignore"
        )
    with zipfile.ZipFile(revised_path, "r") as archive:
        revised_names = archive.namelist()

    assert "word/comments.xml" in original_names
    assert "word/comments.xml" not in revised_names
    assert "风险" in original_comments_xml
    assert "建议" in original_comments_xml


def test_write_docx_comments_accepts_anchor_text_alias(tmp_path, monkeypatch):
    docx_module = pytest.importorskip("docx")

    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    target_path = tmp_path / "contract_old.docx"
    document = docx_module.Document()
    document.add_paragraph("Payment is due within 30 days.")
    document.save(target_path)

    result = task_tools.write_docx_comments(
        "contract_old.docx",
        comments_json=[
            {
                "anchor_text": "30 days",
                "comment": "Other version: 15 days. Current version: 30 days.",
            }
        ],
        differences_detected=1,
    )
    payload = json.loads(result)

    assert payload["success"] is True
    assert payload["annotations_added"] == 1
    with zipfile.ZipFile(target_path, "r") as archive:
        assert "word/comments.xml" in archive.namelist()


def test_write_docx_comments_appends_after_existing_comments(tmp_path, monkeypatch):
    docx_module = pytest.importorskip("docx")

    import app.core.agent.task_tools as task_tools

    monkeypatch.setattr(task_tools, "_WORKSPACE_ROOT", str(tmp_path))

    docx_path = tmp_path / "contract_comments.docx"
    doc = docx_module.Document()
    doc.add_paragraph("付款期限为30日。")
    doc.add_paragraph("终止通知期为30日。")
    doc.save(docx_path)

    first = json.loads(
        task_tools.write_docx_comments(
            "contract_comments.docx",
            comments_json=[
                {"原文片段": "付款期限为30日", "批注内容": "旧批注"}
            ],
        )
    )
    second = json.loads(
        task_tools.write_docx_comments(
            "contract_comments.docx",
            comments_json=[
                {"原文片段": "终止通知期为30日", "批注内容": "新批注：风险与建议"}
            ],
        )
    )

    with zipfile.ZipFile(docx_path, "r") as archive:
        names = archive.namelist()
        comments_xml = archive.read("word/comments.xml").decode(
            "utf-8", errors="ignore"
        )
        document_xml = archive.read("word/document.xml").decode(
            "utf-8", errors="ignore"
        )

    assert first["annotations_added"] == 1
    assert second["annotations_added"] == 1
    assert names.count("word/comments.xml") == 1
    assert comments_xml.count("<w:comment ") == 2
    assert 'w:id="1"' in comments_xml
    assert 'w:id="2"' in comments_xml
    assert "新批注" in comments_xml
    assert document_xml.count("commentRangeStart") == 2


def test_normalize_docx_review_clear_scope_accepts_annotation_synonyms():
    import app.core.agent.task_tools as task_tools

    assert task_tools._normalize_docx_review_clear_scope("标注") == "comments"
    assert task_tools._normalize_docx_review_clear_scope("annotation") == "comments"


def test_annotate_file_pdf_docx_requirement_uses_bridge_streaming_tool_result(
    tmp_path, monkeypatch
):
    import app.core.agent.file_task_doc_annotate_bridge as bridge
    from app.core.agent.file_task_contract import (
        FileTaskToolStreamChunk,
        FileTaskToolStreamResult,
    )
    from app.core.agent.file_task_tool_gateway import (
        FileTaskToolContext,
        FileTaskToolGateway,
    )

    pdf_path = tmp_path / "source.pdf"
    pdf_path.write_bytes(b"%PDF-1.4\n")
    docx_path = tmp_path / "translation.docx"
    docx_path.write_bytes(b"PK\x03\x04")

    captured = {}

    def fake_stream_request_as_tool(request, *, workspace_root="", gemini_client=None):
        captured["request"] = request
        captured["workspace_root"] = workspace_root
        captured["gemini_client"] = gemini_client
        return FileTaskToolStreamResult(
            chunks=[
                FileTaskToolStreamChunk(
                    kind="event",
                    event_type="plan.confirmed",
                    payload={"summary": "按 3 批执行"},
                ),
                FileTaskToolStreamChunk(
                    kind="result",
                    payload={
                        "success": True,
                        "summary": "文件较大，已生成 3 批执行计划，等待确认开始第 1/3 批。",
                        "awaiting_confirmation": True,
                    },
                ),
            ]
        )

    monkeypatch.setattr(bridge, "stream_request_as_tool", fake_stream_request_as_tool)

    gateway = FileTaskToolGateway(
        context=FileTaskToolContext(
            task_files=[
                {"path": str(pdf_path), "type": "pdf"},
                {"path": str(docx_path), "type": "docx", "target": True},
            ],
            workspace_root=str(tmp_path),
            gemini_client="gemini-client",
            request_context={
                "task": "PDF是原文，docx文件是现有翻译稿。文件较大，请拆成多个分段来处理。",
                "target_path": str(docx_path),
                "options": {},
                "model_mode": "cloud",
                "model_id": "gemini-2.5-pro",
            },
        )
    )

    result = gateway.execute(
        "annotate_file",
        {"path": str(docx_path), "annotations": "[]", "requirement": "根据原文审校译稿并拆分执行"},
    )

    assert isinstance(result, FileTaskToolStreamResult)
    assert captured["workspace_root"] == str(tmp_path)
    assert captured["gemini_client"] == "gemini-client"
    assert captured["request"].target_path == str(docx_path)
    assert any(file_info.type == "pdf" for file_info in captured["request"].files)
    assert any(
        file_info.type == "docx" and file_info.target
        for file_info in captured["request"].files
    )


def test_file_snapshot_treats_missing_new_target_as_empty_snapshot(tmp_path):
    from app.core.file.multi_file_coordinator import FileSnapshot

    missing_path = tmp_path / "new_target.txt"
    snapshot = FileSnapshot.from_file(str(missing_path))

    assert snapshot.path == str(missing_path)
    assert snapshot.content == ""
    assert snapshot.content_hash
