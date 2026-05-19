import json
import stat


def test_add_pptx_slides_accepts_list_content_and_bullet_dicts(tmp_path):
    from pptx import Presentation

    from app.core.agent.task_tools import add_pptx_slides

    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "原始页"
    presentation.save(pptx_path)

    payload = json.loads(add_pptx_slides(
        str(pptx_path),
        slides=[
            {"title": "总结一", "content": ["市场需求明确", "替代成本是关键"]},
            {"title": {"text": "总结二"}, "bullets": [{"text": "本地文件交付"}, {"content": "高质量生成"}]},
            {"title": "总结三", "content": {"points": ["下一步做规格核验", "确认客户使用场景"]}},
        ],
    ))

    assert payload["success"] is True
    assert payload["slides_added"] == 3
    assert payload["total_slides"] == 4

    saved = Presentation(str(pptx_path))
    all_text = "\n".join(shape.text for slide in saved.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "总结一" in all_text
    assert "市场需求明确" in all_text
    assert "本地文件交付" in all_text
    assert "确认客户使用场景" in all_text


def test_design_pptx_theme_layout_applies_theme_without_changing_slide_count(tmp_path):
    from pptx import Presentation

    from app.core.agent.task_tools import design_pptx_theme_layout

    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    for title, body in [
        ("产品路线", "统一主题\n明确重点\n保留原有内容"),
        ("商业计划", "市场判断\n收入模型\n下一步行动"),
    ]:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = title
        slide.placeholders[1].text = body
    presentation.save(pptx_path)

    payload = json.loads(design_pptx_theme_layout(
        str(pptx_path),
        style_brief="科技感但适合商业 BP",
        density="compact",
    ))

    assert payload["success"] is True
    assert payload["operation"] == "design_pptx_theme_layout"
    assert payload["slides_designed"] == 2
    assert payload["total_slides"] == 2
    assert payload["theme_name"] == "科技深色"
    assert payload["layout_strategy"] == "safe_placeholder_grid"
    assert payload["text_shapes_styled"] >= 4

    saved = Presentation(str(pptx_path))
    assert len(saved.slides) == 2
    all_text = "\n".join(shape.text for slide in saved.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "产品路线" in all_text
    assert "下一步行动" in all_text
    assert any(
        shape.name == "KOTO_THEME_ACCENT_BAR"
        for slide in saved.slides
        for shape in slide.shapes
    )


def test_add_pptx_slides_clears_readonly_target_before_save(tmp_path):
    from pptx import Presentation

    from app.core.agent.task_tools import add_pptx_slides

    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "原始页"
    presentation.save(pptx_path)
    pptx_path.chmod(stat.S_IREAD)

    payload = json.loads(add_pptx_slides(
        str(pptx_path),
        slides=[{"title": "总结页", "content": ["第一点", "第二点"]}],
    ))

    assert payload["success"] is True
    assert payload["slides_added"] == 1
    assert "自动移除只读属性" in payload["warning"]

    saved = Presentation(str(pptx_path))
    assert len(saved.slides) == 2
    all_text = "\n".join(shape.text for slide in saved.slides for shape in slide.shapes if hasattr(shape, "text"))
    assert "总结页" in all_text
    assert "第一点" in all_text
