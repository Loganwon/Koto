# -*- coding: utf-8 -*-
"""
Integration tests for two PPTX fixes:

1. **Blank PPTX seed**: `_seed_new_file` now creates a title slide with
   placeholder shapes instead of an empty presentation with zero slides.

2. **Ctrl+S auto_save crash**: The `auto_save` endpoint for PPTX now uses
   `_apply_edits` (from `pptx_editor.py`) when it receives the rich geometry
   canvas format `{slides: [{shapes: [{paragraphs: [...]}]}]}`, matching
   the `save_file` / download path.  Previously it called the old file-parser
   PPTX exporter which
   expected a flat `"text"` key, causing:
     - `'str' object has no attribute 'get'` crash when iterating slides, OR
     - silent text loss because `shape_entry.get("text", "")` returned `""`.
"""

from __future__ import annotations

import io
import os
import uuid
from pathlib import Path

import pytest

# ── Helpers ──────────────────────────────────────────────────────────────────


def _make_pptx_with_text(text: str = "Hello World") -> bytes:
    """Create a minimal PPTX with one slide containing a title text box."""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    # Title placeholder
    title = slide.placeholders[0]
    title.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pptx_with_rich_text() -> bytes:
    """Create one text box with two paragraphs and multiple formatted runs."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    first = shape.text_frame.paragraphs[0]
    run = first.add_run()
    run.text = "Keep "
    run.font.bold = True
    run.font.size = Pt(20)
    first.add_run().text = "DELETE"
    second = shape.text_frame.add_paragraph()
    second.add_run().text = "Remove this paragraph"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _parse_pptx_geometry_data(pptx_bytes: bytes) -> dict:
    """Parse PPTX bytes via parse_pptx_geometry (the backend parser)."""
    import tempfile

    from app.core.file.file_parser import parse_pptx_geometry

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(pptx_bytes)
        f.flush()
        data = parse_pptx_geometry(f.name)
    os.unlink(f.name)
    return data


def _read_pptx_texts(pptx_bytes: bytes) -> list[str]:
    """Read all text from a PPTX file's shapes, returns list of shape texts."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(pptx_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return texts


# ── Fixture: Flask test client ───────────────────────────────────────────────


@pytest.fixture(scope="module")
def wa_client(tmp_path_factory):
    """Flask test client backed by temporary workspace/tmp directories."""
    os.environ.setdefault("KOTO_AUTH_ENABLED", "false")

    tmp_root = tmp_path_factory.mktemp("pptx_save_root")
    tmp_dir = tmp_root / "tmp"
    workspace_dir = tmp_root / "workspace"
    tmp_dir.mkdir(parents=True)
    workspace_dir.mkdir(parents=True)

    import web.blueprints.workspace_assistant as _wa
    import web.shared as _shared

    original_tmp = _wa._TMP_DIR
    original_ws = getattr(_shared, "WORKSPACE_DIR", None)
    _wa._TMP_DIR = tmp_dir
    _shared.WORKSPACE_DIR = str(workspace_dir)

    from flask import Flask

    from web.blueprints.pptx_editor import pptx_editor_bp
    from web.blueprints.workspace_assistant import workspace_assistant_bp

    app = Flask(__name__)
    app.register_blueprint(workspace_assistant_bp)
    app.register_blueprint(pptx_editor_bp)
    app.config["TESTING"] = True

    with app.test_client() as client:
        yield client, tmp_dir, workspace_dir

    _wa._TMP_DIR = original_tmp
    if original_ws is not None:
        _shared.WORKSPACE_DIR = original_ws


# ═════════════════════════════════════════════════════════════════════════════


class TestPptxTextDeletionExport:
    """Deleted paragraphs and runs must not come back after PPTX export."""

    def _edited_payload(self) -> tuple[bytes, dict, dict]:
        raw = _make_pptx_with_rich_text()
        data = _parse_pptx_geometry_data(raw)
        shape = next(
            shape
            for shape in data["slides"][0]["shapes"]
            if shape.get("has_text")
            and "DELETE" in "".join(
                str(run.get("text", ""))
                for paragraph in shape.get("paragraphs", [])
                for run in paragraph.get("runs", [])
            )
        )
        return raw, data, shape

    def test_deleting_complete_run_and_paragraph_is_exact(self):
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN
        from web.blueprints.pptx_editor import _apply_edits

        raw, data, shape = self._edited_payload()
        kept_run = dict(shape["paragraphs"][0]["runs"][0])
        kept_run["text"] = "Only this remains"
        shape["paragraphs"] = [{"align": "RIGHT", "runs": [kept_run]}]

        exported = _apply_edits(raw, data["slides"])
        prs = Presentation(io.BytesIO(exported))
        text_shape = next(item for item in prs.slides[0].shapes if item.has_text_frame)

        assert text_shape.text == "Only this remains"
        assert len(text_shape.text_frame.paragraphs) == 1
        assert len(text_shape.text_frame.paragraphs[0].runs) == 1
        assert text_shape.text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT
        assert text_shape.text_frame.paragraphs[0].runs[0].font.bold is True

    def test_deleting_all_text_produces_an_empty_text_frame(self):
        from pptx import Presentation
        from web.blueprints.pptx_editor import _apply_edits

        raw, data, shape = self._edited_payload()
        shape["paragraphs"] = []

        exported = _apply_edits(raw, data["slides"])
        prs = Presentation(io.BytesIO(exported))
        text_shape = next(item for item in prs.slides[0].shapes if item.has_text_frame)

        assert text_shape.text == ""
        assert len(text_shape.text_frame.paragraphs) == 1

    def test_new_runs_are_added_instead_of_silently_dropped(self):
        from pptx import Presentation
        from web.blueprints.pptx_editor import _apply_edits

        raw, data, shape = self._edited_payload()
        template = dict(shape["paragraphs"][0]["runs"][0])
        shape["paragraphs"] = [{
            "align": "LEFT",
            "runs": [
                {**template, "text": "A"},
                {**template, "text": "B", "bold": False},
                {**template, "text": "C"},
            ],
        }]

        exported = _apply_edits(raw, data["slides"])
        prs = Presentation(io.BytesIO(exported))
        text_shape = next(item for item in prs.slides[0].shapes if item.has_text_frame)

        assert [run.text for run in text_shape.text_frame.paragraphs[0].runs] == [
            "A", "B", "C"
        ]
# A. Blank PPTX seed fix
# ═════════════════════════════════════════════════════════════════════════════


class TestSeedNewFilePptx:
    """_seed_new_file must create a PPTX with a title slide, not a blank file."""

    def test_seeded_pptx_has_at_least_one_slide(self, wa_client):
        _, _, workspace_dir = wa_client
        from web.blueprints.workspace_assistant import _seed_new_file

        target = workspace_dir / "seed_test.pptx"
        _seed_new_file(target)
        assert target.stat().st_size > 0

        from pptx import Presentation

        prs = Presentation(str(target))
        assert len(prs.slides) >= 1, "Seeded PPTX must have at least one slide"

    def test_seeded_pptx_title_slide_has_placeholders(self, wa_client):
        """Seeded PPTX must have title + content shapes with matching formatting."""
        _, _, workspace_dir = wa_client
        from web.blueprints.workspace_assistant import _seed_new_file

        target = workspace_dir / "seed_placeholder.pptx"
        _seed_new_file(target)

        from pptx import Presentation

        prs = Presentation(str(target))
        slide = prs.slides[0]
        # Must have at least 2 text shapes (title + content)
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        assert (
            len(text_shapes) >= 2
        ), "Seeded slide must have at least title + content shapes"
        # Verify shapes contain visible text runs (not empty placeholders)
        texts = []
        for s in text_shapes:
            for p in s.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        texts.append(r.text)
        assert (
            len(texts) >= 2
        ), f"Seeded slide must have visible text in both shapes, got: {texts}"

    def test_seeded_pptx_parseable_by_geometry_parser(self, wa_client):
        """parse_pptx_geometry must succeed on a seeded PPTX."""
        _, _, workspace_dir = wa_client
        from web.blueprints.workspace_assistant import _seed_new_file

        target = workspace_dir / "seed_parse.pptx"
        _seed_new_file(target)

        from app.core.file.file_parser import parse_pptx_geometry

        data = parse_pptx_geometry(str(target))
        assert "slides" in data
        assert len(data["slides"]) >= 1

    def test_create_file_endpoint_creates_nonblank_pptx(self, wa_client):
        """POST /api/v1/workspace/create_file with .pptx extension."""
        client, _, workspace_dir = wa_client
        resp = client.post(
            "/api/v1/workspace/create_file",
            json={"folder": "", "name": "new_presentation.pptx"},
        )
        assert resp.status_code == 200
        body = resp.get_json()
        assert body["ok"] is True

        pptx_path = workspace_dir / "new_presentation.pptx"
        assert pptx_path.exists()
        assert pptx_path.stat().st_size > 0

        from pptx import Presentation

        prs = Presentation(str(pptx_path))
        assert len(prs.slides) >= 1, "Created PPTX must have a title slide"

    def test_open_seeded_pptx_returns_shapes(self, wa_client):
        """Opening a seeded PPTX via open_file_by_path should return shapes."""
        client, _, workspace_dir = wa_client
        from web.blueprints.workspace_assistant import _seed_new_file

        target = workspace_dir / "seed_open_test.pptx"
        _seed_new_file(target)

        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "seed_open_test.pptx"},
        )
        assert resp.status_code == 200
        body = resp.get_json()
        assert body["file_type"] == "pptx"
        slides = body["data"]["slides"]
        assert len(slides) >= 1, "Parsed seeded PPTX must have at least 1 slide"
        # Title slide should have shapes (title + subtitle placeholders)
        shapes = slides[0].get("shapes", [])
        assert len(shapes) >= 1, "Title slide must have at least one shape"


# ═════════════════════════════════════════════════════════════════════════════
# B. PPTX auto_save with rich paragraph/run format
# ═════════════════════════════════════════════════════════════════════════════


class TestPptxAutoSaveRichFormat:
    """
    auto_save must handle the geometry canvas format with paragraphs/runs.

    Previously, the old file-parser PPTX exporter only looked for shape_entry.get("text", ""),
    which returned "" because the frontend sends paragraphs with runs.
    Now auto_save uses _apply_edits for the rich format.
    """

    def _upload_and_get_id(self, client) -> tuple[str, dict]:
        """Upload a PPTX, return (file_id, parsed_data)."""
        pptx_bytes = _make_pptx_with_text("Original Title")
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx_bytes), "test_save.pptx")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 200, f"Upload failed: {resp.get_json()}"
        body = resp.get_json()
        return body["file_id"], body["data"]

    def test_auto_save_with_rich_data_returns_200(self, wa_client):
        """auto_save must not crash with rich paragraph/run data."""
        client, _, _ = wa_client
        file_id, data = self._upload_and_get_id(client)

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": file_id,
                "ws_source_path": "",
                "explicit": False,
                "data": data,
            },
        )
        assert resp.status_code == 200, f"auto_save crashed: {resp.get_json()}"
        assert resp.get_json()["ok"] is True

    def test_auto_save_preserves_text_content(self, wa_client):
        """Text in shapes must survive a round-trip through auto_save."""
        client, _, _ = wa_client
        file_id, data = self._upload_and_get_id(client)

        # Verify the parsed data has shapes with text
        slides = data.get("slides", [])
        assert len(slides) >= 1
        title_shape = None
        for shape in slides[0].get("shapes", []):
            if shape.get("has_text") and shape.get("paragraphs"):
                for para in shape["paragraphs"]:
                    for run in para.get("runs", []):
                        if run.get("text", "").strip():
                            title_shape = shape
                            break
        assert title_shape is not None, "Expected at least one shape with text"

        # Save with the parsed data (simulates frontend serialize → POST)
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": file_id,
                "explicit": True,
                "data": data,
            },
        )
        assert resp.status_code == 200

        # Read back the saved PPTX via /raw/ and verify text survived
        raw_resp = client.get(f"/api/v1/workspace/raw/{file_id}")
        assert raw_resp.status_code == 200
        saved_texts = _read_pptx_texts(raw_resp.data)
        all_text = " ".join(saved_texts)
        assert (
            "Original Title" in all_text
        ), f"Text lost after auto_save! Got: {saved_texts}"

    def test_auto_save_explicit_writes_workspace_file(self, wa_client):
        """explicit=True with ws_source_path must write to workspace dir."""
        client, _, workspace_dir = wa_client
        file_id, data = self._upload_and_get_id(client)

        ws_filename = "explicit_save_test.pptx"
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": file_id,
                "ws_source_path": ws_filename,
                "explicit": True,
                "data": data,
            },
        )
        assert resp.status_code == 200
        body = resp.get_json()
        assert body.get("src_written") is True

        ws_file = workspace_dir / ws_filename
        assert ws_file.exists(), "Workspace file must be created on explicit save"
        assert ws_file.stat().st_size > 0

    def test_auto_save_with_edited_text(self, wa_client):
        """Editing text in the data and saving must persist the change."""
        client, _, _ = wa_client
        file_id, data = self._upload_and_get_id(client)

        # Modify the first text shape's text
        slides = data.get("slides", [])
        modified = False
        for shape in slides[0].get("shapes", []):
            if shape.get("has_text") and shape.get("paragraphs"):
                for para in shape["paragraphs"]:
                    for run in para.get("runs", []):
                        run["text"] = "Modified Title Text"
                        modified = True
                        break
                    if modified:
                        break
            if modified:
                break

        assert modified, "Should have found a text shape to modify"

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": file_id,
                "explicit": True,
                "data": data,
            },
        )
        assert resp.status_code == 200

        # Read back and verify the edit was saved
        raw_resp = client.get(f"/api/v1/workspace/raw/{file_id}")
        saved_texts = _read_pptx_texts(raw_resp.data)
        all_text = " ".join(saved_texts)
        assert (
            "Modified Title Text" in all_text
        ), f"Edited text not persisted! Got: {saved_texts}"

    def test_auto_save_with_empty_slides(self, wa_client):
        """auto_save must handle an empty slides array gracefully."""
        client, _, _ = wa_client
        file_id, _ = self._upload_and_get_id(client)

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": file_id,
                "data": {
                    "slideWidthEmu": 9144000,
                    "slideHeightEmu": 6858000,
                    "slides": [],
                },
            },
        )
        assert resp.status_code == 200

    def test_auto_save_rejects_legacy_text_card_payload(self, wa_client):
        """Workspace PPTX save must not silently accept the retired text-card format."""
        client, _, _ = wa_client
        file_id, data = self._upload_and_get_id(client)
        shape_id = None
        for shape in data["slides"][0].get("shapes", []):
            if shape.get("has_text"):
                shape_id = shape.get("id")
                break
        assert shape_id is not None

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": file_id,
                "data": [
                    {
                        "slide_index": 0,
                        "texts": [{"shape_id": shape_id, "text": "Legacy Text"}],
                    }
                ],
            },
        )
        assert resp.status_code == 400
        assert "rich slides" in resp.get_json()["error"]

    def test_auto_save_consecutive_saves_produce_different_bytes(self, wa_client):
        """Multiple saves with different content must produce different bytes."""
        client, _, _ = wa_client
        file_id, data = self._upload_and_get_id(client)

        # First save with original data
        client.post(
            "/api/v1/workspace/auto_save",
            json={"file_type": "pptx", "file_id": file_id, "data": data},
        )
        raw1 = client.get(f"/api/v1/workspace/raw/{file_id}").data

        # Second save with modified data
        slides = data.get("slides", [])
        for shape in slides[0].get("shapes", []):
            if shape.get("has_text") and shape.get("paragraphs"):
                for para in shape["paragraphs"]:
                    for run in para.get("runs", []):
                        run["text"] = "Completely Different Text For Save 2"
                        break
                break

        client.post(
            "/api/v1/workspace/auto_save",
            json={"file_type": "pptx", "file_id": file_id, "data": data},
        )
        raw2 = client.get(f"/api/v1/workspace/raw/{file_id}").data

        assert (
            raw1 != raw2
        ), "Consecutive saves with different text must produce different bytes"

    def test_auto_save_outside_workspace_returns_200_not_403(self, wa_client):
        """
        When ws_source_path is an absolute path outside the workspace (e.g. a
        file opened from Downloads via openBrowserFile), auto_save must NOT
        return 403.  It should save to tmp and return 200.
        If the absolute path's parent directory exists, it also writes back
        to the original file (src_written=True).
        """
        client, tmp_dir, workspace_dir = wa_client

        # Create a real external directory so the write-back can succeed
        import tempfile

        ext_dir = Path(tempfile.mkdtemp(prefix="koto_ext_"))
        ext_file = ext_dir / "external.pptx"
        ext_file.write_bytes(_make_pptx_with_text("External Original"))

        file_id, data = self._upload_and_get_id(client)

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": file_id,
                "ws_source_path": str(ext_file),
                "explicit": True,
                "data": data,
            },
        )
        assert resp.status_code == 200, (
            f"auto_save must not 403 for external paths! Got {resp.status_code}: "
            f"{resp.get_json()}"
        )
        body = resp.get_json()
        assert body["ok"] is True
        # The write-back to the external file should succeed
        assert body["src_written"] is True
        # The external file should have been updated
        assert ext_file.stat().st_size > 0

        # Cleanup
        import shutil

        shutil.rmtree(ext_dir, ignore_errors=True)

    def test_auto_save_outside_workspace_still_writes_tmp(self, wa_client):
        """
        Even when ws_source_path is outside workspace, the tmp file must be
        written so /raw/<file_id> still works.
        """
        import platform

        client, tmp_dir, workspace_dir = wa_client
        file_id, data = self._upload_and_get_id(client)

        # Use a non-existent parent dir so write-back is skipped,
        # but tmp must still be written.  Path must be *absolute* so
        # it doesn't hit the relative-path 403 guard.
        if platform.system() == "Windows":
            outside_path = r"C:\no_such_dir_xyz\report.pptx"
        else:
            outside_path = "/some/nonexistent/dir/report.pptx"
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": file_id,
                "ws_source_path": outside_path,
                "explicit": True,
                "data": data,
            },
        )
        assert resp.status_code == 200

        # Verify the tmp file exists and is non-empty
        raw_resp = client.get(f"/api/v1/workspace/raw/{file_id}")
        assert raw_resp.status_code == 200
        assert len(raw_resp.data) > 0

    def test_auto_save_path_traversal_still_blocked_for_relative(self, wa_client):
        """
        Relative traversal paths like '../../etc/passwd' must be hard-rejected
        with a 403 — they are never valid external source paths.
        """
        client, _, _ = wa_client
        file_id, data = self._upload_and_get_id(client)

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": file_id,
                "ws_source_path": "../../etc/passwd.pptx",
                "explicit": True,
                "data": data,
            },
        )
        assert (
            resp.status_code == 403
        ), f"Relative traversal must be blocked with 403, got {resp.status_code}"


# ═════════════════════════════════════════════════════════════════════════════
