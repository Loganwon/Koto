# -*- coding: utf-8 -*-
"""
End-to-end regression tests that simulate real user workflows.

Each test replays the *exact HTTP call sequence* the browser makes and
asserts on the *exact invariant* that was broken by a real bug.  If a
future change re-introduces any of these bugs, the corresponding test
will fail with a message describing what went wrong.

Bug history these tests are designed to catch:

  BUG-1  external absolute paths must use open_abs_file instead of the
         workspace-relative open_file_by_path route.

  BUG-2  auto_save with explicit=True and an absolute ws_source_path
         (external file) silently skipped the write-back to the original
         file, so closing and reopening showed stale content.

  BUG-3  auto_save called the old file-parser PPTX exporter for the rich geometry format
         (paragraphs/runs) instead of _apply_edits(), causing either a
         crash or silent text loss.

  BUG-4  _seed_new_file created a PPTX with zero slides (bare
         Presentation()), so opening the seeded file returned an empty
         slides array and the canvas was blank.

  BUG-5  the unsafe absolute-byte endpoint was retired; absolute files now
         go through parsed open_abs_file.

  BUG-6  Timer-based (non-explicit) auto_save must NOT touch the
         workspace/source file; only Ctrl+S (explicit=True) should.

  BUG-7  Two files open at the same time could have their save data
         cross-contaminated if file_id handling is wrong.
"""

from __future__ import annotations

import io
import os
from pathlib import Path

import pytest

# ═══════════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════════


def _make_pptx_bytes(title: str = "Hello") -> bytes:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.placeholders[0].text = title
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_docx_bytes(text: str = "Hello World") -> bytes:
    from docx import Document

    doc = Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_all_texts(raw: bytes) -> list[str]:
    """Return every shape's full text from raw PPTX bytes."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return out


def _pptx_slide_count(raw: bytes) -> int:
    from pptx import Presentation

    return len(Presentation(io.BytesIO(raw)).slides)


def _runs_text(data: dict) -> str:
    """Concatenate all run texts from parsed geometry data."""
    parts = []
    for slide in data.get("slides", []):
        for shape in slide.get("shapes", []):
            for para in shape.get("paragraphs", []):
                for run in para.get("runs", []):
                    t = run.get("text", "")
                    if t.strip():
                        parts.append(t)
    return " ".join(parts)


def _set_first_run(data: dict, new_text: str) -> bool:
    """Overwrite the first non-empty run.  Returns False if none found."""
    for slide in data.get("slides", []):
        for shape in slide.get("shapes", []):
            for para in shape.get("paragraphs", []):
                for run in para.get("runs", []):
                    if run.get("text", "").strip():
                        run["text"] = new_text
                        return True
    return False


# ═══════════════════════════════════════════════════════════════════════════════
# Fixture — isolated Flask test client with fake workspace + external dir
# ═══════════════════════════════════════════════════════════════════════════════


@pytest.fixture(scope="module")
def env(tmp_path_factory):
    """(client, tmp_dir, workspace_dir, external_dir)"""
    os.environ.setdefault("KOTO_AUTH_ENABLED", "false")

    root = tmp_path_factory.mktemp("wf")
    tmp_dir = root / "tmp"
    ws_dir = root / "workspace"
    ext_dir = root / "external"  # stands for C:\Users\…\Downloads
    for d in (tmp_dir, ws_dir, ext_dir):
        d.mkdir()

    import web.blueprints.workspace_assistant as _wa
    import web.shared as _shared

    orig_tmp = _wa._TMP_DIR
    orig_ws = getattr(_shared, "WORKSPACE_DIR", None)
    _wa._TMP_DIR = tmp_dir
    _shared.WORKSPACE_DIR = str(ws_dir)

    from flask import Flask

    from web.blueprints.pptx_editor import pptx_editor_bp
    from web.blueprints.workspace_assistant import workspace_assistant_bp

    app = Flask(__name__)
    app.register_blueprint(workspace_assistant_bp)
    app.register_blueprint(pptx_editor_bp)
    app.config["TESTING"] = True

    with app.test_client() as client:
        yield client, tmp_dir, ws_dir, ext_dir

    _wa._TMP_DIR = orig_tmp
    if orig_ws is not None:
        _shared.WORKSPACE_DIR = orig_ws


# ═══════════════════════════════════════════════════════════════════════════════
# 1.  BUG-1 — external absolute paths use open_abs_file
# ═══════════════════════════════════════════════════════════════════════════════


class TestBug1_OpenExternalPath:
    """
    BEFORE FIX:  external files were routed through workspace-relative
    open_file_by_path or raw-byte serving flows, which muddied source-path
    tracking.

    AFTER FIX:  open_abs_file handles absolute paths whose parent exists and
    extension is allowed.
    """

    def test_absolute_external_path_returns_200(self, env):
        """Regression: this returned 403 before the fix."""
        client, _, _, ext_dir = env
        f = ext_dir / "from_downloads.pptx"
        f.write_bytes(_make_pptx_bytes("External"))

        resp = client.post(
            "/api/v1/workspace/open_abs_file",
            json={"path": str(f)},
        )
        assert resp.status_code == 200, (
            f"BUG-1 regression: open_abs_file must accept absolute "
            f"external paths, but returned {resp.status_code}: "
            f"{resp.get_json()}"
        )

    def test_external_file_content_is_actually_parsed(self, env):
        """Not just 200 — the returned data must contain the file's text."""
        client, _, _, ext_dir = env
        f = ext_dir / "content_check.pptx"
        f.write_bytes(_make_pptx_bytes("AlphaContent"))

        body = client.post(
            "/api/v1/workspace/open_abs_file",
            json={"path": str(f)},
        ).get_json()
        text = _runs_text(body["data"])
        assert "AlphaContent" in text, (
            f"BUG-1 regression: external file opened but content not parsed. "
            f"Got: {text!r}"
        )

    def test_traversal_path_still_blocked(self, env):
        """Security: ../../etc/passwd must NOT slip through."""
        client, _, _, _ = env
        resp = client.post(
            "/api/v1/workspace/open_abs_file",
            json={"path": "../../etc/passwd"},
        )
        assert resp.status_code in (
            403,
            404,
        ), f"Traversal not blocked — got {resp.status_code}"

    def test_bad_extension_blocked_even_if_path_exists(self, env):
        """Absolute path with disallowed extension must be rejected."""
        client, _, _, ext_dir = env
        bad = ext_dir / "virus.exe"
        bad.write_bytes(b"MZ")
        resp = client.post(
            "/api/v1/workspace/open_abs_file",
            json={"path": str(bad)},
        )
        assert resp.status_code in (400, 403)

    def test_nonexistent_external_file_returns_404(self, env):
        client, _, _, ext_dir = env
        resp = client.post(
            "/api/v1/workspace/open_abs_file",
            json={"path": str(ext_dir / "ghost.pptx")},
        )
        assert resp.status_code == 404


# ═══════════════════════════════════════════════════════════════════════════════
# 2.  BUG-2 — auto_save must write back to external files on Ctrl+S
# ═══════════════════════════════════════════════════════════════════════════════


class TestBug2_ExternalSaveWriteBack:
    """
    BEFORE FIX:  auto_save with ws_source_path pointing outside the
    workspace silently set src_path=None and only wrote to tmp.
    The original file on disk was never updated.

    AFTER FIX:  If the path is an absolute path with an existing parent
    and an allowed extension, auto_save writes back to it.
    """

    def _open_external(self, client, path: Path):
        resp = client.post(
            "/api/v1/workspace/open_abs_file",
            json={"path": str(path)},
        )
        assert resp.status_code == 200
        body = resp.get_json()
        return body["file_id"], body["data"]

    def test_ctrl_s_updates_external_file_on_disk(self, env):
        """The #1 user-reported bug: Ctrl+S, close, reopen → content gone."""
        client, _, _, ext_dir = env
        f = ext_dir / "save_check.pptx"
        f.write_bytes(_make_pptx_bytes("Before Edit"))

        fid, data = self._open_external(client, f)
        assert _set_first_run(data, "After Edit")

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": fid,
                "ws_source_path": str(f),
                "explicit": True,
                "data": data,
            },
        )
        assert resp.status_code == 200
        body = resp.get_json()
        assert body["src_written"] is True, (
            "BUG-2 regression: auto_save returned ok but did NOT write "
            "back to the original external file (src_written=False)"
        )

        # Verify the actual bytes on disk changed
        disk_texts = _pptx_all_texts(f.read_bytes())
        assert "After Edit" in " ".join(disk_texts), (
            f"BUG-2 regression: src_written=True but the file on disk "
            f"still has the old content: {disk_texts}"
        )

    def test_full_cycle_open_edit_save_close_reopen(self, env):
        """
        The exact sequence: open → edit → Ctrl+S → close tab → reopen.
        This is the complete reproduction of the reported bug.
        """
        client, _, _, ext_dir = env
        f = ext_dir / "full_cycle.pptx"
        f.write_bytes(_make_pptx_bytes("Original"))
        original_bytes = f.read_bytes()

        # Step 1: Open
        fid1, data1 = self._open_external(client, f)
        assert "Original" in _runs_text(data1)

        # Step 2: Edit + Ctrl+S
        assert _set_first_run(data1, "Edited By User")
        client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": fid1,
                "ws_source_path": str(f),
                "explicit": True,
                "data": data1,
            },
        )

        # Step 3: Close tab (nothing to do server-side)

        # Verify file on disk actually changed
        assert (
            f.read_bytes() != original_bytes
        ), "BUG-2: file on disk unchanged after Ctrl+S"

        # Step 4: Reopen (new file_id, reads from disk)
        fid2, data2 = self._open_external(client, f)
        assert fid2 != fid1, "Reopen should produce a new file_id"
        text2 = _runs_text(data2)
        assert (
            "Edited By User" in text2
        ), f"BUG-2: content lost after close → reopen.  Got: {text2!r}"

    def test_nonexistent_parent_dir_skips_writeback_gracefully(self, env):
        """Path with non-existent parent → skip write, but don't crash."""
        import platform

        client, _, _, _ = env
        pptx = _make_pptx_bytes("Skip")
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx), "skip.pptx")},
            content_type="multipart/form-data",
        )
        fid = resp.get_json()["file_id"]
        data = resp.get_json()["data"]

        # Must be *absolute* (so it doesn't hit the relative-path 403 guard)
        # but with a parent directory that doesn't exist.
        if platform.system() == "Windows":
            bogus_path = r"C:\no_such_dir_xyz\file.pptx"
        else:
            bogus_path = "/no/such/dir/file.pptx"

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": fid,
                "ws_source_path": bogus_path,
                "explicit": True,
                "data": data,
            },
        )
        assert resp.status_code == 200
        assert (
            resp.get_json()["src_written"] is False
        ), "Write-back should be skipped for non-existent parent dir"


# ═══════════════════════════════════════════════════════════════════════════════
# 3.  BUG-3 — auto_save must use _apply_edits for rich PPTX format
# ═══════════════════════════════════════════════════════════════════════════════


class TestBug3_RichFormatAutoSave:
    """
    BEFORE FIX:  auto_save always called the old file-parser PPTX exporter, which iterated
    shapes looking for shape.get("text", "").  The frontend sends the
    rich format {paragraphs: [{runs: [{text:…}]}]}, so "text" key was
    absent → all text silently dropped, or crash.

    AFTER FIX:  auto_save detects the rich format (data has "slides"
    key) and uses _apply_edits from pptx_editor.py.
    """

    def _open_pptx_via_upload(self, client, title="Test"):
        pptx = _make_pptx_bytes(title)
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx), "test.pptx")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 200
        body = resp.get_json()
        return body["file_id"], body["data"]

    def test_rich_format_does_not_crash(self, env):
        """Before fix: 500 error or 'str has no attribute get'."""
        client, _, _, _ = env
        fid, data = self._open_pptx_via_upload(client, "NoCrash")

        # data already has the rich format from parse_pptx_geometry
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={"file_type": "pptx", "file_id": fid, "data": data},
        )
        assert resp.status_code == 200, (
            f"BUG-3 regression: auto_save crashed on rich format: " f"{resp.get_json()}"
        )

    def test_rich_format_preserves_text(self, env):
        """Before fix: text silently dropped because the old file-parser PPTX exporter
        used shape.get('text','') which was empty."""
        client, _, _, _ = env
        fid, data = self._open_pptx_via_upload(client, "Preserve Me")

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={"file_type": "pptx", "file_id": fid, "data": data},
        )
        assert resp.status_code == 200

        # Read the saved bytes and check the text survived
        raw_resp = client.get(f"/api/v1/workspace/raw/{fid}")
        assert raw_resp.status_code == 200
        texts = _pptx_all_texts(raw_resp.data)
        assert "Preserve Me" in " ".join(texts), (
            f"BUG-3 regression: text lost after auto_save with rich format. "
            f"Shapes contain: {texts}"
        )

    def test_edited_text_in_rich_format_persists(self, env):
        """Edit a run text in rich format, save, verify the edit stuck."""
        client, _, _, _ = env
        fid, data = self._open_pptx_via_upload(client, "Original")

        assert _set_first_run(data, "Changed Text")
        client.post(
            "/api/v1/workspace/auto_save",
            json={"file_type": "pptx", "file_id": fid, "data": data},
        )

        texts = _pptx_all_texts(client.get(f"/api/v1/workspace/raw/{fid}").data)
        assert "Changed Text" in " ".join(
            texts
        ), f"BUG-3: edited text not written to PPTX.  Got: {texts}"
        assert "Original" not in " ".join(
            texts
        ), f"BUG-3: old text still present after edit.  Got: {texts}"


# ═══════════════════════════════════════════════════════════════════════════════
# 4.  BUG-4 — _seed_new_file must create a non-empty, openable PPTX
# ═══════════════════════════════════════════════════════════════════════════════


class TestBug4_SeedNewFile:
    """
    BEFORE FIX:  _seed_new_file called Presentation() and saved with
    zero slides.  parse_pptx_geometry returned {slides: []}, the canvas
    was blank, and the user saw an empty screen after creating a file.

    AFTER FIX:  Seeded PPTX has at least one slide with title + content
    shapes, matching the JS pptxAddSlide() layout.
    """

    def test_seeded_pptx_has_at_least_one_slide(self, env):
        _, _, ws_dir, _ = env
        from web.blueprints.workspace_assistant import _seed_new_file

        f = ws_dir / "seed_slide_count.pptx"
        _seed_new_file(f)
        count = _pptx_slide_count(f.read_bytes())
        assert (
            count >= 1
        ), f"BUG-4 regression: seeded PPTX has {count} slides (need ≥ 1)"

    def test_seeded_pptx_has_text_shapes(self, env):
        """Blank Presentation() has 0 shapes — we need at least title+body."""
        _, _, ws_dir, _ = env
        from web.blueprints.workspace_assistant import _seed_new_file

        f = ws_dir / "seed_shapes.pptx"
        _seed_new_file(f)
        texts = _pptx_all_texts(f.read_bytes())
        assert len(texts) >= 2, (
            f"BUG-4 regression: seeded PPTX needs title + content shapes. "
            f"Got {len(texts)} text shapes: {texts}"
        )

    def test_create_then_open_returns_non_empty_slides(self, env):
        """Full flow: create_file → open_file_by_path → slides non-empty."""
        client, _, _, _ = env
        client.post(
            "/api/v1/workspace/create_file",
            json={"folder": "", "name": "seed_open.pptx"},
        )
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "seed_open.pptx"},
        )
        assert resp.status_code == 200
        slides = resp.get_json()["data"]["slides"]
        assert len(slides) >= 1, "BUG-4: created PPTX opens with zero slides"
        # Must have shapes with text (not blank placeholders)
        shapes = slides[0].get("shapes", [])
        assert any(
            s.get("has_text") for s in shapes
        ), "BUG-4: first slide has no text shapes"

    def test_zero_byte_pptx_auto_repaired_on_open(self, env):
        """Legacy 0-byte files should be auto-seeded when opened."""
        client, _, ws_dir, _ = env
        f = ws_dir / "zero_byte.pptx"
        f.write_bytes(b"")

        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "zero_byte.pptx"},
        )
        assert resp.status_code == 200, (
            f"BUG-4: 0-byte PPTX should auto-repair, got "
            f"{resp.status_code}: {resp.get_json()}"
        )
        assert len(resp.get_json()["data"]["slides"]) >= 1


# ═══════════════════════════════════════════════════════════════════════════════
# 5.  BUG-6 — timer auto_save must NOT write to workspace / source
# ═══════════════════════════════════════════════════════════════════════════════


class TestBug6_TimerVsExplicitSave:
    """
    Timer-fired auto_save (explicit=False) must ONLY update the tmp file.
    If it wrote to the workspace file, it would overwrite content even
    when the user hasn't deliberately pressed Ctrl+S.
    """

    def test_non_explicit_save_does_not_touch_workspace_file(self, env):
        client, _, ws_dir, _ = env
        name = "timer_guard.pptx"
        client.post(
            "/api/v1/workspace/create_file",
            json={"folder": "", "name": name},
        )
        original_bytes = (ws_dir / name).read_bytes()

        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": name},
        )
        fid = resp.get_json()["file_id"]
        data = resp.get_json()["data"]

        _set_first_run(data, "Timer Sneaky Edit")
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": fid,
                "ws_source_path": name,
                "explicit": False,  # timer, not Ctrl+S
                "data": data,
            },
        )
        assert resp.status_code == 200
        assert resp.get_json()["src_written"] is False

        assert (
            ws_dir / name
        ).read_bytes() == original_bytes, (
            "BUG-6: timer auto_save modified the workspace file!"
        )

    def test_explicit_save_does_update_workspace_file(self, env):
        """Contrast: Ctrl+S (explicit=True) MUST update the workspace file."""
        client, _, ws_dir, _ = env
        name = "explicit_guard.pptx"
        client.post(
            "/api/v1/workspace/create_file",
            json={"folder": "", "name": name},
        )
        original_bytes = (ws_dir / name).read_bytes()

        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": name},
        )
        fid = resp.get_json()["file_id"]
        data = resp.get_json()["data"]

        _set_first_run(data, "Explicit Save Edit")
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "pptx",
                "file_id": fid,
                "ws_source_path": name,
                "explicit": True,  # Ctrl+S
                "data": data,
            },
        )
        assert resp.status_code == 200
        assert resp.get_json()["src_written"] is True

        new_bytes = (ws_dir / name).read_bytes()
        assert (
            new_bytes != original_bytes
        ), "BUG-6: explicit save did NOT update the workspace file!"
        assert "Explicit Save Edit" in " ".join(_pptx_all_texts(new_bytes))


# ═══════════════════════════════════════════════════════════════════════════════
# 6.  BUG-7 — concurrent files must not cross-contaminate
# ═══════════════════════════════════════════════════════════════════════════════


class TestBug7_CrossContamination:

    def test_two_files_stay_independent(self, env):
        client, _, ws_dir, _ = env
        for n in ("iso_a.pptx", "iso_b.pptx"):
            client.post(
                "/api/v1/workspace/create_file",
                json={"folder": "", "name": n},
            )

        rA = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "iso_a.pptx"},
        ).get_json()
        rB = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "iso_b.pptx"},
        ).get_json()

        fidA, dA = rA["file_id"], rA["data"]
        fidB, dB = rB["file_id"], rB["data"]
        assert fidA != fidB

        _set_first_run(dA, "ONLY IN A")
        _set_first_run(dB, "ONLY IN B")

        for fid, data, name in [(fidA, dA, "iso_a.pptx"), (fidB, dB, "iso_b.pptx")]:
            client.post(
                "/api/v1/workspace/auto_save",
                json={
                    "file_type": "pptx",
                    "file_id": fid,
                    "ws_source_path": name,
                    "explicit": True,
                    "data": data,
                },
            )

        textsA = " ".join(_pptx_all_texts((ws_dir / "iso_a.pptx").read_bytes()))
        textsB = " ".join(_pptx_all_texts((ws_dir / "iso_b.pptx").read_bytes()))

        assert (
            "ONLY IN A" in textsA and "ONLY IN B" not in textsA
        ), f"Cross-contamination in file A: {textsA!r}"
        assert (
            "ONLY IN B" in textsB and "ONLY IN A" not in textsB
        ), f"Cross-contamination in file B: {textsB!r}"


# ═══════════════════════════════════════════════════════════════════════════════
# 7.  Full workflows — DOCX and XLSX (same patterns apply)
# ═══════════════════════════════════════════════════════════════════════════════


class TestDocxWorkflow:
    """DOCX: create → open → edit → Ctrl+S → close → reopen → verify."""

    def test_workspace_docx_survives_round_trip(self, env):
        client, _, ws_dir, _ = env
        client.post(
            "/api/v1/workspace/create_file",
            json={"folder": "", "name": "notes.docx"},
        )
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "notes.docx"},
        )
        fid = resp.get_json()["file_id"]

        new_html = "<p>Important meeting notes 20260409</p>"
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "docx",
                "file_id": fid,
                "ws_source_path": "notes.docx",
                "explicit": True,
                "data": new_html,
            },
        )
        assert resp.status_code == 200
        assert resp.get_json()["src_written"] is True

        # Reopen
        resp2 = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "notes.docx"},
        )
        html = resp2.get_json()["data"].get("html", "")
        assert "meeting notes" in html.lower(), f"DOCX content lost: {html[:200]}"

    def test_external_docx_survives_round_trip(self, env):
        client, _, _, ext_dir = env
        f = ext_dir / "letter.docx"
        f.write_bytes(_make_docx_bytes("Dear Alice"))

        resp = client.post(
            "/api/v1/workspace/open_abs_file",
            json={"path": str(f)},
        )
        assert resp.status_code == 200
        fid = resp.get_json()["file_id"]

        client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "docx",
                "file_id": fid,
                "ws_source_path": str(f),
                "explicit": True,
                "data": "<p>Dear Bob — updated letter content</p>",
            },
        )

        resp2 = client.post(
            "/api/v1/workspace/open_abs_file",
            json={"path": str(f)},
        )
        html = resp2.get_json()["data"].get("html", "")
        assert (
            "updated letter" in html.lower()
        ), f"External DOCX content lost: {html[:200]}"


class TestXlsxWorkflow:

    def test_workspace_xlsx_survives_round_trip(self, env):
        """XLSX round-trip: create → open → save → reopen without crash."""
        client, _, ws_dir, _ = env
        client.post(
            "/api/v1/workspace/create_file",
            json={"folder": "", "name": "budget.xlsx"},
        )
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "budget.xlsx"},
        )
        assert resp.status_code == 200
        fid = resp.get_json()["file_id"]
        data = resp.get_json()["data"]

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "xlsx",
                "file_id": fid,
                "ws_source_path": "budget.xlsx",
                "explicit": True,
                "data": data,
            },
        )
        assert resp.status_code == 200
        assert resp.get_json()["src_written"] is True

        # Reopen must not crash
        resp2 = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "budget.xlsx"},
        )
        assert resp2.status_code == 200
        assert resp2.get_json()["file_type"] == "xlsx"


# ═══════════════════════════════════════════════════════════════════════════════
# 8.  FS-browser create in external dir → full round-trip
# ═══════════════════════════════════════════════════════════════════════════════


class TestFsBrowserWorkflow:

    def test_absolute_fs_create_route_validates_payload(self, env):
        """The local browser create route exists but validates its JSON payload."""
        client, _, _, _ = env

        resp = client.post("/api/v1/fs/" + "create_" + "file", json={})
        assert resp.status_code == 400


# ═══════════════════════════════════════════════════════════════════════════════
# 9.  Data-format contract: parse → serialize → save → re-parse must be stable
# ═══════════════════════════════════════════════════════════════════════════════


class TestDataFormatContract:
    """
    The geometry format flowing between backend and frontend must be
    self-consistent:  parse_pptx_geometry → (frontend edits) → auto_save
    → re-parse must not lose structure.
    """

    def test_slide_count_preserved_through_save(self, env):
        """Saving must not add or remove slides."""
        client, _, _, _ = env
        pptx = _make_pptx_bytes("Stability")
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx), "stable.pptx")},
            content_type="multipart/form-data",
        )
        fid = resp.get_json()["file_id"]
        data = resp.get_json()["data"]
        n_before = len(data["slides"])

        client.post(
            "/api/v1/workspace/auto_save",
            json={"file_type": "pptx", "file_id": fid, "data": data},
        )

        raw = client.get(f"/api/v1/workspace/raw/{fid}").data
        n_after = _pptx_slide_count(raw)
        assert n_after == n_before, f"Slide count changed: {n_before} → {n_after}"

    def test_shape_text_survives_parse_save_reparse(self, env):
        """Parse → save unmodified → re-parse must have identical text."""
        client, _, _, _ = env
        pptx = _make_pptx_bytes("Round Trip Text 12345")
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx), "rt.pptx")},
            content_type="multipart/form-data",
        )
        fid = resp.get_json()["file_id"]
        data = resp.get_json()["data"]
        text_before = _runs_text(data)

        # Save without editing
        client.post(
            "/api/v1/workspace/auto_save",
            json={"file_type": "pptx", "file_id": fid, "data": data},
        )

        # Re-parse the saved file
        raw = client.get(f"/api/v1/workspace/raw/{fid}").data
        import tempfile

        from app.core.file.file_parser import parse_pptx_geometry

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tf:
            tf.write(raw)
            tf.flush()
            data2 = parse_pptx_geometry(tf.name)
        os.unlink(tf.name)

        text_after = _runs_text(data2)
        assert text_before == text_after, (
            f"Text changed through unmodified save!\n"
            f"  Before: {text_before!r}\n"
            f"  After:  {text_after!r}"
        )

    def test_rich_format_has_required_keys(self, env):
        """Parsed geometry data must contain the keys auto_save depends on."""
        client, _, _, _ = env
        pptx = _make_pptx_bytes("Key Check")
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx), "keys.pptx")},
            content_type="multipart/form-data",
        )
        data = resp.get_json()["data"]

        # auto_save checks isinstance(data, dict) and "slides" in data
        assert isinstance(data, dict), "Parsed PPTX data must be a dict"
        assert "slides" in data, "Parsed data must have 'slides' key"
        assert isinstance(data["slides"], list)

        for slide in data["slides"]:
            assert "shapes" in slide, "Each slide must have 'shapes'"
            for shape in slide["shapes"]:
                if shape.get("has_text"):
                    assert (
                        "paragraphs" in shape
                    ), "Text shape must have 'paragraphs' key"
                    for para in shape["paragraphs"]:
                        assert "runs" in para, "Paragraph must have 'runs' key"
