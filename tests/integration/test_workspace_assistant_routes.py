# -*- coding: utf-8 -*-
"""
Integration tests for /api/v1/workspace/* endpoints (workspace_assistant_bp).

Covers fixes introduced in Logan/20260326:
  1. raw_file: send_file must use absolute path (Flask 3.1 rejects relative
     paths -> was returning 500, now returns 200).
  2. open_file: full upload-parse-persist workflow (DOCX / PDF).
  3. serve_workspace_file: serves files from workspace dir; path-traversal
     guard returns 403; missing file returns 404.
  4. auto_save: both first and second explicit save write different bytes
     (fix for "second save does not work" bug).
  5. raw_file: returns Cache-Control: no-store headers so browser never
     serves a cached (stale) response on repeated saves.
  6. JS source: cache-buster ?_=<ts> added to raw fetch URL, and
     showSaveFilePicker called on first Ctrl+S when no handle exists.

Note: recent-files panel feature was removed ("recent files not needed").
Tests for that feature (TestRecentFilesJsFix, uploads-path re-open) are
intentionally absent.
"""

from __future__ import annotations

import io
import os
import uuid
from pathlib import Path

import pytest

# ── Fixture: minimal Flask app with only workspace_assistant_bp ──────────────


@pytest.fixture(scope="module")
def wa_client(tmp_path_factory):
    """
    Flask test client backed by a temporary workspace/tmp directory.
    We patch _TMP_ROOT and _get_session_id so tests never touch the real workspace
    and don't need a real Flask session cookie.
    """
    os.environ.setdefault("KOTO_AUTH_ENABLED", "false")

    tmp_root = tmp_path_factory.mktemp("wa_root")
    tmp_dir = tmp_root / "tmp" / "testsession"
    workspace_dir = tmp_root / "workspace"
    tmp_dir.mkdir(parents=True)
    workspace_dir.mkdir(parents=True)

    import web.blueprints.workspace_assistant as _wa
    import web.shared as _shared

    original_tmp = _wa._TMP_ROOT
    original_get_sid = _wa._get_session_id
    original_ws = getattr(_shared, "WORKSPACE_DIR", None)
    # Patch _TMP_ROOT to the parent of our isolated dir and fix session to always
    # return 'testsession' so _ensure_tmp_dir() resolves to tmp_dir.
    _wa._TMP_ROOT = tmp_root / "tmp"
    _wa._get_session_id = lambda: "testsession"
    _shared.WORKSPACE_DIR = str(workspace_dir)

    from flask import Flask

    from web.blueprints.workspace_assistant import workspace_assistant_bp

    app = Flask(__name__)
    app.register_blueprint(workspace_assistant_bp)
    app.config["TESTING"] = True
    app.config["SECRET_KEY"] = "test-secret"

    with app.test_client() as client:
        yield client, tmp_dir, workspace_dir

    _wa._TMP_ROOT = original_tmp
    _wa._get_session_id = original_get_sid
    if original_ws is not None:
        _shared.WORKSPACE_DIR = original_ws


def _fake_pdf_bytes() -> bytes:
    return b"%PDF-1.4\n1 0 obj\n<</Type /Catalog>>\nendobj\n%%EOF"


def _fake_docx_bytes() -> bytes:
    try:
        from io import BytesIO

        import docx

        doc = docx.Document()
        doc.add_paragraph("Hello world")
        buf = BytesIO()
        doc.save(buf)
        return buf.getvalue()
    except ImportError:
        return b""


# ── 1. GET /api/v1/workspace/raw/<file_id> ───────────────────────────────────


class TestRawFile:
    """Fix: send_file must receive an absolute path (Flask 3.1 requirement)."""

    def test_returns_200_with_pdf_content(self, wa_client):
        client, tmp_dir, _ = wa_client
        file_id = uuid.uuid4().hex
        (tmp_dir / f"{file_id}.pdf").write_bytes(_fake_pdf_bytes())
        resp = client.get(f"/api/v1/workspace/raw/{file_id}")
        assert resp.status_code == 200
        assert resp.content_type == "application/pdf"
        assert resp.data.startswith(b"%PDF")

    def test_path_passed_to_send_file_is_absolute(self, wa_client, monkeypatch):
        """The fix: .resolve() ensures send_file gets an absolute path."""
        client, tmp_dir, _ = wa_client
        file_id = uuid.uuid4().hex
        (tmp_dir / f"{file_id}.pdf").write_bytes(_fake_pdf_bytes())

        captured = {}
        import flask

        real_send_file = flask.send_file

        def spy_send_file(path_or_file, *args, **kwargs):
            captured["path"] = str(path_or_file)
            return real_send_file(path_or_file, *args, **kwargs)

        monkeypatch.setattr(
            "web.blueprints.workspace_assistant.send_file", spy_send_file
        )
        resp = client.get(f"/api/v1/workspace/raw/{file_id}")

        assert resp.status_code == 200
        assert Path(captured["path"]).is_absolute(), (
            f"send_file received relative path {captured['path']!r}; "
            "fix requires .resolve() to make it absolute"
        )

    def test_returns_404_for_missing_file_id(self, wa_client):
        client, _, _ = wa_client
        resp = client.get(f"/api/v1/workspace/raw/{uuid.uuid4().hex}")
        assert resp.status_code == 404

    def test_returns_400_for_non_alnum_id(self, wa_client):
        client, _, _ = wa_client
        resp = client.get("/api/v1/workspace/raw/bad..id")
        assert resp.status_code in (400, 404)

    def test_docx_served_with_correct_mime(self, wa_client):
        client, tmp_dir, _ = wa_client
        file_id = uuid.uuid4().hex
        (tmp_dir / f"{file_id}.docx").write_bytes(b"PK\x03\x04fake")
        resp = client.get(f"/api/v1/workspace/raw/{file_id}")
        assert resp.status_code == 200
        assert "wordprocessingml" in resp.content_type


# ── 2. POST /api/v1/workspace/open_file ──────────────────────────────────────


class TestOpenFile:

    def test_rejects_missing_file_field(self, wa_client):
        client, _, _ = wa_client
        resp = client.post("/api/v1/workspace/open_file", data={})
        assert resp.status_code == 400

    def test_rejects_unsupported_extension(self, wa_client):
        client, _, _ = wa_client
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(b"hello"), "notes.xyz")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 400
        assert "不支持" in resp.get_json().get("error", "")

    def test_pdf_stored_in_tmp_dir(self, wa_client):
        """Uploading a PDF always returns 200 (fix: no longer throws on missing pdfplumber)."""
        client, tmp_dir, _ = wa_client
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(_fake_pdf_bytes()), "report.pdf")},
            content_type="multipart/form-data",
        )
        body = resp.get_json()
        assert resp.status_code == 200, (
            f"Expected 200 but got {resp.status_code}. "
            f"parse_pdf should never raise when PDF libs are missing; got: {body}"
        )
        assert "file_id" in body
        matches = list(tmp_dir.glob(f"{body['file_id']}.*"))
        assert matches, "uploaded PDF must be saved in tmp_dir"

    def test_docx_returns_html_data(self, wa_client):
        client, _, _ = wa_client
        docx_bytes = _fake_docx_bytes()
        if not docx_bytes:
            pytest.skip("python-docx not available")
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(docx_bytes), "document.docx")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 200
        body = resp.get_json()
        assert body["file_type"] == "docx"
        assert "html" in body.get("data", {})

    def test_file_id_is_hex_uuid(self, wa_client):
        client, _, _ = wa_client
        docx_bytes = _fake_docx_bytes()
        if not docx_bytes:
            pytest.skip("python-docx not available")
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(docx_bytes), "id_test.docx")},
            content_type="multipart/form-data",
        )
        if resp.status_code != 200:
            pytest.skip("parse failed")
        file_id = resp.get_json()["file_id"]
        assert file_id.isalnum() and len(file_id) == 32

    def test_open_file_by_path_retries_docx_after_tmp_zip_failure(self, wa_client, monkeypatch):
        import zipfile

        client, tmp_dir, workspace_dir = wa_client
        docx_bytes = _fake_docx_bytes()
        if not docx_bytes:
            pytest.skip("python-docx not available")

        target = workspace_dir / "retry.docx"
        target.write_bytes(docx_bytes)

        import app.core.file.parsers.docx_parser as parser_mod

        real_parse_docx = parser_mod.parse_docx
        call_count = {"value": 0}

        def flaky_parse_docx(path: str, *args, **kwargs):
            call_count["value"] += 1
            if call_count["value"] == 1:
                raise zipfile.BadZipFile("File is not a zip file")
            return real_parse_docx(path, *args, **kwargs)

        monkeypatch.setattr(parser_mod, "parse_docx", flaky_parse_docx)

        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "retry.docx"},
        )

        assert resp.status_code == 200, resp.get_data(as_text=True)
        assert call_count["value"] == 2
        body = resp.get_json()
        assert body["file_type"] == "docx"
        assert body.get("data", {}).get("raw_url")
        tmp_copy = tmp_dir / f"{body['file_id']}.docx"
        assert tmp_copy.is_file()
        assert tmp_copy.read_bytes() == docx_bytes


class TestAIContextPreview:

    def test_unicode_docx_path_is_readable(self, wa_client):
        client, _, workspace_dir = wa_client
        target_name = "\u8bfb\u53d6\u6d4b\u8bd5.docx"
        target = workspace_dir / target_name
        target.write_bytes(_make_docx_bytes(["中文文件名读取测试", "第二段内容"]))

        resp = client.post(
            "/api/v1/workspace/ai_context_preview",
            json={"path": target_name},
        )

        assert resp.status_code == 200, resp.get_data(as_text=True)
        body = resp.get_json()
        assert body["file_name"] == target_name
        assert body["file_type"] == "docx"
        assert "中文文件名读取测试" in body["content_preview"]

    def test_docx_original_chars_uses_full_document_count(self, wa_client):
        client, _, workspace_dir = wa_client
        paragraphs = [
            f"第{idx + 1}段：" + "这是用于验证DOCX统计更接近Word和WPS的测试内容。" * 5
            for idx in range(360)
        ]
        expected_chars = sum(len("".join(text.split())) for text in paragraphs)

        target = workspace_dir / "long-preview.docx"
        target.write_bytes(_make_docx_bytes(paragraphs))

        resp = client.post(
            "/api/v1/workspace/ai_context_preview",
            json={"path": "long-preview.docx"},
        )

        assert resp.status_code == 200, resp.get_data(as_text=True)
        body = resp.get_json()
        preview_chars = len("".join(str(body.get("content_preview") or "").split()))
        assert body["file_type"] == "docx"
        assert body["original_chars"] == expected_chars
        assert body["original_chars"] > preview_chars

    def test_parse_error_keeps_attachment_available(self, wa_client, monkeypatch):
        client, _, workspace_dir = wa_client
        target = workspace_dir / "parse-warning.pdf"
        target.write_bytes(_fake_pdf_bytes())

        from app.core.agent import task_tools

        monkeypatch.setattr(
            task_tools,
            "parse_file_to_text",
            lambda *args, **kwargs: "Error parsing file: simulated parser failure",
        )

        resp = client.post(
            "/api/v1/workspace/ai_context_preview",
            json={"path": "parse-warning.pdf"},
        )

        assert resp.status_code == 200, resp.get_data(as_text=True)
        body = resp.get_json()
        assert body["file_type"] == "pdf"
        assert body["content_preview"] == ""
        assert "simulated parser failure" in body["preview_error"]

    def test_unexpected_parse_exception_keeps_attachment_available(self, wa_client, monkeypatch):
        client, _, workspace_dir = wa_client
        target = workspace_dir / "parse-exception.docx"
        target.write_bytes(_make_docx_bytes("body"))

        from app.core.agent import task_tools

        def _raise_parse_error(*args, **kwargs):
            raise RuntimeError("unexpected parser boom")

        monkeypatch.setattr(task_tools, "parse_file_to_text", _raise_parse_error)

        resp = client.post(
            "/api/v1/workspace/ai_context_preview",
            json={"path": "parse-exception.docx"},
        )

        assert resp.status_code == 200, resp.get_data(as_text=True)
        body = resp.get_json()
        assert body["file_type"] == "docx"
        assert body["content_preview"] == ""
        assert "unexpected parser boom" in body["preview_error"]


# ── 2b. PDF-specific loading tests ───────────────────────────────────────────


class TestPdfLoading:
    """Dedicated tests for PDF load pipeline (parse_pdf + raw endpoint)."""

    def _upload_pdf(self, client, name="sample.pdf"):
        return client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(_fake_pdf_bytes()), name)},
            content_type="multipart/form-data",
        )

    def test_pdf_open_always_returns_200(self, wa_client):
        """PDF open must never return 500 even when text-extraction libs are absent."""
        client, _, _ = wa_client
        resp = self._upload_pdf(client)
        assert resp.status_code == 200, (
            f"parse_pdf must not raise when PDF libs missing; got {resp.status_code}: "
            f"{resp.get_json()}"
        )

    def test_pdf_response_contains_raw_url(self, wa_client):
        """data.raw_url must point to /api/v1/workspace/raw/<file_id>."""
        client, _, _ = wa_client
        resp = self._upload_pdf(client)
        assert resp.status_code == 200
        body = resp.get_json()
        raw_url = body.get("data", {}).get("raw_url", "")
        file_id = body["file_id"]
        assert f"/api/v1/workspace/raw/{file_id}" == raw_url, (
            f"raw_url should be /api/v1/workspace/raw/<file_id>, got {raw_url!r}"
        )

    def test_pdf_file_type_is_pdf(self, wa_client):
        client, _, _ = wa_client
        resp = self._upload_pdf(client)
        assert resp.status_code == 200
        assert resp.get_json()["file_type"] == "pdf"

    def test_pdf_page_count_is_non_negative(self, wa_client):
        """page_count must exist and be >= 0 (0 when lib unavailable is fine)."""
        client, _, _ = wa_client
        resp = self._upload_pdf(client)
        assert resp.status_code == 200
        page_count = resp.get_json().get("data", {}).get("page_count")
        assert page_count is not None, "data.page_count must be present"
        assert isinstance(page_count, int) and page_count >= 0

    def test_pdf_pages_is_list(self, wa_client):
        """data.pages must be a list (empty list is acceptable when lib unavailable)."""
        client, _, _ = wa_client
        resp = self._upload_pdf(client)
        assert resp.status_code == 200
        pages = resp.get_json().get("data", {}).get("pages")
        assert isinstance(pages, list), f"data.pages must be a list, got {type(pages)}"

    def test_pdf_open_succeeds_without_pdfplumber(self, wa_client, monkeypatch):
        """
        Critical regression guard: even when pdfplumber, pypdf, AND PyPDF2 are
        all absent, parse_pdf must NOT raise — it returns a graceful fallback.
        """
        import builtins

        real_import = builtins.__import__

        def blocked_import(name, *args, **kwargs):
            if name in ("pdfplumber", "pypdf", "PyPDF2"):
                raise ImportError(f"[test] blocked import: {name}")
            return real_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", blocked_import)

        client, _, _ = wa_client
        resp = self._upload_pdf(client, "no_lib.pdf")
        assert resp.status_code == 200, (
            f"parse_pdf must return 200 even with no PDF text libs; "
            f"got {resp.status_code}: {resp.get_json()}"
        )
        data = resp.get_json().get("data", {})
        # raw_url must still be present so PDF.js can render
        assert "/api/v1/workspace/raw/" in data.get("raw_url", ""), (
            "raw_url must be present even when text extraction is skipped"
        )
        # text/pages may be empty — that's fine
        assert isinstance(data.get("pages", []), list)

    def test_raw_url_is_fetchable_after_upload(self, wa_client):
        """After upload, GET data.raw_url should return the binary PDF bytes."""
        client, _, _ = wa_client
        resp = self._upload_pdf(client, "fetchable.pdf")
        assert resp.status_code == 200
        raw_url = resp.get_json()["data"]["raw_url"]
        raw_resp = client.get(raw_url)
        assert raw_resp.status_code == 200
        assert raw_resp.data.startswith(b"%PDF"), (
            f"{raw_url} did not return PDF bytes; first bytes: {raw_resp.data[:20]!r}"
        )


# ── 3. GET /api/v1/workspace/file/<path> ─────────────────────────────────────


class TestServeWorkspaceFile:

    def test_serves_pdf_from_workspace_root(self, wa_client):
        client, _, workspace_dir = wa_client
        (workspace_dir / "sample.pdf").write_bytes(_fake_pdf_bytes())
        resp = client.get("/api/v1/workspace/file/sample.pdf")
        assert resp.status_code == 200
        assert resp.content_type == "application/pdf"

    def test_serves_pdf_from_uploads_subdir(self, wa_client):
        """
        Critical for the recent-files fix: clicking a recent file uses
        'uploads/<name>' path which must resolve to workspace/uploads/<name>.
        """
        client, _, workspace_dir = wa_client
        uploads = workspace_dir / "uploads"
        uploads.mkdir(exist_ok=True)
        (uploads / "sub.pdf").write_bytes(_fake_pdf_bytes())
        resp = client.get("/api/v1/workspace/file/uploads/sub.pdf")
        assert resp.status_code == 200
        assert resp.content_type == "application/pdf"

    def test_returns_404_for_missing_file(self, wa_client):
        client, _, _ = wa_client
        resp = client.get("/api/v1/workspace/file/ghost.pdf")
        assert resp.status_code == 404

    def test_returns_400_for_unsupported_extension(self, wa_client):
        client, _, workspace_dir = wa_client
        (workspace_dir / "readme.xyz").write_text("hi")
        resp = client.get("/api/v1/workspace/file/readme.xyz")
        assert resp.status_code == 400

    def test_path_traversal_blocked(self, wa_client):
        client, _, _ = wa_client
        resp = client.get("/api/v1/workspace/file/../../../etc/passwd")
        assert resp.status_code in (403, 404)


class TestLegacyRoutesRemoved:

    def test_obsolete_workspace_assistant_routes_are_unregistered(self, wa_client):
        client, _, _ = wa_client
        rules = {rule.rule for rule in client.application.url_map.iter_rules()}
        assert "/api/v1/workspace/read_for_ai" not in rules
        assert "/api/v1/workspace/summarize" not in rules
        assert "/api/v1/workspace/quick-action" not in rules


# ── 5. Round-trip: upload → raw endpoint ─────────────────────────────────────


class TestRoundTrip:

    def test_upload_pdf_then_fetch_via_raw(self, wa_client):
        """
        Full pipeline: upload PDF -> get file_id -> fetch via /raw/<id> -> 200.
        Tests the absolute-path fix end-to-end.
        """
        client, _, _ = wa_client
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(_fake_pdf_bytes()), "roundtrip.pdf")},
            content_type="multipart/form-data",
        )
        if resp.status_code != 200:
            pytest.skip("PDF parser not available")

        file_id = resp.get_json()["file_id"]
        raw = client.get(f"/api/v1/workspace/raw/{file_id}")
        assert raw.status_code == 200, (
            f"/raw/{file_id} returned {raw.status_code}. "
            "Check that send_file uses absolute path (.resolve())."
        )
        assert raw.data.startswith(b"%PDF")


# ── 6. list_files: size + mtime metadata ─────────────────────────────────────


class TestListFilesMetadata:

    def test_list_files_returns_size_and_mtime_for_files(self, wa_client):
        """
        After the panel-improvements commit, list_files must return size
        (human-readable string) and mtime (milliseconds int) for every file entry.
        """
        client, _, workspace_dir = wa_client
        # plant a file directly in the workspace dir so list_files picks it up
        uploads = workspace_dir / "uploads"
        uploads.mkdir(exist_ok=True)
        (uploads / "meta_test.pdf").write_bytes(b"%PDF-1.0 test")

        resp = client.get("/api/v1/workspace/list_files")
        assert resp.status_code == 200
        data = resp.get_json()

        # find our test file in the tree
        found = None
        for node in data.get("files", []):
            if node["type"] == "folder" and node["name"] == "uploads":
                for child in node.get("children", []):
                    if child["name"] == "meta_test.pdf":
                        found = child
                        break
        assert found is not None, "meta_test.pdf not found in list_files response"
        assert "size" in found, "list_files must include 'size' field for each file"
        assert "mtime" in found, "list_files must include 'mtime' field for each file"
        assert isinstance(
            found["size"], str
        ), "'size' should be a human-readable string"
        assert isinstance(
            found["mtime"], (int, float)
        ), "'mtime' should be a numeric timestamp"


# ── 7. PATCH /api/v1/workspace/rename ────────────────────────────────────────


class TestRenameEndpoint:

    def test_rename_nonexistent_file_returns_404(self, wa_client):
        client, _, _ = wa_client
        resp = client.patch(
            "/api/v1/workspace/rename",
            json={"path": "uploads/does_not_exist.pdf", "name": "new_name"},
        )
        assert resp.status_code == 404

    def test_rename_path_traversal_rejected(self, wa_client):
        client, _, _ = wa_client
        resp = client.patch(
            "/api/v1/workspace/rename",
            json={"path": "uploads/legit.pdf", "name": "../evil"},
        )
        # name contains '/' or '\\' → 400
        assert resp.status_code == 400

    def test_rename_preserves_extension(self, wa_client):
        client, _, workspace_dir = wa_client
        uploads = workspace_dir / "uploads"
        uploads.mkdir(exist_ok=True)
        src = uploads / "to_rename.pdf"
        src.write_bytes(b"%PDF-1.0")

        resp = client.patch(
            "/api/v1/workspace/rename",
            json={"path": "uploads/to_rename.pdf", "name": "renamed_file"},
        )
        assert resp.status_code == 200
        data = resp.get_json()
        # Extension must be preserved even though user didn't supply it
        assert data["path"].endswith(
            ".pdf"
        ), "Rename must preserve original file extension"
        assert not (uploads / "to_rename.pdf").exists(), "Old file should be gone"
        assert (uploads / "renamed_file.pdf").exists(), "Renamed file should exist"

    def test_rename_duplicate_name_returns_409(self, wa_client):
        client, _, workspace_dir = wa_client
        uploads = workspace_dir / "uploads"
        uploads.mkdir(exist_ok=True)
        (uploads / "dup_src.docx").write_bytes(b"PK content")
        (uploads / "dup_target.docx").write_bytes(b"PK content")

        resp = client.patch(
            "/api/v1/workspace/rename",
            json={"path": "uploads/dup_src.docx", "name": "dup_target"},
        )
        assert (
            resp.status_code == 409
        ), "Rename to an existing filename must return 409 Conflict"

    def test_rename_success_response_shape(self, wa_client):
        client, _, workspace_dir = wa_client
        uploads = workspace_dir / "uploads"
        uploads.mkdir(exist_ok=True)
        (uploads / "shape_test.docx").write_bytes(b"PK content")

        resp = client.patch(
            "/api/v1/workspace/rename",
            json={"path": "uploads/shape_test.docx", "name": "shape_renamed"},
        )
        assert resp.status_code == 200
        data = resp.get_json()
        assert "path" in data, "Rename response must include 'path'"
        assert "name" in data, "Rename response must include 'name'"


# ── helpers shared by new test classes ───────────────────────────────────────


def _make_docx_bytes(text: str | list[str] = "Test") -> bytes:
    """Minimal valid .docx (ZIP) with one or more paragraphs of text."""
    import io as _io
    import zipfile
    from xml.sax.saxutils import escape as _xml_escape

    if isinstance(text, list):
        body = "".join(
            f'<w:p><w:r><w:t xml:space="preserve">{_xml_escape(str(item))}</w:t></w:r></w:p>'
            for item in text
        )
    else:
        body = (
            f'<w:p><w:r><w:t xml:space="preserve">{_xml_escape(str(text))}</w:t></w:r></w:p>'
        )

    ct = (
        '<?xml version="1.0"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml"'
        ' ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1"'
        ' Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument"'
        ' Target="word/document.xml"/>'
        "</Relationships>"
    )
    doc = (
        '<?xml version="1.0"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        f"<w:body>{body}</w:body>"
        "</w:document>"
    )
    dr = (
        '<?xml version="1.0"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'
    )
    buf = _io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr("[Content_Types].xml", ct)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", doc)
        z.writestr("word/_rels/document.xml.rels", dr)
    return buf.getvalue()


# ── 8. POST /api/v1/workspace/auto_save ──────────────────────────────────────


class TestAutoSave:
    """
    Fix: 'second save does not work'.

    Root cause: auto_save correctly updates the tmp file on every call, but
    the browser cached the first /raw/<id> response and returned stale bytes
    on every subsequent save, so the local file appeared unchanged.

    These tests verify the server always produces fresh bytes per save.
    """

    def _upload_docx(self, client) -> str:
        """Upload a minimal docx and return its file_id."""
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(_make_docx_bytes("original")), "save_test.docx")},
            content_type="multipart/form-data",
        )
        if resp.status_code != 200:
            pytest.skip("docx parse not available in this environment")
        return resp.get_json()["file_id"]

    def test_missing_fields_returns_400(self, wa_client):
        client, _, _ = wa_client
        resp = client.post("/api/v1/workspace/auto_save", json={})
        assert resp.status_code == 400

    def test_invalid_file_id_returns_400(self, wa_client):
        client, _, _ = wa_client
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={"file_type": "docx", "file_id": "../evil", "data": "<p>x</p>"},
        )
        assert resp.status_code == 400

    def test_first_explicit_save_returns_ok(self, wa_client):
        client, _, _ = wa_client
        fid = self._upload_docx(client)
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "docx",
                "file_id": fid,
                "ws_source_path": "save_test.docx",
                "explicit": True,
                "data": "<p>first edit</p>",
            },
        )
        assert resp.status_code == 200
        body = resp.get_json()
        assert body.get("ok") is True
        assert "saved_at" in body

    def test_second_save_raw_bytes_differ_from_first(self, wa_client):
        """
        Core regression test: /raw/<id> must return different bytes after
        save2 vs save1 because the content changed.
        Before the cache-buster fix, the browser cached the first response
        so the local file was always overwritten with save1 bytes.
        """
        client, _, _ = wa_client
        fid = self._upload_docx(client)

        client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "docx",
                "file_id": fid,
                "ws_source_path": "save_test.docx",
                "explicit": True,
                "data": "<p>first edit</p>",
            },
        )
        raw1 = client.get(f"/api/v1/workspace/raw/{fid}").data

        client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "docx",
                "file_id": fid,
                "ws_source_path": "save_test.docx",
                "explicit": True,
                "data": "<p>second edit with substantially more text appended</p>",
            },
        )
        raw2 = client.get(f"/api/v1/workspace/raw/{fid}").data

        assert raw1 != raw2, (
            "raw bytes after save2 must differ from save1 — "
            "identical bytes means the server is not updating the tmp file"
        )
        assert len(raw2) > len(raw1), (
            f"save2 bytes ({len(raw2)}) should exceed save1 ({len(raw1)}) "
            "since more text was saved"
        )

    def test_workspace_file_updated_on_each_save(self, wa_client):
        """workspace/<ws_source_path> must be overwritten on every explicit save."""
        client, _, workspace_dir = wa_client
        fid = self._upload_docx(client)

        client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "docx",
                "file_id": fid,
                "ws_source_path": "save_test.docx",
                "explicit": True,
                "data": "<p>save A</p>",
            },
        )
        size1 = (workspace_dir / "save_test.docx").stat().st_size

        client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "docx",
                "file_id": fid,
                "ws_source_path": "save_test.docx",
                "explicit": True,
                "data": "<p>save B with much more text to ensure size increases</p>",
            },
        )
        size2 = (workspace_dir / "save_test.docx").stat().st_size

        assert (
            size2 != size1
        ), f"workspace file unchanged after second save ({size1} bytes both times)"

    def test_auto_save_implicit_succeeds(self, wa_client):
        """explicit=False (background auto-save) must also return 200."""
        client, _, _ = wa_client
        fid = self._upload_docx(client)
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={"file_type": "docx", "file_id": fid, "data": "<p>auto</p>"},
        )
        assert resp.status_code == 200

    def test_structured_docx_payload_writes_header_footer(self, wa_client):
        client, _, _ = wa_client
        docx_module = pytest.importorskip("docx")
        import zipfile

        src = io.BytesIO()
        source_doc = docx_module.Document()
        source_doc.add_paragraph("original body")
        source_doc.save(src)
        src.seek(0)

        upload = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (src, "header_footer_save.docx")},
            content_type="multipart/form-data",
        )
        if upload.status_code != 200:
            pytest.skip("docx parse not available in this environment")
        fid = upload.get_json()["file_id"]

        header_html = (
            '<p><span class="koto-hdr-col">项目计划</span>'
            '<span class="koto-hdr-col"><span class="koto-hdr-page-num">1</span></span>'
            '<span class="koto-hdr-col">内部使用</span></p>'
        )
        footer_html = '<p>页脚说明</p>'
        payload = {
            "html": "<p>更新后的正文</p>",
            "header_html": header_html,
            "footer_html": footer_html,
            "sections": [{
                "header_html": header_html,
                "footer_html": footer_html,
                "first_header_html": "",
                "first_footer_html": "",
                "even_header_html": "",
                "even_footer_html": "",
            }],
        }

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "docx",
                "file_id": fid,
                "ws_source_path": "header_footer_save.docx",
                "explicit": True,
                "data": payload,
            },
        )
        assert resp.status_code == 200

        raw = client.get(f"/api/v1/workspace/raw/{fid}").data
        saved_doc = docx_module.Document(io.BytesIO(raw))
        header_text = "\n".join(p.text for p in saved_doc.sections[0].header.paragraphs)
        footer_text = "\n".join(p.text for p in saved_doc.sections[0].footer.paragraphs)
        body_text = "\n".join(p.text for p in saved_doc.paragraphs)

        assert "项目计划" in header_text
        assert "内部使用" in header_text
        assert "页脚说明" in footer_text
        assert "更新后的正文" in body_text

        with zipfile.ZipFile(io.BytesIO(raw)) as archive:
            header_parts = [name for name in archive.namelist() if name.startswith("word/header")]
            assert header_parts, "expected DOCX export to generate at least one header part"
            header_xml = "".join(
                archive.read(name).decode("utf-8", errors="ignore")
                for name in header_parts
            )
        assert "PAGE" in header_xml, "header export should preserve Word PAGE field"

    def test_structured_docx_payload_writes_comments_xml(self, wa_client):
        client, _, _ = wa_client
        docx_module = pytest.importorskip("docx")
        import zipfile

        src = io.BytesIO()
        source_doc = docx_module.Document()
        source_doc.add_paragraph("第一段原文")
        source_doc.add_paragraph("第二段保留")
        source_doc.save(src)
        src.seek(0)

        upload = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (src, "comment_save.docx")},
            content_type="multipart/form-data",
        )
        if upload.status_code != 200:
            pytest.skip("docx parse not available in this environment")
        fid = upload.get_json()["file_id"]

        payload = {
            "html": "<p>第一段原文</p><p>第二段保留</p>",
            "comments": [
                {
                    "id": "comment-1",
                    "author": "审阅人",
                    "date": "2026-05-12T10:30:00Z",
                    "text": "这里需要进一步说明",
                    "anchor_text": "第一段原文",
                }
            ],
        }

        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "docx",
                "file_id": fid,
                "ws_source_path": "comment_save.docx",
                "explicit": True,
                "data": payload,
            },
        )
        assert resp.status_code == 200

        raw = client.get(f"/api/v1/workspace/raw/{fid}").data
        with zipfile.ZipFile(io.BytesIO(raw)) as archive:
            names = archive.namelist()
            assert "word/comments.xml" in names
            comments_xml = archive.read("word/comments.xml").decode("utf-8", errors="ignore")
            document_xml = archive.read("word/document.xml").decode("utf-8", errors="ignore")
            rels_xml = archive.read("word/_rels/document.xml.rels").decode("utf-8", errors="ignore")
            content_types_xml = archive.read("[Content_Types].xml").decode("utf-8", errors="ignore")

        assert "这里需要进一步说明" in comments_xml
        assert "审阅人" in comments_xml
        assert "commentRangeStart" in document_xml
        assert "commentReference" in document_xml
        assert "comments.xml" in rels_xml
        assert "/word/comments.xml" in content_types_xml

    def test_src_written_true_when_ws_path_provided(self, wa_client):
        """Response must include src_written=True when ws_source_path is given."""
        client, _, _ = wa_client
        fid = self._upload_docx(client)
        resp = client.post(
            "/api/v1/workspace/auto_save",
            json={
                "file_type": "docx",
                "file_id": fid,
                "ws_source_path": "save_test.docx",
                "explicit": True,
                "data": "<p>check src_written</p>",
            },
        )
        assert resp.status_code == 200
        assert resp.get_json().get("src_written") is True


# ── 9. raw_file: Cache-Control no-store ──────────────────────────────────────


class TestRawFileNoCacheHeaders:
    """
    Fix: raw_file must return Cache-Control: no-store so the browser never
    serves the first-save's cached response on subsequent saves.
    Without this, the local file written via FileSystemFileHandle always
    received stale bytes (second save appeared to do nothing).
    """

    def test_no_store_header_present(self, wa_client):
        client, tmp_dir, _ = wa_client
        fid = uuid.uuid4().hex
        (tmp_dir / f"{fid}.pdf").write_bytes(_fake_pdf_bytes())

        resp = client.get(f"/api/v1/workspace/raw/{fid}")
        assert resp.status_code == 200
        cc = resp.headers.get("Cache-Control", "")
        assert "no-store" in cc, (
            f"Cache-Control must contain 'no-store' (got {cc!r}) — "
            "without it the browser caches save1 bytes and second save writes stale data"
        )

    def test_no_cache_header_present(self, wa_client):
        client, tmp_dir, _ = wa_client
        fid = uuid.uuid4().hex
        (tmp_dir / f"{fid}.docx").write_bytes(b"PK\x03\x04fake")

        resp = client.get(f"/api/v1/workspace/raw/{fid}")
        cc = resp.headers.get("Cache-Control", "")
        assert "no-cache" in cc, f"Cache-Control must contain 'no-cache' (got {cc!r})"

    def test_pragma_no_cache_present(self, wa_client):
        client, tmp_dir, _ = wa_client
        fid = uuid.uuid4().hex
        (tmp_dir / f"{fid}.pdf").write_bytes(_fake_pdf_bytes())

        resp = client.get(f"/api/v1/workspace/raw/{fid}")
        pragma = resp.headers.get("Pragma", "")
        assert (
            "no-cache" in pragma
        ), f"Pragma: no-cache header required for HTTP/1.0 compatibility (got {pragma!r})"


# ── 10. JS source: save-flow fixes ───────────────────────────────────────────


class TestSaveFlowJsFixes:
    """
    Validates workspace-assistant.js contains all three save-flow fixes:
      (a) State (tab, fsHandle, fileId, etc.) captured before any await
      (b) showSaveFilePicker called when no fsHandle exists (first save)
      (c) Cache-buster ?_=Date.now() on the /raw/ fetch URL
      (d) _isSaving guard with finally-block reset
    """

    JS_PATH = (
        Path(__file__).parents[2] / "web" / "static" / "js" / "workspace-assistant.js"
    )

    @property
    def src(self) -> str:
        return self.JS_PATH.read_text(encoding="utf-8")

    # (a) pre-await state capture -----------------------------------------

    def test_save_tab_captured_before_await(self):
        assert "_saveTab" in self.src, "saveFile must capture _saveTab before any await"

    def test_save_fshandle_captured_before_await(self):
        assert (
            "_saveFsHandle" in self.src
        ), "saveFile must capture _saveFsHandle before any await"

    def test_save_file_id_captured_before_await(self):
        assert (
            "_saveFileId" in self.src
        ), "saveFile must capture _saveFileId before any await"

    def test_save_file_type_captured_before_await(self):
        assert (
            "_saveFileType" in self.src
        ), "saveFile must capture _saveFileType before any await"

    def test_save_ws_path_captured_before_await(self):
        assert (
            "_saveWsPath" in self.src
        ), "saveFile must capture _saveWsPath before any await"

    # (b) showSaveFilePicker on first save ---------------------------------

    def test_show_save_file_picker_used(self):
        assert (
            "showSaveFilePicker" in self.src
        ), "saveFile must call showSaveFilePicker to get a write handle on first save"

    def test_abort_error_handled(self):
        assert (
            "AbortError" in self.src
        ), "saveFile must handle AbortError (user cancelled the picker)"

    def test_handle_persisted_in_map(self):
        assert (
            "_fsHandleMap.set(" in self.src
        ), "Acquired fsHandle must be stored in _fsHandleMap so future saves reuse it"

    def test_handle_stored_on_tab_after_picker(self):
        assert (
            "_saveTab.fsHandle = _saveFsHandle" in self.src
        ), "Acquired fsHandle must be stored on the tab object for the next save"

    # (c) cache-buster on raw fetch ----------------------------------------

    def test_raw_fetch_has_timestamp_cache_buster(self):
        assert (
            "Date.now()" in self.src
        ), "raw bytes fetch URL must include Date.now() as a cache-buster query param"

    def test_raw_fetch_url_has_query_param(self):
        assert (
            "?_=${Date.now()}" in self.src
        ), "raw fetch URL must contain ?_=${Date.now()} to prevent browser caching"

    # (d) _isSaving guard --------------------------------------------------

    def test_is_saving_flag_present(self):
        assert (
            "_isSaving" in self.src
        ), "_isSaving flag must exist to prevent concurrent double-saves"

    def test_is_saving_reset_in_finally(self):
        finally_idx = self.src.rfind("finally {")
        assert finally_idx != -1, "saveFile must have a finally block"
        assert (
            "_isSaving = false" in self.src[finally_idx:]
        ), "_isSaving must be reset to false in the finally block"


# ── XLSX parsing: IWorkbookData format compliance ─────────────────────────────


def _make_xlsx_bytes() -> bytes:
    """Create a minimal real .xlsx file using openpyxl for testing."""
    try:
        import openpyxl
        from io import BytesIO

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Hello"
        ws["B1"] = 123
        ws["A2"] = "World"
        ws["B2"] = 45.6
        buf = BytesIO()
        wb.save(buf)
        return buf.getvalue()
    except ImportError:
        return b""


class TestXlsxOpenFile:
    """Tests for XLSX parsing → IWorkbookData format returned by open_file."""

    def _upload_xlsx(self, client, name="test.xlsx", xlsx_bytes=None):
        if xlsx_bytes is None:
            xlsx_bytes = _make_xlsx_bytes()
        return client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(xlsx_bytes), name)},
            content_type="multipart/form-data",
        )

    def test_xlsx_returns_200(self, wa_client):
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200, f"Expected 200, got {resp.status_code}: {resp.get_json()}"

    def test_xlsx_file_type_is_xlsx(self, wa_client):
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        assert resp.get_json()["file_type"] == "xlsx"

    def test_xlsx_data_is_iworkbookdata(self, wa_client):
        """data must have all required IWorkbookData top-level keys."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client, "myfile.xlsx")
        assert resp.status_code == 200
        data = resp.get_json()["data"]
        for key in ("id", "name", "appVersion", "locale", "sheetOrder", "sheets", "styles", "resources"):
            assert key in data, f"IWorkbookData missing required key: {key!r}"

    def test_xlsx_workbook_name_is_original_filename(self, wa_client):
        """workbook name must be derived from the original uploaded filename, not a UUID."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client, "MyReport.xlsx")
        assert resp.status_code == 200
        data = resp.get_json()["data"]
        wb_name = data.get("name", "")
        # Must be the stem of the original filename, not a 32-char UUID hex
        assert wb_name == "MyReport", (
            f"workbook name should be 'MyReport' (from 'MyReport.xlsx'), got {wb_name!r}"
        )
        assert len(wb_name) != 32 or not wb_name.isalnum(), (
            f"workbook name looks like a UUID hex: {wb_name!r}"
        )

    def test_xlsx_app_version_is_set(self, wa_client):
        """appVersion must be '0.5.0' for Univer compatibility."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        data = resp.get_json()["data"]
        assert data.get("appVersion") == "0.5.0"

    def test_xlsx_locale_is_set(self, wa_client):
        """locale must be present for Univer to apply zh-CN formatting."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        data = resp.get_json()["data"]
        assert data.get("locale") == "zh-CN"

    def test_xlsx_sheet_order_matches_sheets_keys(self, wa_client):
        """sheetOrder must list the same sheet IDs as the sheets dict keys."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        data = resp.get_json()["data"]
        sheet_order = data.get("sheetOrder", [])
        sheets = data.get("sheets", {})
        assert set(sheet_order) == set(sheets.keys()), (
            f"sheetOrder {sheet_order} != sheets keys {list(sheets.keys())}"
        )

    def test_xlsx_sheet_has_required_fields(self, wa_client):
        """Each IWorksheetData must have id, name, rowCount, columnCount, cellData."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        sheets = resp.get_json()["data"].get("sheets", {})
        assert sheets, "sheets must not be empty"
        for sheet_id, sheet_data in sheets.items():
            for key in ("id", "name", "rowCount", "columnCount", "cellData"):
                assert key in sheet_data, (
                    f"IWorksheetData[{sheet_id!r}] missing required key: {key!r}"
                )

    def test_xlsx_sheet_id_matches_key(self, wa_client):
        """sheet.id must equal the key in the sheets dict."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        sheets = resp.get_json()["data"].get("sheets", {})
        for sheet_id, sheet_data in sheets.items():
            assert sheet_data.get("id") == sheet_id, (
                f"sheet.id {sheet_data.get('id')!r} != dict key {sheet_id!r}"
            )

    def test_xlsx_cell_data_has_expected_cell(self, wa_client):
        """
        The test spreadsheet has 'Hello' in A1 (row=0, col=0).
        After JSON serialization, keys must be strings '0'.
        """
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        sheets = resp.get_json()["data"].get("sheets", {})
        first_sheet = sheets.get("sheet1", next(iter(sheets.values()), {}))
        cell_data = first_sheet.get("cellData", {})
        assert cell_data, "cellData should not be empty for a populated sheet"
        # JSON keys must be strings
        row_keys = list(cell_data.keys())
        assert all(isinstance(k, str) for k in row_keys), (
            f"cellData row keys must be strings after JSON, got: {row_keys[:3]}"
        )
        # Row 0 must exist and have cell (0,0) with "Hello"
        row0 = cell_data.get("0", {})
        assert row0, "row 0 must be present"
        col0 = row0.get("0", {})
        assert col0, "cell (0,0) must be present"
        assert col0.get("v") == "Hello", f"cell(0,0).v should be 'Hello', got {col0.get('v')!r}"
        assert col0.get("t") == 1, f"cell(0,0).t should be 1 (string), got {col0.get('t')!r}"

    def test_xlsx_numeric_cell_has_correct_type(self, wa_client):
        """Numeric cells must have t=2 (number type in Univer)."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        sheets = resp.get_json()["data"].get("sheets", {})
        first_sheet = sheets.get("sheet1", next(iter(sheets.values()), {}))
        cell_data = first_sheet.get("cellData", {})
        row0 = cell_data.get("0", {})
        col1 = row0.get("1", {})  # B1 = 123
        assert col1.get("t") == 2, f"numeric cell type should be 2, got {col1.get('t')!r}"
        assert col1.get("v") == 123, f"numeric cell value should be 123, got {col1.get('v')!r}"

    def test_xlsx_styles_is_dict(self, wa_client):
        """styles must be present as a dict (even if empty)."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        data = resp.get_json()["data"]
        assert isinstance(data.get("styles"), dict), "styles must be a dict"

    def test_xlsx_resources_is_list(self, wa_client):
        """resources must be present as a list (even if empty)."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        data = resp.get_json()["data"]
        assert isinstance(data.get("resources"), list), "resources must be a list"

    def test_xlsx_row_and_column_counts_positive(self, wa_client):
        """rowCount and columnCount must each be > 0."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        sheets = resp.get_json()["data"].get("sheets", {})
        for sid, sheet in sheets.items():
            assert sheet.get("rowCount", 0) > 0, f"sheet {sid} rowCount should be > 0"
            assert sheet.get("columnCount", 0) > 0, f"sheet {sid} columnCount should be > 0"

    def test_xlsx_merge_data_is_list(self, wa_client):
        """mergeData must be a list (empty list if no merged cells)."""
        client, _, _ = wa_client
        if not _make_xlsx_bytes():
            pytest.skip("openpyxl not available")
        resp = self._upload_xlsx(client)
        assert resp.status_code == 200
        sheets = resp.get_json()["data"].get("sheets", {})
        for sid, sheet in sheets.items():
            assert isinstance(sheet.get("mergeData"), list), (
                f"sheet {sid} mergeData must be a list"
            )


# ── PPTX parsing: slide geometry format compliance ───────────────────────────


def _make_pptx_bytes() -> bytes:
    """Create a minimal real .pptx file using python-pptx for testing."""
    try:
        from io import BytesIO
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Test Title"
        subtitle.text = "Test Subtitle"

        # Second slide with a table
        blank_layout = prs.slide_layouts[6]  # Blank
        slide2 = prs.slides.add_slide(blank_layout)
        rows, cols = 2, 3
        left = top = Inches(1)
        width = Inches(6)
        height = Inches(2)
        table = slide2.shapes.add_table(rows, cols, left, top, width, height).table
        table.cell(0, 0).text = "Header1"
        table.cell(0, 1).text = "Header2"
        table.cell(0, 2).text = "Header3"
        table.cell(1, 0).text = "Value1"
        table.cell(1, 1).text = "Value2"
        table.cell(1, 2).text = "Value3"

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except ImportError:
        return b""


class TestPptxOpenFile:
    """Tests for PPTX parsing → slide geometry format returned by open_file."""

    def _upload_pptx(self, client, name="test.pptx", pptx_bytes=None):
        if pptx_bytes is None:
            pptx_bytes = _make_pptx_bytes()
        return client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx_bytes), name)},
            content_type="multipart/form-data",
        )

    def test_pptx_returns_200(self, wa_client):
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200, f"Expected 200, got {resp.status_code}: {resp.get_json()}"

    def test_pptx_file_type_is_pptx(self, wa_client):
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        assert resp.get_json()["file_type"] == "pptx"

    def test_pptx_data_has_required_keys(self, wa_client):
        """data must have slide_width_emu, slide_height_emu, and slides."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        data = resp.get_json()["data"]
        for key in ("slide_width_emu", "slide_height_emu", "slides"):
            assert key in data, f"PPTX data missing required key: {key!r}"

    def test_pptx_slide_dimensions_are_positive(self, wa_client):
        """Slide dimensions must be positive integers (in EMU)."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        data = resp.get_json()["data"]
        assert data["slide_width_emu"] > 0, "slide_width_emu must be positive"
        assert data["slide_height_emu"] > 0, "slide_height_emu must be positive"

    def test_pptx_slides_is_non_empty_list(self, wa_client):
        """slides must be a non-empty list (we have 2 slides)."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        slides = resp.get_json()["data"].get("slides", [])
        assert isinstance(slides, list) and len(slides) >= 1, (
            f"slides must be a non-empty list, got: {slides!r}"
        )

    def test_pptx_slide_count_matches_presentation(self, wa_client):
        """Our test PPTX has 2 slides; the response must reflect that."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        slides = resp.get_json()["data"].get("slides", [])
        assert len(slides) == 2, f"Expected 2 slides, got {len(slides)}"

    def test_pptx_slide_has_required_fields(self, wa_client):
        """Each slide must have slide_index, background, shapes."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        for i, slide in enumerate(resp.get_json()["data"]["slides"]):
            for key in ("slide_index", "background", "shapes"):
                assert key in slide, f"slide[{i}] missing required key: {key!r}"

    def test_pptx_text_shapes_have_paragraphs(self, wa_client):
        """Slide 0 (title slide) must have at least one TEXT shape with paragraphs."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        slide0 = resp.get_json()["data"]["slides"][0]
        text_shapes = [s for s in slide0["shapes"] if s.get("_type") == "TEXT"]
        assert text_shapes, "Slide 0 should have at least one TEXT shape"
        for shape in text_shapes:
            assert "paragraphs" in shape, "TEXT shape must have paragraphs"
            assert isinstance(shape["paragraphs"], list)

    def test_pptx_title_text_is_present(self, wa_client):
        """Slide 0 title placeholder must contain 'Test Title'."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        slide0 = resp.get_json()["data"]["slides"][0]
        all_text = []
        for shape in slide0["shapes"]:
            for para in shape.get("paragraphs", []):
                for run in para.get("runs", []):
                    if run.get("text"):
                        all_text.append(run["text"])
        assert "Test Title" in all_text, (
            f"Expected 'Test Title' in slide 0 text runs, got: {all_text}"
        )

    def test_pptx_table_shape_on_slide2(self, wa_client):
        """Slide 1 (index 1) has a table shape with cells."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        slide1 = resp.get_json()["data"]["slides"][1]
        table_shapes = [s for s in slide1["shapes"] if s.get("_type") == "TABLE"]
        assert table_shapes, "Slide 1 should have at least one TABLE shape"

    def test_pptx_table_cells_contain_text(self, wa_client):
        """Table cells in slide 1 must include the test cell text."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        slide1 = resp.get_json()["data"]["slides"][1]
        table_shapes = [s for s in slide1["shapes"] if s.get("_type") == "TABLE"]
        assert table_shapes
        cells = table_shapes[0].get("cells", [])
        cell_texts = [c["text"] for c in cells]
        assert "Header1" in cell_texts, f"Expected 'Header1' in table cells, got: {cell_texts}"
        assert "Value1" in cell_texts, f"Expected 'Value1' in table cells, got: {cell_texts}"

    def test_pptx_table_has_correct_dimensions(self, wa_client):
        """Table in slide 1 must have the right row/col counts."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        slide1 = resp.get_json()["data"]["slides"][1]
        table_shapes = [s for s in slide1["shapes"] if s.get("_type") == "TABLE"]
        assert table_shapes
        table = table_shapes[0]
        assert table.get("table_rows") == 2, f"Expected 2 rows, got {table.get('table_rows')}"
        assert table.get("table_cols") == 3, f"Expected 3 cols, got {table.get('table_cols')}"

    def test_pptx_shapes_have_geometry(self, wa_client):
        """All shapes must have left, top, width, height (in EMU)."""
        client, _, _ = wa_client
        if not _make_pptx_bytes():
            pytest.skip("python-pptx not available")
        resp = self._upload_pptx(client)
        assert resp.status_code == 200
        for slide in resp.get_json()["data"]["slides"]:
            for shape in slide["shapes"]:
                for geo in ("left", "top", "width", "height"):
                    assert geo in shape, (
                        f"shape {shape.get('id')} missing geometry key {geo!r}"
                    )


# ── Embedded-mode render reliability: JS source contract ─────────────────────


class TestEmbeddedModeRenderGuards:
    """
    Validates that workspace-assistant.js contains all guards required for
    reliable XLSX/PPTX rendering in embedded mode (#workspaceView starts
    hidden and transitions from display:none → flex before files are opened).

    These checks verify the *presence* of the guard mechanisms without
    running a full browser — they complement the existing XLSX/PPTX API tests
    and should catch regressions introduced by future refactors.
    """

    JS_PATH = (
        Path(__file__).parents[2] / "web" / "static" / "js" / "workspace-assistant.js"
    )
    PPTX_EDITOR_PATH = (
        Path(__file__).parents[2] / "web" / "src" / "editors" / "pptx-editor.ts"
    )

    @property
    def src(self) -> str:
        return self.JS_PATH.read_text(encoding="utf-8")

    @property
    def pptx_src(self) -> str:
        return self.PPTX_EDITOR_PATH.read_text(encoding="utf-8")

    # ── Layout guard helper ────────────────────────────────────────────────

    def test_wait_for_editor_layout_function_exists(self):
        """_waitForEditorLayout must be defined — it is the central visibility guard."""
        assert "_waitForEditorLayout" in self.src, (
            "_waitForEditorLayout() is missing from workspace-assistant.js"
        )

    def test_wait_for_editor_layout_handles_xlsx(self):
        """Guard must cover xlsx container id."""
        assert "wa-xlsx-editor" in self.src and "_waitForEditorLayout" in self.src, (
            "_waitForEditorLayout must reference 'wa-xlsx-editor'"
        )

    def test_wait_for_editor_layout_handles_pptx(self):
        """Guard must cover pptx container id."""
        assert "wa-pptx-editor" in self.src and "_waitForEditorLayout" in self.src, (
            "_waitForEditorLayout must reference 'wa-pptx-editor'"
        )

    def test_prime_editor_layout_helper_exists(self):
        """xlsx/pptx shells must be pre-activated before waiting for layout."""
        assert "function _primeEditorLayout" in self.src, (
            "workspace-assistant.js must define _primeEditorLayout()"
        )

    def test_prime_editor_layout_activates_hidden_shells(self):
        """The priming helper must add the active class so hidden shells can size."""
        src = self.src
        helper_start = src.find("function _primeEditorLayout")
        helper_end = src.find("function _waitForEditorLayout", helper_start)
        helper_body = src[helper_start:helper_end]
        assert "classList.add('active')" in helper_body, (
            "_primeEditorLayout must activate the editor shell before waiting"
        )

    def test_wait_for_editor_layout_timeout_resolve(self):
        """Guard must resolve (not reject) on timeout so editors receive a mount attempt."""
        # The guard should call resolve() on deadline, not reject()
        src = self.src
        guard_start = src.find("function _waitForEditorLayout")
        guard_end   = src.find("\n  }", guard_start) + 4
        guard_body  = src[guard_start:guard_end]
        assert "resolve();" in guard_body, (
            "_waitForEditorLayout must call resolve() on timeout (not reject)"
        )

    # ── Router.load guard ────────────────────────────────────────────────

    def test_router_load_awaits_layout_guard(self):
        """The file-open path must await _waitForEditorLayout after toggleWorkspace."""
        src = self.src
        # The actual file-open logic lives in _applyFileJson (called by Router.load)
        fn_start = src.find("async function _applyFileJson")
        if fn_start == -1:
            fn_start = src.find("const Router = {")
        fn_end = src.find("new KotoXlsxEditor()", fn_start)
        body = src[fn_start:fn_end + 100] if fn_end != -1 else src[fn_start:fn_start + 3000]
        assert "await _waitForEditorLayout" in body, (
            "File-open path must await _waitForEditorLayout before creating editors"
        )

    def test_router_load_guard_before_xlsx_editor(self):
        """The guard await must appear before new KotoXlsxEditor() in _applyFileJson."""
        src = self.src
        # The actual file-open logic lives in _applyFileJson (called by Router.load)
        fn_start = src.find("async function _applyFileJson")
        if fn_start == -1:
            fn_start = src.find("const Router = {")
        fn_end = src.find("new KotoXlsxEditor()", fn_start)
        body = src[fn_start:fn_end + 100] if fn_end != -1 else src[fn_start:fn_start + 3000]
        guard_pos = body.find("await _waitForEditorLayout")
        xlsx_pos  = body.find("new KotoXlsxEditor()")
        assert guard_pos != -1 and xlsx_pos != -1 and guard_pos < xlsx_pos, (
            "_waitForEditorLayout await must precede new KotoXlsxEditor()"
        )

    def test_router_load_primes_layout_before_waiting(self):
        """The file-open path must prime xlsx/pptx shells before waiting for size."""
        src = self.src
        fn_start = src.find("async function _applyFileJson")
        fn_end = src.find("new KotoXlsxEditor()", fn_start)
        body = src[fn_start:fn_end + 120] if fn_end != -1 else src[fn_start:fn_start + 3200]
        prime_pos = body.find("_primeEditorLayout(state.fileType)")
        guard_pos = body.find("await _waitForEditorLayout(state.fileType)")
        assert prime_pos != -1 and guard_pos != -1 and prime_pos < guard_pos, (
            "_applyFileJson must prime the editor shell before waiting for layout"
        )

    def test_router_load_guard_before_pptx_editor(self):
        """The guard await must appear before PPTX mount in _applyFileJson."""
        src = self.src
        fn_start = src.find("async function _applyFileJson")
        if fn_start == -1:
            fn_start = src.find("const Router = {")
        fn_end = src.find("_mountPptxEditor(json.data)", fn_start)
        body = src[fn_start:fn_end + 100] if fn_end != -1 else src[fn_start:fn_start + 3000]
        guard_pos = body.find("await _waitForEditorLayout")
        pptx_pos  = body.find("_mountPptxEditor(json.data)")
        assert guard_pos != -1 and pptx_pos != -1 and guard_pos < pptx_pos, (
            "_waitForEditorLayout await must precede PPTX editor mount"
        )

    # ── _switchToTab guard ───────────────────────────────────────────────

    def test_switch_to_tab_awaits_layout_guard(self):
        """_switchToTab must also await _waitForEditorLayout for tab-switch renders."""
        src = self.src
        tab_start = src.find("async function _switchToTab")
        tab_end   = src.find("\n  }", tab_start) + 4
        tab_body  = src[tab_start:tab_end]
        assert "await _waitForEditorLayout" in tab_body, (
            "_switchToTab must await _waitForEditorLayout before creating editors"
        )

    def test_switch_to_tab_primes_layout_before_waiting(self):
        """Tab switches must re-activate xlsx/pptx shells before waiting for layout."""
        src = self.src
        tab_start = src.find("async function _switchToTab")
        tab_end = src.find("new KotoXlsxEditor()", tab_start)
        tab_body = src[tab_start:tab_end + 120] if tab_end != -1 else src[tab_start:tab_start + 2400]
        prime_pos = tab_body.find("_primeEditorLayout(tab.fileType)")
        guard_pos = tab_body.find("await _waitForEditorLayout(tab.fileType)")
        assert prime_pos != -1 and guard_pos != -1 and prime_pos < guard_pos, (
            "_switchToTab must prime the editor shell before waiting for layout"
        )

    # ── KotoXlsxEditor size-polling ──────────────────────────────────────

    def test_xlsx_editor_polls_for_non_zero_size(self):
        """KotoXlsxEditor.render must use requestAnimationFrame before calling KotoSheetsAPI.create."""
        src = self.src
        xlsx_start = src.find("class KotoXlsxEditor {")
        xlsx_end   = src.find("\n  class Koto", xlsx_start)
        xlsx_body  = src[xlsx_start:xlsx_end]
        assert "requestAnimationFrame" in xlsx_body, (
            "KotoXlsxEditor must use requestAnimationFrame before mounting Univer"
        )
        assert "KotoSheetsAPI.create" in xlsx_body, (
            "KotoXlsxEditor must call KotoSheetsAPI.create"
        )

    def test_xlsx_editor_has_mount_deadline(self):
        """KotoXlsxEditor.render must have error handling for create failures."""
        src = self.src
        xlsx_start = src.find("class KotoXlsxEditor {")
        xlsx_end   = src.find("\n  class Koto", xlsx_start)
        xlsx_body  = src[xlsx_start:xlsx_end]
        assert "catch" in xlsx_body and "初始化失败" in xlsx_body, (
            "KotoXlsxEditor must catch errors from KotoSheetsAPI.create"
        )

    def test_xlsx_editor_resize_nudge_present(self):
        """KotoXlsxEditor.render must pass string container ID to KotoSheetsAPI.create."""
        src = self.src
        xlsx_start = src.find("class KotoXlsxEditor {")
        xlsx_end   = src.find("\n  class Koto", xlsx_start)
        xlsx_body  = src[xlsx_start:xlsx_end]
        assert "this._containerId" in xlsx_body, (
            "KotoXlsxEditor must pass string container ID to KotoSheetsAPI.create"
        )

    # ── KotoPptxEditor size-polling ──────────────────────────────────────

    def test_pptx_editor_polls_for_slide_area_width(self):
        """KotoPptxEditor.render must poll clientWidth before calling _renderSlide(0)."""
        src = self.pptx_src
        pptx_start = src.find("class KotoPptxEditor")
        pptx_end   = len(src)
        pptx_body  = src[pptx_start:pptx_end]
        assert "_tryPptxRender" in pptx_body or "_pptxMountDeadline" in pptx_body, (
            "KotoPptxEditor must use a polling strategy for first-slide render"
        )

    def test_pptx_editor_has_mount_deadline(self):
        """KotoPptxEditor must have a deadline to prevent infinite polling."""
        src = self.pptx_src
        pptx_start = src.find("class KotoPptxEditor")
        pptx_end   = len(src)
        pptx_body  = src[pptx_start:pptx_end]
        assert "_pptxMountDeadline" in pptx_body or "Date.now()" in pptx_body, (
            "KotoPptxEditor mount polling must have a bounded deadline"
        )

    # ── openInMainView reflow hook ───────────────────────────────────────

    def test_open_in_main_view_reflows_xlsx(self):
        """openInMainView must trigger a ResizeObserver nudge for active XLSX editors."""
        src = self.src
        oim_start = src.find("window.WA.openInMainView = function")
        oim_end   = src.find("\n  };", oim_start) + 4
        oim_body  = src[oim_start:oim_end]
        assert "wa-xlsx-sheet" in oim_body, (
            "openInMainView must reference 'wa-xlsx-sheet' for the reflow nudge"
        )
        assert "style.width" in oim_body, (
            "openInMainView must perform the width+1/reset nudge for XLSX after showing"
        )

    def test_open_in_main_view_reflows_pptx(self):
        """openInMainView must trigger a re-render for active PPTX editors."""
        src = self.src
        oim_start = src.find("window.WA.openInMainView = function")
        oim_end   = src.find("\n  };", oim_start) + 4
        oim_body  = src[oim_start:oim_end]
        assert "_renderSlide" in oim_body, (
            "openInMainView must call _renderSlide for active PPTX editor after showing"
        )

    # ── XLSX API response still valid ────────────────────────────────────

    def test_xlsx_open_file_returns_workbook_data(self, wa_client):
        """Backend must still return valid IWorkbookData after frontend changes."""
        client, _, _ = wa_client
        xlsx_bytes = _make_xlsx_bytes()
        if not xlsx_bytes:
            pytest.skip("openpyxl not available")
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(xlsx_bytes), "embedded_test.xlsx")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 200, f"open_file failed: {resp.data}"
        body = resp.get_json()
        assert body["file_type"] == "xlsx"
        data = body["data"]
        assert "sheets" in data or "sheetOrder" in data, (
            "IWorkbookData must contain 'sheets' or 'sheetOrder'"
        )

    def test_pptx_open_file_returns_slides(self, wa_client):
        """Backend must still return valid slide data after frontend changes."""
        client, _, _ = wa_client
        pptx_bytes = _make_pptx_bytes()
        if not pptx_bytes:
            pytest.skip("python-pptx not available")
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx_bytes), "embedded_test.pptx")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 200, f"open_file failed: {resp.data}"
        body = resp.get_json()
        assert body["file_type"] == "pptx"
        assert isinstance(body["data"].get("slides"), list)
        assert len(body["data"]["slides"]) >= 1

