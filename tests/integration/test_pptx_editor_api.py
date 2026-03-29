# -*- coding: utf-8 -*-
"""
Integration tests for the PPTX editor API blueprint.

Tests the full HTTP lifecycle: upload → list → get → save → download → delete.
Uses a dedicated Flask test app with only the pptx_editor_bp registered.
"""

from __future__ import annotations

import io
import json
import os

import pytest

# ── Fixture: minimal Flask app with pptx_editor_bp ───────────────────────────


@pytest.fixture(scope="module")
def pptx_app(tmp_path_factory):
    """Isolated Flask app for PPTX editor tests; workspace → tmp_path."""
    root = pytest.importorskip("flask")  # skip whole module if Flask absent
    import sys
    from pathlib import Path

    repo_root = str(Path(__file__).resolve().parents[2])
    if repo_root not in sys.path:
        sys.path.insert(0, repo_root)

    tmpdir = tmp_path_factory.mktemp("pptx_workspace")
    os.environ["KOTO_WORKSPACE"] = str(tmpdir)

    from flask import Flask
    from web.blueprints.pptx_editor import pptx_editor_bp

    app = Flask(__name__)
    app.register_blueprint(pptx_editor_bp)
    app.config["TESTING"] = True
    yield app

    os.environ.pop("KOTO_WORKSPACE", None)


@pytest.fixture(scope="module")
def pptx_client(pptx_app):
    return pptx_app.test_client()


# ── PPTX builder helper ───────────────────────────────────────────────────────


def _make_pptx_bytes(slides: list[list[str]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for texts in slides:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        for i, text in enumerate(texts):
            if i == 0 and slide.shapes.title:
                slide.shapes.title.text = text
            else:
                txBox = slide.shapes.add_textbox(
                    Inches(1), Inches(2 + i * 0.5), Inches(8), Inches(0.5)
                )
                txBox.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Upload ────────────────────────────────────────────────────────────────────


@pytest.mark.integration
class TestPptxUpload:
    def test_upload_returns_201(self, pptx_client):
        raw = _make_pptx_bytes([["Hello World", "Body text"]])
        resp = pptx_client.post(
            "/api/pptx/upload",
            data={"file": (io.BytesIO(raw), "test.pptx")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 201
        data = resp.get_json()
        assert "id" in data
        assert data["slideCount"] >= 1

    def test_upload_rejects_wrong_extension(self, pptx_client):
        resp = pptx_client.post(
            "/api/pptx/upload",
            data={"file": (io.BytesIO(b"not a pptx"), "file.docx")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 400

    def test_upload_rejects_no_file(self, pptx_client):
        resp = pptx_client.post(
            "/api/pptx/upload", data={}, content_type="multipart/form-data"
        )
        assert resp.status_code == 400

    def test_upload_pptm_accepted(self, pptx_client):
        raw = _make_pptx_bytes([["PPTM Slide"]])
        resp = pptx_client.post(
            "/api/pptx/upload",
            data={"file": (io.BytesIO(raw), "macro.pptm")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 201


# ── Full lifecycle ────────────────────────────────────────────────────────────


@pytest.mark.integration
class TestPptxLifecycle:
    @pytest.fixture(autouse=True)
    def _upload(self, pptx_client):
        """Upload a PPTX once; provide self.file_id and self.client."""
        raw = _make_pptx_bytes([["Original Title", "Original Body"], ["Slide Two"]])
        resp = pptx_client.post(
            "/api/pptx/upload",
            data={"file": (io.BytesIO(raw), "lifecycle.pptx")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 201
        self.file_id = resp.get_json()["id"]
        self.client = pptx_client

    def test_get_returns_slides(self):
        resp = self.client.get(f"/api/pptx/{self.file_id}")
        assert resp.status_code == 200
        data = resp.get_json()
        assert "slides" in data
        assert len(data["slides"]) == 2

    def test_list_contains_file(self):
        resp = self.client.get("/api/pptx/list")
        assert resp.status_code == 200
        ids = [f["id"] for f in resp.get_json()["files"]]
        assert self.file_id in ids

    def test_save_updates_text(self):
        # Get current data
        data = self.client.get(f"/api/pptx/{self.file_id}").get_json()
        # Modify first text run on slide 0
        for shape in data["slides"][0]["shapes"]:
            if shape["has_text"]:
                for para in shape["paragraphs"]:
                    if para["runs"]:
                        para["runs"][0]["text"] = "EDITED TITLE"
                        break
                break

        resp = self.client.put(
            f"/api/pptx/{self.file_id}",
            json={"slides": data["slides"]},
        )
        assert resp.status_code == 200
        assert resp.get_json()["ok"] is True

        # Reload and verify
        updated = self.client.get(f"/api/pptx/{self.file_id}").get_json()
        texts = [
            run["text"]
            for shape in updated["slides"][0]["shapes"]
            for para in shape.get("paragraphs", [])
            for run in para.get("runs", [])
        ]
        assert "EDITED TITLE" in texts

    def test_download_is_valid_pptx(self):
        from pptx import Presentation

        resp = self.client.get(f"/api/pptx/{self.file_id}/download")
        assert resp.status_code == 200
        assert "application/vnd.openxmlformats" in resp.content_type
        prs = Presentation(io.BytesIO(resp.data))
        assert len(prs.slides) == 2

    def test_download_contains_edited_text(self):
        from pptx import Presentation

        # Save an edit then download
        data = self.client.get(f"/api/pptx/{self.file_id}").get_json()
        for shape in data["slides"][0]["shapes"]:
            if shape["has_text"]:
                for para in shape["paragraphs"]:
                    if para["runs"]:
                        para["runs"][0]["text"] = "DOWNLOAD_CHECK"
                        break
                break
        self.client.put(f"/api/pptx/{self.file_id}", json={"slides": data["slides"]})

        resp = self.client.get(f"/api/pptx/{self.file_id}/download")
        prs = Presentation(io.BytesIO(resp.data))
        all_texts = " ".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if hasattr(shape, "text_frame")
        )
        assert "DOWNLOAD_CHECK" in all_texts

    def test_delete_removes_file(self):
        resp = self.client.delete(f"/api/pptx/{self.file_id}")
        assert resp.status_code == 200
        assert resp.get_json()["ok"] is True

        get_resp = self.client.get(f"/api/pptx/{self.file_id}")
        assert get_resp.status_code == 404


# ── Edge cases ────────────────────────────────────────────────────────────────


@pytest.mark.integration
class TestPptxEdgeCases:
    def test_get_nonexistent(self, pptx_client):
        resp = pptx_client.get("/api/pptx/doesnotexist000")
        assert resp.status_code == 404

    def test_save_nonexistent(self, pptx_client):
        resp = pptx_client.put("/api/pptx/doesnotexist000", json={"slides": []})
        assert resp.status_code == 404

    def test_delete_nonexistent(self, pptx_client):
        resp = pptx_client.delete("/api/pptx/doesnotexist000")
        assert resp.status_code == 404

    def test_list_empty_initially(self, tmp_path_factory, monkeypatch):
        """A workspace with no files → list returns empty array."""
        import sys
        from pathlib import Path

        repo_root = str(Path(__file__).resolve().parents[2])
        if repo_root not in sys.path:
            sys.path.insert(0, repo_root)

        tmpdir = tmp_path_factory.mktemp("empty_ws")
        # Patch the store dir directly on the module to isolate from other tests
        import web.blueprints.pptx_editor as _mod

        original = _mod._STORE_DIR
        _mod._STORE_DIR = str(tmpdir / "pptx-files")
        os.makedirs(_mod._STORE_DIR, exist_ok=True)
        try:
            from flask import Flask

            app2 = Flask(__name__ + "_empty2")
            app2.register_blueprint(_mod.pptx_editor_bp)
            app2.config["TESTING"] = True
            with app2.test_client() as c:
                resp = c.get("/api/pptx/list")
                assert resp.status_code == 200
                assert resp.get_json()["files"] == []
        finally:
            _mod._STORE_DIR = original
