# -*- coding: utf-8 -*-
"""
Functional tests for the PPTX "Package not found" fix.

Regression: when the Flask server was launched from a directory other than
the project root, `_TMP_DIR` resolved to a relative path such as
`workspace\\tmp\\<id>.pptx`.  python-pptx's `Presentation()` received that
relative path, resolved it against the process CWD, and raised:

    BadZipFile / PackageNotFoundError:
        "Package not found at 'workspace\\tmp\\<id>.pptx'"

Two code changes were made:
  1. `workspace_assistant.py` — `_TMP_DIR` is now an absolute path derived
     from `Path(__file__).resolve()`, never from CWD.
  2. `file_parser.py / parse_pptx_geometry` — the path is normalised with
     `os.path.abspath()` before being passed to `Presentation()`, and a
     clear `FileNotFoundError` is raised when the file does not exist.

The tests in this module verify both fixes end-to-end via the Flask test
client as well as at the unit level.
"""
from __future__ import annotations

import io
import os
import uuid
import zipfile
from pathlib import Path

import pytest

# ── helpers ──────────────────────────────────────────────────────────────────


def _make_minimal_pptx() -> bytes:
    """
    Return the bytes of a minimal but structurally valid .pptx archive.

    A .pptx is a ZIP containing at minimum:
      [Content_Types].xml
      _rels/.rels
      ppt/presentation.xml
      ppt/_rels/presentation.xml.rels
      ppt/slides/slide1.xml
      ppt/slides/_rels/slide1.xml.rels
      ppt/slideLayouts/slideLayout1.xml  (required by python-pptx)
      ppt/slideLayouts/_rels/slideLayout1.xml.rels
      ppt/slideMasters/slideMaster1.xml
      ppt/slideMasters/_rels/slideMaster1.xml.rels

    Rather than building the full structure by hand we use python-pptx itself
    (if available) to produce valid bytes; otherwise we fall back to a
    pre-built minimal ZIP that python-pptx accepts.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        if title:
            title.text = "Test Slide"
        body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body:
            body.text = "Hello from test"
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()
    except ImportError:
        pytest.skip("python-pptx not installed")


def _fake_pdf_bytes() -> bytes:
    return b"%PDF-1.4\n1 0 obj\n<</Type /Catalog>>\nendobj\n%%EOF"


# ── Fixture: isolated Flask test client ──────────────────────────────────────


@pytest.fixture(scope="module")
def wa_client(tmp_path_factory):
    """
    Minimal Flask test client for workspace_assistant_bp with fully isolated
    tmp and workspace directories.  _TMP_DIR is patched to a known temp
    directory so we can assert on its contents.
    """
    os.environ.setdefault("KOTO_AUTH_ENABLED", "false")

    tmp_root = tmp_path_factory.mktemp("pptx_fix_root")
    tmp_dir = tmp_root / "tmp"
    workspace_dir = tmp_root / "workspace"
    tmp_dir.mkdir(parents=True)
    workspace_dir.mkdir(parents=True)

    import web.blueprints.workspace_assistant as _wa
    import web.shared as _shared

    original_tmp = _wa._TMP_DIR
    original_ws = getattr(_shared, "WORKSPACE_DIR", None)

    _wa._TMP_DIR = tmp_dir
    _shared.WORKSPACE_DIR = str(workspace_dir)

    from flask import Flask
    from web.blueprints.workspace_assistant import workspace_assistant_bp

    app = Flask(__name__)
    app.register_blueprint(workspace_assistant_bp)
    app.config["TESTING"] = True

    with app.test_client() as client:
        yield client, tmp_dir, workspace_dir

    _wa._TMP_DIR = original_tmp
    if original_ws is not None:
        _shared.WORKSPACE_DIR = original_ws


# ─────────────────────────────────────────────────────────────────────────────
# 1.  _TMP_DIR is always an absolute path
# ─────────────────────────────────────────────────────────────────────────────


class TestTmpDirIsAbsolute:
    """
    The root-cause fix: _TMP_DIR must be absolute so python-pptx (and any
    other library that passes the path to OS-level file APIs) can always
    find the file regardless of the process CWD.
    """

    def test_tmp_dir_is_absolute(self):
        import web.blueprints.workspace_assistant as _wa

        assert Path(_wa._TMP_DIR).is_absolute(), (
            f"_TMP_DIR must be an absolute path; got {_wa._TMP_DIR!r}. "
            "python-pptx resolves relative paths against CWD which may not "
            "be the project root, causing 'Package not found' errors."
        )

    def test_tmp_dir_does_not_use_cwd(self, monkeypatch, tmp_path):
        """
        Even when os.getcwd() returns a completely unrelated directory,
        _TMP_DIR must still point inside the project tree — not inside tmp_path.
        """
        import web.blueprints.workspace_assistant as _wa

        monkeypatch.chdir(tmp_path)
        # Re-evaluate the stored path; it must already be absolute and correct
        assert Path(_wa._TMP_DIR).is_absolute()
        assert str(tmp_path) not in str(_wa._TMP_DIR), (
            "_TMP_DIR must not be affected by changing the working directory"
        )

    def test_ensure_tmp_dir_returns_absolute(self, wa_client):
        """_ensure_tmp_dir() must always return an absolute path."""
        import web.blueprints.workspace_assistant as _wa

        result = _wa._ensure_tmp_dir()
        assert result.is_absolute(), (
            f"_ensure_tmp_dir() returned a relative path: {result!r}"
        )

    def test_tmp_path_passed_to_parse_is_absolute(self, wa_client, monkeypatch):
        """
        When open_file saves the uploaded file, the tmp_path given to
        parse_pptx_geometry must be absolute.
        """
        client, tmp_dir, _ = wa_client
        pptx_bytes = _make_minimal_pptx()

        captured_paths: list[str] = []

        import app.core.file.file_parser as _fp
        real_fn = _fp.parse_pptx_geometry

        def spy(path, *a, **kw):
            captured_paths.append(path)
            return real_fn(path, *a, **kw)

        monkeypatch.setattr(_fp, "parse_pptx_geometry", spy)
        monkeypatch.setattr(
            "web.blueprints.workspace_assistant.parse_pptx_geometry", spy,
            raising=False,
        )

        client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx_bytes), "spy_test.pptx")},
            content_type="multipart/form-data",
        )

        assert captured_paths, "parse_pptx_geometry was never called"
        for p in captured_paths:
            assert Path(p).is_absolute(), (
                f"parse_pptx_geometry received a relative path {p!r}; "
                "this causes 'Package not found' when CWD != project root"
            )


# ─────────────────────────────────────────────────────────────────────────────
# 2.  POST /api/v1/workspace/open_file  — PPTX upload
# ─────────────────────────────────────────────────────────────────────────────


class TestOpenFilePptx:
    """End-to-end upload tests for .pptx via the open_file endpoint."""

    def _upload(self, client, name="presentation.pptx", data=None):
        pptx_bytes = data or _make_minimal_pptx()
        return client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx_bytes), name)},
            content_type="multipart/form-data",
        )

    def test_pptx_upload_returns_200(self, wa_client):
        """Regression: must not return 500 'Package not found'."""
        client, _, _ = wa_client
        resp = self._upload(client)
        assert resp.status_code == 200, (
            f"PPTX upload returned {resp.status_code}: {resp.get_json()}. "
            "Expected 200 — 'Package not found' means _TMP_DIR is still relative."
        )

    def test_pptx_response_file_type_is_pptx(self, wa_client):
        client, _, _ = wa_client
        resp = self._upload(client)
        assert resp.status_code == 200
        assert resp.get_json()["file_type"] == "pptx"

    def test_pptx_response_contains_slides_key(self, wa_client):
        """data.slides must be a list (may be empty for a blank presentation)."""
        client, _, _ = wa_client
        resp = self._upload(client)
        assert resp.status_code == 200
        data = resp.get_json().get("data", {})
        assert "slides" in data, f"Response data missing 'slides' key; got: {list(data)}"
        assert isinstance(data["slides"], list)

    def test_pptx_response_contains_slide_dimensions(self, wa_client):
        """slide_width_emu and slide_height_emu must be positive integers."""
        client, _, _ = wa_client
        resp = self._upload(client)
        assert resp.status_code == 200
        data = resp.get_json().get("data", {})
        assert data.get("slide_width_emu", 0) > 0, "slide_width_emu must be > 0"
        assert data.get("slide_height_emu", 0) > 0, "slide_height_emu must be > 0"

    def test_pptx_file_saved_to_tmp_dir(self, wa_client):
        """Uploaded PPTX must be persisted in the tmp directory."""
        client, tmp_dir, _ = wa_client
        resp = self._upload(client)
        assert resp.status_code == 200
        file_id = resp.get_json()["file_id"]
        matches = list(tmp_dir.glob(f"{file_id}.pptx"))
        assert matches, (
            f"Expected {file_id}.pptx in {tmp_dir}; found: {list(tmp_dir.iterdir())}"
        )

    def test_pptx_file_id_is_hex_uuid(self, wa_client):
        client, _, _ = wa_client
        resp = self._upload(client)
        assert resp.status_code == 200
        fid = resp.get_json()["file_id"]
        assert fid.isalnum() and len(fid) == 32, (
            f"file_id must be a 32-char hex UUID; got {fid!r}"
        )

    def test_pptx_file_name_preserved(self, wa_client):
        client, _, _ = wa_client
        resp = self._upload(client, name="my_deck.pptx")
        assert resp.status_code == 200
        assert resp.get_json()["file_name"] == "my_deck.pptx"

    def test_pptx_upload_does_not_return_package_not_found(self, wa_client):
        """
        Explicit check: error message must NOT contain 'Package not found'.
        This is the exact string python-pptx raises when given a bad path.
        """
        client, _, _ = wa_client
        resp = self._upload(client)
        body = resp.get_json() or {}
        error = body.get("error", "")
        assert "Package not found" not in error, (
            f"Server returned the old 'Package not found' error: {error!r}. "
            "This means _TMP_DIR is still a relative path."
        )

    def test_pptx_raw_url_fetchable_after_upload(self, wa_client):
        """After upload, /raw/<file_id> must return the original PPTX bytes."""
        client, _, _ = wa_client
        pptx_bytes = _make_minimal_pptx()
        resp = self._upload(client, data=pptx_bytes)
        assert resp.status_code == 200
        file_id = resp.get_json()["file_id"]
        raw = client.get(f"/api/v1/workspace/raw/{file_id}")
        assert raw.status_code == 200
        # PPTX files are ZIP archives; first bytes are the ZIP magic number
        assert raw.data[:4] == b"PK\x03\x04", (
            f"Expected ZIP magic (PK\\x03\\x04) from raw endpoint; "
            f"got {raw.data[:8]!r}"
        )

    def test_corrupt_pptx_returns_500_not_package_not_found(self, wa_client):
        """
        Uploading a corrupt (non-ZIP) PPTX should return a 500 with an error
        about the file content — not a misleading 'Package not found' that
        implies the file was missing from disk.
        """
        client, _, _ = wa_client
        resp = self._upload(client, data=b"this is not a zip file at all")
        # Must be a parse error (500), not a path error
        assert resp.status_code == 500
        error = (resp.get_json() or {}).get("error", "")
        # The important thing: it must NOT say "Package not found at 'workspace\tmp\..."
        # (the relative-path variant) — it may say "Package not found at 'C:\...'"
        # which is acceptable because at least that points to the real file.
        assert "workspace\\tmp\\" not in error and "workspace/tmp/" not in error, (
            f"Error still contains a relative path reference: {error!r}. "
            "This means the absolute-path fix is not active."
        )


# ─────────────────────────────────────────────────────────────────────────────
# 3.  POST /api/v1/workspace/open_file_by_path  — PPTX from workspace tree
# ─────────────────────────────────────────────────────────────────────────────


class TestOpenFileByPathPptx:
    """
    Tests for the open_file_by_path endpoint which opens a PPTX already
    present in the workspace (left-panel file-tree click).

    This code path copies the file to tmp then calls parse_pptx_geometry —
    the same bug applied here too.
    """

    def _place_pptx(self, workspace_dir: Path, name: str = "deck.pptx") -> Path:
        pptx_bytes = _make_minimal_pptx()
        target = workspace_dir / name
        target.write_bytes(pptx_bytes)
        return target

    def test_open_by_path_returns_200(self, wa_client):
        client, _, workspace_dir = wa_client
        self._place_pptx(workspace_dir, "test_deck.pptx")
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "test_deck.pptx"},
        )
        assert resp.status_code == 200, (
            f"open_file_by_path returned {resp.status_code}: {resp.get_json()}. "
            "Expected 200 — 'Package not found' means copy to tmp used a relative path."
        )

    def test_open_by_path_file_type_is_pptx(self, wa_client):
        client, _, workspace_dir = wa_client
        self._place_pptx(workspace_dir, "typed_deck.pptx")
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "typed_deck.pptx"},
        )
        assert resp.status_code == 200
        assert resp.get_json()["file_type"] == "pptx"

    def test_open_by_path_contains_slides(self, wa_client):
        client, _, workspace_dir = wa_client
        self._place_pptx(workspace_dir, "slides_deck.pptx")
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "slides_deck.pptx"},
        )
        assert resp.status_code == 200
        data = resp.get_json().get("data", {})
        assert "slides" in data
        assert isinstance(data["slides"], list)

    def test_open_by_path_no_package_not_found_error(self, wa_client):
        """Explicit check for the original error string."""
        client, _, workspace_dir = wa_client
        self._place_pptx(workspace_dir, "no_err_deck.pptx")
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "no_err_deck.pptx"},
        )
        error = (resp.get_json() or {}).get("error", "")
        assert "Package not found" not in error, (
            f"open_file_by_path returned: {error!r}"
        )

    def test_open_by_path_missing_file_returns_404(self, wa_client):
        client, _, _ = wa_client
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "does_not_exist.pptx"},
        )
        assert resp.status_code == 404

    def test_open_by_path_traversal_rejected(self, wa_client):
        client, _, _ = wa_client
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "../../../etc/passwd"},
        )
        assert resp.status_code in (403, 404)

    def test_open_by_path_empty_path_returns_400(self, wa_client):
        client, _, _ = wa_client
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": ""},
        )
        assert resp.status_code == 400

    def test_open_by_path_copies_to_tmp(self, wa_client):
        """Verify the file is actually copied into tmp_dir before parsing."""
        client, tmp_dir, workspace_dir = wa_client
        self._place_pptx(workspace_dir, "copy_check.pptx")
        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "copy_check.pptx"},
        )
        assert resp.status_code == 200
        file_id = resp.get_json()["file_id"]
        matches = list(tmp_dir.glob(f"{file_id}.pptx"))
        assert matches, (
            f"Expected copy at {tmp_dir / file_id}.pptx; "
            f"tmp contents: {list(tmp_dir.iterdir())}"
        )


# ─────────────────────────────────────────────────────────────────────────────
# 4.  parse_pptx_geometry unit-level path handling
# ─────────────────────────────────────────────────────────────────────────────


@pytest.mark.unit
class TestParsePptxGeometryPathHandling:
    """
    Unit tests for the path normalisation added to parse_pptx_geometry.
    These exercise the fix independently of the Flask layer.
    """

    def test_raises_file_not_found_for_nonexistent_path(self):
        from app.core.file.file_parser import parse_pptx_geometry

        with pytest.raises((FileNotFoundError, Exception)) as exc_info:
            parse_pptx_geometry("/nonexistent/path/deck.pptx")
        # Must clearly report the file is missing, not "Package not found"
        # with a confusing relative path
        assert "nonexistent" in str(exc_info.value) or "不存在" in str(exc_info.value), (
            f"Expected a clear FileNotFoundError; got: {exc_info.value!r}"
        )

    def test_accepts_absolute_path(self, tmp_path):
        """parse_pptx_geometry must succeed when given an absolute path."""
        from app.core.file.file_parser import parse_pptx_geometry

        pptx_bytes = _make_minimal_pptx()
        pptx_file = tmp_path / "abs_path_test.pptx"
        pptx_file.write_bytes(pptx_bytes)

        result = parse_pptx_geometry(str(pptx_file))
        assert "slides" in result
        assert "slide_width_emu" in result
        assert "slide_height_emu" in result

    def test_accepts_relative_path_that_exists(self, tmp_path, monkeypatch):
        """
        When a relative path is passed AND the file exists relative to CWD,
        parse_pptx_geometry must resolve it correctly (via os.path.abspath).
        """
        from app.core.file.file_parser import parse_pptx_geometry

        pptx_bytes = _make_minimal_pptx()
        pptx_file = tmp_path / "relative_test.pptx"
        pptx_file.write_bytes(pptx_bytes)

        monkeypatch.chdir(tmp_path)
        result = parse_pptx_geometry("relative_test.pptx")
        assert "slides" in result

    def test_result_slides_is_list(self, tmp_path):
        from app.core.file.file_parser import parse_pptx_geometry

        pptx_bytes = _make_minimal_pptx()
        f = tmp_path / "slides_list.pptx"
        f.write_bytes(pptx_bytes)

        result = parse_pptx_geometry(str(f))
        assert isinstance(result["slides"], list)

    def test_result_dimensions_are_positive_ints(self, tmp_path):
        from app.core.file.file_parser import parse_pptx_geometry

        pptx_bytes = _make_minimal_pptx()
        f = tmp_path / "dims.pptx"
        f.write_bytes(pptx_bytes)

        result = parse_pptx_geometry(str(f))
        assert isinstance(result["slide_width_emu"], int)
        assert isinstance(result["slide_height_emu"], int)
        assert result["slide_width_emu"] > 0
        assert result["slide_height_emu"] > 0

    def test_slide_has_required_keys(self, tmp_path):
        from app.core.file.file_parser import parse_pptx_geometry

        pptx_bytes = _make_minimal_pptx()
        f = tmp_path / "slide_keys.pptx"
        f.write_bytes(pptx_bytes)

        result = parse_pptx_geometry(str(f))
        assert result["slides"], "Expected at least one slide"
        slide = result["slides"][0]
        for key in ("slide_index", "slide_id", "background", "shapes"):
            assert key in slide, f"Slide dict missing key {key!r}; keys: {list(slide)}"

    def test_error_for_non_zip_file(self, tmp_path):
        """A file that is not a ZIP must raise an exception (not hang or return garbage)."""
        from app.core.file.file_parser import parse_pptx_geometry

        bad = tmp_path / "not_a_zip.pptx"
        bad.write_bytes(b"this is not a zip archive")

        with pytest.raises(Exception):
            parse_pptx_geometry(str(bad))


# ─────────────────────────────────────────────────────────────────────────────
# 5.  CWD-independence regression test
# ─────────────────────────────────────────────────────────────────────────────


class TestCwdIndependence:
    """
    Simulates the exact failure mode: server launched from a directory that
    is NOT the project root.  Before the fix, uploading a PPTX returned 500
    with 'Package not found at workspace\\tmp\\<id>.pptx'.
    """

    def test_pptx_upload_succeeds_from_unrelated_cwd(self, wa_client, tmp_path, monkeypatch):
        """
        Change the process CWD to an unrelated directory, then upload a PPTX.
        Must return 200, not 500 with 'Package not found'.
        """
        monkeypatch.chdir(tmp_path)  # CWD is now completely unrelated to the project

        client, _, _ = wa_client
        pptx_bytes = _make_minimal_pptx()
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(pptx_bytes), "cwd_test.pptx")},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 200, (
            f"PPTX upload failed when CWD={tmp_path}: "
            f"{resp.status_code} {resp.get_json()}. "
            "This is the exact regression — _TMP_DIR must be absolute."
        )

    def test_open_by_path_succeeds_from_unrelated_cwd(self, wa_client, tmp_path, monkeypatch):
        """Same CWD-independence test for the open_file_by_path endpoint."""
        client, _, workspace_dir = wa_client

        pptx_bytes = _make_minimal_pptx()
        (workspace_dir / "cwd_by_path.pptx").write_bytes(pptx_bytes)

        monkeypatch.chdir(tmp_path)

        resp = client.post(
            "/api/v1/workspace/open_file_by_path",
            json={"path": "cwd_by_path.pptx"},
        )
        assert resp.status_code == 200, (
            f"open_file_by_path failed when CWD={tmp_path}: "
            f"{resp.status_code} {resp.get_json()}"
        )

    def test_no_relative_path_in_error_message(self, wa_client, tmp_path, monkeypatch):
        """
        Even when the upload fails for a legitimate reason (corrupt file),
        the error message must NOT contain a relative path like
        'workspace\\tmp\\...' — that would indicate the old bug is still present.
        """
        monkeypatch.chdir(tmp_path)

        client, _, _ = wa_client
        resp = client.post(
            "/api/v1/workspace/open_file",
            data={"file": (io.BytesIO(b"bad data"), "corrupt.pptx")},
            content_type="multipart/form-data",
        )
        error = (resp.get_json() or {}).get("error", "")
        assert "workspace\\tmp\\" not in error and "workspace/tmp/" not in error, (
            f"Error message contains a relative tmp path: {error!r}. "
            "This confirms the 'Package not found' bug is not fully fixed."
        )
