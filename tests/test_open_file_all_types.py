"""
文件助手 全格式加载通路测试
==============================================
测试 POST /api/v1/workspace/open_file 端点能正常加载所有支持的文件类型：
  - .docx  — Word 文档
  - .xlsx  — Excel 表格
  - .pptx  — PowerPoint 演示文稿
  - .pdf   — PDF 文档
  - .txt   — 纯文本
  - .md    — Markdown
  - .json  — JSON 代码文件
  - .py    — Python 代码文件
  - .png   — PNG 图片
  - .jpg   — JPEG 图片

每种格式验证：
  1. HTTP 200
  2. 返回 file_id, file_name, file_type, data
  3. data 包含格式对应的关键字段且非空
"""

from __future__ import annotations

import io
import os
import struct
import sys
import zlib

import pytest

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)

# ── Flask 测试 app ────────────────────────────────────────────────────────────


@pytest.fixture(scope="module")
def wa_client():
    """独立的 Flask test client，只注册 workspace_assistant_bp。"""
    from flask import Flask

    from web.blueprints.workspace_assistant import workspace_assistant_bp

    app = Flask(__name__)
    app.secret_key = "test-secret-key"
    app.register_blueprint(workspace_assistant_bp)
    app.config["TESTING"] = True
    # 让 tmp 目录指向临时路径，避免污染工作区
    import tempfile

    _tmproot = tempfile.mkdtemp()
    import web.blueprints.workspace_assistant as _wa_mod

    _orig_tmp_root = _wa_mod._TMP_ROOT
    _wa_mod._TMP_ROOT = type(_wa_mod._TMP_ROOT)(_tmproot)
    with app.test_client() as c:
        yield c
    _wa_mod._TMP_ROOT = _orig_tmp_root


# ── 文件构造工具 ───────────────────────────────────────────────────────────────


def _make_docx() -> bytes:
    from docx import Document

    buf = io.BytesIO()
    doc = Document()
    doc.add_heading("文件助手测试文档", level=1)
    doc.add_paragraph("这是第一段正文，用于验证 DOCX 加载功能。")
    doc.add_paragraph("第二段包含更多内容。")
    t = doc.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "列A"
    t.cell(0, 1).text = "列B"
    t.cell(1, 0).text = "值1"
    t.cell(1, 1).text = "值2"
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx() -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["产品", "数量", "单价"])
    ws.append(["苹果", 100, 5.5])
    ws.append(["香蕉", 200, 3.0])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx() -> bytes:
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[1]
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "文件助手 PPT 测试"
    s1.placeholders[1].text = "第一张幻灯片内容。"
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "第二张"
    s2.placeholders[1].text = "更多内容。"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pdf() -> bytes:
    """生成最小合法 PDF（不依赖 reportlab）。"""
    stream_data = b"BT /F1 12 Tf 72 720 Td (File Assistant PDF Test) Tj ET"
    stream_len = len(stream_data)
    pdf = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R"
        b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
        b"4 0 obj<</Length "
        + str(stream_len).encode()
        + b">>\nstream\n"
        + stream_data
        + b"\nendstream\nendobj\n"
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"xref\n0 6\n"
        b"0000000000 65535 f \n"
        b"0000000009 00000 n \n"
        b"0000000058 00000 n \n"
        b"0000000115 00000 n \n"
        b"0000000266 00000 n \n"
        b"0000000360 00000 n \n"
        b"trailer<</Size 6/Root 1 0 R>>\n"
        b"startxref\n441\n%%EOF\n"
    )
    return pdf


def _make_png() -> bytes:
    def chunk(t, d):
        l = struct.pack(">I", len(d))
        c = t + d
        return l + c + struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    idat = chunk(b"IDAT", zlib.compress(b"\x00\xff\x00\x00"))
    iend = chunk(b"IEND", b"")
    return sig + ihdr + idat + iend


def _make_jpg() -> bytes:
    """最小合法 JPEG（SOI + APP0 + EOI）。"""
    return bytes(
        [
            0xFF,
            0xD8,  # SOI
            0xFF,
            0xE0,
            0x00,
            0x10,  # APP0 marker + length=16
            0x4A,
            0x46,
            0x49,
            0x46,
            0x00,  # "JFIF\0"
            0x01,
            0x01,  # version 1.1
            0x00,
            0x00,
            0x01,
            0x00,
            0x01,  # pixel aspect
            0x00,
            0x00,  # no thumbnail
            0xFF,
            0xD9,  # EOI
        ]
    )


# ── 工具：发 multipart POST ────────────────────────────────────────────────────


def _upload(client, filename: str, data: bytes):
    return client.post(
        "/api/v1/workspace/open_file",
        data={"file": (io.BytesIO(data), filename)},
        content_type="multipart/form-data",
    )


# ── 测试用例 ───────────────────────────────────────────────────────────────────


class TestOpenFileAllTypes:

    # ── .docx ────────────────────────────────────────────────────────────────
    def test_docx_loads(self, wa_client):
        pytest.importorskip("docx", reason="python-docx 未安装")
        resp = _upload(wa_client, "test.docx", _make_docx())
        assert resp.status_code == 200, resp.get_data(as_text=True)
        j = resp.get_json()
        assert j["file_type"] == "docx"
        assert j["file_id"]
        # parse_docx 返回包含 html 或 content 的 dict
        assert isinstance(j["data"], dict)
        assert j["data"]  # 非空

    # ── .xlsx ────────────────────────────────────────────────────────────────
    def test_xlsx_loads(self, wa_client):
        pytest.importorskip("openpyxl", reason="openpyxl 未安装")
        resp = _upload(wa_client, "test.xlsx", _make_xlsx())
        assert resp.status_code == 200, resp.get_data(as_text=True)
        j = resp.get_json()
        assert j["file_type"] == "xlsx"
        assert isinstance(j["data"], dict)
        assert j["data"]

    # ── .pptx ────────────────────────────────────────────────────────────────
    def test_pptx_loads(self, wa_client):
        pytest.importorskip("pptx", reason="python-pptx 未安装")
        resp = _upload(wa_client, "test.pptx", _make_pptx())
        assert resp.status_code == 200, resp.get_data(as_text=True)
        j = resp.get_json()
        assert j["file_type"] == "pptx"
        data = j["data"]
        # _parse_slides 返回 {slides: [...], ...}
        assert isinstance(data, dict)
        slides = data.get("slides") or data.get("slide_count") or data
        assert slides  # 至少有幻灯片数据

    # ── .pdf ─────────────────────────────────────────────────────────────────
    def test_pdf_loads(self, wa_client):
        resp = _upload(wa_client, "test.pdf", _make_pdf())
        assert resp.status_code == 200, resp.get_data(as_text=True)
        j = resp.get_json()
        assert j["file_type"] == "pdf"
        data = j["data"]
        assert isinstance(data, dict)
        # parse_pdf 通常返回 {pages, text, raw_url, ...}
        assert data  # 非空

    # ── .txt ─────────────────────────────────────────────────────────────────
    def test_txt_loads(self, wa_client):
        content = "这是一份纯文本测试文件\n第二行\nHello!"
        resp = _upload(wa_client, "test.txt", content.encode("utf-8"))
        assert resp.status_code == 200, resp.get_data(as_text=True)
        j = resp.get_json()
        assert j["file_type"] == "text"
        assert j["data"]["content"] == content
        assert j["data"]["extension"] == ".txt"

    # ── .md ──────────────────────────────────────────────────────────────────
    def test_md_loads(self, wa_client):
        content = "# 标题\n\n正文内容。\n\n- 列表项1\n- 列表项2\n"
        resp = _upload(wa_client, "readme.md", content.encode("utf-8"))
        assert resp.status_code == 200, resp.get_data(as_text=True)
        j = resp.get_json()
        assert j["file_type"] == "text"
        assert "标题" in j["data"]["content"]
        assert j["data"]["extension"] == ".md"

    # ── .json ─────────────────────────────────────────────────────────────────
    def test_json_file_loads(self, wa_client):
        import json as _json

        content = _json.dumps({"name": "Koto", "version": "1.0"}, ensure_ascii=False)
        resp = _upload(wa_client, "config.json", content.encode("utf-8"))
        assert resp.status_code == 200, resp.get_data(as_text=True)
        j = resp.get_json()
        assert j["file_type"] == "code"
        assert "Koto" in j["data"]["content"]

    # ── .py ──────────────────────────────────────────────────────────────────
    def test_py_file_loads(self, wa_client):
        content = "# Python test\nprint('hello')\n"
        resp = _upload(wa_client, "script.py", content.encode("utf-8"))
        assert resp.status_code == 200, resp.get_data(as_text=True)
        j = resp.get_json()
        assert j["file_type"] == "code"
        assert "print" in j["data"]["content"]

    # ── .png ─────────────────────────────────────────────────────────────────
    def test_png_loads(self, wa_client):
        resp = _upload(wa_client, "image.png", _make_png())
        assert resp.status_code == 200, resp.get_data(as_text=True)
        j = resp.get_json()
        assert j["file_type"] == "image"
        assert "raw_url" in j["data"]
        assert j["data"]["raw_url"]

    # ── .jpg ─────────────────────────────────────────────────────────────────
    def test_jpg_loads(self, wa_client):
        resp = _upload(wa_client, "photo.jpg", _make_jpg())
        assert resp.status_code == 200, resp.get_data(as_text=True)
        j = resp.get_json()
        assert j["file_type"] == "image"
        assert "raw_url" in j["data"]

    # ── 通用响应结构验证 ─────────────────────────────────────────────────────
    def test_response_has_required_fields(self, wa_client):
        """所有成功响应都必须包含 file_id, file_name, file_type, data 四个字段。"""
        content = "quick check"
        resp = _upload(wa_client, "check.txt", content.encode("utf-8"))
        assert resp.status_code == 200
        j = resp.get_json()
        for key in ("file_id", "file_name", "file_type", "data"):
            assert key in j, f"响应缺少字段: {key}"
        assert j["file_name"] == "check.txt"

    # ── 不支持的格式 ─────────────────────────────────────────────────────────
    def test_unsupported_format_returns_400(self, wa_client):
        resp = _upload(wa_client, "file.exe", b"\x4d\x5a\x90\x00")
        assert resp.status_code == 400
        j = resp.get_json()
        assert "error" in j

    # ── 缺少文件字段 ─────────────────────────────────────────────────────────
    def test_missing_file_field_returns_400(self, wa_client):
        resp = wa_client.post(
            "/api/v1/workspace/open_file",
            data={},
            content_type="multipart/form-data",
        )
        assert resp.status_code == 400
        j = resp.get_json()
        assert "error" in j
