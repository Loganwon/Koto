#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Skills 功能验证测试脚本
验证 10 个新文件技能引用的每个工具/路径是否真实可用。
运行方式: .venv\Scripts\python.exe scripts\test_file_skills.py
"""
import os, sys, json, traceback, shutil, tempfile
from pathlib import Path

# 确保项目根目录在 sys.path
ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

PASS = "✅ PASS"
FAIL = "❌ FAIL"
WARN = "⚠️ WARN"

results = []

def report(label, status, detail=""):
    icon = {"pass": PASS, "fail": FAIL, "warn": WARN}.get(status, status)
    msg = f"{icon}  {label}"
    if detail:
        msg += f"\n       {detail}"
    print(msg)
    results.append((label, status, detail))

# ─── 准备测试文件 ────────────────────────────────────────────────────────────
tmp_dir = Path(tempfile.mkdtemp(prefix="koto_skill_test_"))
print(f"\n🔧 临时测试目录: {tmp_dir}\n{'='*60}\n")

# 创建 TXT
txt_path = tmp_dir / "test_report.txt"
txt_path.write_text(
    "# 2026年Q1销售报告\n\n## 数据汇总\n总销售额: 1,250,000 元\n"
    "同比增长: 15.3%\n\n## 区域分析\n华东区: 450,000 元\n华南区: 380,000 元\n"
    "华北区: 420,000 元\n\n## 问题与风险\n供应链延迟影响3月出货量\n客户投诉率上升2%",
    encoding="utf-8"
)

# 创建 CSV
csv_path = tmp_dir / "sales_data.csv"
csv_path.write_text(
    "月份,销售额,目标,完成率,备注\n"
    "1月,380000,350000,108.6%,超额完成\n"
    "2月,420000,400000,105.0%,正常\n"
    "3月,450000,500000,90.0%,供应链延迟\n"
    "合计,1250000,1250000,100.0%,\n",
    encoding="utf-8-sig"
)

# 创建 DOCX (需要 python-docx)
docx_path = tmp_dir / "test_document.docx"
try:
    from docx import Document as _Doc
    d = _Doc()
    d.add_heading("测试合同文档", 0)
    d.add_heading("第一章 总则", 1)
    d.add_paragraph("本合同由甲乙双方依据相关法律法规签订，具有法律效力。")
    d.add_heading("第二章 付款条款", 1)
    d.add_paragraph("合同总金额：人民币 200,000 元整。")
    t = d.add_table(rows=1, cols=3)
    t.style = 'Table Grid'
    h = t.rows[0].cells
    h[0].text = "期次"; h[1].text = "金额"; h[2].text = "时间"
    r = t.add_row().cells
    r[0].text = "首付款"; r[1].text = "60,000元"; r[2].text = "合同签订后3日内"
    d.save(str(docx_path))
    report("创建 DOCX 测试文件", "pass", str(docx_path))
except Exception as e:
    docx_path = None
    report("创建 DOCX 测试文件", "fail", str(e))

# 创建 XLSX (需要 openpyxl)
xlsx_path = tmp_dir / "test_excel.xlsx"
try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "销售数据"
    headers = ["月份", "销售额", "目标", "完成率"]
    fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    for i, h in enumerate(headers, 1):
        c = ws.cell(1, i, h)
        c.fill = fill
        c.font = Font(bold=True, color="FFFFFF")
    data = [("1月", 380000, 350000, "108.6%"), ("2月", 420000, 400000, "105.0%"), ("3月", 450000, 500000, "90.0%")]
    for row_i, row in enumerate(data, 2):
        for col_i, val in enumerate(row, 1):
            ws.cell(row_i, col_i, val)
    wb.save(str(xlsx_path))
    report("创建 XLSX 测试文件", "pass", str(xlsx_path))
except Exception as e:
    xlsx_path = None
    report("创建 XLSX 测试文件", "fail", str(e))

# 创建 PPTX (需要 python-pptx)
pptx_path = tmp_dir / "test_presentation.pptx"
try:
    from pptx import Presentation as _Prs
    from pptx.util import Pt
    prs = _Prs()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "2026年战略规划"
    slide.placeholders[1].text = "打造行业领先的AI助手"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "三大核心目标"
    slide2.placeholders[1].text = "1. 用户增长 50%\n2. 技术升级迭代\n3. 成本优化 20%"
    prs.save(str(pptx_path))
    report("创建 PPTX 测试文件", "pass", str(pptx_path))
except Exception as e:
    pptx_path = None
    report("创建 PPTX 测试文件", "fail", str(e))

print()

# ─── TEST 1: read_file_snippet ────────────────────────────────────────────────
print("=" * 60)
print("TEST GROUP 1: read_file_snippet (pdf_reader / multi_format_reader)")
print("=" * 60)

from app.core.file.file_tools import FileToolsPlugin
tools = FileToolsPlugin()

# TXT
result = tools.read_file_snippet(str(txt_path), max_chars=1000)
if "销售额" in result and "错误" not in result:
    report("read_file_snippet(.txt)", "pass", f"读取 {len(result)} 字符，包含期望内容")
else:
    report("read_file_snippet(.txt)", "fail", result[:200])

# CSV
result = tools.read_file_snippet(str(csv_path), max_chars=1000)
if "销售额" in result or "月份" in result:
    report("read_file_snippet(.csv)", "pass", f"读取 {len(result)} 字符")
else:
    report("read_file_snippet(.csv)", "fail", result[:200])

# DOCX
if docx_path:
    result = tools.read_file_snippet(str(docx_path), max_chars=2000)
    if "合同" in result or "总则" in result:
        report("read_file_snippet(.docx)", "pass", f"读取 {len(result)} 字符")
    elif "错误" in result:
        report("read_file_snippet(.docx)", "fail", result[:200])
    else:
        report("read_file_snippet(.docx)", "warn", f"内容可能空: {result[:100]}")

# XLSX
if xlsx_path:
    result = tools.read_file_snippet(str(xlsx_path), max_chars=2000)
    if "月份" in result or "销售额" in result or "380000" in result:
        report("read_file_snippet(.xlsx)", "pass", f"读取 {len(result)} 字符")
    elif "错误" in result:
        report("read_file_snippet(.xlsx)", "fail", result[:200])
    else:
        report("read_file_snippet(.xlsx)", "warn", f"内容: {result[:100]}")

# PPTX (skill 说不支持，应该返回空/提示)
if pptx_path:
    result = tools.read_file_snippet(str(pptx_path), max_chars=2000)
    if "无法提取" in result or result.strip() == "" or "pptx" in result.lower():
        report("read_file_snippet(.pptx) → 不支持(符合预期)", "pass",
               f"返回: '{result[:80]}'")
    elif "战略" in result or "核心" in result:
        report("read_file_snippet(.pptx) → 意外支持!", "warn",
               f"skill 文档说不支持，但实际读出: {result[:100]}")
    else:
        report("read_file_snippet(.pptx)", "warn", f"返回: {result[:100]}")

# 不存在的文件
result = tools.read_file_snippet(str(tmp_dir / "nonexistent.pdf"), max_chars=1000)
if "不存在" in result or "错误" in result:
    report("read_file_snippet(不存在文件) → 正确报错", "pass", result[:80])
else:
    report("read_file_snippet(不存在文件) → 未报错", "fail", result[:80])

# max_chars 最小值钳位验证（实现将 max_chars 钳制到 [100, 10000]，所以 50 → 100）
result = tools.read_file_snippet(str(txt_path), max_chars=50)
if 100 <= len(result) <= 200:  # 50 应被钳制到 100，输出长度应在合理范围内
    report("read_file_snippet(max_chars=50) → 最小值钳位到100(符合预期)", "pass", f"实际长度: {len(result)}")
else:
    report("read_file_snippet(max_chars=50) → 钳位行为异常", "fail", f"实际长度: {len(result)}")

print()

# ─── TEST 2: find_file ─────────────────────────────────────────────────────────
print("=" * 60)
print("TEST GROUP 2: find_file (multi_format_reader / multi_doc_synthesis)")
print("=" * 60)

result = tools.find_file(query="销售报告")
if isinstance(result, str) and len(result) > 0:
    if "未找到" in result:
        report("find_file('销售报告') → 无索引(正常)", "warn",
               "FileRegistry 未扫描临时目录，这是预期行为")
    else:
        report("find_file('销售报告')", "pass", result[:150])
else:
    report("find_file('销售报告')", "fail", str(result)[:100])

# 尝试查找工作区里的真实文件
result2 = tools.find_file(query="config", limit=3)
if isinstance(result2, str):
    report("find_file('config', limit=3)", "pass", result2[:150])
else:
    report("find_file('config')", "fail", str(result2)[:100])

print()

# ─── TEST 3: summarize_file ──────────────────────────────────────────────────
print("=" * 60)
print("TEST GROUP 3: summarize_file (long_doc_parser / pdf_reader)")
print("=" * 60)

# summarize_file 依赖 LLM，这里只测试它能不能被调用、错误处理是否合理
result = tools.summarize_file(str(txt_path), focus="销售数据和风险")
if "LLM 摘要失败" in result or "摘要" in result.lower() or "报告" in result:
    # 两种预期结果：LLM成功返回摘要 OR LLM不可用但有回退
    if "LLM 摘要失败" in result:
        report("summarize_file (LLM不可用→回退预览)", "warn",
               f"回退内容长度: {len(result)}，skill 说这是兜底行为，符合预期")
    else:
        report("summarize_file (LLM摘要成功)", "pass", f"返回 {len(result)} 字符摘要")
elif "不存在" in result or "失败" in result:
    report("summarize_file (错误)", "fail", result[:200])
else:
    report("summarize_file (未知返回)", "warn", result[:200])

# 测试不存在文件
result2 = tools.summarize_file(str(tmp_dir / "fake.pdf"))
if "不存在" in result2 or "错误" in result2:
    report("summarize_file(不存在文件) → 正确报错", "pass")
else:
    report("summarize_file(不存在文件) → 未报错", "fail", result2[:100])

print()

# ─── TEST 4: PPTX 提取 (python-pptx CODER 路径) ───────────────────────────
print("=" * 60)
print("TEST GROUP 4: PPTX 文本提取 CODER 路径 (multi_format_reader skill)")
print("=" * 60)

if pptx_path:
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        extracted = []
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    extracted.append(f"[页{i}] {shape.text.strip()}")
        combined = "\n".join(extracted)
        if "战略规划" in combined and "核心目标" in combined:
            report("python-pptx 提取 PPTX 文字", "pass",
                   f"提取 {len(extracted)} 个文本块:\n       " + "\n       ".join(extracted))
        else:
            report("python-pptx 提取 PPTX 文字", "fail", f"未找到期望内容: {combined[:100]}")
    except Exception as e:
        report("python-pptx 提取 PPTX 文字", "fail", str(e))
else:
    report("python-pptx 提取 PPTX (跳过, pptx未创建)", "warn")

print()

# ─── TEST 5: DOCX 表格提取 (table_extractor skill) ─────────────────────────
print("=" * 60)
print("TEST GROUP 5: DOCX 表格提取 CODER 路径 (table_extractor skill)")
print("=" * 60)

if docx_path:
    try:
        from docx import Document
        doc = Document(str(docx_path))
        table_rows = []
        for t_idx, table in enumerate(doc.tables):
            for row in table.rows:
                row_text = ' | '.join(c.text.strip() for c in row.cells)
                table_rows.append(f"[表{t_idx+1}] {row_text}")
        if table_rows and "期次" in "\n".join(table_rows):
            report("python-docx 提取 DOCX 表格", "pass",
                   f"提取 {len(table_rows)} 行:\n       " + "\n       ".join(table_rows))
        else:
            report("python-docx 提取 DOCX 表格", "fail",
                   f"未找到期望内容: {table_rows[:3]}")
    except Exception as e:
        report("python-docx 提取 DOCX 表格", "fail", str(e))
else:
    report("python-docx DOCX 表格提取 (跳过)", "warn")

print()

# ─── TEST 6: Excel 生成 (excel_generator_pro skill 代码模板) ────────────────
print("=" * 60)
print("TEST GROUP 6: Excel 生成代码模板 (excel_generator_pro skill)")
print("=" * 60)

try:
    import openpyxl
    from openpyxl.styles import Font as _Font, PatternFill as _Fill, Alignment as _Align, Border as _Border, Side as _Side

    wb2 = openpyxl.Workbook()
    ws2 = wb2.active
    ws2.title = "季度报告"

    headers = ["部门", "Q1销售额", "Q2销售额", "增长率"]
    header_fill = _Fill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
    header_font = _Font(bold=True, color="FFFFFF", size=12)
    border = _Border(left=_Side(style='thin'), right=_Side(style='thin'),
                     top=_Side(style='thin'), bottom=_Side(style='thin'))
    for col_i, h in enumerate(headers, 1):
        c = ws2.cell(1, col_i, h)
        c.fill = header_fill; c.font = header_font
        c.alignment = _Align(horizontal='center'); c.border = border

    data = [("华东区", 450000, 520000, "15.6%"), ("华南区", 380000, 410000, "8.2%"),
            ("华北区", 420000, 480000, "14.3%")]
    alt_fill = _Fill(start_color="EBF3FB", end_color="EBF3FB", fill_type="solid")
    for row_i, row in enumerate(data, 2):
        for col_i, val in enumerate(row, 1):
            c = ws2.cell(row_i, col_i, val)
            c.border = border
            if row_i % 2 == 0: c.fill = alt_fill

    for col in ws2.columns:
        max_len = max((len(str(c.value or '')) for c in col), default=8)
        ws2.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)
    ws2.freeze_panes = 'A2'

    out_xlsx = tmp_dir / "generated_report.xlsx"
    wb2.save(str(out_xlsx))

    # 验证生成结果
    wb_verify = openpyxl.load_workbook(str(out_xlsx))
    ws_verify = wb_verify.active
    assert ws_verify.cell(1, 1).value == "部门", "标题行错误"
    assert ws_verify.cell(2, 1).value == "华东区", "数据行错误"
    assert ws_verify.freeze_panes == "A2", "冻结行错误"
    report("Excel 生成模板 (openpyxl)", "pass",
           f"文件大小: {out_xlsx.stat().st_size} bytes，标题/数据/冻结窗格均正确")
except Exception as e:
    report("Excel 生成模板 (openpyxl)", "fail", traceback.format_exc(limit=3))

print()

# ─── TEST 7: Word 生成 (docx_generator_pro skill 代码模板) ──────────────────
print("=" * 60)
print("TEST GROUP 7: Word 生成代码模板 (docx_generator_pro skill)")
print("=" * 60)

try:
    from docx import Document as _Doc2
    from docx.shared import Pt as _Pt, Cm as _Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc2 = _Doc2()
    section = doc2.sections[0]
    section.page_width = _Cm(21.0); section.page_height = _Cm(29.7)
    section.left_margin = section.right_margin = _Cm(2.54)

    title = doc2.add_heading('项目进展汇报（测试文档）', level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc2.add_heading('1. 项目背景', level=1)
    para = doc2.add_paragraph('本项目于2026年1月启动，旨在提升系统智能化水平。')
    para.paragraph_format.first_line_indent = _Pt(24)
    doc2.add_paragraph('• 核心功能开发完成', style='List Bullet')
    doc2.add_paragraph('• 性能优化进行中', style='List Bullet')

    table = doc2.add_table(rows=1, cols=3)
    table.style = 'Table Grid'
    hdr = table.rows[0].cells
    hdr[0].text = '里程碑'; hdr[1].text = '计划时间'; hdr[2].text = '状态'
    row = table.add_row().cells
    row[0].text = '功能开发'; row[1].text = '2026-01'; row[2].text = '✅已完成'

    out_docx = tmp_dir / "generated_report.docx"
    doc2.save(str(out_docx))

    # 验证
    verify_doc = _Doc2(str(out_docx))
    headings = [p.text for p in verify_doc.paragraphs if p.style.name.startswith('Heading')]
    assert len(verify_doc.tables) >= 1, "表格未生成"
    assert any("项目" in h for h in headings), "标题未生成"
    report("Word 生成模板 (python-docx)", "pass",
           f"文件 {out_docx.stat().st_size} bytes，{len(headings)} 个标题，{len(verify_doc.tables)} 个表格")
except Exception as e:
    report("Word 生成模板 (python-docx)", "fail", traceback.format_exc(limit=3))

print()

# ─── TEST 8: PDF 生成 (pdf_generator_pro skill 代码模板) ────────────────────
print("=" * 60)
print("TEST GROUP 8: PDF 生成代码模板 (pdf_generator_pro skill)")
print("=" * 60)

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont

    pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))

    out_pdf = tmp_dir / "generated_report.pdf"
    doc3 = SimpleDocTemplate(str(out_pdf), pagesize=A4,
                              leftMargin=2.5*cm, rightMargin=2.5*cm,
                              topMargin=2.5*cm, bottomMargin=2.5*cm)
    title_style = ParagraphStyle('T', fontName='STSong-Light', fontSize=18,
                                  spaceAfter=12, alignment=1)
    body_style = ParagraphStyle('B', fontName='STSong-Light', fontSize=12,
                                 leading=20, spaceAfter=8)
    story = [
        Paragraph('2026年度工作总结报告（测试）', title_style),
        Spacer(1, 0.5*cm),
        Paragraph('本文档为技能测试自动生成。包含中文内容验证。', body_style),
        Paragraph('测试结论：中文字体渲染正常。', body_style),
    ]
    doc3.build(story)
    assert out_pdf.exists() and out_pdf.stat().st_size > 1000, "PDF文件太小或未生成"
    report("PDF 生成模板 (reportlab + STSong-Light 中文字体)", "pass",
           f"文件大小: {out_pdf.stat().st_size} bytes")
except ImportError as e:
    report("PDF 生成 (reportlab 未安装)", "warn",
           f"skill 引用了 reportlab，但本机未安装: {e}\n       "
           "→ 需在 requirements.txt 中添加 reportlab，或改用 fpdf2")
except Exception as e:
    report("PDF 生成 (reportlab)", "fail", traceback.format_exc(limit=3))

# fpdf2 备选方案
try:
    from fpdf import FPDF
    pdf2 = FPDF(); pdf2.add_page(); pdf2.set_font('Helvetica', size=12)
    pdf2.cell(0, 10, 'PDF generation fallback test (English only)', new_x='LMARGIN', new_y='NEXT')
    out_pdf2 = tmp_dir / "fallback.pdf"
    pdf2.output(str(out_pdf2))
    report("PDF 生成备选方案 (fpdf2 英文)", "pass", f"文件大小: {out_pdf2.stat().st_size} bytes")
except ImportError:
    report("PDF 生成备选 (fpdf2 未安装)", "warn", "需安装: pip install fpdf2")
except Exception as e:
    report("PDF 生成备选 (fpdf2)", "fail", str(e))

print()

# ─── TEST 9: generate_ppt 工具可用性 ─────────────────────────────────────────
print("=" * 60)
print("TEST GROUP 9: generate_ppt 工具注册状态 (ppt_generator_pro skill)")
print("=" * 60)

try:
    from app.core.agent.tool_registry import ToolRegistry
    reg = ToolRegistry()
    tools_list = reg.list_tools()
    tool_names = [t.get("name", "") if isinstance(t, dict) else getattr(t, "name", "") for t in tools_list]
    if "generate_ppt" in tool_names:
        report("generate_ppt 已注册到 ToolRegistry", "pass")
    else:
        report("generate_ppt 未在 ToolRegistry (仅 LanggraphWorkflow 内部用)", "warn",
               f"已注册工具: {[n for n in tool_names if n][:10]}\n"
               "       → skill 中的 generate_ppt() 写法对 agent 不可直接调用，需修正为工作流触发")
except ImportError as e:
    report("ToolRegistry 导入失败", "warn", str(e))
except Exception as e:
    report("generate_ppt 工具检查异常", "warn", traceback.format_exc(limit=2))

print()

# ─── TEST 10: SkillAutoMatcher 新规则 ───────────────────────────────────────
print("=" * 60)
print("TEST GROUP 10: SkillAutoMatcher 新规则 (关键词触发)")
print("=" * 60)

try:
    from app.core.skills.skill_auto_matcher import SkillAutoMatcher

    test_cases = [
        ("帮我读取这个pdf文件",     ["pdf_reader"]),
        ("把这个Excel分析一下",    ["spreadsheet_analyst"]),
        ("生成一份ppt",            ["ppt_generator_pro"]),
        ("生成word文档",           ["docx_generator_pro"]),
        ("做一个excel表格",        ["excel_generator_pro"]),
        ("把这几份文件对比一下",    ["multi_doc_synthesis"]),
        ("提取文档里的表格",       ["table_extractor"]),
        ("生成pdf报告",            ["pdf_generator_pro"]),
        ("读取文件内容",           ["multi_format_reader"]),
        ("这个文档太长了分段分析",  ["long_doc_parser"]),
        ("帮我把这个word翻译成英文", ["docx_translator"]),
        ("翻译这个docx文档",       ["docx_translator"]),
        ("文档翻译成日文",         ["docx_translator"]),
    ]

    hit = 0
    for user_input, expected_skills in test_cases:
        matched = SkillAutoMatcher._match_with_patterns(user_input, [
            {"id": s, "name": s} for s in [
                "pdf_reader", "multi_format_reader", "long_doc_parser",
                "spreadsheet_analyst", "multi_doc_synthesis", "ppt_generator_pro",
                "excel_generator_pro", "docx_generator_pro", "pdf_generator_pro",
                "table_extractor", "docx_translator",
            ]
        ])
        for expected in expected_skills:
            if expected in matched:
                report(f"  AutoMatcher: '{user_input}' → {expected}", "pass")
                hit += 1
            else:
                report(f"  AutoMatcher: '{user_input}' → {expected}", "fail",
                       f"实际匹配: {matched}")

    print(f"\n  命中率: {hit}/{len(test_cases)}")
except Exception as e:
    report("SkillAutoMatcher 测试失败", "fail", traceback.format_exc(limit=3))

print()

# ─── 汇总 ────────────────────────────────────────────────────────────────────
print("=" * 60)
print("测试汇总")
print("=" * 60)
pass_n = sum(1 for _, s, _ in results if s == "pass")
fail_n = sum(1 for _, s, _ in results if s == "fail")
warn_n = sum(1 for _, s, _ in results if s == "warn")
print(f"  {PASS}: {pass_n}   {FAIL}: {fail_n}   {WARN}: {warn_n}")
if fail_n:
    print("\n失败项：")
    for label, status, detail in results:
        if status == "fail":
            print(f"  • {label}")
            if detail:
                print(f"    {detail[:120]}")
if warn_n:
    print("\n注意项：")
    for label, status, detail in results:
        if status == "warn":
            print(f"  • {label}")
            if detail:
                print(f"    {detail[:120]}")

# 清理
shutil.rmtree(tmp_dir, ignore_errors=True)
print(f"\n🧹 临时目录已清理: {tmp_dir}")
print("\n完成。")
sys.exit(0 if fail_n == 0 else 1)
