#!/usr/bin/env python3
"""Deep XML inspection of PPTX text elements."""
import sys
sys.path.insert(0, '.')

from pptx import Presentation
import lxml.etree as etree

_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

prs = Presentation('workspace/AI Agent.pptx')
all_slides = list(prs.slides)

# Look at slide 4 which should have content
slide = all_slides[4]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    if 'xlabel' in shape.name.lower() or '内容' in shape.name:
        print(f"Shape: {shape.name}")
        for pi, para in enumerate(shape.text_frame.paragraphs[:2]):
            # Print full XML of paragraph
            xml_bytes = etree.tostring(para._p, pretty_print=True, encoding='unicode')
            print(f"  Para {pi} XML (first 800 chars):")
            print(xml_bytes[:800])
            print()
        break

# Try slide 3 content placeholder
print("=== SLIDE 3 CONTENT PLACEHOLDER ===")
slide = all_slides[3]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    print(f"\nShape: {shape.name!r}, text={shape.text_frame.text!r}")
    # Find the paragraph with actual content
    for pi, para in enumerate(shape.text_frame.paragraphs):
        if para.text.strip():
            xml_bytes = etree.tostring(para._p, pretty_print=True, encoding='unicode')
            print(f"  Para {pi} with text {para.text!r}:")
            print(xml_bytes[:600])
            break
    # If none have text, show the first one
    paras = shape.text_frame.paragraphs
    if paras:
        xml_bytes = etree.tostring(paras[0]._p, pretty_print=True, encoding='unicode')
        print(f"  Para 0 XML:")
        print(xml_bytes[:600])

# Check the <a:t> elements directly
print()
print("=== RAW T-ELEMENT INSPECTION ===")
slide = all_slides[4]
for shape in slide.shapes:
    if not shape.has_text_frame:
        continue
    for pi, para in enumerate(shape.text_frame.paragraphs[:2]):
        for child in para._p:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag == 'r':
                t_el = child.find(f"{{{_NS}}}t")
                if t_el is not None:
                    print(f"  <a:t> text={t_el.text!r}, len(text)={len(t_el.text) if t_el.text else 'None'}")
                    # Check if text is in another namespace
                    print(f"  <a:t> tag={t_el.tag!r}")
                    print(f"  <a:t> attrib={dict(t_el.attrib)!r}")
                    # Look for any child text
                    for child2 in t_el:
                        print(f"  <a:t> child: {child2.tag!r} text={child2.text!r}")
        break
    break
