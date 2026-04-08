#!/usr/bin/env python3
"""Deep diagnostic of PPTX text extraction."""
import sys
sys.path.insert(0, '.')

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("python-pptx not installed")
    sys.exit(1)

prs = Presentation('workspace/AI Agent.pptx')
print(f"Loaded presentation: {len(prs.slides)} slides")
print(f"Slide size: {prs.slide_width} x {prs.slide_height}")

_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

def extract_text_direct(shape):
    """Direct text extraction using python-pptx high-level API."""
    if not shape.has_text_frame:
        return None
    texts = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            texts.append(repr(run.text))
    return texts

def extract_text_xml(shape):
    """XML-level text extraction."""
    if not shape.has_text_frame:
        return None
    texts = []
    for para in shape.text_frame.paragraphs:
        p_el = para._p
        # All <a:r><a:t> elements
        for child in p_el:
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('r', 'fld'):
                t_el = child.find(f"{{{_NS}}}t")
                if t_el is not None:
                    texts.append(repr(t_el.text))
    return texts

# Check first 3 slides
all_slides = list(prs.slides)
for slide_idx, slide in enumerate(all_slides[:5]):
    print(f"\n=== SLIDE {slide_idx} ===")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        hi_level = extract_text_direct(shape)
        xml_level = extract_text_xml(shape)
        para_text = shape.text_frame.text  # Full text
        
        print(f"  Shape: {shape.name!r}")
        print(f"    shape.text_frame.text = {para_text!r}")
        print(f"    high-level .runs = {hi_level}")
        print(f"    xml-level <a:r><a:t> = {xml_level}")
        
        # Check if text is in paragraphs
        paras = shape.text_frame.paragraphs
        for pi, para in enumerate(paras[:3]):
            ptext = para.text
            if ptext:
                print(f"    para[{pi}].text = {ptext!r}")
            # Check raw XML
            import lxml.etree as etree
            xml_str = etree.tostring(para._p, encoding='unicode')
            if len(xml_str) > 50:  # Only show non-trivial paragraphs
                print(f"    para[{pi}] XML len={len(xml_str)}: {xml_str[:200]!r}")
