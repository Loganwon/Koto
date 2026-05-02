#!/usr/bin/env python3
"""Check if PPTX text is in slide layout/master, notes, or elsewhere."""
import sys
sys.path.insert(0, '.')

from pptx import Presentation
from pptx.util import Pt
import lxml.etree as etree
import zipfile

prs = Presentation('workspace/AI Agent.pptx')
all_slides = list(prs.slides)

print("=== CHECKING SLIDE NOTES ===")
slide = all_slides[4]
if slide.has_notes_slide:
    notes = slide.notes_slide
    print("Notes:", notes.notes_text_frame.text[:200])
else:
    print("No notes")

print()
print("=== CHECKING SLIDE LAYOUT ===")
layout = slide.slide_layout
for ph in layout.placeholders:
    if ph.has_text_frame and ph.text_frame.text.strip():
        print(f"Layout PH {ph.placeholder_format.idx}: {ph.text_frame.text[:100]!r}")

print()
print("=== CHECKING SLIDE MASTER ===")
master = slide.slide_layout.slide_master
for ph in master.placeholders:
    if ph.has_text_frame and ph.text_frame.text.strip():
        print(f"Master PH {ph.placeholder_format.idx}: {ph.text_frame.text[:100]!r}")

print()
print("=== READING RAW PPTX ZIP ===")
with zipfile.ZipFile('workspace/AI Agent.pptx', 'r') as z:
    names = z.namelist()
    slide_xmls = [n for n in names if n.startswith('ppt/slides/slide') and not 'rels' in n]
    slide_xmls.sort()
    print("Slide files:", slide_xmls[:5])
    
    # Read slide 5 raw XML
    if 'ppt/slides/slide5.xml' in names:
        raw = z.read('ppt/slides/slide5.xml').decode('utf-8')
        # Look for any text between <a:t> tags
        import re
        t_matches = re.findall(r'<a:t[^>]*>([^<]*)</a:t>', raw)
        print(f"\nSlide 5 <a:t> text content matches: {t_matches[:20]!r}")
        print(f"Total <a:t> matches: {len(t_matches)}")
        
        # Also check for <p:sp> shapes with text
        sp_texts = re.findall(r'<a:t[^/]', raw)
        print(f"Non-empty <a:t> tags: {len(sp_texts)}")
        
        # Find all text between all opening/closing tags in the file
        all_text_chunks = re.findall(r'>([^<]+)<', raw)
        real_text = [t for t in all_text_chunks if t.strip() and not t.strip().startswith('?')]
        print(f"\nAll non-whitespace text between ANY tags in slide5.xml: {real_text[:30]!r}")

    # Check slide 1  
    if 'ppt/slides/slide1.xml' in names:
        raw = z.read('ppt/slides/slide1.xml').decode('utf-8')
        t_matches = re.findall(r'<a:t[^>]*>([^<]*)</a:t>', raw)
        print(f"\nSlide 1 <a:t> text matches: {t_matches[:20]!r}")
        all_text_chunks = re.findall(r'>([^<]+)<', raw)
        real_text = [t for t in all_text_chunks if t.strip()]
        print(f"All non-whitespace text in slide1.xml: {real_text[:30]!r}")
        
print()
print("=== CHECK ALL SLIDES FOR ANY TEXT ===")
with zipfile.ZipFile('workspace/AI Agent.pptx', 'r') as z:
    for slide_xml in slide_xmls[:17]:
        raw = z.read(slide_xml).decode('utf-8')
        t_matches = re.findall(r'<a:t[^>]*>([^<]*)</a:t>', raw)
        non_empty = [t for t in t_matches if t.strip()]
        print(f"{slide_xml}: <a:t> total={len(t_matches)}, non-empty={len(non_empty)}, first few={non_empty[:3]!r}")
