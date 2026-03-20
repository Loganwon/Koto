#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
文档对比与总结器 - 支持多格式（PDF/DOCX/XLSX/PPTX/TXT/MD）、
多文档横向对比矩阵、AI 语义分析 prompt 构建。
"""

import os
import difflib
from typing import Dict, List, Any, Optional
from datetime import datetime
import logging

logger = logging.getLogger(__name__)

# ── 参数 ──────────────────────────────────────────────────────────────────────
_MAX_CHARS_PER_DOC = 6000   # 每份文档送 LLM 的最大字符数
_MAX_DOCS = 5               # 同时对比最多多少份文档


class DocumentComparator:
    """多文档对比与总结器（文本 diff + AI 语义分析）"""

    SUPPORTED_FORMATS = {
        '.txt', '.md', '.markdown',
        '.docx', '.doc',
        '.pdf',
        '.xlsx', '.xls',
        '.pptx', '.ppt',
    }

    def __init__(self):
        # 向后兼容
        self.supported_formats = list(self.SUPPORTED_FORMATS)

    # ── 公共接口 ──────────────────────────────────────────────────────────────

    def compare_documents(self, file_a: str, file_b: str,
                          output_format: str = "markdown") -> Dict[str, Any]:
        """对比两个文档（向后兼容原始接口）。"""
        if not os.path.exists(file_a) or not os.path.exists(file_b):
            return {"success": False, "error": "文件不存在"}

        text_a = self._read_file(file_a)
        text_b = self._read_file(file_b)

        if text_a is None or text_b is None:
            return {"success": False, "error": "无法读取文件内容"}

        lines_a = text_a.splitlines()
        lines_b = text_b.splitlines()

        changes = self._analyze_changes(lines_a, lines_b)
        summary = self._generate_summary(changes)

        if output_format == "markdown":
            diff_output = self._format_diff_markdown(lines_a, lines_b)
        elif output_format == "html":
            diff_output = self._format_diff_html(lines_a, lines_b)
        elif output_format == "inline_json":
            diff_output = self._format_diff_inline_json(lines_a, lines_b)
        else:
            diff_output = '\n'.join(
                difflib.unified_diff(lines_a, lines_b, lineterm='')
            )

        return {
            "success": True,
            "file_a": file_a,
            "file_b": file_b,
            "changes": changes,
            "summary": summary,
            "diff": diff_output,
            "timestamp": datetime.now().isoformat(),
        }

    def compare_multiple(self, file_paths: List[str],
                         output_format: str = "inline_json") -> Dict[str, Any]:
        """
        N 个文档横向对比：返回所有两两 diff 矩阵 + 全局统计。

        Args:
            file_paths: 文档路径列表（2-5 个）
            output_format: "inline_json" | "markdown" | "html"
        """
        if len(file_paths) < 2:
            return {"success": False, "error": "至少需要两个文件"}
        if len(file_paths) > _MAX_DOCS:
            return {"success": False, "error": f"最多支持 {_MAX_DOCS} 个文件同时对比"}

        # 读取所有文档
        docs: List[Dict] = []
        for path in file_paths:
            if not os.path.exists(path):
                return {"success": False, "error": f"文件不存在: {path}"}
            text = self._read_file(path)
            if text is None:
                return {"success": False, "error": f"无法读取文件: {os.path.basename(path)}"}
            docs.append({
                "name": os.path.basename(path),
                "path": path,
                "content": text,
                "char_count": len(text),
                "content_preview": text[:300].replace('\n', ' '),
            })

        n = len(docs)
        matrix: List[Dict] = []
        similarities: List[float] = []

        for i in range(n):
            for j in range(i + 1, n):
                lines_i = docs[i]["content"].splitlines()
                lines_j = docs[j]["content"].splitlines()
                changes = self._analyze_changes(lines_i, lines_j)
                summary = self._generate_summary(changes)

                if output_format == "markdown":
                    diff_out = self._format_diff_markdown(lines_i, lines_j)
                elif output_format == "html":
                    diff_out = self._format_diff_html(lines_i, lines_j)
                else:
                    diff_out = self._format_diff_inline_json(lines_i, lines_j)

                matrix.append({
                    "doc_a": {"index": i, "name": docs[i]["name"]},
                    "doc_b": {"index": j, "name": docs[j]["name"]},
                    "changes": changes,
                    "summary": summary,
                    "diff": diff_out,
                })
                similarities.append(changes["similarity"])

        if matrix:
            max_sim = max(matrix, key=lambda p: p["changes"]["similarity"])
            min_sim = min(matrix, key=lambda p: p["changes"]["similarity"])
            avg_sim = round(sum(similarities) / len(similarities), 2)
        else:
            max_sim = min_sim = None
            avg_sim = 0.0

        return {
            "success": True,
            "files": [
                {
                    "name": d["name"],
                    "path": d["path"],
                    "char_count": d["char_count"],
                    "content_preview": d["content_preview"],
                }
                for d in docs
            ],
            "matrix": matrix,
            "global_stats": {
                "file_count": n,
                "pair_count": len(matrix),
                "avg_similarity": avg_sim,
                "most_similar": {
                    "pair": f"{max_sim['doc_a']['name']} vs {max_sim['doc_b']['name']}",
                    "similarity": max_sim["changes"]["similarity"],
                } if max_sim else None,
                "most_different": {
                    "pair": f"{min_sim['doc_a']['name']} vs {min_sim['doc_b']['name']}",
                    "similarity": min_sim["changes"]["similarity"],
                } if min_sim else None,
            },
        }

    def build_ai_prompt(self, file_paths: List[str],
                        focus: str = "general") -> Optional[str]:
        """
        构建用于 AI 语义对比的 prompt（由 Flask 路由注入 LLM 调用）。

        Args:
            file_paths: 文档路径列表
            focus: "general" | "argument" | "data" | "structure"

        Returns:
            prompt 字符串，或 None（文件读取失败）
        """
        docs: List[Dict] = []
        for path in file_paths:
            text = self._read_file(path)
            if text is None:
                logger.warning(f"[DocComparator] 跳过无法读取的文件: {path}")
                continue
            truncated = text[:_MAX_CHARS_PER_DOC]
            if len(text) > _MAX_CHARS_PER_DOC:
                truncated += "\n\n[...内容已截断，仅显示前段...]"
            docs.append({"name": os.path.basename(path), "content": truncated})

        if len(docs) < 2:
            return None

        doc_sections = "\n\n".join(
            f"【文档{chr(65 + i)}】{d['name']}\n{'─' * 50}\n{d['content']}"
            for i, d in enumerate(docs)
        )

        focus_instruction = {
            "general":   "综合对比这些文档的主题、观点、结构和关键信息的异同",
            "argument":  "重点对比各文档的核心论点、论据和结论是否一致或矛盾",
            "data":      "重点对比各文档中的数据、数字、日期、统计信息的差异",
            "structure": "重点对比各文档的章节结构、逻辑框架和信息组织方式",
        }.get(focus, "综合对比这些文档的主题、观点、结构和关键信息的异同")

        return f"""你是专业的文档分析师。请{focus_instruction}。

{doc_sections}

请按以下结构输出分析报告（使用 Markdown 格式）：

## 📋 文档概览
（逐一简介每份文档的主题和核心内容，1-2句）

## 🔍 主要差异
（列出最显著的 3-5 个差异点，每点说明哪份文档如何、哪份文档如何）

## ✅ 共同点
（所有文档一致或高度相似的内容）

## ⚡ 关键冲突
（如果存在相互矛盾的信息，明确指出；若无则写"未发现明显冲突"）

## 💡 综合结论
（基于对比给出综合性判断和建议）
"""

    def compare_versions(self, file_paths: List[str]) -> Dict[str, Any]:
        """对比多个版本的文档（向后兼容）。"""
        if len(file_paths) < 2:
            return {"success": False, "error": "至少需要两个文件"}

        versions = []
        for i in range(len(file_paths) - 1):
            result = self.compare_documents(file_paths[i], file_paths[i + 1])
            if result["success"]:
                versions.append({
                    "from": os.path.basename(file_paths[i]),
                    "to": os.path.basename(file_paths[i + 1]),
                    "summary": result["summary"],
                })

        return {
            "success": True,
            "total_versions": len(file_paths),
            "comparisons": len(versions),
            "versions": versions,
        }

    def generate_change_log(self, comparisons: List[Dict[str, Any]],
                            output_file: str) -> str:
        """生成变更日志 Markdown 文件。"""
        lines = ["# 文档变更日志",
                 f"\n生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"]

        for i, comp in enumerate(comparisons, 1):
            lines.append(f"## 版本 {i}: {comp['file_a']} → {comp['file_b']}")
            lines.append(f"\n{comp['summary']}\n")

            for section_key, label in [("additions", "新增内容"),
                                        ("deletions", "删除内容")]:
                section = comp["changes"].get(section_key, {})
                cnt = section.get("count", 0)
                if cnt:
                    lines.append(f"### {label}")
                    prefix = "+" if section_key == "additions" else "-"
                    for line in section.get("lines", [])[:10]:
                        lines.append(f"{prefix} {line}")
                    if cnt > 10:
                        lines.append(f"... 还有 {cnt - 10} 行")
                    lines.append("")

        with open(output_file, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        return output_file

    # ── 私有辅助 ─────────────────────────────────────────────────────────────

    def _analyze_changes(self, lines_a: List[str],
                         lines_b: List[str]) -> Dict[str, Any]:
        matcher = difflib.SequenceMatcher(None, lines_a, lines_b)
        additions, deletions, modifications = [], [], []

        for tag, i1, i2, j1, j2 in matcher.get_opcodes():
            if tag == 'insert':
                additions.extend(lines_b[j1:j2])
            elif tag == 'delete':
                deletions.extend(lines_a[i1:i2])
            elif tag == 'replace':
                modifications.append({"old": lines_a[i1:i2], "new": lines_b[j1:j2]})

        text_a = '\n'.join(lines_a)
        text_b = '\n'.join(lines_b)

        return {
            "additions": {"count": len(additions), "lines": additions},
            "deletions": {"count": len(deletions), "lines": deletions},
            "modifications": {"count": len(modifications), "details": modifications},
            "similarity": round(matcher.ratio() * 100, 2),
            "char_diff": len(text_b) - len(text_a),
            "line_diff": len(lines_b) - len(lines_a),
        }

    def _generate_summary(self, changes: Dict[str, Any]) -> str:
        sim = changes["similarity"]
        if sim >= 95:
            label = "✅ 文档变化很小"
        elif sim >= 80:
            label = "📝 文档有适度修改"
        elif sim >= 50:
            label = "⚠️ 文档有较大变化"
        else:
            label = "🔄 文档被大幅改写"

        lines = [label, f"- 相似度: {sim}%"]
        if changes["additions"]["count"]:
            lines.append(f"- 新增: {changes['additions']['count']} 行")
        if changes["deletions"]["count"]:
            lines.append(f"- 删除: {changes['deletions']['count']} 行")
        if changes["modifications"]["count"]:
            lines.append(f"- 修改: {changes['modifications']['count']} 处")
        cd = changes["char_diff"]
        if cd > 0:
            lines.append(f"- 内容增加: +{cd} 字符")
        elif cd < 0:
            lines.append(f"- 内容减少: {cd} 字符")
        return '\n'.join(lines)

    def _format_diff_inline_json(self, lines_a: List[str],
                                 lines_b: List[str]) -> List[Dict]:
        """
        返回结构化 diff 列表，供前端逐块高亮渲染。
        每个元素: {"type": "equal"|"insert"|"delete"|"replace",
                   "lines_a": [...], "lines_b": [...]}
        """
        result = []
        matcher = difflib.SequenceMatcher(None, lines_a, lines_b, autojunk=False)
        for tag, i1, i2, j1, j2 in matcher.get_opcodes():
            result.append({
                "type": tag,
                "lines_a": lines_a[i1:i2],
                "lines_b": lines_b[j1:j2],
            })
        return result

    def _format_diff_markdown(self, lines_a: List[str],
                               lines_b: List[str]) -> str:
        matcher = difflib.SequenceMatcher(None, lines_a, lines_b)
        output = ["# 文档对比\n"]

        for tag, i1, i2, j1, j2 in matcher.get_opcodes():
            if tag == 'equal':
                span = i2 - i1
                if span > 4:
                    for line in lines_a[i1:i1 + 2]:
                        output.append(f"  {line}")
                    output.append(f"  ... ({span - 4} 行相同内容省略) ...")
                    for line in lines_a[i2 - 2:i2]:
                        output.append(f"  {line}")
                else:
                    for line in lines_a[i1:i2]:
                        output.append(f"  {line}")
            elif tag == 'delete':
                output.append("\n**删除:**")
                for line in lines_a[i1:i2]:
                    output.append(f"- ~~{line}~~")
            elif tag == 'insert':
                output.append("\n**新增:**")
                for line in lines_b[j1:j2]:
                    output.append(f"+ **{line}**")
            elif tag == 'replace':
                output.append("\n**修改:**")
                output.append("原文:")
                for line in lines_a[i1:i2]:
                    output.append(f"- ~~{line}~~")
                output.append("改为:")
                for line in lines_b[j1:j2]:
                    output.append(f"+ **{line}**")

        return '\n'.join(output)

    def _format_diff_html(self, lines_a: List[str], lines_b: List[str]) -> str:
        return difflib.HtmlDiff().make_file(
            lines_a, lines_b, fromdesc="原始版本", todesc="修改版本"
        )

    def _read_file(self, file_path: str) -> Optional[str]:
        """
        读取文件内容，返回纯文本字符串。
        支持: TXT, MD, DOCX, PDF, XLSX, PPTX。
        失败返回 None。
        """
        ext = os.path.splitext(file_path)[1].lower()
        try:
            # ── 纯文本 ──
            if ext in {'.txt', '.md', '.markdown'}:
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    return f.read()

            # ── Word ──
            elif ext in {'.docx', '.doc'}:
                from web.file_parser import FileParser
                result = FileParser.parse_file(file_path)
                return result.get("content") if result.get("success") else None

            # ── PDF ──
            elif ext == '.pdf':
                from web.file_parser import FileParser
                result = FileParser.parse_file(file_path)
                return result.get("content") if result.get("success") else None

            # ── Excel ──
            elif ext in {'.xlsx', '.xls'}:
                return self._read_xlsx(file_path)

            # ── PowerPoint ──
            elif ext in {'.pptx', '.ppt'}:
                return self._read_pptx(file_path)

            else:
                logger.warning(f"[DocComparator] 不支持的格式: {ext}")
                return None

        except Exception as e:
            logger.warning(f"[DocComparator] 读取文件失败 {file_path}: {e}")
            return None

    def _read_xlsx(self, file_path: str) -> Optional[str]:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    row_text = "\t".join(
                        str(cell) if cell is not None else "" for cell in row
                    )
                    if row_text.strip():
                        lines.append(row_text)
            return '\n'.join(lines)
        except ImportError:
            logger.warning("[DocComparator] openpyxl 未安装，无法读取 Excel 文件")
            return None

    def _read_pptx(self, file_path: str) -> Optional[str]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            lines = []
            for i, slide in enumerate(prs.slides, 1):
                lines.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                lines.append(text)
            return '\n'.join(lines)
        except ImportError:
            logger.warning("[DocComparator] python-pptx 未安装，无法读取 PPT 文件")
            return None
