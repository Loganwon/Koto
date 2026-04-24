# Copyright (C) 2024-2026 Koto AI. All rights reserved.
# SPDX-License-Identifier: LicenseRef-Koto-Proprietary
"""
PPTX File Editor API — /api/pptx

Supports uploading an existing .pptx file, editing slide text in-place
(preserving all original formatting, images, and theme), and downloading
the result.  Slides data is persisted as JSON in workspace/pptx-files/.

Routes:
  POST   /api/pptx/upload          — Upload .pptx, parse, store
  GET    /api/pptx/list            — List all uploaded PPTX files
  GET    /api/pptx/<id>            — Get parsed slides data
  PUT    /api/pptx/<id>            — Save edited slides data
  GET    /api/pptx/<id>/download   — Export & download modified .pptx
  DELETE /api/pptx/<id>            — Delete a PPTX session
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import time
import uuid

from flask import Blueprint, Response, jsonify, request, send_file

_logger = logging.getLogger("koto.routes.pptx_editor")

pptx_editor_bp = Blueprint("pptx_editor", __name__, url_prefix="/api/pptx")

# ── Storage ──────────────────────────────────────────────────────────────────

_STORE_DIR: str | None = None


def _get_store_dir() -> str:
    global _STORE_DIR
    if _STORE_DIR is None:
        from web.shared import WORKSPACE_DIR

        _STORE_DIR = os.path.join(WORKSPACE_DIR, "pptx-files")
    os.makedirs(_STORE_DIR, exist_ok=True)
    return _STORE_DIR


def _meta_path(file_id: str) -> str:
    safe = re.sub(r"[^a-zA-Z0-9_-]", "", file_id)
    if not safe:
        raise ValueError("Invalid PPTX file ID")
    return os.path.join(_get_store_dir(), f"{safe}.json")


def _orig_path(file_id: str) -> str:
    safe = re.sub(r"[^a-zA-Z0-9_-]", "", file_id)
    return os.path.join(_get_store_dir(), f"{safe}_original.pptx")


def _new_id() -> str:
    return uuid.uuid4().hex[:12]


def _read_meta(file_id: str) -> dict | None:
    path = _meta_path(file_id)
    if not os.path.exists(path):
        return None
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def _write_meta(meta: dict) -> None:
    with open(_meta_path(meta["id"]), "w", encoding="utf-8") as f:
        json.dump(meta, f, ensure_ascii=False, indent=2)


# ── PPTX parsing ─────────────────────────────────────────────────────────────


def _parse_slides(raw_bytes: bytes) -> dict:
    """
    Parse a .pptx binary blob into structured slide data using the geometry parser.
    """
    import io
    from app.core.file.file_parser import parse_pptx_geometry

    parsed = parse_pptx_geometry(io.BytesIO(raw_bytes))

    # Backward-compatibility for older editor/tests expecting `index` and
    # shape `type` while the geometry parser now emits `slide_index`/`_type`.
    slides = parsed.get("slides", []) if isinstance(parsed, dict) else []
    for i, slide in enumerate(slides):
        if not isinstance(slide, dict):
            continue
        slide.setdefault("index", slide.get("slide_index", i))
        shapes = slide.get("shapes", [])
        for shape in shapes:
            if isinstance(shape, dict) and "type" not in shape and "_type" in shape:
                shape["type"] = shape["_type"]

    return parsed


# ── PPTX export (in-place text update) ───────────────────────────────────────

def _apply_edits(orig_bytes: bytes, slides_edits: list[dict]) -> bytes:
    """
    Apply text and formatting edits to an existing .pptx, preserving everything else.

    `slides_edits` is the `slides` array produced by _parse_slides().
    Runs are matched by (slide_index, shape_id, paragraph_index, run_index).
    Writes back: text, bold, italic, underline, color, size.
    Paragraph alignment is also applied where present.
    Runs not present in edits are left untouched.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    _ALIGN_MAP = {
        "LEFT": PP_ALIGN.LEFT,
        "CENTER": PP_ALIGN.CENTER,
        "RIGHT": PP_ALIGN.RIGHT,
        "JUSTIFY": PP_ALIGN.JUSTIFY,
        "DISTRIBUTE": PP_ALIGN.DISTRIBUTE,
        "THAI_DISTRIBUTE": PP_ALIGN.THAI_DISTRIBUTE,
    }

    prs = Presentation(io.BytesIO(orig_bytes))
    slide_map = {edit.get("index", edit.get("slide_index", i)): edit for i, edit in enumerate(slides_edits)}

    for slide_idx, slide in enumerate(prs.slides):
        edit_slide = slide_map.get(slide_idx)
        if not edit_slide:
            continue

        # ── Apply slide background ──────────────────────────────────────────
        # Write background color changes back to PPTX XML.
        # Image backgrounds are stored as data URIs in the JSON; for the PPTX
        # export we convert them back to an embedded image relationship.
        try:
            bg_color = edit_slide.get("background")
            bg_image = edit_slide.get("backgroundImage")
            if bg_image and bg_image.startswith("data:"):
                # Decode data URI and embed as blipFill in slide background
                import base64 as _b64
                from pptx.oxml.ns import qn as _qn
                from lxml import etree as _et
                header, data = bg_image.split(",", 1)
                img_bytes = _b64.b64decode(data)
                mime = header.split(":")[1].split(";")[0]
                ext_map = {"image/jpeg": "jpg", "image/png": "png", "image/gif": "gif",
                           "image/webp": "webp", "image/bmp": "bmp"}
                img_ext = ext_map.get(mime, "png")
                # Add image as a media part inside the PPTX package
                # Use python-pptx's internal Part API (wrapped in try/except for safety)
                from pptx.opc.part import Part as _Part
                from pptx.opc.constants import RELATIONSHIP_TYPE as _RT
                part_name = f"/ppt/media/koto_bg_{slide_idx}.{img_ext}"
                img_part = _Part(part_name, mime, img_bytes)
                rId = slide.part.relate_to(img_part, _RT.IMAGE)
                # Build <p:bg><p:bgPr><a:blipFill>...</a:blipFill>...</p:bgPr></p:bg>
                bg_xml = (
                    f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
                    f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
                    f' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
                    f'<p:bgPr><a:blipFill dpi="0" rotWithShape="1">'
                    f'<a:blip r:embed="{rId}"/>'
                    f'<a:stretch><a:fillRect/></a:stretch>'
                    f'</a:blipFill>'
                    f'<a:effectLst/></p:bgPr></p:bg>'
                )
                new_bg = _et.fromstring(bg_xml)
                sp_tree = slide.shapes._spTree
                cSld = sp_tree.getparent()
                old_bg = cSld.find(_qn('p:bg'))
                if old_bg is not None:
                    cSld.remove(old_bg)
                cSld.insert(0, new_bg)
            elif bg_color and bg_color.startswith("#") and len(bg_color) == 7:
                # Solid color background
                from pptx.oxml.ns import qn as _qn
                from lxml import etree as _et
                hex_val = bg_color.lstrip("#")
                bg_xml = (
                    f'<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
                    f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    f'<p:bgPr><a:solidFill>'
                    f'<a:srgbClr val="{hex_val}"/>'
                    f'</a:solidFill>'
                    f'<a:effectLst/></p:bgPr></p:bg>'
                )
                new_bg = _et.fromstring(bg_xml)
                sp_tree = slide.shapes._spTree
                cSld = sp_tree.getparent()
                old_bg = cSld.find(_qn('p:bg'))
                if old_bg is not None:
                    cSld.remove(old_bg)
                cSld.insert(0, new_bg)
        except Exception as _bg_exc:
            _logger.debug("Could not write background for slide %d: %s", slide_idx, _bg_exc)

        shape_map = {s["id"]: s for s in edit_slide.get("shapes", [])}

        # Keep track of existing ids from the file
        file_shape_ids = [s.shape_id for s in slide.shapes]
        
        # Traverse backwards to safely delete
        for shape_id in reversed(file_shape_ids):
            shape = next((s for s in slide.shapes if s.shape_id == shape_id), None)
            if not shape: continue
            
            edit_shape = shape_map.get(shape_id)
            if not edit_shape:
                # Shape was deleted from the frontend, remove it from slide
                try:
                    shape.element.getparent().remove(shape.element)
                except Exception:
                    pass
                continue
                
            # Update shape geometry if present
            try:
                if "left" in edit_shape: shape.left = int(edit_shape["left"])
                if "top" in edit_shape: shape.top = int(edit_shape["top"])
                if "width" in edit_shape: shape.width = int(edit_shape["width"])
                if "height" in edit_shape: shape.height = int(edit_shape["height"])
            except Exception:
                pass

            if not edit_shape.get("has_text"):
                continue
            if not getattr(shape, "has_text_frame", False):
                continue

            edit_paras = edit_shape.get("paragraphs", [])
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                if p_idx >= len(edit_paras):
                    break
                ep = edit_paras[p_idx]

                # Paragraph alignment
                align_str = ep.get("align", "").upper()
                if align_str in _ALIGN_MAP:
                    try:
                        para.alignment = _ALIGN_MAP[align_str]
                    except Exception:
                        pass

                edit_runs = ep.get("runs", [])
                orig_runs = list(para.runs)
                for r_idx, run in enumerate(orig_runs):
                    if r_idx >= len(edit_runs):
                        # Clear text from original runs that no longer have an edit
                        try:
                            run.text = ""
                        except Exception:
                            pass
                        continue
                    er = edit_runs[r_idx]

                    # Text
                    new_text = er.get("text", run.text)
                    if new_text != run.text:
                        run.text = new_text

                    # Bold
                    if "bold" in er:
                        try:
                            run.font.bold = er["bold"]
                        except Exception:
                            pass

                    # Italic
                    if "italic" in er:
                        try:
                            run.font.italic = er["italic"]
                        except Exception:
                            pass

                    # Underline
                    if "underline" in er:
                        try:
                            run.font.underline = er["underline"]
                        except Exception:
                            pass

                    # Font size
                    if "size" in er and er["size"]:
                        try:
                            run.font.size = Pt(float(er["size"]))
                        except Exception:
                            pass

                    # Font color
                    if "color" in er and er["color"]:
                        try:
                            hex_color = er["color"].lstrip("#")
                            r_val = int(hex_color[0:2], 16)
                            g_val = int(hex_color[2:4], 16)
                            b_val = int(hex_color[4:6], 16)
                            run.font.color.rgb = RGBColor(r_val, g_val, b_val)
                        except Exception:
                            pass

        # ── New shapes (negative IDs = inserted on frontend) ─────────────────
        existing_ids = {s.shape_id for s in slide.shapes}
        for edit_shape in edit_slide.get("shapes", []):
            try:
                sid = int(edit_shape.get("id", 0))
            except (ValueError, TypeError):
                sid = -1
            if sid >= 0 and sid in existing_ids:
                continue  # already handled above
            if not edit_shape.get("has_text"):
                continue
            from pptx.util import Emu
            txBox = slide.shapes.add_textbox(
                Emu(edit_shape.get("left", 0)),
                Emu(edit_shape.get("top", 0)),
                Emu(edit_shape.get("width", 2743200)),
                Emu(edit_shape.get("height", 914400)),
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for p_idx, ep in enumerate(edit_shape.get("paragraphs", [])):
                para = tf.paragraphs[0] if p_idx == 0 else tf.add_paragraph()
                align_str = ep.get("align", "LEFT").upper()
                if align_str in _ALIGN_MAP:
                    try:
                        para.alignment = _ALIGN_MAP[align_str]
                    except Exception:
                        pass
                for er in ep.get("runs", []):
                    run = para.add_run()
                    run.text = er.get("text", "")
                    try:
                        if er.get("bold"):      run.font.bold = True
                        if er.get("italic"):    run.font.italic = True
                        if er.get("underline"): run.font.underline = True
                        if er.get("size"):      run.font.size = Pt(float(er["size"]))
                        if er.get("color"):
                            h = er["color"].lstrip("#")
                            run.font.color.rgb = RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
                    except Exception:
                        pass

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Routes ────────────────────────────────────────────────────────────────────


@pptx_editor_bp.route("/upload", methods=["POST"])
def upload_pptx() -> Response:
    """Upload a .pptx/.ppt file, parse it, and create an editing session."""
    if "file" not in request.files:
        return jsonify({"error": "No file provided"}), 400
    f = request.files["file"]
    if not f.filename:
        return jsonify({"error": "Empty filename"}), 400

    ext = os.path.splitext(f.filename)[1].lower()
    if ext not in (".pptx", ".ppt", ".pptm"):
        return jsonify({"error": "Only .pptx / .pptm / .ppt files are supported"}), 400

    raw = f.read()
    try:
        parsed = _parse_slides(raw)
    except Exception as exc:
        _logger.error("Failed to parse %s: %s", f.filename, exc)
        return jsonify({"error": f"无法解析 PPTX 文件: {exc}"}), 400

    file_id = _new_id()
    now = time.strftime("%Y-%m-%dT%H:%M:%S")
    doc_name = os.path.splitext(f.filename)[0][:200]

    # Persist original bytes for lossless export
    with open(_orig_path(file_id), "wb") as fp:
        fp.write(raw)

    meta = {
        "id": file_id,
        "name": doc_name,
        "originalFilename": f.filename,
        "createdAt": now,
        "updatedAt": now,
        "slideWidthEmu": parsed["slide_width_emu"],
        "slideHeightEmu": parsed["slide_height_emu"],
        "slides": parsed["slides"],
    }
    _write_meta(meta)
    _logger.info(
        "PPTX imported: %s → %s (%d slides)", f.filename, file_id, len(parsed["slides"])
    )
    return (
        jsonify(
            {
                "id": file_id,
                "name": doc_name,
                "slideCount": len(parsed["slides"]),
                "slideWidthEmu": parsed["slide_width_emu"],
                "slideHeightEmu": parsed["slide_height_emu"],
            }
        ),
        201,
    )


@pptx_editor_bp.route("/list", methods=["GET"])
def list_pptx() -> Response:
    """List all uploaded PPTX editing sessions."""
    store = _get_store_dir()
    items = []
    for fname in sorted(os.listdir(store)):
        if not fname.endswith(".json"):
            continue
        try:
            with open(os.path.join(store, fname), "r", encoding="utf-8") as fp:
                meta = json.load(fp)
            items.append(
                {
                    "id": meta["id"],
                    "name": meta.get("name", ""),
                    "originalFilename": meta.get("originalFilename", ""),
                    "slideCount": len(meta.get("slides", [])),
                    "createdAt": meta.get("createdAt", ""),
                    "updatedAt": meta.get("updatedAt", ""),
                }
            )
        except Exception as e:
            _logger.warning("Bad meta file %s: %s", fname, e)
    items.sort(key=lambda x: x.get("updatedAt", ""), reverse=True)
    return jsonify({"files": items})


@pptx_editor_bp.route("/<file_id>", methods=["GET"])
def get_pptx(file_id: str) -> Response:
    """Return the full slides data for editing."""
    meta = _read_meta(file_id)
    if not meta:
        return jsonify({"error": "Not found"}), 404
    return jsonify(meta)


@pptx_editor_bp.route("/<file_id>", methods=["PUT"])
def save_pptx(file_id: str) -> Response:
    """
    Save edited slides data.

    Accepts a JSON body with:
      { "slides": [...],  "name": "optional new name" }
    """
    meta = _read_meta(file_id)
    if not meta:
        return jsonify({"error": "Not found"}), 404

    data = request.get_json(silent=True) or {}
    if "slides" in data:
        meta["slides"] = data["slides"]
    if "name" in data:
        meta["name"] = str(data["name"]).strip()[:200]
    meta["updatedAt"] = time.strftime("%Y-%m-%dT%H:%M:%S")
    _write_meta(meta)
    return jsonify({"ok": True})


@pptx_editor_bp.route("/<file_id>/download", methods=["GET"])
def download_pptx(file_id: str) -> Response:
    """
    Export the edited slides back into the original .pptx and stream it.

    All formatting, images, animations and themes from the original file are
    preserved; only the text runs that were changed are written back.
    """
    meta = _read_meta(file_id)
    if not meta:
        return jsonify({"error": "Not found"}), 404

    orig = _orig_path(file_id)
    if not os.path.exists(orig):
        return jsonify({"error": "Original file missing"}), 500

    with open(orig, "rb") as fp:
        orig_bytes = fp.read()

    try:
        edited_bytes = _apply_edits(orig_bytes, meta.get("slides", []))
    except Exception as exc:
        _logger.error("Export failed for %s: %s", file_id, exc)
        return jsonify({"error": f"导出失败: {exc}"}), 500

    safe_name = re.sub(r'[\\/*?:"<>|]', "_", meta.get("name", "presentation"))
    download_name = f"{safe_name}_edited.pptx"

    return send_file(
        io.BytesIO(edited_bytes),
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=download_name,
    )


@pptx_editor_bp.route("/<file_id>", methods=["DELETE"])
def delete_pptx(file_id: str) -> Response:
    """Delete a PPTX editing session and its stored files."""
    mp = _meta_path(file_id)
    op = _orig_path(file_id)
    if not os.path.exists(mp):
        return jsonify({"error": "Not found"}), 404
    try:
        os.remove(mp)
    except OSError as e:
        _logger.warning("Could not remove meta %s: %s", mp, e)
    if os.path.exists(op):
        try:
            os.remove(op)
        except OSError as e:
            _logger.warning("Could not remove original %s: %s", op, e)
    return jsonify({"ok": True})
