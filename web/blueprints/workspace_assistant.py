# Copyright (C) 2024-2026 Koto AI. All rights reserved.
# SPDX-License-Identifier: LicenseRef-Koto-Proprietary
"""
Koto unified file-workstation BFF.

Retention note:
  The module name `workspace_assistant` is legacy, but the blueprint is still
  active runtime code. It owns the file-workstation API used by the unified
  `/` shell. Do not delete or rename it until every registration, isolated
  Flask test, and `/api/v1/workspace/*` caller has migrated to a new BFF module.

Routes:
  POST /api/v1/workspace/open_file   — 上传并解析文件，返回标准化 JSON
  GET  /api/v1/workspace/raw/<id>    — 返回暂存的原始文件字节（供 PDF.js 渲染）
  POST /api/v1/workspace/save_file   — 接收编辑数据，导出原格式文件并触发下载
"""

from __future__ import annotations

import logging
import os
import uuid
from pathlib import Path

from flask import Blueprint, Response, current_app, jsonify, request, send_file, session, stream_with_context
from app.core.file.file_parser import export_docx, export_xlsx
from app.core.file_assistant import (
    ALLOWED_EXTENSIONS,
    IMAGE_EXTENSIONS,
    TEXT_EXTENSIONS,
    AutoSavePermissionError,
    AutoSavePersistenceService,
    FileContextPreviewService,
    FileAssistantService,
    FileTooLargeError,
    OpenFileByPathService,
    OpenFileCopyError,
    OpenFileEmptyError,
    OpenFileInConfigError,
    OpenFileNotFoundError,
    OpenFilePermissionError,
    OpenFileUnsupportedTypeError,
    PptxPreflightError,
    PptxPreflightService,
    TempFileInvalidIdError,
    TempFileNotFoundError,
    UnsupportedFileTypeError,
    UploadedOpenFileService,
    WorkspaceFileDownloadService,
    WorkspaceFileNotFoundError,
    WorkspaceFilePermissionError,
    WorkspaceFileUnsupportedTypeError,
    WorkspaceFsError,
    WorkspaceFsService,
    WorkspaceTempStore,
    WorkspaceTreeService,
)
from web.shared import PROJECT_ROOT

logger = logging.getLogger(__name__)

workspace_assistant_bp = Blueprint("workspace_assistant", __name__)
_FILE_ASSISTANT = FileAssistantService()
_AUTO_SAVE_PERSISTENCE = AutoSavePersistenceService()
_OPEN_FILE_BY_PATH = OpenFileByPathService()
_UPLOADED_OPEN_FILE = UploadedOpenFileService()
_FILE_CONTEXT_PREVIEW = FileContextPreviewService()
_WORKSPACE_TEMP_STORE = WorkspaceTempStore()
_WORKSPACE_TREE = WorkspaceTreeService()
_WORKSPACE_FILE_DOWNLOAD = WorkspaceFileDownloadService()
_PPTX_PREFLIGHT = PptxPreflightService()
_WORKSPACE_FS = WorkspaceFsService()

# ─── Critical static asset check ─────────────────────────────────────────────
# Prevents silent failures when univer-dist bundle files are missing.
_STATIC_ROOT = Path(__file__).resolve().parent.parent / "static"
_CRITICAL_ASSETS = [
    _STATIC_ROOT / "univer-dist" / "assets" / "sheets-main.js",
    _STATIC_ROOT / "univer-dist" / "assets" / "sheets-main.css",
    _STATIC_ROOT / "js" / "build" / "workspace-bundle.js",
    _STATIC_ROOT / "js" / "build" / "review-bundle.js",
]

def _check_critical_assets() -> list[str]:
    """Return list of missing critical asset paths."""
    missing = []
    for p in _CRITICAL_ASSETS:
        if not p.exists():
            missing.append(str(p))
    return missing

_missing_at_startup = _check_critical_assets()
if _missing_at_startup:
    logger.error(
        "⚠️  关键静态资源缺失！Excel 加载将失败。缺失文件: %s\n"
        "修复方法: 在 web/univer-editor/ 目录下执行 npm run build 重新构建，"
        "或使用 git checkout <commit> -- web/static/univer-dist/ 恢复。",
        ", ".join(_missing_at_startup),
    )

# 临时文件存储目录根（绝对路径，兼容旧的 _TMP_DIR 测试/调用方）
_DEFAULT_TMP_DIR = (Path(PROJECT_ROOT) / "workspace" / "tmp").resolve()
_TMP_DIR = _DEFAULT_TMP_DIR
_TMP_ROOT = _DEFAULT_TMP_DIR

# 纯文本 / 代码文件后缀（直接读取 UTF-8 内容）
_TEXT_EXTS = TEXT_EXTENSIONS

# 图片文件后缀
_IMAGE_EXTS = IMAGE_EXTENSIONS

# 允许上传的文件后缀
_ALLOWED_EXT = ALLOWED_EXTENSIONS


def _get_session_id() -> str:
    """Return a per-browser session ID, creating one if absent.

    This is the only isolation guarantor between users on a shared instance.
    The ID is stored in a signed Flask session cookie so it survives page reloads
    without a database.
    """
    if not current_app.secret_key:
        return "default"
    try:
        sid = session.get("ws_session_id")
        if not sid:
            sid = uuid.uuid4().hex
            session["ws_session_id"] = sid
            session.permanent = True
    except RuntimeError:
        return "default"
    return sid


def _current_tmp_root() -> Path:
    """Return the effective temp root, honoring both legacy and current overrides."""
    tmp_dir = Path(globals().get("_TMP_DIR", _DEFAULT_TMP_DIR))
    tmp_root = Path(globals().get("_TMP_ROOT", _DEFAULT_TMP_DIR))

    if tmp_root != _DEFAULT_TMP_DIR:
        return tmp_root.resolve()
    if tmp_dir != _DEFAULT_TMP_DIR:
        return tmp_dir.resolve()
    return _DEFAULT_TMP_DIR


def _ensure_tmp_dir() -> Path:
    """Return an isolated tmp directory for the current browser session."""
    sid = _get_session_id()
    tmp_dir = _current_tmp_root() / sid
    tmp_dir.mkdir(parents=True, exist_ok=True)
    return tmp_dir


def _tmp_workspace_relpath(file_id: str, ext: str) -> str:
    """Return the workspace-relative path for the current session temp file."""
    sid = _get_session_id()
    return f"tmp/{sid}/{file_id}{ext}"


def _export_workspace_pptx(file_id: str, data) -> bytes:
    """Export the current PPTX editor payload through the rich slide format only."""
    if not file_id or not file_id.isalnum():
        raise ValueError("PPTX 保存需要有效的 file_id")
    if not isinstance(data, dict) or not isinstance(data.get("slides"), list):
        raise ValueError("PPTX 保存仅支持当前 rich slides 数据格式，请重新打开文件后再保存")

    tmp_dir = _ensure_tmp_dir()
    matches = list(tmp_dir.glob(f"{file_id}.pptx"))
    if not matches:
        raise FileNotFoundError("原始 PPTX 文件不存在或已过期")

    from web.blueprints.pptx_editor import _apply_edits as _pptx_apply

    orig_bytes = matches[0].read_bytes()
    return _pptx_apply(orig_bytes, data["slides"])


def _minimal_pdf_bytes() -> bytes:
    """Return a tiny valid one-page PDF placeholder."""
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R >>",
        b"<< /Length 0 >>\nstream\n\nendstream",
    ]
    chunks = [b"%PDF-1.4\n"]
    offsets = [0]
    for idx, obj in enumerate(objects, start=1):
        offsets.append(sum(len(chunk) for chunk in chunks))
        chunks.append(f"{idx} 0 obj\n".encode("ascii"))
        chunks.append(obj)
        chunks.append(b"\nendobj\n")
    xref_offset = sum(len(chunk) for chunk in chunks)
    chunks.append(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    chunks.append(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        chunks.append(f"{offset:010d} 00000 n \n".encode("ascii"))
    chunks.append(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode("ascii")
    )
    return b"".join(chunks)


def _seed_new_file(target: Path) -> None:
    """Create a minimal valid file for formats that cannot be 0-byte."""
    target = Path(target)
    target.parent.mkdir(parents=True, exist_ok=True)
    ext = target.suffix.lower()

    if ext == ".docx":
        from docx import Document

        doc = Document()
        doc.add_paragraph("Koto 新建文档")
        doc.save(str(target))
    elif ext == ".xlsx":
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Koto 新建表格"
        wb.save(str(target))
    elif ext == ".pptx":
        from pptx import Presentation

        prs = Presentation()
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = "Koto 新建演示文稿"
        placeholders = [shape for shape in slide.placeholders if shape != slide.shapes.title]
        if placeholders:
            placeholders[0].text = "在此输入内容"
        prs.save(str(target))
    elif ext == ".pdf":
        target.write_bytes(_minimal_pdf_bytes())
    else:
        target.touch()


def _repair_zero_byte_office_file(target: Path) -> bool:
    """Seed legacy 0-byte workspace documents before parsing them."""
    target = Path(target)
    if target.suffix.lower() not in {".docx", ".xlsx", ".pptx", ".pdf"}:
        return False
    try:
        if not target.is_file() or target.stat().st_size != 0:
            return False
        _seed_new_file(target)
        logger.info("[WorkspaceAssistant] 已修复 0 字节办公文件: %s", target)
        return True
    except Exception as exc:
        logger.warning("[WorkspaceAssistant] 修复 0 字节办公文件失败 %s: %s", target, exc)
        return False


def cleanup_tmp_dir(max_age_hours: int = 24) -> int:
    """Remove stale or 0-byte temp files and return the number removed."""
    import time

    tmp_root = _current_tmp_root()
    if not tmp_root.exists():
        return 0

    removed = 0
    max_age_seconds = max(max_age_hours, 0) * 3600
    now = time.time()

    for path in list(tmp_root.rglob("*")):
        if not path.is_file():
            continue
        try:
            stat = path.stat()
            is_zero_byte = stat.st_size == 0
            is_expired = max_age_seconds > 0 and (now - stat.st_mtime) > max_age_seconds
            if is_zero_byte or is_expired:
                path.unlink(missing_ok=True)
                removed += 1
        except OSError:
            continue

    for path in sorted(tmp_root.rglob("*"), key=lambda p: len(p.parts), reverse=True):
        if path == tmp_root or not path.is_dir():
            continue
        try:
            path.rmdir()
        except OSError:
            pass

    return removed


def _tmp_file_path(file_id: str, ext: str) -> Path:
    return _ensure_tmp_dir() / f"{file_id}{ext}"


def _copy_file_to_tmp_with_validation(src_path: Path, tmp_path: Path, *, ext: str, attempts: int = 2) -> None:
    import shutil
    import zipfile

    src_path = Path(src_path).resolve()
    tmp_path = Path(tmp_path).resolve()
    last_error: Exception | None = None

    for attempt in range(1, max(attempts, 1) + 1):
        try:
            tmp_path.parent.mkdir(parents=True, exist_ok=True)
            tmp_path.unlink(missing_ok=True)
        except Exception:
            pass

        try:
            if attempt == 1:
                shutil.copy2(str(src_path), str(tmp_path))
            else:
                tmp_path.write_bytes(src_path.read_bytes())

            if not tmp_path.is_file():
                raise FileNotFoundError(f"临时副本不存在: {tmp_path}")

            src_size = src_path.stat().st_size
            tmp_size = tmp_path.stat().st_size
            if src_size != tmp_size:
                raise RuntimeError(
                    f"临时副本大小异常: src={src_size} tmp={tmp_size}"
                )

            if ext.lower() in {".docx", ".xlsx", ".pptx"} and src_size > 0:
                if not zipfile.is_zipfile(tmp_path):
                    raise zipfile.BadZipFile(
                        f"临时副本不是合法的 {ext} ZIP 文件: {tmp_path}"
                    )
            return
        except Exception as exc:
            last_error = exc
            logger.warning(
                "[WorkspaceAssistant] tmp copy validation failed %s -> %s (attempt %d/%d): %s",
                src_path,
                tmp_path,
                attempt,
                max(attempts, 1),
                exc,
            )

    raise RuntimeError(f"临时副本创建失败: {last_error}")


def _build_workspace_capability_profile(*, file_type: str, path: str = "", name: str = "") -> dict:
    from app.core.agent.file_task_capability import build_file_capability_profile

    return build_file_capability_profile(file_type=file_type, path=path, name=name)


def _workspace_open_payload(
    *,
    file_id: str,
    file_type: str,
    data: dict,
    file_name: str | None = None,
    ws_source_path: str | None = None,
    temp_path: str | None = None,
    source_path: str = "",
) -> dict:
    payload = {
        "file_id": file_id,
        "file_type": file_type,
        "data": data,
        "capability_profile": _build_workspace_capability_profile(
            file_type=file_type,
            path=source_path,
            name=file_name or "",
        ),
    }
    if file_name is not None:
        payload["file_name"] = file_name
    if ws_source_path is not None:
        payload["ws_source_path"] = ws_source_path
    if temp_path is not None:
        payload["temp_path"] = temp_path
    return payload


# ─────────────────────────────────────────────────────────────────────────────
# GET /api/v1/workspace/asset_health — verify critical static assets exist
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/asset_health")
def asset_health():
    """Quick check that all critical frontend assets are present on disk."""
    missing = _check_critical_assets()
    if missing:
        return jsonify({"ok": False, "missing": missing}), 500
    return jsonify({"ok": True, "missing": []})


# ─────────────────────────────────────────────────────────────────────────────
# GET /api/v1/workspace/list_files
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/current_dir")
def get_current_workspace_dir():
    """Return the current workspace root path and display name."""
    from web.shared import WORKSPACE_DIR

    p = Path(WORKSPACE_DIR).resolve()
    return jsonify({"path": str(p), "name": p.name})


@workspace_assistant_bp.route("/api/v1/workspace/list_files")
def list_workspace_files():
    """
    获取工作区的文件树结构（全文件类型）。
    过滤掉隐藏文件和系统文件夹，以树状 JSON 返回。
    supported=true 表示 Koto 可直接打开和编辑该文件。
    """
    from web.shared import WORKSPACE_DIR

    root_path = Path(WORKSPACE_DIR).resolve()
    cleanup_tmp_dir()
    return jsonify(
        _WORKSPACE_TREE.build_workspace_tree(
            root_path=root_path,
            allowed_extensions=_ALLOWED_EXT,
            skip_names=WorkspaceTreeService.DEFAULT_SKIP_NAMES | {"ppt_sessions"},
        )
    )


# ─────────────────────────────────────────────────────────────────────────────
# GET /api/v1/workspace/file/<path:filepath>
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/file/<path:filepath>")
def serve_workspace_file(filepath: str):
    """
    从工作区目录安全地提供文件下载（用于左侧文件树点击打开）。
    防范路径遍历攻击，仅返回支持格式的文件。
    """
    from web.shared import WORKSPACE_DIR

    try:
        served = _WORKSPACE_FILE_DOWNLOAD.serve_file(
            workspace_dir=WORKSPACE_DIR,
            filepath=filepath,
            allowed_extensions=_ALLOWED_EXT,
        )
    except WorkspaceFilePermissionError:
        return jsonify({"error": "路径不合法"}), 403
    except WorkspaceFileNotFoundError:
        return jsonify({"error": "文件不存在"}), 404
    except WorkspaceFileUnsupportedTypeError:
        return jsonify({"error": "不支持的文件类型"}), 400

    return send_file(
        str(served.path),
        mimetype=served.mime_type,
        as_attachment=False,
        download_name=served.download_name,
    )


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/open_file_by_path
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/open_file_by_path", methods=["POST"])
def open_file_by_path():
    """
    直接从工作区路径打开并解析文件，无需先下载再上传。
    比 open_file 更高效，用于工作区文件树点击打开。
    Body (JSON): {"path": "relative/path/to/file.docx"}
    Response: 同 open_file
    """
    body = request.get_json(force=True, silent=True) or {}
    raw_path = (body.get("path") or "").strip()
    if not raw_path:
        return jsonify({"error": "缺少 path 字段"}), 400

    file_id = uuid.uuid4().hex
    try:
        from web.shared import WORKSPACE_DIR

        prepared = _OPEN_FILE_BY_PATH.prepare(
            raw_path=raw_path,
            workspace_dir=WORKSPACE_DIR,
            app_config_dir=_APP_CONFIG_DIR,
            tmp_dir=_ensure_tmp_dir(),
            file_id=file_id,
            allowed_extensions=_ALLOWED_EXT,
            fs_guard=_fs_guard,
            repair_zero_byte_file=_repair_zero_byte_office_file,
            copy_to_tmp=_copy_file_to_tmp_with_validation,
        )
    except OpenFileInConfigError:
        return jsonify({"error": "不允许访问应用配置目录"}), 403
    except OpenFilePermissionError:
        return jsonify({"error": "路径不合法"}), 403
    except OpenFileNotFoundError:
        return jsonify({"error": "文件不存在"}), 404
    except OpenFileUnsupportedTypeError as e:
        return jsonify({"error": str(e)}), 400
    except OpenFileCopyError as e:
        return jsonify({"error": str(e)}), 500

    target = prepared.target_path
    tmp_path = prepared.tmp_path
    ext = prepared.extension

    try:
        parsed = _FILE_ASSISTANT.parse_editor_file(
            tmp_path,
            file_id=file_id,
            display_name=target.name,
            text_source_path=target,
            source_path=target,
            docx_copy_to_tmp=_copy_file_to_tmp_with_validation,
        )
        file_type = parsed.file_type
        data = parsed.data
        if file_type == "docx":
            html_len = len(data.get("html", ""))
            logger.info(
                "[open_file_by_path] %s 解析成功, HTML=%dKB, messages=%s",
                target.name,
                html_len // 1024,
                data.get("messages", []),
            )

    except FileTooLargeError as e:
        return jsonify({"error": str(e)}), e.status_code
    except UnsupportedFileTypeError as e:
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        logger.error(f"[WorkspaceAssistant] 解析失败 {target.name}: {e}", exc_info=True)
        try:
            tmp_path.unlink(missing_ok=True)
        except Exception:
            pass
        return jsonify({"error": f"文件解析失败: {str(e)}"}), 500

    return jsonify(
        _workspace_open_payload(
            file_id=file_id,
            file_name=target.name,
            file_type=file_type,
            ws_source_path=raw_path,
            temp_path=_tmp_workspace_relpath(file_id, ext),
            data=data,
            source_path=str(target),
        )
    )


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/open_file
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/open_file", methods=["POST"])
def open_file():
    """
    接受 multipart/form-data 上传，解析文件并返回适合前端多态渲染器的标准 JSON。
    Body: file=<binary>
    Response:
      {"file_id": str, "file_name": str, "file_type": str, "data": <parsed>}
    """
    if "file" not in request.files:
        return jsonify({"error": "缺少 file 字段"}), 400

    uploaded = request.files["file"]
    file_id = uuid.uuid4().hex
    try:
        prepared = _UPLOADED_OPEN_FILE.prepare(
            original_name=uploaded.filename or "unknown",
            tmp_dir=_ensure_tmp_dir(),
            file_id=file_id,
            allowed_extensions=_ALLOWED_EXT,
            save_upload=lambda target: uploaded.save(str(target)),
        )
    except OpenFileUnsupportedTypeError as e:
        return jsonify({"error": str(e)}), 400
    except OpenFileEmptyError as e:
        return jsonify({"error": str(e)}), 400

    original_name = prepared.original_name
    tmp_path = prepared.tmp_path
    ext = prepared.extension

    # 文件只暂存在 tmp 目录，不立即写入工作区。
    # 用户显式保存后才会写入 WORKSPACE_DIR（由 auto_save explicit=true 处理）。
    ws_path = request.form.get("ws_path", "").strip()

    try:
        parsed = _FILE_ASSISTANT.parse_editor_file(
            tmp_path,
            file_id=file_id,
            display_name=original_name,
            docx_copy_to_tmp=_copy_file_to_tmp_with_validation,
        )
        file_type = parsed.file_type
        data = parsed.data

    except FileTooLargeError as e:
        return jsonify({"error": str(e)}), e.status_code
    except UnsupportedFileTypeError as e:
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        logger.error(
            f"[WorkspaceAssistant] 解析失败 {original_name}: {e}", exc_info=True
        )
        # 清理临时文件
        try:
            tmp_path.unlink(missing_ok=True)
        except Exception:
            pass
        return jsonify({"error": f"文件解析失败: {str(e)}"}), 500

    return jsonify(
        _workspace_open_payload(
            file_id=file_id,
            file_name=original_name,
            file_type=file_type,
            ws_source_path=ws_path or "",
            temp_path=_tmp_workspace_relpath(file_id, ext),
            data=data,
            source_path=ws_path or original_name,
        )
    )


@workspace_assistant_bp.route("/api/v1/workspace/ai_context_preview", methods=["POST"])
def ai_context_preview():
    """Return a lightweight text preview for AI file attachments.

    Unlike editor open routes, this endpoint avoids shipping full rich parser
    payloads through the WebView. It is used by the AI side-panel attachment
    flow, which only needs a bounded text preview.
    """
    body = request.get_json(force=True, silent=True) or {}
    raw_path = (body.get("path") or "").strip()
    if not raw_path:
        return jsonify({"error": "缺少 path 字段"}), 400

    try:
        preview_limit = int(body.get("max_chars") or 12_000)
    except (TypeError, ValueError):
        preview_limit = 12_000

    try:
        from web.shared import WORKSPACE_DIR

        preview = _FILE_CONTEXT_PREVIEW.build(
            raw_path=raw_path,
            workspace_dir=WORKSPACE_DIR,
            app_config_dir=_APP_CONFIG_DIR,
            allowed_extensions=_ALLOWED_EXT,
            fs_guard=_fs_guard,
            preview_limit=preview_limit,
        )
    except OpenFileInConfigError:
        return jsonify({"error": "不允许访问应用配置目录"}), 403
    except OpenFilePermissionError:
        return jsonify({"error": "路径不合法"}), 403
    except OpenFileNotFoundError:
        return jsonify({"error": "文件不存在"}), 404
    except OpenFileUnsupportedTypeError as e:
        return jsonify({"error": str(e)}), 400

    return jsonify(
        {
            "path": preview.path,
            "file_name": preview.file_name,
            "file_type": preview.file_type,
            "content_preview": preview.content_preview,
            "original_chars": preview.original_chars,
            "preview_error": preview.preview_error,
        }
    )


# ─────────────────────────────────────────────────────────────────────────────
# GET /api/v1/workspace/raw/<file_id>
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/raw/<file_id>")
def raw_file(file_id: str):
    """
    返回暂存的原始文件（用于 PDF.js 直接渲染）。
    file_id 只允许十六进制字符，防止路径遍历。
    """
    try:
        raw = _WORKSPACE_TEMP_STORE.raw_file(tmp_dir=_ensure_tmp_dir(), file_id=file_id)
    except TempFileInvalidIdError:
        return jsonify({"error": "无效的 file_id"}), 400
    except TempFileNotFoundError:
        return jsonify({"error": "文件不存在或已过期"}), 404

    resp = send_file(str(raw.path), mimetype=raw.mime_type)
    # Prevent browser from caching — each save produces new bytes at the same URL
    for key, value in raw.headers.items():
        resp.headers[key] = value
    return resp


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/save_file
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/save_file", methods=["POST"])
def save_file():
    """
    接收编辑后数据，导出为原格式二进制并触发浏览器下载。
    Body (JSON):
      {"file_type": "docx"|"xlsx"|"pptx",
       "file_id": str,          # PPTX 需要原始文件 ID
       "data": <editor_payload>,
       "file_name": str}        # 可选，用于下载文件名
    """
    body = request.get_json(force=True, silent=True) or {}
    file_type = body.get("file_type", "").lower()
    file_id = body.get("file_id", "")
    data = body.get("data")
    file_name = body.get("file_name", f"koto_export.{file_type}")

    if not file_type or data is None:
        return jsonify({"error": "缺少 file_type 或 data 字段"}), 400

    try:
        exported = _FILE_ASSISTANT.export_editor_file(
            file_type=file_type,
            file_id=file_id,
            data=data,
            file_name=file_name,
            tmp_dir=_ensure_tmp_dir(),
            pptx_exporter=_export_workspace_pptx,
        )

    except FileNotFoundError as e:
        return jsonify({"error": str(e)}), 404
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        logger.error(f"[WorkspaceAssistant] 导出失败 {file_type}: {e}", exc_info=True)
        return jsonify({"error": f"导出失败: {str(e)}"}), 500

    import io

    return send_file(
        io.BytesIO(exported.raw_bytes),
        mimetype=exported.mime,
        as_attachment=True,
        download_name=exported.file_name,
    )


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/upload_image
# GET  /api/v1/workspace/tmp_image/<session_id>/<filename>
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/upload_image", methods=["POST"])
def upload_image():
    """
    Upload a chart/image data URI, save it as a session-scoped temp file, and
    return a stable server URL.

    Converting the multi-MB base64 string to a short server URL before inserting
    into WangEditor/Slate prevents the virtual DOM from bloating (which freezes
    the browser) and keeps the exported DOCX clean (the export pipeline can fetch
    the bytes from disk instead of embedding a giant base64 inline attribute).

    Body (JSON): {"data": "data:image/png;base64,..."}
    Returns:     {"ok": true, "url": "/api/v1/workspace/tmp_image/{sid}/{uuid}.png"}
    """
    import base64 as _b64
    import mimetypes as _mime

    body = request.get_json(force=True, silent=True) or {}
    data_uri = body.get("data", "")
    if not isinstance(data_uri, str) or not data_uri.startswith("data:image/"):
        return jsonify({"error": "无效的图片数据 (须为 data:image/... URI)"}), 400

    # Parse: data:image/png;base64,<b64_payload>
    try:
        header, _, b64_str = data_uri.partition(",")
        mime = header.split(":")[1].split(";")[0]   # e.g. "image/png"
        img_bytes = _b64.b64decode(b64_str)
    except Exception:
        return jsonify({"error": "图片数据解码失败"}), 400

    if len(img_bytes) > 10 * 1024 * 1024:
        return jsonify({"error": "图片大小超过 10 MB 限制"}), 413

    # Map MIME type to a safe file extension
    ext = _mime.guess_extension(mime) or ".png"
    if ext in (".jpe", ".jfif"):
        ext = ".jpg"
    if ext not in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"):
        ext = ".png"

    # Save to session-scoped tmp/images/ directory
    img_id = uuid.uuid4().hex
    sid = _get_session_id()
    img_dir = _current_tmp_root() / sid / "images"
    img_dir.mkdir(parents=True, exist_ok=True)
    img_path = img_dir / f"{img_id}{ext}"
    img_path.write_bytes(img_bytes)

    url = f"/api/v1/workspace/tmp_image/{sid}/{img_id}{ext}"
    logger.debug("[upload_image] saved %d bytes → %s", len(img_bytes), url)
    return jsonify({"ok": True, "url": url})


@workspace_assistant_bp.route("/api/v1/workspace/save_to_workspace", methods=["POST"])
def save_to_workspace():
    """
    Save an AI-generated asset (image or file) directly into the workspace
    directory so it appears in the file tree.

    Body (JSON) — image mode:
      {"type": "image",
       "src_url": "/api/v1/workspace/tmp_image/<sid>/<uuid.png>",
       "filename": "chart.png"}

    Body (JSON) — file mode:
      {"type": "file",
       "data": "<base64-encoded bytes>",
       "filename": "修改后_foo.docx"}

    Returns: {"ok": true, "ws_path": "images/chart.png"}
    """
    import base64 as _b64
    import shutil

    body = request.get_json(force=True, silent=True) or {}
    asset_type = body.get("type", "")
    filename = (body.get("filename") or "").strip()

    # Validate filename: no path separators
    if not filename or "/" in filename or "\\" in filename or ".." in filename:
        return jsonify({"error": "无效的文件名"}), 400

    from web.shared import WORKSPACE_DIR
    ws_root = Path(WORKSPACE_DIR).resolve()

    if asset_type == "image":
        src_url = body.get("src_url", "")
        # Expected: /api/v1/workspace/tmp_image/<session_id>/<filename>
        prefix = "/api/v1/workspace/tmp_image/"
        if not src_url.startswith(prefix):
            return jsonify({"error": "无效的图片 URL"}), 400
        parts = src_url[len(prefix):].split("/")
        if len(parts) != 2:
            return jsonify({"error": "无效的图片路径"}), 400
        sid, img_fname = parts
        # Validate both components
        if len(sid) != 32 or not all(c in "0123456789abcdef" for c in sid):
            return jsonify({"error": "无效的 session_id"}), 400
        if not img_fname or "/" in img_fname or "\\" in img_fname or ".." in img_fname:
            return jsonify({"error": "无效的图片文件名"}), 400
        tmp_root = _current_tmp_root().resolve()
        src_path = (tmp_root / sid / "images" / img_fname).resolve()
        try:
            src_path.relative_to(tmp_root)
        except ValueError:
            return jsonify({"error": "路径非法"}), 403
        if not src_path.is_file():
            return jsonify({"error": "临时图片不存在或已过期"}), 404
        dest_dir = ws_root / "images"
        dest_dir.mkdir(parents=True, exist_ok=True)
        dest = _unique_path(dest_dir, filename)
        shutil.copy2(str(src_path), str(dest))

    elif asset_type == "file":
        data_b64 = body.get("data", "")
        if not data_b64:
            return jsonify({"error": "缺少 data 字段"}), 400
        try:
            file_bytes = _b64.b64decode(data_b64)
        except Exception:
            return jsonify({"error": "数据解码失败"}), 400
        if len(file_bytes) > 50 * 1024 * 1024:
            return jsonify({"error": "文件大小超过 50 MB"}), 413
        ext = Path(filename).suffix.lower()
        if ext not in _ALLOWED_EXT:
            return jsonify({"error": f"不支持的格式: {ext}"}), 400
        dest = _unique_path(ws_root, filename)
        dest.write_bytes(file_bytes)

    else:
        return jsonify({"error": "type 须为 image 或 file"}), 400

    ws_rel = str(dest.relative_to(ws_root)).replace("\\", "/")
    logger.info("[save_to_workspace] saved %s → %s", asset_type, ws_rel)
    return jsonify({"ok": True, "ws_path": ws_rel})


def _unique_path(directory: Path, filename: str) -> Path:
    """Return a non-colliding path inside directory, appending _1, _2 … as needed."""
    candidate = directory / filename
    if not candidate.exists():
        return candidate
    stem = Path(filename).stem
    suffix = Path(filename).suffix
    i = 1
    while True:
        candidate = directory / f"{stem}_{i}{suffix}"
        if not candidate.exists():
            return candidate
        i += 1


@workspace_assistant_bp.route("/api/v1/workspace/tmp_image/<session_id>/<filename>")
def serve_tmp_image(session_id: str, filename: str):
    """
    Serve a previously uploaded temp image file by its session-scoped URL.
    Both session_id and filename are validated to prevent path traversal.
    """
    # Validate session_id: must be exactly 32 lowercase hex characters (uuid4().hex)
    if len(session_id) != 32 or not all(c in "0123456789abcdef" for c in session_id):
        return jsonify({"error": "无效的 session_id"}), 400
    # Validate filename: no path separators or parent-dir references
    if not filename or "/" in filename or "\\" in filename or ".." in filename:
        return jsonify({"error": "无效的文件名"}), 400

    tmp_root = _current_tmp_root().resolve()
    img_path = (tmp_root / session_id / "images" / filename).resolve()
    # Path-traversal guard (belt-and-suspenders)
    try:
        img_path.relative_to(tmp_root)
    except ValueError:
        return jsonify({"error": "路径非法"}), 403

    if not img_path.is_file():
        return jsonify({"error": "图片不存在或已过期"}), 404

    _mime_map = {
        ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".gif": "image/gif", ".webp": "image/webp", ".svg": "image/svg+xml",
    }
    mime = _mime_map.get(img_path.suffix.lower(), "image/png")
    resp = send_file(str(img_path), mimetype=mime)
    resp.headers["Cache-Control"] = "public, max-age=3600"
    return resp


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/replace_image
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/replace_image", methods=["POST"])
def replace_image():
    """
    Replace a picture shape in a PPTX with a new uploaded image.
    Overwrites the tmp file in-place so subsequent save_file/auto_save exports
    use the new image without requiring a re-upload.

    Body: multipart/form-data
      file_id     — str (hex)
      slide_index — int
      shape_id    — int
      image       — file  (image/*)

    Returns: {"ok": true, "image_b64": "data:image/...;base64,..."}
    """
    file_id = request.form.get("file_id", "").strip()
    slide_index_str = request.form.get("slide_index", "").strip()
    shape_id_str = request.form.get("shape_id", "").strip()

    if not file_id or not file_id.isalnum():
        return jsonify({"error": "无效的 file_id"}), 400
    if not slide_index_str.isdigit() or not shape_id_str.isdigit():
        return jsonify({"error": "slide_index 和 shape_id 必须是整数"}), 400
    if "image" not in request.files:
        return jsonify({"error": "缺少 image 字段"}), 400

    img_file = request.files["image"]
    content_type = img_file.content_type or ""
    if not content_type.startswith("image/"):
        return jsonify({"error": "仅支持图片格式 (image/*)"}), 400

    img_bytes = img_file.read()
    if not img_bytes:
        return jsonify({"error": "图片文件为空"}), 400

    slide_index = int(slide_index_str)
    shape_id = int(shape_id_str)

    tmp_dir = _ensure_tmp_dir()
    matches = list(tmp_dir.glob(f"{file_id}.pptx"))
    if not matches:
        return jsonify({"error": "原始 PPTX 文件不存在或已过期"}), 404

    pptx_path = str(matches[0])
    try:
        import io as _io

        from pptx import Presentation

        prs = Presentation(pptx_path)
        if slide_index >= len(prs.slides):
            return jsonify({"error": "幻灯片序号超出范围"}), 400

        slide = prs.slides[slide_index]
        target_shape = None
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                target_shape = shape
                break

        if target_shape is None:
            return jsonify({"error": "未找到指定形状"}), 404

        left = target_shape.left
        top = target_shape.top
        width = target_shape.width
        height = target_shape.height

        # Remove old picture element from slide XML
        sp_elem = target_shape._element
        sp_elem.getparent().remove(sp_elem)

        # Insert new picture at the same position/size
        slide.shapes.add_picture(_io.BytesIO(img_bytes), left, top, width, height)

        # Overwrite tmp file in-place so later export/auto_save uses updated image
        prs.save(pptx_path)

    except Exception as e:
        logger.error(f"[WorkspaceAssistant] replace_image 失败: {e}", exc_info=True)
        return jsonify({"error": f"替换失败: {str(e)}"}), 500

    import base64 as _b64

    b64 = _b64.b64encode(img_bytes).decode("ascii")
    return jsonify({"ok": True, "image_b64": f"data:{content_type};base64,{b64}"})


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/auto_save
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/auto_save", methods=["POST"])
def auto_save():
    """
    Silently save edited content back to the workspace tmp file AND the original
    workspace source file (so switching back to the file shows latest changes).
    No download — just persists the current state so it survives a reload.
    Body (JSON):
      {"file_type": "docx"|"xlsx"|"pptx",
       "file_id": str,
       "ws_source_path": str,   # optional: workspace-relative path to overwrite
       "data": <editor_payload>}
    Returns: {"ok": true, "saved_at": "<ISO timestamp>"}
    """
    body = request.get_json(force=True, silent=True) or {}
    file_type = body.get("file_type", "").lower()
    file_id = body.get("file_id", "")
    ws_source_path = body.get("ws_source_path", "")  # e.g. "foo.docx"
    data = body.get("data")

    if not file_type or not file_id or data is None:
        return jsonify({"error": "缺少必要字段"}), 400
    if not file_id.isalnum():
        return jsonify({"error": "无效的 file_id"}), 400

    explicit = body.get("explicit", False)
    data_len = len(data) if isinstance(data, str) else (len(str(data)) if data else 0)
    logger.info(
        "[auto_save] explicit=%s file_id=%s...%s data_len=%d preview=%.120s",
        explicit,
        file_id[:8],
        file_id[-4:],
        data_len,
        (data[:120] if isinstance(data, str) else str(data)[:120]),
    )

    try:
        tmp_dir = _ensure_tmp_dir()
        auto_save_name = f"koto_autosave.{file_type}"
        if file_type in ("text", "code"):
            existing = [f for f in tmp_dir.glob(f"{file_id}.*") if f.suffix.lower() in _TEXT_EXTS]
            auto_save_name = f"koto_autosave{existing[0].suffix.lower() if existing else '.txt'}"

        exported = _FILE_ASSISTANT.export_editor_file(
            file_type=file_type,
            file_id=file_id,
            data=data,
            file_name=auto_save_name,
            tmp_dir=tmp_dir,
            pptx_exporter=_export_workspace_pptx,
        )
        raw_bytes = exported.raw_bytes
        suffix = exported.suffix

    except FileNotFoundError as e:
        return jsonify({"error": str(e)}), 404
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except Exception as e:
        logger.error(
            "[WorkspaceAssistant] auto_save 失败 %s: %s", file_type, e, exc_info=True
        )
        return jsonify({"error": f"自动保存失败: {str(e)}"}), 500

    try:
        from web.shared import WORKSPACE_DIR

        persisted = _AUTO_SAVE_PERSISTENCE.persist(
            tmp_dir=_ensure_tmp_dir(),
            file_id=file_id,
            raw_bytes=raw_bytes,
            suffix=suffix,
            explicit=bool(explicit),
            ws_source_path=ws_source_path,
            workspace_dir=WORKSPACE_DIR,
            allowed_extensions=_ALLOWED_EXT,
            fs_guard=_fs_guard,
        )
    except AutoSavePermissionError:
        return jsonify({"error": "路径不合法"}), 403
    except Exception as e:
        logger.warning(
            "[WorkspaceAssistant] auto_save: could not write source file: %s", e
        )
        if explicit:
            return jsonify({"error": f"保存失败: {str(e)}"}), 500
        return jsonify({"error": f"自动保存失败: {str(e)}"}), 500

    return jsonify(
        {
            "ok": True,
            "saved_at": persisted.saved_at,
            "src_written": persisted.src_written,
        }
    )


# ─────────────────────────────────────────────────────────────────────────────
# GET /api/v1/workspace/versions   — list version snapshots for a file
# POST /api/v1/workspace/restore-version — restore a snapshot
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/versions", methods=["GET"])
def list_versions():
    """
    列出文件的历史版本快照。
    Query: path=relative/path/to/file.docx
    Returns: {"versions": [{name, snap_path, saved_at, size_bytes}]}
    """
    from web.shared import WORKSPACE_DIR

    rel = request.args.get("path", "").strip()
    if not rel:
        return jsonify({"error": "缺少 path 参数"}), 400

    ws_root = Path(WORKSPACE_DIR).resolve()
    src = ws_root.joinpath(rel).resolve()
    try:
        src.relative_to(ws_root)
    except ValueError:
        return jsonify({"error": "路径非法"}), 400

    snap_dir = src.parent / ".koto_versions" / src.stem
    if not snap_dir.is_dir():
        return jsonify({"versions": []})

    suffix = src.suffix.lower()
    snaps = sorted(snap_dir.glob(f"*{suffix}"), reverse=True)
    result = []
    for s in snaps[:20]:
        try:
            stat = s.stat()
            result.append({
                "name": s.name,
                "snap_path": str(s),
                "saved_at": s.stem.replace("_", " "),
                "size_bytes": stat.st_size,
            })
        except Exception:
            pass
    return jsonify({"versions": result})


@workspace_assistant_bp.route("/api/v1/workspace/restore-version", methods=["POST"])
def restore_version():
    """
    将版本快照恢复为当前文件。
    Body JSON: { "snap_path": "abs/path/to/snapshot.docx", "target_path": "relative/target.docx" }
    """
    import shutil
    from web.shared import WORKSPACE_DIR

    body = request.get_json(force=True, silent=True) or {}
    snap_path_str = body.get("snap_path", "").strip()
    target_rel = body.get("target_path", "").strip()
    if not snap_path_str or not target_rel:
        return jsonify({"error": "缺少 snap_path 或 target_path"}), 400

    snap = Path(snap_path_str)
    if not snap.is_file():
        return jsonify({"error": "快照文件不存在"}), 404

    ws_root = Path(WORKSPACE_DIR).resolve()
    target = ws_root.joinpath(target_rel).resolve()
    try:
        target.relative_to(ws_root)
    except ValueError:
        return jsonify({"error": "target_path 非法"}), 400
    if target.suffix.lower() not in _ALLOWED_EXT:
        return jsonify({"error": "不支持的格式"}), 400

    try:
        shutil.copy2(str(snap), str(target))
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    return jsonify({"ok": True, "restored_to": str(target)})


@workspace_assistant_bp.route("/api/v1/workspace/checkpoint", methods=["POST"])
def create_checkpoint():
    """
    Create a pre-agent checkpoint for a file (Step 4.3).
    Body JSON: { "path": "relative/path/to/file.docx", "label": "agent_pre" }
    Returns: { "ok": true, "snap_path": "..." }
    """
    import shutil
    from web.shared import WORKSPACE_DIR

    body = request.get_json(force=True, silent=True) or {}
    rel = body.get("path", "").strip()
    label = body.get("label", "agent_pre")
    if not rel:
        return jsonify({"error": "缺少 path 参数"}), 400

    ws_root = Path(WORKSPACE_DIR).resolve()
    src = ws_root.joinpath(rel).resolve()
    try:
        src.relative_to(ws_root)
    except ValueError:
        return jsonify({"error": "路径非法"}), 400

    if not src.is_file():
        return jsonify({"error": "文件不存在"}), 404

    suffix = src.suffix.lower()
    snap_dir = src.parent / ".koto_versions" / src.stem
    snap_dir.mkdir(parents=True, exist_ok=True)
    ts_str = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    snap_path = snap_dir / f"{ts_str}_{label}{suffix}"
    try:
        shutil.copy2(str(src), str(snap_path))
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    return jsonify({"ok": True, "snap_path": str(snap_path), "target_path": rel})


# ─────────────────────────────────────────────────────────────────────────────
# DELETE /api/v1/workspace/file  ?path=<relative_path>
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/file", methods=["DELETE"])
def delete_workspace_file():
    """
    删除工作区中的一个文件（用于左侧面板的 × 按钮）。
    Query param:  path=relative/path/to/file
    """
    from web.shared import WORKSPACE_DIR

    filepath = request.args.get("path", "").strip()
    try:
        _WORKSPACE_FS.delete_file(
            workspace_dir=WORKSPACE_DIR,
            rel_path=filepath,
            allowed_extensions=_ALLOWED_EXT,
        )
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code

    return jsonify({"ok": True})


# ─────────────────────────────────────────────────────────────────────────────
# PATCH /api/v1/workspace/rename
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/rename", methods=["PATCH"])
def rename_workspace_file():
    """
    重命名工作区文件。
    Body (JSON): {"path": "uploads/old.docx", "name": "new_name.docx"}
    Extension must stay the same; new_name must not contain path separators.
    """
    from web.shared import WORKSPACE_DIR

    body = request.get_json(silent=True) or {}
    old_path = body.get("path", "").strip()
    new_name = body.get("name", "").strip()
    try:
        result = _WORKSPACE_FS.rename(
            workspace_dir=WORKSPACE_DIR,
            rel_path=old_path,
            new_name=new_name,
        )
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code

    logger.info("[WorkspaceAssistant] 重命名: %s -> %s", old_path, result.path)
    return jsonify({"ok": True, "path": result.path, "name": result.name})


# ─────────────────────────────────────────────────────────────────────────────
# DELETE /api/v1/workspace/folder
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/folder", methods=["DELETE"])
def delete_workspace_folder():
    """
    递归删除工作区中的一个文件夹。
    Query param:  path=relative/path/to/folder
    """
    from web.shared import WORKSPACE_DIR

    folderpath = request.args.get("path", "").strip()
    try:
        _WORKSPACE_FS.delete_folder(
            workspace_dir=WORKSPACE_DIR,
            rel_path=folderpath,
        )
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code

    return jsonify({"ok": True})


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/create_file
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/create_file", methods=["POST"])
def create_workspace_file():
    """
    在工作区指定目录创建一个新文件（空文件）。
    Body (JSON): {"folder": "relative/path", "name": "filename.txt"}
    folder 为 "" 时在工作区根目录创建。
    """
    from web.shared import WORKSPACE_DIR

    body = request.get_json(force=True, silent=True) or {}
    folder = (body.get("folder") or "").strip().strip("/")
    name = (body.get("name") or "").strip()

    try:
        result = _WORKSPACE_FS.create_file(
            workspace_dir=WORKSPACE_DIR,
            folder=folder,
            name=name,
            allowed_extensions=_ALLOWED_EXT,
            seed_file=_seed_new_file,
        )
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code

    logger.info("[WorkspaceAssistant] 创建文件: %s", result.path)
    return jsonify({"ok": True, "path": result.path, "name": result.name})


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/create_folder
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/create_folder", methods=["POST"])
def create_workspace_folder():
    """
    在工作区指定目录创建一个新文件夹。
    Body (JSON): {"parent": "relative/path", "name": "foldername"}
    parent 为 "" 时在工作区根目录创建。
    """
    from web.shared import WORKSPACE_DIR

    body = request.get_json(force=True, silent=True) or {}
    parent_rel = (body.get("parent") or "").strip().strip("/")
    name = (body.get("name") or "").strip()

    try:
        result = _WORKSPACE_FS.create_folder(
            workspace_dir=WORKSPACE_DIR,
            parent_rel=parent_rel,
            name=name,
        )
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code

    logger.info("[WorkspaceAssistant] 创建文件夹: %s", result.path)
    return jsonify({"ok": True, "path": result.path, "name": result.name})


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/set_workspace_dir
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/set_workspace_dir", methods=["POST"])
def set_workspace_dir_endpoint():
    """
    将工作区根目录切换到指定的本地文件夹。
    Body (JSON): {"path": "/absolute/path/to/folder"}
    持久化到 config/user_settings.json 并立即生效（无需重启）。
    """
    import json as _json

    body = request.get_json(force=True, silent=True) or {}
    new_path = (body.get("path") or "").strip()
    if not new_path:
        return jsonify({"error": "缺少 path 字段"}), 400

    target = Path(new_path).resolve()
    if not _fs_guard(target):
        return jsonify({"error": "不允许将系统路径设为工作区"}), 403
    if not target.exists():
        try:
            target.mkdir(parents=True, exist_ok=True)
        except Exception as e:
            return jsonify({"error": f"无法创建目录: {e}"}), 400
    if not target.is_dir():
        return jsonify({"error": "路径不是文件夹"}), 400

    # Persist to user_settings.json
    from web.shared import clear_user_settings_cache, get_user_settings_path
    settings_path = Path(get_user_settings_path())
    try:
        try:
            with open(settings_path, "r", encoding="utf-8") as f:
                settings = _json.load(f)
        except Exception:
            settings = {}
        settings.setdefault("storage", {})["workspace_dir"] = str(target)
        with open(settings_path, "w", encoding="utf-8") as f:
            _json.dump(settings, f, ensure_ascii=False, indent=2)
    except Exception as e:
        return jsonify({"error": f"设置保存失败: {e}"}), 500

    # Invalidate cache and update live module variable (no restart needed)
    clear_user_settings_cache()
    import web.shared as _shared
    _shared.WORKSPACE_DIR = str(target)

    logger.info("[WorkspaceAssistant] 工作区已切换: %s", target)
    return jsonify({"ok": True, "path": str(target), "name": target.name})


# ─────────────────────────────────────────────────────────────────────────────
# GET /api/v1/workspace/ollama-status
# Check if local Ollama is running and which models are available.
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/ollama-status")
def ollama_status():
    """Return {running: bool, model: str|null, models: [str]}."""
    try:
        import requests as _req

        r = _req.get(
            "http://127.0.0.1:11434/api/tags",
            timeout=3,
            proxies={"http": None, "https": None},  # bypass system proxy for localhost
        )
        data = r.json()
        models = [m["name"] for m in data.get("models", [])]
        if not models:
            return jsonify({"running": True, "model": None, "models": []})

        # 1. Prefer the model configured in user_settings.json
        try:
            from web.shared import PROJECT_ROOT as _PR
            import json as _js, os as _os
            _cfg_path = _os.path.join(_PR, "config", "user_settings.json")
            with open(_cfg_path, "r", encoding="utf-8") as _f:
                _cfg = _js.load(_f)
            _configured = (_cfg.get("local_model") or _cfg.get("ai", {}).get("local_model") or "").strip()
            if _configured and _configured in models:
                return jsonify({"running": True, "model": _configured, "models": models})
        except Exception:
            pass

        # 2. Fall back to size-based preference (include 9b)
        preferred = next(
            (
                m for m in models
                if any(k in m.lower() for k in ("9b", "7b", "8b", "13b", "14b", "32b", "70b"))
            ),
            models[0],
        )
        return jsonify({"running": True, "model": preferred, "models": models})
    except Exception:
        return jsonify({"running": False, "model": None, "models": []})


# ─── Browse local filesystem (lazy) ──────────────────────────────────────────

@workspace_assistant_bp.route("/api/v1/workspace/browse_local")
def browse_local():
    """
    Lazy filesystem browser — no WORKSPACE_DIR restriction.
    No path  → drives + quick-access locations (Windows) / home shortcuts (other).
    With path → all files and folders at that absolute path.
    """
    import sys as _sys

    path = request.args.get("path", "").strip()

    _openable = frozenset(_ALLOWED_EXT)

    if not path:
        # Root level: drives + quick-access locations
        try:
            from web.blueprints.workspace import _list_drives, _quick_access_locations
            return jsonify({
                "is_root": True,
                "drives": _list_drives(),
                "quick_access": _quick_access_locations(),
            })
        except Exception as e:
            # Fallback: just return home dir
            home = Path.home()
            return jsonify({
                "is_root": True,
                "quick_access": [{"name": "主目录", "path": str(home), "type": "quick"}],
                "drives": [],
            })

    target = Path(path).resolve()
    if not target.exists():
        return jsonify({"error": "路径不存在"}), 404
    if not target.is_dir():
        return jsonify({"error": "不是文件夹"}), 400

    _SKIP = {
        "__pycache__", "node_modules", ".git", ".venv", "venv",
        "$RECYCLE.BIN", "System Volume Information",
    }

    entries: list[dict] = []
    try:
        for p in sorted(target.iterdir(), key=lambda x: (x.is_file(), x.name.lower())):
            if p.name.startswith(".") or p.name in _SKIP:
                continue
            if p.is_dir():
                entries.append({"name": p.name, "path": str(p), "type": "folder"})
            elif p.is_file():
                entries.append(
                    _WORKSPACE_TREE.file_entry(
                        p,
                        str(p),
                        _openable,
                    )
                )
    except PermissionError:
        return jsonify({"error": "无访问权限", "entries": []}), 403

    parent = str(target.parent)
    if parent == str(target):
        parent = None  # Drive root (e.g. C:\)

    return jsonify({
        "is_root": False,
        "current": str(target),
        "parent": parent,
        "entries": entries,
    })


# Application config directory — must never be served over the API
_APP_CONFIG_DIR = (Path(__file__).resolve().parents[2] / "config").resolve()


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/open_abs_file
# ─────────────────────────────────────────────────────────────────────────────

@workspace_assistant_bp.route("/api/v1/workspace/open_abs_file", methods=["POST"])
def open_abs_file():
    """
    Parse a file by absolute path — server reads directly from disk, no browser round-trip.
    This is the parsed-file route for browser-opened absolute paths.
    Includes fast ZIP pre-scan to reject PPTX/DOCX containing embedded video before
    any heavy parsing begins.

    Body (JSON): {"path": "/absolute/path/to/file.pptx"}
    Response: same format as open_file
    """
    body = request.get_json(force=True, silent=True) or {}
    abs_path = (body.get("path") or "").strip()
    if not abs_path:
        return jsonify({"error": "缺少 path 字段"}), 400

    file_id = uuid.uuid4().hex
    try:
        from web.shared import WORKSPACE_DIR

        prepared = _OPEN_FILE_BY_PATH.prepare(
            raw_path=str(Path(abs_path).resolve()),
            workspace_dir=WORKSPACE_DIR,
            app_config_dir=_APP_CONFIG_DIR,
            tmp_dir=_ensure_tmp_dir(),
            file_id=file_id,
            allowed_extensions=_ALLOWED_EXT,
            fs_guard=_fs_guard,
            repair_zero_byte_file=_repair_zero_byte_office_file,
            copy_to_tmp=_copy_file_to_tmp_with_validation,
            pre_copy_check=_PPTX_PREFLIGHT.check,
            copy_error_prefix="文件读取失败",
            allow_external_absolute=True,
        )
    except OpenFileInConfigError:
        return jsonify({"error": "不允许访问应用配置目录"}), 403
    except OpenFilePermissionError:
        return jsonify({"error": "路径不合法"}), 403
    except OpenFileNotFoundError:
        return jsonify({"error": "文件不存在"}), 404
    except OpenFileUnsupportedTypeError as e:
        return jsonify({"error": str(e)}), 400
    except PptxPreflightError as e:
        return jsonify({"error": str(e)}), e.status_code
    except OpenFileCopyError as e:
        return jsonify({"error": str(e)}), 500

    target = prepared.target_path
    tmp_path = prepared.tmp_path
    ext = prepared.extension

    try:
        parsed = _FILE_ASSISTANT.parse_editor_file(
            tmp_path,
            file_id=file_id,
            display_name=target.name,
            text_source_path=target,
            source_path=target,
            docx_copy_to_tmp=_copy_file_to_tmp_with_validation,
        )
        file_type = parsed.file_type
        data = parsed.data
        if file_type in ("text", "code") and isinstance(data, dict):
            data["extension"] = ext.lstrip(".")

    except Exception as e:
        logger.error("[open_abs_file] 解析失败 %s: %s", target.name, e, exc_info=True)
        try:
            tmp_path.unlink(missing_ok=True)
        except Exception:
            pass
        return jsonify({"error": f"文件解析失败: {str(e)}"}), 500

    return jsonify(
        _workspace_open_payload(
            file_id=file_id,
            file_name=target.name,
            file_type=file_type,
            data=data,
            source_path=str(target),
        )
    )


@workspace_assistant_bp.route("/api/v1/workspace/docx_full", methods=["POST"])
def load_full_docx():
    """Hydrate a preview-opened DOCX to its full HTML payload."""
    body = request.get_json(force=True, silent=True) or {}
    file_id = (body.get("file_id") or "").strip()
    if not file_id:
        return jsonify({"error": "缺少 file_id 字段"}), 400

    tmp_path = _tmp_file_path(file_id, ".docx")
    if not tmp_path.is_file():
        return jsonify({"error": "DOCX 临时文件不存在或已过期"}), 404

    try:
        data = _FILE_ASSISTANT.load_full_docx(tmp_path, file_id=file_id)
    except Exception as exc:
        logger.error("[load_full_docx] 完整解析失败 %s: %s", file_id, exc, exc_info=True)
        return jsonify({"error": f"DOCX 完整加载失败: {exc}"}), 500

    return jsonify(
        _workspace_open_payload(
            file_id=file_id,
            file_type="docx",
            data=data,
            source_path=str(tmp_path),
        )
    )


_FS_PROTECTED = {
    "windows", "program files", "program files (x86)", "programdata",
    "system volume information", "$recycle.bin",
}


def _fs_guard(p: Path) -> bool:
    """Return True if the path is safe to operate on (not a system root/dir)."""
    parts = {pp.lower() for pp in p.parts}
    if parts & _FS_PROTECTED:
        return False
    # Reject drive roots on Windows (e.g. C:\)
    if p == p.parent:
        return False
    # Protect application config directory (contains JWT secrets, token data, etc.)
    try:
        p.relative_to(_APP_CONFIG_DIR)
        return False
    except ValueError:
        pass
    return True


@workspace_assistant_bp.route("/api/v1/fs/create_file", methods=["POST"])
def fs_create_file():
    """
    Create a file by absolute parent path for the lazy filesystem browser.
    Body (JSON): {"parent": "<absolute folder>", "name": "filename.docx"}
    """
    body = request.get_json(force=True, silent=True) or {}
    parent_raw = (body.get("parent") or "").strip()
    name = (body.get("name") or "").strip()

    try:
        result = _WORKSPACE_FS.create_absolute_file(
            parent_raw=parent_raw,
            name=name,
            allowed_extensions=_ALLOWED_EXT,
            seed_file=_seed_new_file,
            path_guard=_fs_guard,
        )
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code

    logger.info("[Browser] 创建文件: %s", result.path)
    return jsonify({"ok": True, "path": result.path, "name": result.name})


@workspace_assistant_bp.route("/api/v1/fs/create_folder", methods=["POST"])
def fs_create_folder():
    """
    Create a folder by absolute parent path for the lazy filesystem browser.
    Body (JSON): {"parent": "<absolute folder>", "name": "foldername"}
    """
    body = request.get_json(force=True, silent=True) or {}
    parent_raw = (body.get("parent") or "").strip()
    name = (body.get("name") or "").strip()

    try:
        result = _WORKSPACE_FS.create_absolute_folder(
            parent_raw=parent_raw,
            name=name,
            path_guard=_fs_guard,
        )
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code

    logger.info("[Browser] 创建文件夹: %s", result.path)
    return jsonify({"ok": True, "path": result.path, "name": result.name})


@workspace_assistant_bp.route("/api/v1/workspace/fs_delete", methods=["DELETE"])
def fs_delete():
    """
    Delete any file or folder by absolute path from the filesystem browser.
    Query param: path=<abs_path>
    """
    path = request.args.get("path", "").strip()
    try:
        _WORKSPACE_FS.delete_absolute_path(raw_path=path, path_guard=_fs_guard)
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code
    logger.info("[Browser] 删除: %s", path)
    return jsonify({"ok": True})


@workspace_assistant_bp.route("/api/v1/workspace/fs_rename", methods=["PATCH"])
def fs_rename():
    """
    Rename file or folder by absolute path.
    Body (JSON): {"path": "<abs>", "name": "<new_name>"}
    """
    body = request.get_json(silent=True) or {}
    path = (body.get("path") or "").strip()
    new_name = (body.get("name") or "").strip()

    try:
        result = _WORKSPACE_FS.rename_absolute_path(
            raw_path=path,
            new_name=new_name,
            path_guard=_fs_guard,
        )
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code

    logger.info("[Browser] 重命名: %s -> %s", path, result.path)
    return jsonify({"ok": True, "name": result.name, "path": result.path})


@workspace_assistant_bp.route("/api/v1/workspace/fs_copy", methods=["POST"])
def fs_copy():
    """
    Copy or move a file/folder to a destination directory.
    Body (JSON): {"src": "<abs>", "dst_dir": "<abs_dir>", "move": false}
    """
    body = request.get_json(silent=True) or {}
    src = (body.get("src") or "").strip()
    dst_dir = (body.get("dst_dir") or "").strip()
    do_move = bool(body.get("move", False))

    try:
        result = _WORKSPACE_FS.copy_or_move_absolute_path(
            src_raw=src,
            dst_dir_raw=dst_dir,
            move=do_move,
            path_guard=_fs_guard,
        )
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code

    logger.info("[Browser] %s: %s -> %s", "移动" if do_move else "复制", src, result.path)
    return jsonify({"ok": True, "name": result.name, "path": result.path})


@workspace_assistant_bp.route("/api/v1/workspace/upload-to-folder", methods=["POST"])
def upload_to_folder():
    """
    Receive one or more uploaded files and write them into a filesystem-browser folder.
    Form fields: dest_dir=<absolute folder>, file=<one or more uploads>
    """
    dest_dir = (request.form.get("dest_dir") or "").strip()
    uploaded_files = request.files.getlist("file")
    try:
        saved = _WORKSPACE_FS.upload_to_absolute_folder(
            dest_dir_raw=dest_dir,
            uploaded_files=uploaded_files,
            path_guard=_fs_guard,
        )
    except WorkspaceFsError as exc:
        return jsonify({"error": str(exc)}), exc.status_code

    logger.info("[Browser] upload-to-folder: %s files -> %s", len(saved), dest_dir)
    return jsonify(
        {
            "ok": True,
            "saved": [{"name": item.name, "path": item.path} for item in saved],
        }
    )


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/patch_file
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/patch_file", methods=["POST"])
def patch_file():
    """
    Apply text-replacement proposals to an existing file and return the patched
    file for browser download.  Used by the multi-document content-sync feature
    so AI-generated proposals can be exported without modifying disk.

    Body (JSON):
      {
        "path":      "/absolute/or/workspace-relative/path/to/file.docx",
        "proposals": [
          {"original_text": "...", "proposed_text": "..."},
          ...
        ]
      }

    Supported formats: .docx, .txt, .md
    Returns: patched file binary (Content-Disposition: attachment)
    """
    import io as _io

    body = request.get_json(force=True, silent=True) or {}
    raw_path = (body.get("path") or "").strip()
    proposals = body.get("proposals") or []

    if not raw_path:
        return jsonify({"error": "缺少 path 字段"}), 400
    if not proposals:
        return jsonify({"error": "缺少 proposals 字段"}), 400

    # Resolve to an absolute path; try workspace-relative first, then absolute
    target: "Path | None" = None
    candidate = Path(raw_path)
    if candidate.is_absolute() and candidate.is_file():
        target = candidate.resolve()
    else:
        try:
            from web.shared import WORKSPACE_DIR
            ws_root = Path(WORKSPACE_DIR).resolve()
            rel_try = ws_root.joinpath(raw_path).resolve()
            if rel_try.is_file():
                target = rel_try
        except Exception:
            pass

    if target is None:
        return jsonify({"error": "文件不存在或路径无效"}), 404

    ext = target.suffix.lower()
    file_name = target.name

    if ext not in (".docx", ".txt", ".md"):
        return jsonify({"error": f"暂不支持对 {ext} 格式进行文本修补，请使用 .docx / .txt / .md"}), 400

    # Only include well-formed proposals
    clean = [
        p for p in proposals
        if isinstance(p, dict)
        and (p.get("original_text") or "").strip()
        and (p.get("proposed_text") or "").strip()
    ]
    if not clean:
        return jsonify({"error": "proposals 中没有有效的修改条目"}), 400

    # ── DOCX: python-docx paragraph / table cell replacement ─────────────────
    if ext == ".docx":
        try:
            from docx import Document

            doc = Document(str(target))

            def _replace_in_para(para, orig: str, new: str) -> bool:
                """Replace *orig* with *new* inside *para*, preserving run structure."""
                full = para.text
                if orig not in full:
                    return False
                # Fast path: orig lives entirely inside a single run
                for run in para.runs:
                    if orig in run.text:
                        run.text = run.text.replace(orig, new)
                        return True
                # Slow path: orig spans multiple runs — rebuild first run, clear rest
                new_full = full.replace(orig, new, 1)
                if para.runs:
                    para.runs[0].text = new_full
                    for run in para.runs[1:]:
                        run.text = ""
                return True

            def _replace_all(orig: str, new: str):
                for para in doc.paragraphs:
                    _replace_in_para(para, orig, new)
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                _replace_in_para(para, orig, new)

            for p in clean:
                orig = p["original_text"].strip()
                new  = p["proposed_text"].strip()
                _replace_all(orig, new)

            buf = _io.BytesIO()
            doc.save(buf)
            buf.seek(0)
            return send_file(
                buf,
                mimetype="application/"
                         "vnd.openxmlformats-officedocument.wordprocessingml.document",
                as_attachment=True,
                download_name=f"修改后_{file_name}",
            )
        except ImportError:
            return jsonify({"error": "python-docx 未安装，请执行: pip install python-docx"}), 500
        except Exception as exc:
            logger.error("[patch_file] DOCX 修补失败: %s", exc, exc_info=True)
            return jsonify({"error": f"DOCX 修补失败: {str(exc)}"}), 500

    # ── TXT / MD: plain string replacement ───────────────────────────────────
    try:
        content = target.read_text(encoding="utf-8", errors="replace")
        for p in clean:
            content = content.replace(p["original_text"].strip(), p["proposed_text"].strip(), 1)
        buf = _io.BytesIO(content.encode("utf-8"))
        mime = "text/markdown; charset=utf-8" if ext == ".md" else "text/plain; charset=utf-8"
        return send_file(buf, mimetype=mime, as_attachment=True, download_name=f"修改后_{file_name}")
    except Exception as exc:
        logger.error("[patch_file] TXT/MD 修补失败: %s", exc, exc_info=True)
        return jsonify({"error": f"文本修补失败: {str(exc)}"}), 500


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/audio_overview
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/audio_overview", methods=["POST"])
def audio_overview():
    """
    Generate a two-host podcast audio overview from a set of files.

    Body JSON:
      { "files": [{"name": "...", "content": "..."}], "session_id": "..." }

    SSE stream:
      {"event": "script",    "data": [{speaker, text}, ...]}
      {"event": "progress",  "data": "合成音频…"}
      {"event": "audio_url", "data": "/static/audio_cache/podcast_xxx.mp3"}
      {"event": "error",     "data": "…"}
    """
    import asyncio
    import uuid as _uuid

    body = request.get_json(force=True, silent=True) or {}
    files = body.get("files") or []
    if not files:
        return jsonify({"error": "缺少 files 字段"}), 400

    combined_text = "\n\n".join(
        f"=== {f.get('name', '文件')} ===\n{f.get('content', '')}"
        for f in files
    )[:20000]

    session_id = body.get("session_id") or _uuid.uuid4().hex[:12]

    def _generate():
        import json as _json

        try:
            from web.audio_overview import AudioOverviewGenerator
            from web.runtime_context import get_client, get_model_map

            client = get_client()

            # Pick a suitable model
            _model = get_model_map().get("CHAT") or "gemini-2.5-flash-lite"
            if _model.startswith("deep-research"):
                _model = "gemini-2.5-flash-lite"

            # Build a simple wrapper that AudioOverviewGenerator expects
            class _ModelAdapter:
                def generate_content(self, prompt):
                    resp = client.models.generate_content(
                        model=_model, contents=prompt
                    )
                    return resp

            gen = AudioOverviewGenerator(
                output_dir=os.path.join(
                    os.path.dirname(os.path.dirname(__file__)), "static", "audio_cache"
                )
            )

            loop = asyncio.new_event_loop()
            script = loop.run_until_complete(
                gen.generate_script(combined_text, _ModelAdapter())
            )
            if not script:
                yield f"data: {_json.dumps({'event': 'error', 'data': '脚本生成失败，请重试'}, ensure_ascii=False)}\n\n"
                return

            yield f"data: {_json.dumps({'event': 'script', 'data': script}, ensure_ascii=False)}\n\n"

            # Attempt TTS synthesis
            try:
                audio_path = loop.run_until_complete(
                    gen.synthesize_audio(script, session_id)
                )
                loop.close()
                if audio_path and os.path.exists(audio_path):
                    audio_url = "/static/audio_cache/" + os.path.basename(audio_path)
                    yield f"data: {_json.dumps({'event': 'audio_url', 'data': audio_url}, ensure_ascii=False)}\n\n"
                else:
                    yield f"data: {_json.dumps({'event': 'audio_url', 'data': None}, ensure_ascii=False)}\n\n"
            except Exception as tts_err:
                loop.close()
                logger.warning("[audio_overview] TTS 失败: %s", tts_err)
                yield f"data: {_json.dumps({'event': 'audio_url', 'data': None}, ensure_ascii=False)}\n\n"

        except Exception as exc:
            logger.error("[audio_overview] 失败: %s", exc, exc_info=True)
            yield f"data: {_json.dumps({'event': 'error', 'data': str(exc)}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(_generate()),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/notebook_guide
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/notebook_guide", methods=["POST"])
def notebook_guide():
    """
    Generate a 4-section study guide (学习包) from attached files.

    Body JSON:
      { "files": [{"name": "...", "content": "..."}] }

    SSE stream (one event per section):
      {"section": "summary",  "content": "..."}
      {"section": "points",   "content": "..."}
      {"section": "faq",      "content": "..."}
      {"section": "glossary", "content": "..."}
      {"section": "done"}
      {"section": "error",    "content": "..."}
    """
    body = request.get_json(force=True, silent=True) or {}
    files = body.get("files") or []
    if not files:
        return jsonify({"error": "缺少 files 字段"}), 400

    combined_text = "\n\n".join(
        f"=== {f.get('name', '文件')} ===\n{f.get('content', '')}"
        for f in files
    )[:24000]

    SECTIONS = [
        (
            "summary",
            "执行摘要",
            "请用200-300字对以下资料进行执行摘要，抓住核心结论和关键数据，不要逐条列点。",
        ),
        (
            "points",
            "关键要点",
            "请从以下资料中提炼5-8条关键要点，每条以「·」开头，包含具体数据或结论，不要泛泛而谈。",
        ),
        (
            "faq",
            "常见问答",
            "请根据以下资料生成5个读者最可能提出的问题及详细解答，格式：Q: 问题\nA: 解答",
        ),
        (
            "glossary",
            "核心词汇",
            "请从以下资料中提取8-12个专业术语或核心概念，每个词汇后附一句简洁定义，格式：**词汇** — 定义",
        ),
    ]

    def _generate():
        import json as _json

        try:
            from web.runtime_context import get_client, get_model_map

            client = get_client()
            _model = get_model_map().get("CHAT") or "gemini-2.5-flash-lite"
            if _model.startswith("deep-research"):
                _model = "gemini-2.5-flash-lite"

            for sec_key, sec_label, sec_prompt in SECTIONS:
                full_prompt = (
                    f"{sec_prompt}\n\n"
                    f"资料内容（共 {len(files)} 个文件）:\n{combined_text}"
                )
                try:
                    resp = client.models.generate_content(
                        model=_model, contents=full_prompt
                    )
                    content = (getattr(resp, "text", None) or "").strip()
                    if not content:
                        content = "（AI 暂无回复）"
                except Exception as sec_err:
                    content = f"（生成失败: {sec_err}）"

                yield f"data: {_json.dumps({'section': sec_key, 'label': sec_label, 'content': content}, ensure_ascii=False)}\n\n"

            yield f"data: {_json.dumps({'section': 'done'}, ensure_ascii=False)}\n\n"

        except Exception as exc:
            logger.error("[notebook_guide] 失败: %s", exc, exc_info=True)
            yield f"data: {_json.dumps({'section': 'error', 'content': str(exc)}, ensure_ascii=False)}\n\n"

    return Response(
        stream_with_context(_generate()),
        mimetype="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
        },
    )


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/pdf/save_annotations
# POST /api/v1/workspace/pdf/load_annotations/<file_id>
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/pdf/save_annotations", methods=["POST"])
def pdf_save_annotations():
    """
    Embed client-side annotations into the PDF binary and return the modified file
    for download.

    Body (JSON):
      {
        "file_id": str,          — tmp file identifier (hex)
        "annotations": [...],    — array of annotation objects from KotoPdfViewer
        "filename": str          — original filename (used for download name)
      }

    Returns the annotated PDF as an attachment.
    """
    body = request.get_json(silent=True) or {}
    file_id = str(body.get("file_id", ""))
    annotations = body.get("annotations", [])
    filename = body.get("filename", "annotated.pdf")

    # Sanitize file_id: only allow alphanumeric
    if not file_id.isalnum():
        return jsonify({"error": "无效的 file_id"}), 400

    tmp_dir = _ensure_tmp_dir()
    matches = list(tmp_dir.glob(f"{file_id}.pdf"))
    if not matches:
        # Also try without extension
        matches = list(tmp_dir.glob(f"{file_id}.*"))
        matches = [m for m in matches if m.suffix.lower() == ".pdf"]
    if not matches:
        return jsonify({"error": "找不到 PDF 文件"}), 404

    pdf_path = str(matches[0])

    try:
        from web.pdf_annotator import embed_annotations
        pdf_bytes = embed_annotations(pdf_path, annotations)
    except Exception as exc:
        logger.error("[pdf_save_annotations] 注释嵌入失败: %s", exc, exc_info=True)
        return jsonify({"error": f"注释嵌入失败: {str(exc)}"}), 500

    import io as _io
    buf = _io.BytesIO(pdf_bytes)
    safe_name = filename if filename.lower().endswith(".pdf") else filename + ".pdf"
    return send_file(
        buf,
        mimetype="application/pdf",
        as_attachment=True,
        download_name=safe_name,
    )


@workspace_assistant_bp.route(
    "/api/v1/workspace/pdf/load_annotations/<file_id>", methods=["GET"]
)
def pdf_load_annotations(file_id: str):
    """
    Read annotations embedded in a cached PDF and return them as JSON.

    Returns:
      {"annotations": [...]}
    """
    if not file_id.isalnum():
        return jsonify({"error": "无效的 file_id"}), 400

    tmp_dir = _ensure_tmp_dir()
    matches = [m for m in tmp_dir.glob(f"{file_id}.*") if m.suffix.lower() == ".pdf"]
    if not matches:
        return jsonify({"error": "找不到 PDF 文件"}), 404

    pdf_path = str(matches[0])

    try:
        from web.pdf_annotator import read_annotations
        annotations = read_annotations(pdf_path)
    except Exception as exc:
        logger.error("[pdf_load_annotations] 批注读取失败: %s", exc, exc_info=True)
        return jsonify({"error": f"批注读取失败: {str(exc)}"}), 500

    return jsonify({"annotations": annotations})


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/pdf/page_ops
# ─────────────────────────────────────────────────────────────────────────────


@workspace_assistant_bp.route("/api/v1/workspace/pdf/page_ops", methods=["POST"])
def pdf_page_ops():
    """
    Reconstruct a PDF with pages in the requested order / rotation.
    Used by the Page Manager for reorder, rotate, delete, and split (export subset).

    Body (JSON):
      {
        "file_id": str,   — tmp file identifier (hex)
        "pages": [{"orig_page": int, "rotation": int}, ...]
      }

    Returns the new PDF as an attachment.
    """
    body = request.get_json(silent=True) or {}
    file_id = str(body.get("file_id", ""))
    pages = body.get("pages", [])

    if not file_id.isalnum():
        return jsonify({"error": "无效的 file_id"}), 400
    if not pages:
        return jsonify({"error": "pages 不能为空"}), 400

    tmp_dir = _ensure_tmp_dir()
    matches = [m for m in tmp_dir.glob(f"{file_id}.*") if m.suffix.lower() == ".pdf"]
    if not matches:
        return jsonify({"error": "找不到 PDF 文件"}), 404

    pdf_path = str(matches[0])
    orig_name = matches[0].name

    try:
        from web.pdf_annotator import apply_page_ops
        pdf_bytes = apply_page_ops(pdf_path, pages)
    except Exception as exc:
        logger.error("[pdf_page_ops] 页面操作失败: %s", exc, exc_info=True)
        return jsonify({"error": f"页面操作失败: {str(exc)}"}), 500

    import io as _io
    buf = _io.BytesIO(pdf_bytes)
    return send_file(
        buf,
        mimetype="application/pdf",
        as_attachment=True,
        download_name=orig_name,
    )


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/pdf/convert
# ─────────────────────────────────────────────────────────────────────────────

_PDF_CONVERT_MIME = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


@workspace_assistant_bp.route("/api/v1/workspace/pdf/convert", methods=["POST"])
def pdf_convert():
    """
    Convert a cached PDF to another format.

    Body (JSON):
      {
        "file_id"      : str,   — tmp file identifier (hex)
        "target_format": str,   — "docx" | "xlsx" | "pptx"
        "filename"     : str    — original filename (used for download name, optional)
      }

    Returns the converted file as an attachment.
    """
    body = request.get_json(silent=True) or {}
    file_id = str(body.get("file_id", ""))
    target_fmt = str(body.get("target_format", "")).lower().lstrip(".")
    filename = body.get("filename", "converted")

    if not file_id.isalnum():
        return jsonify({"error": "无效的 file_id"}), 400
    if target_fmt not in _PDF_CONVERT_MIME:
        return jsonify({"error": f"不支持的目标格式 '{target_fmt}'，支持：docx / xlsx / pptx"}), 400

    tmp_dir = _ensure_tmp_dir()
    matches = [m for m in tmp_dir.glob(f"{file_id}.*") if m.suffix.lower() == ".pdf"]
    if not matches:
        return jsonify({"error": "找不到 PDF 文件"}), 404

    pdf_path = str(matches[0])
    stem = Path(filename).stem or "converted"
    download_name = f"{stem}.{target_fmt}"
    warning = ""

    try:
        from web.pdf_annotator import pdf_to_docx, pdf_to_xlsx, pdf_to_pptx
        if target_fmt == "docx":
            file_bytes, warning = pdf_to_docx(pdf_path)
        elif target_fmt == "xlsx":
            file_bytes = pdf_to_xlsx(pdf_path)
        elif target_fmt == "pptx":
            file_bytes = pdf_to_pptx(pdf_path)
        else:
            return jsonify({"error": "内部错误"}), 500
    except Exception as exc:
        logger.error("[pdf_convert] 格式转换失败: %s", exc, exc_info=True)
        return jsonify({"error": f"格式转换失败: {str(exc)}"}), 500

    import io as _io
    resp = send_file(
        _io.BytesIO(file_bytes),
        mimetype=_PDF_CONVERT_MIME[target_fmt],
        as_attachment=True,
        download_name=download_name,
    )
    if warning:
        resp.headers["X-Koto-Warning"] = warning
    return resp


# ─────────────────────────────────────────────────────────────────────────────
# POST /api/v1/workspace/pdf/remove_watermark
# ─────────────────────────────────────────────────────────────────────────────

@workspace_assistant_bp.route("/api/v1/workspace/pdf/remove_watermark", methods=["POST"])
def pdf_remove_watermark():
    """
    AI-assisted watermark removal.
    Body: {"file_id": str, "use_ai": bool}
    Returns the cleaned PDF as an attachment.
    """
    body = request.get_json(silent=True) or {}
    file_id = str(body.get("file_id", ""))
    use_ai  = bool(body.get("use_ai", True))

    # Validate file_id: must be alphanumeric/hyphen/underscore only (prevents path traversal)
    if not file_id or not file_id.replace("-", "").replace("_", "").isalnum():
        return jsonify({"error": "无效的 file_id"}), 400

    tmp_dir = _ensure_tmp_dir()
    matches = [m for m in tmp_dir.glob(f"{file_id}.*") if m.suffix.lower() == ".pdf"]
    if not matches:
        return jsonify({"error": "找不到对应的 PDF 文件，请重新打开"}), 404

    pdf_path = str(matches[0])

    api_key = None
    if use_ai:
        import os
        api_key = os.getenv("GEMINI_API_KEY") or os.getenv("API_KEY")

    try:
        from web.pdf_annotator import remove_watermark
        pdf_bytes, removed_count, method_used = remove_watermark(
            pdf_path, use_ai=use_ai, api_key=api_key
        )
    except Exception as exc:
        logger.error("[pdf_remove_watermark] 失败: %s", exc, exc_info=True)
        return jsonify({"error": f"去水印失败: {str(exc)}"}), 500

    import io as _io
    orig_stem = matches[0].stem
    resp = send_file(
        _io.BytesIO(pdf_bytes),
        mimetype="application/pdf",
        as_attachment=True,
        download_name=f"{orig_stem}_去水印.pdf",
    )
    resp.headers["X-Koto-Removed-Count"] = str(removed_count)
    resp.headers["X-Koto-Method"] = method_used
    return resp
