"""
pdf_annotator.py — Embed / read PDF annotations, page ops, and format conversion.

Annotation types supported (matching KotoPdfViewer JS):
  highlight, underline, strikethrough, note (free-text), draw (ink)

All coordinates are expressed as **PDF points** from the bottom-left origin
(standard PDF coordinate system).  The frontend converts CSS pixels to PDF
points before calling save_annotations.

Format conversion:
  pdf_to_docx  — via doc_converter (LibreOffice → pypdf text fallback)
  pdf_to_xlsx  — via pdfplumber table extraction → openpyxl
  pdf_to_pptx  — via PyMuPDF page images → python-pptx slides
"""

from __future__ import annotations

import logging
from io import BytesIO
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────

def embed_annotations(pdf_path: str, annotations: list[dict]) -> bytes:
    """
    Read *pdf_path*, embed *annotations*, and return the modified PDF bytes.

    Each annotation dict must have:
        page   : int   — 1-based page number
        type   : str   — "highlight" | "underline" | "strikethrough" | "note" | "draw"
        color  : str   — hex color, e.g. "#FFFF00"
        rects  : list  — list of [x1, y1, x2, y2] in PDF points (for text annots)
        content: str   — optional text (used for notes)
        inkList: list  — for draw: list of point arrays [[x,y], ...]

    Returns the bytes of the annotated PDF.
    """
    try:
        from pypdf import PdfReader, PdfWriter  # type: ignore
        from pypdf.generic import (  # type: ignore
            ArrayObject, DictionaryObject, FloatObject,
            NameObject, NumberObject, TextStringObject,
        )
    except ImportError as exc:
        raise RuntimeError("pypdf is required for annotation embedding") from exc

    with open(pdf_path, "rb") as fh:
        reader = PdfReader(fh)
        writer = PdfWriter()
        writer.clone_reader_document_root(reader)

        for annot in annotations:
            page_num = int(annot.get("page", 1))
            if page_num < 1 or page_num > len(writer.pages):
                logger.warning(f"[PdfAnnotator] 无效页码 {page_num}，跳过")
                continue
            page = writer.pages[page_num - 1]
            annot_obj = _build_annot_object(annot, page, writer)
            if annot_obj is None:
                continue
            if "/Annots" not in page:
                page[NameObject("/Annots")] = ArrayObject()
            page["/Annots"].append(writer._add_object(annot_obj))  # type: ignore[attr-defined]

    out = BytesIO()
    writer.write(out)
    return out.getvalue()


def read_annotations(pdf_path: str) -> list[dict]:
    """
    Read annotations embedded in *pdf_path* and return them as a list of dicts
    compatible with the KotoPdfViewer annotation format.
    """
    result: list[dict] = []
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        logger.warning("[PdfAnnotator] pypdf not installed — cannot read annotations")
        return result

    try:
        with open(pdf_path, "rb") as fh:
            reader = PdfReader(fh)
            for page_num, page in enumerate(reader.pages, start=1):
                annots = page.get("/Annots")
                if not annots:
                    continue
                for annot_ref in annots:
                    try:
                        annot = annot_ref.get_object() if hasattr(annot_ref, "get_object") else annot_ref
                        parsed = _parse_annot_object(annot, page_num)
                        if parsed:
                            result.append(parsed)
                    except Exception as e:
                        logger.debug(f"[PdfAnnotator] 解析批注失败: {e}")
    except Exception as e:
        logger.error(f"[PdfAnnotator] 读取 PDF 失败: {e}")
    return result


def apply_page_ops(pdf_path: str, pages: list[dict]) -> bytes:
    """
    Reconstruct a PDF from *pdf_path* using the given page specification.

    *pages* is an ordered list of dicts (from the Page Manager):
        [{"orig_page": int, "rotation": int}, ...]

    - Pages are included in the requested order.
    - Pages not listed are excluded (delete support).
    - Rotation of 90 / 180 / 270 is applied to each page.

    Returns the bytes of the new PDF.
    """
    try:
        from pypdf import PdfReader, PdfWriter  # type: ignore
    except ImportError as exc:
        raise RuntimeError("pypdf is required for page operations") from exc

    with open(pdf_path, "rb") as fh:
        reader = PdfReader(fh)
        total = len(reader.pages)
        writer = PdfWriter()

        for spec in pages:
            orig = int(spec.get("orig_page", 0))
            rotation = int(spec.get("rotation", 0))
            if orig < 1 or orig > total:
                logger.warning(f"[PdfAnnotator] page_ops: 无效页码 {orig}，跳过")
                continue
            page = reader.pages[orig - 1]
            if rotation:
                # pypdf page.rotate() takes clockwise degrees (90, 180, 270)
                page = page.rotate(rotation % 360)
            writer.add_page(page)

    out = BytesIO()
    writer.write(out)
    return out.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# Format conversion
# ─────────────────────────────────────────────────────────────────────────────

def pdf_to_docx(pdf_path: str) -> tuple[bytes, str]:
    """
    Convert PDF to an editable DOCX while retaining page boundaries, basic text
    styling, alignment, and embedded images. Scanned pages fall back to a page
    image instead of failing with an empty document.
    """
    try:
        import pymupdf as fitz  # type: ignore
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required for PDF→DOCX conversion") from exc
    try:
        from docx import Document  # type: ignore
        from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore
        from docx.shared import Inches, Pt, RGBColor  # type: ignore
    except ImportError as exc:
        raise RuntimeError("python-docx is required for PDF→DOCX conversion") from exc

    document = Document()
    section = document.sections[0]
    section.top_margin = Pt(36)
    section.bottom_margin = Pt(36)
    section.left_margin = Pt(36)
    section.right_margin = Pt(36)
    image_fallback_pages = 0

    with fitz.open(pdf_path) as pdf:
        if len(pdf):
            first_rect = pdf[0].rect
            section.page_width = Pt(float(first_rect.width))
            section.page_height = Pt(float(first_rect.height))

        for page_index, page in enumerate(pdf):
            blocks = sorted(
                page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE).get("blocks", []),
                key=lambda block: (round(float(block.get("bbox", (0, 0, 0, 0))[1]), 1),
                                   float(block.get("bbox", (0, 0, 0, 0))[0])),
            )
            emitted = False
            for block in blocks:
                block_type = block.get("type")
                if block_type == 0:
                    for line in block.get("lines", []):
                        spans = line.get("spans", [])
                        text = "".join(str(span.get("text", "")) for span in spans)
                        if not text.strip():
                            continue
                        paragraph = document.add_paragraph()
                        bbox = line.get("bbox", block.get("bbox", (0, 0, 0, 0)))
                        page_width = max(float(page.rect.width), 1.0)
                        left_gap = float(bbox[0])
                        right_gap = page_width - float(bbox[2])
                        if abs(left_gap - right_gap) < page_width * 0.06:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        elif left_gap > page_width * 0.55:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                        else:
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                        for span in spans:
                            value = str(span.get("text", ""))
                            if not value:
                                continue
                            run = paragraph.add_run(value)
                            font_name = str(span.get("font", ""))
                            run.bold = "bold" in font_name.lower()
                            run.italic = any(token in font_name.lower() for token in ("italic", "oblique"))
                            size = float(span.get("size", 0) or 0)
                            if size:
                                run.font.size = Pt(max(1.0, min(size, 200.0)))
                            if font_name:
                                run.font.name = font_name
                            color = int(span.get("color", 0) or 0)
                            run.font.color.rgb = RGBColor(
                                (color >> 16) & 0xFF,
                                (color >> 8) & 0xFF,
                                color & 0xFF,
                            )
                        emitted = True
                elif block_type == 1 and block.get("image"):
                    try:
                        available_width = max(float(section.page_width - section.left_margin - section.right_margin), 1)
                        document.add_picture(
                            BytesIO(block["image"]),
                            width=Inches(min(available_width / 914400.0, 7.5)),
                        )
                        emitted = True
                    except Exception:
                        logger.debug("Could not embed PDF image block", exc_info=True)

            if not emitted:
                pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                available_width = max(float(section.page_width - section.left_margin - section.right_margin), 1)
                document.add_picture(
                    BytesIO(pix.tobytes("png")),
                    width=Inches(min(available_width / 914400.0, 7.5)),
                )
                image_fallback_pages += 1
            if page_index < len(pdf) - 1:
                document.add_page_break()

    output = BytesIO()
    document.save(output)
    warning = ""
    if image_fallback_pages:
        warning = f"{image_fallback_pages} 页未检测到可编辑文字，已按页面图像保留。"
    return output.getvalue(), warning


def pdf_to_xlsx(pdf_path: str) -> bytes:
    """
    Extract tables from PDF pages and write to XLSX.
    Falls back to plain text rows when no tables are detected.
    Requires pdfplumber and openpyxl.
    """
    try:
        import pdfplumber  # type: ignore
    except ImportError as exc:
        raise RuntimeError("pdfplumber is required for PDF→XLSX conversion") from exc
    try:
        from openpyxl import Workbook  # type: ignore
    except ImportError as exc:
        raise RuntimeError("openpyxl is required for PDF→XLSX conversion") from exc

    from openpyxl.styles import Alignment, Font, PatternFill  # type: ignore
    from openpyxl.utils import get_column_letter  # type: ignore

    wb = Workbook()
    wb.remove(wb.active)

    with pdfplumber.open(pdf_path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            ws = wb.create_sheet(title=f"第{page_num}页")
            current_row = 1
            tables = page.extract_tables() or []
            if tables:
                for table_index, table in enumerate(tables, start=1):
                    if len(tables) > 1:
                        ws.cell(current_row, 1, f"表格 {table_index}").font = Font(bold=True)
                        current_row += 1
                    for row_index, row in enumerate(table):
                        for col_index, value in enumerate(row or [], start=1):
                            cell = ws.cell(current_row, col_index, "" if value is None else str(value))
                            cell.alignment = Alignment(
                                horizontal="center" if row_index == 0 else "left",
                                vertical="top",
                                wrap_text=True,
                            )
                            if row_index == 0:
                                cell.font = Font(bold=True)
                                cell.fill = PatternFill("solid", fgColor="E8EEF7")
                        current_row += 1
                    current_row += 1
            else:
                text = page.extract_text() or ""
                if text.strip():
                    for line_number, line in enumerate(text.splitlines(), start=1):
                        ws.cell(current_row, 1, line_number)
                        content_cell = ws.cell(current_row, 2, line)
                        content_cell.alignment = Alignment(horizontal="left", vertical="top", wrap_text=True)
                        current_row += 1
                    ws.column_dimensions["A"].width = 8

            for column in range(1, ws.max_column + 1):
                max_length = max(
                    (len(str(ws.cell(row, column).value or "")) for row in range(1, ws.max_row + 1)),
                    default=0,
                )
                ws.column_dimensions[get_column_letter(column)].width = min(max(max_length + 2, 8), 60)
            ws.freeze_panes = "A2" if ws.max_row > 1 else None

    if not wb.worksheets:
        wb.create_sheet("内容")

    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def pdf_to_pptx(pdf_path: str) -> bytes:
    """
    Convert each PDF page to a full-page image slide in PPTX.
    Requires PyMuPDF (fitz) and python-pptx.
    """
    try:
        import pymupdf as fitz  # type: ignore
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required for PDF→PPTX conversion") from exc
    try:
        from pptx import Presentation  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        from pptx.util import Inches  # type: ignore
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PDF→PPTX conversion") from exc

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    with fitz.open(pdf_path) as doc:
        if len(doc):
            first = doc[0].rect
            ratio = max(float(first.width), 1.0) / max(float(first.height), 1.0)
            if ratio >= 1:
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(13.333 / ratio)
            else:
                prs.slide_height = Inches(13.333)
                prs.slide_width = Inches(13.333 * ratio)

        for page_obj in doc:
            pix = page_obj.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
            img_bytes = pix.tobytes("png")
            slide = prs.slides.add_slide(blank_layout)
            page_ratio = max(float(page_obj.rect.width), 1.0) / max(float(page_obj.rect.height), 1.0)
            slide_ratio = float(prs.slide_width) / max(float(prs.slide_height), 1.0)
            if page_ratio >= slide_ratio:
                width = prs.slide_width
                height = int(width / page_ratio)
                left, top = 0, int((prs.slide_height - height) / 2)
            else:
                height = prs.slide_height
                width = int(height * page_ratio)
                left, top = int((prs.slide_width - width) / 2), 0
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
            slide.shapes.add_picture(BytesIO(img_bytes), left, top, width, height)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def remove_watermark(pdf_path: str, use_ai: bool = True,
                     api_key: str | None = None) -> tuple[bytes, int, str]:
    """
    Attempt to remove watermarks from a PDF.

    Deterministic local strategies:
      1. remove PDF watermark annotations;
      2. redact repeated, light, large, diagonal, or keyword-matched text;
      3. remove repeated central image XObjects used as image watermarks.

    ``use_ai`` and ``api_key`` remain as compatibility parameters only. The
    implementation intentionally does not send document content to a cloud API.

    Returns (pdf_bytes, n_regions_removed, method_used).
    """
    try:
        import pymupdf as fitz  # type: ignore
    except ImportError as exc:
        raise RuntimeError("PyMuPDF is required for watermark removal") from exc

    import re

    doc = fitz.open(pdf_path)
    page_count = len(doc)
    removed = 0

    if page_count == 0:
        data = doc.tobytes()
        doc.close()
        return data, 0, "none"

    WATERMARK_KEYWORDS = [
        "confidential", "机密", "draft", "草稿", "internal", "内部",
        "sample", "watermark", "水印", "copyright", "版权",
        "do not copy", "禁止复制", "仅供内部", "for internal use",
    ]
    def _normalized_text(value: str) -> str:
        return re.sub(r"\s+", " ", value).strip().casefold()

    # ── Strategy 1: explicit watermark annotations ───────────────────────────
    for page in doc:
        annotations = list(page.annots() or [])
        for annotation in annotations:
            type_name = str((annotation.type or (None, ""))[1] or "").casefold()
            content = _normalized_text((annotation.info or {}).get("content", ""))
            if "watermark" in type_name or "水印" in content:
                page.delete_annot(annotation)
                removed += 1

    # ── Strategy 2: structural text detection ────────────────────────────────
    sample_pages = min(20, page_count)
    candidate_map: dict[str, dict[str, Any]] = {}

    for pg_idx in range(sample_pages):
        page = doc[pg_idx]
        try:
            blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        except Exception:
            continue
        for block in blocks:
            if block.get("type") != 0:
                continue
            for line in block.get("lines", []):
                direction = line.get("dir", (1.0, 0.0))
                is_diagonal = abs(float(direction[1] or 0.0)) > 0.15
                for span in line.get("spans", []):
                    txt = span.get("text", "").strip()
                    if not txt:
                        continue
                    # Light colour ≈ watermark
                    c = span.get("color", 0)
                    r_ch = (c >> 16) & 0xFF
                    g_ch = (c >> 8) & 0xFF
                    b_ch = c & 0xFF
                    alpha = int(span.get("alpha", 255) or 255)
                    font_size = float(span.get("size", 0) or 0)
                    is_light = r_ch > 170 and g_ch > 170 and b_ch > 170
                    txt_lower = _normalized_text(txt)
                    is_keyword = any(kw in txt_lower for kw in WATERMARK_KEYWORDS)
                    looks_like_watermark = (
                        is_keyword
                        or (is_light and (font_size >= 14 or alpha < 220))
                        or (is_diagonal and font_size >= 18)
                    )
                    if looks_like_watermark:
                        key = txt_lower[:120]
                        entry = candidate_map.setdefault(
                            key, {"text": txt, "pages": set(), "keyword": False}
                        )
                        entry["pages"].add(pg_idx)
                        entry["keyword"] = bool(entry["keyword"] or is_keyword)

    threshold = 1 if sample_pages == 1 else max(2, (sample_pages + 1) // 2)
    watermark_texts = {
        value["text"]
        for value in candidate_map.values()
        if len(value["pages"]) >= threshold
        or (sample_pages == 1 and value["keyword"])
    }

    if watermark_texts:
        for pg_idx in range(page_count):
            page = doc[pg_idx]
            areas: list = []
            for wt in watermark_texts:
                areas.extend(page.search_for(wt, quads=False))
            for area in areas:
                page.add_redact_annot(area, fill=(1, 1, 1))
            if areas:
                try:
                    page.apply_redactions(images=0, graphics=0)
                except TypeError:
                    page.apply_redactions()
                removed += len(areas)

    # ── Strategy 3: repeated central image XObjects ──────────────────────────
    image_candidates: dict[int, dict[str, Any]] = {}
    for pg_idx in range(sample_pages):
        page = doc[pg_idx]
        page_area = max(float(page.rect.width * page.rect.height), 1.0)
        for info in page.get_image_info(xrefs=True):
            xref = int(info.get("xref", 0) or 0)
            bbox = fitz.Rect(info.get("bbox", (0, 0, 0, 0)))
            ratio = float(bbox.get_area()) / page_area
            center = (bbox.x0 + bbox.x1) / 2.0, (bbox.y0 + bbox.y1) / 2.0
            is_central = (
                page.rect.width * 0.2 <= center[0] <= page.rect.width * 0.8
                and page.rect.height * 0.2 <= center[1] <= page.rect.height * 0.8
            )
            if xref > 0 and is_central and 0.05 <= ratio <= 0.8:
                entry = image_candidates.setdefault(xref, {"pages": set(), "count": 0})
                entry["pages"].add(pg_idx)
                entry["count"] += 1

    image_threshold = 1 if sample_pages == 1 else max(2, (sample_pages + 1) // 2)
    for xref, candidate in image_candidates.items():
        # For a single-page document, only remove an image with a soft mask;
        # otherwise a central photo would be indistinguishable from a watermark.
        has_soft_mask = any(
            image[0] == xref and int(image[1] or 0) > 0
            for page in doc
            for image in page.get_images(full=True)
        )
        if len(candidate["pages"]) < image_threshold:
            continue
        if sample_pages == 1 and not has_soft_mask:
            continue
        try:
            for page in doc:
                page.delete_image(xref)
            removed += int(candidate["count"])
        except Exception:
            logger.debug("Could not remove repeated PDF image xref=%s", xref, exc_info=True)

    buf = BytesIO()
    doc.save(buf, garbage=4, deflate=True)
    doc.close()
    return buf.getvalue(), removed, "structural" if removed else "none"


# ─────────────────────────────────────────────────────────────────────────────
# Internal helpers
# ─────────────────────────────────────────────────────────────────────────────

_SUBTYPE_MAP = {
    "highlight":     "/Highlight",
    "underline":     "/Underline",
    "strikethrough": "/StrikeOut",
    "note":          "/FreeText",
    "draw":          "/Ink",
    "rect":          "/Square",
    "ellipse":       "/Circle",
    "line":          "/Line",
    "arrow":         "/Line",
    "textbox":       "/FreeText",
}

_SUBTYPE_REVERSE: dict[str, str] = {v: k for k, v in _SUBTYPE_MAP.items()}


def _hex_to_pdf_color(hex_color: str) -> list[float]:
    """Convert '#RRGGBB' to [r, g, b] in [0.0, 1.0] range."""
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return [1.0, 1.0, 0.0]  # yellow fallback
    r = int(h[0:2], 16) / 255.0
    g = int(h[2:4], 16) / 255.0
    b = int(h[4:6], 16) / 255.0
    return [r, g, b]


def _pdf_color_to_hex(color_array: Any) -> str:
    """Convert a PDF color array to '#RRGGBB'."""
    try:
        vals = [float(v) for v in color_array]
        if len(vals) == 3:
            r, g, b = [int(v * 255) for v in vals]
            return "#{:02X}{:02X}{:02X}".format(r, g, b)
    except Exception:
        pass
    return "#FFFF00"


def _page_dimensions(page: Any) -> tuple[float, float]:
    """Return page width/height in PDF points."""
    try:
        box = page.mediabox
        return float(box.width), float(box.height)
    except Exception:
        return 612.0, 792.0


def _frontend_rect_to_pdf(rect: Any, annot: dict, page: Any) -> list[float] | None:
    """Convert a frontend top-left CSS rect to a PDF bottom-left point rect."""
    if isinstance(rect, dict):
        try:
            x = float(rect.get("x", 0))
            y = float(rect.get("y", 0))
            w = float(rect.get("w", rect.get("width", 0)))
            h = float(rect.get("h", rect.get("height", 0)))
        except Exception:
            return None
        if w <= 0 or h <= 0:
            return None
        pdf_w, pdf_h = _page_dimensions(page)
        css_w = float(annot.get("pageWidth") or annot.get("page_width") or pdf_w)
        css_h = float(annot.get("pageHeight") or annot.get("page_height") or pdf_h)
        sx = pdf_w / css_w if css_w > 0 else 1.0
        sy = pdf_h / css_h if css_h > 0 else 1.0
        x1 = x * sx
        x2 = (x + w) * sx
        y1 = pdf_h - ((y + h) * sy)
        y2 = pdf_h - (y * sy)
        return [x1, y1, x2, y2]
    if isinstance(rect, (list, tuple)) and len(rect) >= 4:
        try:
            return [float(rect[0]), float(rect[1]), float(rect[2]), float(rect[3])]
        except Exception:
            return None
    return None


def _normalize_annotation_rects(annot: dict, page: Any) -> list[list[float]]:
    """Normalize supported frontend/PDF rect formats to PDF point rectangles."""
    rects = annot.get("rects", [])
    if not isinstance(rects, list):
        return []
    normalized: list[list[float]] = []
    for rect in rects:
        converted = _frontend_rect_to_pdf(rect, annot, page)
        if converted is not None:
            normalized.append(converted)
    return normalized


def _build_annot_object(annot: dict, page: Any, writer: Any) -> Any | None:
    """Build a pypdf DictionaryObject for the given annotation dict."""
    try:
        from pypdf.generic import (  # type: ignore
            ArrayObject, DictionaryObject, FloatObject, NameObject,
            NumberObject, RectangleObject, TextStringObject,
        )
    except ImportError:
        return None

    annot_type = annot.get("type", "highlight")
    subtype = _SUBTYPE_MAP.get(annot_type)
    if not subtype:
        logger.warning(f"[PdfAnnotator] 未知批注类型: {annot_type}")
        return None

    color_vals = _hex_to_pdf_color(annot.get("color", "#FFFF00"))
    pdf_color = ArrayObject([FloatObject(v) for v in color_vals])

    # Bounding rect (union of all rects, or first rect)
    rects: list = _normalize_annotation_rects(annot, page)
    if rects:
        all_x1 = [r[0] for r in rects]
        all_y1 = [r[1] for r in rects]
        all_x2 = [r[2] for r in rects]
        all_y2 = [r[3] for r in rects]
        bbox = [min(all_x1), min(all_y1), max(all_x2), max(all_y2)]
    else:
        bbox = [0, 0, 100, 20]

    obj = DictionaryObject()
    obj[NameObject("/Type")] = NameObject("/Annot")
    obj[NameObject("/Subtype")] = NameObject(subtype)
    obj[NameObject("/Rect")] = ArrayObject([FloatObject(v) for v in bbox])
    obj[NameObject("/C")] = pdf_color
    obj[NameObject("/F")] = NumberObject(4)  # Print flag

    if annot.get("content"):
        obj[NameObject("/Contents")] = TextStringObject(annot["content"])

    # QuadPoints for highlight / underline / strikethrough
    if annot_type in ("highlight", "underline", "strikethrough") and rects:
        quad_points: list[float] = []
        for r in rects:
            x1, y1, x2, y2 = r[0], r[1], r[2], r[3]
            # PDF QuadPoints order: bottom-left, bottom-right, top-left, top-right
            quad_points += [x1, y1, x2, y1, x1, y2, x2, y2]
        obj[NameObject("/QuadPoints")] = ArrayObject(
            [FloatObject(v) for v in quad_points]
        )

    # InkList for draw annotations
    if annot_type == "draw":
        ink_list_raw: list = annot.get("inkList", [])
        ink_array = ArrayObject()
        for stroke in ink_list_raw:
            stroke_pts: list[float] = []
            for pt in stroke:
                stroke_pts += [float(pt[0]), float(pt[1])]
            ink_array.append(ArrayObject([FloatObject(v) for v in stroke_pts]))
        obj[NameObject("/InkList")] = ink_array

    # Shape annotations: override /Rect and add shape-specific keys
    if annot_type == "rect":
        x, y, w, h = (annot.get("x", 0), annot.get("y", 0),
                      annot.get("w", 50), annot.get("h", 30))
        obj[NameObject("/Rect")] = ArrayObject(
            [FloatObject(x), FloatObject(y), FloatObject(x + w), FloatObject(y + h)])
    elif annot_type == "ellipse":
        cx, cy = annot.get("cx", 0), annot.get("cy", 0)
        rx, ry = annot.get("rx", 25), annot.get("ry", 15)
        obj[NameObject("/Rect")] = ArrayObject(
            [FloatObject(cx - rx), FloatObject(cy - ry), FloatObject(cx + rx), FloatObject(cy + ry)])
    elif annot_type in ("line", "arrow"):
        x1, y1 = annot.get("x1", 0), annot.get("y1", 0)
        x2, y2 = annot.get("x2", 100), annot.get("y2", 0)
        obj[NameObject("/Rect")] = ArrayObject([
            FloatObject(min(x1, x2)), FloatObject(min(y1, y2)),
            FloatObject(max(x1, x2)), FloatObject(max(y1, y2))])
        obj[NameObject("/L")] = ArrayObject(
            [FloatObject(x1), FloatObject(y1), FloatObject(x2), FloatObject(y2)])
        if annot_type == "arrow":
            obj[NameObject("/LE")] = ArrayObject(
                [NameObject("/None"), NameObject("/OpenArrow")])
    elif annot_type == "textbox":
        x, y, w, h = (annot.get("x", 0), annot.get("y", 0),
                      annot.get("w", 120), annot.get("h", 30))
        obj[NameObject("/Rect")] = ArrayObject(
            [FloatObject(x), FloatObject(y), FloatObject(x + w), FloatObject(y + h)])
        text = annot.get("text", "")
        if text:
            obj[NameObject("/Contents")] = TextStringObject(text)
        font_size = annot.get("fontSize", 14)
        r, g, b = color_vals
        obj[NameObject("/DA")] = TextStringObject(f"/Helvetica {font_size} Tf {r:.2f} {g:.2f} {b:.2f} rg")

    # Border style (line width) for shape and draw annotations
    if annot_type in ("rect", "ellipse", "line", "arrow", "draw"):
        lw = float(annot.get("lineWidth", 2))
        obj[NameObject("/BS")] = DictionaryObject({
            NameObject("/W"): FloatObject(lw),
            NameObject("/S"): NameObject("/S"),
        })

    return obj


def _parse_annot_object(annot: Any, page_num: int) -> dict | None:
    """Convert a pypdf annotation DictionaryObject back to a frontend-compatible dict."""
    try:
        subtype = str(annot.get("/Subtype", ""))
        annot_type = _SUBTYPE_REVERSE.get(subtype)
        if not annot_type:
            return None

        rect_raw = annot.get("/Rect")
        rect = [float(v) for v in rect_raw] if rect_raw else [0, 0, 0, 0]

        color_raw = annot.get("/C")
        color = _pdf_color_to_hex(color_raw) if color_raw else "#FFFF00"
        content = str(annot.get("/Contents", "") or "")

        result: dict[str, Any] = {
            "id": f"pdf-{page_num}-{int(rect[0])}-{int(rect[1])}",
            "type": annot_type,
            "page": page_num,
            "color": color,
            "content": content,
            "rects": [rect],
        }

        # Recover individual rects from QuadPoints
        qp_raw = annot.get("/QuadPoints")
        if qp_raw:
            qp = [float(v) for v in qp_raw]
            rects: list[list[float]] = []
            for i in range(0, len(qp), 8):
                if i + 7 < len(qp):
                    xs = [qp[i], qp[i+2], qp[i+4], qp[i+6]]
                    ys = [qp[i+1], qp[i+3], qp[i+5], qp[i+7]]
                    rects.append([min(xs), min(ys), max(xs), max(ys)])
            if rects:
                result["rects"] = rects

        # Reconstruct ink strokes
        if annot_type == "draw":
            il_raw = annot.get("/InkList")
            if il_raw:
                ink_list: list[list[list[float]]] = []
                for stroke in il_raw:
                    pts = [float(v) for v in stroke]
                    ink_list.append([[pts[j], pts[j+1]] for j in range(0, len(pts), 2)])
                result["inkList"] = ink_list

        return result
    except Exception as e:
        logger.debug(f"[PdfAnnotator] _parse_annot_object 失败: {e}")
        return None
