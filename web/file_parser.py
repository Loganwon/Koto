#!/usr/bin/env python
# -*- coding: utf-8 -*-
# Copyright (C) 2024-2026 Koto AI. All rights reserved.
# SPDX-License-Identifier: LicenseRef-Koto-Proprietary
"""
文件解析器 - 支持 PDF/DOCX/TXT/MD 文件读取
用于 PPT 生成的多源文件融合
"""

import logging
import os
from pathlib import Path
from typing import Dict, List, Optional

logger = logging.getLogger(__name__)


class FileParser:
    """多格式文件解析器"""

    SUPPORTED_FORMATS = [
        ".pdf",
        ".docx",
        ".doc",
        ".txt",
        ".md",
        ".markdown",
        ".pptx",
        ".pptm",
        ".ppt",
        ".xlsx",
        ".xls",
        ".csv",
    ]
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB 上限
    MAX_CONTENT_LENGTH = 100000  # 提取最多 10 万字符

    @staticmethod
    def parse_file(file_path: str) -> Optional[Dict[str, any]]:
        """
        解析单个文件

        Args:
            file_path: 文件路径

        Returns:
            {
                "success": bool,
                "filename": str,
                "format": str,
                "content": str,  # 提取的文本内容
                "char_count": int,
                "error": str (if failed)
            }
        """
        if not os.path.exists(file_path):
            return {"success": False, "error": f"文件不存在: {file_path}"}

        file_size = os.path.getsize(file_path)
        if file_size > FileParser.MAX_FILE_SIZE:
            return {
                "success": False,
                "error": f"文件过大 ({file_size/1024/1024:.1f}MB > 50MB)",
            }

        file_ext = Path(file_path).suffix.lower()
        filename = os.path.basename(file_path)

        if file_ext not in FileParser.SUPPORTED_FORMATS:
            return {"success": False, "error": f"不支持的格式: {file_ext}"}

        try:
            if file_ext == ".pdf":
                content = FileParser._parse_pdf(file_path)
            elif file_ext in [".docx", ".doc"]:
                content = FileParser._parse_docx(file_path)
            elif file_ext in [".pptx", ".pptm", ".ppt"]:
                content = FileParser._parse_pptx(file_path)
            elif file_ext in [".xlsx", ".xls"]:
                content = FileParser._parse_xlsx(file_path)
            elif file_ext == ".csv":
                content = FileParser._parse_csv(file_path)
            elif file_ext in [".txt", ".md", ".markdown"]:
                content = FileParser._parse_text(file_path)
            else:
                return {"success": False, "error": "未知格式"}

            # 截断超长内容
            if len(content) > FileParser.MAX_CONTENT_LENGTH:
                content = content[: FileParser.MAX_CONTENT_LENGTH] + "\n\n[内容已截断]"

            return {
                "success": True,
                "filename": filename,
                "format": file_ext.lstrip("."),
                "content": content,
                "char_count": len(content),
            }

        except Exception as e:
            return {
                "success": False,
                "filename": filename,
                "error": f"解析失败: {str(e)}",
            }

    @staticmethod
    def _parse_pdf(file_path: str) -> str:
        """PDF 文本提取（尝试 pypdf → PyPDF2 → pdfplumber）"""
        content = []

        # 优先: pypdf（PyPDF2 的继任者，纯Python，已作为 pypdf 包发布）
        for pkg_name, mod_name in [("pypdf", "pypdf"), ("PyPDF2", "PyPDF2")]:
            try:
                mod = __import__(mod_name)
                PdfReader = getattr(mod, "PdfReader")
                with open(file_path, "rb") as f:
                    reader = PdfReader(f)
                    for page_num, page in enumerate(reader.pages):
                        text = page.extract_text() or ""
                        if text.strip():
                            content.append(f"[第 {page_num + 1} 页]\n{text}")
                if content:
                    return "\n\n".join(content)
            except ImportError:
                continue
            except Exception:
                break  # 库可用但解析失败，尝试 pdfplumber

        # 回退: pdfplumber
        try:
            import pdfplumber

            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        content.append(f"[第 {page_num + 1} 页]\n{text}")
            return "\n\n".join(content)
        except ImportError:
            raise ImportError("需要安装 pypdf 或 pdfplumber: pip install pypdf")

    @staticmethod
    def _parse_docx(file_path: str) -> str:
        """DOCX 文本提取"""
        try:
            from docx import Document
        except ImportError:
            raise ImportError("需要安装 python-docx: pip install python-docx")

        doc = Document(file_path)
        content = []

        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text)

        # 也提取表格，且带上有格式的Markdown
        for table in doc.tables:
            table_content = []
            for i, row in enumerate(table.rows):
                row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                table_content.append("| " + " | ".join(row_data) + " |")
                if i == 0:
                    table_content.append("|" + "|".join(["---" for _ in row_data]) + "|")
            if table_content:
                content.append("\n" + "\n".join(table_content))

        return "\n".join(content)

    @staticmethod
    def _parse_pptx(file_path: str) -> str:
        """PPTX/PPTM/PPT 文本提取（逐页提取标题、正文和备注）"""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("需要安装 python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        content = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "has_table") and shape.has_table:
                    for row_idx, row in enumerate(shape.table.rows):
                        row_data = [cell.text_frame.text.replace("\n", " ").strip() for cell in row.cells]
                        slide_texts.append("| " + " | ".join(row_data) + " |")
                        if row_idx == 0:
                            slide_texts.append("|" + "|".join(["---" for _ in row_data]) + "|")
                            
                if not hasattr(shape, "text"):
                    continue
                text = shape.text.strip()
                if not text:
                    continue
                # 收集文本框内每一段，保留层次
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        pt = para.text.strip()
                        if pt:
                            slide_texts.append(pt)
                else:
                    slide_texts.append(text)
            if slide_texts:
                content.append(f"[第 {idx} 页]\n" + "\n".join(slide_texts))
            # 提取演讲者备注
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf:
                    notes = notes_tf.text.strip()
                    if notes:
                        content.append(f"[第 {idx} 页·备注]\n{notes}")
        return "\n\n".join(content)

    @staticmethod
    def build_pptx_slides(file_path: str) -> dict | None:
        """
        Parse a PPTX file into structured slide data for the visual canvas editor.

        Returns a dict with:
          slide_width_emu, slide_height_emu – original EMU dimensions
          slides – list of slides, each containing:
            background – hex colour string (default #FFFFFF)
            shapes – list of shape objects:
              id, name, type, left/top/width/height (EMU), fill, z_order,
              has_text, paragraphs (list of paragraph objects with runs)
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return None

        import html as _html

        try:
            prs = Presentation(file_path)
        except Exception:
            return None

        slide_w = prs.slide_width or 9144000
        slide_h = prs.slide_height or 6858000

        slides_data = []
        for slide_idx, slide in enumerate(prs.slides):
            # Background fill
            bg_hex = "#FFFFFF"
            try:
                bg_fill = slide.background.fill
                if bg_fill.type is not None and str(bg_fill.type) in ("SOLID", "1"):
                    try:
                        bg_hex = "#{:06x}".format(int(bg_fill.fore_color.rgb))
                    except Exception:
                        pass
            except Exception:
                pass

            shapes_data = []
            for z_idx, shape in enumerate(slide.shapes):
                s: dict = {
                    "id": shape.shape_id,
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "left": shape.left or 0,
                    "top": shape.top or 0,
                    "width": shape.width or 0,
                    "height": shape.height or 0,
                    "z_order": z_idx,
                    "has_text": False,
                    "fill": None,
                }

                # Shape fill colour
                try:
                    fill = shape.fill
                    if fill.type is not None and str(fill.type) in ("SOLID", "1"):
                        s["fill"] = "#{:06x}".format(int(fill.fore_color.rgb))
                except Exception:
                    pass

                # Text frame
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    s["has_text"] = True
                    paras = []
                    for para in shape.text_frame.paragraphs:
                        align_name = "LEFT"
                        try:
                            if para.alignment:
                                align_name = para.alignment.name
                        except Exception:
                            pass
                        p_obj: dict = {"align": align_name, "runs": []}
                        for run in para.runs:
                            r: dict = {"text": run.text}
                            try:
                                if run.font.size:
                                    r["size"] = round(run.font.size.pt, 1)
                            except Exception:
                                pass
                            try:
                                if run.font.bold:
                                    r["bold"] = True
                            except Exception:
                                pass
                            try:
                                if run.font.italic:
                                    r["italic"] = True
                            except Exception:
                                pass
                            try:
                                if run.font.underline:
                                    r["underline"] = True
                            except Exception:
                                pass
                            try:
                                if run.font.color and run.font.color.type is not None:
                                    r["color"] = "#{:06x}".format(int(run.font.color.rgb))
                            except Exception:
                                pass
                            p_obj["runs"].append(r)
                        paras.append(p_obj)
                    s["paragraphs"] = paras

                shapes_data.append(s)

            slides_data.append({
                "background": bg_hex,
                "shapes": shapes_data,
            })

        return {
            "slide_width_emu": slide_w,
            "slide_height_emu": slide_h,
            "slides": slides_data,
        }

    @staticmethod
    def _parse_pptx_html(file_path: str) -> str:
        """PPTX → rich HTML 渲染（幻灯片卡片式展示）"""
        try:
            from pptx import Presentation
        except ImportError:
            return ""

        import html as html_mod

        prs = Presentation(file_path)
        parts = []
        for idx, slide in enumerate(prs.slides, 1):
            parts.append(f'<div class="wa-slide-card">')
            parts.append(f'<div class="wa-slide-header">幻灯片 {idx}</div>')
            parts.append('<div class="wa-slide-body">')

            for shape in slide.shapes:
                # Tables
                if hasattr(shape, "has_table") and shape.has_table:
                    parts.append('<table class="wa-doc-table">')
                    for ri, row in enumerate(shape.table.rows):
                        tag = "th" if ri == 0 else "td"
                        cells = "".join(
                            f"<{tag}>{html_mod.escape(cell.text_frame.text.strip())}</{tag}>"
                            for cell in row.cells
                        )
                        parts.append(f"<tr>{cells}</tr>")
                    parts.append("</table>")
                    continue

                if not hasattr(shape, "text_frame"):
                    continue
                tf = shape.text_frame
                for para in tf.paragraphs:
                    pt = para.text.strip()
                    if not pt:
                        continue
                    lvl = para.level or 0
                    # Inline formatting
                    spans = []
                    for run in para.runs:
                        t = html_mod.escape(run.text)
                        if not t:
                            continue
                        if run.font.bold:
                            t = f"<strong>{t}</strong>"
                        if run.font.italic:
                            t = f"<em>{t}</em>"
                        spans.append(t)
                    text_html = "".join(spans) if spans else html_mod.escape(pt)
                    if lvl == 0 and para == tf.paragraphs[0]:
                        parts.append(f"<h4>{text_html}</h4>")
                    elif lvl > 0:
                        parts.append(f'<p style="margin-left:{lvl * 20}px;">• {text_html}</p>')
                    else:
                        parts.append(f"<p>{text_html}</p>")

            parts.append("</div>")  # slide-body

            # Speaker notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf and notes_tf.text.strip():
                    parts.append(
                        f'<div class="wa-slide-notes">'
                        f'<span class="wa-slide-notes-label">备注</span> '
                        f'{html_mod.escape(notes_tf.text.strip())}'
                        f'</div>'
                    )

            parts.append("</div>")  # slide-card
        return "\n".join(parts)

    @staticmethod
    def _parse_xlsx(file_path: str) -> str:
        """XLSX/XLS 文本提取（以 Markdown 表格格式输出）"""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("需要安装 openpyxl: pip install openpyxl")

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        content = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(max_row=500, values_only=True))
            if not rows:
                continue
            content.append(f"[工作表: {sheet_name}]")
            for i, row in enumerate(rows):
                cells = [str(c).strip() if c is not None else "" for c in row]
                content.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    content.append("|" + "|".join(["---" for _ in cells]) + "|")
            content.append("")
        wb.close()
        return "\n".join(content)

    @staticmethod
    def _parse_xlsx_html(file_path: str) -> str:
        """XLSX → HTML table 渲染（供前端直接展示）"""
        try:
            import openpyxl
        except ImportError:
            return ""

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        html_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(max_row=500, values_only=True))
            if not rows:
                continue
            html_parts.append(f'<h3 class="wa-sheet-title">{sheet_name}</h3>')
            html_parts.append('<div class="wa-table-wrap"><table class="wa-spreadsheet">')
            for i, row in enumerate(rows):
                tag = "th" if i == 0 else "td"
                cells = "".join(
                    f"<{tag}>{str(c) if c is not None else ''}</{tag}>"
                    for c in row
                )
                html_parts.append(f"<tr>{cells}</tr>")
            html_parts.append("</table></div>")
        wb.close()
        return "\n".join(html_parts)

    @staticmethod
    def _parse_csv(file_path: str) -> str:
        """CSV 文本提取"""
        import csv

        content = []
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= 500:
                    content.append("[... 截断至 500 行]")
                    break
                cells = [c.strip() for c in row]
                content.append("| " + " | ".join(cells) + " |")
                if i == 0:
                    content.append("|" + "|".join(["---" for _ in cells]) + "|")
        return "\n".join(content)

    @staticmethod
    def _parse_csv_html(file_path: str) -> str:
        """CSV → HTML table 渲染"""
        import csv

        html_parts = ['<div class="wa-table-wrap"><table class="wa-spreadsheet">']
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= 500:
                    break
                tag = "th" if i == 0 else "td"
                cells = "".join(f"<{tag}>{c}</{tag}>" for c in row)
                html_parts.append(f"<tr>{cells}</tr>")
        html_parts.append("</table></div>")
        return "\n".join(html_parts)

    @staticmethod
    def _parse_docx_html(file_path: str) -> str:
        """DOCX → rich HTML 渲染（保留标题/加粗/斜体/表格格式，按文档顺序交错）"""
        try:
            from docx import Document
            from docx.oxml.ns import qn
        except ImportError:
            return ""

        import html as html_mod

        doc = Document(file_path)
        parts = []

        # Walk through the document body XML children in order
        # to correctly interleave paragraphs and tables
        for child in doc.element.body:
            if child.tag == qn('w:p'):
                # Paragraph element — find the matching Paragraph object
                from docx.text.paragraph import Paragraph
                para = Paragraph(child, doc)
                if not para.text.strip():
                    continue
                style_name = (para.style.name or "").lower()
                if "heading 1" in style_name:
                    parts.append(f"<h1>{html_mod.escape(para.text)}</h1>")
                elif "heading 2" in style_name:
                    parts.append(f"<h2>{html_mod.escape(para.text)}</h2>")
                elif "heading 3" in style_name:
                    parts.append(f"<h3>{html_mod.escape(para.text)}</h3>")
                elif "heading" in style_name:
                    parts.append(f"<h4>{html_mod.escape(para.text)}</h4>")
                else:
                    spans = []
                    for run in para.runs:
                        text = html_mod.escape(run.text)
                        if not text:
                            continue
                        if run.bold:
                            text = f"<strong>{text}</strong>"
                        if run.italic:
                            text = f"<em>{text}</em>"
                        if run.underline:
                            text = f"<u>{text}</u>"
                        spans.append(text)
                    if spans:
                        parts.append(f"<p>{''.join(spans)}</p>")

            elif child.tag == qn('w:tbl'):
                # Table element — find the matching Table object
                from docx.table import Table as DocxTable
                table = DocxTable(child, doc)
                parts.append('<table class="wa-doc-table">')
                for i, row in enumerate(table.rows):
                    tag = "th" if i == 0 else "td"
                    cells = "".join(
                        f"<{tag}>{html_mod.escape(cell.text.strip())}</{tag}>"
                        for cell in row.cells
                    )
                    parts.append(f"<tr>{cells}</tr>")
                parts.append("</table>")

        return "\n".join(parts)

    # ─── Edit Model Builders ────────────────────────────────────────────────

    @staticmethod
    def build_edit_model(file_path: str) -> Optional[Dict]:
        """
        Build a structured edit model for in-browser editing.
        Returns { format, capabilities, blocks[], warnings[] } or None
        if format is not editable.
        """
        suffix = Path(file_path).suffix.lower()
        if suffix in (".docx", ".doc"):
            return FileParser._build_docx_edit_model(file_path)
        # Future: .pptx, .xlsx
        return None

    @staticmethod
    def _build_docx_edit_model(file_path: str) -> Optional[Dict]:
        """Walk DOCX body XML children → structured JSON blocks for editing."""
        try:
            from docx import Document
            from docx.oxml.ns import qn
        except ImportError:
            return None

        doc = Document(file_path)
        blocks = []
        idx = 0

        for child in doc.element.body:
            if child.tag == qn("w:p"):
                from docx.text.paragraph import Paragraph

                para = Paragraph(child, doc)
                style_name = (para.style.name or "").lower()
                if "heading" in style_name:
                    level = 1
                    for lv in ("heading 1", "heading 2", "heading 3", "heading 4"):
                        if lv in style_name:
                            level = int(lv[-1])
                            break
                    blocks.append(
                        {
                            "idx": idx,
                            "type": "heading",
                            "level": level,
                            "text": para.text,
                        }
                    )
                else:
                    runs = []
                    for run in para.runs:
                        runs.append(
                            {
                                "text": run.text,
                                "bold": bool(run.bold),
                                "italic": bool(run.italic),
                                "underline": bool(run.underline),
                            }
                        )
                    blocks.append(
                        {"idx": idx, "type": "paragraph", "text": para.text, "runs": runs}
                    )
                idx += 1

            elif child.tag == qn("w:tbl"):
                from docx.table import Table as DocxTable

                table = DocxTable(child, doc)
                rows = []
                for row in table.rows:
                    rows.append([cell.text for cell in row.cells])
                blocks.append({"idx": idx, "type": "table", "rows": rows})
                idx += 1

        return {
            "format": "docx",
            "capabilities": ["edit_text", "edit_table"],
            "blocks": blocks,
            "warnings": [],
        }

    # ─── Univer Snapshot Builders ───────────────────────────────────────────

    @staticmethod
    def build_univer_doc_snapshot(file_path: str) -> Optional[Dict]:
        """
        Convert DOCX → Univer IDocumentData snapshot.
        Builds dataStream, textRuns, and paragraphs arrays that Univer Docs
        expects inside body.
        """
        try:
            from docx import Document
            from docx.oxml.ns import qn
        except ImportError:
            return None

        doc = Document(file_path)

        data_stream = ""
        text_runs = []       # {st, ed, ts}
        paragraphs = []      # {startIndex}
        cursor = 0           # current char offset in dataStream

        for child in doc.element.body:
            if child.tag == qn("w:p"):
                from docx.text.paragraph import Paragraph
                para = Paragraph(child, doc)
                para_text = para.text or ""

                paragraphs.append({"startIndex": cursor})

                for run in para.runs:
                    run_text = run.text or ""
                    if not run_text:
                        continue
                    st = cursor
                    cursor += len(run_text)
                    ed = cursor

                    ts = {}
                    if run.bold:
                        ts["bl"] = 1
                    if run.italic:
                        ts["it"] = 1
                    if run.underline:
                        ts["ul"] = {"s": 1}
                    if run.font and run.font.size:
                        # python-docx font.size is in EMU; convert to pt
                        ts["fs"] = round(run.font.size.pt)

                    text_runs.append({"st": st, "ed": ed, "ts": ts})

                # Paragraph separator (\n in Univer dataStream)
                data_stream += para_text + "\n"
                cursor = len(data_stream)

            elif child.tag == qn("w:tbl"):
                # Tables are flattened as text paragraphs for Univer Docs
                from docx.table import Table as DocxTable
                table = DocxTable(child, doc)
                for row in table.rows:
                    row_text = "\t".join(cell.text.strip() for cell in row.cells)
                    paragraphs.append({"startIndex": cursor})
                    data_stream += row_text + "\n"
                    cursor = len(data_stream)

        # Final section break required by Univer
        data_stream += "\x00"

        return {
            "id": "doc-snapshot",
            "body": {
                "dataStream": data_stream,
                "textRuns": text_runs,
                "paragraphs": paragraphs,
            },
            "documentStyle": {
                "pageSize": {"width": 595.28, "height": 841.89},
                "marginTop": 72,
                "marginBottom": 72,
                "marginLeft": 90,
                "marginRight": 90,
            },
        }

    @staticmethod
    def build_univer_sheet_snapshot(file_path: str) -> Optional[Dict]:
        """
        Convert XLSX → Univer IWorkbookData snapshot.
        Builds sheets with cellData in {row: {col: {v}}} format.
        """
        suffix = Path(file_path).suffix.lower()

        if suffix == ".csv":
            return FileParser._build_univer_sheet_from_csv(file_path)

        try:
            import openpyxl
        except ImportError:
            return None

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        sheets = {}
        sheet_order = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            cell_data = {}
            max_row = 0
            max_col = 0

            for ri, row in enumerate(ws.iter_rows(max_row=2000, values_only=False)):
                row_data = {}
                for ci, cell in enumerate(row):
                    if cell.value is not None:
                        cv = {"v": cell.value}
                        # Preserve number format if set
                        if cell.number_format and cell.number_format != "General":
                            cv["s"] = {"n": {"pattern": cell.number_format}}
                        row_data[ci] = cv
                        if ci > max_col:
                            max_col = ci
                if row_data:
                    cell_data[ri] = row_data
                    if ri > max_row:
                        max_row = ri

            sheet_id = f"sheet_{abs(hash(sheet_name)) % 10**8}"
            sheets[sheet_id] = {
                "id": sheet_id,
                "name": sheet_name,
                "cellData": cell_data,
                "rowCount": max(max_row + 50, 100),
                "columnCount": max(max_col + 10, 26),
            }
            sheet_order.append(sheet_id)

        wb.close()

        return {
            "id": "workbook-snapshot",
            "sheets": sheets,
            "sheetOrder": sheet_order,
        }

    @staticmethod
    def _build_univer_sheet_from_csv(file_path: str) -> Optional[Dict]:
        """Convert CSV → Univer IWorkbookData snapshot."""
        import csv

        cell_data = {}
        max_col = 0

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for ri, row in enumerate(reader):
                if ri >= 2000:
                    break
                row_data = {}
                for ci, val in enumerate(row):
                    val = val.strip()
                    if val:
                        # Try to preserve numbers
                        try:
                            num = float(val)
                            row_data[ci] = {"v": num}
                        except ValueError:
                            row_data[ci] = {"v": val}
                        if ci > max_col:
                            max_col = ci
                if row_data:
                    cell_data[ri] = row_data

        sheet_id = "sheet_csv"
        return {
            "id": "workbook-snapshot",
            "sheets": {
                sheet_id: {
                    "id": sheet_id,
                    "name": "Sheet1",
                    "cellData": cell_data,
                    "rowCount": max(len(cell_data) + 50, 100),
                    "columnCount": max(max_col + 10, 26),
                }
            },
            "sheetOrder": [sheet_id],
        }

    # ── Reverse compilers: Univer snapshot → file ─────────────────────────

    @staticmethod
    def build_docx_from_univer_snapshot(snapshot: Dict) -> Optional[str]:
        """
        Reverse compile a Univer IDocumentData snapshot back to a .docx file.
        Returns the path to a temporary .docx file, or None on failure.
        """
        try:
            from docx import Document
            from docx.shared import Pt
        except ImportError:
            logger.warning("python-docx not available for reverse compile")
            return None

        try:
            body = snapshot.get("body") or {}
            data_stream = body.get("dataStream", "")
            text_runs = body.get("textRuns") or []
            # paragraphs list not strictly needed — we split dataStream by \n

            # Remove trailing \x00 section break if present
            if data_stream.endswith("\x00"):
                data_stream = data_stream[:-1]

            # Build a sorted list of formatting spans for fast lookup
            # Each textRun: {st, ed, ts: {bl, it, ul, fs, ...}}
            text_runs_sorted = sorted(text_runs, key=lambda r: r.get("st", 0))

            doc = Document()

            # Split by paragraph separator
            para_texts = data_stream.split("\n")

            for para_text in para_texts:
                if para_text == "" and para_texts[-1] == para_text:
                    # Skip trailing empty paragraph from split
                    break
                para = doc.add_paragraph()
                if not para_text:
                    # Empty paragraph
                    continue

                # Calculate the absolute offset of this paragraph in dataStream
                # We need to track cumulative offset as we go
                pass  # handled below

            # Re-do with proper offset tracking
            for p in doc.paragraphs:
                # Remove the paragraphs we just added — we'll redo properly
                p._element.getparent().remove(p._element)

            cursor = 0
            for para_text in para_texts:
                if not para_text and cursor >= len(data_stream):
                    break

                para_start = cursor
                para_end = cursor + len(para_text)

                para = doc.add_paragraph()

                if not para_text:
                    cursor = para_end + 1  # +1 for \n
                    continue

                # Collect text_runs that overlap [para_start, para_end)
                overlapping = []
                for tr in text_runs_sorted:
                    tr_st = tr.get("st", 0)
                    tr_ed = tr.get("ed", 0)
                    if tr_ed <= para_start:
                        continue
                    if tr_st >= para_end:
                        break
                    overlapping.append(tr)

                if not overlapping:
                    # No formatting — single plain run
                    para.add_run(para_text)
                else:
                    # Build runs by splitting text according to formatting spans
                    pos = para_start
                    for tr in overlapping:
                        tr_st = max(tr.get("st", 0), para_start)
                        tr_ed = min(tr.get("ed", 0), para_end)

                        # Gap before this formatted run
                        if tr_st > pos:
                            gap_text = data_stream[pos:tr_st]
                            if gap_text:
                                para.add_run(gap_text)

                        # Formatted run
                        run_text = data_stream[tr_st:tr_ed]
                        if run_text:
                            run = para.add_run(run_text)
                            ts = tr.get("ts") or {}
                            if ts.get("bl"):
                                run.bold = True
                            if ts.get("it"):
                                run.italic = True
                            if ts.get("ul"):
                                run.underline = True
                            fs = ts.get("fs")
                            if fs and isinstance(fs, (int, float)) and fs > 0:
                                run.font.size = Pt(fs)

                        pos = tr_ed

                    # Trailing unformatted text
                    if pos < para_end:
                        trailing = data_stream[pos:para_end]
                        if trailing:
                            para.add_run(trailing)

                cursor = para_end + 1  # +1 for \n separator

            import tempfile
            tmp = tempfile.NamedTemporaryFile(
                suffix=".docx", delete=False, prefix="koto_univer_"
            )
            tmp.close()
            doc.save(tmp.name)
            return tmp.name

        except Exception:
            logger.exception("build_docx_from_univer_snapshot failed")
            return None

    @staticmethod
    def build_xlsx_from_univer_snapshot(snapshot: Dict) -> Optional[str]:
        """
        Reverse compile a Univer IWorkbookData snapshot back to an .xlsx file.
        Returns the path to a temporary .xlsx file, or None on failure.
        """
        try:
            import openpyxl
        except ImportError:
            logger.warning("openpyxl not available for reverse compile")
            return None

        try:
            wb = openpyxl.Workbook()
            # Remove the default sheet — we'll create our own
            if wb.worksheets:
                wb.remove(wb.worksheets[0])

            sheets_data = snapshot.get("sheets") or {}
            sheet_order = snapshot.get("sheetOrder") or list(sheets_data.keys())

            for sheet_id in sheet_order:
                sheet_info = sheets_data.get(sheet_id)
                if not sheet_info:
                    continue
                ws = wb.create_sheet(title=sheet_info.get("name", "Sheet"))
                cell_data = sheet_info.get("cellData") or {}

                for row_key, row_dict in cell_data.items():
                    ri = int(row_key)
                    if not isinstance(row_dict, dict):
                        continue
                    for col_key, cell_obj in row_dict.items():
                        ci = int(col_key)
                        if not isinstance(cell_obj, dict):
                            continue
                        value = cell_obj.get("v")
                        if value is not None:
                            ws.cell(row=ri + 1, column=ci + 1, value=value)

            # Ensure at least one sheet
            if not wb.worksheets:
                wb.create_sheet(title="Sheet1")

            import tempfile
            tmp = tempfile.NamedTemporaryFile(
                suffix=".xlsx", delete=False, prefix="koto_univer_"
            )
            tmp.close()
            wb.save(tmp.name)
            return tmp.name

        except Exception:
            logger.exception("build_xlsx_from_univer_snapshot failed")
            return None

    @staticmethod
    def _parse_text(file_path: str) -> str:
        """纯文本/Markdown 读取"""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    @staticmethod
    def batch_parse(file_paths: List[str]) -> List[Dict]:
        """
        批量解析多个文件，融合为统一格式

        Returns:
            [
                {
                    "filename": str,
                    "format": str,
                    "content": str,
                    "success": bool,
                    "error": str (if failed)
                },
                ...
            ]
        """
        results = []
        for path in file_paths:
            result = FileParser.parse_file(path)
            results.append(result)
        return results

    @staticmethod
    def merge_contents(parse_results: List[Dict]) -> str:
        """
        将多个文件的内容合并为统一的参考材料格式

        Args:
            parse_results: batch_parse 的返回结果

        Returns:
            合并后的文本（带来源标记）
        """
        merged = []
        for i, result in enumerate(parse_results, 1):
            if result.get("success"):
                filename = result.get("filename", f"文件{i}")
                format_type = result.get("format", "unknown")
                content = result.get("content", "")

                merged.append(
                    f"【来源文件 {i}】{filename} ({format_type})\n"
                    f"{'=' * 60}\n"
                    f"{content}\n"
                    f"{'=' * 60}\n"
                )

        return "\n\n".join(merged)

    @staticmethod
    def sanitize_file_path(file_path: str) -> Optional[str]:
        """
        检查文件路径是否安全（防止路径遍历攻击）

        返回规范化的绝对路径，或 None if 不安全
        """
        try:
            abs_path = os.path.abspath(file_path)

            # 确保文件在允许的目录内
            allowed_dirs = [
                os.path.abspath(os.path.dirname(__file__)),  # web/ 目录
                os.path.abspath(
                    os.path.join(os.path.dirname(__file__), "..")
                ),  # 项目根目录
            ]

            if not any(abs_path.startswith(d) for d in allowed_dirs):
                return None

            return abs_path
        except (OSError, ValueError) as e:
            logger.debug("Path validation failed for %s: %s", file_path, e)
            return None
